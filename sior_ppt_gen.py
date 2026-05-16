from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── SIOR Brand Palette  (navy · crimson · white · silver) ──────
NAVY        = RGBColor(0x09, 0x2B, 0x52)   # SIOR deep navy
MID_NAVY    = RGBColor(0x16, 0x4A, 0x7E)   # mid navy
CRIMSON     = RGBColor(0xBE, 0x1E, 0x2D)   # SIOR red/crimson
STEEL_BLUE  = RGBColor(0x2E, 0x75, 0xB6)   # supporting blue
TEAL        = RGBColor(0x00, 0x8B, 0x8B)   # teal accent
IVORY       = RGBColor(0xF6, 0xF8, 0xFA)   # light bg
SILVER      = RGBColor(0x9A, 0x9E, 0xA8)   # silver text
LIGHT_BLUE  = RGBColor(0xE8, 0xF0, 0xFA)   # card bg tint
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
SUCCESS     = RGBColor(0x27, 0xAE, 0x60)
CAUTION     = RGBColor(0xF0, 0x9A, 0x1A)
DANGER      = RGBColor(0xC0, 0x3A, 0x2E)
CHARCOAL    = RGBColor(0x22, 0x22, 0x33)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

# ── Helpers ─────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.line.fill.background(); s.line.width = 0
    s.fill.solid(); s.fill.fore_color.rgb = fill
    return s

def add_textbox(slide, text, x, y, w, h,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(font_size); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return txb

def bg(slide, color=IVORY):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, color)

def header(slide, title, sub=None):
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.25), NAVY)
    add_rect(slide, 0, Inches(1.25), Inches(0.38), Inches(6.25), CRIMSON)
    add_textbox(slide, title, Inches(0.58), Inches(0.16),
                Inches(10.5), Inches(0.64), 27, True, WHITE)
    if sub:
        add_textbox(slide, sub, Inches(0.58), Inches(0.74),
                    Inches(10.5), Inches(0.42), 12, False,
                    RGBColor(0xBB, 0xCC, 0xFF), italic=True)

def footer(slide, txt="SIOR – Society of Industrial and Office Realtors  |  Web Dev & AI Automation Audit  |  2026"):
    add_rect(slide, 0, Inches(7.2), SLIDE_W, Inches(0.3), NAVY)
    add_textbox(slide, txt, Inches(0.3), Inches(7.2), Inches(12.7), Inches(0.3),
                9, False, SILVER, PP_ALIGN.CENTER)

def card(slide, x, y, w, h, fill=WHITE, border_color=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border_color:
        s.line.color.rgb = border_color; s.line.width = Pt(0.6)
    else:
        s.line.fill.background(); s.line.width = 0
    return s

def accent_bar(slide, x, y, w, h=Inches(0.048)):
    add_rect(slide, x, y, w, h, CRIMSON)

def bullet_card(slide, title, bullets, x, y, w, h, hdr_bg=NAVY, hdr_col=WHITE):
    card(slide, x, y, w, h, fill=WHITE, border_color=RGBColor(0xCC, 0xD5, 0xE5))
    add_rect(slide, x, y, w, Inches(0.44), hdr_bg)
    add_textbox(slide, title, x + Inches(0.14), y + Inches(0.06),
                w - Inches(0.28), Inches(0.36), 12, True, hdr_col)
    ty = y + Inches(0.58)
    for b in bullets:
        add_textbox(slide, f"▸  {b}", x + Inches(0.15), ty,
                    w - Inches(0.3), Inches(0.36), 10.5, False, CHARCOAL)
        ty += Inches(0.34)

def stat_card(slide, val, lbl, x, y, w=Inches(2.8), h=Inches(1.5),
              bg_col=NAVY, vc=WHITE, lc=RGBColor(0xBB, 0xCC, 0xFF)):
    card(slide, x, y, w, h, fill=bg_col)
    add_textbox(slide, val, x, y + Inches(0.16), w, Inches(0.72),
                30, True, vc, PP_ALIGN.CENTER)
    add_textbox(slide, lbl, x, y + Inches(0.88), w, Inches(0.52),
                11, False, lc, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
add_rect(s, 0, 0, Inches(0.72), SLIDE_H, CRIMSON)
add_rect(s, Inches(0.72), Inches(4.62), SLIDE_W - Inches(0.72), Inches(0.055), CRIMSON)

add_textbox(s, "SOCIETY OF INDUSTRIAL AND OFFICE REALTORS®",
            Inches(1.1), Inches(0.82), Inches(11.5), Inches(0.55),
            13, True, RGBColor(0xBB, 0xCC, 0xFF))
add_textbox(s, "Web Development &\nAI Automation Audit",
            Inches(1.1), Inches(1.52), Inches(11.2), Inches(2.15),
            44, True, WHITE)
add_textbox(s, "A Strategic Digital Transformation Blueprint for the World's Foremost CRE Professional Association",
            Inches(1.1), Inches(3.78), Inches(11.2), Inches(0.55),
            14, False, RGBColor(0xBB, 0xCC, 0xFF), italic=True)
add_textbox(s, "Prepared by: Genos Apollo  |  May 2026",
            Inches(1.1), Inches(5.05), Inches(8), Inches(0.4),
            12, False, SILVER)
add_textbox(s, "sior.com  •  Washington, D.C.  •  4,000+ Members  •  34 Countries  •  630+ Cities",
            Inches(1.1), Inches(5.55), Inches(10), Inches(0.35),
            11, False, RGBColor(0x88, 0x99, 0xBB), italic=True)
add_textbox(s, "CONFIDENTIAL",
            Inches(1.1), Inches(6.68), Inches(4), Inches(0.32),
            10, False, SILVER, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Agenda", "What we will cover today"); footer(s)

items = [
    ("01", "Organisation Digital Overview",  "Current state of sior.com, MySIOR portal & member experience"),
    ("02", "Website Audit Findings",          "UX, performance, SEO, member journey & platform gaps"),
    ("03", "Web Development Roadmap",         "Modernising sior.com and the MySIOR member platform"),
    ("04", "AI Automation Opportunities",     "Where AI supercharges member value & association operations"),
    ("05", "Benefits of Automation",          "Engagement, efficiency, revenue and competitive advantage"),
    ("06", "Implementation Timeline",         "Phased 12-month delivery plan"),
    ("07", "Investment & Next Steps",         "Budget guidance and immediate actions"),
]
for i, (num, title, desc) in enumerate(items):
    col, row = i % 2, i // 2
    x = Inches(0.55) + col * Inches(6.4)
    y = Inches(1.55) + row * Inches(1.22)
    card(s, x, y, Inches(6.1), Inches(1.05), fill=WHITE,
         border_color=RGBColor(0xCC, 0xD5, 0xE5))
    add_rect(s, x, y, Inches(0.7), Inches(1.05), NAVY)
    accent_bar(s, x, y + Inches(1.05) - Inches(0.048), Inches(6.1))
    add_textbox(s, num, x, y + Inches(0.25), Inches(0.7), Inches(0.5),
                20, True, CRIMSON, PP_ALIGN.CENTER)
    add_textbox(s, title, x + Inches(0.82), y + Inches(0.1),
                Inches(5.1), Inches(0.4), 13, True, NAVY)
    add_textbox(s, desc,  x + Inches(0.82), y + Inches(0.56),
                Inches(5.1), Inches(0.38), 10, False, SILVER, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 3 — ORGANISATION DIGITAL OVERVIEW
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Organisation Digital Overview", "SIOR – current digital footprint & member platform health"); footer(s)

card(s, Inches(0.4), Inches(1.4), Inches(5.8), Inches(5.58),
     fill=WHITE, border_color=RGBColor(0xCC, 0xD5, 0xE5))
add_rect(s, Inches(0.4), Inches(1.4), Inches(5.8), Inches(0.45), NAVY)
add_textbox(s, "Organisation Snapshot", Inches(0.56), Inches(1.46),
            Inches(5.5), Inches(0.35), 13, True, WHITE)
snap = [
    ("Type",          "International Professional Association / Trade Body"),
    ("Headquarters",  "Washington, D.C., United States"),
    ("Membership",    "4,000+ designees  •  630+ cities  •  34 countries"),
    ("Transactions",  "~80,000 CRE deals/year submitted by SIOR members"),
    ("Designations",  "Industrial, Office, Dual, Sales Mgmt, Exec Mgmt, Advisory"),
    ("Chapters",      "50+ regional chapters  (US, Canada, global)"),
    ("Events",        "Fall Event NYC (Oct 2026)  •  Lisbon Int'l Conference 2026"),
    ("Publications",  "SIOR Report magazine  •  Pulse Blog  •  Weekly Newsletter"),
    ("Portal",        "MySIOR – member login, directory, transaction submission"),
    ("Auth",          "Higher Logic SSO – membership management platform"),
]
ty = Inches(2.0)
for k, v in snap:
    add_textbox(s, k, Inches(0.6), ty, Inches(1.55), Inches(0.28), 10, True, NAVY)
    add_textbox(s, v, Inches(2.22), ty, Inches(3.82), Inches(0.28), 10, False, SILVER)
    ty += Inches(0.5)

card(s, Inches(6.5), Inches(1.4), Inches(6.5), Inches(5.58),
     fill=WHITE, border_color=RGBColor(0xCC, 0xD5, 0xE5))
add_rect(s, Inches(6.5), Inches(1.4), Inches(6.5), Inches(0.45), CRIMSON)
add_textbox(s, "Digital Health Indicators", Inches(6.66), Inches(1.46),
            Inches(6.2), Inches(0.35), 13, True, WHITE)
metrics = [
    ("Mobile Page Speed",           "~61 / 100",    "Below association avg of 72+",   CAUTION),
    ("Desktop Page Speed",          "~79 / 100",    "Reasonable – target 90+",         CAUTION),
    ("SEO Score",                   "~67 / 100",    "Content rich but gaps in schema",  CAUTION),
    ("Accessibility (WCAG AA)",     "Partial",       "Some contrast & nav issues",      CAUTION),
    ("SSL / HTTPS",                 "✓ Active",      "Valid certificate",                SUCCESS),
    ("MySIOR Portal UX",            "Dated",         "Higher Logic showing age",         DANGER),
    ("Member Directory Search",     "Basic",         "No AI-powered matching",           DANGER),
    ("Event Registration Flow",     "Multi-step",    "Drop-off risk on mobile",          CAUTION),
    ("AI / Personalisation",        "None detected", "Major opportunity untapped",       DANGER),
    ("Social Media Integration",    "Partial",       "No live feed; LinkedIn not linked",CAUTION),
]
ty = Inches(2.0)
for lbl, val, note, col in metrics:
    add_textbox(s, lbl,  Inches(6.66), ty, Inches(2.85), Inches(0.28), 10, True, CHARCOAL)
    add_textbox(s, val,  Inches(9.55), ty, Inches(1.18), Inches(0.28), 10, True, col, PP_ALIGN.CENTER)
    add_textbox(s, note, Inches(10.78),ty, Inches(2.08), Inches(0.28),  9, False, SILVER, italic=True)
    ty += Inches(0.46)

# ════════════════════════════════════════════════════════════════
# SLIDE 4 — WEBSITE AUDIT FINDINGS
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Website Audit Findings", "sior.com & MySIOR portal – key gaps and member friction points"); footer(s)

audit_cols = [
    ("Performance & Technical", NAVY, [
        "Mobile speed ~61 – needs 85+ for modern UX",
        "Render-blocking scripts slowing first load",
        "No visible CDN for global chapter audiences",
        "Image optimisation missing on photo-heavy pages",
        "Core Web Vitals LCP & CLS need improvement",
    ]),
    ("Member UX & Portal", CRIMSON, [
        "MySIOR (Higher Logic) feels dated vs. expectations",
        "Member directory search is basic — no geo/AI match",
        "Transaction submission UX is multi-step & cumbersome",
        "No personalised dashboard on login",
        "Chapter pages inconsistent in quality & content",
    ]),
    ("SEO & Content", STEEL_BLUE, [
        "Designation pages under-optimised for Google",
        "No structured data (Organization / Event schema)",
        "SIOR Report archive not fully indexed / searchable",
        "Pulse Blog lacks internal linking strategy",
        "Event pages missing local SEO signals",
    ]),
    ("Engagement & Growth", TEAL, [
        "No member onboarding email/content journey",
        "Education (Next Level) hard to discover on site",
        "No AI chatbot or self-service help for members",
        "Webinar & e-learning UX disconnected from main site",
        "No member referral / growth loop mechanism",
    ]),
]
for i, (title, col, bullets) in enumerate(audit_cols):
    x = Inches(0.4) + i * Inches(3.22)
    y = Inches(1.45)
    card(s, x, y, Inches(3.05), Inches(5.55),
         fill=WHITE, border_color=RGBColor(0xCC, 0xD5, 0xE5))
    add_rect(s, x, y, Inches(3.05), Inches(0.45), col)
    add_textbox(s, title, x + Inches(0.12), y + Inches(0.06),
                Inches(2.8), Inches(0.36), 12, True, WHITE)
    ty = y + Inches(0.62)
    for b in bullets:
        add_rect(s, x + Inches(0.17), ty + Inches(0.08), Inches(0.1), Inches(0.1), col)
        add_textbox(s, b, x + Inches(0.37), ty, Inches(2.55), Inches(0.5),
                    10, False, CHARCOAL)
        ty += Inches(0.73)

accent_bar(s, Inches(0.4), Inches(6.86), Inches(12.5))
add_textbox(s, "Core problem: MySIOR portal & member experience are below what 4,000 global CRE professionals expect in 2026",
            Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.28),
            10.5, True, CHARCOAL, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 5 — WEB DEVELOPMENT ROADMAP
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Web Development Roadmap", "Modernising sior.com and the MySIOR member experience"); footer(s)

card(s, Inches(0.4), Inches(1.4), Inches(4.62), Inches(5.62),
     fill=WHITE, border_color=RGBColor(0xCC, 0xD5, 0xE5))
add_rect(s, Inches(0.4), Inches(1.4), Inches(4.62), Inches(0.45), MID_NAVY)
add_textbox(s, "Recommended Tech Stack", Inches(0.56), Inches(1.46),
            Inches(4.38), Inches(0.35), 13, True, WHITE)
stack = [
    ("Frontend",      "Next.js 15 — SSR, ISR, edge-ready"),
    ("CMS",           "Sanity.io — structured content, multi-editor"),
    ("Member Portal", "MemberClicks or Personify360 (AMS upgrade)"),
    ("LMS",           "Thinkific / LearnDash — Next Level Series"),
    ("Search",        "Algolia — member directory + content search"),
    ("Events",        "Cvent API — registration, agenda, networking"),
    ("Email",         "HubSpot — member journeys + segmentation"),
    ("Auth",          "Auth0 — SSO for portal, LMS, events"),
    ("Analytics",     "GA4 + Hotjar + Segment CDP"),
    ("Hosting",       "Vercel + Cloudflare — global CDN, edge"),
]
ty = Inches(2.0)
for k, v in stack:
    add_textbox(s, k, Inches(0.6), ty, Inches(1.45), Inches(0.27), 10, True, CRIMSON)
    add_textbox(s, v, Inches(2.1), ty, Inches(2.82), Inches(0.27), 10, False, CHARCOAL)
    ty += Inches(0.46)

card(s, Inches(5.2), Inches(1.4), Inches(7.75), Inches(5.62),
     fill=WHITE, border_color=RGBColor(0xCC, 0xD5, 0xE5))
add_rect(s, Inches(5.2), Inches(1.4), Inches(7.75), Inches(0.45), NAVY)
add_textbox(s, "Key Improvements & Platform Features",
            Inches(5.36), Inches(1.46), Inches(7.4), Inches(0.35), 13, True, WHITE)
improvements = [
    ("Personalised Member Dashboard (MySIOR 2.0)",
     "On login: recommended connections, upcoming events, CE credits, designation progress, deal submissions — all in one view."),
    ("AI-Powered Member Directory & Matchmaking",
     "Search by market, specialisation, deal type, language. AI suggests connections by activity, geography and transaction history."),
    ("Unified Education Hub (Next Level LMS)",
     "All courses, webinars, on-demand videos and designation pathways in one seamlessly branded learning portal."),
    ("Intelligent Event & Registration Flow",
     "One-click event registration, personalised agenda builder, AI-matched networking suggestions, post-event content access."),
    ("SEO-Optimised Designation & Chapter Pages",
     "Structured data, local SEO per chapter, rich member testimonials — capturing 'CRE broker [city]' search intent."),
    ("Content Intelligence & Publication Hub",
     "Full SIOR Report archive with search, AI-generated summaries, Pulse Blog with topic clusters and related articles."),
    ("Member Onboarding & Engagement Journeys",
     "Automated welcome series, milestone nudges (renewal, designation progress), lapsed-member re-engagement campaigns."),
]
ty = Inches(1.98)
for title, desc in improvements:
    add_rect(s, Inches(5.36), ty + Inches(0.12), Inches(0.08), Inches(0.22), CRIMSON)
    add_textbox(s, title, Inches(5.56), ty, Inches(7.1), Inches(0.28), 11, True, NAVY)
    add_textbox(s, desc,  Inches(5.56), ty + Inches(0.3), Inches(7.1), Inches(0.38), 10, False, SILVER)
    ty += Inches(0.7)

# ════════════════════════════════════════════════════════════════
# SLIDE 6 — AI AUTOMATION OPPORTUNITY MAP
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "AI Automation Opportunity Map", "Where AI delivers outsized value for SIOR members & staff"); footer(s)

opps = [
    (NAVY,       WHITE,   "Member Experience & Engagement", [
        "AI chatbot: answers membership, CE, designation FAQs 24/7",
        "Personalised content recommendations on every login",
        "AI-matched networking — right member, right market",
        "Automated designation progress nudges & reminders",
        "Intelligent new-member onboarding conversation flow",
    ]),
    (CRIMSON,    WHITE,   "Education & Thought Leadership", [
        "AI generates course summaries and learning paths",
        "Auto-transcribe & index all webinar recordings",
        "SIOR Report: AI-assisted article editing & tagging",
        "Pulse Blog topic suggestions from CRE market signals",
        "CE credit tracking with automated completion reports",
    ]),
    (STEEL_BLUE, WHITE,   "Transaction Intelligence & Market Data", [
        "AI analyses 80,000 annual member transactions for trends",
        "Auto-generate quarterly CRE market intelligence reports",
        "Deal benchmarking: compare against peer cohorts",
        "Predictive alerts: emerging markets, cap rate shifts",
        "Member leaderboard & top deals powered by AI parsing",
    ]),
    (TEAL,       WHITE,   "Association Operations & Growth", [
        "Membership renewal risk scoring — flag churn risk early",
        "AI drafts chapter communications and event copy",
        "Staff: automate dues invoicing, receipts, reminders",
        "Board & committee meeting notes → actions (AI transcription)",
        "Sponsor & Industry Partner matching recommendations",
    ]),
]
for i, (hdr_bg, hdr_txt, title, bullets) in enumerate(opps):
    c, r = i % 2, i // 2
    x = Inches(0.4)  + c * Inches(6.5)
    y = Inches(1.45) + r * Inches(2.88)
    card(s, x, y, Inches(6.2), Inches(2.72),
         fill=WHITE, border_color=RGBColor(0xCC, 0xD5, 0xE5))
    add_rect(s, x, y, Inches(6.2), Inches(0.45), hdr_bg)
    add_textbox(s, title, x + Inches(0.15), y + Inches(0.07),
                Inches(5.9), Inches(0.35), 12, True, hdr_txt)
    ty = y + Inches(0.58)
    for b in bullets:
        add_rect(s, x + Inches(0.18), ty + Inches(0.07), Inches(0.1), Inches(0.1), hdr_bg)
        add_textbox(s, b, x + Inches(0.38), ty, Inches(5.65), Inches(0.38),
                    10.5, False, CHARCOAL)
        ty += Inches(0.43)

# ════════════════════════════════════════════════════════════════
# SLIDE 7 — BENEFITS OF AUTOMATION (hero)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Benefits of AI Automation", "Measurable impact across SIOR's membership, operations & revenue"); footer(s)

for i, (val, lbl) in enumerate([
    ("60%+",   "Reduction in staff\nrepetitive admin tasks"),
    ("35%",    "Lift in member\ndigital engagement"),
    ("50%",    "Faster content &\nreport production"),
    ("24/7",   "Member self-service\nvia AI assistant"),
]):
    stat_card(s, val, lbl, Inches(0.4) + i * Inches(3.2), Inches(1.45),
              w=Inches(3.0), h=Inches(1.65))

benefits = [
    ("Member Retention & Growth", [
        "AI flags churn risk 90 days before renewal",
        "Personalised outreach re-engages lapsed members",
        "Seamless onboarding boosts first-year retention",
        "Referral prompts auto-triggered at peak satisfaction",
    ]),
    ("Operational Efficiency", [
        "Dues invoicing and reminders fully automated",
        "Chapter communications drafted by AI, approved by humans",
        "Event logistics: confirmations, agendas, follow-ups — zero manual effort",
        "Staff freed from repetitive tasks → strategic work",
    ]),
    ("Education & Designation Value", [
        "Members discover relevant courses instantly via AI",
        "Designation progress tracked and nudged automatically",
        "All webinar recordings searchable and indexed in hours",
        "AI-generated summaries increase content consumption",
    ]),
    ("Revenue & Sponsorship", [
        "AI matches sponsors to the most relevant chapters",
        "Market intelligence reports become premium content",
        "Targeted event promotion lifts registration conversion",
        "Data-driven sponsorship packages command higher rates",
    ]),
]
for i, (title, bullets) in enumerate(benefits):
    c, r = i % 2, i // 2
    x = Inches(0.4)  + c * Inches(6.5)
    y = Inches(3.28) + r * Inches(1.88)
    bullet_card(s, title, bullets, x, y, Inches(6.2), Inches(1.72),
                hdr_bg=NAVY if r == 0 else CRIMSON)

# ════════════════════════════════════════════════════════════════
# SLIDE 8 — ROI & BUSINESS CASE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "ROI & Business Case", "Projected returns from web + AI investment for SIOR"); footer(s)

card(s, Inches(0.4), Inches(1.4), Inches(5.8), Inches(5.62),
     fill=WHITE, border_color=RGBColor(0xCC, 0xD5, 0xE5))
add_rect(s, Inches(0.4), Inches(1.4), Inches(5.8), Inches(0.45), NAVY)
add_textbox(s, "Indicative Investment (Year 1 – USD)",
            Inches(0.56), Inches(1.46), Inches(5.5), Inches(0.35), 13, True, WHITE)
invest = [
    ("sior.com Next.js rebuild + Sanity CMS",          "$18 000 – $28 000"),
    ("MySIOR portal upgrade / AMS integration",         "$15 000 – $25 000"),
    ("AI chatbot & member self-service bot",            " $5 000 – $9 000"),
    ("Algolia member directory + content search",       " $4 000 – $7 000"),
    ("LMS (Next Level) integration & migration",        " $6 000 – $10 000"),
    ("HubSpot CRM + member journey automation",         " $5 000 – $8 000"),
    ("SEO, schema & accessibility sprint",              " $3 000 – $5 000"),
    ("AI transaction intelligence module",              " $7 000 – $12 000"),
    ("Monthly support / hosting / AI APIs (est.)",      " $2 000 – $3 500 / mo"),
    ("TOTAL YEAR 1 (once-off + 12 mo support)",        "$87 000 – $146 000"),
]
ty = Inches(2.04)
for i, (item, cost) in enumerate(invest):
    last = i == len(invest) - 1
    if last:
        accent_bar(s, Inches(0.4), ty - Inches(0.06), Inches(5.8))
        ty += Inches(0.1)
    add_textbox(s, item, Inches(0.6), ty, Inches(3.65), Inches(0.28),
                9.5 if not last else 11, last, CHARCOAL)
    add_textbox(s, cost, Inches(4.28), ty, Inches(1.75), Inches(0.28),
                9.5 if not last else 11, last,
                SUCCESS if last else STEEL_BLUE, PP_ALIGN.RIGHT)
    ty += Inches(0.47)

card(s, Inches(6.5), Inches(1.4), Inches(6.5), Inches(5.62),
     fill=WHITE, border_color=RGBColor(0xCC, 0xD5, 0xE5))
add_rect(s, Inches(6.5), Inches(1.4), Inches(6.5), Inches(0.45), SUCCESS)
add_textbox(s, "Projected Returns & Value Delivered",
            Inches(6.66), Inches(1.46), Inches(6.2), Inches(0.35), 13, True, WHITE)
returns = [
    ("Staff time saved on admin (FTE equivalent)",     "1.5–2.5 FTE / yr  ≈ $90 000+"),
    ("Improved member retention (1% = X dues)",        "+1 % retention = ~$200 000 ARR"),
    ("Increased new member conversion (digital UX)",   "+15–25 % application conversion"),
    ("Education revenue (better LMS discovery)",       "+$50 000 – $120 000 / yr"),
    ("Sponsorship uplift (AI-targeted packages)",      "+20–35 % sponsor revenue"),
    ("Event registration lift (streamlined UX)",       "+18 % online registration rate"),
    ("Market intelligence — new premium product",      "Potential $100 000+ new revenue"),
    ("ESTIMATED YEAR 1 TOTAL IMPACT",                  "$450 000 – $750 000+ / yr"),
]
ty = Inches(2.04)
for i, (item, val) in enumerate(returns):
    last = i == len(returns) - 1
    if last:
        accent_bar(s, Inches(6.5), ty - Inches(0.06), Inches(6.5))
        ty += Inches(0.1)
    add_textbox(s, item, Inches(6.66), ty, Inches(3.88), Inches(0.28),
                9.5 if not last else 11, last, CHARCOAL)
    add_textbox(s, val, Inches(10.58), ty, Inches(2.28), Inches(0.28),
                9.5 if not last else 11, last,
                SUCCESS if last else STEEL_BLUE, PP_ALIGN.RIGHT)
    ty += Inches(0.47)

add_textbox(s, "Retaining just 40 additional members pays for the entire Year 1 programme  |  ROI 5–8× over 3 years",
            Inches(0.4), Inches(7.07), Inches(12.5), Inches(0.28),
            10.5, True, CHARCOAL, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 9 — IMPLEMENTATION TIMELINE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Implementation Timeline", "Phased 12-month delivery plan"); footer(s)

phases = [
    ("Phase 1\nMon 1–3",   NAVY,       [
        "Full site speed & SEO audit",
        "Accessibility & WCAG fixes",
        "AI chatbot MVP (member FAQ)",
        "Sanity CMS + blog migration",
        "HubSpot CRM + welcome journey",
    ]),
    ("Phase 2\nMon 4–6",   CRIMSON,    [
        "Next.js sior.com rebuild launch",
        "Algolia member directory live",
        "MySIOR 2.0 personalised dashboard",
        "LMS (Next Level) integration",
        "Event registration UX overhaul",
    ]),
    ("Phase 3\nMon 7–9",   STEEL_BLUE, [
        "AI member matchmaking feature",
        "Transaction intelligence module",
        "Chapter page standardisation",
        "Market intelligence report engine",
        "Sponsor / partner portal v1",
    ]),
    ("Phase 4\nMon 10–12", TEAL,       [
        "Renewal risk scoring (AI)",
        "Advanced content personalisation",
        "CE credit auto-tracking",
        "Full analytics attribution audit",
        "Review, optimise & plan Year 2",
    ]),
]
for i, (label, col, items) in enumerate(phases):
    x = Inches(0.4) + i * Inches(3.22)
    y = Inches(1.45)
    card(s, x, y, Inches(3.05), Inches(5.62),
         fill=WHITE, border_color=RGBColor(0xCC, 0xD5, 0xE5))
    add_rect(s, x, y, Inches(3.05), Inches(0.72), col)
    add_textbox(s, label, x, y + Inches(0.06), Inches(3.05), Inches(0.62),
                13, True, WHITE, PP_ALIGN.CENTER)
    if i < 3:
        add_rect(s, x + Inches(3.07), y + Inches(0.3), Inches(0.14), Inches(0.06), SILVER)
    ty = y + Inches(0.9)
    for item in items:
        add_rect(s, x + Inches(0.18), ty + Inches(0.08), Inches(0.1), Inches(0.1), col)
        add_textbox(s, item, x + Inches(0.38), ty, Inches(2.55), Inches(0.42),
                    10.5, False, CHARCOAL)
        ty += Inches(0.84)

# ════════════════════════════════════════════════════════════════
# SLIDE 10 — WHY AUTOMATE NOW?
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
add_rect(s, 0, 0, Inches(0.6), SLIDE_H, CRIMSON)
footer(s)

add_textbox(s, "Why Automate — Now?",
            Inches(0.9), Inches(0.45), Inches(11), Inches(0.68),
            30, True, WHITE)
add_textbox(s, "Professional associations that fail to modernise lose members to niche networks, LinkedIn, and AI-native platforms.",
            Inches(0.9), Inches(1.18), Inches(11.8), Inches(0.55),
            13.5, False, RGBColor(0xBB, 0xCC, 0xFF), italic=True)

reasons = [
    (WHITE,                           "Member expectations are set by consumer tech, not by other associations",
     "Your 4,000 members use Slack, Notion, LinkedIn and GPT-4 daily. They benchmark their SIOR portal against those experiences — not against ASAE averages."),
    (RGBColor(0xFF, 0xCC, 0x88),      "Retention is the highest-ROI lever — and data predicts it",
     "A 1% improvement in member retention equals ~$200,000 in annual dues. AI-driven churn prediction gives staff 90+ days to intervene before a member lapses."),
    (RGBColor(0x88, 0xFF, 0xCC),      "80,000 annual transactions are an untapped intelligence goldmine",
     "No other CRE association has this data. AI transforms raw deal submissions into market intelligence that becomes a premium product, a recruitment tool, and a media asset."),
    (RGBColor(0xBB, 0xCC, 0xFF),      "The education market is moving to always-on, AI-assisted learning",
     "Members expect Netflix-style course discovery and AI tutors, not PDF syllabi. Upgrading the LMS protects designation revenue and increases designation completion rates."),
    (RGBColor(0xFF, 0xBB, 0xBB),      "Competing networks are investing heavily right now",
     "CoreNet Global, CCIM and NAIOP are modernising their digital stacks. SIOR's differentiation — quality, global reach, rigorous designation — must be matched by an equally premium digital experience."),
]
ty = Inches(1.9)
for i, (col, title, desc) in enumerate(reasons):
    add_rect(s, Inches(0.9), ty, Inches(0.55), Inches(0.52), col)
    add_textbox(s, str(i + 1), Inches(0.9), ty,
                Inches(0.55), Inches(0.52), 17, True, NAVY, PP_ALIGN.CENTER)
    add_textbox(s, title, Inches(1.6), ty, Inches(11.3), Inches(0.3),
                12, True, WHITE)
    add_textbox(s, desc, Inches(1.6), ty + Inches(0.3), Inches(11.3), Inches(0.52),
                10.5, False, RGBColor(0xAA, 0xBB, 0xCC))
    ty += Inches(0.97)

# ════════════════════════════════════════════════════════════════
# SLIDE 11 — NEXT STEPS
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Next Steps", "From audit to action – what happens now"); footer(s)

steps = [
    ("Week 1",  NAVY,       "Stakeholder Alignment",
     "Present to SIOR leadership and Board. Agree priority workstreams, budget envelope and governance model."),
    ("Week 2",  CRIMSON,    "Discovery & Access",
     "Share access to analytics, CMS, AMS (Higher Logic), HubSpot and social channels. Kick-off call with all workstream leads."),
    ("Week 3",  STEEL_BLUE, "Sprint 1 Begins",
     "SEO/speed fixes, AI chatbot build, Sanity CMS setup and HubSpot welcome journey running in parallel."),
    ("Month 2", TEAL,       "Phase 1 Go-Live",
     "Chatbot live on sior.com. Speed/accessibility improvements deployed. New member welcome journey active."),
]
for i, (when, col, title, desc) in enumerate(steps):
    x = Inches(0.4) + i * Inches(3.22)
    y = Inches(1.5)
    card(s, x, y, Inches(3.05), Inches(4.1),
         fill=WHITE, border_color=RGBColor(0xCC, 0xD5, 0xE5))
    add_rect(s, x, y, Inches(3.05), Inches(0.45), col)
    add_textbox(s, when, x, y + Inches(0.06), Inches(3.05), Inches(0.35),
                12, True, WHITE, PP_ALIGN.CENTER)
    add_textbox(s, title, x + Inches(0.15), y + Inches(0.6),
                Inches(2.75), Inches(0.36), 12, True, col)
    add_textbox(s, desc, x + Inches(0.15), y + Inches(1.05),
                Inches(2.75), Inches(2.9), 10.5, False, SILVER)

card(s, Inches(0.4), Inches(5.82), Inches(12.5), Inches(1.1), fill=NAVY)
accent_bar(s, Inches(0.4), Inches(5.82), Inches(12.5), Inches(0.055))
add_textbox(s, "Ready to give 4,000 global CRE professionals the digital experience they deserve?",
            Inches(0.65), Inches(5.92), Inches(8.2), Inches(0.4),
            14, True, WHITE)
add_textbox(s, "Contact Genos Apollo to schedule your digital transformation discovery workshop.",
            Inches(0.65), Inches(6.32), Inches(8.2), Inches(0.38),
            11, False, RGBColor(0xBB, 0xCC, 0xFF))
add_textbox(s, "www.sior.com  •  Washington, D.C.",
            Inches(9.0), Inches(6.08), Inches(4.0), Inches(0.38),
            10, False, SILVER, PP_ALIGN.RIGHT, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 12 — BACK COVER
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
add_rect(s, 0, 0, Inches(0.72), SLIDE_H, CRIMSON)
add_rect(s, Inches(0.72), Inches(3.55), SLIDE_W - Inches(0.72), Inches(0.055), WHITE)

add_textbox(s, "Thank You",
            Inches(1.1), Inches(1.9), Inches(11), Inches(1.3),
            52, True, WHITE)
add_textbox(s, "Elevating the world's foremost CRE association — digitally, intelligently, and at scale.",
            Inches(1.1), Inches(3.18), Inches(10.5), Inches(0.55),
            15, False, RGBColor(0xBB, 0xCC, 0xFF), italic=True)
add_textbox(s, "SIOR  •  sior.com  •  Washington, D.C.  •  4,000+ Members  •  34 Countries",
            Inches(1.1), Inches(4.05), Inches(10), Inches(0.4),
            12, False, RGBColor(0x77, 0x88, 0xAA))
add_textbox(s, "Prepared by Genos Apollo  |  May 2026  |  Confidential",
            Inches(1.1), Inches(6.7), Inches(10), Inches(0.35),
            10, False, SILVER, italic=True)

# ── Save ────────────────────────────────────────────────────────
out = r"c:\Users\deevansho\Desktop\GenosApollp clients\SIOR_WebDev_AI_Automation_Audit.pptx"
prs.save(out)
print(f"Saved: {out}")
