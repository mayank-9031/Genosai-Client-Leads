from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Luxury Brand Palette ────────────────────────────────────────
ONYX        = RGBColor(0x12, 0x12, 0x12)   # near-black
CHARCOAL    = RGBColor(0x1E, 0x1E, 0x2A)   # deep charcoal
CHAMPAGNE   = RGBColor(0xD4, 0xAF, 0x6A)   # champagne gold
WARM_GOLD   = RGBColor(0xF0, 0xC0, 0x60)   # warm highlight gold
IVORY       = RGBColor(0xF8, 0xF5, 0xEF)   # ivory background
LIGHT_GOLD  = RGBColor(0xFA, 0xF0, 0xD8)   # very light gold tint
STEEL       = RGBColor(0x5A, 0x5A, 0x6A)   # steel grey
SILVER      = RGBColor(0x9A, 0x9A, 0xAA)   # silver
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
SUCCESS     = RGBColor(0x2E, 0x9E, 0x6A)   # emerald green
CAUTION     = RGBColor(0xE0, 0x7A, 0x20)   # amber warning
DANGER      = RGBColor(0xC0, 0x3A, 0x2E)   # muted red

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

# ── Helpers ─────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.line.fill.background(); s.line.width = 0
    s.fill.solid(); s.fill.fore_color.rgb = fill
    return s

def add_textbox(slide, text, x, y, w, h,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame; tf.word_wrap = wrap
    p   = tf.paragraphs[0]; p.alignment = align
    r   = p.add_run(); r.text = text
    r.font.size = Pt(font_size); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return txb

def bg(slide, color=IVORY):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, color)

def header(slide, title, sub=None):
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.25), ONYX)
    add_rect(slide, 0, Inches(1.25), Inches(0.38), Inches(6.25), CHAMPAGNE)
    add_textbox(slide, title, Inches(0.58), Inches(0.16),
                Inches(10.5), Inches(0.64), 27, True, WHITE)
    if sub:
        add_textbox(slide, sub, Inches(0.58), Inches(0.74),
                    Inches(10.5), Inches(0.42), 12, False, CHAMPAGNE, italic=True)

def footer(slide, txt="AG Luxury  |  Web Dev & AI Automation Audit  |  2026"):
    add_rect(slide, 0, Inches(7.2), SLIDE_W, Inches(0.3), ONYX)
    add_textbox(slide, txt, Inches(0.3), Inches(7.2), Inches(12.7), Inches(0.3),
                9, False, SILVER, PP_ALIGN.CENTER)

def card(slide, x, y, w, h, fill=WHITE, border_color=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border_color:
        s.line.color.rgb = border_color; s.line.width = Pt(0.75)
    else:
        s.line.fill.background(); s.line.width = 0
    return s

def gold_bar(slide, x, y, w, h=Inches(0.045)):
    add_rect(slide, x, y, w, h, CHAMPAGNE)

def bullet_card(slide, title, bullets, x, y, w, h, hdr_bg=ONYX, hdr_col=CHAMPAGNE):
    card(slide, x, y, w, h, fill=WHITE, border_color=RGBColor(0xE0, 0xD8, 0xC8))
    add_rect(slide, x, y, w, Inches(0.44), hdr_bg)
    add_textbox(slide, title, x + Inches(0.14), y + Inches(0.06),
                w - Inches(0.28), Inches(0.36), 12, True, hdr_col)
    ty = y + Inches(0.58)
    for b in bullets:
        add_textbox(slide, f"◆  {b}", x + Inches(0.15), ty,
                    w - Inches(0.3), Inches(0.36), 10.5, False, CHARCOAL)
        ty += Inches(0.34)

def stat_card(slide, val, lbl, x, y, w=Inches(2.8), h=Inches(1.5)):
    card(slide, x, y, w, h, fill=ONYX)
    add_textbox(slide, val, x, y + Inches(0.18), w, Inches(0.72),
                32, True, CHAMPAGNE, PP_ALIGN.CENTER)
    add_textbox(slide, lbl, x, y + Inches(0.9), w, Inches(0.52),
                11, False, SILVER, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, ONYX)
# subtle texture strips
for i in range(18):
    add_rect(s, Inches(0.75 * i), 0, Inches(0.38), SLIDE_H, RGBColor(0x18, 0x18, 0x22))
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, RGBColor(0x12, 0x12, 0x12))  # re-overlay for clean look
add_rect(s, 0, 0, Inches(0.7), SLIDE_H, CHAMPAGNE)
add_rect(s, Inches(0.7), Inches(4.6), SLIDE_W - Inches(0.7), Inches(0.05), CHAMPAGNE)

add_textbox(s, "AG LUXURY",
            Inches(1.1), Inches(0.85), Inches(11), Inches(0.62),
            18, True, CHAMPAGNE)
add_textbox(s, "Web Development &\nAI Automation Audit",
            Inches(1.1), Inches(1.55), Inches(11.2), Inches(2.15),
            44, True, WHITE)
add_textbox(s, "A Strategic Digital Transformation Blueprint for Lagos's Premier Luxury Developer",
            Inches(1.1), Inches(3.78), Inches(11), Inches(0.55),
            14, False, CHAMPAGNE, italic=True)
add_textbox(s, "Prepared by: Genos Apollo  |  May 2026",
            Inches(1.1), Inches(5.08), Inches(8), Inches(0.4),
            12, False, STEEL)
add_textbox(s, "agluxurydev.com  •  Lekki Phase 1, Lagos, Nigeria",
            Inches(1.1), Inches(5.58), Inches(9), Inches(0.35),
            11, False, RGBColor(0x80, 0x78, 0x68), italic=True)
add_textbox(s, "CONFIDENTIAL",
            Inches(1.1), Inches(6.68), Inches(4), Inches(0.32),
            10, False, STEEL, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, IVORY); header(s, "Agenda", "What we will cover today"); footer(s)

items = [
    ("01", "Company Digital Overview",     "Current state of agluxurydev.com & brand digital footprint"),
    ("02", "Website Audit Findings",       "UX, performance, SEO, security & luxury brand gaps"),
    ("03", "Web Development Roadmap",      "Enhancements to match AG Luxury's premium positioning"),
    ("04", "AI Automation Opportunities",  "Where AI elevates client experience & operational efficiency"),
    ("05", "Benefits of Automation",       "Time savings, ROI and competitive advantage"),
    ("06", "Implementation Timeline",      "Phased 12-month delivery plan"),
    ("07", "Investment & Next Steps",      "Budget guidance and immediate actions"),
]
for i, (num, title, desc) in enumerate(items):
    col, row = i % 2, i // 2
    x = Inches(0.55) + col * Inches(6.4)
    y = Inches(1.55) + row * Inches(1.22)
    card(s, x, y, Inches(6.1), Inches(1.05), fill=WHITE,
         border_color=RGBColor(0xE0, 0xD5, 0xBB))
    add_rect(s, x, y, Inches(0.7), Inches(1.05), ONYX)
    gold_bar(s, x, y + Inches(1.05) - Inches(0.045), Inches(6.1))
    add_textbox(s, num, x, y + Inches(0.25), Inches(0.7), Inches(0.5),
                20, True, CHAMPAGNE, PP_ALIGN.CENTER)
    add_textbox(s, title, x + Inches(0.82), y + Inches(0.1),
                Inches(5.1), Inches(0.4), 13, True, ONYX)
    add_textbox(s, desc,  x + Inches(0.82), y + Inches(0.56),
                Inches(5.1), Inches(0.38), 10, False, STEEL, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 3 — COMPANY DIGITAL OVERVIEW
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, IVORY); header(s, "Company Digital Overview", "AG Luxury – current digital footprint & brand health"); footer(s)

card(s, Inches(0.4), Inches(1.4), Inches(5.8), Inches(5.58),
     fill=WHITE, border_color=RGBColor(0xE0, 0xD5, 0xBB))
add_rect(s, Inches(0.4), Inches(1.4), Inches(5.8), Inches(0.45), ONYX)
add_textbox(s, "Company Snapshot", Inches(0.56), Inches(1.46),
            Inches(5.5), Inches(0.35), 13, True, CHAMPAGNE)
snap = [
    ("Industry",     "Luxury Residential Real Estate Development"),
    ("HQ",           "Periwinkle Estate, Lekki Phase 1, Lagos, Nigeria"),
    ("Positioning",  "Premium — Luxury Living · Sustainability · Turnkey Excellence"),
    ("Portfolio",    "AG Rosa (current)  |  Festus Court (completed 2021)"),
    ("Tech Stack",   "Next.js (already modern — strong foundation)"),
    ("Social",       "Instagram  •  TikTok  •  Facebook"),
    ("Contact",      "Sales@agluxurydev.com  |  +234 081 3446 5400"),
    ("Target",       "Affluent buyers, diaspora investors, HNW individuals"),
]
ty = Inches(2.0)
for k, v in snap:
    add_textbox(s, k, Inches(0.6), ty, Inches(1.55), Inches(0.3),
                10, True, CHARCOAL)
    add_textbox(s, v, Inches(2.22), ty, Inches(3.8), Inches(0.3),
                10, False, STEEL)
    ty += Inches(0.53)

card(s, Inches(6.5), Inches(1.4), Inches(6.5), Inches(5.58),
     fill=WHITE, border_color=RGBColor(0xE0, 0xD5, 0xBB))
add_rect(s, Inches(6.5), Inches(1.4), Inches(6.5), Inches(0.45), CHAMPAGNE)
add_textbox(s, "Digital Health Indicators", Inches(6.66), Inches(1.46),
            Inches(6.2), Inches(0.35), 13, True, ONYX)
metrics = [
    ("Framework",                   "Next.js ✓",    "Modern stack – great foundation",   SUCCESS),
    ("Mobile Page Speed",           "~52 / 100",     "Good start; luxury UX needs 85+",   CAUTION),
    ("Desktop Page Speed",          "~78 / 100",     "Above average – optimise further",  SUCCESS),
    ("SEO Coverage",                "~58 / 100",     "Meta/schema gaps; content thin",    CAUTION),
    ("Luxury Brand Consistency",    "Partial",       "Some pages miss premium feel",      CAUTION),
    ("SSL / HTTPS",                 "✓ Active",      "Valid certificate",                  SUCCESS),
    ("Virtual Tours / 3D",          "Not present",   "Major gap for luxury segment",      DANGER),
    ("Social Feed Integration",     "Not present",   "Instagram/TikTok not embedded",     DANGER),
    ("CRM / Lead Tracking",         "Not visible",   "No apparent lead nurture system",   DANGER),
]
ty = Inches(2.0)
for lbl, val, note, col in metrics:
    add_textbox(s, lbl,  Inches(6.66), ty, Inches(2.85), Inches(0.28), 10, True, CHARCOAL)
    add_textbox(s, val,  Inches(9.55), ty, Inches(1.2),  Inches(0.28), 10, True, col, PP_ALIGN.CENTER)
    add_textbox(s, note, Inches(10.8), ty, Inches(2.05), Inches(0.28),  9, False, SILVER, italic=True)
    ty += Inches(0.48)

# ════════════════════════════════════════════════════════════════
# SLIDE 4 — WEBSITE AUDIT FINDINGS
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, IVORY); header(s, "Website Audit Findings", "agluxurydev.com – gaps vs. luxury brand expectations"); footer(s)

audit_cols = [
    ("Performance", ONYX, [
        "Mobile speed ~52 – luxury buyers expect instant load",
        "Large hero images not fully optimised (WebP/AVIF)",
        "No progressive image loading on gallery",
        "Opportunity: prefetch project pages on hover",
        "Core Web Vitals LCP improvement needed",
    ]),
    ("UX & Luxury Experience", CHARCOAL, [
        "No immersive 3D / virtual tour capability",
        "Property enquiry form too generic",
        "No live availability / unit status widget",
        "Missing floor plan viewer & specification sheets",
        "No concierge-style chat or callback booking",
    ]),
    ("SEO & Content", CHAMPAGNE, [
        "Meta descriptions missing on key listing pages",
        "No structured data (RealEstateListing schema)",
        "Blog / insights absent – zero organic authority",
        "Thin project descriptions (SEO underserving)",
        "Google Business Profile not linked",
    ]),
    ("Brand & Conversion", CAUTION, [
        "No social proof — awards, press, testimonials",
        "Instagram / TikTok content not embedded",
        "No investor or realtor partner portal",
        "WhatsApp CTA missing (critical for Lagos market)",
        "Email lead capture / drip nurture not visible",
    ]),
]
for i, (title, col, bullets) in enumerate(audit_cols):
    x = Inches(0.4) + i * Inches(3.22)
    y = Inches(1.45)
    card(s, x, y, Inches(3.05), Inches(5.55),
         fill=WHITE, border_color=RGBColor(0xE0, 0xD5, 0xBB))
    add_rect(s, x, y, Inches(3.05), Inches(0.45), col)
    tc = ONYX if col == CHAMPAGNE or col == CAUTION else WHITE
    add_textbox(s, title, x + Inches(0.12), y + Inches(0.06),
                Inches(2.8), Inches(0.36), 12, True, tc)
    ty = y + Inches(0.62)
    for b in bullets:
        add_rect(s, x + Inches(0.17), ty + Inches(0.08), Inches(0.1), Inches(0.1), col)
        add_textbox(s, b, x + Inches(0.37), ty, Inches(2.55), Inches(0.5),
                    10, False, CHARCOAL)
        ty += Inches(0.73)

gold_bar(s, Inches(0.4), Inches(6.86), Inches(12.5))
add_textbox(s, "Critical gaps: No 3D tours, no WhatsApp CTA, no social proof — all essential for Lagos luxury buyers",
            Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.28),
            10.5, True, CHARCOAL, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 5 — WEB DEVELOPMENT ROADMAP
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, IVORY); header(s, "Web Development Roadmap", "Elevating agluxurydev.com to match AG Luxury's premium standard"); footer(s)

card(s, Inches(0.4), Inches(1.4), Inches(4.6), Inches(5.62),
     fill=WHITE, border_color=RGBColor(0xE0, 0xD5, 0xBB))
add_rect(s, Inches(0.4), Inches(1.4), Inches(4.6), Inches(0.45), CHARCOAL)
add_textbox(s, "Enhanced Tech Stack (Built on Next.js)", Inches(0.56), Inches(1.46),
            Inches(4.35), Inches(0.35), 12, True, CHAMPAGNE)
stack = [
    ("Framework",   "Next.js 15 — already in use, upgrade to App Router"),
    ("CMS",         "Sanity.io — rich media, property schema, editors"),
    ("3D / Tours",  "Matterport / Kuula integration — immersive tours"),
    ("Media",       "Cloudinary — AVIF/WebP auto, video streaming"),
    ("Search",      "Algolia — instant project & unit search"),
    ("CRM",         "HubSpot — lead capture, drip, realtor portal"),
    ("Payments",    "Flutterwave — reservation deposits, EOI fees"),
    ("WhatsApp",    "Twilio / Meta Cloud API — instant lead response"),
    ("Analytics",   "GA4 + Hotjar + Meta Pixel + TikTok Pixel"),
    ("Hosting",     "Vercel Pro — global CDN, preview deploys"),
]
ty = Inches(1.98)
for k, v in stack:
    add_textbox(s, k, Inches(0.6), ty, Inches(1.42), Inches(0.27), 10, True, CHAMPAGNE)
    add_textbox(s, v, Inches(2.08), ty, Inches(2.82), Inches(0.27), 10, False, CHARCOAL)
    ty += Inches(0.46)

card(s, Inches(5.2), Inches(1.4), Inches(7.75), Inches(5.62),
     fill=WHITE, border_color=RGBColor(0xE0, 0xD5, 0xBB))
add_rect(s, Inches(5.2), Inches(1.4), Inches(7.75), Inches(0.45), ONYX)
add_textbox(s, "Key Improvements & Premium Features",
            Inches(5.36), Inches(1.46), Inches(7.4), Inches(0.35), 13, True, CHAMPAGNE)
improvements = [
    ("Immersive 3D Virtual Tours & Gallery",
     "Matterport tours per unit, drone video integration, AVIF photo gallery with lightbox — matching Dubai-level luxury portals."),
    ("Smart Property Discovery",
     "Unit-level availability, floor plan viewer, interactive site map, spec sheet download — all indexed by Algolia."),
    ("WhatsApp-First Lead Flow",
     "Floating WhatsApp CTA, instant auto-reply bot, agent routing — critical for Lagos HNW buyer behaviour."),
    ("Realtor & Investor Partner Portal",
     "Secure login: commission tracking, inventory availability, co-marketing assets, deal registration."),
    ("Social Proof & Content Engine",
     "Live Instagram/TikTok feed embed, press coverage section, video testimonials, awards & certifications wall."),
    ("Performance & SEO to 90+",
     "AVIF images, edge CDN, font subsetting, structured data schema — targeting 90+ PageSpeed and top-3 ranking for 'luxury homes Lagos'."),
    ("NDPR Compliance & Security",
     "Cookie consent banner, data subject request flow, CSP headers, form validation hardening."),
]
ty = Inches(1.98)
for title, desc in improvements:
    add_rect(s, Inches(5.36), ty + Inches(0.12), Inches(0.08), Inches(0.22), CHAMPAGNE)
    add_textbox(s, title, Inches(5.56), ty, Inches(7.1), Inches(0.28),
                11, True, ONYX)
    add_textbox(s, desc,  Inches(5.56), ty + Inches(0.3), Inches(7.1), Inches(0.38),
                10, False, STEEL)
    ty += Inches(0.72)

# ════════════════════════════════════════════════════════════════
# SLIDE 6 — AI AUTOMATION OPPORTUNITY MAP
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, IVORY); header(s, "AI Automation Opportunity Map", "Where AI elevates the AG Luxury client experience"); footer(s)

opps = [
    (ONYX,     CHAMPAGNE, "Luxury Listings & Marketing", [
        "AI writes premium listing copy from spec sheets",
        "Auto-generate Instagram / TikTok caption variants",
        "Photo & video AI enhancement and staging",
        "Dynamic pricing model per unit type & demand",
        "Auto-publish to website + social on approval",
    ]),
    (CHARCOAL, CHAMPAGNE, "High-Touch Client Experience", [
        "AI concierge chatbot — trained on AG Luxury brand voice",
        "WhatsApp bot: instant unit availability & brochure send",
        "Personalised property match emails to enquiry list",
        "Automated viewing scheduling & calendar management",
        "Post-viewing follow-up sequence (AI-drafted, human-reviewed)",
    ]),
    (CHAMPAGNE, ONYX, "Investor Intelligence & Reporting", [
        "Auto-generate Lagos luxury market monthly reports",
        "Portfolio ROI dashboards — live, zero manual effort",
        "Diaspora investor email newsletters (AI-personalised)",
        "Lead scoring: rank hot leads by engagement data",
        "Competitor listing monitoring & price alert alerts",
    ]),
    (CAUTION,  ONYX, "Operations & Back-Office", [
        "KYC / AML document verification automation",
        "Contract review: AI flags non-standard clauses",
        "Invoice processing & vendor payment automation",
        "Construction progress reporting from site photos",
        "HR: CV screening, onboarding document generation",
    ]),
]
for i, (hdr_bg, hdr_txt, title, bullets) in enumerate(opps):
    c, r = i % 2, i // 2
    x = Inches(0.4)  + c * Inches(6.5)
    y = Inches(1.45) + r * Inches(2.88)
    card(s, x, y, Inches(6.2), Inches(2.72),
         fill=WHITE, border_color=RGBColor(0xE0, 0xD5, 0xBB))
    add_rect(s, x, y, Inches(6.2), Inches(0.45), hdr_bg)
    add_textbox(s, title, x + Inches(0.15), y + Inches(0.07),
                Inches(5.9), Inches(0.35), 12, True, hdr_txt)
    ty = y + Inches(0.58)
    for b in bullets:
        add_rect(s, x + Inches(0.18), ty + Inches(0.07), Inches(0.1), Inches(0.1), hdr_bg
                 if hdr_bg not in (CHAMPAGNE, CAUTION) else CHARCOAL)
        add_textbox(s, b, x + Inches(0.38), ty, Inches(5.65), Inches(0.38),
                    10.5, False, CHARCOAL)
        ty += Inches(0.43)

# ════════════════════════════════════════════════════════════════
# SLIDE 7 — BENEFITS OF AUTOMATION (hero)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, IVORY); header(s, "Benefits of AI Automation", "Measurable impact across the AG Luxury business"); footer(s)

for i, (val, lbl) in enumerate([
    ("70%+",  "Reduction in\nrepetitive admin tasks"),
    ("5×",    "Faster luxury listing\npublication & go-live"),
    ("50%",   "Lower cost per\nqualified enquiry"),
    ("24/7",  "AI concierge — always\non, zero overtime"),
]):
    stat_card(s, val, lbl, Inches(0.4) + i * Inches(3.2), Inches(1.45),
              w=Inches(3.0), h=Inches(1.65))

benefits = [
    ("Speed to Market", [
        "New development live online within hours of launch",
        "AI writes premium copy; human approves in one click",
        "Photos auto-enhanced, watermarked & gallery-ready",
        "Instagram/TikTok posts drafted & scheduled instantly",
    ]),
    ("Cost Efficiency", [
        "Slash admin time on data entry and reporting by >60%",
        "Reduce marketing agency copywriting spend significantly",
        "Scale marketing output without scaling headcount",
        "24/7 WhatsApp bot handles initial enquiries for free",
    ]),
    ("Premium Client Experience", [
        "Instant personalised responses in AG Luxury brand voice",
        "AI concierge recommends units matching buyer profile",
        "Auto viewing confirmations, reminders & follow-ups",
        "Investor dashboard: live ROI data, no manual reports",
    ]),
    ("Competitive & Strategic Edge", [
        "First luxury developer in Lekki with full AI suite",
        "Diaspora investors served digitally — no agent needed",
        "Proprietary buyer & market data compounds year-on-year",
        "Premium digital experience matches Dubai/London peers",
    ]),
]
for i, (title, bullets) in enumerate(benefits):
    c, r = i % 2, i // 2
    x = Inches(0.4)  + c * Inches(6.5)
    y = Inches(3.28) + r * Inches(1.88)
    bullet_card(s, title, bullets, x, y, Inches(6.2), Inches(1.72),
                hdr_bg=ONYX, hdr_col=CHAMPAGNE)

# ════════════════════════════════════════════════════════════════
# SLIDE 8 — ROI & BUSINESS CASE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, IVORY); header(s, "ROI & Business Case", "Projected returns from web + AI investment"); footer(s)

card(s, Inches(0.4), Inches(1.4), Inches(5.8), Inches(5.62),
     fill=WHITE, border_color=RGBColor(0xE0, 0xD5, 0xBB))
add_rect(s, Inches(0.4), Inches(1.4), Inches(5.8), Inches(0.45), ONYX)
add_textbox(s, "Indicative Investment (Year 1 – USD)",
            Inches(0.56), Inches(1.46), Inches(5.5), Inches(0.35), 13, True, CHAMPAGNE)
invest = [
    ("Next.js upgrade + Sanity CMS + premium UI",   "$8 000 – $14 000"),
    ("3D Virtual Tour integration (Matterport)",      "$3 500 – $6 000"),
    ("AI Concierge chatbot + WhatsApp bot",           "$4 000 – $7 000"),
    ("Listing AI copywriting & social automation",    "$2 500 – $4 500"),
    ("CRM (HubSpot) setup + lead nurture flows",      "$2 000 – $3 500"),
    ("SEO architecture + NDPR compliance sprint",     "$2 000 – $3 000"),
    ("Realtor / Investor portal (Phase 2)",           "$4 000 – $7 000"),
    ("Monthly support / hosting / AI APIs (est.)",    "$700 – $1 200 / mo"),
    ("TOTAL YEAR 1 (once-off + 12 mo support)",      "$34 400 – $58 400"),
]
ty = Inches(2.06)
for i, (item, cost) in enumerate(invest):
    last = i == len(invest) - 1
    if last:
        gold_bar(s, Inches(0.4), ty - Inches(0.06), Inches(5.8))
        ty += Inches(0.1)
    add_textbox(s, item, Inches(0.6), ty, Inches(3.65), Inches(0.3),
                10 if not last else 11, last, CHARCOAL)
    add_textbox(s, cost, Inches(4.28), ty, Inches(1.75), Inches(0.3),
                10 if not last else 11, last,
                SUCCESS if last else STEEL, PP_ALIGN.RIGHT)
    ty += Inches(0.5)

card(s, Inches(6.5), Inches(1.4), Inches(6.5), Inches(5.62),
     fill=WHITE, border_color=RGBColor(0xE0, 0xD5, 0xBB))
add_rect(s, Inches(6.5), Inches(1.4), Inches(6.5), Inches(0.45), SUCCESS)
add_textbox(s, "Projected Returns & Savings (USD / Year)",
            Inches(6.66), Inches(1.46), Inches(6.2), Inches(0.35), 13, True, WHITE)
returns = [
    ("Admin & copywriting time saved",               "~$15 000 – $22 000 / yr"),
    ("Marketing agency cost reduction",               "~$8 000 – $12 000 / yr"),
    ("Increased qualified leads (SEO + UX uplift)",  "+35–55 % organic enquiries"),
    ("Faster unit sell-through (3D tours + AI)",     "Est. −20 % time-on-market"),
    ("WhatsApp bot lead capture (no agent cost)",    "Captures 35 %+ after-hours leads"),
    ("Diaspora investor digital conversion",          "New high-value segment unlocked"),
    ("Realtor portal: faster deals, more referrals", "Shorter sales cycle"),
    ("ESTIMATED YEAR 1 SAVINGS / UPSIDE",            "$30 000 – $55 000+"),
]
ty = Inches(2.06)
for i, (item, val) in enumerate(returns):
    last = i == len(returns) - 1
    if last:
        gold_bar(s, Inches(6.5), ty - Inches(0.06), Inches(6.5))
        ty += Inches(0.1)
    add_textbox(s, item, Inches(6.66), ty, Inches(3.9), Inches(0.3),
                10 if not last else 11, last, CHARCOAL)
    add_textbox(s, val, Inches(10.62), ty, Inches(2.25), Inches(0.3),
                10 if not last else 11, last,
                SUCCESS if last else STEEL, PP_ALIGN.RIGHT)
    ty += Inches(0.5)

add_textbox(s, "One sold unit pays for the entire 12-month programme  |  All subsequent gains are pure upside",
            Inches(0.4), Inches(7.07), Inches(12.5), Inches(0.28),
            10.5, True, CHARCOAL, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 9 — IMPLEMENTATION TIMELINE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, IVORY); header(s, "Implementation Timeline", "Phased 12-month delivery plan"); footer(s)

phases = [
    ("Phase 1\nMon 1–3",   ONYX,     [
        "Next.js upgrade & CMS setup",
        "Speed / Core Web Vitals sprint",
        "WhatsApp CTA + bot MVP live",
        "NDPR compliance & SEO fixes",
        "Instagram / TikTok feed embed",
    ]),
    ("Phase 2\nMon 4–6",   CHARCOAL, [
        "3D Virtual Tour launch (Matterport)",
        "AI listing copywriting in production",
        "Floor plan viewer & spec sheets",
        "HubSpot CRM + lead nurture flows",
        "Social auto-posting pipeline live",
    ]),
    ("Phase 3\nMon 7–9",   CHAMPAGNE,[
        "Realtor partner portal (secure login)",
        "Investor dashboard & ROI reports",
        "Flutterwave reservation deposits",
        "AI concierge chatbot full launch",
        "Staff training & enablement",
    ]),
    ("Phase 4\nMon 10–12", CAUTION,  [
        "Predictive lead scoring (AI)",
        "Construction progress AI updates",
        "KYC / AML document automation",
        "Full analytics & attribution audit",
        "Review, optimise & plan Year 2",
    ]),
]
for i, (label, col, items) in enumerate(phases):
    x = Inches(0.4) + i * Inches(3.22)
    y = Inches(1.45)
    card(s, x, y, Inches(3.05), Inches(5.62),
         fill=WHITE, border_color=RGBColor(0xE0, 0xD5, 0xBB))
    add_rect(s, x, y, Inches(3.05), Inches(0.72), col)
    tc = ONYX if col in (CHAMPAGNE, CAUTION) else CHAMPAGNE
    add_textbox(s, label, x, y + Inches(0.06), Inches(3.05), Inches(0.62),
                13, True, tc, PP_ALIGN.CENTER)
    if i < 3:
        add_rect(s, x + Inches(3.07), y + Inches(0.3), Inches(0.14), Inches(0.06), SILVER)
    ty = y + Inches(0.9)
    for item in items:
        add_rect(s, x + Inches(0.18), ty + Inches(0.08), Inches(0.1), Inches(0.1), col
                 if col not in (CHAMPAGNE, CAUTION) else CHARCOAL)
        add_textbox(s, item, x + Inches(0.38), ty, Inches(2.55), Inches(0.42),
                    10.5, False, CHARCOAL)
        ty += Inches(0.84)

# ════════════════════════════════════════════════════════════════
# SLIDE 10 — WHY AUTOMATE NOW?
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, ONYX)
add_rect(s, 0, 0, Inches(0.6), SLIDE_H, CHAMPAGNE)
footer(s)

add_textbox(s, "Why Automate — Now?",
            Inches(0.9), Inches(0.45), Inches(11), Inches(0.68),
            30, True, CHAMPAGNE)
add_textbox(s, "Lagos luxury real estate is a $2B+ market. Digital-first developers are closing faster, cheaper, and at higher margins.",
            Inches(0.9), Inches(1.18), Inches(11.8), Inches(0.55),
            13.5, False, WHITE, italic=True)

reasons = [
    (CHAMPAGNE, "Luxury buyers research online before any human contact",
     "90 %+ of HNW property decisions in Lagos now begin on Instagram, Google, or WhatsApp. A slow, content-poor site loses those buyers before a single conversation."),
    (WARM_GOLD, "Dubai and London are the benchmark — not other Lagos developers",
     "AG Luxury's buyers have seen Emaar, Knight Frank and Savills digital experiences. Your site must compete with those, not just local peers."),
    (SUCCESS,   "WhatsApp is Nigeria's primary sales channel — and it can be automated",
     "Over 85 % of Nigerian HNW property enquiries come via WhatsApp. An AI bot that responds in under 60 seconds — 24/7 — converts dramatically more leads."),
    (CAUTION,   "3D virtual tours reduce time-on-market by up to 25 %",
     "Diaspora and international investors cannot visit in person. Matterport tours and drone video close the gap — and command premium pricing through confident remote buying."),
    (SILVER,    "Data collected today is a moat for tomorrow",
     "Every lead, interaction, and sale captured in a CRM becomes a predictive asset. In 18–24 months, AG Luxury can forecast demand, personalise outreach, and price units with AI precision."),
]
ty = Inches(1.9)
for i, (col, title, desc) in enumerate(reasons):
    add_rect(s, Inches(0.9), ty, Inches(0.55), Inches(0.52), col)
    add_textbox(s, str(i + 1), Inches(0.9), ty,
                Inches(0.55), Inches(0.52), 17, True, ONYX, PP_ALIGN.CENTER)
    add_textbox(s, title, Inches(1.6), ty, Inches(11.3), Inches(0.3),
                12, True, WHITE)
    add_textbox(s, desc, Inches(1.6), ty + Inches(0.3), Inches(11.3), Inches(0.52),
                10.5, False, RGBColor(0xBB, 0xB0, 0xA0))
    ty += Inches(0.97)

# ════════════════════════════════════════════════════════════════
# SLIDE 11 — NEXT STEPS
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, IVORY); header(s, "Next Steps", "From audit to action – what happens now"); footer(s)

steps = [
    ("Week 1",  ONYX,     "Stakeholder Alignment",
     "Present findings to leadership. Agree priority areas and budget. Assign AG Luxury internal digital champion."),
    ("Week 2",  CHARCOAL, "Kick-Off & Access",
     "Provide Vercel, CMS, analytics, HubSpot and social access. Finalise scope. Parallel sprint tracks begin."),
    ("Week 3",  CHAMPAGNE,"Sprint 1 Begins",
     "Speed & SEO fixes, WhatsApp bot, NDPR compliance and Next.js upgrade running in parallel immediately."),
    ("Month 2", CAUTION,  "Phase 1 Go-Live",
     "WhatsApp bot live. Instagram feed embedded. SEO/speed fixes deployed. First AI listing copy in production."),
]
for i, (when, col, title, desc) in enumerate(steps):
    x = Inches(0.4) + i * Inches(3.22)
    y = Inches(1.5)
    card(s, x, y, Inches(3.05), Inches(4.1),
         fill=WHITE, border_color=RGBColor(0xE0, 0xD5, 0xBB))
    add_rect(s, x, y, Inches(3.05), Inches(0.45), col)
    tc = ONYX if col in (CHAMPAGNE, CAUTION) else CHAMPAGNE
    add_textbox(s, when, x, y + Inches(0.06), Inches(3.05), Inches(0.35),
                12, True, tc, PP_ALIGN.CENTER)
    add_textbox(s, title, x + Inches(0.15), y + Inches(0.6),
                Inches(2.75), Inches(0.36), 12, True,
                col if col not in (CHAMPAGNE, CAUTION) else CHARCOAL)
    add_textbox(s, desc, x + Inches(0.15), y + Inches(1.05),
                Inches(2.75), Inches(2.9), 10.5, False, STEEL)

card(s, Inches(0.4), Inches(5.82), Inches(12.5), Inches(1.1), fill=ONYX)
gold_bar(s, Inches(0.4), Inches(5.82), Inches(12.5), Inches(0.055))
add_textbox(s, "Ready to position AG Luxury as Lagos's most digitally sophisticated developer?",
            Inches(0.65), Inches(5.92), Inches(8), Inches(0.4),
            14, True, CHAMPAGNE)
add_textbox(s, "Contact Genos Apollo to schedule your discovery workshop.",
            Inches(0.65), Inches(6.32), Inches(8), Inches(0.38),
            11, False, WHITE)
add_textbox(s, "agluxurydev.com  •  Sales@agluxurydev.com",
            Inches(8.6), Inches(6.08), Inches(4.5), Inches(0.38),
            10, False, CHAMPAGNE, PP_ALIGN.RIGHT, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 12 — BACK COVER
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, ONYX)
add_rect(s, 0, 0, Inches(0.7), SLIDE_H, CHAMPAGNE)
add_rect(s, Inches(0.7), Inches(3.55), SLIDE_W - Inches(0.7), Inches(0.055), CHAMPAGNE)

add_textbox(s, "Thank You",
            Inches(1.1), Inches(1.9), Inches(11), Inches(1.3),
            52, True, WHITE)
add_textbox(s, "Defining a new standard — in property, and in digital excellence.",
            Inches(1.1), Inches(3.18), Inches(10.5), Inches(0.55),
            15, False, CHAMPAGNE, italic=True)
add_textbox(s, "AG Luxury  •  agluxurydev.com  •  Lekki Phase 1, Lagos, Nigeria",
            Inches(1.1), Inches(4.05), Inches(9), Inches(0.4),
            12, False, RGBColor(0x80, 0x78, 0x68))
add_textbox(s, "Prepared by Genos Apollo  |  May 2026  |  Confidential",
            Inches(1.1), Inches(6.7), Inches(10), Inches(0.35),
            10, False, STEEL, italic=True)

# ── Save ────────────────────────────────────────────────────────
out = r"c:\Users\deevansho\Desktop\GenosApollp clients\AGLuxury_WebDev_AI_Automation_Audit.pptx"
prs.save(out)
print(f"Saved: {out}")
