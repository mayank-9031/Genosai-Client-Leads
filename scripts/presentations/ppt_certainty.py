import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.parent / "utils"))
from ppt_theme import (
    NAVY, DARK_BLUE, ACCENT, GOLD, RED, GREEN, PURPLE, ORANGE, TEAL, CYAN,
    WHITE, LGREY, MGREY, W, H,
    bg, rect, tb, aline, bullets, tag,
    slide_header, slide_footer, score_badge,
    stat_cards, top_bar_card, phase_card, before_after, genos_logo_block,
)

ROOT = Path(__file__).parent.parent.parent
prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]

# SLIDE 1 - COVER
s1 = prs.slides.add_slide(BLANK)
bg(s1)
rect(s1, 0, 0, Inches(0.12), H, TEAL)
rect(s1, Inches(8.5), 0, Inches(4.83), Inches(2.2), DARK_BLUE)
tb(s1, "GENOS AI", Inches(9.2), Inches(0.4), Inches(3.5), Inches(0.8),
   size=28, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
tb(s1, "www.genosai.tech", Inches(9.2), Inches(1.05), Inches(3.5), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
tb(s1, "CERTAINTY HOME LENDING", Inches(0.55), Inches(1.6),
   Inches(9), Inches(0.6), size=13, bold=False, color=TEAL)
tb(s1, "Website & AI Growth\nAudit Report", Inches(0.55), Inches(2.05),
   Inches(9), Inches(1.8), size=48, bold=True, color=WHITE)
aline(s1, Inches(3.85), color=TEAL, w=Inches(3))
tb(s1, "Prepared exclusively for Franco Terango, CEO",
   Inches(0.55), Inches(4.15), Inches(9), Inches(0.5), size=14, color=LGREY)
rect(s1, 0, Inches(6.5), W, Inches(1.0), DARK_BLUE)
tb(s1, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699",
   Inches(0.55), Inches(6.6), Inches(9), Inches(0.5), size=12, color=MGREY)
tb(s1, "May 2026  |  CONFIDENTIAL", Inches(9.5), Inches(6.6), Inches(3.3), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
rect(s1, Inches(10.5), Inches(3.3), Inches(2.3), Inches(2.5), DARK_BLUE)
tb(s1, "AUDIT SCORE", Inches(10.55), Inches(3.4), Inches(2.2), Inches(0.5),
   size=11, bold=True, color=MGREY, align=PP_ALIGN.CENTER)
tb(s1, "6", Inches(10.55), Inches(3.75), Inches(2.2), Inches(1.1),
   size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
tb(s1, "/ 10", Inches(10.55), Inches(4.7), Inches(2.2), Inches(0.4),
   size=18, color=LGREY, align=PP_ALIGN.CENTER)
tag(s1, "HIGH PRIORITY", Inches(10.75), Inches(5.1), fill=RED, size=10)


# SLIDE 2 - AT A GLANCE
s2 = prs.slides.add_slide(BLANK)
bg(s2)
rect(s2, 0, 0, Inches(0.12), H, TEAL)
tb(s2, "AT A GLANCE", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=TEAL)
tb(s2, "Franco Terango and Certainty Home Lending",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=30, bold=True, color=WHITE)
aline(s2, Inches(1.35), color=TEAL)

stats = [
    ("25 yrs", "Mortgage banking experience"),
    ("34 states\n+ DC", "Licensed nationwide"),
    ("200+", "Employees"),
    ("10,671", "Reviews at 4.9 stars"),
    ("5x", "MortgageCX Best-In-Class\n2018-2022"),
]
cw = Inches(2.3); ch = Inches(2.2); gap = Inches(0.18)
for i, (num, label) in enumerate(stats):
    l = Inches(0.55) + i * (cw + gap)
    rect(s2, l, Inches(1.8), cw, ch, DARK_BLUE)
    tb(s2, num, l + Inches(0.1), Inches(1.95), cw - Inches(0.2), Inches(1.1),
       size=20, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    tb(s2, label, l + Inches(0.1), Inches(3.05), cw - Inches(0.2), Inches(0.7),
       size=11, color=LGREY, align=PP_ALIGN.CENTER)

rect(s2, Inches(0.55), Inches(4.2), Inches(12.23), Inches(0.6), DARK_BLUE)
tb(s2, "DBA of Guaranteed Rate (#2 US Retail Mortgage Lender)  |  Founded 2000 as WR Starkey Mortgage  |  Conventional / FHA / VA / Jumbo / HELOC / Renovation / Construction",
   Inches(0.65), Inches(4.28), Inches(12), Inches(0.45), size=11, color=LGREY, align=PP_ALIGN.CENTER)

tb(s2, "FRANCO'S BACKGROUND", Inches(0.55), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=TEAL)
bullets(s2, [
    "CEO since April 2023 - 35+ years in financial services",
    "BA Economics, George Mason University 1987",
    "Board of Trustees, Greater Los Angeles Zoo Association since 2019",
    "Published: HousingWire, Scotsman Guide, MBA NewsLink"],
        Inches(0.55), Inches(5.35), Inches(5.8), Inches(1.5), size=12, color=LGREY)

tb(s2, "WHAT CERTAINTY OFFERS", Inches(6.8), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=TEAL)
bullets(s2, [
    "Purchase, refinance, HELOC across all major loan types",
    "Spanish-language digital mortgage (launched March 2026)",
    "Agent Advantage portal for Realtor partner relationships",
    "Parent Rate Intelligence AI platform (launched Oct 2024)"],
        Inches(6.8), Inches(5.35), Inches(5.8), Inches(1.5), size=12, color=LGREY)


# SLIDE 3 - STRENGTHS
s3 = prs.slides.add_slide(BLANK)
bg(s3)
rect(s3, 0, 0, Inches(0.12), H, GREEN)
tb(s3, "WHAT'S WORKING", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s3, "Strong Foundations - and a CEO Who Already Believes in AI",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=30, bold=True, color=WHITE)
aline(s3, Inches(1.35), color=GREEN)

strengths = [
    ("Franco Has Already Done the Ideological Work on AI in Mortgage",
     "He wrote in HousingWire that the LO role will soon be fundamentally different because of technology. He wrote in Scotsman Guide that user-friendly tools reducing administrative burden are the defining factor in LO recruitment and retention. He has publicly stated that real-time decision-making must occur across the organization - not just at the top. The concept sell is already done. The conversation with Franco is about specifics, not whether AI matters."),
    ("Parent Rate's $100M AI Infrastructure Investment Is Already in Place",
     "Rate Intelligence (launched October 2024) reduces mortgage approval from weeks to minutes. SmartUnderwrite provides real-time income and employment verification. The Same Day Mortgage program closed $1.1B in its pilot year. Certainty as a DBA benefits from this infrastructure investment. The gap is not at the underwriting layer - it is at the customer-facing layer: chatbot, social proof, rate alerts, and post-close automation that Certainty's website and brand do not yet have."),
    ("10,671 Reviews at 4.9 Stars Is a Trust Asset Most Lenders Cannot Match",
     "The volume and rating on RaveCapture represents years of deliberate service investment. It is one of the most powerful trust signals in mortgage lending - a space where referrals and reputation drive the majority of loan volume. The asset exists. The gap is deployment: it is not visible to a prospect landing on certaintyhomelending.com for the first time."),
    ("Spanish-Language Digital Mortgage Shows Active UX Investment Mindset",
     "Launching a complete Spanish-language digital mortgage experience in March 2026 is not a conservative move. It requires design, translation, legal review, and platform integration. It demonstrates that Franco and the leadership team invest in borrower experience as a competitive lever. An AI pre-qualification chatbot and embedded review module are natural extensions of that same investment mindset."),
    ("34-State Coverage and 200-Plus Employees Means Scale to Absorb Automation-Generated Volume",
     "A lender operating at this scale with active executive hiring in 2024-2025 is in expansion mode. Automation-generated lead volume will be absorbed by the existing LO team - not create a headcount problem. This is the right moment to build the customer-facing AI layer: the team is large enough to handle what better digital infrastructure generates."),
]
for i, (title, desc) in enumerate(strengths):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s3, Inches(0.55), t, Inches(0.06), Inches(0.7), GREEN)
    tb(s3, title, Inches(0.75), t, Inches(11.8), Inches(0.38), size=14, bold=True, color=WHITE)
    tb(s3, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# SLIDE 4 - THE THREE GAPS
s4 = prs.slides.add_slide(BLANK)
bg(s4)
rect(s4, 0, 0, Inches(0.12), H, RED)
tb(s4, "THE THREE GAPS", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=RED)
tb(s4, "10,671 Five-Star Reviews. Zero of Them on the Homepage.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=28, bold=True, color=WHITE)
aline(s4, Inches(1.35), color=RED)

gap_cards = [
    ("GAP 1\nBURIED\nSOCIAL PROOF", RED,
     "The trust signal that should be closing prospects",
     [
         "10,671 reviews at 4.9 stars on RaveCapture - not a single one visible on the homepage",
         "Prospects landing on certaintyhomelending.com see a calculator and a directory button",
         "Review embedded on the page that describes the loan converts at 2-3x the rate of the same page without it",
         "Competitors who embed testimonials above the fold are winning on trust before the LO conversation starts",
     ]),
    ("GAP 2\nZERO AFTER-HOURS\nCONVERSION", GOLD,
     "Where Sunday evening traffic goes",
     [
         "No chatbot, no live chat - a prospect at 9pm on a Sunday hits a static page and leaves",
         "AI lead response tools answer in under 90 seconds and pre-qualify before a human is needed",
         "Industry data: lenders using AI chatbots convert 37% more ready-to-apply leads within 60 days",
         "For a 34-state operation, weekend and after-hours traffic is not a rounding error - it is a substantial pipeline",
     ]),
    ("GAP 3\nNO PIPELINE FOR\nNOT-YET-READY LEADS", ACCENT,
     "What happens to prospects who are 6 months away",
     [
         "No rate alert subscription for prospects who want to transact when rates drop",
         "No email nurture capture for visitors who research but are not ready to apply",
         "No post-close referral sequence - warm referrals solicited manually or not at all",
         "A prospect who visits today and buys in 6 months will use whichever lender stayed in touch",
     ]),
]
cw3 = Inches(3.95); gap3i = Inches(0.18)
for i, (title, color, subtitle, pts) in enumerate(gap_cards):
    l = Inches(0.55) + i * (cw3 + gap3i)
    rect(s4, l, Inches(1.8), cw3, Inches(0.42), color)
    tb(s4, title, l + Inches(0.12), Inches(1.85), cw3 - Inches(0.24), Inches(0.36),
       size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    rect(s4, l, Inches(2.22), cw3, Inches(4.7), DARK_BLUE)
    tb(s4, subtitle, l + Inches(0.12), Inches(2.3), cw3 - Inches(0.24), Inches(0.35),
       size=10, italic=True, color=color)
    bullets(s4, pts, l + Inches(0.12), Inches(2.72), cw3 - Inches(0.24), Inches(3.9),
            size=12, color=LGREY, dot=color)

rect(s4, Inches(0.55), Inches(7.05), Inches(12.23), Inches(0.3), DARK_BLUE)
tb(s4, "The reviews are there. The traffic is there. The gaps are in what the website does with both.",
   Inches(0.7), Inches(7.08), Inches(12), Inches(0.25), size=11, italic=True, color=GOLD)


# SLIDE 5 - WEBSITE DESIGN AUDIT
s5 = prs.slides.add_slide(BLANK)
bg(s5)
rect(s5, 0, 0, Inches(0.12), H, PURPLE)
tb(s5, "WEBSITE DESIGN AUDIT", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=PURPLE)
tb(s5, "A Professional Site That Stops Short of Converting.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.8), size=32, bold=True, color=WHITE)
aline(s5, Inches(1.45), color=PURPLE, w=Inches(1.5))

rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(3.6), DARK_BLUE)
rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(0.42), RED)
tb(s5, "CERTAINTYHOMELENDING.COM TODAY", Inches(0.65), Inches(1.92), Inches(3.6), Inches(0.38),
   size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "No reviews embedded - 10,671 stars invisible to visitors",
    "Find a Loan Officer = directory dead end, no intelligent routing",
    "No chatbot or live chat - zero after-hours engagement",
    "No rate quote or rate alert tool for top-of-funnel visitors",
    "No video content or personalized LO video messaging",
    "No email nurture capture for not-yet-ready prospects",
    "Blog/content: press releases only, no SEO strategy",
    "Agent portal exists but has no AI or co-marketing features",
], Inches(0.7), Inches(2.42), Inches(3.6), Inches(2.9), size=11, color=LGREY, dot=RED)

tb(s5, "->", Inches(4.45), Inches(3.6), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(3.6), DARK_BLUE)
rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(0.42), GOLD)
tb(s5, "WHAT A CONVERTING LENDER SITE DOES", Inches(5.1), Inches(1.92), Inches(3.4), Inches(0.38),
   size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Reviews above the fold - trust signal visible before scroll",
    "AI chatbot: pre-qualifies and routes to right LO with context",
    "Rate alert subscription for not-yet-ready prospects",
    "Live rate quote or affordability calculator top of funnel",
    "LO video messaging for personalized prospect outreach",
    "Email nurture capture with automated drip for researchers",
    "SEO content hub targeting state-specific loan keywords",
    "Agent portal with co-marketing automation and lead sharing",
], Inches(5.1), Inches(2.42), Inches(3.4), Inches(2.9), size=11, color=LGREY, dot=GOLD)

tb(s5, "->", Inches(8.68), Inches(3.6), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(3.6), DARK_BLUE)
rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(0.42), GREEN)
tb(s5, "WHAT GENOS AI BUILDS", Inches(9.35), Inches(1.92), Inches(3.4), Inches(0.38),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Reviews embedded above fold on homepage + loan pages",
    "AI chatbot: 5-7 pre-qual questions, intelligent LO routing",
    "Rate alert subscription with automated drip on rate drops",
    "SEO content engine: state-specific + loan-type keywords",
    "Post-close referral and review sequence (D+3/7/30)",
    "Email nurture for research-stage prospects not yet applying",
    "LO video messaging capability for personalized outreach",
    "Agent portal AI layer: co-marketing and lead sharing",
], Inches(9.35), Inches(2.42), Inches(3.4), Inches(2.9), size=11, color=LGREY, dot=GREEN)

rect(s5, Inches(0.55), Inches(5.72), Inches(12.23), Inches(1.28), DARK_BLUE)
tb(s5, "CONVERSION BENCHMARK:", Inches(0.7), Inches(5.78), Inches(2.6), Inches(0.38),
   size=11, bold=True, color=PURPLE)
stats2 = [
    ("10,671", "five-star reviews not\nshowing on homepage"),
    ("0 sec", "chatbot response time\nvs. Monday morning"),
    ("37%", "more ready-to-apply leads\nwith AI lead response"),
    ("34 states", "of weekend traffic\nwith no after-hours capture"),
]
sw = Inches(2.7); sl = Inches(3.2)
for i, (num, label) in enumerate(stats2):
    xl = sl + i * sw
    tb(s5, num, xl, Inches(5.75), Inches(2.5), Inches(0.55),
       size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s5, label, xl, Inches(6.22), Inches(2.5), Inches(0.7),
       size=10, color=LGREY, align=PP_ALIGN.CENTER)


# SLIDE 6 - AI AUTOMATION STACK
s6 = prs.slides.add_slide(BLANK)
bg(s6)
rect(s6, 0, 0, Inches(0.12), H, ACCENT)
tb(s6, "AI AUTOMATION STACK", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s6, "Four Layers That Work While the LO Team Focuses on Relationships.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=28, bold=True, color=WHITE)
aline(s6, Inches(1.35), color=ACCENT)

opps = [
    ("AI Pre-Qualification\nChatbot", TEAL,
     "Replaces the Find a Loan Officer directory dead end. A prospect answers 5-7 qualifying questions - purchase or refi, loan type, credit range, loan amount, state, timeline. The system routes to the right LO with full context. After-hours prospects are captured and queued. Unqualified traffic is filtered without an LO spending time on it. Runs 24/7 across all 34 licensed states."),
    ("Rate Alert\nSubscription System", GOLD,
     "Visitors who are researching but not yet ready to transact subscribe to a rate watch. When rates drop to their threshold, they receive an automated email or SMS with a personalised affordability scenario. Franco's stated strategic concern is the market rebound - this is the pipeline infrastructure for when it comes. Prospects who subscribed 4 months ago become active leads the day rates move."),
    ("Post-Close Review\nand Referral Sequence", GREEN,
     "Certainty already has 10,671 reviews. The question is whether the next 10,000 come through systematic automation or manual follow-up. An automated sequence at Day 3, Day 7, and Day 30 post-close solicits the review at the emotional high point and then asks: do you know anyone buying or refinancing? Warm referrals generated at close are the highest-converting lead source in mortgage."),
    ("LO Prospecting\nAutomation for Agents", PURPLE,
     "Certainty's model is Realtor-referred. The Agent Advantage portal exists but has no AI layer. An automated system that monitors Realtor production data, sends co-marketing content on a cadence, and reminds LOs to follow up with top-producing agents at the right moment is a force multiplier for the existing agent relationship model - without adding headcount to the BD function."),
]
cw2 = Inches(2.95); gap2 = Inches(0.16)
for i, (title, color, desc) in enumerate(opps):
    l2 = Inches(0.55) + i * (cw2 + gap2)
    rect(s6, l2, Inches(1.8), cw2, Inches(0.06), color)
    rect(s6, l2, Inches(1.86), cw2, Inches(4.7), DARK_BLUE)
    tb(s6, title, l2 + Inches(0.12), Inches(1.98), cw2 - Inches(0.24), Inches(0.65),
       size=13, bold=True, color=color)
    tb(s6, desc, l2 + Inches(0.12), Inches(2.65), cw2 - Inches(0.24), Inches(3.65),
       size=11, color=LGREY)

rect(s6, Inches(0.55), Inches(6.7), Inches(12.23), Inches(0.45), DARK_BLUE)
tb(s6, "The LO's job is relationships and closings. The automation layer handles everything that doesn't require a licensed professional.",
   Inches(0.7), Inches(6.75), Inches(12), Inches(0.35), size=11, italic=True, color=GOLD)


# SLIDE 7 - WHAT THE INDUSTRY IS DOING (benchmark data - Franco speaks this language)
s7 = prs.slides.add_slide(BLANK)
bg(s7)
rect(s7, 0, 0, Inches(0.12), H, ORANGE)
tb(s7, "INDUSTRY BENCHMARKS", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ORANGE)
tb(s7, "What the Leading Lenders Are Already Running.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s7, Inches(1.35), color=ORANGE, w=Inches(1.8))

# Top stat bar
rect(s7, Inches(0.55), Inches(1.65), Inches(12.23), Inches(1.55), DARK_BLUE)
bm_stats = [
    ("< 90 sec", "AI lead response time\nvs. industry avg 12 hrs"),
    ("37%", "more ready-to-apply leads\nwith AI chatbot qualification"),
    ("55%", "of lenders starting/expanding\nAI trials in 2025 (Fannie Mae)"),
    ("$100M", "Rate Intelligence AI\ninvestment by parent company"),
    ("D+3/7/30", "post-close sequence timing\nfor highest referral conversion"),
]
ssw = Inches(2.35)
for i, (num, label) in enumerate(bm_stats):
    xl = Inches(0.7) + i * ssw
    tb(s7, num, xl, Inches(1.73), Inches(2.2), Inches(0.65),
       size=20, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    tb(s7, label, xl, Inches(2.33), Inches(2.2), Inches(0.5),
       size=10, color=LGREY, align=PP_ALIGN.CENTER)

# Two-column: what top lenders have vs. what Certainty has
rect(s7, Inches(0.55), Inches(3.4), Inches(5.9), Inches(3.25), DARK_BLUE)
tb(s7, "WHAT TOP LENDERS ARE DEPLOYING", Inches(0.7), Inches(3.47), Inches(5.6), Inches(0.35),
   size=11, bold=True, color=ORANGE)
bullets(s7, [
    "AI chatbot pre-qualification with intelligent LO routing (Structurely, Conversica, custom)",
    "Rate alert and market update automation for pipeline prospects (Bonzo, Mortgage Coach)",
    "Post-close review and referral sequences driving 20-30% of new loan volume (RaveCapture + CRM)",
    "SEO content engines publishing 50-200 articles/month on state and loan-type keywords (Rocket, LoanDepot)",
    "LO video messaging for pre-approval and rate lock communication (BombBomb, Vidyard)",
], Inches(0.7), Inches(3.85), Inches(5.6), Inches(2.55), size=12, color=LGREY, dot=ORANGE)

rect(s7, Inches(6.65), Inches(3.4), Inches(6.13), Inches(3.25), DARK_BLUE)
tb(s7, "WHAT CERTAINTY HAS TODAY", Inches(6.8), Inches(3.47), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=RED)
bullets(s7, [
    "No AI chatbot - Find a Loan Officer directory replaces qualification layer",
    "No rate alert subscription - not-yet-ready prospects have no re-engagement mechanism",
    "Post-close reviews exist (10,671) but solicitation sequence is not systematized",
    "No SEO content strategy - press releases only, long-tail keywords unaddressed across 34 states",
    "No LO video messaging capability at the brand level",
], Inches(6.8), Inches(3.85), Inches(5.8), Inches(2.55), size=12, color=LGREY, dot=RED)

rect(s7, Inches(0.55), Inches(6.8), Inches(12.23), Inches(0.38), DARK_BLUE)
tb(s7, "Parent Rate is investing $100M in AI at the underwriting layer. The customer-facing layer is where Certainty builds its own brand advantage.",
   Inches(0.7), Inches(6.85), Inches(12), Inches(0.28), size=11, italic=True, color=ORANGE)


# SLIDE 8 - IMPLEMENTATION PLAN
s8 = prs.slides.add_slide(BLANK)
bg(s8)
rect(s8, 0, 0, Inches(0.12), H, GREEN)
tb(s8, "IMPLEMENTATION PLAN", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s8, "A Phased Build - Start with One Market, Prove It, Scale.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=30, bold=True, color=WHITE)
aline(s8, Inches(1.35), color=GREEN)

phases = [
    ("PHASE 1", "Weeks 1-3", GOLD, [
        ("Reviews Embedded on Homepage", "RaveCapture review module integrated above the fold. Star rating, count, and rotating testimonials visible before prospect scrolls. Immediate trust signal for all existing traffic."),
        ("AI Chatbot Live", "Pre-qualification flow trained on Certainty's loan types, 34-state geography, and LO roster. After-hours capture active. Routes to correct LO with prospect context on submission."),
        ("Rate Alert Subscription", "Email capture with rate threshold input live on homepage and loan pages. Prospect enters target rate - automated alert fires when market hits it."),
    ]),
    ("PHASE 2", "Weeks 4-6", ACCENT, [
        ("Post-Close Sequence Active", "D+3 thank you and review request. D+7 review reminder. D+30 referral ask. All loans closed from Phase 1 go into the sequence immediately."),
        ("LO Video Messaging Deployed", "Video messaging capability live for pilot group of 10 LOs. Pre-approval delivery, rate lock confirmation, and market update templates provided."),
        ("Agent Portal AI Layer", "Co-marketing automation and production-triggered LO follow-up live in Agent Advantage. Pilot with top 20 Realtor relationships."),
    ]),
    ("PHASE 3", "Weeks 7-12", GREEN, [
        ("SEO Content Engine Live", "First 20 state-specific and loan-type articles published. Targeting VA loans Texas, HELOC rates California, FHA first-time buyer Illinois, and equivalent across top 10 licensed states."),
        ("Chatbot Optimisation", "Pre-qualification routing refined based on 6 weeks of live data. Conversion rate per loan type tracked. LO routing accuracy verified against closed loan data."),
        ("Full-Team Rollout Ready", "Phase 1-2 pilot results packaged for Franco to present internally. Video messaging, agent automation, and SEO scaled to full 34-state LO team."),
    ]),
]
pw = Inches(3.9); gap4 = Inches(0.18)
for i, (phase, timing, color, items) in enumerate(phases):
    l3 = Inches(0.55) + i * (pw + gap4)
    rect(s8, l3, Inches(1.8), pw, Inches(0.52), color)
    tb(s8, phase, l3 + Inches(0.12), Inches(1.84), pw * 0.5 - Inches(0.12), Inches(0.42),
       size=13, bold=True, color=NAVY)
    tb(s8, timing, l3 + pw * 0.5, Inches(1.88), pw * 0.5 - Inches(0.15), Inches(0.38),
       size=11, color=NAVY, align=PP_ALIGN.RIGHT)
    rect(s8, l3, Inches(2.32), pw, Inches(4.6), DARK_BLUE)
    y_item = Inches(2.42)
    for title, desc in items:
        rect(s8, l3 + Inches(0.12), y_item + Inches(0.06), Inches(0.05), Inches(0.52), color)
        tb(s8, title, l3 + Inches(0.25), y_item, pw - Inches(0.38), Inches(0.32),
           size=12, bold=True, color=color)
        tb(s8, desc, l3 + Inches(0.25), y_item + Inches(0.3), pw - Inches(0.38), Inches(0.45),
           size=11, color=LGREY)
        y_item += Inches(0.95)


# SLIDE 9 - BEFORE vs AFTER
s9 = prs.slides.add_slide(BLANK)
bg(s9)
rect(s9, 0, 0, Inches(0.12), H, ACCENT)
tb(s9, "BEFORE vs AFTER", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s9, "What Actually Changes at Certainty in 90 Days",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=30, bold=True, color=WHITE)
aline(s9, Inches(1.35))

comparisons = [
    ("Homepage Trust Signal",
     "10,671 reviews at 4.9 stars sit on RaveCapture. A prospect landing on certaintyhomelending.com sees none of them.",
     "Reviews embedded above the fold. Star rating and testimonials visible before scroll on every loan type page. Trust gap closed before the LO conversation starts."),
    ("After-Hours Lead Capture",
     "A prospect landing on a Sunday evening hits a static page. No chatbot. No capture. They leave and a competitor's AI follows up within 90 seconds.",
     "AI chatbot pre-qualifies, captures contact, and queues for Monday morning LO follow-up with full context. Zero after-hours traffic goes cold."),
    ("Not-Yet-Ready Pipeline",
     "A researcher who visits today but buys in 6 months has no connection to Certainty unless they call back. Most do not.",
     "Rate alert subscriber list grows from every site visit. Automated email fires when their threshold rate hits. They re-engage as an active, warm lead."),
    ("Post-Close Referrals",
     "Referral solicitation depends on LO memory and manual follow-up. The most valuable moment (emotional high at close) passes without a systematic ask.",
     "D+3, D+7, D+30 automated sequence. Review solicited at peak satisfaction. Referral asked at D+30. Warm referral volume measurable and growing monthly."),
    ("LO Administrative Load",
     "LOs spend significant time on pre-qualification questions, chasing documents, and manually nurturing relationships that automation handles at other lenders.",
     "Chatbot pre-qualifies. Rate alerts nurture. Post-close sequences handle referrals. LOs focus on relationships and closings - which is where their expertise creates value."),
]

lw = Inches(5.4)
rect(s9, Inches(0.55), Inches(1.65), lw, Inches(0.35), RED)
tb(s9, "TODAY", Inches(0.65), Inches(1.68), lw - Inches(0.2), Inches(0.28),
   size=11, bold=True, color=WHITE)
rect(s9, Inches(6.45), Inches(1.65), lw + Inches(0.43), Inches(0.35), GREEN)
tb(s9, "AFTER 90 DAYS", Inches(6.55), Inches(1.68), lw, Inches(0.28),
   size=11, bold=True, color=WHITE)

for i, (topic, before, after) in enumerate(comparisons):
    t = Inches(2.1) + i * Inches(1.02)
    rect(s9, Inches(0.55), t, Inches(12.23), Inches(0.22), DARK_BLUE)
    tb(s9, topic.upper(), Inches(0.65), t + Inches(0.02), Inches(12), Inches(0.2),
       size=9, bold=True, color=MGREY)
    tb(s9, before, Inches(0.65), t + Inches(0.24), lw - Inches(0.2), Inches(0.66),
       size=11, color=LGREY)
    tb(s9, after, Inches(6.55), t + Inches(0.24), lw, Inches(0.66),
       size=11, color=WHITE)


# SLIDE 10 - NEXT STEPS
s10 = prs.slides.add_slide(BLANK)
bg(s10)
rect(s10, 0, 0, Inches(0.12), H, TEAL)
rect(s10, Inches(7.5), 0, Inches(5.83), H, DARK_BLUE)
tb(s10, "NEXT STEPS", Inches(0.55), Inches(0.3), Inches(6.5), Inches(0.5),
   size=11, bold=True, color=TEAL)
tb(s10, "Let's Talk for\n15 Minutes.",
   Inches(0.55), Inches(0.8), Inches(6.8), Inches(2.0), size=48, bold=True, color=WHITE)
aline(s10, Inches(2.8), color=TEAL, w=Inches(2.5))
tb(s10, "No pitch deck required. Reply or book a call.\nWe'll walk through what fits Certainty Home Lending.",
   Inches(0.55), Inches(3.05), Inches(6.5), Inches(0.8), size=14, color=LGREY)

steps = [
    "Reply to this email - we'll set up a time",
    "Pick one gap to pilot first (Franco's preferred approach)",
    "Walk through the audit findings in 15 minutes",
    "You decide what, if anything, makes sense to move on",
]
bullets(s10, steps, Inches(0.55), Inches(3.95), Inches(6.5), Inches(2.0), size=14, color=WHITE, dot=TEAL)

rect(s10, 0, Inches(6.5), Inches(7.5), Inches(1.0), DARK_BLUE)
tb(s10, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699  |  www.genosai.tech",
   Inches(0.55), Inches(6.6), Inches(6.8), Inches(0.5), size=12, color=MGREY)

tb(s10, "THREE GAPS.\nTHREE FIXES.", Inches(7.85), Inches(0.6), Inches(5.0), Inches(1.1),
   size=22, bold=True, color=TEAL)
aline(s10, Inches(1.65), color=TEAL)

summary_items = [
    ("Buried Reviews", "Embed 10,671 five-star reviews above the fold on homepage"),
    ("Zero After-Hours", "AI chatbot pre-qualifies and routes - 24/7 across 34 states"),
    ("No Pipeline Nurture", "Rate alert subscription + D+3/7/30 post-close sequences"),
    ("No SEO Content", "State-specific + loan-type keyword targeting across licensed states"),
    ("Agent Portal Gap", "AI co-marketing layer on Agent Advantage for Realtor partners"),
    ("LO Admin Load", "Automation removes pre-qual, nurture, and referral manual work"),
]
for i, (label, val) in enumerate(summary_items):
    t2 = Inches(1.85) + i * Inches(0.88)
    rect(s10, Inches(7.85), t2, Inches(1.5), Inches(0.72), NAVY)
    tb(s10, label, Inches(7.9), t2 + Inches(0.05), Inches(1.4), Inches(0.32),
       size=10, bold=True, color=TEAL)
    tb(s10, val, Inches(9.5), t2 + Inches(0.05), Inches(3.6), Inches(0.65),
       size=11, color=LGREY)


prs.save(ROOT / "output" / "decks" / "CertaintyHomeLending_Audit_GenosAI.pptx")
print("Saved: CertaintyHomeLending_Audit_GenosAI.pptx")
