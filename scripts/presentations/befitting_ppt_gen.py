from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand colours (Befitting Group – emerald green + gold) ─────
DARK_GREEN  = RGBColor(0x05, 0x4A, 0x29)   # deep emerald
MID_GREEN   = RGBColor(0x0A, 0x7A, 0x45)   # mid green
LIGHT_GREEN = RGBColor(0xD9, 0xF0, 0xE5)   # soft mint bg
GOLD        = RGBColor(0xC8, 0xA0, 0x32)   # gold
AMBER       = RGBColor(0xF5, 0xA6, 0x23)   # warm amber accent
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xF5, 0xF7, 0xF5)
MID_GREY    = RGBColor(0x6B, 0x7B, 0x6E)
NAVY        = RGBColor(0x0D, 0x2B, 0x1E)
BLUE        = RGBColor(0x23, 0x7A, 0xB5)
SUCCESS     = RGBColor(0x27, 0xAE, 0x60)
DANGER      = RGBColor(0xE7, 0x4C, 0x3C)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

# ── Helpers ────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.line.fill.background(); s.line.width = 0
    s.fill.solid(); s.fill.fore_color.rgb = fill
    return s

def add_textbox(slide, text, x, y, w, h,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame; tf.word_wrap = wrap
    p   = tf.paragraphs[0]; p.alignment = align
    r   = p.add_run(); r.text = text
    r.font.size = Pt(font_size); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return txb

def bg(slide, color=LIGHT_GREY):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, color)

def header(slide, title, sub=None):
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.25), DARK_GREEN)
    add_rect(slide, 0, Inches(1.25), Inches(0.38), Inches(6.25), GOLD)
    add_textbox(slide, title, Inches(0.58), Inches(0.17),
                Inches(10), Inches(0.62), 27, True, WHITE)
    if sub:
        add_textbox(slide, sub, Inches(0.58), Inches(0.73),
                    Inches(10), Inches(0.42), 12, False, GOLD, italic=True)

def footer(slide, txt="Befitting Group  |  Web Dev & AI Automation Audit  |  2026"):
    add_rect(slide, 0, Inches(7.2), SLIDE_W, Inches(0.3), DARK_GREEN)
    add_textbox(slide, txt, Inches(0.3), Inches(7.2), Inches(12.7), Inches(0.3),
                9, False, WHITE, PP_ALIGN.CENTER)

def card(slide, x, y, w, h, fill=WHITE, border=True):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = RGBColor(0xCC, 0xDD, 0xCC); s.line.width = Pt(0.5)
    else:
        s.line.fill.background(); s.line.width = 0
    return s

def bullet_card(slide, title, bullets, x, y, w, h, hdr_bg=DARK_GREEN):
    card(slide, x, y, w, h)
    add_rect(slide, x, y, w, Inches(0.44), hdr_bg)
    add_textbox(slide, title, x + Inches(0.12), y + Inches(0.05),
                w - Inches(0.24), Inches(0.36), 12, True, WHITE)
    ty = y + Inches(0.55)
    for b in bullets:
        add_textbox(slide, f"▸  {b}", x + Inches(0.15), ty,
                    w - Inches(0.3), Inches(0.36), 10.5, False, NAVY)
        ty += Inches(0.34)

def stat_card(slide, val, lbl, x, y, w=Inches(2.8), h=Inches(1.5),
              bg_col=DARK_GREEN, vc=GOLD, lc=WHITE):
    card(slide, x, y, w, h, fill=bg_col, border=False)
    add_textbox(slide, val, x, y + Inches(0.18), w, Inches(0.7),
                32, True, vc, PP_ALIGN.CENTER)
    add_textbox(slide, lbl, x, y + Inches(0.9), w, Inches(0.5),
                11, False, lc, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, DARK_GREEN)
add_rect(s, 0, 0, Inches(0.75), SLIDE_H, GOLD)
add_rect(s, Inches(0.75), Inches(4.55), SLIDE_W - Inches(0.75), Inches(0.06), GOLD)

add_textbox(s, "BEFITTING GROUP",
            Inches(1.1), Inches(0.9), Inches(11), Inches(0.6),
            16, True, GOLD)
add_textbox(s, "Web Development &\nAI Automation Audit",
            Inches(1.1), Inches(1.6), Inches(11), Inches(2.1),
            44, True, WHITE)
add_textbox(s, "Strategic Digital Transformation Roadmap for a Pan-African Property Leader",
            Inches(1.1), Inches(3.75), Inches(11), Inches(0.55),
            15, False, AMBER, italic=True)
add_textbox(s, "Prepared by: Genos Apollo  |  May 2026",
            Inches(1.1), Inches(5.05), Inches(8), Inches(0.4),
            12, False, MID_GREY)
add_textbox(s, "befittingproperties.com  •  Lagos, Nigeria  •  Pan-African Operations",
            Inches(1.1), Inches(5.55), Inches(9), Inches(0.35),
            11, False, RGBColor(0x88, 0xAA, 0x88), italic=True)
add_textbox(s, "CONFIDENTIAL",
            Inches(1.1), Inches(6.65), Inches(4), Inches(0.35),
            10, False, MID_GREY, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Agenda", "What we will cover today"); footer(s)

items = [
    ("01", "Company Digital Overview",     "Current state of befittingproperties.com & digital footprint"),
    ("02", "Website Audit Findings",       "UX, performance, SEO, security & compliance gaps"),
    ("03", "Web Development Roadmap",      "Modernisation blueprint & tech stack recommendations"),
    ("04", "AI Automation Opportunities",  "Where AI drives efficiency, revenue & scale for Befitting"),
    ("05", "Benefits of Automation",       "Time savings, ROI and competitive edge"),
    ("06", "Implementation Timeline",      "Phased 12-month delivery plan"),
    ("07", "Investment & Next Steps",      "Budget guidance and immediate actions"),
]
for i, (num, title, desc) in enumerate(items):
    col, row = i % 2, i // 2
    x = Inches(0.55) + col * Inches(6.4)
    y = Inches(1.55) + row * Inches(1.22)
    card(s, x, y, Inches(6.1), Inches(1.05))
    add_rect(s, x, y, Inches(0.7), Inches(1.05), DARK_GREEN)
    add_textbox(s, num, x, y + Inches(0.25), Inches(0.7), Inches(0.5),
                20, True, GOLD, PP_ALIGN.CENTER)
    add_textbox(s, title, x + Inches(0.8), y + Inches(0.1),
                Inches(5.1), Inches(0.4), 13, True, DARK_GREEN)
    add_textbox(s, desc,  x + Inches(0.8), y + Inches(0.54),
                Inches(5.1), Inches(0.4), 10, False, MID_GREY, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 3 — COMPANY DIGITAL OVERVIEW
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Company Digital Overview", "Befitting Group – current digital footprint"); footer(s)

# Left – company snapshot
card(s, Inches(0.4), Inches(1.4), Inches(5.8), Inches(5.55))
add_rect(s, Inches(0.4), Inches(1.4), Inches(5.8), Inches(0.45), DARK_GREEN)
add_textbox(s, "Company Snapshot", Inches(0.56), Inches(1.45),
            Inches(5.5), Inches(0.36), 13, True, WHITE)
snap = [
    ("Industry",      "Real Estate Development, Property Management & Investment"),
    ("Headquarters",  "Lagos, Nigeria (Pan-African operations)"),
    ("Market",        "Residential, Commercial, Office & Mixed-Use Developments"),
    ("Experience",    "12+ years in African real estate sector"),
    ("Website",       "befittingproperties.com  (primary digital channel)"),
    ("Mobile App",    "Befitting Properties – available on App Store (iOS)"),
    ("LinkedIn",      "Befitting Group – active company presence"),
    ("Social",        "Instagram  @befittingproperties  |  Facebook active"),
]
ty = Inches(2.0)
for k, v in snap:
    add_textbox(s, k, Inches(0.6), ty, Inches(1.55), Inches(0.3),
                10, True, DARK_GREEN)
    add_textbox(s, v, Inches(2.2), ty, Inches(3.8), Inches(0.3),
                10, False, MID_GREY)
    ty += Inches(0.53)

# Right – digital health
card(s, Inches(6.5), Inches(1.4), Inches(6.5), Inches(5.55))
add_rect(s, Inches(6.5), Inches(1.4), Inches(6.5), Inches(0.45), GOLD)
add_textbox(s, "Digital Health Indicators", Inches(6.66), Inches(1.45),
            Inches(6.2), Inches(0.36), 13, True, WHITE)
metrics = [
    ("Website Page Speed (Mobile)",   "~38 / 100",  "Needs urgent optimisation",  DANGER),
    ("Website Page Speed (Desktop)",  "~65 / 100",  "Below industry benchmark",    GOLD),
    ("SEO Coverage",                  "~55 / 100",  "Thin content, missing schema", GOLD),
    ("Mobile Responsiveness",         "Partial",    "Some breakpoint issues",       GOLD),
    ("SSL / HTTPS",                   "✓ Active",   "Valid certificate in place",   SUCCESS),
    ("Mobile App (iOS)",              "Live",       "App Store – property search",  SUCCESS),
    ("Social Media Integration",      "Minimal",    "No live feed on website",      DANGER),
    ("AI Features (current)",         "Virtual Tours","Limited AI deployment",      GOLD),
]
ty = Inches(2.0)
for lbl, val, note, col in metrics:
    add_textbox(s, lbl,  Inches(6.66), ty, Inches(2.9), Inches(0.3), 10, True, NAVY)
    add_textbox(s, val,  Inches(9.6),  ty, Inches(1.2), Inches(0.3), 10, True, col, PP_ALIGN.CENTER)
    add_textbox(s, note, Inches(10.85),ty, Inches(2.0), Inches(0.3),  9, False, MID_GREY, italic=True)
    ty += Inches(0.52)

# ════════════════════════════════════════════════════════════════
# SLIDE 4 — WEBSITE AUDIT FINDINGS
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Website Audit Findings", "befittingproperties.com – identified gaps & opportunities"); footer(s)

audit_cols = [
    ("Performance", MID_GREEN, [
        "Mobile load time estimated >5 s",
        "Heavy unoptimised image assets",
        "No CDN – assets from single origin",
        "JS render-blocking on first paint",
        "No lazy-loading on listing images",
    ]),
    ("UX / Design", DARK_GREEN, [
        "Navigation hierarchy unclear",
        "Property search limited & slow",
        "No saved searches / wishlists",
        "CTA buttons lack visual weight",
        "Inconsistent mobile typography",
    ]),
    ("SEO & Content", GOLD, [
        "Meta descriptions absent >65 % pages",
        "No structured data / schema markup",
        "Blog / insights section thin / empty",
        "Listing descriptions too brief for SEO",
        "No Google Business Profile linked",
    ]),
    ("Security & Legal", DANGER, [
        "Cookie consent banner absent",
        "NDPR / data privacy notice weak",
        "Form validation needs hardening",
        "Missing security headers (CSP, X-Frame)",
        "No HTTPS enforcement on all sub-paths",
    ]),
]
for i, (title, col, bullets) in enumerate(audit_cols):
    x = Inches(0.4) + i * Inches(3.22)
    y = Inches(1.45)
    card(s, x, y, Inches(3.05), Inches(5.5))
    add_rect(s, x, y, Inches(3.05), Inches(0.45), col)
    add_textbox(s, title, x + Inches(0.12), y + Inches(0.06),
                Inches(2.8), Inches(0.36), 12, True, WHITE)
    ty = y + Inches(0.62)
    for b in bullets:
        add_rect(s, x + Inches(0.17), ty + Inches(0.08), Inches(0.1), Inches(0.1), col)
        add_textbox(s, b, x + Inches(0.37), ty, Inches(2.55), Inches(0.5),
                    10, False, NAVY)
        ty += Inches(0.72)

add_rect(s, Inches(0.4), Inches(6.83), Inches(12.5), Inches(0.055), GOLD)
add_textbox(s, "Top Priority: Mobile performance + SEO content + NDPR/cookie compliance – highest ROI quickwins",
            Inches(0.4), Inches(6.88), Inches(12.5), Inches(0.28),
            10.5, True, NAVY, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 5 — WEB DEVELOPMENT ROADMAP
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Web Development Roadmap", "Modernising befittingproperties.com for performance & growth"); footer(s)

card(s, Inches(0.4), Inches(1.4), Inches(4.55), Inches(5.6))
add_rect(s, Inches(0.4), Inches(1.4), Inches(4.55), Inches(0.45), BLUE)
add_textbox(s, "Recommended Tech Stack", Inches(0.56), Inches(1.45),
            Inches(4.3), Inches(0.36), 13, True, WHITE)
stack = [
    ("Frontend",    "Next.js 15 (React) – SSR & static generation"),
    ("CMS",         "Contentful / Sanity.io – headless, developer-friendly"),
    ("Hosting",     "Vercel + AWS CloudFront CDN (Nigeria edge nodes)"),
    ("Search",      "Algolia – instant property faceted search"),
    ("Maps",        "Google Maps API + geo-tagged listings"),
    ("Mobile",      "React Native – unified iOS/Android app"),
    ("Payments",    "Flutterwave / Paystack integration"),
    ("Analytics",   "GA4 + Hotjar heatmaps + Meta Pixel"),
    ("Auth",        "Auth0 – investor & tenant portal login"),
]
ty = Inches(2.0)
for k, v in stack:
    add_textbox(s, k, Inches(0.6), ty, Inches(1.45), Inches(0.28), 10, True, BLUE)
    add_textbox(s, v, Inches(2.1), ty, Inches(2.75), Inches(0.28), 10, False, NAVY)
    ty += Inches(0.49)

card(s, Inches(5.2), Inches(1.4), Inches(7.75), Inches(5.6))
add_rect(s, Inches(5.2), Inches(1.4), Inches(7.75), Inches(0.45), DARK_GREEN)
add_textbox(s, "Key Improvements & New Features",
            Inches(5.36), Inches(1.45), Inches(7.4), Inches(0.36), 13, True, WHITE)
improvements = [
    ("Speed & Core Web Vitals (Target 90+ Score)",
     "WebP images, code-split JS, CDN edge delivery, lazy-load listings, pre-fetch on hover."),
    ("Advanced Property Search & Discovery",
     "Faceted filters (type, price, size, location). Map-based browsing. Saved searches & alerts."),
    ("Investor & Tenant Self-Service Portal",
     "Secure login: documents, rent history, maintenance requests, ROI dashboards."),
    ("Integrated Payment Gateway",
     "Flutterwave/Paystack for rent, deposits, and investment subscriptions — in-app and web."),
    ("SEO Architecture & Schema Markup",
     "Property schema, local SEO per Lagos/Abuja node, rich snippets, auto sitemaps."),
    ("NDPR / Cookie Compliance",
     "Consent banner, data subject request flow, privacy policy modernisation."),
    ("Social & App Deep-Link Integration",
     "Instagram feed on site, auto-post new listings, deep links from web → mobile app."),
]
ty = Inches(1.98)
for title, desc in improvements:
    add_rect(s, Inches(5.36), ty + Inches(0.12), Inches(0.08), Inches(0.22), GOLD)
    add_textbox(s, title, Inches(5.56), ty, Inches(7.1), Inches(0.28),
                11, True, DARK_GREEN)
    add_textbox(s, desc,  Inches(5.56), ty + Inches(0.3), Inches(7.1), Inches(0.35),
                10, False, MID_GREY)
    ty += Inches(0.72)

# ════════════════════════════════════════════════════════════════
# SLIDE 6 — AI AUTOMATION OPPORTUNITY MAP
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "AI Automation Opportunity Map", "Where AI delivers immediate value for Befitting Group"); footer(s)

opps = [
    (DARK_GREEN, "Property Listings & Marketing", [
        "AI auto-generates listing copy from spec sheet",
        "Photo enhancement, tagging & floor plan AI",
        "Dynamic pricing engine (demand-driven rents)",
        "Bulk import & normalise data from Excel/PDF",
        "Auto-post new listings to Instagram & Facebook",
    ]),
    (GOLD, "Client & Tenant Communication", [
        "24/7 AI chatbot on website + WhatsApp bot",
        "Auto-draft tenancy renewal & payment reminders",
        "Maintenance ticket classification & routing",
        "Personalised property recommendations via email",
        "Multi-language support (English, Pidgin, French)",
    ]),
    (BLUE, "Market Intelligence & Analytics", [
        "Auto-scrape & summarise competitor listings",
        "Weekly AI-generated Lagos/Abuja market reports",
        "Predictive vacancy & demand forecasting",
        "Portfolio performance dashboards (real-time)",
        "Investor ROI reporting – fully automated",
    ]),
    (MID_GREEN, "Operations & Back-Office Automation", [
        "Invoice & accounts payable processing",
        "AI contract review – clause flagging & alerts",
        "HR: CV screening, onboarding document generation",
        "Meeting notes → action items (AI transcription)",
        "KYC / AML document verification automation",
    ]),
]
for i, (col, title, bullets) in enumerate(opps):
    c, r = i % 2, i // 2
    x = Inches(0.4)  + c * Inches(6.5)
    y = Inches(1.45) + r * Inches(2.9)
    card(s, x, y, Inches(6.2), Inches(2.7))
    add_rect(s, x, y, Inches(6.2), Inches(0.45), col)
    add_textbox(s, title, x + Inches(0.15), y + Inches(0.07),
                Inches(5.9), Inches(0.36), 12, True, WHITE)
    ty = y + Inches(0.58)
    for b in bullets:
        add_rect(s, x + Inches(0.18), ty + Inches(0.07), Inches(0.1), Inches(0.1), col)
        add_textbox(s, b, x + Inches(0.38), ty, Inches(5.65), Inches(0.38),
                    10.5, False, NAVY)
        ty += Inches(0.43)

# ════════════════════════════════════════════════════════════════
# SLIDE 7 — BENEFITS OF AUTOMATION (hero)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Benefits of AI Automation", "Measurable impact across the Befitting Group business"); footer(s)

for i, (val, lbl) in enumerate([
    ("65–75%",  "Reduction in\nrepetitive admin tasks"),
    ("4×",      "Faster property\nlisting publication"),
    ("45%",     "Lower cost per\nclient enquiry"),
    ("24/7",    "Client service via\nAI chatbot (0 overtime)"),
]):
    stat_card(s, val, lbl, Inches(0.4) + i * Inches(3.2), Inches(1.45),
              w=Inches(3.0), h=Inches(1.6))

benefits = [
    ("Speed to Market", [
        "New listings live in minutes, not days",
        "AI writes SEO-optimised description from a spec",
        "Photos auto-tagged, watermarked & resized",
        "Social posts drafted & scheduled automatically",
    ]),
    ("Cost Efficiency", [
        "Cut admin staff data-entry time by >55 %",
        "Fewer errors = fewer costly tenant disputes",
        "Lower dependence on external marketing agencies",
        "Scale listing volume without scaling headcount",
    ]),
    ("Client Experience", [
        "Instant 24/7 WhatsApp & web chatbot responses",
        "Personalised property match notifications",
        "Auto viewing confirmations, reminders & follow-ups",
        "Self-service portal for rent & maintenance",
    ]),
    ("Competitive & Strategic Edge", [
        "AI-driven Lagos/Abuja market intelligence",
        "First-mover in Nigerian PropTech AI adoption",
        "Attract diaspora investors with automated reports",
        "Proprietary data asset builds year-on-year",
    ]),
]
for i, (title, bullets) in enumerate(benefits):
    c, r = i % 2, i // 2
    x = Inches(0.4)  + c * Inches(6.5)
    y = Inches(3.22) + r * Inches(1.87)
    bullet_card(s, title, bullets, x, y, Inches(6.2), Inches(1.7),
                hdr_bg=DARK_GREEN if r == 0 else MID_GREEN)

# ════════════════════════════════════════════════════════════════
# SLIDE 8 — ROI & BUSINESS CASE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "ROI & Business Case", "Projected returns from web + AI investment"); footer(s)

card(s, Inches(0.4), Inches(1.4), Inches(5.8), Inches(5.62))
add_rect(s, Inches(0.4), Inches(1.4), Inches(5.8), Inches(0.45), DARK_GREEN)
add_textbox(s, "Indicative Investment (Year 1 – USD)",
            Inches(0.56), Inches(1.45), Inches(5.5), Inches(0.36), 13, True, WHITE)
invest = [
    ("Website Rebuild (Next.js + Headless CMS)",  "$12 000 – $18 000"),
    ("AI Chatbot & WhatsApp Automation",           " $4 500 – $7 500"),
    ("Listing Automation & AI Copywriting Stack",  " $3 000 – $5 500"),
    ("Mobile App Enhancement (React Native)",      " $6 000 – $10 000"),
    ("SEO, NDPR Compliance & Speed Sprint",        " $2 500 – $4 000"),
    ("Training, Docs & Change Management",         " $1 500 – $2 500"),
    ("Monthly support / hosting / AI APIs (est.)", " $800 – $1 400 / mo"),
    ("TOTAL YEAR 1 (once-off + 12 mo support)",   "$39 100 – $64 300"),
]
ty = Inches(2.06)
for i, (item, cost) in enumerate(invest):
    last = i == len(invest) - 1
    if last:
        add_rect(s, Inches(0.4), ty - Inches(0.06), Inches(5.8), Inches(0.06), GOLD)
        ty += Inches(0.1)
    add_textbox(s, item, Inches(0.6), ty, Inches(3.6), Inches(0.3),
                10 if not last else 11, last, NAVY)
    add_textbox(s, cost, Inches(4.25), ty, Inches(1.75), Inches(0.3),
                10 if not last else 11, last,
                SUCCESS if last else MID_GREY, PP_ALIGN.RIGHT)
    ty += Inches(0.52)

card(s, Inches(6.5), Inches(1.4), Inches(6.5), Inches(5.62))
add_rect(s, Inches(6.5), Inches(1.4), Inches(6.5), Inches(0.45), SUCCESS)
add_textbox(s, "Projected Returns & Savings (USD / Year)",
            Inches(6.66), Inches(1.45), Inches(6.2), Inches(0.36), 13, True, WHITE)
returns = [
    ("Admin time saved (FTE equivalent)",           "1–2 FTE / yr  ≈  $18 000+"),
    ("Reduced marketing agency dependency",          "~$8 000 / year"),
    ("Increased organic leads (SEO + UX uplift)",   "+30–45 % qualified enquiries"),
    ("Faster listing → tenancy conversion",         "Est. −12 days avg. vacancy"),
    ("Chatbot deflects call centre volume",         "−35 % inbound call load"),
    ("Investor portal subscription / fee revenue",  "Potential $6 000+ / year"),
    ("Diaspora investor reach (digital-first)",     "New market segment unlocked"),
    ("ESTIMATED YEAR 1 SAVINGS / UPSIDE",           "$32 000 – $55 000+"),
]
ty = Inches(2.06)
for i, (item, val) in enumerate(returns):
    last = i == len(returns) - 1
    if last:
        add_rect(s, Inches(6.5), ty - Inches(0.06), Inches(6.5), Inches(0.06), GOLD)
        ty += Inches(0.1)
    add_textbox(s, item, Inches(6.66), ty, Inches(3.95), Inches(0.3),
                10 if not last else 11, last, NAVY)
    add_textbox(s, val, Inches(10.65), ty, Inches(2.2), Inches(0.3),
                10 if not last else 11, last,
                SUCCESS if last else MID_GREY, PP_ALIGN.RIGHT)
    ty += Inches(0.52)

add_textbox(s, "Estimated breakeven: 10–16 months  |  Automation savings compound every year thereafter",
            Inches(0.4), Inches(7.07), Inches(12.5), Inches(0.28),
            10.5, True, DARK_GREEN, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 9 — IMPLEMENTATION TIMELINE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Implementation Timeline", "Phased 12-month delivery plan"); footer(s)

phases = [
    ("Phase 1\nMon 1–3",  GOLD,      [
        "Full website performance audit",
        "Speed & SEO quick-win sprint",
        "NDPR / cookie consent go-live",
        "AI chatbot MVP launch on site",
        "WhatsApp bot pilot for enquiries",
    ]),
    ("Phase 2\nMon 4–6",  DARK_GREEN,[
        "Next.js rebuild & CMS migration",
        "Advanced property search (Algolia)",
        "Listing automation + AI copywriting",
        "Instagram auto-post integration",
        "Mobile app v2 (React Native uplift)",
    ]),
    ("Phase 3\nMon 7–9",  BLUE,      [
        "Investor & tenant portal launch",
        "Payment gateway (Flutterwave) live",
        "AI market intelligence reports",
        "Invoice & contract automation",
        "Staff enablement & training",
    ]),
    ("Phase 4\nMon 10–12",MID_GREEN, [
        "Predictive vacancy analytics",
        "Advanced chatbot fine-tuning",
        "ERP / accounting API integration",
        "KYC / AML document AI",
        "Full audit, review & optimise",
    ]),
]
for i, (label, col, items) in enumerate(phases):
    x = Inches(0.4) + i * Inches(3.22)
    y = Inches(1.45)
    card(s, x, y, Inches(3.05), Inches(5.6))
    add_rect(s, x, y, Inches(3.05), Inches(0.72), col)
    add_textbox(s, label, x, y + Inches(0.06), Inches(3.05), Inches(0.62),
                13, True, WHITE, PP_ALIGN.CENTER)
    if i < 3:
        add_rect(s, x + Inches(3.07), y + Inches(0.3), Inches(0.14), Inches(0.06), MID_GREY)
    ty = y + Inches(0.9)
    for item in items:
        add_rect(s, x + Inches(0.18), ty + Inches(0.08), Inches(0.1), Inches(0.1), col)
        add_textbox(s, item, x + Inches(0.38), ty, Inches(2.55), Inches(0.42),
                    10.5, False, NAVY)
        ty += Inches(0.84)

# ════════════════════════════════════════════════════════════════
# SLIDE 10 — WHY AUTOMATE NOW?
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, DARK_GREEN)
add_rect(s, 0, 0, Inches(0.55), SLIDE_H, GOLD)
footer(s)

add_textbox(s, "Why Automate — Now?",
            Inches(0.85), Inches(0.48), Inches(11), Inches(0.65),
            30, True, GOLD)
add_textbox(s, "Nigerian PropTech is growing at 23 % CAGR. The window to lead — not follow — is open today.",
            Inches(0.85), Inches(1.18), Inches(11.5), Inches(0.55),
            13.5, False, WHITE, italic=True)

reasons = [
    (GOLD,     "Nigerian PropTech is accelerating fast",
     "Competitors are deploying AI-powered search, chatbots and virtual tours right now. First movers capture SEO authority and brand preference that late adopters cannot recover for years."),
    (AMBER,    "Labour costs and inflation keep rising",
     "Every manual process — data entry, email follow-up, report writing — costs more each year. Automation freezes that cost curve and redeploys talent to relationship and revenue work."),
    (SUCCESS,  "Clients expect instant, 24/7 responses",
     "Over 70 % of African property enquirers move to the first agent who responds. An AI WhatsApp bot means Befitting is always first — without overtime or weekend pay."),
    (BLUE,     "AI tools are now genuinely affordable",
     "GPT-4o API costs have dropped >90 % since 2023. Enterprise-grade automation is no longer reserved for REITs — it is accessible to ambitious mid-market property groups today."),
    (MID_GREEN,"Data compounds in competitive value",
     "Every automated interaction generates structured data. Over 2–3 years, Befitting builds a proprietary Lagos/Abuja dataset for pricing, forecasting, and investor reporting that no competitor can replicate."),
]
ty = Inches(1.9)
for i, (col, title, desc) in enumerate(reasons):
    add_rect(s, Inches(0.85), ty, Inches(0.55), Inches(0.52), col)
    add_textbox(s, str(i + 1), Inches(0.85), ty,
                Inches(0.55), Inches(0.52), 17, True, DARK_GREEN, PP_ALIGN.CENTER)
    add_textbox(s, title, Inches(1.55), ty, Inches(11.2), Inches(0.3),
                12, True, WHITE)
    add_textbox(s, desc, Inches(1.55), ty + Inches(0.3), Inches(11.3), Inches(0.52),
                10.5, False, RGBColor(0xB0, 0xD8, 0xBB))
    ty += Inches(0.98)

# ════════════════════════════════════════════════════════════════
# SLIDE 11 — NEXT STEPS
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Next Steps", "From audit to action – what happens now"); footer(s)

steps = [
    ("Week 1",  GOLD,      "Stakeholder Alignment",
     "Present findings to leadership. Agree priority areas and budget envelope. Assign internal digital champion."),
    ("Week 2",  DARK_GREEN,"Kick-Off & Access",
     "Provide hosting, CMS, analytics and social access. Finalise project charter. Daily standups begin."),
    ("Week 3",  BLUE,      "Sprint 1 Begins",
     "Speed fixes, NDPR compliance and AI chatbot integration start simultaneously across parallel tracks."),
    ("Month 2", MID_GREEN, "Phase 1 Go-Live",
     "Chatbot + WhatsApp bot live. SEO/speed fixes deployed. First AI-generated listings in production."),
]
for i, (when, col, title, desc) in enumerate(steps):
    x = Inches(0.4) + i * Inches(3.22)
    y = Inches(1.5)
    card(s, x, y, Inches(3.05), Inches(4.1))
    add_rect(s, x, y, Inches(3.05), Inches(0.45), col)
    add_textbox(s, when, x, y + Inches(0.06), Inches(3.05), Inches(0.36),
                12, True, WHITE, PP_ALIGN.CENTER)
    add_textbox(s, title, x + Inches(0.15), y + Inches(0.6),
                Inches(2.75), Inches(0.36), 12, True, col)
    add_textbox(s, desc, x + Inches(0.15), y + Inches(1.05),
                Inches(2.75), Inches(2.9), 10.5, False, NAVY)

card(s, Inches(0.4), Inches(5.82), Inches(12.5), Inches(1.1), fill=DARK_GREEN, border=False)
add_textbox(s, "Ready to transform Befitting Group's digital presence?",
            Inches(0.65), Inches(5.9), Inches(7.5), Inches(0.4),
            14, True, GOLD)
add_textbox(s, "Contact Genos Apollo to schedule your discovery workshop and technical deep-dive.",
            Inches(0.65), Inches(6.3), Inches(7.5), Inches(0.38),
            11, False, WHITE)
add_textbox(s, "befittingproperties.com  •  @befittingproperties",
            Inches(8.5), Inches(6.05), Inches(4.5), Inches(0.38),
            10, False, GOLD, PP_ALIGN.RIGHT, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 12 — BACK COVER
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, DARK_GREEN)
add_rect(s, 0, 0, Inches(0.75), SLIDE_H, GOLD)
add_rect(s, Inches(0.75), Inches(3.55), SLIDE_W - Inches(0.75), Inches(0.06), GOLD)

add_textbox(s, "Thank You",
            Inches(1.1), Inches(1.95), Inches(11), Inches(1.25),
            52, True, WHITE)
add_textbox(s, "Building Africa's property future — smarter, faster, and more connected.",
            Inches(1.1), Inches(3.15), Inches(10.5), Inches(0.55),
            15, False, AMBER, italic=True)
add_textbox(s, "Befitting Group  •  befittingproperties.com  •  Lagos, Nigeria",
            Inches(1.1), Inches(4.05), Inches(9), Inches(0.4),
            12, False, RGBColor(0x88, 0xBB, 0x99))
add_textbox(s, "Prepared by Genos Apollo  |  May 2026  |  Confidential",
            Inches(1.1), Inches(6.72), Inches(10), Inches(0.35),
            10, False, MID_GREY, italic=True)

# ── Save ──────────────────────────────────────────────────────
out = r"c:\Users\deevansho\Desktop\GenosApollp clients\Befitting_WebDev_AI_Automation_Audit.pptx"
prs.save(out)
print(f"Saved: {out}")
