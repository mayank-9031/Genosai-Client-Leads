import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.parent / "utils"))
from ppt_theme import (
    NAVY, DARK_BLUE, ACCENT, GOLD, RED, GREEN, PURPLE, ORANGE, TEAL, CYAN,
    WHITE, LGREY, MGREY, W, H,
    bg, rect, tb, aline, bullets, tag,
    slide_header, slide_footer, score_badge,
    stat_cards, top_bar_card, phase_card, before_after, genos_logo_block,
)

ROOT = Path(__file__).parent.parent.parent
prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]

# SLIDE 1 - COVER
s1 = prs.slides.add_slide(BLANK)
bg(s1)
rect(s1, 0, 0, Inches(0.12), H, TEAL)
rect(s1, Inches(8.5), 0, Inches(4.83), Inches(2.2), DARK_BLUE)
tb(s1, "GENOS AI", Inches(9.2), Inches(0.4), Inches(3.5), Inches(0.8),
   size=28, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
tb(s1, "www.genosai.tech", Inches(9.2), Inches(1.05), Inches(3.5), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
tb(s1, "HOUSING FOR COMMUNITIES, INC.", Inches(0.55), Inches(1.6),
   Inches(9), Inches(0.6), size=13, bold=False, color=TEAL)
tb(s1, "Website & AI Audit\nfor a 4-Person Mission", Inches(0.55), Inches(2.05),
   Inches(9), Inches(1.8), size=44, bold=True, color=WHITE)
aline(s1, Inches(3.85), color=TEAL, w=Inches(3))
tb(s1, "Prepared exclusively for Andrea N. Frymire, CCIM - President & CEO",
   Inches(0.55), Inches(4.15), Inches(9), Inches(0.5), size=14, color=LGREY)
rect(s1, 0, Inches(6.5), W, Inches(1.0), DARK_BLUE)
tb(s1, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699",
   Inches(0.55), Inches(6.6), Inches(9), Inches(0.5), size=12, color=MGREY)
tb(s1, "May 11, 2026  |  CONFIDENTIAL", Inches(9.5), Inches(6.6), Inches(3.3), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
rect(s1, Inches(10.5), Inches(3.3), Inches(2.3), Inches(2.5), DARK_BLUE)
tb(s1, "AUDIT SCORE", Inches(10.55), Inches(3.4), Inches(2.2), Inches(0.5),
   size=11, bold=True, color=MGREY, align=PP_ALIGN.CENTER)
tb(s1, "5", Inches(10.55), Inches(3.75), Inches(2.2), Inches(1.1),
   size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
tb(s1, "/ 10", Inches(10.55), Inches(4.7), Inches(2.2), Inches(0.4),
   size=18, color=LGREY, align=PP_ALIGN.CENTER)
tag(s1, "HIGH PRIORITY", Inches(10.75), Inches(5.1), fill=RED, size=10)


# SLIDE 2 - AT A GLANCE
s2 = prs.slides.add_slide(BLANK)
bg(s2)
rect(s2, 0, 0, Inches(0.12), H, TEAL)
tb(s2, "AT A GLANCE", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=TEAL)
tb(s2, "Andrea Frymire and Housing for Communities",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=30, bold=True, color=WHITE)
aline(s2, Inches(1.35), color=TEAL)

stats = [
    ("25+ yrs", "In commercial real estate\n& affordable housing"),
    ("4 people", "Running training,\nconsulting, development"),
    ("3 pillars", "Training | Consulting\n| Development"),
    ("8 days", "Until Senior Housing\nConference (May 19)"),
    ("CCIM", "+ 2021 CRE Advocate\nof the Year"),
]
cw = Inches(2.3); ch = Inches(2.2); gap = Inches(0.18)
for i, (num, label) in enumerate(stats):
    l = Inches(0.55) + i * (cw + gap)
    rect(s2, l, Inches(1.8), cw, ch, DARK_BLUE)
    tb(s2, num, l + Inches(0.1), Inches(1.95), cw - Inches(0.2), Inches(1.1),
       size=20, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    tb(s2, label, l + Inches(0.1), Inches(3.05), cw - Inches(0.2), Inches(0.7),
       size=11, color=LGREY, align=PP_ALIGN.CENTER)

rect(s2, Inches(0.55), Inches(4.2), Inches(12.23), Inches(0.6), DARK_BLUE)
tb(s2, "501(c)(3) non-profit  |  Founded 2023, OKC  |  Walmart & Sam's Club Spark Good Local Grants  |  Partners: Earthly Dwelling, Housing Forward, NCOA",
   Inches(0.65), Inches(4.28), Inches(12), Inches(0.45), size=11, color=LGREY, align=PP_ALIGN.CENTER)

tb(s2, "ANDREA'S BACKGROUND", Inches(0.55), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=TEAL)
bullets(s2, [
    "President & CEO since 2023 - 25+ years CRE & affordable housing",
    "BS Finance Radford U; Master's Real Estate; CCIM designation",
    "Founder & Inaugural Board President, Oklahoma Coalition for Affordable Housing",
    "Board President again 2025 for Coalition's 10-year anniversary"],
        Inches(0.55), Inches(5.35), Inches(5.8), Inches(1.5), size=12, color=LGREY)

tb(s2, "WHAT HOUSING FOR COMMUNITIES DOES", Inches(6.8), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=TEAL)
bullets(s2, [
    "Training: workshops, CEUs for housing professionals",
    "Consulting: affordable housing strategy & development",
    "Development: production & preservation of affordable units",
    "Senior Housing Conference May 19, 2026 - Metro Tech OKC"],
        Inches(6.8), Inches(5.35), Inches(5.8), Inches(1.5), size=12, color=LGREY)


# SLIDE 3 - STRENGTHS
s3 = prs.slides.add_slide(BLANK)
bg(s3)
rect(s3, 0, 0, Inches(0.12), H, GREEN)
tb(s3, "WHAT'S WORKING", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s3, "A Mission-Driven Organization with Real Credibility - and a Conference in 8 Days",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.8), size=24, bold=True, color=WHITE)
aline(s3, Inches(1.45), color=GREEN)

strengths = [
    ("Andrea's Credibility Opens Doors That Small Non-Profits Cannot",
     "25 years in commercial real estate, CCIM designation, Master's in Real Estate, Broker/Owner of Oklahoma Investment Realty, Founder and Board President of the Oklahoma Coalition for Affordable Housing - now in its 10-year anniversary year. 2021 Commercial Real Estate Advocate of the Year. Andrea has helped form 8+ real estate non-profits. The credibility ceiling for Housing for Communities is much higher than the current website signals - the website is the bottleneck, not the brand."),
    ("The Senior Housing Conference Is a Real Asset, Not a One-Off",
     "Safe, Stable & Supported on May 19 at Metro Tech OKC brings together housing professionals, service providers, advocates, developers, and community leaders for one day with CEUs, networking, and policy content. The AGING20 promo, the early-bird ladder, and the speaker line-up all show operational discipline. This is a platform that should pay back into Training and Consulting pipeline year-round - if the website infrastructure captures and nurtures attendees afterward."),
    ("Walmart & Sam's Club Spark Good Grant Recognition Validates the Mission",
     "Receiving Spark Good Local Grants from both Walmart Neighborhood Market Midwest City and Sam's Club Distribution Center in the same year is not a coincidence - it is corporate validation of the work. Most 2-year-old non-profits do not have this kind of credibility from major corporate funders. This is a fundraising story the website is currently not telling."),
    ("Three Service Pillars Are Clearly Defined",
     "Training, Consulting, Development. Most non-profits struggle to articulate even one of these. Having three with clear scope is a strategic strength - and it makes the website automation work much more straightforward. Each pillar has a different prospect, a different qualifying question, a different routing path. The structure exists; it just is not surfaced through the form layer."),
    ("Cross-Organization Reach via the OK Coalition Multiplier",
     "Andrea is simultaneously CEO of Housing for Communities AND Board President of the Oklahoma Coalition for Affordable Housing during its 10-year anniversary year. That is two organizations of audience reach for one website upgrade. Coalition events, members, and policy momentum can be cross-promoted through Housing for Communities' digital channels with no incremental marketing spend."),
]
for i, (title, desc) in enumerate(strengths):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s3, Inches(0.55), t, Inches(0.06), Inches(0.7), GREEN)
    tb(s3, title, Inches(0.75), t, Inches(11.8), Inches(0.38), size=13, bold=True, color=WHITE)
    tb(s3, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.6), size=11, color=LGREY)


# SLIDE 4 - THE THREE GAPS
s4 = prs.slides.add_slide(BLANK)
bg(s4)
rect(s4, 0, 0, Inches(0.12), H, RED)
tb(s4, "THE THREE GAPS", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=RED)
tb(s4, "A 4-Person Team. Three Service Lines. One Generic Inbox.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=26, bold=True, color=WHITE)
aline(s4, Inches(1.35), color=RED)

gap_cards = [
    ("GAP 1\nINBOX BOTTLENECK\n(EVENT URGENCY)", RED,
     "Where the May 19 conference inquiries land",
     [
         "8 days from this audit, a major event - questions about AGING20, CEU paperwork, parking, agenda all hit info@",
         "No chatbot, no FAQ, no self-service - every question requires Madeline or Bailey to type a reply",
         "Routine conference questions absorb hours of team capacity that should go to attendee experience and content",
         "Same gap will repeat at every future conference - and there's a 2027 to plan",
     ]),
    ("GAP 2\nONE FORM FOR\nTHREE BUSINESSES", GOLD,
     "Training, Consulting, Development - all the same intake",
     [
         "4-field contact form (first, last, email, message) handles training enrollment, consulting RFPs, and development partnerships identically",
         "No qualifying detail captured - every inquiry needs a back-and-forth before it can be routed",
         "Bailey (Development) and Andrea (Consulting) get inquiries with no project type, timeline, geography, or budget signal",
         "Qualified consulting leads get the same treatment as a CEU question",
     ]),
    ("GAP 3\nMISSION CREDIBILITY\nUNDERLEVERAGED", ACCENT,
     "Andrea's voice and the donor story off the site",
     [
         "Spark Good grants from Walmart and Sam's Club - no donor CTA, no recurring giving flow, no donor wall on the site",
         "Andrea's CCIM expertise, founding role, and 2021 Advocate award - barely visible as web content beyond a board page",
         "No newsletter capture, no policy alert signup - email list compounds across years if started today",
         "Coalition anniversary year provides cross-promotion fuel that the site cannot currently deploy",
     ]),
]
cw3 = Inches(3.95); gap3i = Inches(0.18)
for i, (title, color, subtitle, pts) in enumerate(gap_cards):
    l = Inches(0.55) + i * (cw3 + gap3i)
    rect(s4, l, Inches(1.8), cw3, Inches(0.5), color)
    tb(s4, title, l + Inches(0.12), Inches(1.85), cw3 - Inches(0.24), Inches(0.42),
       size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    rect(s4, l, Inches(2.3), cw3, Inches(4.6), DARK_BLUE)
    tb(s4, subtitle, l + Inches(0.12), Inches(2.38), cw3 - Inches(0.24), Inches(0.35),
       size=10, italic=True, color=color)
    bullets(s4, pts, l + Inches(0.12), Inches(2.78), cw3 - Inches(0.24), Inches(3.9),
            size=11, color=LGREY, dot=color)

rect(s4, Inches(0.55), Inches(7.05), Inches(12.23), Inches(0.3), DARK_BLUE)
tb(s4, "The team isn't small for the size of the mission. The website is small for the size of the team.",
   Inches(0.7), Inches(7.08), Inches(12), Inches(0.25), size=11, italic=True, color=GOLD)


# SLIDE 5 - WEBSITE DESIGN AUDIT
s5 = prs.slides.add_slide(BLANK)
bg(s5)
rect(s5, 0, 0, Inches(0.12), H, PURPLE)
tb(s5, "WEBSITE DESIGN AUDIT", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=PURPLE)
tb(s5, "From Brochure Site to Operational Front Door.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.8), size=32, bold=True, color=WHITE)
aline(s5, Inches(1.45), color=PURPLE, w=Inches(1.5))

rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(3.6), DARK_BLUE)
rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(0.42), RED)
tb(s5, "HOUSINGFORCOMMUNITIES.ORG TODAY", Inches(0.65), Inches(1.92), Inches(3.6), Inches(0.38),
   size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "4-field contact form, no qualification, no routing",
    "No AI chatbot or live chat - all inquiries hit info@",
    "No FAQ or self-service knowledge base",
    "No conference registration self-service on-site",
    "No donor CTA, no recurring giving flow",
    "No newsletter signup, no policy alerts",
    "No automated CEU certificate delivery",
    "Andrea's thought leadership barely surfaced",
], Inches(0.7), Inches(2.42), Inches(3.6), Inches(2.9), size=11, color=LGREY, dot=RED)

tb(s5, "->", Inches(4.45), Inches(3.6), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(3.6), DARK_BLUE)
rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(0.42), GOLD)
tb(s5, "WHAT MISSION-DRIVEN SITES DO", Inches(5.1), Inches(1.92), Inches(3.4), Inches(0.38),
   size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "AI agent answers FAQ + service questions 24/7",
    "Service-specific intake forms with qualifying detail",
    "Donor flow with one-time + recurring giving",
    "Event registration & post-event automation",
    "Newsletter capture with welcome sequence",
    "Automated CEU + certificate delivery",
    "Thought leadership content engine indexed for SEO",
    "Multilingual access where the population requires it",
], Inches(5.1), Inches(2.42), Inches(3.4), Inches(2.9), size=11, color=LGREY, dot=GOLD)

tb(s5, "->", Inches(8.68), Inches(3.6), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(3.6), DARK_BLUE)
rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(0.42), GREEN)
tb(s5, "WHAT GENOS AI BUILDS", Inches(9.35), Inches(1.92), Inches(3.4), Inches(0.38),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Conference AI agent live before May 19 (rapid deploy)",
    "3-track intake: Training, Consulting, Development",
    "Donor + recurring giving flow with Spark Good story",
    "Newsletter automation: OK housing policy welcome series",
    "Automated CEU certificate delivery to attendees",
    "Andrea voice-aligned thought leadership engine",
    "Spanish-language site version via AI translation",
    "Coalition cross-promotion module for anniversary year",
], Inches(9.35), Inches(2.42), Inches(3.4), Inches(2.9), size=11, color=LGREY, dot=GREEN)

rect(s5, Inches(0.55), Inches(5.72), Inches(12.23), Inches(1.28), DARK_BLUE)
tb(s5, "MISSION OPERATIONS BENCHMARK:", Inches(0.7), Inches(5.78), Inches(3.5), Inches(0.38),
   size=11, bold=True, color=PURPLE)
stats2 = [
    ("70-80%", "of routine inquiries\nabsorbed by AI agent"),
    ("3 tracks", "Training | Consulting\n| Development intake"),
    ("8 days", "until conference -\nagent live in week 1"),
    ("4 people", "team capacity preserved\nfor mission work"),
]
sw = Inches(2.7); sl = Inches(3.85)
for i, (num, label) in enumerate(stats2):
    xl = sl + i * sw * 0.95
    tb(s5, num, xl, Inches(5.75), Inches(2.5), Inches(0.55),
       size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s5, label, xl, Inches(6.22), Inches(2.5), Inches(0.7),
       size=10, color=LGREY, align=PP_ALIGN.CENTER)


# SLIDE 6 - AI AUTOMATION STACK
s6 = prs.slides.add_slide(BLANK)
bg(s6)
rect(s6, 0, 0, Inches(0.12), H, ACCENT)
tb(s6, "AI AUTOMATION STACK", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s6, "Four Layers That Buy Andrea's Team Their Time Back.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=26, bold=True, color=WHITE)
aline(s6, Inches(1.35), color=ACCENT)

opps = [
    ("Conference & Service\nAI Agent", TEAL,
     "Trained on the Senior Housing conference (agenda, AGING20 promo, CEU paperwork, parking, registration link) and the three service lines. Absorbs 70-80% of routine questions in real time. Routes the 20% needing humans with full context already captured. Live before May 19. Then expanded to handle Training/Consulting/Development questions year-round. One agent. Every visitor. 24/7."),
    ("Three-Track\nIntake Forms", GOLD,
     "Training: course interest, CEU need, organisation, role. Consulting: project type, geography, timeline, budget range. Development: partnership type, unit count, location, funding stage. Each form routes to Bailey (Development) or Andrea (Consulting) or Madeline (Training) with full qualifying detail. The first email back can be substantive instead of a triage question."),
    ("Donor & Newsletter\nAutomation", GREEN,
     "Spark Good Local Grants from Walmart and Sam's Club give credibility most 2-year-old non-profits cannot match. A donor CTA with one-time and recurring giving converts that credibility into capacity. Newsletter signup with an Oklahoma housing policy welcome sequence builds an email list that compounds: every visitor today is a year-round audience, not a one-time inquiry."),
    ("Thought Leadership\nContent Engine", PURPLE,
     "Andrea has 25 years of CRE expertise, a CCIM, founder credibility, and an Advocate-of-the-Year award. Almost none of that lives as indexed web content today. An AI-assisted content engine that takes Andrea's voice (extracted from her writing and speaking) and produces Oklahoma-housing-policy and senior-housing articles compounds her authority across years - and seeds the Coalition anniversary year with material."),
]
cw2 = Inches(2.95); gap2 = Inches(0.16)
for i, (title, color, desc) in enumerate(opps):
    l2 = Inches(0.55) + i * (cw2 + gap2)
    rect(s6, l2, Inches(1.8), cw2, Inches(0.06), color)
    rect(s6, l2, Inches(1.86), cw2, Inches(4.7), DARK_BLUE)
    tb(s6, title, l2 + Inches(0.12), Inches(1.98), cw2 - Inches(0.24), Inches(0.65),
       size=13, bold=True, color=color)
    tb(s6, desc, l2 + Inches(0.12), Inches(2.65), cw2 - Inches(0.24), Inches(3.65),
       size=11, color=LGREY)

rect(s6, Inches(0.55), Inches(6.7), Inches(12.23), Inches(0.45), DARK_BLUE)
tb(s6, "The team's job is housing policy and consulting. The automation layer handles everything that doesn't require a human housing expert.",
   Inches(0.7), Inches(6.75), Inches(12), Inches(0.35), size=11, italic=True, color=GOLD)


# SLIDE 7 - CONFERENCE TIMELINE URGENCY
s7 = prs.slides.add_slide(BLANK)
bg(s7)
rect(s7, 0, 0, Inches(0.12), H, ORANGE)
tb(s7, "CONFERENCE TIMELINE", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ORANGE)
tb(s7, "8 Days. What Could Be Live Before May 19.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s7, Inches(1.35), color=ORANGE, w=Inches(1.8))

rect(s7, Inches(0.55), Inches(1.65), Inches(12.23), Inches(1.55), DARK_BLUE)
bm_stats = [
    ("8 days", "Until conference\n(May 11 -> May 19)"),
    ("70-80%", "Routine inquiries\nabsorbed by agent"),
    ("24/7", "After-hours capture\nincluding weekends"),
    ("0 hours", "Of Madeline's time\non FAQ replies"),
    ("1 conf", "Becomes a year-round\nlead engine"),
]
ssw = Inches(2.35)
for i, (num, label) in enumerate(bm_stats):
    xl = Inches(0.7) + i * ssw
    tb(s7, num, xl, Inches(1.73), Inches(2.2), Inches(0.65),
       size=20, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    tb(s7, label, xl, Inches(2.33), Inches(2.2), Inches(0.5),
       size=10, color=LGREY, align=PP_ALIGN.CENTER)

rect(s7, Inches(0.55), Inches(3.4), Inches(5.9), Inches(3.25), DARK_BLUE)
tb(s7, "WHAT THE AGENT HANDLES BY MAY 19", Inches(0.7), Inches(3.47), Inches(5.6), Inches(0.35),
   size=11, bold=True, color=ORANGE)
bullets(s7, [
    "Is there still space available for the May 19 conference?",
    "What's the registration link / where do I pay?",
    "AGING20 promo code - how do I apply it?",
    "Does the conference offer continuing education credits?",
    "Where is Metro Tech / what's the parking situation?",
    "Who are the speakers / what's the agenda?",
    "Can I get a group rate for my organization?",
], Inches(0.7), Inches(3.85), Inches(5.6), Inches(2.55), size=11, color=LGREY, dot=ORANGE)

rect(s7, Inches(6.65), Inches(3.4), Inches(6.13), Inches(3.25), DARK_BLUE)
tb(s7, "WHAT THE AGENT ROUTES TO HUMANS", Inches(6.8), Inches(3.47), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=GREEN)
bullets(s7, [
    "Sponsorship inquiries - routed to Andrea with sponsor profile captured",
    "Press / media questions - routed to Madeline with publication context",
    "Speaker requests for future events - routed with topic & credentials",
    "Consulting RFPs that come in via conference traffic - routed to Andrea",
    "Development partnership pitches - routed to Bailey with project detail",
    "Custom training requests for organizations - routed to Madeline",
    "Anything not covered by the agent's knowledge - escalated with summary",
], Inches(6.8), Inches(3.85), Inches(5.8), Inches(2.55), size=11, color=LGREY, dot=GREEN)

rect(s7, Inches(0.55), Inches(6.8), Inches(12.23), Inches(0.38), DARK_BLUE)
tb(s7, "Most non-profits regret the conference inbox 48 hours before the event. The fix takes a week, not a month.",
   Inches(0.7), Inches(6.85), Inches(12), Inches(0.28), size=11, italic=True, color=ORANGE)


# SLIDE 8 - IMPLEMENTATION PLAN
s8 = prs.slides.add_slide(BLANK)
bg(s8)
rect(s8, 0, 0, Inches(0.12), H, GREEN)
tb(s8, "IMPLEMENTATION PLAN", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s8, "Conference First. Then Year-Round Pipeline. Then Scale.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=26, bold=True, color=WHITE)
aline(s8, Inches(1.35), color=GREEN)

phases = [
    ("PHASE 1", "Week 1 (Pre-Conf)", GOLD, [
        ("Conference Agent Live", "AI agent deployed on homepage trained on May 19 conference details, AGING20 promo, CEU info, registration, parking, agenda. Absorbs majority of inquiries in the 8 days before the event."),
        ("Routing Logic Configured", "Sponsorship, press, custom training routes set up. Madeline, Bailey, Andrea each receive only inquiries that need them - with full context already captured."),
        ("Team Training & Refinement", "Team monitors agent responses for first 48 hours. Refines knowledge base. By Day 3, agent runs autonomously through conference week."),
    ]),
    ("PHASE 2", "Weeks 2-4 (Post-Conf)", ACCENT, [
        ("Three-Track Intake Forms", "Training, Consulting, Development each get a qualifying intake form. Conference attendees become year-round Training and Consulting prospects."),
        ("CEU Certificate Automation", "Automated CEU certificate delivery to all conference attendees. Post-event satisfaction survey. Recording access link automation."),
        ("Newsletter Welcome Series", "OK housing policy welcome sequence live. Conference attendee email list seeds the year-round nurture engine."),
    ]),
    ("PHASE 3", "Months 2-3", GREEN, [
        ("Donor + Recurring Giving", "Donor CTA live with Spark Good grant credibility woven into the story. Recurring monthly giving option active."),
        ("Thought Leadership Engine", "First Andrea-voice articles on OK senior housing policy and affordable housing finance published and indexed. Coalition cross-promotion active."),
        ("Spanish-Language Site", "AI-translated Spanish version of core service pages live. Coalition anniversary year landing pages built. 2027 conference early-interest pipeline already capturing leads."),
    ]),
]
pw = Inches(3.9); gap4 = Inches(0.18)
for i, (phase, timing, color, items) in enumerate(phases):
    l3 = Inches(0.55) + i * (pw + gap4)
    rect(s8, l3, Inches(1.8), pw, Inches(0.52), color)
    tb(s8, phase, l3 + Inches(0.12), Inches(1.84), pw * 0.5 - Inches(0.12), Inches(0.42),
       size=13, bold=True, color=NAVY)
    tb(s8, timing, l3 + pw * 0.45, Inches(1.88), pw * 0.55 - Inches(0.15), Inches(0.38),
       size=11, color=NAVY, align=PP_ALIGN.RIGHT)
    rect(s8, l3, Inches(2.32), pw, Inches(4.6), DARK_BLUE)
    y_item = Inches(2.42)
    for title, desc in items:
        rect(s8, l3 + Inches(0.12), y_item + Inches(0.06), Inches(0.05), Inches(0.52), color)
        tb(s8, title, l3 + Inches(0.25), y_item, pw - Inches(0.38), Inches(0.32),
           size=12, bold=True, color=color)
        tb(s8, desc, l3 + Inches(0.25), y_item + Inches(0.3), pw - Inches(0.38), Inches(0.45),
           size=11, color=LGREY)
        y_item += Inches(0.95)


# SLIDE 9 - BEFORE vs AFTER
s9 = prs.slides.add_slide(BLANK)
bg(s9)
rect(s9, 0, 0, Inches(0.12), H, ACCENT)
tb(s9, "BEFORE vs AFTER", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s9, "What Actually Changes for Housing for Communities in 90 Days",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=26, bold=True, color=WHITE)
aline(s9, Inches(1.35))

comparisons = [
    ("Conference Inquiries",
     "Every question - registration, AGING20, CEU, parking, agenda - hits info@ and waits for Madeline or Bailey to type a reply. Capacity drains in the week before the event.",
     "AI agent answers 70-80% in real time. Sponsorship and press routed to humans with context. Madeline focused on attendee experience, not inbox triage."),
    ("Service Line Routing",
     "Training enrollment, consulting RFP, and development partnership all use the same 4-field form. Bailey and Andrea get inquiries with no qualifying detail.",
     "Three-track intake captures project type, timeline, budget range, geography. First reply is substantive - no triage round trip."),
    ("Donor Engagement",
     "Spark Good grants from Walmart and Sam's Club are credibility the website does not surface. No donor CTA, no recurring giving option, no donor wall.",
     "Donor flow live with grant-credibility story. Recurring monthly giving option drives compounding revenue. Donor recognition surfaced on the site."),
    ("Thought Leadership",
     "Andrea's CCIM expertise, founder role, and Advocate-of-Year award barely surface as web content. Most expertise stays on her LinkedIn.",
     "Andrea-voice content engine produces indexed articles on OK housing policy and senior housing. Coalition anniversary year cross-promoted. Authority compounds across years."),
    ("Year-Round Pipeline",
     "Conference is a one-time event. Attendee data sits in a registration export. No nurture, no list-building, no pathway to Training or Consulting.",
     "Conference attendee list seeds newsletter, Training enrollment, and Consulting pipeline. One conference becomes a year-round lead engine."),
]

lw = Inches(5.4)
rect(s9, Inches(0.55), Inches(1.65), lw, Inches(0.35), RED)
tb(s9, "TODAY", Inches(0.65), Inches(1.68), lw - Inches(0.2), Inches(0.28),
   size=11, bold=True, color=WHITE)
rect(s9, Inches(6.45), Inches(1.65), lw + Inches(0.43), Inches(0.35), GREEN)
tb(s9, "AFTER 90 DAYS", Inches(6.55), Inches(1.68), lw, Inches(0.28),
   size=11, bold=True, color=WHITE)

for i, (topic, before, after) in enumerate(comparisons):
    t = Inches(2.1) + i * Inches(1.02)
    rect(s9, Inches(0.55), t, Inches(12.23), Inches(0.22), DARK_BLUE)
    tb(s9, topic.upper(), Inches(0.65), t + Inches(0.02), Inches(12), Inches(0.2),
       size=9, bold=True, color=MGREY)
    tb(s9, before, Inches(0.65), t + Inches(0.24), lw - Inches(0.2), Inches(0.66),
       size=11, color=LGREY)
    tb(s9, after, Inches(6.55), t + Inches(0.24), lw, Inches(0.66),
       size=11, color=WHITE)


# SLIDE 10 - NEXT STEPS
s10 = prs.slides.add_slide(BLANK)
bg(s10)
rect(s10, 0, 0, Inches(0.12), H, TEAL)
rect(s10, Inches(7.5), 0, Inches(5.83), H, DARK_BLUE)
tb(s10, "NEXT STEPS", Inches(0.55), Inches(0.3), Inches(6.5), Inches(0.5),
   size=11, bold=True, color=TEAL)
tb(s10, "Let's Talk for\n15 Minutes.",
   Inches(0.55), Inches(0.8), Inches(6.8), Inches(2.0), size=48, bold=True, color=WHITE)
aline(s10, Inches(2.8), color=TEAL, w=Inches(2.5))
tb(s10, "Conference is May 19. We can have the agent live\nbefore then. The rest can come after the dust settles.",
   Inches(0.55), Inches(3.05), Inches(6.5), Inches(0.8), size=14, color=LGREY)

steps = [
    "Reply to this email - we'll set up a time",
    "Pre-conference: ship the conference AI agent in week 1",
    "Post-conference: 3-track intake + CEU automation in weeks 2-4",
    "Months 2-3: donor flow, thought leadership, Spanish version",
]
bullets(s10, steps, Inches(0.55), Inches(3.95), Inches(6.5), Inches(2.0), size=14, color=WHITE, dot=TEAL)

rect(s10, 0, Inches(6.5), Inches(7.5), Inches(1.0), DARK_BLUE)
tb(s10, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699  |  www.genosai.tech",
   Inches(0.55), Inches(6.6), Inches(6.8), Inches(0.5), size=11, color=MGREY)

tb(s10, "THREE GAPS.\nTHREE FIXES.", Inches(7.85), Inches(0.6), Inches(5.0), Inches(1.1),
   size=22, bold=True, color=TEAL)
aline(s10, Inches(1.65), color=TEAL)

summary_items = [
    ("Inbox Bottleneck", "AI agent absorbs 70-80% of routine inquiries in real time"),
    ("One-Form Triage", "3-track intake: Training, Consulting, Development with routing"),
    ("Donor Story Off-Site", "Donor CTA + recurring giving with Spark Good credibility"),
    ("Thought Leadership", "Andrea-voice content engine indexed for OK housing keywords"),
    ("Conference One-Off", "Year-round nurture: attendee list seeds Training & Consulting"),
    ("Coalition Anniversary", "Cross-promotion module for OK Coalition 10-year year"),
]
for i, (label, val) in enumerate(summary_items):
    t2 = Inches(1.85) + i * Inches(0.88)
    rect(s10, Inches(7.85), t2, Inches(1.5), Inches(0.72), NAVY)
    tb(s10, label, Inches(7.9), t2 + Inches(0.05), Inches(1.4), Inches(0.32),
       size=10, bold=True, color=TEAL)
    tb(s10, val, Inches(9.5), t2 + Inches(0.05), Inches(3.6), Inches(0.65),
       size=11, color=LGREY)


prs.save(ROOT / "output" / "decks" / "HousingForCommunities_Audit_GenosAI.pptx")
print("Saved: HousingForCommunities_Audit_GenosAI.pptx")
