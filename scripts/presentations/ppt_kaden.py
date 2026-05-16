import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.parent / "utils"))
from ppt_theme import (
    NAVY, DARK_BLUE, ACCENT, GOLD, RED, GREEN, PURPLE, ORANGE, TEAL, CYAN,
    WHITE, LGREY, MGREY, W, H,
    bg, rect, tb, aline, bullets, tag,
    slide_header, slide_footer, score_badge,
    stat_cards, top_bar_card, phase_card, before_after, genos_logo_block,
)

ROOT = Path(__file__).parent.parent.parent
prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]
# Client-specific colors
DGREY    = RGBColor(0x2C, 0x38, 0x4A)
AMBER    = RGBColor(0xF3, 0x9C, 0x12)
BLUE_ACC = RGBColor(0x0A, 0x7B, 0xF3)


# ── Local helpers ────────────────────────────────────────────────────────────

def header(slide, title, subtitle, accent=CYAN):
    rect(slide, 0, 0, Inches(0.12), H, accent)
    tb(slide, title.upper(), Inches(0.55 + 0.12), Inches(0.3),
       Inches(12), Inches(0.45), size=11, bold=True, color=accent)
    tb(slide, subtitle, Inches(0.55 + 0.12), Inches(0.72),
       Inches(12 - 0.12), Inches(0.65), size=22, bold=True, color=WHITE)
    aline(slide, Inches(1.38), color=accent, w=Inches(1.5))


def footer(slide,
           text='Genos AI  |  hello@genosai.tech  |  www.genosai.tech  |  +91 638-714-2699'):
    rect(slide, 0, Inches(7.05), W, Inches(0.45), DARK_BLUE)
    tb(slide, text, Inches(0.55), Inches(7.1), Inches(12.4), Inches(0.35),
       size=9, color=MGREY, align=PP_ALIGN.CENTER)


def card(slide, x, y, w, h, fill=DARK_BLUE):
    return rect(slide, x, y, w, h, fill)


def stat_card(slide, val, label, x, y, w=Inches(3.0), h=Inches(1.6), accent=CYAN):
    rect(slide, x, y, w, h, DARK_BLUE)
    rect(slide, x, y, w, Inches(0.05), accent)
    tb(slide, str(val), x + Inches(0.1), y + Inches(0.1),
       w - Inches(0.2), Inches(0.9), size=28, bold=True, color=accent,
       align=PP_ALIGN.CENTER)
    tb(slide, label, x + Inches(0.1), y + Inches(1.1),
       w - Inches(0.2), Inches(0.42), size=10, color=LGREY, align=PP_ALIGN.CENTER)


def bullet_card(slide, title, items, x, y, w, h,
                accent=CYAN, title_bg=None):
    if title_bg is None:
        title_bg = DARK_BLUE
    rect(slide, x, y, w, h, DARK_BLUE)
    rect(slide, x, y, w, Inches(0.42), title_bg)
    rect(slide, x, y, Inches(0.05), h, accent)
    tb(slide, title, x + Inches(0.1), y + Inches(0.05),
       w - Inches(0.15), Inches(0.35), size=11, bold=True, color=WHITE)
    bx = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.48),
                                  w - Inches(0.15), h - Inches(0.55))
    tf = bx.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r0 = p.add_run(); r0.text = '  '; r0.font.size = Pt(9)
        r1 = p.add_run(); r1.text = '● '; r1.font.size = Pt(9)
        r1.font.color.rgb = accent; r1.font.bold = True
        r2 = p.add_run(); r2.text = str(item)
        r2.font.size = Pt(10.5); r2.font.color.rgb = LGREY
        if i < len(items) - 1:
            p.space_after = Pt(5)




# SLIDE 1 – TITLE
# ════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
bg(s1, NAVY)
rect(s1, 0, 0, Inches(0.7), H, GOLD)
rect(s1, Inches(0.7), Inches(4.5), W - Inches(0.7), Inches(0.06), GOLD)

tb(s1, "KADEN INVESTMENT",
   Inches(1.1), Inches(1.0), Inches(11), Inches(0.7),
   size=16, bold=True, color=GOLD)
tb(s1, "Web Development &\nAI Automation Audit",
   Inches(1.1), Inches(1.7), Inches(11), Inches(2.0),
   size=44, bold=True, color=WHITE)
tb(s1, "Strategic Review & Digital Transformation Roadmap",
   Inches(1.1), Inches(3.72), Inches(10), Inches(0.5),
   size=16, color=GOLD, italic=True)
tb(s1, "Prepared by: Genos Apollo  |  May 2026",
   Inches(1.1), Inches(5.0), Inches(8), Inches(0.4),
   size=12, color=MGREY)
tb(s1, "CONFIDENTIAL",
   Inches(1.1), Inches(6.6), Inches(4), Inches(0.35),
   size=10, color=MGREY, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 2 – AGENDA
# ════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
bg(s2)
header(s2, "Agenda", "What we will cover today")
footer(s2)

agenda = [
    ("01", "Company Digital Overview",    "Current state of kaden-sa.com and digital footprint"),
    ("02", "Website Audit Findings",      "UX, performance, SEO and accessibility gaps"),
    ("03", "Web Development Roadmap",     "Recommended improvements & modern tech stack"),
    ("04", "AI Automation Opportunity",   "Where AI can drive efficiency at Kaden"),
    ("05", "Benefits of Automation",      "ROI, time savings and competitive advantage"),
    ("06", "Implementation Timeline",     "Phased rollout plan across 12 months"),
    ("07", "Investment & Next Steps",     "Budget guidance and immediate actions"),
]

for i, (num, title, desc) in enumerate(agenda):
    col = i % 2
    row = i // 2
    x = Inches(0.55) + col * Inches(6.4)
    y = Inches(1.55) + row * Inches(1.22)
    card(s2, x, y, Inches(6.1), Inches(1.05))
    rect(s2, x, y, Inches(0.7), Inches(1.05), NAVY)
    tb(s2, num, x, y + Inches(0.25), Inches(0.7), Inches(0.5),
       size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s2, title, x + Inches(0.8), y + Inches(0.1),
       Inches(5.0), Inches(0.4), size=13, bold=True, color=NAVY)
    tb(s2, desc, x + Inches(0.8), y + Inches(0.52),
       Inches(5.0), Inches(0.42), size=10.5, color=MGREY, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 3 – COMPANY DIGITAL OVERVIEW
# ════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
bg(s3)
header(s3, "Company Digital Overview", "Kaden Investment – current digital footprint")
footer(s3)

# Left – company snapshot
card(s3, Inches(0.4), Inches(1.4), Inches(5.8), Inches(5.5))
rect(s3, Inches(0.4), Inches(1.4), Inches(5.8), Inches(0.45), NAVY)
tb(s3, "Company Snapshot", Inches(0.55), Inches(1.44),
   Inches(5.5), Inches(0.38), size=13, bold=True, color=WHITE)

snapshot = [
    ("Industry",      "Real Estate Development & Asset Management"),
    ("Geography",     "Saudi Arabia (Riyadh, Jeddah & key cities)"),
    ("Established",   "2015 – Vision 2030 aligned portfolio"),
    ("Segments",      "Mixed-Use | Logistics | Retail | Waterfront | Commercial"),
    ("Website",       "www.kaden-sa.com  (Squarespace, Arabic-first)"),
    ("Social",        "LinkedIn + Twitter/X – active company pages"),
    ("CMS",           "Squarespace – limited customisation & API access"),
    ("Mobile",        "Responsive design; RTL Arabic + English versions"),
]
ty = Inches(2.02)
for k, v in snapshot:
    tb(s3, k, Inches(0.6), ty, Inches(1.55), Inches(0.3),
       size=10.5, bold=True, color=NAVY)
    tb(s3, v, Inches(2.2), ty, Inches(3.8), Inches(0.3),
       size=10.5, color=MGREY)
    ty += Inches(0.52)

# Right – digital health metrics
card(s3, Inches(6.5), Inches(1.4), Inches(6.5), Inches(5.5))
rect(s3, Inches(6.5), Inches(1.4), Inches(6.5), Inches(0.45), GOLD)
tb(s3, "Key Digital Health Indicators", Inches(6.65), Inches(1.44),
   Inches(6.2), Inches(0.38), size=13, bold=True, color=WHITE)

metrics = [
    ("Mobile Page Speed",        "44 / 100",   "Google PageSpeed – below benchmark",    RED),
    ("Desktop Speed",            "69 / 100",   "PageSpeed Insights estimate",            AMBER),
    ("SEO Score",                "58 / 100",   "Limited schema & meta coverage",         AMBER),
    ("Accessibility (WCAG)",     "Partial",    "RTL + EN gaps; AA compliance needed",    AMBER),
    ("SSL / HTTPS",              "✓ Active",   "Valid certificate in place",             GREEN),
    ("Bilingual Optimisation",   "Basic",      "Arabic primary; EN content incomplete",  AMBER),
    ("AI / Chatbot",             "None",       "No automated enquiry handling",          RED),
    ("Social Integration",       "Limited",    "No live feed or post automation",        RED),
]
ty2 = Inches(2.02)
for label, val, note, col in metrics:
    tb(s3, label, Inches(6.65), ty2, Inches(2.7), Inches(0.3),
       size=10.5, bold=True, color=NAVY)
    tb(s3, val,   Inches(9.4), ty2, Inches(1.25), Inches(0.3),
       size=10.5, bold=True, color=col, align=PP_ALIGN.CENTER)
    tb(s3, note,  Inches(10.7), ty2, Inches(2.1), Inches(0.3),
       size=9, color=MGREY, italic=True)
    ty2 += Inches(0.52)

# ════════════════════════════════════════════════════════════════
# SLIDE 4 – WEBSITE AUDIT FINDINGS
# ════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
bg(s4)
header(s4, "Website Audit Findings", "kaden-sa.com – identified gaps and opportunities")
footer(s4)

audit_cards = [
    ("Performance",
     ["Mobile load >4.5s (target <2.5s)", "Squarespace CDN slower than custom stack",
      "Heavy image carousels not lazy-loaded", "No server-side rendering or caching layer"]),
    ("UX / Design",
     ["English navigation incomplete vs Arabic", "No advanced property search / filters",
      "Investor portal links to SharePoint only", "CTA buttons lack visual hierarchy"]),
    ("SEO",
     ["Missing meta descriptions on 70%+ pages", "No structured data / schema markup",
      "Bilingual hreflang tags absent", "Blog/news content sparse – low domain authority"]),
    ("Security & Compliance",
     ["No PDPL-compliant cookie consent banner", "Contact form validation needs hardening",
      "Missing security headers (CSP, HSTS)", "Employee portal exposed via direct SharePoint URL"]),
]

for i, (title, bullets) in enumerate(audit_cards):
    x = Inches(0.4) + i * Inches(3.2)
    y = Inches(1.45)
    card(s4, x, y, Inches(3.05), Inches(5.3))
    rect(s4, x, y, Inches(3.05), Inches(0.45), NAVY)
    tb(s4, title, x + Inches(0.1), y + Inches(0.06),
       Inches(2.85), Inches(0.36), size=12, bold=True, color=WHITE)
    ty = y + Inches(0.62)
    for b in bullets:
        rect(s4, x + Inches(0.15), ty + Inches(0.07), Inches(0.1), Inches(0.1), GOLD)
        tb(s4, b, x + Inches(0.35), ty, Inches(2.6), Inches(0.55),
           size=10, color=NAVY)
        ty += Inches(0.72)

rect(s4, Inches(0.4), Inches(6.8), Inches(12.5), Inches(0.06), GOLD)
tb(s4, "Priority: Performance + bilingual SEO + PDPL compliance upgrades deliver fastest ROI on kaden-sa.com",
   Inches(0.4), Inches(6.86), Inches(12.5), Inches(0.3),
   size=10.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 5 – WEB DEVELOPMENT ROADMAP
# ════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
bg(s5)
header(s5, "Web Development Roadmap", "Modernising kaden-sa.com for performance and global reach")
footer(s5)

# Left – tech stack
card(s5, Inches(0.4), Inches(1.4), Inches(4.5), Inches(5.6))
rect(s5, Inches(0.4), Inches(1.4), Inches(4.5), Inches(0.45), TEAL)
tb(s5, "Recommended Tech Stack", Inches(0.55), Inches(1.44),
   Inches(4.2), Inches(0.38), size=13, bold=True, color=WHITE)

stack = [
    ("Frontend",      "Next.js 15 (React) – SSR, bilingual (AR/EN)"),
    ("CMS",           "Sanity.io – headless, API-first, RTL support"),
    ("Hosting",       "Vercel + Cloudflare CDN (GCC edge nodes)"),
    ("Search",        "Algolia – Arabic-language property search"),
    ("Maps",          "Google Maps API + project geo-pins"),
    ("Forms",         "React Hook Form + Zod + PDPL consent"),
    ("Analytics",     "GA4 + Microsoft Clarity heatmaps"),
    ("Auth",          "Azure AD B2C – investor & tenant portal"),
]
ty = Inches(2.02)
for k, v in stack:
    tb(s5, k, Inches(0.6), ty, Inches(1.4), Inches(0.3),
       size=10, bold=True, color=TEAL)
    tb(s5, v, Inches(2.05), ty, Inches(2.75), Inches(0.3),
       size=10, color=NAVY)
    ty += Inches(0.5)

# Right – improvements
card(s5, Inches(5.2), Inches(1.4), Inches(7.75), Inches(5.6))
rect(s5, Inches(5.2), Inches(1.4), Inches(7.75), Inches(0.45), NAVY)
tb(s5, "Key Improvements & Features",
   Inches(5.35), Inches(1.44), Inches(7.4), Inches(0.38),
   size=13, bold=True, color=WHITE)

improvements = [
    ("Speed & Core Web Vitals",
     "Target 90+ PageSpeed. GCC-based CDN, WebP imagery, SSR caching, lazy-load carousels."),
    ("Bilingual (AR / EN) Excellence",
     "Full RTL + LTR parity. hreflang tags, bilingual sitemap, Arabic SEO keyword strategy."),
    ("Advanced Project & Property Search",
     "Faceted filters: location, type, size, status. Interactive map-based exploration."),
    ("Investor & Tenant Portal",
     "Secure login via Azure AD: lease docs, unit reports, payment history, maintenance requests."),
    ("PDPL Cookie Consent & Security",
     "Saudi PDPL-compliant consent banner, CSP headers, secure form validation, HSTS."),
    ("Social & Media Automation",
     "New project launches auto-shared to LinkedIn + X. News syndicated from headless CMS."),
]
ty2 = Inches(1.98)
for title, desc in improvements:
    rect(s5, Inches(5.35), ty2 + Inches(0.1), Inches(0.08), Inches(0.25), GOLD)
    tb(s5, title, Inches(5.55), ty2, Inches(7.1), Inches(0.28),
       size=11, bold=True, color=NAVY)
    tb(s5, desc, Inches(5.55), ty2 + Inches(0.3), Inches(7.1), Inches(0.35),
       size=10, color=MGREY)
    ty2 += Inches(0.78)

# ════════════════════════════════════════════════════════════════
# SLIDE 6 – AI AUTOMATION OPPORTUNITY MAP
# ════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
bg(s6)
header(s6, "AI Automation Opportunity Map", "Where AI delivers immediate value for Kaden")
footer(s6)

opps = [
    (NAVY,     "Project Listings & Content",
     ["Auto-generate bilingual (AR/EN) project descriptions", "AI-powered photo tagging & virtual tour support",
      "Dynamic unit availability updates from Excel/ERP", "Bulk listing import & normalisation"]),
    (GOLD,     "Investor & Client Communications",
     ["AI chatbot (24/7 Arabic + English enquiries)", "Auto-draft lease renewal & unit offers",
      "Maintenance ticket classification & routing", "WhatsApp + email follow-up sequences"]),
    (TEAL,     "Market Intelligence & Reporting",
     ["Auto-scrape & summarise Saudi real estate trends", "Weekly AI-generated Vision 2030 alignment reports",
      "Portfolio performance dashboards (live data)", "Predictive vacancy & pricing alerts"]),
    (GREEN,    "Operations & Administration",
     ["Invoice processing & AP automation (SAR)", "Arabic contract review & clause flagging",
      "Meeting transcription → bilingual action items", "HR: CV screening & onboarding in Arabic/EN"]),
]

for i, (col, title, bullets) in enumerate(opps):
    c = i % 2
    r = i // 2
    x = Inches(0.4) + c * Inches(6.5)
    y = Inches(1.45) + r * Inches(2.85)
    card(s6, x, y, Inches(6.2), Inches(2.65))
    rect(s6, x, y, Inches(6.2), Inches(0.45), col)
    tb(s6, title, x + Inches(0.15), y + Inches(0.06),
       Inches(5.9), Inches(0.36), size=13, bold=True, color=WHITE)
    ty = y + Inches(0.62)
    for b in bullets:
        rect(s6, x + Inches(0.18), ty + Inches(0.08), Inches(0.1), Inches(0.1), col)
        tb(s6, b, x + Inches(0.38), ty, Inches(5.65), Inches(0.38),
           size=10.5, color=NAVY)
        ty += Inches(0.48)

# ════════════════════════════════════════════════════════════════
# SLIDE 7 – BENEFITS OF AI AUTOMATION
# ════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
bg(s7)
header(s7, "Benefits of AI Automation", "Measurable impact across Kaden's operations")
footer(s7)

stats = [
    ("60–70%",   "Reduction in\nrepetitive admin tasks"),
    ("3×",       "Faster project\nlisting publication"),
    ("40%",      "Lower cost per\nclient enquiry"),
    ("24/7",     "Arabic + English\nAI client service"),
]
for i, (val, label) in enumerate(stats):
    stat_card(s7, val, label,
              Inches(0.4) + i * Inches(3.2), Inches(1.45),
              w=Inches(3.0), h=Inches(1.6))

benefits = [
    ("Speed to Market",
     ["Listings live in minutes, not hours", "AI writes bilingual descriptions from spec sheet",
      "Photos auto-tagged and resized", "Social post drafted & scheduled automatically"]),
    ("Cost Efficiency",
     ["Reduce staff time on data entry by >50%", "Lower translation and copywriting costs",
      "Decrease marketing agency dependency", "Scale output without scaling headcount"]),
    ("Investor & Client Experience",
     ["Instant 24/7 responses via AI chatbot (AR/EN)", "Personalised project recommendations",
      "Automated viewing confirmations & reminders", "Proactive lease renewal outreach"]),
    ("Competitive Edge",
     ["AI-driven Saudi market intelligence reports", "First-mover advantage in KSA prop-tech",
      "Attract tech-savvy Vision 2030 investors", "Data-driven portfolio decisions"]),
]

for i, (title, bullets) in enumerate(benefits):
    c = i % 2
    r = i // 2
    x = Inches(0.4) + c * Inches(6.5)
    y = Inches(3.2) + r * Inches(1.85)
    bullet_card(s7, title, bullets, x, y, Inches(6.2), Inches(1.65),
                title_bg=NAVY if r == 0 else TEAL)

# ════════════════════════════════════════════════════════════════
# SLIDE 8 – ROI & BUSINESS CASE
# ════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
bg(s8)
header(s8, "ROI & Business Case", "Projected returns from web + AI investment (SAR)")
footer(s8)

# Left – investment
card(s8, Inches(0.4), Inches(1.4), Inches(5.8), Inches(5.6))
rect(s8, Inches(0.4), Inches(1.4), Inches(5.8), Inches(0.45), NAVY)
tb(s8, "Indicative Investment (Year 1 – SAR)",
   Inches(0.55), Inches(1.44), Inches(5.5), Inches(0.38),
   size=13, bold=True, color=WHITE)

invest_rows = [
    ("Website Rebuild (Next.js + Sanity CMS)",     "SAR 80,000 – SAR 130,000"),
    ("Bilingual (AR/EN) SEO & PDPL Sprint",        "SAR 18,000 – SAR 32,000"),
    ("AI Chatbot & Automation Stack",              "SAR 35,000 – SAR 60,000"),
    ("Listing Automation & AI Copywriting",        "SAR 20,000 – SAR 38,000"),
    ("Investor / Tenant Portal (Azure AD)",        "SAR 30,000 – SAR 55,000"),
    ("Training & Change Management",               "SAR 10,000 – SAR 18,000"),
    ("Ongoing monthly support (est.)",             "SAR 8,000 – SAR 14,000 / mo"),
    ("TOTAL YEAR 1 (once-off + 12 mo)",           "SAR 289,000 – SAR 501,000"),
]
ty = Inches(2.0)
for i, (item, cost) in enumerate(invest_rows):
    is_total = i == len(invest_rows) - 1
    if is_total:
        rect(s8, Inches(0.4), ty - Inches(0.05), Inches(5.8), Inches(0.06), GOLD)
        ty += Inches(0.12)
        tb(s8, item, Inches(0.6), ty, Inches(3.5), Inches(0.35),
           size=11, bold=True, color=NAVY)
        tb(s8, cost, Inches(4.05), ty, Inches(1.95), Inches(0.35),
           size=11, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)
    else:
        tb(s8, item, Inches(0.6), ty, Inches(3.5), Inches(0.3),
           size=10, color=NAVY)
        tb(s8, cost, Inches(4.05), ty, Inches(1.95), Inches(0.3),
           size=10, color=MGREY, align=PP_ALIGN.RIGHT)
    ty += Inches(0.5)

# Right – returns
card(s8, Inches(6.5), Inches(1.4), Inches(6.5), Inches(5.6))
rect(s8, Inches(6.5), Inches(1.4), Inches(6.5), Inches(0.45), GREEN)
tb(s8, "Projected Returns & Savings (SAR)",
   Inches(6.65), Inches(1.44), Inches(6.2), Inches(0.38),
   size=13, bold=True, color=WHITE)

returns = [
    ("Admin time saved (FTE equivalent)",           "1.5–2 FTE/yr  ≈  SAR 180,000+"),
    ("Reduced translation & copywriting costs",     "~SAR 60,000 / year"),
    ("Increased qualified investor enquiries (SEO)","+ 25–40% organic enquiries"),
    ("Faster listing → lease / sale cycle",         "Est. −10 days vacancy per unit"),
    ("Reduced call volume via AI chatbot",          "−30% inbound call load"),
    ("Investor portal: reduced manual reporting",   "SAR 40,000+ saved / year"),
    ("ESTIMATED YEAR 1 SAVINGS / UPSIDE",          "SAR 280,000 – SAR 500,000+"),
]
ty2 = Inches(2.0)
for i, (item, val) in enumerate(returns):
    is_total = i == len(returns) - 1
    if is_total:
        rect(s8, Inches(6.5), ty2 - Inches(0.05), Inches(6.5), Inches(0.06), GOLD)
        ty2 += Inches(0.12)
        tb(s8, item, Inches(6.65), ty2, Inches(3.9), Inches(0.35),
           size=11, bold=True, color=NAVY)
        tb(s8, val, Inches(10.6), ty2, Inches(2.2), Inches(0.35),
           size=11, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)
    else:
        tb(s8, item, Inches(6.65), ty2, Inches(3.9), Inches(0.3),
           size=10, color=NAVY)
        tb(s8, val, Inches(10.6), ty2, Inches(2.2), Inches(0.3),
           size=10, color=MGREY, align=PP_ALIGN.RIGHT)
    ty2 += Inches(0.5)

tb(s8, "Breakeven estimated within 8–14 months  |  3-year NPV strongly positive",
   Inches(0.4), Inches(7.07), Inches(12.5), Inches(0.28),
   size=10.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 9 – IMPLEMENTATION TIMELINE
# ════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(BLANK)
bg(s9)
header(s9, "Implementation Timeline", "Phased 12-month delivery plan")
footer(s9)

phases = [
    ("Phase 1\nMon 1–3",  GOLD,     [
        "Website audit & bilingual SEO sprint",
        "PDPL cookie consent go-live",
        "AI chatbot MVP on kaden-sa.com",
        "Project auto-description pilot (AR/EN)",
    ]),
    ("Phase 2\nMon 4–6",  NAVY,     [
        "Full Next.js + Sanity CMS rebuild",
        "Algolia bilingual property search",
        "LinkedIn + X auto-posting setup",
        "Security headers & PDPL hardening",
    ]),
    ("Phase 3\nMon 7–9",  TEAL,     [
        "Investor / Tenant portal v1 (Azure AD)",
        "AI market intelligence reports (KSA)",
        "Invoice & contract automation",
        "Staff training & enablement",
    ]),
    ("Phase 4\nMon 10–12", GREEN,   [
        "Predictive analytics dashboard",
        "Advanced chatbot (Arabic fine-tuning)",
        "ERP / SAP integration layer",
        "Full audit, review & optimise",
    ]),
]

for i, (label, col, items) in enumerate(phases):
    x = Inches(0.4) + i * Inches(3.2)
    y = Inches(1.45)
    card(s9, x, y, Inches(3.0), Inches(5.55))
    rect(s9, x, y, Inches(3.0), Inches(0.7), col)
    tb(s9, label, x, y + Inches(0.06), Inches(3.0), Inches(0.62),
       size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 3:
        rect(s9, x + Inches(3.02), y + Inches(0.3), Inches(0.16), Inches(0.06), MGREY)
    ty = y + Inches(0.88)
    for item in items:
        rect(s9, x + Inches(0.18), ty + Inches(0.08), Inches(0.1), Inches(0.1), col)
        tb(s9, item, x + Inches(0.38), ty, Inches(2.5), Inches(0.42),
           size=10.5, color=NAVY)
        ty += Inches(0.9)

# ════════════════════════════════════════════════════════════════
# SLIDE 10 – WHY AUTOMATE NOW?
# ════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(BLANK)
bg(s10, NAVY)
rect(s10, 0, 0, Inches(0.5), H, GOLD)
footer(s10)

tb(s10, "Why Automate — Now?",
   Inches(0.8), Inches(0.5), Inches(11), Inches(0.65),
   size=30, bold=True, color=GOLD)
tb(s10, "Saudi Arabia's real estate sector is digitalising at speed. The window for first-mover advantage is open — but not for long.",
   Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.55),
   size=14, color=WHITE, italic=True)

reasons = [
    ("Vision 2030 is accelerating digital adoption",
     "The Saudi government is actively pushing prop-tech, smart cities, and digital infrastructure. Kaden can position as a flagship digital-first developer and attract premium investor attention."),
    ("Competitors are moving fast",
     "Major KSA developers are investing in AI-powered portals and automated marketing. Early movers capture SEO authority, brand preference, and cost advantages that are very difficult to recover later."),
    ("Bilingual clients expect instant, personalised responses",
     "Arabic and English-speaking investors expect immediate answers. An AI chatbot means Kaden is first to respond — always — without overtime pay or language barriers."),
    ("AI tools are now Gulf-market ready",
     "Arabic-language LLMs (GPT-4o, Claude 3.5) have matured significantly. Enterprise-grade bilingual automation is now accessible and cost-effective for mid-size developers."),
    ("Data compounds in value",
     "Every automated interaction generates structured data. Over 12–24 months, Kaden builds a proprietary dataset for pricing, forecasting, and investor reporting that competitors cannot replicate."),
]
ty = Inches(1.95)
for i, (title, desc) in enumerate(reasons):
    rect(s10, Inches(0.8), ty, Inches(0.55), Inches(0.55),
         GOLD if i % 2 == 0 else TEAL)
    tb(s10, str(i + 1), Inches(0.8), ty, Inches(0.55), Inches(0.55),
       size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    tb(s10, title, Inches(1.5), ty, Inches(11), Inches(0.3),
       size=12, bold=True, color=WHITE)
    tb(s10, desc, Inches(1.5), ty + Inches(0.3), Inches(11.3), Inches(0.55),
       size=10.5, color=RGBColor(0xBB, 0xCC, 0xDD))
    ty += Inches(0.98)

# ════════════════════════════════════════════════════════════════
# SLIDE 11 – NEXT STEPS
# ════════════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(BLANK)
bg(s11)
header(s11, "Next Steps", "From audit to action — what happens now")
footer(s11)

steps = [
    ("Week 1",  GOLD,     "Stakeholder Alignment",
     "Present findings to Kaden leadership. Agree priority areas and budget envelope. Sign off on Phase 1 scope."),
    ("Week 2",  NAVY,     "Kick-Off & Access",
     "Provide CMS, hosting and analytics access. Assign internal champion. Finalise project charter & NDA."),
    ("Week 3",  TEAL,     "Sprint 1 Begins",
     "PDPL fix, bilingual SEO and AI chatbot integration start simultaneously. Daily stand-ups (AR/EN)."),
    ("Month 2", GREEN,    "Phase 1 Go-Live",
     "Chatbot live on kaden-sa.com. SEO fixes deployed. First AI-generated bilingual project descriptions in production."),
]

for i, (when, col, title, desc) in enumerate(steps):
    x = Inches(0.4) + i * Inches(3.2)
    y = Inches(1.5)
    card(s11, x, y, Inches(3.0), Inches(4.0))
    rect(s11, x, y, Inches(3.0), Inches(0.45), col)
    tb(s11, when, x, y + Inches(0.06), Inches(3.0), Inches(0.36),
       size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s11, title, x + Inches(0.15), y + Inches(0.58),
       Inches(2.7), Inches(0.36), size=12, bold=True, color=col)
    tb(s11, desc, x + Inches(0.15), y + Inches(1.0),
       Inches(2.7), Inches(2.8), size=10.5, color=NAVY)

# CTA box
card(s11, Inches(0.4), Inches(5.8), Inches(12.5), Inches(1.1), fill=NAVY)
tb(s11, "Ready to transform Kaden's digital presence?",
   Inches(0.6), Inches(5.88), Inches(7), Inches(0.4),
   size=14, bold=True, color=GOLD)
tb(s11, "Contact Genos Apollo to schedule your discovery workshop.",
   Inches(0.6), Inches(6.28), Inches(7), Inches(0.38),
   size=11, color=WHITE)
tb(s11, "www.kaden-sa.com  |  linkedin.com/company/kaden-investment",
   Inches(8.2), Inches(6.0), Inches(4.5), Inches(0.38),
   size=10, color=GOLD, italic=True, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════
# SLIDE 12 – THANK YOU / BACK COVER
# ════════════════════════════════════════════════════════════════
s12 = prs.slides.add_slide(BLANK)
bg(s12, NAVY)
rect(s12, 0, 0, Inches(0.7), H, GOLD)
rect(s12, Inches(0.7), Inches(3.5), W - Inches(0.7), Inches(0.06), GOLD)

tb(s12, "Thank You",
   Inches(1.1), Inches(2.0), Inches(11), Inches(1.2),
   size=52, bold=True, color=WHITE)
tb(s12, "Building tomorrow's cities through digital intelligence.",
   Inches(1.1), Inches(3.1), Inches(10), Inches(0.55),
   size=16, color=GOLD, italic=True)
tb(s12, "Kaden Investment  ·  www.kaden-sa.com",
   Inches(1.1), Inches(4.0), Inches(8), Inches(0.4),
   size=12, color=RGBColor(0xAA, 0xBB, 0xCC))
tb(s12, "Prepared by Genos Apollo  |  May 2026  |  Confidential",
   Inches(1.1), Inches(6.7), Inches(10), Inches(0.35),
   size=10, color=MGREY)

# ── Save ────────────────────────────────────────────────────────
OUT = ROOT / "output" / "decks" / "Kaden_WebDev_AI_Automation_Audit_GenosAI.pptx"
prs.save(str(OUT))
print(f"Saved: {OUT}")
