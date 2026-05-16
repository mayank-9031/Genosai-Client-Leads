import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.parent / "utils"))
from ppt_theme import (
    NAVY, DARK_BLUE, ACCENT, GOLD, RED, GREEN, PURPLE, ORANGE, TEAL, CYAN,
    WHITE, LGREY, MGREY, W, H,
    bg, rect, tb, aline, bullets, tag,
    slide_header, slide_footer, score_badge,
    stat_cards, top_bar_card, phase_card, before_after, genos_logo_block,
)

ROOT = Path(__file__).parent.parent.parent
prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── SLIDE 1 — COVER ──────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(BLANK)
bg(s1)
rect(s1, 0, 0, Inches(0.12), H, CYAN)
rect(s1, Inches(8.5), 0, Inches(4.83), Inches(2.2), DARK_BLUE)
tb(s1, "GENOS AI", Inches(9.2), Inches(0.4), Inches(3.5), Inches(0.8),
   size=28, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
tb(s1, "www.genosai.tech", Inches(9.2), Inches(1.05), Inches(3.5), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
tb(s1, "AQUA TECHNIC — HYDERABAD", Inches(0.55), Inches(1.6),
   Inches(9), Inches(0.6), size=13, color=CYAN)
tb(s1, "Website & AI Automation\nAudit Report", Inches(0.55), Inches(2.05),
   Inches(9), Inches(1.8), size=44, bold=True, color=WHITE)
aline(s1, Inches(3.85), color=CYAN, w=Inches(3))
tb(s1, "Prepared exclusively for Anurag Saxena, CEO — India's Largest Pool Builder",
   Inches(0.55), Inches(4.15), Inches(9), Inches(0.5), size=13, color=LGREY)
rect(s1, 0, Inches(6.5), W, Inches(1.0), DARK_BLUE)
tb(s1, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699",
   Inches(0.55), Inches(6.6), Inches(9), Inches(0.5), size=12, color=MGREY)
tb(s1, "May 2026  |  CONFIDENTIAL", Inches(9.5), Inches(6.6), Inches(3.3), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
rect(s1, Inches(10.5), Inches(3.3), Inches(2.3), Inches(2.5), DARK_BLUE)
tb(s1, "AUDIT SCORE", Inches(10.55), Inches(3.4), Inches(2.2), Inches(0.5),
   size=11, bold=True, color=MGREY, align=PP_ALIGN.CENTER)
tb(s1, "5", Inches(10.55), Inches(3.75), Inches(2.2), Inches(1.1),
   size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
tb(s1, "/ 10", Inches(10.55), Inches(4.7), Inches(2.2), Inches(0.4),
   size=18, color=LGREY, align=PP_ALIGN.CENTER)
tag(s1, "HIGH PRIORITY", Inches(10.75), Inches(5.1), fill=RED, size=10)


# ── SLIDE 2 — AT A GLANCE ────────────────────────────────────────────────────
s2 = prs.slides.add_slide(BLANK)
bg(s2)
rect(s2, 0, 0, Inches(0.12), H, CYAN)
tb(s2, "AT A GLANCE", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=CYAN)
tb(s2, "Anurag Saxena and Aqua Technic",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=30, bold=True, color=WHITE)
aline(s2, Inches(1.35), color=CYAN)

stats = [
    ("505", "Projects completed\nsince 2009"),
    ("490", "Happy maintenance\nclients (AMC base)"),
    ("12+ yrs", "In swimming pool\nconstruction"),
    ("5", "Top pool builder\nawards"),
    ("3 mkts", "India, Gulf\n& Australia"),
]
cw = Inches(2.3); ch = Inches(2.2); gap = Inches(0.18)
for i, (num, label) in enumerate(stats):
    l = Inches(0.55) + i * (cw + gap)
    rect(s2, l, Inches(1.8), cw, ch, DARK_BLUE)
    tb(s2, num, l + Inches(0.1), Inches(1.95), cw - Inches(0.2), Inches(1.1),
       size=22, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
    tb(s2, label, l + Inches(0.1), Inches(3.05), cw - Inches(0.2), Inches(0.7),
       size=11, color=LGREY, align=PP_ALIGN.CENTER)

rect(s2, Inches(0.55), Inches(4.2), Inches(12.23), Inches(0.6), DARK_BLUE)
tb(s2, "Hyderabad, Telangana  |  enquiry@aquatechnic.co.in  |  +91-9391355379  |  Certified Technician Based Pool Service Company",
   Inches(0.65), Inches(4.28), Inches(12), Inches(0.45), size=11, color=LGREY, align=PP_ALIGN.CENTER)

tb(s2, "SERVICES OFFERED", Inches(0.55), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=CYAN)
bullets(s2, [
    "Turnkey swimming pool construction (above + below ground)",
    "Pool renovation, electrical, mechanical, waterproofing",
    "Daily maintenance, repair, water features, hydrotherapy",
    "Franchise operations across India"],
        Inches(0.55), Inches(5.35), Inches(5.8), Inches(1.5), size=12, color=LGREY)

tb(s2, "INDUSTRIES SERVED", Inches(6.8), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=CYAN)
bullets(s2, [
    "Hotels & resorts (primary high-value segment)",
    "Sports facilities and wellness centres",
    "Residential (high-end homeowners)",
    "International: Gulf and Australia projects"],
        Inches(6.8), Inches(5.35), Inches(5.8), Inches(1.5), size=12, color=LGREY)


# ── SLIDE 3 — WHAT GENOS AI DOES ─────────────────────────────────────────────
s3 = prs.slides.add_slide(BLANK)
bg(s3)
rect(s3, 0, 0, Inches(0.12), H, ACCENT)
tb(s3, "ABOUT GENOS AI", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s3, "AI Automation for Service & Construction Businesses.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=30, bold=True, color=WHITE)
aline(s3, Inches(1.35), color=ACCENT)

tb(s3, "What We Build", Inches(0.55), Inches(1.6), Inches(12.5), Inches(0.4),
   size=14, bold=True, color=ACCENT)

services = [
    ("WhatsApp AI Agents", TEAL,
     "24/7 AI agents on WhatsApp that answer inquiries, qualify leads by project type and budget, and book site visits automatically. The most effective channel for the Indian market."),
    ("Lead Qualification\n& Nurture", GOLD,
     "Automated email and WhatsApp sequences that re-engage prospects who inquired but didn't convert. Most high-value construction contracts close after 3-5 follow-ups, not the first contact."),
    ("AMC & Maintenance\nAutomation", GREEN,
     "Automated renewal reminders, service scheduling, and client communication for recurring maintenance clients. Reduces churn, removes manual follow-up from the team's plate."),
    ("Website AI Chatbot", PURPLE,
     "AI chatbot on your website that answers cost questions, qualifies project scope, and routes hotel clients vs residential vs international - with full context before any human is involved."),
]
cw2 = Inches(2.95); gap2 = Inches(0.16)
for i, (title, color, desc) in enumerate(services):
    l2 = Inches(0.55) + i * (cw2 + gap2)
    rect(s3, l2, Inches(2.1), cw2, Inches(0.06), color)
    rect(s3, l2, Inches(2.16), cw2, Inches(3.5), DARK_BLUE)
    tb(s3, title, l2 + Inches(0.12), Inches(2.28), cw2 - Inches(0.24), Inches(0.6),
       size=13, bold=True, color=color)
    tb(s3, desc, l2 + Inches(0.12), Inches(2.9), cw2 - Inches(0.24), Inches(2.6),
       size=11, color=LGREY)

rect(s3, Inches(0.55), Inches(5.85), Inches(12.23), Inches(1.3), DARK_BLUE)
tb(s3, "WHY GENOS AI FOR AQUA TECHNIC", Inches(0.7), Inches(5.92), Inches(5.5), Inches(0.35),
   size=11, bold=True, color=ACCENT)
bullets(s3, [
    "Built for Indian market: WhatsApp-first, multilingual, India-timezone aware",
    "Construction industry experience: understands project qualification, site visit workflows, AMC cycles",
    "Rapid deployment: WhatsApp agent live in under 2 weeks",
    "ROI-focused: every automation tied to leads captured, AMC renewals saved, or response time improved"],
        Inches(0.7), Inches(6.28), Inches(5.8), Inches(0.8), size=11, color=LGREY, dot=ACCENT)

tb(s3, "CLIENTS WE SERVE", Inches(7.0), Inches(5.92), Inches(5.5), Inches(0.35),
   size=11, bold=True, color=ACCENT)
bullets(s3, [
    "Construction & engineering firms managing large project pipelines",
    "Service businesses with recurring maintenance/AMC client bases",
    "Real estate developers and hospitality businesses",
    "Any business losing leads to slow WhatsApp or email response"],
        Inches(7.0), Inches(6.28), Inches(5.8), Inches(0.8), size=11, color=LGREY, dot=ACCENT)


# ── SLIDE 4 — THE THREE GAPS ─────────────────────────────────────────────────
s4 = prs.slides.add_slide(BLANK)
bg(s4)
rect(s4, 0, 0, Inches(0.12), H, RED)
tb(s4, "THE THREE GAPS", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=RED)
tb(s4, "505 Projects Built. Weekend Inquiries Still Hitting a Static Form.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=26, bold=True, color=WHITE)
aline(s4, Inches(1.35), color=RED)

gap_cards = [
    ("GAP 1\nNO WHATSAPP\nRESPONSE", RED,
     "The #1 channel in India — completely absent",
     [
         "No WhatsApp chat widget on the site — hotel developers browse on Sunday and WhatsApp 3 builders simultaneously",
         "The one who responds first in under 5 minutes wins the site visit in 80%+ of cases",
         "Aqua Technic's response: fill a form, wait for a call on Monday",
         "Competitors with WhatsApp AI agents respond in under 60 seconds, 24/7",
     ]),
    ("GAP 2\n490 AMC CLIENTS\nMANUAL ONLY", GOLD,
     "Largest recurring revenue base managed without automation",
     [
         "490 maintenance clients = significant AMC renewal revenue each year",
         "No automated renewal reminders via WhatsApp or email at 30/7/1 days before expiry",
         "No online service scheduling — clients call or email to book maintenance visits",
         "Each lapsed AMC is direct revenue loss that automation prevents",
     ]),
    ("GAP 3\nNO COST ESTIMATOR\n+ CREDIBILITY BUGS", ACCENT,
     "First question unanswered. Site has visible test entries.",
     [
         "Most common first question: 'How much does a swimming pool cost?' — gets no immediate answer",
         "A cost estimator tool (pool size, type, features) qualifies serious buyers and captures contact",
         "Navigation has placeholder/test entries visible to all site visitors — active credibility damage",
         "No post-project review system despite 505 completed projects — social proof wasted",
     ]),
]
cw3 = Inches(3.95); gap3i = Inches(0.18)
for i, (title, color, subtitle, pts) in enumerate(gap_cards):
    l = Inches(0.55) + i * (cw3 + gap3i)
    rect(s4, l, Inches(1.8), cw3, Inches(0.5), color)
    tb(s4, title, l + Inches(0.12), Inches(1.85), cw3 - Inches(0.24), Inches(0.42),
       size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    rect(s4, l, Inches(2.3), cw3, Inches(4.6), DARK_BLUE)
    tb(s4, subtitle, l + Inches(0.12), Inches(2.38), cw3 - Inches(0.24), Inches(0.35),
       size=10, italic=True, color=color)
    bullets(s4, pts, l + Inches(0.12), Inches(2.78), cw3 - Inches(0.24), Inches(3.9),
            size=11, color=LGREY, dot=color)

rect(s4, Inches(0.55), Inches(7.05), Inches(12.23), Inches(0.3), DARK_BLUE)
tb(s4, "Aqua Technic built India's largest pool portfolio manually. The next 500 projects need automation to compete at the same speed as tech-enabled rivals.",
   Inches(0.7), Inches(7.08), Inches(12), Inches(0.25), size=11, italic=True, color=GOLD)


# ── SLIDE 5 — WEBSITE AUDIT ──────────────────────────────────────────────────
s5 = prs.slides.add_slide(BLANK)
bg(s5)
rect(s5, 0, 0, Inches(0.12), H, PURPLE)
tb(s5, "WEBSITE DESIGN AUDIT", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=PURPLE)
tb(s5, "Established Track Record. Outdated Digital Infrastructure.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.8), size=28, bold=True, color=WHITE)
aline(s5, Inches(1.45), color=PURPLE, w=Inches(1.5))

rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(3.6), DARK_BLUE)
rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(0.42), RED)
tb(s5, "AQUATECHNIC.CO.IN TODAY", Inches(0.65), Inches(1.92), Inches(3.6), Inches(0.38),
   size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "No WhatsApp chat widget",
    "Contact form only — no automated follow-up",
    "Navigation has test/placeholder entries",
    "No pool cost estimator or quote tool",
    "No automated AMC renewal reminders",
    "No post-project review collection",
    "No filterable project gallery",
    "No separate hotel vs residential intake",
], Inches(0.7), Inches(2.42), Inches(3.6), Inches(2.9), size=11, color=LGREY, dot=RED)

tb(s5, "->", Inches(4.45), Inches(3.6), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(3.6), DARK_BLUE)
rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(0.42), GOLD)
tb(s5, "WHAT TOP INDIAN BUILDERS DO", Inches(5.1), Inches(1.92), Inches(3.4), Inches(0.38),
   size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "WhatsApp AI agent responds in under 60 seconds",
    "Cost estimator captures serious buyers immediately",
    "Automated follow-up within 5 minutes of inquiry",
    "AMC renewals managed via automated WhatsApp",
    "Post-project review request within 48 hrs of completion",
    "Project gallery filterable by type and scale",
    "Hotel-specific landing page with hospitality case studies",
    "Client portal for maintenance history and scheduling",
], Inches(5.1), Inches(2.42), Inches(3.4), Inches(2.9), size=11, color=LGREY, dot=GOLD)

tb(s5, "->", Inches(8.68), Inches(3.6), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(3.6), DARK_BLUE)
rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(0.42), GREEN)
tb(s5, "WHAT GENOS AI BUILDS", Inches(9.35), Inches(1.92), Inches(3.4), Inches(0.38),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "WhatsApp AI agent: qualify + book site visits 24/7",
    "Pool cost estimator with WhatsApp/email gate",
    "5-min automated follow-up to every form inquiry",
    "AMC renewal: WhatsApp at 30/7/1 days + pay link",
    "Post-project review automation via WhatsApp",
    "Filterable gallery: hotel, residential, water features",
    "Hotel + residential separate landing pages",
    "Fix navigation placeholders (week 1, immediate)",
], Inches(9.35), Inches(2.42), Inches(3.4), Inches(2.9), size=11, color=LGREY, dot=GREEN)

rect(s5, Inches(0.55), Inches(5.72), Inches(12.23), Inches(1.28), DARK_BLUE)
tb(s5, "INDIA MARKET BENCHMARKS:", Inches(0.7), Inches(5.78), Inches(2.8), Inches(0.38),
   size=11, bold=True, color=PURPLE)
bstats = [
    ("< 60 sec", "WhatsApp response time\nto win the site visit"),
    ("490", "AMC clients with\nzero renewal automation"),
    ("505", "Projects with no\nauto review collection"),
    ("80%", "Indian B2B buyers prefer\nWhatsApp over email"),
]
sw = Inches(2.55); sl = Inches(3.65)
for i, (num, label) in enumerate(bstats):
    xl = sl + i * sw
    tb(s5, num, xl, Inches(5.75), Inches(2.4), Inches(0.55),
       size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s5, label, xl, Inches(6.22), Inches(2.4), Inches(0.7),
       size=10, color=LGREY, align=PP_ALIGN.CENTER)


# ── SLIDE 6 — AI AUTOMATION STACK ────────────────────────────────────────────
s6 = prs.slides.add_slide(BLANK)
bg(s6)
rect(s6, 0, 0, Inches(0.12), H, ACCENT)
tb(s6, "AI AUTOMATION STACK", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s6, "Four Layers Built for India's Largest Pool Builder.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=28, bold=True, color=WHITE)
aline(s6, Inches(1.35), color=ACCENT)

opps = [
    ("WhatsApp AI\nAgent", TEAL,
     "Responds to every inquiry on WhatsApp within 60 seconds, 24/7. Asks 5 qualifying questions: pool type, size, hotel or residential, location (city), timeline. Serious buyers are connected to Anurag's team with full project context. Time-wasters are filtered before any human is involved. Hotel developers browsing on Sunday get a response before your competitors finish their morning chai."),
    ("Pool Cost\nEstimator Tool", GOLD,
     "Most prospects want a ballpark before they call. An interactive estimator on the homepage asks pool dimensions, type (lap pool, infinity, family), finish, and add-ons. Before showing the estimate, it captures WhatsApp number or email. The prospect gets their number. Aqua Technic gets a qualified lead with project scope already defined. High-intent traffic converts instead of bouncing."),
    ("AMC Renewal\nAutomation", GREEN,
     "490 maintenance clients. Each AMC renewal is revenue. Automated WhatsApp messages at 30 days, 7 days, and 1 day before expiry — with a payment link — recover renewals that would otherwise require a manual follow-up call. Service scheduling via WhatsApp removes phone tag from the maintenance team. Aqua Technic's recurring revenue base becomes a managed asset, not a manual chore."),
    ("Post-Project\nReview & Referral", PURPLE,
     "505 completed projects with no systematic review collection. A WhatsApp message sent 48 hours after project handover asks for a Google review and a referral. Hotel and resort clients who are happy at handover are the highest-converting referral source in the construction industry. Catching them at the peak satisfaction moment — before they move to the next project — is the difference between 5 reviews and 500."),
]
cw2 = Inches(2.95); gap2 = Inches(0.16)
for i, (title, color, desc) in enumerate(opps):
    l2 = Inches(0.55) + i * (cw2 + gap2)
    rect(s6, l2, Inches(1.8), cw2, Inches(0.06), color)
    rect(s6, l2, Inches(1.86), cw2, Inches(4.7), DARK_BLUE)
    tb(s6, title, l2 + Inches(0.12), Inches(1.98), cw2 - Inches(0.24), Inches(0.65),
       size=13, bold=True, color=color)
    tb(s6, desc, l2 + Inches(0.12), Inches(2.65), cw2 - Inches(0.24), Inches(3.65),
       size=11, color=LGREY)

rect(s6, Inches(0.55), Inches(6.7), Inches(12.23), Inches(0.45), DARK_BLUE)
tb(s6, "Anurag's team builds the pools. Genos AI handles every touchpoint before the first site visit and after every handover.",
   Inches(0.7), Inches(6.75), Inches(12), Inches(0.35), size=11, italic=True, color=GOLD)


# ── SLIDE 7 — WHATSAPP FIRST: THE INDIA ARGUMENT ─────────────────────────────
s7 = prs.slides.add_slide(BLANK)
bg(s7)
rect(s7, 0, 0, Inches(0.12), H, ORANGE)
tb(s7, "WHY WHATSAPP FIRST", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ORANGE)
tb(s7, "India's Hotel Developers Don't Fill Forms. They WhatsApp.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=28, bold=True, color=WHITE)
aline(s7, Inches(1.35), color=ORANGE, w=Inches(2.5))

rect(s7, Inches(0.55), Inches(1.65), Inches(12.23), Inches(1.55), DARK_BLUE)
india_stats = [
    ("500M+", "WhatsApp users\nin India (2026)"),
    ("80%", "B2B buyers prefer\nWhatsApp over email"),
    ("< 5 min", "Response window to\nwin a site visit"),
    ("24/7", "When hotel developers\nactually browse"),
    ("3 rivals", "Messaged simultaneously\nwith Aqua Technic"),
]
ssw = Inches(2.35)
for i, (num, label) in enumerate(india_stats):
    xl = Inches(0.7) + i * ssw
    tb(s7, num, xl, Inches(1.73), Inches(2.2), Inches(0.65),
       size=20, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    tb(s7, label, xl, Inches(2.33), Inches(2.2), Inches(0.5),
       size=10, color=LGREY, align=PP_ALIGN.CENTER)

rect(s7, Inches(0.55), Inches(3.4), Inches(5.9), Inches(3.25), DARK_BLUE)
tb(s7, "HOW A HOTEL DEVELOPER INQUIRES TODAY", Inches(0.7), Inches(3.47), Inches(5.6), Inches(0.35),
   size=11, bold=True, color=RED)
bullets(s7, [
    "Browses aquatechnic.co.in on a Sunday afternoon",
    "Wants a quote for a rooftop infinity pool at a new hotel",
    "Fills a 4-field contact form and submits",
    "Simultaneously WhatsApps 3 competing builders",
    "Gets a reply from Competitor A within 4 minutes",
    "Competitor A books a site visit for Tuesday morning",
    "Aqua Technic's reply arrives Monday — site visit already booked with rival",
], Inches(0.7), Inches(3.85), Inches(5.6), Inches(2.55), size=11, color=LGREY, dot=RED)

rect(s7, Inches(6.65), Inches(3.4), Inches(6.13), Inches(3.25), DARK_BLUE)
tb(s7, "WITH GENOS AI WHATSAPP AGENT", Inches(6.8), Inches(3.47), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=GREEN)
bullets(s7, [
    "Developer clicks WhatsApp widget on site at 3pm Sunday",
    "AI agent responds within 60 seconds: 'Tell me about your project'",
    "5 qualifying questions: type, size, location, timeline, budget range",
    "Agent sends a ballpark range and books a call for Monday 10am",
    "Anurag's team receives: project brief, contact, preferred time",
    "Monday morning: Aqua Technic calls with full context ready",
    "Site visit confirmed before the developer has heard from Competitor A",
], Inches(6.8), Inches(3.85), Inches(5.8), Inches(2.55), size=11, color=LGREY, dot=GREEN)

rect(s7, Inches(0.55), Inches(6.8), Inches(12.23), Inches(0.38), DARK_BLUE)
tb(s7, "The first credible response wins the site visit. The site visit wins the contract. Speed is the competitive advantage.",
   Inches(0.7), Inches(6.85), Inches(12), Inches(0.28), size=11, italic=True, color=ORANGE)


# ── SLIDE 8 — IMPLEMENTATION PLAN ────────────────────────────────────────────
s8 = prs.slides.add_slide(BLANK)
bg(s8)
rect(s8, 0, 0, Inches(0.12), H, GREEN)
tb(s8, "IMPLEMENTATION PLAN", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s8, "WhatsApp First. AMC Second. Reviews and SEO Third.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=28, bold=True, color=WHITE)
aline(s8, Inches(1.35), color=GREEN)

phases = [
    ("PHASE 1", "Weeks 1-2", GOLD, [
        ("Fix Navigation Placeholders", "Remove all test/placeholder entries from navigation immediately. Credibility restored before anything else goes live."),
        ("WhatsApp AI Agent Live", "Agent trained on pool types, cost ranges, project qualification. Widget on homepage and all service pages. Responds within 60 seconds, 24/7."),
        ("AMC Renewal Automation", "WhatsApp reminders at 30/7/1 days before expiry with payment link deployed for all 490 active maintenance clients."),
    ]),
    ("PHASE 2", "Weeks 3-5", ACCENT, [
        ("Pool Cost Estimator", "Interactive cost estimator with WhatsApp/email gate before results. Hotel and residential paths with separate qualifying questions."),
        ("Post-Project Review Automation", "WhatsApp message at 48 hrs post-handover requesting Google review and referral. Deployed for all future project completions."),
        ("Hotel vs Residential Landing Pages", "Separate pages for hospitality clients (hotel case studies, ROI framing) and residential (design, family, lifestyle framing)."),
    ]),
    ("PHASE 3", "Months 2-3", GREEN, [
        ("SEO Content Engine", "10-15 articles: swimming pool construction Hyderabad, hotel pool builders India, pool maintenance Hyderabad, rooftop infinity pool cost India."),
        ("Filterable Project Gallery", "Gallery with 50+ completed projects filterable by type: hotel, residential, sports, water features. Hotel developers evaluate by portfolio."),
        ("Gulf & Australia Routing", "International inquiry handling with timezone-aware follow-up and separate qualification for Gulf and Australia project inquiries."),
    ]),
]
pw = Inches(3.9); gap4 = Inches(0.18)
for i, (phase, timing, color, items) in enumerate(phases):
    l3 = Inches(0.55) + i * (pw + gap4)
    rect(s8, l3, Inches(1.8), pw, Inches(0.52), color)
    tb(s8, phase, l3 + Inches(0.12), Inches(1.84), pw * 0.5 - Inches(0.12), Inches(0.42),
       size=13, bold=True, color=NAVY)
    tb(s8, timing, l3 + pw * 0.45, Inches(1.88), pw * 0.55 - Inches(0.15), Inches(0.38),
       size=11, color=NAVY, align=PP_ALIGN.RIGHT)
    rect(s8, l3, Inches(2.32), pw, Inches(4.6), DARK_BLUE)
    y_item = Inches(2.42)
    for title, desc in items:
        rect(s8, l3 + Inches(0.12), y_item + Inches(0.06), Inches(0.05), Inches(0.52), color)
        tb(s8, title, l3 + Inches(0.25), y_item, pw - Inches(0.38), Inches(0.32),
           size=12, bold=True, color=color)
        tb(s8, desc, l3 + Inches(0.25), y_item + Inches(0.3), pw - Inches(0.38), Inches(0.45),
           size=11, color=LGREY)
        y_item += Inches(0.95)


# ── SLIDE 9 — BEFORE vs AFTER ────────────────────────────────────────────────
s9 = prs.slides.add_slide(BLANK)
bg(s9)
rect(s9, 0, 0, Inches(0.12), H, ACCENT)
tb(s9, "BEFORE vs AFTER", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s9, "What Actually Changes for Aqua Technic in 90 Days",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=28, bold=True, color=WHITE)
aline(s9, Inches(1.35))

comparisons = [
    ("Hotel Developer Inquiry",
     "Developer WhatsApps 3 pool builders on Sunday. Aqua Technic's form reply arrives Monday. Competitor already has the site visit booked.",
     "WhatsApp AI agent responds within 60 seconds. 5 qualifying questions. Site visit booked for Monday morning before Anurag's team starts their day."),
    ("AMC Renewal (490 Clients)",
     "Maintenance contracts expire. Clients are called manually or not at all. Each lapsed AMC is revenue lost and a relationship that weakens.",
     "WhatsApp reminders at 30/7/1 days. Payment link in message. Renewals confirmed without a single manual call. 490 clients managed automatically."),
    ("Cost Estimation",
     "Prospect asks 'how much for a pool?' Site has no answer. They leave, search Google, and find a competitor with an instant estimate tool.",
     "Interactive estimator on homepage. Pool type, size, finish, add-ons. Email/WhatsApp captured before results shown. Qualified lead with full project scope."),
    ("Post-Project Reviews",
     "505 completed projects. Review collection depends on whether the project manager remembers to ask. Aqua Technic has a fraction of the reviews its track record deserves.",
     "WhatsApp message 48 hrs after handover. Google review requested at peak client satisfaction. Referral asked at 30 days. Review volume compounds monthly."),
    ("Navigation Credibility",
     "Placeholder and test entries visible in site navigation. Every hotel developer who visits sees them and doubts the professionalism of the operation behind the site.",
     "Navigation cleaned in Week 1. Professional, complete, no stray entries. The site matches the 12-year track record it represents."),
]

lw = Inches(5.4)
rect(s9, Inches(0.55), Inches(1.65), lw, Inches(0.35), RED)
tb(s9, "TODAY", Inches(0.65), Inches(1.68), lw - Inches(0.2), Inches(0.28),
   size=11, bold=True, color=WHITE)
rect(s9, Inches(6.45), Inches(1.65), lw + Inches(0.43), Inches(0.35), GREEN)
tb(s9, "AFTER 90 DAYS", Inches(6.55), Inches(1.68), lw, Inches(0.28),
   size=11, bold=True, color=WHITE)

for i, (topic, before, after) in enumerate(comparisons):
    t = Inches(2.1) + i * Inches(1.02)
    rect(s9, Inches(0.55), t, Inches(12.23), Inches(0.22), DARK_BLUE)
    tb(s9, topic.upper(), Inches(0.65), t + Inches(0.02), Inches(12), Inches(0.2),
       size=9, bold=True, color=MGREY)
    tb(s9, before, Inches(0.65), t + Inches(0.24), lw - Inches(0.2), Inches(0.66),
       size=11, color=LGREY)
    tb(s9, after, Inches(6.55), t + Inches(0.24), lw, Inches(0.66),
       size=11, color=WHITE)


# ── SLIDE 10 — NEXT STEPS ────────────────────────────────────────────────────
s10 = prs.slides.add_slide(BLANK)
bg(s10)
rect(s10, 0, 0, Inches(0.12), H, CYAN)
rect(s10, Inches(7.5), 0, Inches(5.83), H, DARK_BLUE)
tb(s10, "NEXT STEPS", Inches(0.55), Inches(0.3), Inches(6.5), Inches(0.5),
   size=11, bold=True, color=CYAN)
tb(s10, "Let's Talk for\n15 Minutes.",
   Inches(0.55), Inches(0.8), Inches(6.8), Inches(2.0), size=48, bold=True, color=WHITE)
aline(s10, Inches(2.8), color=CYAN, w=Inches(2.5))
tb(s10, "No pitch needed. Call or WhatsApp.\nWe'll start with the WhatsApp agent — live in 2 weeks.",
   Inches(0.55), Inches(3.05), Inches(6.5), Inches(0.8), size=14, color=LGREY)

steps = [
    "Call or WhatsApp Rohan: +91 638-714-2699",
    "Week 1-2: WhatsApp AI agent + AMC automation live",
    "Week 3-5: Cost estimator + review automation + landing pages",
    "Month 2-3: SEO content + gallery + Gulf/Australia routing",
]
bullets(s10, steps, Inches(0.55), Inches(3.95), Inches(6.5), Inches(2.0), size=14, color=WHITE, dot=CYAN)

rect(s10, 0, Inches(6.5), Inches(7.5), Inches(1.0), DARK_BLUE)
tb(s10, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699  |  www.genosai.tech",
   Inches(0.55), Inches(6.6), Inches(6.8), Inches(0.5), size=11, color=MGREY)

tb(s10, "THREE GAPS.\nTHREE FIXES.", Inches(7.85), Inches(0.6), Inches(5.0), Inches(1.1),
   size=22, bold=True, color=CYAN)
aline(s10, Inches(1.65), color=CYAN)

summary_items = [
    ("No WhatsApp", "AI agent: 60-sec response, 24/7, qualifies + books"),
    ("490 AMC Manual", "Automated renewal reminders + payment link"),
    ("No Cost Tool", "Pool estimator: qualify + capture before results"),
    ("505 No Reviews", "Post-handover WhatsApp: Google review + referral"),
    ("Nav Bugs", "Placeholders fixed in Week 1 — immediate credibility"),
    ("No SEO", "Hotel + residential + maintenance keyword content"),
]
for i, (label, val) in enumerate(summary_items):
    t2 = Inches(1.85) + i * Inches(0.88)
    rect(s10, Inches(7.85), t2, Inches(1.5), Inches(0.72), NAVY)
    tb(s10, label, Inches(7.9), t2 + Inches(0.05), Inches(1.4), Inches(0.32),
       size=10, bold=True, color=CYAN)
    tb(s10, val, Inches(9.5), t2 + Inches(0.05), Inches(3.6), Inches(0.65),
       size=11, color=LGREY)


prs.save(ROOT / "output" / "decks" / "AquaTechnic_Audit_GenosAI.pptx")
print("Saved: AquaTechnic_Audit_GenosAI.pptx")
