from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── BT Industries Brand Palette  (charcoal · industrial orange · steel) ──
CHARCOAL    = RGBColor(0x1A, 0x1E, 0x26)   # near-black charcoal
DARK_GREY   = RGBColor(0x2C, 0x31, 0x3A)   # dark panel grey
ORANGE      = RGBColor(0xE8, 0x65, 0x0A)   # BT Industries orange
AMBER       = RGBColor(0xF5, 0xA6, 0x23)   # warm amber highlight
STEEL       = RGBColor(0x2E, 0x6E, 0xA6)   # steel blue
TEAL        = RGBColor(0x0D, 0x8A, 0x72)   # teal accent
IVORY       = RGBColor(0xF6, 0xF7, 0xF4)   # warm light bg
SILVER      = RGBColor(0x8A, 0x90, 0x9E)   # silver text
LIGHT_GREY  = RGBColor(0xE8, 0xEA, 0xEE)   # card border
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
SUCCESS     = RGBColor(0x27, 0xAE, 0x60)
CAUTION     = RGBColor(0xF0, 0x9A, 0x1A)
DANGER      = RGBColor(0xC0, 0x38, 0x28)
MID_TEXT    = RGBColor(0x3A, 0x3E, 0x48)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

# ── Helpers ─────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.line.fill.background(); s.line.width = 0
    s.fill.solid(); s.fill.fore_color.rgb = fill
    return s

def add_textbox(slide, text, x, y, w, h,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(font_size); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return txb

def bg(slide, color=IVORY):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, color)

def header(slide, title, sub=None):
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.25), CHARCOAL)
    add_rect(slide, 0, Inches(1.25), Inches(0.38), Inches(6.25), ORANGE)
    add_textbox(slide, title, Inches(0.58), Inches(0.16),
                Inches(10.5), Inches(0.64), 27, True, WHITE)
    if sub:
        add_textbox(slide, sub, Inches(0.58), Inches(0.74),
                    Inches(10.5), Inches(0.42), 12, False, AMBER, italic=True)

def footer(slide,
           txt="BT Industries  |  Web Dev & AI Automation Audit  |  2026"):
    add_rect(slide, 0, Inches(7.2), SLIDE_W, Inches(0.3), CHARCOAL)
    add_textbox(slide, txt, Inches(0.3), Inches(7.2),
                Inches(12.7), Inches(0.3), 9, False, SILVER, PP_ALIGN.CENTER)

def card(slide, x, y, w, h, fill=WHITE, border=True):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = LIGHT_GREY; s.line.width = Pt(0.6)
    else:
        s.line.fill.background(); s.line.width = 0
    return s

def accent_bar(slide, x, y, w, h=Inches(0.048)):
    add_rect(slide, x, y, w, h, ORANGE)

def bullet_card(slide, title, bullets, x, y, w, h,
                hdr_bg=CHARCOAL, hdr_col=ORANGE):
    card(slide, x, y, w, h)
    add_rect(slide, x, y, w, Inches(0.44), hdr_bg)
    add_textbox(slide, title, x + Inches(0.14), y + Inches(0.06),
                w - Inches(0.28), Inches(0.36), 12, True, hdr_col)
    ty = y + Inches(0.58)
    for b in bullets:
        add_textbox(slide, f"▸  {b}", x + Inches(0.15), ty,
                    w - Inches(0.3), Inches(0.36), 10.5, False, MID_TEXT)
        ty += Inches(0.34)

def stat_card(slide, val, lbl, x, y, w=Inches(2.8), h=Inches(1.5)):
    card(slide, x, y, w, h, fill=CHARCOAL, border=False)
    add_textbox(slide, val, x, y + Inches(0.16), w, Inches(0.72),
                30, True, ORANGE, PP_ALIGN.CENTER)
    add_textbox(slide, lbl, x, y + Inches(0.88), w, Inches(0.52),
                11, False, SILVER, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, CHARCOAL)
add_rect(s, 0, 0, Inches(0.72), SLIDE_H, ORANGE)
add_rect(s, Inches(0.72), Inches(4.65), SLIDE_W - Inches(0.72), Inches(0.055), ORANGE)

add_textbox(s, "BT INDUSTRIES  ·  INTEGRATED SOLUTIONS PARTNER",
            Inches(1.1), Inches(0.82), Inches(11.5), Inches(0.55),
            13, True, AMBER)
add_textbox(s, "Web Development &\nAI Automation Audit",
            Inches(1.1), Inches(1.52), Inches(11.2), Inches(2.15),
            44, True, WHITE)
add_textbox(s, "A Strategic Digital Transformation Blueprint for 65 Years of Supply Chain Excellence",
            Inches(1.1), Inches(3.78), Inches(11.2), Inches(0.55),
            14, False, AMBER, italic=True)
add_textbox(s, "Prepared by: Genos Apollo  |  May 2026",
            Inches(1.1), Inches(5.05), Inches(8), Inches(0.4),
            12, False, SILVER)
add_textbox(s, "btindustries.co.za  •  South Africa  •  ISO 9001  •  Level 2 B-BBEE",
            Inches(1.1), Inches(5.55), Inches(10), Inches(0.35),
            11, False, RGBColor(0x80, 0x86, 0x8E), italic=True)
add_textbox(s, "CONFIDENTIAL",
            Inches(1.1), Inches(6.68), Inches(4), Inches(0.32),
            10, False, SILVER, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Agenda", "What we will cover today"); footer(s)

items = [
    ("01", "Company Digital Overview",    "Current state of btindustries.co.za & digital presence"),
    ("02", "Website Audit Findings",      "UX, performance, SEO, B2B portal & content gaps"),
    ("03", "Web Development Roadmap",     "Modernising the site for B2B customer experience"),
    ("04", "AI Automation Opportunities", "Where AI transforms supply chain, ops & client service"),
    ("05", "Benefits of Automation",      "Efficiency, cost savings & competitive advantage"),
    ("06", "Implementation Timeline",     "Phased 12-month delivery plan"),
    ("07", "Investment & Next Steps",     "Budget guidance and immediate actions"),
]
for i, (num, title, desc) in enumerate(items):
    col, row = i % 2, i // 2
    x = Inches(0.55) + col * Inches(6.4)
    y = Inches(1.55) + row * Inches(1.22)
    card(s, x, y, Inches(6.1), Inches(1.05))
    add_rect(s, x, y, Inches(0.7), Inches(1.05), CHARCOAL)
    accent_bar(s, x, y + Inches(1.05) - Inches(0.048), Inches(6.1))
    add_textbox(s, num, x, y + Inches(0.25), Inches(0.7), Inches(0.5),
                20, True, ORANGE, PP_ALIGN.CENTER)
    add_textbox(s, title, x + Inches(0.82), y + Inches(0.1),
                Inches(5.1), Inches(0.4), 13, True, CHARCOAL)
    add_textbox(s, desc,  x + Inches(0.82), y + Inches(0.56),
                Inches(5.1), Inches(0.38), 10, False, SILVER, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 3 — COMPANY DIGITAL OVERVIEW
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Company Digital Overview",
       "BT Industries – current digital footprint & B2B platform health")
footer(s)

card(s, Inches(0.4), Inches(1.4), Inches(5.8), Inches(5.58))
add_rect(s, Inches(0.4), Inches(1.4), Inches(5.8), Inches(0.45), CHARCOAL)
add_textbox(s, "Company Snapshot", Inches(0.56), Inches(1.46),
            Inches(5.5), Inches(0.35), 13, True, ORANGE)
snap = [
    ("Heritage",      "65+ years in South African distribution & supply chain"),
    ("Positioning",   "Integrated Solutions Partner — 'Confidence of Supply'"),
    ("Sectors",       "Automotive/OEM  ·  Industrial  ·  Aftermarket  ·  Mining"),
    ("New Sector",    "Renewable Energy / Solar (DoGo Power partnership, 2026)"),
    ("Speciality",    "Fasteners, components, engineered parts, supply chain mgmt"),
    ("Certifications","ISO 9001  ·  Level 2 B-BBEE Contributor"),
    ("Website",       "btindustries.co.za  (primary digital channel)"),
    ("LinkedIn",      "BT Industries Africa – active company presence"),
    ("Target",        "OEMs, manufacturers, mining ops, solar developers"),
]
ty = Inches(2.02)
for k, v in snap:
    add_textbox(s, k, Inches(0.6), ty, Inches(1.55), Inches(0.28), 10, True, CHARCOAL)
    add_textbox(s, v, Inches(2.22), ty, Inches(3.82), Inches(0.28), 10, False, SILVER)
    ty += Inches(0.5)

card(s, Inches(6.5), Inches(1.4), Inches(6.5), Inches(5.58))
add_rect(s, Inches(6.5), Inches(1.4), Inches(6.5), Inches(0.45), ORANGE)
add_textbox(s, "Digital Health Indicators", Inches(6.66), Inches(1.46),
            Inches(6.2), Inches(0.35), 13, True, WHITE)
metrics = [
    ("Mobile Page Speed",          "~55 / 100",    "Needs improvement for field use",   CAUTION),
    ("Desktop Page Speed",         "~74 / 100",    "Reasonable – target 90+",            CAUTION),
    ("SEO Score",                   "~60 / 100",   "Sector/product pages under-optimised",CAUTION),
    ("SSL / HTTPS",                "✓ Active",     "Valid certificate in place",          SUCCESS),
    ("B2B Customer Portal",        "Not present",  "No self-service ordering / tracking", DANGER),
    ("Product / Parts Catalogue",  "Not present",  "No searchable online catalogue",      DANGER),
    ("RFQ / Quoting Online",       "Not present",  "Enquiry form only — no automation",   DANGER),
    ("Live Inventory Visibility",  "Not present",  "Critical gap for OEM customers",      DANGER),
    ("AI / Personalisation",       "None detected","Major opportunity",                   DANGER),
    ("Social Media Integration",   "Minimal",      "LinkedIn content not embedded",       CAUTION),
]
ty = Inches(2.02)
for lbl, val, note, col in metrics:
    add_textbox(s, lbl,  Inches(6.66), ty, Inches(2.85), Inches(0.28), 10, True, MID_TEXT)
    add_textbox(s, val,  Inches(9.55), ty, Inches(1.18), Inches(0.28), 10, True, col, PP_ALIGN.CENTER)
    add_textbox(s, note, Inches(10.78),ty, Inches(2.08), Inches(0.28),  9, False, SILVER, italic=True)
    ty += Inches(0.46)

# ════════════════════════════════════════════════════════════════
# SLIDE 4 — WEBSITE AUDIT FINDINGS
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Website Audit Findings",
       "btindustries.co.za – key gaps vs. modern B2B industrial expectations")
footer(s)

audit_cols = [
    ("Performance & Technical", CHARCOAL, [
        "Mobile speed ~55 – poor for field/plant engineers",
        "No CDN – assets served from single SA origin only",
        "Render-blocking scripts slow time-to-interactive",
        "Images not optimised (WebP/AVIF missing)",
        "Core Web Vitals LCP & INP need improvement",
    ]),
    ("B2B Experience", ORANGE, [
        "No customer self-service portal (major gap)",
        "No searchable product/parts catalogue online",
        "RFQ process is manual form only — slow & opaque",
        "No order status tracking for existing clients",
        "No account manager contact directly on product pages",
    ]),
    ("SEO & Content", STEEL, [
        "Sector and product pages thin on keyword content",
        "No structured data (Product / Organization schema)",
        "Blog/news section under-utilised for thought leadership",
        "Missing local SEO signals for Gauteng / Cape Town / Durban",
        "DoGo Power / solar story not amplified for new market",
    ]),
    ("Trust & Conversion", TEAL, [
        "ISO 9001 & B-BBEE credentials buried on site",
        "No case studies or client success stories visible",
        "No sustainability / ESG section despite solar pivot",
        "65-year heritage story not told compellingly",
        "No live chat / WhatsApp CTA for urgent B2B enquiries",
    ]),
]
for i, (title, col, bullets) in enumerate(audit_cols):
    x = Inches(0.4) + i * Inches(3.22)
    y = Inches(1.45)
    card(s, x, y, Inches(3.05), Inches(5.55))
    add_rect(s, x, y, Inches(3.05), Inches(0.45), col)
    tc = CHARCOAL if col in (ORANGE, AMBER) else WHITE
    add_textbox(s, title, x + Inches(0.12), y + Inches(0.06),
                Inches(2.8), Inches(0.36), 12, True, tc)
    ty = y + Inches(0.62)
    for b in bullets:
        add_rect(s, x + Inches(0.17), ty + Inches(0.08), Inches(0.1), Inches(0.1), col)
        add_textbox(s, b, x + Inches(0.37), ty, Inches(2.55), Inches(0.5),
                    10, False, MID_TEXT)
        ty += Inches(0.73)

accent_bar(s, Inches(0.4), Inches(6.86), Inches(12.5))
add_textbox(s,
            "Biggest gap: No B2B self-service portal — OEM customers expect real-time inventory, orders & RFQs online in 2026",
            Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.28),
            10.5, True, MID_TEXT, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 5 — WEB DEVELOPMENT ROADMAP
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Web Development Roadmap",
       "Building a world-class B2B digital platform for BT Industries")
footer(s)

card(s, Inches(0.4), Inches(1.4), Inches(4.62), Inches(5.62))
add_rect(s, Inches(0.4), Inches(1.4), Inches(4.62), Inches(0.45), DARK_GREY)
add_textbox(s, "Recommended Tech Stack", Inches(0.56), Inches(1.46),
            Inches(4.38), Inches(0.35), 13, True, ORANGE)
stack = [
    ("Frontend",    "Next.js 15 — fast, SEO-ready, B2B UX"),
    ("CMS",         "Sanity.io — product data, news, sector pages"),
    ("B2B Portal",  "Custom React dashboard — orders, RFQ, inventory"),
    ("Catalogue",   "Elasticsearch — 10,000s of SKUs, instant search"),
    ("ERP Link",    "REST / GraphQL API to SAP / Sage / SYSPRO"),
    ("CRM",         "HubSpot or Salesforce — B2B lead & account mgmt"),
    ("Quotes",      "Quoter / Custom RFQ engine — auto-pricing rules"),
    ("WhatsApp",    "Twilio / Meta Cloud API — B2B instant comms"),
    ("Analytics",   "GA4 + Segment CDP + HubSpot revenue attribution"),
    ("Hosting",     "Vercel + Cloudflare CDN — SA + global edge"),
]
ty = Inches(2.0)
for k, v in stack:
    add_textbox(s, k, Inches(0.6), ty, Inches(1.45), Inches(0.27), 10, True, ORANGE)
    add_textbox(s, v, Inches(2.1), ty, Inches(2.82), Inches(0.27), 10, False, MID_TEXT)
    ty += Inches(0.46)

card(s, Inches(5.2), Inches(1.4), Inches(7.75), Inches(5.62))
add_rect(s, Inches(5.2), Inches(1.4), Inches(7.75), Inches(0.45), CHARCOAL)
add_textbox(s, "Key Improvements & Platform Features",
            Inches(5.36), Inches(1.46), Inches(7.4), Inches(0.35), 13, True, ORANGE)
improvements = [
    ("B2B Self-Service Customer Portal",
     "Secure client login: live inventory levels, order history, delivery tracking, invoice download, account statements — zero phone calls needed."),
    ("Searchable Product & Parts Catalogue",
     "100,000+ SKUs searchable by part number, spec, application, sector. Cross-reference by OEM code. PDF spec sheet download per item."),
    ("Automated RFQ & Digital Quoting Engine",
     "Customers submit RFQs online. System auto-prices standard items, routes complex queries to the right account manager with SLA timers."),
    ("ERP / Inventory API Integration",
     "Real-time stock availability fed from SAP/Sage/SYSPRO into the web portal — customers see live stock before they enquire."),
    ("Sector & Solution Landing Pages (SEO)",
     "Dedicated, keyword-rich pages per sector: Automotive, Mining, Renewable Energy, Aftermarket — each with case studies and schema markup."),
    ("65-Year Heritage & Sustainability Story",
     "Compelling About section: timeline, milestones, ISO 9001, B-BBEE Level 2, DoGo Power solar partnership, ESG commitments."),
    ("WhatsApp Business & Live Chat",
     "Floating WhatsApp CTA — critical for SA B2B — plus Intercom live chat for RFQ follow-up and urgent supply queries."),
]
ty = Inches(1.98)
for title, desc in improvements:
    add_rect(s, Inches(5.36), ty + Inches(0.12), Inches(0.08), Inches(0.22), ORANGE)
    add_textbox(s, title, Inches(5.56), ty, Inches(7.1), Inches(0.28), 11, True, CHARCOAL)
    add_textbox(s, desc,  Inches(5.56), ty + Inches(0.3), Inches(7.1), Inches(0.38), 10, False, SILVER)
    ty += Inches(0.7)

# ════════════════════════════════════════════════════════════════
# SLIDE 6 — AI AUTOMATION OPPORTUNITY MAP
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "AI Automation Opportunity Map",
       "Where AI transforms BT Industries' supply chain, operations & client service")
footer(s)

opps = [
    (CHARCOAL, ORANGE, "Supply Chain & Inventory Intelligence", [
        "AI demand forecasting — predict stock needs per client & season",
        "Automated re-order triggers based on usage patterns",
        "Supplier lead time risk alerts & mitigation recommendations",
        "Multi-warehouse inventory optimisation across SA locations",
        "Solar/renewable demand spike detection (DoGo Power pipeline)",
    ]),
    (ORANGE, CHARCOAL, "B2B Customer Experience & Sales", [
        "AI chatbot: 24/7 product enquiries, spec lookups, RFQ status",
        "WhatsApp bot: instant stock check & order confirmation",
        "Personalised product recommendations per sector / account",
        "Automated quote follow-up sequences (CRM-triggered)",
        "Account health scoring — flag at-risk or growth accounts",
    ]),
    (STEEL, WHITE, "Procurement & Vendor Management", [
        "AI-powered supplier performance scoring & benchmarking",
        "Automated PO generation when stock hits reorder point",
        "Invoice processing & 3-way match (PO / GRN / invoice)",
        "Contract review: flag non-standard terms or price deviations",
        "Freight & clearing document automation (import logistics)",
    ]),
    (TEAL, WHITE, "Quality, Compliance & Reporting", [
        "ISO 9001 audit trail automation & non-conformance tracking",
        "B-BBEE scorecard data collection & reporting automation",
        "AI-generated monthly operations & client performance reports",
        "Defect image recognition at goods receiving (computer vision)",
        "Regulatory compliance alerts (SABS, customs, SARS changes)",
    ]),
]
for i, (hdr_bg, hdr_txt, title, bullets) in enumerate(opps):
    c, r = i % 2, i // 2
    x = Inches(0.4)  + c * Inches(6.5)
    y = Inches(1.45) + r * Inches(2.88)
    card(s, x, y, Inches(6.2), Inches(2.72))
    add_rect(s, x, y, Inches(6.2), Inches(0.45), hdr_bg)
    add_textbox(s, title, x + Inches(0.15), y + Inches(0.07),
                Inches(5.9), Inches(0.35), 12, True, hdr_txt)
    ty = y + Inches(0.58)
    for b in bullets:
        add_rect(s, x + Inches(0.18), ty + Inches(0.07),
                 Inches(0.1), Inches(0.1), ORANGE)
        add_textbox(s, b, x + Inches(0.38), ty, Inches(5.65), Inches(0.38),
                    10.5, False, MID_TEXT)
        ty += Inches(0.43)

# ════════════════════════════════════════════════════════════════
# SLIDE 7 — BENEFITS OF AUTOMATION (hero)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Benefits of AI Automation",
       "Measurable impact across BT Industries' operations, supply chain & revenue")
footer(s)

for i, (val, lbl) in enumerate([
    ("55–65%", "Reduction in\nmanual admin tasks"),
    ("30%",    "Lower inventory\nholding costs (AI forecasting)"),
    ("3×",     "Faster RFQ-to-order\nconversion cycle"),
    ("24/7",   "B2B customer service\nvia AI & WhatsApp bot"),
]):
    stat_card(s, val, lbl, Inches(0.4) + i * Inches(3.2), Inches(1.45),
              w=Inches(3.0), h=Inches(1.65))

benefits = [
    ("Supply Continuity & Reliability", [
        "AI forecasting prevents stock-outs for OEM clients",
        "Supplier risk alerts keep production lines running",
        "Automated re-orders eliminate manual monitoring",
        "'Confidence of Supply' delivered by data, not guesswork",
    ]),
    ("Operational Cost Efficiency", [
        "Slash admin time on POs, invoices & reporting by >55%",
        "Reduce excess inventory with demand-driven ordering",
        "Automate ISO 9001 compliance reporting — zero manual effort",
        "Free up key staff for strategic account management",
    ]),
    ("B2B Customer Experience", [
        "Instant 24/7 stock checks via WhatsApp or web portal",
        "Self-service RFQ + order tracking — no phone queues",
        "Personalised account dashboards for each OEM client",
        "Automated delivery updates and proof-of-delivery docs",
    ]),
    ("Revenue Growth & New Markets", [
        "Digital catalogue opens door to long-tail & SME clients",
        "Solar / renewable energy sector served digitally at scale",
        "Data-driven upsell: AI suggests complementary products",
        "National reach beyond current SA footprint via e-commerce",
    ]),
]
for i, (title, bullets) in enumerate(benefits):
    c, r = i % 2, i // 2
    x = Inches(0.4)  + c * Inches(6.5)
    y = Inches(3.28) + r * Inches(1.88)
    bullet_card(s, title, bullets, x, y, Inches(6.2), Inches(1.72),
                hdr_bg=CHARCOAL if r == 0 else DARK_GREY, hdr_col=ORANGE)

# ════════════════════════════════════════════════════════════════
# SLIDE 8 — ROI & BUSINESS CASE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "ROI & Business Case",
       "Projected returns from web + AI investment for BT Industries")
footer(s)

card(s, Inches(0.4), Inches(1.4), Inches(5.8), Inches(5.62))
add_rect(s, Inches(0.4), Inches(1.4), Inches(5.8), Inches(0.45), CHARCOAL)
add_textbox(s, "Indicative Investment (Year 1 – ZAR)",
            Inches(0.56), Inches(1.46), Inches(5.5), Inches(0.35), 13, True, ORANGE)
invest = [
    ("Next.js website rebuild + Sanity CMS",            "R 200 000 – R 320 000"),
    ("B2B Customer Portal (orders, RFQ, tracking)",     "R 250 000 – R 420 000"),
    ("Product catalogue + Elasticsearch search",        "R  80 000 – R 140 000"),
    ("ERP / Inventory API integration",                  "R 120 000 – R 200 000"),
    ("AI chatbot + WhatsApp B2B bot",                   "R  60 000 – R 100 000"),
    ("AI demand forecasting module",                     "R  90 000 – R 160 000"),
    ("SEO, schema & content sprint",                    "R  40 000 – R  70 000"),
    ("HubSpot CRM + quote automation",                  "R  50 000 – R  90 000"),
    ("Monthly support / hosting / AI APIs (est.)",      "R  20 000 – R  35 000 / mo"),
    ("TOTAL YEAR 1 (once-off + 12 mo support)",        "R 1.13M – R 1.92M"),
]
ty = Inches(2.04)
for i, (item, cost) in enumerate(invest):
    last = i == len(invest) - 1
    if last:
        accent_bar(s, Inches(0.4), ty - Inches(0.06), Inches(5.8))
        ty += Inches(0.1)
    add_textbox(s, item, Inches(0.6), ty, Inches(3.65), Inches(0.28),
                9.5 if not last else 11, last, MID_TEXT)
    add_textbox(s, cost, Inches(4.28), ty, Inches(1.75), Inches(0.28),
                9.5 if not last else 11, last,
                SUCCESS if last else STEEL, PP_ALIGN.RIGHT)
    ty += Inches(0.46)

card(s, Inches(6.5), Inches(1.4), Inches(6.5), Inches(5.62))
add_rect(s, Inches(6.5), Inches(1.4), Inches(6.5), Inches(0.45), SUCCESS)
add_textbox(s, "Projected Returns & Savings (ZAR / Year)",
            Inches(6.66), Inches(1.46), Inches(6.2), Inches(0.35), 13, True, WHITE)
returns = [
    ("Admin & procurement staff time saved",            "1.5–3 FTE / yr  ≈ R 500 000+"),
    ("Inventory cost reduction (AI forecasting)",       "~R 300 000 – R 600 000 / yr"),
    ("Faster RFQ → order cycle (reduced lost deals)",   "+20–30% quote conversion"),
    ("Self-service portal reduces call centre load",    "−40% inbound query volume"),
    ("New SME & solar clients via digital catalogue",   "Significant new revenue stream"),
    ("Reduced stock-out incidents for OEM clients",     "Prevents penalty clauses"),
    ("B-BBEE & ISO reporting automation savings",       "~R 80 000 – R 150 000 / yr"),
    ("ESTIMATED YEAR 1 SAVINGS / UPSIDE",              "R 1.2M – R 2.5M+"),
]
ty = Inches(2.04)
for i, (item, val) in enumerate(returns):
    last = i == len(returns) - 1
    if last:
        accent_bar(s, Inches(6.5), ty - Inches(0.06), Inches(6.5))
        ty += Inches(0.1)
    add_textbox(s, item, Inches(6.66), ty, Inches(3.88), Inches(0.28),
                9.5 if not last else 11, last, MID_TEXT)
    add_textbox(s, val, Inches(10.58), ty, Inches(2.28), Inches(0.28),
                9.5 if not last else 11, last,
                SUCCESS if last else STEEL, PP_ALIGN.RIGHT)
    ty += Inches(0.46)

add_textbox(s,
            "Preventing a single major OEM stock-out event pays for the entire programme  |  AI forecasting ROI positive within 6 months",
            Inches(0.4), Inches(7.07), Inches(12.5), Inches(0.28),
            10.5, True, MID_TEXT, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 9 — IMPLEMENTATION TIMELINE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Implementation Timeline", "Phased 12-month delivery plan")
footer(s)

phases = [
    ("Phase 1\nMon 1–3",   CHARCOAL, [
        "Site speed & SEO audit + quick fixes",
        "WhatsApp CTA & live chat live",
        "Sanity CMS + blog / news migration",
        "HubSpot CRM + RFQ lead routing",
        "Sector landing pages (SEO sprint)",
    ]),
    ("Phase 2\nMon 4–6",   ORANGE,   [
        "Next.js full website rebuild launch",
        "Searchable product catalogue (MVP)",
        "Digital RFQ & quoting engine live",
        "AI chatbot launch on site",
        "ERP inventory API integration (MVP)",
    ]),
    ("Phase 3\nMon 7–9",   STEEL,    [
        "B2B customer self-service portal",
        "Live inventory visibility per account",
        "Order tracking + delivery updates",
        "AI demand forecasting module live",
        "Automated PO / re-order triggers",
    ]),
    ("Phase 4\nMon 10–12", TEAL,     [
        "AI supplier risk scoring",
        "ISO 9001 compliance automation",
        "B-BBEE reporting automation",
        "Computer vision quality check pilot",
        "Full analytics audit & Year 2 plan",
    ]),
]
for i, (label, col, items) in enumerate(phases):
    x = Inches(0.4) + i * Inches(3.22)
    y = Inches(1.45)
    card(s, x, y, Inches(3.05), Inches(5.62))
    add_rect(s, x, y, Inches(3.05), Inches(0.72), col)
    tc = CHARCOAL if col == ORANGE else WHITE
    add_textbox(s, label, x, y + Inches(0.06), Inches(3.05), Inches(0.62),
                13, True, tc, PP_ALIGN.CENTER)
    if i < 3:
        add_rect(s, x + Inches(3.07), y + Inches(0.3),
                 Inches(0.14), Inches(0.06), SILVER)
    ty = y + Inches(0.9)
    for item in items:
        add_rect(s, x + Inches(0.18), ty + Inches(0.08),
                 Inches(0.1), Inches(0.1), col)
        add_textbox(s, item, x + Inches(0.38), ty, Inches(2.55), Inches(0.42),
                    10.5, False, MID_TEXT)
        ty += Inches(0.84)

# ════════════════════════════════════════════════════════════════
# SLIDE 10 — WHY AUTOMATE NOW?
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, CHARCOAL)
add_rect(s, 0, 0, Inches(0.6), SLIDE_H, ORANGE)
footer(s)

add_textbox(s, "Why Automate — Now?",
            Inches(0.9), Inches(0.45), Inches(11), Inches(0.68),
            30, True, ORANGE)
add_textbox(s,
            "Industrial distribution is being disrupted globally. Digital-first competitors are winning contracts before a salesperson ever makes a call.",
            Inches(0.9), Inches(1.18), Inches(11.8), Inches(0.55),
            13.5, False, AMBER, italic=True)

reasons = [
    (ORANGE,  "OEM procurement is moving fully digital — right now",
     "South Africa's automotive and mining OEMs are mandating supplier portal integration and EDI connectivity. Distributors without digital order channels risk losing preferred supplier status within 24 months."),
    (AMBER,   "Supply chain disruption is the new normal — AI is the antidote",
     "Post-COVID and geopolitical volatility have made demand forecasting critical. AI-driven inventory models reduce stock-out risk by 40–60% while cutting excess carrying costs simultaneously."),
    (SUCCESS, "The solar / renewable sector is moving fast — digital wins the deals",
     "The DoGo Power partnership opens a high-growth market. Solar developers and EPCs issue RFQs digitally and make decisions in days. A modern digital quoting platform captures this revenue that a phone-and-email process misses."),
    (STEEL,   "65 years of data is an untapped competitive weapon",
     "BT Industries holds decades of transaction, component, and supply chain data. AI unlocks it for demand forecasting, supplier scoring, and market intelligence that no competitor can replicate."),
    (TEAL,    "B-BBEE Level 2 and ISO 9001 are advantages — only if clients can see them",
     "Certification and transformation credentials are increasingly part of procurement scoring. A modern, content-rich digital presence ensures these differentiators are visible, verifiable and compelling at the point of evaluation."),
]
ty = Inches(1.9)
for i, (col, title, desc) in enumerate(reasons):
    add_rect(s, Inches(0.9), ty, Inches(0.55), Inches(0.52), col)
    add_textbox(s, str(i + 1), Inches(0.9), ty,
                Inches(0.55), Inches(0.52), 17, True, CHARCOAL, PP_ALIGN.CENTER)
    add_textbox(s, title, Inches(1.6), ty, Inches(11.3), Inches(0.3),
                12, True, WHITE)
    add_textbox(s, desc, Inches(1.6), ty + Inches(0.3), Inches(11.3), Inches(0.52),
                10.5, False, RGBColor(0xAA, 0xA8, 0xA0))
    ty += Inches(0.97)

# ════════════════════════════════════════════════════════════════
# SLIDE 11 — NEXT STEPS
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Next Steps", "From audit to action – what happens now")
footer(s)

steps = [
    ("Week 1",  CHARCOAL, "Stakeholder Alignment",
     "Present findings to BT Industries leadership. Agree priority workstreams, budget envelope and ERP/system access requirements."),
    ("Week 2",  ORANGE,   "Discovery & Access",
     "Share access to website CMS, ERP/inventory system, CRM and analytics. Technical architecture review and data mapping session."),
    ("Week 3",  STEEL,    "Sprint 1 Begins",
     "Speed & SEO fixes, WhatsApp CTA, HubSpot CRM setup and sector landing pages running in parallel across dedicated tracks."),
    ("Month 2", TEAL,     "Phase 1 Go-Live",
     "Improved site live. RFQ digital routing active. WhatsApp bot handling stock enquiries. HubSpot welcome sequence for new leads."),
]
for i, (when, col, title, desc) in enumerate(steps):
    x = Inches(0.4) + i * Inches(3.22)
    y = Inches(1.5)
    card(s, x, y, Inches(3.05), Inches(4.1))
    add_rect(s, x, y, Inches(3.05), Inches(0.45), col)
    tc = CHARCOAL if col == ORANGE else WHITE
    add_textbox(s, when, x, y + Inches(0.06), Inches(3.05), Inches(0.35),
                12, True, tc, PP_ALIGN.CENTER)
    add_textbox(s, title, x + Inches(0.15), y + Inches(0.6),
                Inches(2.75), Inches(0.36), 12, True,
                ORANGE if col != ORANGE else CHARCOAL)
    add_textbox(s, desc, x + Inches(0.15), y + Inches(1.05),
                Inches(2.75), Inches(2.9), 10.5, False, SILVER)

card(s, Inches(0.4), Inches(5.82), Inches(12.5), Inches(1.1),
     fill=CHARCOAL, border=False)
accent_bar(s, Inches(0.4), Inches(5.82), Inches(12.5), Inches(0.055))
add_textbox(s, "Ready to give your OEM clients the Confidence of Supply they expect — digitally?",
            Inches(0.65), Inches(5.92), Inches(8.5), Inches(0.4),
            14, True, ORANGE)
add_textbox(s, "Contact Genos Apollo to schedule your discovery workshop and ERP integration scoping session.",
            Inches(0.65), Inches(6.32), Inches(8.5), Inches(0.38),
            11, False, WHITE)
add_textbox(s, "btindustries.co.za",
            Inches(9.2), Inches(6.1), Inches(3.8), Inches(0.35),
            10, False, AMBER, PP_ALIGN.RIGHT, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 12 — BACK COVER
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, CHARCOAL)
add_rect(s, 0, 0, Inches(0.72), SLIDE_H, ORANGE)
add_rect(s, Inches(0.72), Inches(3.55),
         SLIDE_W - Inches(0.72), Inches(0.055), WHITE)

add_textbox(s, "Thank You",
            Inches(1.1), Inches(1.9), Inches(11), Inches(1.3),
            52, True, WHITE)
add_textbox(s,
            "65 years of supply chain mastery — now powered by data, AI and digital excellence.",
            Inches(1.1), Inches(3.18), Inches(10.5), Inches(0.55),
            15, False, AMBER, italic=True)
add_textbox(s,
            "BT Industries  •  btindustries.co.za  •  ISO 9001  •  Level 2 B-BBEE",
            Inches(1.1), Inches(4.05), Inches(10), Inches(0.4),
            12, False, RGBColor(0x72, 0x78, 0x80))
add_textbox(s, "Prepared by Genos Apollo  |  May 2026  |  Confidential",
            Inches(1.1), Inches(6.7), Inches(10), Inches(0.35),
            10, False, SILVER, italic=True)

# ── Save ────────────────────────────────────────────────────────
out = r"c:\Users\deevansho\Desktop\GenosApollp clients\BTIndustries_WebDev_AI_Automation_Audit.pptx"
prs.save(out)
print(f"Saved: {out}")
