import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.parent / "utils"))
from ppt_theme import (
    NAVY, DARK_BLUE, ACCENT, GOLD, RED, GREEN, PURPLE, ORANGE, TEAL, CYAN,
    WHITE, LGREY, MGREY, W, H,
    bg, rect, tb, aline, bullets, tag,
    slide_header, slide_footer, score_badge,
    stat_cards, top_bar_card, phase_card, before_after, genos_logo_block,
)

ROOT = Path(__file__).parent.parent.parent
prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ══ SLIDE 1 — COVER ══
s1 = prs.slides.add_slide(BLANK)
bg(s1)
rect(s1, 0, 0, Inches(0.12), H, GOLD)
rect(s1, Inches(8.5), 0, Inches(4.83), Inches(2.2), DARK_BLUE)
tb(s1, "GENOS AI", Inches(9.2), Inches(0.4), Inches(3.5), Inches(0.8),
   size=28, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
tb(s1, "www.genosai.tech", Inches(9.2), Inches(1.05), Inches(3.5), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
tb(s1, "BBB OF CHICAGO & NORTHERN ILLINOIS", Inches(0.55), Inches(1.6),
   Inches(9), Inches(0.6), size=13, bold=False, color=GOLD)
tb(s1, "Website & AI Growth\nAudit Report", Inches(0.55), Inches(2.1),
   Inches(9), Inches(1.8), size=48, bold=True, color=WHITE)
aline(s1, Inches(3.85), color=GOLD, w=Inches(3))
tb(s1, "Prepared exclusively for Steve Bernas, President & CEO",
   Inches(0.55), Inches(4.15), Inches(9), Inches(0.5), size=14, color=LGREY)
rect(s1, 0, Inches(6.5), W, Inches(1.0), DARK_BLUE)
tb(s1, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699",
   Inches(0.55), Inches(6.6), Inches(9), Inches(0.5), size=12, color=MGREY)
tb(s1, "May 2026  |  CONFIDENTIAL", Inches(9.5), Inches(6.6), Inches(3.3), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
rect(s1, Inches(10.5), Inches(3.3), Inches(2.3), Inches(2.5), DARK_BLUE)
tb(s1, "AUDIT SCORE", Inches(10.55), Inches(3.4), Inches(2.2), Inches(0.5),
   size=11, bold=True, color=MGREY, align=PP_ALIGN.CENTER)
tb(s1, "6", Inches(10.55), Inches(3.75), Inches(2.2), Inches(1.1),
   size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
tb(s1, "/ 10", Inches(10.55), Inches(4.7), Inches(2.2), Inches(0.4),
   size=18, color=LGREY, align=PP_ALIGN.CENTER)
tag(s1, "MEDIUM PRIORITY", Inches(10.65), Inches(5.1), fill=GOLD, tc=NAVY, size=10)


# ══ SLIDE 2 — AT A GLANCE ══
s2 = prs.slides.add_slide(BLANK)
bg(s2)
rect(s2, 0, 0, Inches(0.12), H, GOLD)
tb(s2, "AT A GLANCE", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GOLD)
tb(s2, "Steve Bernas — Who He Is",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s2, Inches(1.35), color=GOLD)

stats = [
    ("100\nYears", "BBB Chicago Founded 1926"),
    ("9,000+", "Accredited Businesses"),
    ("35+", "Years Steve at BBB"),
    ("Apr 22\n2026", "Centennial at Navy Pier"),
    ("Daily\nHerald", "Regular Columnist"),
]
cw = Inches(2.3); ch = Inches(2.2); gap = Inches(0.18)
for i, (num, label) in enumerate(stats):
    l = Inches(0.55) + i * (cw + gap)
    rect(s2, l, Inches(1.8), cw, ch, DARK_BLUE)
    tb(s2, num, l + Inches(0.1), Inches(1.95), cw - Inches(0.2), Inches(1.1),
       size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s2, label, l + Inches(0.1), Inches(3.05), cw - Inches(0.2), Inches(0.7),
       size=12, color=LGREY, align=PP_ALIGN.CENTER)

rect(s2, Inches(0.55), Inches(4.2), Inches(12.23), Inches(0.6), DARK_BLUE)
tb(s2, "121 W Wacker Dr, Suite 2000, Chicago IL  ·  Centennial: Navy Pier w/ Lester Holt  ·  Media: NBC, ABC, CBS, Daily Herald  ·  Twitter: @SteveBernasBBB",
   Inches(0.65), Inches(4.28), Inches(12), Inches(0.45), size=12, color=LGREY, align=PP_ALIGN.CENTER)

tb(s2, "MISSION", Inches(0.55), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=GOLD)
bullets(s2, ["Advance marketplace trust across Chicago & Northern Illinois",
             "Accreditation, complaint resolution, scam alerts",
             "Consumer education and business standards"],
        Inches(0.55), Inches(5.35), Inches(5.8), Inches(1.5), size=13, color=LGREY)

tb(s2, "EDUCATION & BACKGROUND", Inches(6.8), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=GOLD)
bullets(s2, ["BA Psychology, Loyola / Niles College (1987)",
             "35+ years growing BBB Chicago membership",
             "Regular media voice on consumer protection"],
        Inches(6.8), Inches(5.35), Inches(5.8), Inches(1.5), size=13, color=LGREY)


# ══ SLIDE 3 — WHAT'S WORKING ══
s3 = prs.slides.add_slide(BLANK)
bg(s3)
rect(s3, 0, 0, Inches(0.12), H, GREEN)
tb(s3, "WHAT'S WORKING", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s3, "Strengths to Build On",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s3, Inches(1.35), color=GREEN)

strengths = [
    ("100 Years of Brand Authority",
     "Very few organizations in Chicago can say they've operated continuously for a century. BBB Chicago's 1926 founding, the Navy Pier centennial, and Mayor Brandon Johnson's presence are powerful trust signals that almost no competitor can match."),
    ("Steve Bernas Is a Recognized Public Voice",
     "Regular Daily Herald columns, TV appearances on NBC/ABC/CBS, and Twitter presence give Steve credibility as a consumer protection expert - not just an executive. That personal brand is an underused marketing asset."),
    ("9,000 Accredited Business Network",
     "A base of 9,000 active accredited businesses represents both recurring revenue and a distribution channel. Renewal communications, event invitations, and value-add content can all be automated and personalized at scale."),
    ("Active Scam Alert Publishing",
     "BBB consistently publishes consumer scam warnings across multiple channels. This high-frequency content output is infrastructure for an automated syndication pipeline - written once, distributed everywhere automatically."),
    ("Centennial Media Moment Just Happened",
     "The April 22 Navy Pier centennial with Lester Holt is a once-in-a-century credibility moment. That story should be working on the website, in email sequences, and in social content for months, not just the week of the event."),
]
for i, (title, desc) in enumerate(strengths):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s3, Inches(0.55), t, Inches(0.06), Inches(0.7), GREEN)
    tb(s3, title, Inches(0.75), t, Inches(11.8), Inches(0.38), size=14, bold=True, color=WHITE)
    tb(s3, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# ══ SLIDE 4 — WHERE LEADS ARE LEAKING ══
s4 = prs.slides.add_slide(BLANK)
bg(s4)
rect(s4, 0, 0, Inches(0.12), H, RED)
tb(s4, "WHERE IT'S COSTING YOU", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=RED)
tb(s4, "Operational & Digital Gaps Found in the Audit",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s4, Inches(1.35), color=RED)

issues = [
    ("COMPLAINT TRIAGE IS MANUAL", "HIGH",
     "Thousands of complaints processed by hand means slower resolution times, inconsistent routing, and no ability to detect emerging fraud patterns early. AI triage fixes this - routing by category, flagging patterns, auto-acknowledging intake."),
    ("ACCREDITATION APPLICATIONS SIT IN A QUEUE", "HIGH",
     "Businesses that want to get accredited wait for a manual review workflow. Every day of delay is a day they might lose interest or go to a competitor. Automation speeds review and keeps applicants warm during the process."),
    ("CENTENNIAL CONTENT NOT ON LOCAL PAGES", "HIGH",
     "The Navy Pier centennial happened 2 weeks ago. Mayor Johnson spoke. Lester Holt was there. None of this shows up on bbb.org/chicago. A business owner researching whether to get accredited sees none of that credibility."),
    ("MEMBER RENEWALS ARE NON-PERSONALIZED", "MEDIUM",
     "9,000 accredited businesses get the same renewal communication regardless of industry, tenure, or engagement. Personalized sequences by segment - by industry, years as a member, complaint history - meaningfully reduce lapse rates."),
    ("NO AI CHAT FOR AFTER-HOURS INQUIRIES", "MEDIUM",
     "Businesses deciding whether to seek accreditation often research at night. Consumers with complaints do the same. With no chat widget, those inquiries either wait until Monday or go unanswered entirely."),
]
for i, (title, pri, desc) in enumerate(issues):
    t = Inches(1.7) + i * Inches(1.05)
    pc = RED if pri == "HIGH" else GOLD
    rect(s4, Inches(0.55), t, Inches(0.06), Inches(0.7), pc)
    tb(s4, title, Inches(0.75), t, Inches(9.2), Inches(0.38), size=13, bold=True, color=WHITE)
    tag(s4, pri, Inches(10.1), t + Inches(0.05), fill=pc, size=9)
    tb(s4, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# ══ SLIDE 5 — WEBSITE DESIGN AUDIT ══
s5 = prs.slides.add_slide(BLANK)
bg(s5)
rect(s5, 0, 0, Inches(0.12), H, PURPLE)
tb(s5, "WEBSITE DESIGN AUDIT", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=PURPLE)
tb(s5, "100 Years of Credibility. Local Pages Don't Show It.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.8), size=32, bold=True, color=WHITE)
aline(s5, Inches(1.45), color=PURPLE, w=Inches(1.5))

# 3-column
rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(3.6), DARK_BLUE)
rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(0.42), RED)
tb(s5, "BBB CHICAGO TODAY", Inches(0.65), Inches(1.92), Inches(3.6), Inches(0.38),
   size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Centennial milestone absent from local pages",
    "No video or media coverage embedded",
    "Generic institutional design, low personality",
    "No live chat or AI chatbot widget",
    "Scam alerts published one channel at a time",
    "Accreditation CTA functional but cold",
    "No personalization for business vs consumer",
    "Steve's media coverage not showcased locally",
], Inches(0.7), Inches(2.42), Inches(3.6), Inches(2.9), size=11, color=LGREY, dot=RED)

tb(s5, "->", Inches(4.45), Inches(3.6), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(3.6), DARK_BLUE)
rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(0.42), GOLD)
tb(s5, "LEADING NONPROFITS DO THIS", Inches(5.1), Inches(1.92), Inches(3.4), Inches(0.38),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Milestone moments front and center (AARP, BBB.com national)",
    "CEO personal brand integrated on homepage",
    "Video testimonials from member businesses",
    "AI chat widget trained on org services",
    "Automated content syndication across channels",
    "Personalized UX: consumer vs business visitor",
    "Media coverage ticker or press wall",
    "Event recaps drive post-event SEO",
], Inches(5.1), Inches(2.42), Inches(3.4), Inches(2.9), size=11, color=LGREY, dot=GOLD)

tb(s5, "->", Inches(8.68), Inches(3.6), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(3.6), DARK_BLUE)
rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(0.42), GREEN)
tb(s5, "WHAT GENOS AI BUILDS", Inches(9.35), Inches(1.92), Inches(3.4), Inches(0.38),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Centennial hero section with Navy Pier photography",
    "Steve's Daily Herald columns auto-feed to local pages",
    "Video highlights from media appearances embedded",
    "AI chat for accreditation and complaint intake 24/7",
    "Scam alert auto-syndication: email + social + press",
    "Personalized entry path: consumer vs business",
    "Member testimonial section with BBB reviews",
    "Mobile-first, modern premium design",
], Inches(9.35), Inches(2.42), Inches(3.4), Inches(2.9), size=11, color=LGREY, dot=GREEN)

rect(s5, Inches(0.55), Inches(5.72), Inches(12.23), Inches(1.28), DARK_BLUE)
tb(s5, "WHY IT MATTERS:", Inches(0.7), Inches(5.78), Inches(2.2), Inches(0.38),
   size=11, bold=True, color=PURPLE)
stats2 = [
    ("94%", "of first impressions\nare design-related"),
    ("75%", "of B2B buyers judge\ncredibility by website"),
    ("3 sec", "before a visitor decides\nto stay or leave"),
    ("200%+", "increase in conversions\nwith personalized UX"),
]
sw = Inches(2.7); sl = Inches(2.9)
for i, (num, label) in enumerate(stats2):
    xl = sl + i * sw
    tb(s5, num, xl, Inches(5.75), Inches(2.5), Inches(0.55),
       size=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s5, label, xl, Inches(6.22), Inches(2.5), Inches(0.7),
       size=10, color=LGREY, align=PP_ALIGN.CENTER)


# ══ SLIDE 6 — AI AUTOMATION OPPORTUNITIES ══
s6 = prs.slides.add_slide(BLANK)
bg(s6)
rect(s6, 0, 0, Inches(0.12), H, ACCENT)
tb(s6, "AI AUTOMATION OPPORTUNITIES", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s6, "If Scammers Use AI, BBB Should Too.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=34, bold=True, color=WHITE)
aline(s6, Inches(1.35), color=ACCENT)

opps = [
    ("Complaint Triage AI", GOLD,
     "Route incoming complaints automatically by category (contractor, retail, financial, online), flag emerging fraud patterns before they hit the news cycle, and send instant intake acknowledgments. Your team reviews decisions - not incoming queues."),
    ("Accreditation Processing Automation", ACCENT,
     "Automate application status updates, document requests, and follow-up sequences for businesses in the accreditation pipeline. Every day a business waits is a day they reconsider. Keep them warm and moving."),
    ("Scam Alert Auto-Syndication", GREEN,
     "Steve writes the Daily Herald column. AI handles the rest: email blast to accredited businesses, LinkedIn post, Twitter/X thread, press release draft, and SEO page update - all triggered from one publication. Zero extra staff time."),
    ("Member Renewal Sequences", PURPLE,
     "9,000 accredited businesses segmented by industry, tenure, and engagement. Renewal outreach personalized to each segment - contractors get different messaging than financial services firms. Automated, scheduled, measurable."),
    ("AI Chat - 24/7 Intake", RED,
     "A chat widget trained on BBB Chicago's services, accreditation process, and common consumer questions. Captures business accreditation inquiries and consumer complaint intake at midnight on a Sunday the same as 9am Monday."),
]
cw2 = Inches(2.35); gap2 = Inches(0.14)
for i, (title, color, desc) in enumerate(opps):
    l2 = Inches(0.55) + i * (cw2 + gap2)
    rect(s6, l2, Inches(1.8), cw2, Inches(0.06), color)
    rect(s6, l2, Inches(1.86), cw2, Inches(4.6), DARK_BLUE)
    tb(s6, title, l2 + Inches(0.12), Inches(1.98), cw2 - Inches(0.24), Inches(0.65),
       size=13, bold=True, color=color)
    tb(s6, desc, l2 + Inches(0.12), Inches(2.65), cw2 - Inches(0.24), Inches(3.6),
       size=11, color=LGREY)

rect(s6, Inches(0.55), Inches(6.7), Inches(12.23), Inches(0.45), DARK_BLUE)
tb(s6, "\"With AI, the fraudsters are very convincing and effective.\" — Steve Bernas, Daily Herald, April 2026  |  BBB Chicago should use the same tools.",
   Inches(0.7), Inches(6.75), Inches(12), Inches(0.35), size=11, italic=True, color=GOLD)


# ══ SLIDE 7 — WHAT THIS IS COSTING YOU ══
s7 = prs.slides.add_slide(BLANK)
bg(s7)
rect(s7, 0, 0, Inches(0.12), H, RED)
tb(s7, "WHAT THIS IS COSTING YOU", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=RED)
tb(s7, "The Hidden Cost of Manual Operations",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s7, Inches(1.35), color=RED)

costs = [
    ("Manual triage delays = fraud patterns missed",
     "When complaints are triaged by hand, emerging scam patterns take weeks to surface. AI spots them in days - before they become a headline you're reacting to instead of leading."),
    ("Non-personalized renewals = silent lapse rate",
     "Accredited businesses that don't renew rarely call to say goodbye. They just don't renew. Personalized renewal sequences by industry and tenure reduce this quietly and measurably."),
    ("No chat = accreditation inquiries lost after hours",
     "A contractor researching BBB accreditation at 10pm on a Tuesday finds no way to engage. They either wait, forget, or decide it's not worth the effort. AI chat recovers every one of those conversations."),
    ("Centennial moment not on local pages = wasted credibility",
     "The Navy Pier centennial, Mayor Johnson, Lester Holt - this is the kind of social proof that converts skeptics into members. It happened two weeks ago and it's not visible on bbb.org/chicago. That's a recurring daily cost."),
    ("Scam alerts published manually = slow and inconsistent reach",
     "Steve writes the column. Someone else formats the email. Another person posts on social. It takes hours. Automation publishes across all channels in minutes - and reaches more people each time a new scam alert drops."),
]
for i, (title, desc) in enumerate(costs):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s7, Inches(0.55), t, Inches(0.06), Inches(0.7), RED)
    tb(s7, title, Inches(0.75), t, Inches(11.8), Inches(0.38), size=14, bold=True, color=WHITE)
    tb(s7, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# ══ SLIDE 8 — WHAT WE'D BUILD ══
s8 = prs.slides.add_slide(BLANK)
bg(s8)
rect(s8, 0, 0, Inches(0.12), H, GREEN)
tb(s8, "WHAT WE'D BUILD", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s8, "Genos AI for BBB Chicago",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s8, Inches(1.35), color=GREEN)

builds = [
    ("bbb.org/chicago Rebuild",
     "Centennial hero section front and center. Steve's Daily Herald column feed auto-updating. Media coverage wall (NBC, ABC, CBS). Personalized entry paths for business vs consumer visitors. Modern premium design that reflects 100 years, not 2008."),
    ("Complaint Triage AI",
     "Automated intake acknowledgment, category routing, and fraud pattern flagging. Your team sees a prioritized queue with recommended actions, not a raw inbox. Emerging scam trends surface in days, not weeks."),
    ("Scam Alert Auto-Syndication Pipeline",
     "Steve publishes one piece. AI formats and distributes: email to accredited businesses, LinkedIn, Twitter/X, press release, and SEO page update - all triggered automatically. One click, eight channels."),
    ("Member Renewal Sequences",
     "9,000 businesses segmented by industry, years as member, and engagement. Automated renewal outreach starts 90 days out, with personalized messaging by segment. Lapse rate goes down without adding headcount."),
    ("AI Chat - Accreditation + Complaint Intake",
     "Widget trained on BBB Chicago's full service offering. Captures accreditation interest and complaint intake 24/7. Routes to the right staff queue. No inquiry left waiting until Monday morning."),
]
lw = Inches(5.8); rw = Inches(5.8)
for i, (title, desc) in enumerate(builds):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s8, Inches(0.55), t, Inches(0.06), Inches(0.7), GREEN)
    tb(s8, title, Inches(0.75), t, lw, Inches(0.38), size=14, bold=True, color=WHITE)
    tb(s8, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# ══ SLIDE 9 — 90-DAY OUTCOMES ══
s9 = prs.slides.add_slide(BLANK)
bg(s9)
rect(s9, 0, 0, Inches(0.12), H, ACCENT)
tb(s9, "90-DAY OUTCOMES", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s9, "What Changes in 90 Days",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s9, Inches(1.35))

phases = [
    ("Month 1", GOLD, [
        "AI chat live - capturing accreditation inquiries and complaint intake 24/7",
        "Complaint triage AI deployed - queue organized, patterns flagging",
        "Scam alert auto-syndication live - Steve's next column goes to 8 channels automatically",
        "bbb.org/chicago rebuild begins - centennial content and design in progress",
    ]),
    ("Month 2", ACCENT, [
        "New Chicago pages live - centennial hero, Steve's media coverage, personalized UX",
        "Accreditation processing automation active - applicants stay warm through the queue",
        "Member renewal sequences running - 90-day pre-renewal outreach by segment",
        "First automated scam alert campaign performance data available",
    ]),
    ("Month 3", GREEN, [
        "Renewal velocity measurably up across key industry segments",
        "Complaint resolution time down - triage AI routing vs. manual",
        "Centennial content driving local SEO - Navy Pier coverage still ranking",
        "AI chat converting after-hours accreditation inquiries to active applications",
    ]),
]
pw = Inches(3.9); gap3 = Inches(0.18)
for i, (month, color, pts) in enumerate(phases):
    l3 = Inches(0.55) + i * (pw + gap3)
    rect(s9, l3, Inches(1.8), pw, Inches(0.42), color)
    tb(s9, month.upper(), l3 + Inches(0.1), Inches(1.84), pw - Inches(0.2), Inches(0.36),
       size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    rect(s9, l3, Inches(2.22), pw, Inches(4.7), DARK_BLUE)
    bullets(s9, pts, l3 + Inches(0.15), Inches(2.35), pw - Inches(0.3), Inches(4.4),
            size=12, color=LGREY, dot=color)


# ══ SLIDE 10 — NEXT STEPS ══
s10 = prs.slides.add_slide(BLANK)
bg(s10)
rect(s10, 0, 0, Inches(0.12), H, GOLD)
rect(s10, Inches(7.5), 0, Inches(5.83), H, DARK_BLUE)
tb(s10, "NEXT STEPS", Inches(0.55), Inches(0.3), Inches(6.5), Inches(0.5),
   size=11, bold=True, color=GOLD)
tb(s10, "Let's Talk for\n15 Minutes.",
   Inches(0.55), Inches(0.8), Inches(6.8), Inches(2.0), size=48, bold=True, color=WHITE)
aline(s10, Inches(2.8), color=GOLD, w=Inches(2.5))
tb(s10, "No pitch deck required on your end. Just reply to the email\nor book a call at the link below.",
   Inches(0.55), Inches(3.05), Inches(6.5), Inches(0.8), size=14, color=LGREY)

steps = [
    "Reply to this email and we'll set up a time",
    "Book directly: www.genosai.tech/call",
    "We'll walk through the audit findings in 15 minutes",
    "You decide what, if anything, makes sense to move on",
]
bullets(s10, steps, Inches(0.55), Inches(3.95), Inches(6.5), Inches(2.0), size=14, color=WHITE, dot=GOLD)

rect(s10, 0, Inches(6.5), Inches(7.5), Inches(1.0), DARK_BLUE)
tb(s10, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699  |  www.genosai.tech",
   Inches(0.55), Inches(6.6), Inches(6.8), Inches(0.5), size=12, color=MGREY)

tb(s10, "WHAT WE BUILT FOR\nBBB CHICAGO", Inches(7.85), Inches(0.6), Inches(5.0), Inches(0.9),
   size=20, bold=True, color=GOLD)
aline(s10, Inches(1.45), color=GOLD)

summary_items = [
    ("Website", "Centennial redesign - modern, premium, Steve's media integrated"),
    ("Complaint AI", "Triage, routing, pattern detection"),
    ("Scam Alerts", "Auto-syndication pipeline - 1 click, 8 channels"),
    ("Renewals", "Personalized sequences for 9,000 accredited businesses"),
    ("AI Chat", "24/7 accreditation + complaint intake widget"),
    ("Accreditation", "Application processing and follow-up automation"),
]
for i, (label, val) in enumerate(summary_items):
    t2 = Inches(1.6) + i * Inches(0.88)
    rect(s10, Inches(7.85), t2, Inches(1.5), Inches(0.72), NAVY)
    tb(s10, label, Inches(7.9), t2 + Inches(0.05), Inches(1.4), Inches(0.32),
       size=10, bold=True, color=GOLD)
    tb(s10, val, Inches(9.5), t2 + Inches(0.05), Inches(3.6), Inches(0.65),
       size=11, color=LGREY)


prs.save(ROOT / "output" / "decks" / "BBBChicago_Audit_GenosAI.pptx")
print("Saved: BBBChicago_Audit_GenosAI.pptx")
