import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.parent / "utils"))
from ppt_theme import (
    NAVY, DARK_BLUE, ACCENT, GOLD, RED, GREEN, PURPLE, ORANGE, TEAL, CYAN,
    WHITE, LGREY, MGREY, W, H,
    bg, rect, tb, aline, bullets, tag,
    slide_header, slide_footer, score_badge,
    stat_cards, top_bar_card, phase_card, before_after, genos_logo_block,
)

ROOT = Path(__file__).parent.parent.parent
prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]

# SLIDE 1 - COVER
s1 = prs.slides.add_slide(BLANK)
bg(s1)
rect(s1, 0, 0, Inches(0.12), H, GOLD)
rect(s1, Inches(8.5), 0, Inches(4.83), Inches(2.2), DARK_BLUE)
tb(s1, "GENOS AI", Inches(9.2), Inches(0.4), Inches(3.5), Inches(0.8),
   size=28, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
tb(s1, "www.genosai.tech", Inches(9.2), Inches(1.05), Inches(3.5), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
tb(s1, "PEAK 1031 EXCHANGE", Inches(0.55), Inches(1.6),
   Inches(9), Inches(0.6), size=13, bold=False, color=GOLD)
tb(s1, "AI Automation\nOpportunity Audit", Inches(0.55), Inches(2.1),
   Inches(9), Inches(1.8), size=48, bold=True, color=WHITE)
aline(s1, Inches(3.85), color=GOLD, w=Inches(3))
tb(s1, "Prepared exclusively for Kevin Levine, EVP & Principal",
   Inches(0.55), Inches(4.15), Inches(9), Inches(0.5), size=14, color=LGREY)
rect(s1, 0, Inches(6.5), W, Inches(1.0), DARK_BLUE)
tb(s1, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699",
   Inches(0.55), Inches(6.6), Inches(9), Inches(0.5), size=12, color=MGREY)
tb(s1, "May 2026  |  CONFIDENTIAL", Inches(9.5), Inches(6.6), Inches(3.3), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
rect(s1, Inches(10.5), Inches(3.3), Inches(2.3), Inches(2.5), DARK_BLUE)
tb(s1, "AUTO. SCORE", Inches(10.55), Inches(3.4), Inches(2.2), Inches(0.5),
   size=11, bold=True, color=MGREY, align=PP_ALIGN.CENTER)
tb(s1, "3", Inches(10.55), Inches(3.75), Inches(2.2), Inches(1.1),
   size=72, bold=True, color=RED, align=PP_ALIGN.CENTER)
tb(s1, "/ 10", Inches(10.55), Inches(4.7), Inches(2.2), Inches(0.4),
   size=18, color=LGREY, align=PP_ALIGN.CENTER)
tag(s1, "HIGH PRIORITY", Inches(10.75), Inches(5.1), fill=RED, size=10)
tb(s1, "3 = current automation maturity\n10 = what's achievable", Inches(10.3), Inches(5.5),
   Inches(2.8), Inches(0.6), size=9, color=MGREY, align=PP_ALIGN.CENTER)


# SLIDE 2 - AT A GLANCE
s2 = prs.slides.add_slide(BLANK)
bg(s2)
rect(s2, 0, 0, Inches(0.12), H, GOLD)
tb(s2, "AT A GLANCE", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GOLD)
tb(s2, "Peak 1031 Exchange - The Operation",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s2, Inches(1.35), color=GOLD)

stats = [
    ("2003", "Founded (Peak Corp. 1991)"),
    ("9+", "Team, National Operation"),
    ("45 / 180", "Day IRS Deadlines"),
    ("20+", "Docs Per Exchange"),
    ("4.9", "Stars, 76+ Reviews"),
]
cw = Inches(2.3); ch = Inches(2.2); gap = Inches(0.18)
for i, (num, label) in enumerate(stats):
    l = Inches(0.55) + i * (cw + gap)
    rect(s2, l, Inches(1.8), cw, ch, DARK_BLUE)
    tb(s2, num, l + Inches(0.1), Inches(1.95), cw - Inches(0.2), Inches(1.1),
       size=26, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s2, label, l + Inches(0.1), Inches(3.05), cw - Inches(0.2), Inches(0.7),
       size=12, color=LGREY, align=PP_ALIGN.CENTER)

rect(s2, Inches(0.55), Inches(4.2), Inches(12.23), Inches(0.6), DARK_BLUE)
tb(s2, "Qualified Intermediary  -  All 50 States  -  FEA Member  -  Woodland Hills CA  -  Part of Peak Corporate Network  -  IRC Section 1031",
   Inches(0.65), Inches(4.28), Inches(12), Inches(0.45), size=12, color=LGREY, align=PP_ALIGN.CENTER)

tb(s2, "EXCHANGE TYPES HANDLED", Inches(0.55), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=GOLD)
bullets(s2, ["Delayed exchanges (45-day ID, 180-day close - most common)",
             "Simultaneous, reverse, and improvement exchanges",
             "Personal property exchanges (equipment, vehicles, IP)"],
        Inches(0.55), Inches(5.35), Inches(5.8), Inches(1.5), size=13, color=LGREY)

tb(s2, "THE DEADLINE REALITY", Inches(6.8), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=GOLD)
bullets(s2, ["One missed 45-day identification = client owes full capital gains tax",
             "Client liability on a $2M sale = $200K-$500K+ tax event",
             "9 people managing this risk nationally on manual workflows"],
        Inches(6.8), Inches(5.35), Inches(5.8), Inches(1.5), size=13, color=LGREY)


# SLIDE 3 - WHAT'S WORKING
s3 = prs.slides.add_slide(BLANK)
bg(s3)
rect(s3, 0, 0, Inches(0.12), H, GREEN)
tb(s3, "WHAT'S WORKING", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s3, "Strengths to Build On",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s3, Inches(1.35), color=GREEN)

strengths = [
    ("20-Year Track Record in a Compliance-Critical, Trust-Based Business",
     "Kevin and Steven co-founded Peak 1031 in 2003 and have operated nationally for over two decades. In a business where one operational error can trigger a six-figure tax liability for a client, that track record is the product. The 4.9-star rating across 76 reviews isn't marketing - it's evidence that the team delivers outcomes under deadline pressure consistently."),
    ("Lean 9-Person Team Delivering National Volume Demonstrates Per-Person Productivity",
     "Operating all 50 states at 9 people means each team member is handling a significant transaction load with high accuracy. That's a cultural and operational foundation that automation multiplies rather than replaces. The team doesn't need more people for the same volume - it needs fewer manual hours per exchange to handle more volume at the same quality."),
    ("FEA Membership and Compliance Depth Set a High Bar for Competitors",
     "Federation of Exchange Accommodators membership signals compliance commitment that smaller QIs can't credibly match. The presence of Steven Rosansky as legal counsel in the operation means exchange structuring and compliance edge cases are handled at a depth that's hard to replicate. Automation doesn't change that - it removes the administrative overhead that sits around it."),
    ("Peak Corporate Network Cross-Referral Creates Structural Deal Flow",
     "Being part of a network that spans mortgage, realty, and property management means 1031 referrals come through channels that competitors without a network can't access. That deal flow advantage compounds if the operational capacity to handle increased volume is there. Right now the 9-person ceiling is the constraint, not the referral network."),
    ("Multi-Exchange Type Capability Captures a Larger Addressable Market",
     "Most QIs default to delayed exchanges. Handling simultaneous, reverse, improvement, and personal property exchanges means Peak can take mandates that single-product operators decline. Each of those exchange types has a different documentation path, compliance checklist, and timeline structure - exactly the kind of multi-variant workflow that AI-driven compliance automation handles efficiently."),
]
for i, (title, desc) in enumerate(strengths):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s3, Inches(0.55), t, Inches(0.06), Inches(0.7), GREEN)
    tb(s3, title, Inches(0.75), t, Inches(11.8), Inches(0.38), size=14, bold=True, color=WHITE)
    tb(s3, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# SLIDE 4 - WHERE AUTOMATION IS MISSING
s4 = prs.slides.add_slide(BLANK)
bg(s4)
rect(s4, 0, 0, Inches(0.12), H, RED)
tb(s4, "WHERE AUTOMATION IS MISSING", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=RED)
tb(s4, "Every Exchange Has 10+ Hours of Manual Work That Shouldn't Require Your Team.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=24, bold=True, color=WHITE)
aline(s4, Inches(1.35), color=RED)

issues = [
    ("45/180-DAY DEADLINE TRACKING RUNS ON MANUAL CALENDARS AND SPREADSHEETS", "CRITICAL",
     "Every active exchange has two non-negotiable IRS deadlines. The current tracking system is calendar entries and spreadsheet rows managed manually. One calendar error, one overlooked row, one team member out sick on a critical day - and a client owes full capital gains on their sale. The liability doesn't fall on Peak legally, but the relationship does. This is the highest-risk manual process in the operation and the fastest to automate."),
    ("20+ DOCUMENTS PER EXCHANGE CHASED AND ORGANISED BY EMAIL", "HIGH",
     "Each exchange file requires sale contracts, title reports, purchase agreements, proof of funds, property descriptions, IRS forms, and state filings. The current collection process is outbound email to clients and counterparties, tracking replies manually, organising by folder, and chasing missing items individually. 5-10 hours per exchange for document assembly that an automated collection and classification system reduces to under 1 hour."),
    ("PROPERTY IDENTIFICATION LETTERS VERIFIED AND PREPARED MANUALLY", "HIGH",
     "IRS rules require exact legal descriptions, addresses, and parcel numbers on identification letters. A wrong street number or incomplete legal description voids the exchange. Current process: manual cross-reference against county records and MLS, spreadsheet comparison, phone calls to confirm details, letter drafted from template. An identification bot handles the cross-reference, pre-fills the compliant letter, and flags discrepancies before submission."),
    ("CLIENTS CALL OR EMAIL FOR EXCHANGE STATUS AND TIMELINE QUESTIONS", "MEDIUM",
     "Investor calls asking where their exchange stands, what documents are still needed, when the 180-day deadline falls, and how the identification rules work under their specific scenario are handled by whoever picks up the phone. A client status portal with AI chat deflects routine status and FAQ queries automatically, and escalates only the exceptions that genuinely require a senior exchange officer."),
    ("COMPLIANCE CHECKLISTS GENERATED MANUALLY PER EXCHANGE TYPE", "MEDIUM",
     "Delayed, simultaneous, reverse, and improvement exchanges each have different documentation requirements, verification steps, and compliance paths. Each checklist is produced manually, reviewed by senior staff, and tracked on spreadsheets. Rules-based automation generates the correct checklist for each exchange type at intake, tracks completion automatically, and flags gaps before the file reaches senior review."),
]
for i, (title, pri, desc) in enumerate(issues):
    t = Inches(1.7) + i * Inches(1.05)
    pc = RED if pri in ("CRITICAL", "HIGH") else GOLD
    if pri == "CRITICAL": pc = RGBColor(0xCC, 0x00, 0x00)
    rect(s4, Inches(0.55), t, Inches(0.06), Inches(0.7), pc)
    tb(s4, title, Inches(0.75), t, Inches(9.2), Inches(0.38), size=12, bold=True, color=WHITE)
    tag(s4, pri, Inches(10.1), t + Inches(0.05), fill=pc, size=9)
    tb(s4, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# SLIDE 5 - THE CURRENT VS AUTOMATED EXCHANGE
s5 = prs.slides.add_slide(BLANK)
bg(s5)
rect(s5, 0, 0, Inches(0.12), H, PURPLE)
tb(s5, "PROCESS COMPARISON", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=PURPLE)
tb(s5, "One Exchange. Manual vs. Automated.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.8), size=32, bold=True, color=WHITE)
aline(s5, Inches(1.45), color=PURPLE, w=Inches(1.5))

rect(s5, Inches(0.55), Inches(1.9), Inches(5.8), Inches(4.7), DARK_BLUE)
rect(s5, Inches(0.55), Inches(1.9), Inches(5.8), Inches(0.42), RED)
tb(s5, "PEAK 1031 TODAY - MANUAL PROCESS", Inches(0.65), Inches(1.92), Inches(5.6), Inches(0.38),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Day 0: Sale closes, manual calendar entry for Day 45 and Day 180",
    "Days 1-10: Email client for 20+ documents, chase missing items individually",
    "Days 10-40: Manually cross-reference property ID against county records",
    "Day 40-44: Staff reviews identification letter, manual compliance check",
    "Day 45: Identification letter submitted - tracked on spreadsheet",
    "Days 46-170: Document organisation, fund tracking, client phone updates",
    "Days 170-179: Manual reminder calls/emails, compliance checklist review",
    "Day 180: Closing coordinated manually, Form 8824 prepared by staff",
    "Total manual hours per exchange: 10-15+ hours of administrative work",
], Inches(0.7), Inches(2.42), Inches(5.5), Inches(3.9), size=11, color=LGREY, dot=RED)

tb(s5, "->", Inches(6.45), Inches(3.8), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(7.0), Inches(1.9), Inches(5.78), Inches(4.7), DARK_BLUE)
rect(s5, Inches(7.0), Inches(1.9), Inches(5.78), Inches(0.42), GREEN)
tb(s5, "PEAK 1031 AUTOMATED - GENOS AI", Inches(7.1), Inches(1.92), Inches(5.5), Inches(0.38),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Day 0: Sale closes, system auto-calculates Day 45 and Day 180 alerts",
    "Days 1-10: Auto-request sent to client, docs classified on receipt by OCR",
    "Days 10-40: Identification bot cross-references parcel data, pre-fills letter",
    "Day 40-44: Compliance checklist auto-generated, discrepancies flagged",
    "Day 45: Compliant ID letter ready - staff reviews exceptions only",
    "Days 46-170: Client portal shows status, AI chat handles update requests",
    "Days 170-179: Tiered automated alerts to team and client (7/3/1 day)",
    "Day 180: Form 8824 auto-generated from exchange file, staff reviews",
    "Total manual hours per exchange: under 2 hours - exceptions and decisions only",
], Inches(7.15), Inches(2.42), Inches(5.5), Inches(3.9), size=11, color=LGREY, dot=GREEN)


# SLIDE 6 - AI AUTOMATION STACK
s6 = prs.slides.add_slide(BLANK)
bg(s6)
rect(s6, 0, 0, Inches(0.12), H, ACCENT)
tb(s6, "THE AUTOMATION STACK", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s6, "7 Systems. One Exchange. Zero Missed Deadlines.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s6, Inches(1.35), color=ACCENT)

opps = [
    ("Deadline\nTracking + Alerts", GOLD,
     "Auto-calculates 45 and 180-day deadlines from sale close date. Tiered alerts to team and client at 15, 7, 3, and 1 day out. Escalates to senior staff if any action is still open within 48 hours of a critical date. Calendar integration. Zero manual tracking required on any active exchange."),
    ("AI Document\nCollection + OCR", ACCENT,
     "Automated outbound document requests with a specific checklist per exchange type. Incoming documents classified automatically by type. OCR extracts key data fields. Missing items flagged with a specific list sent to client and team. Full document assembly from 5-10 hours to under 1 hour per exchange."),
    ("Property ID\nVerification Bot", GREEN,
     "Client submits property address. System cross-references parcel data, legal description, and county records. Pre-fills the identification letter with IRS-compliant formatting. Flags any discrepancy before the letter goes to client or staff for final review. Identification errors eliminated before submission."),
    ("Client Portal\n+ AI Chat", RED,
     "Investor self-service for exchange status, document checklist, deadline countdown, and IRS timeline FAQs. AI chat trained on 1031 exchange rules handles routine questions and routes complex scenarios to the right exchange officer. Reduces inbound calls and status emails by 40-50%."),
    ("Compliance\nAutomation", PURPLE,
     "Rules-based system generates the correct compliance checklist for each exchange type at intake - delayed, simultaneous, reverse, improvement. Tracks completion automatically. Flags gaps before the file reaches senior review. Senior staff time spent on exceptions, not checklist generation."),
]
cw2 = Inches(2.35); gap2 = Inches(0.14)
for i, (title, color, desc) in enumerate(opps):
    l2 = Inches(0.55) + i * (cw2 + gap2)
    rect(s6, l2, Inches(1.8), cw2, Inches(0.06), color)
    rect(s6, l2, Inches(1.86), cw2, Inches(4.6), DARK_BLUE)
    tb(s6, title, l2 + Inches(0.12), Inches(1.98), cw2 - Inches(0.24), Inches(0.65),
       size=13, bold=True, color=color)
    tb(s6, desc, l2 + Inches(0.12), Inches(2.65), cw2 - Inches(0.24), Inches(3.6),
       size=11, color=LGREY)

rect(s6, Inches(0.55), Inches(6.7), Inches(12.23), Inches(0.45), DARK_BLUE)
tb(s6, "9 people. National 1031 volume. A 45-day deadline where one missed alert costs a client more than Peak earned on that exchange all year.",
   Inches(0.7), Inches(6.75), Inches(12), Inches(0.35), size=11, italic=True, color=GOLD)


# SLIDE 7 - THE CAPACITY MATH
s7 = prs.slides.add_slide(BLANK)
bg(s7)
rect(s7, 0, 0, Inches(0.12), H, RED)
tb(s7, "THE CAPACITY MATH", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=RED)
tb(s7, "Manual Hours Per Exchange Define the Revenue Ceiling.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=30, bold=True, color=WHITE)
aline(s7, Inches(1.35), color=RED)

costs = [
    ("10-15 manual hours per exchange means 9 people cap at roughly X exchanges per month",
     "If each exchange requires 10-15 hours of administrative work spread across the team, the monthly transaction ceiling is directly tied to headcount. Adding transactions means adding people - or accepting that quality and deadline safety degrade as volume increases. Automation breaks that equation: the same 9 people handling 2-hour exchanges instead of 12-hour exchanges can process 5-6x the volume without a single new hire."),
    ("Deadline liability risk scales linearly with active exchange volume under manual tracking",
     "At 20 active exchanges, a manual calendar has 40 critical dates to track. At 50 active exchanges, it has 100. At 100 exchanges, the probability that one date slips under manual management approaches certainty. Automated deadline tracking doesn't degrade with volume - it handles 500 active exchanges with the same reliability as 5. The risk profile stays flat regardless of growth."),
    ("Document chasing is the highest-volume manual workflow and the most replaceable",
     "Sending document request emails, tracking what's received, chasing what's missing, classifying what arrives, and matching files to the right exchange is the largest single block of administrative time per transaction. None of it requires exchange expertise. All of it can be automated. Freeing that time doesn't just increase volume capacity - it redirects senior staff to the compliance and client relationship work that actually requires their knowledge."),
    ("Form 8824 preparation and fund reconciliation are the two highest-error manual workflows",
     "IRS Form 8824 requires exact transaction data from the exchange file. Manual preparation from multiple sources introduces transcription errors that trigger CPA back-and-forth and sometimes IRS notices. Similarly, escrow fund reconciliation against multiple bank accounts is error-prone under manual matching. Both workflows auto-generate from the exchange file with an automated system, with staff reviewing for exceptions rather than building from scratch."),
    ("The capacity constraint is the growth constraint - automation removes it without proportional cost",
     "Peak 1031's competitive advantage is service quality, track record, and compliance depth - not transaction processing speed. But processing speed is currently the ceiling on growth. Removing 8-10 hours of manual overhead per exchange doesn't change what makes Peak valuable - it removes the constraint that limits how much of that value gets delivered and how many clients benefit from it."),
]
for i, (title, desc) in enumerate(costs):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s7, Inches(0.55), t, Inches(0.06), Inches(0.7), RED)
    tb(s7, title, Inches(0.75), t, Inches(11.8), Inches(0.38), size=14, bold=True, color=WHITE)
    tb(s7, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# SLIDE 8 - WHAT WE'D BUILD
s8 = prs.slides.add_slide(BLANK)
bg(s8)
rect(s8, 0, 0, Inches(0.12), H, GREEN)
tb(s8, "WHAT WE'D BUILD", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s8, "Genos AI for Peak 1031 Exchange",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s8, Inches(1.35), color=GREEN)

builds = [
    ("Automated 45/180-Day Deadline Tracking With Tiered Alert System",
     "Auto-calculates both IRS deadlines from sale close date for every exchange at intake. Sends tiered alerts to team and client at 15, 7, 3, and 1 day out via email and SMS. Escalates to a named senior staff member if a required action is still incomplete within 48 hours of either deadline. Integrates with Outlook and Google Calendar. Zero calendar management required. No active exchange ever has an unmonitored deadline."),
    ("AI Document Collection, Classification, and OCR Extraction",
     "Automated outbound document request with exchange-type-specific checklist sent to client at intake. Incoming documents classified automatically by type - sale contract, title report, purchase agreement, identification form, proof of funds. OCR extracts key data fields (dates, parties, property address, transaction amount). Missing items flagged with a specific list sent to client and team daily until complete. Full document assembly time from 5-10 hours to under 1 hour per exchange."),
    ("Property Identification Verification Bot + IRS-Compliant Letter Generation",
     "Client or title company submits a property address. Bot cross-references parcel data, legal description, and county records via public data APIs. Pre-fills the identification letter with IRS-compliant formatting including exact legal description, address, and parcel number. Flags any discrepancy before the letter is reviewed. Letters generated with correct formatting for all three identification rules (3-property, 200% value, 95% performance). Identification errors eliminated before they reach submission."),
    ("Client Status Portal With AI Chat for Timeline and FAQ Deflection",
     "Investor self-service portal showing exchange timeline, document checklist status, deadline countdown, and funding status. AI chat trained on 1031 exchange rules, Peak's specific process, and IRS FAQ responses. Handles routine investor questions after hours without staff involvement. Routes complex structuring questions and exceptions to the appropriate exchange officer with conversation context. Inbound status calls and emails reduced by 40-50%."),
    ("Compliance Checklist Automation + Form 8824 Auto-Generation",
     "Rules-based system generates the correct compliance checklist for each exchange type at intake - delayed, simultaneous, reverse, improvement. Tracks completion item by item automatically. Flags any gap before the file reaches senior staff review. IRS Form 8824 auto-populates from exchange file data (sale price, basis, replacement property, boot received) and generates a filing-ready PDF for CPA or attorney review. Multi-state variants handled."),
]
for i, (title, desc) in enumerate(builds):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s8, Inches(0.55), t, Inches(0.06), Inches(0.7), GREEN)
    tb(s8, title, Inches(0.75), t, Inches(11.8), Inches(0.38), size=14, bold=True, color=WHITE)
    tb(s8, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# SLIDE 9 - 90-DAY OUTCOMES
s9 = prs.slides.add_slide(BLANK)
bg(s9)
rect(s9, 0, 0, Inches(0.12), H, ACCENT)
tb(s9, "90-DAY OUTCOMES", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s9, "What Changes in 90 Days",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s9, Inches(1.35))

phases = [
    ("Month 1", GOLD, [
        "Deadline tracking live on all active and new exchanges - tiered alerts to team and clients",
        "AI document collection and OCR classification deployed - incoming docs auto-sorted",
        "Missing document alerts running - team notified with specific outstanding items",
        "Client status portal live - investors can check exchange status without calling the office",
    ]),
    ("Month 2", ACCENT, [
        "Property identification verification bot live - parcel cross-reference and letter pre-fill",
        "Compliance checklist automation by exchange type deployed and tracking completion",
        "AI chat handling routine investor questions and IRS FAQ deflection",
        "Staff manual hours per exchange measurably reduced - first throughput increase visible",
    ]),
    ("Month 3", GREEN, [
        "IRS Form 8824 auto-generation from exchange file - CPA review of exceptions only",
        "Full exchange file assembled automatically from intake to close",
        "Manual hours per exchange at or under 2 hours for standard delayed exchanges",
        "Transaction throughput ceiling raised - same team, measurably higher volume capacity",
    ]),
]
pw = Inches(3.9); gap3 = Inches(0.18)
for i, (month, color, pts) in enumerate(phases):
    l3 = Inches(0.55) + i * (pw + gap3)
    rect(s9, l3, Inches(1.8), pw, Inches(0.42), color)
    tb(s9, month.upper(), l3 + Inches(0.1), Inches(1.84), pw - Inches(0.2), Inches(0.36),
       size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    rect(s9, l3, Inches(2.22), pw, Inches(4.7), DARK_BLUE)
    bullets(s9, pts, l3 + Inches(0.15), Inches(2.35), pw - Inches(0.3), Inches(4.4),
            size=12, color=LGREY, dot=color)


# SLIDE 10 - NEXT STEPS
s10 = prs.slides.add_slide(BLANK)
bg(s10)
rect(s10, 0, 0, Inches(0.12), H, GOLD)
rect(s10, Inches(7.5), 0, Inches(5.83), H, DARK_BLUE)
tb(s10, "NEXT STEPS", Inches(0.55), Inches(0.3), Inches(6.5), Inches(0.5),
   size=11, bold=True, color=GOLD)
tb(s10, "Let's Talk for\n15 Minutes.",
   Inches(0.55), Inches(0.8), Inches(6.8), Inches(2.0), size=48, bold=True, color=WHITE)
aline(s10, Inches(2.8), color=GOLD, w=Inches(2.5))
tb(s10, "No pitch deck required. Reply or book a call directly.\nWe'll walk through what fits Peak 1031 Exchange.",
   Inches(0.55), Inches(3.05), Inches(6.5), Inches(0.8), size=14, color=LGREY)

steps = [
    "Reply to this email - we'll set up a time",
    "Book directly: www.genosai.tech/call",
    "Walk through the automation audit in 15 minutes",
    "You decide what, if anything, makes sense to move on",
]
bullets(s10, steps, Inches(0.55), Inches(3.95), Inches(6.5), Inches(2.0), size=14, color=WHITE, dot=GOLD)

rect(s10, 0, Inches(6.5), Inches(7.5), Inches(1.0), DARK_BLUE)
tb(s10, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699  |  www.genosai.tech",
   Inches(0.55), Inches(6.6), Inches(6.8), Inches(0.5), size=12, color=MGREY)

tb(s10, "WHAT WE BUILD FOR\nPEAK 1031", Inches(7.85), Inches(0.6), Inches(5.0), Inches(0.9),
   size=20, bold=True, color=GOLD)
aline(s10, Inches(1.45), color=GOLD)

summary_items = [
    ("Deadlines", "45/180-day auto-tracking, tiered alerts, escalation"),
    ("Documents", "AI collection, OCR classification, missing-item alerts"),
    ("ID Letters", "Parcel verification bot + IRS-compliant letter prefill"),
    ("Client Portal", "Status self-service + AI chat for FAQ deflection"),
    ("Compliance", "Checklist automation by exchange type"),
    ("Form 8824", "Auto-generation from exchange file data"),
]
for i, (label, val) in enumerate(summary_items):
    t2 = Inches(1.6) + i * Inches(0.88)
    rect(s10, Inches(7.85), t2, Inches(1.5), Inches(0.72), NAVY)
    tb(s10, label, Inches(7.9), t2 + Inches(0.05), Inches(1.4), Inches(0.32),
       size=10, bold=True, color=GOLD)
    tb(s10, val, Inches(9.5), t2 + Inches(0.05), Inches(3.6), Inches(0.65),
       size=11, color=LGREY)


prs.save(ROOT / "output" / "decks" / "PeakExchange_Audit_GenosAI.pptx")
print("Saved: PeakExchange_Audit_GenosAI.pptx")
