import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.parent / "utils"))
from ppt_theme import (
    NAVY, DARK_BLUE, ACCENT, GOLD, RED, GREEN, PURPLE, ORANGE, TEAL, CYAN,
    WHITE, LGREY, MGREY, W, H,
    bg, rect, tb, aline, bullets, tag,
    slide_header, slide_footer, score_badge,
    stat_cards, top_bar_card, phase_card, before_after, genos_logo_block,
)

ROOT = Path(__file__).parent.parent.parent
prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]

# SLIDE 1 - COVER
s1 = prs.slides.add_slide(BLANK)
bg(s1)
rect(s1, 0, 0, Inches(0.12), H, GOLD)
rect(s1, Inches(8.5), 0, Inches(4.83), Inches(2.2), DARK_BLUE)
tb(s1, "GENOS AI", Inches(9.2), Inches(0.4), Inches(3.5), Inches(0.8),
   size=28, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
tb(s1, "www.genosai.tech", Inches(9.2), Inches(1.05), Inches(3.5), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
tb(s1, "BRONSHTEYN GMBH", Inches(0.55), Inches(1.6),
   Inches(9), Inches(0.6), size=13, bold=False, color=GOLD)
tb(s1, "Website & AI Growth\nAudit Report", Inches(0.55), Inches(2.05),
   Inches(9), Inches(1.8), size=48, bold=True, color=WHITE)
aline(s1, Inches(3.85), color=GOLD, w=Inches(3))
tb(s1, "Prepared exclusively for Boris Bronshteyn, CEO - Bronshteyn GmbH",
   Inches(0.55), Inches(4.15), Inches(9), Inches(0.5), size=14, color=LGREY)
rect(s1, 0, Inches(6.5), W, Inches(1.0), DARK_BLUE)
tb(s1, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699",
   Inches(0.55), Inches(6.6), Inches(9), Inches(0.5), size=12, color=MGREY)
tb(s1, "May 2026  |  CONFIDENTIAL", Inches(9.5), Inches(6.6), Inches(3.3), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
rect(s1, Inches(10.5), Inches(3.3), Inches(2.3), Inches(2.5), DARK_BLUE)
tb(s1, "AUDIT SCORE", Inches(10.55), Inches(3.4), Inches(2.2), Inches(0.5),
   size=11, bold=True, color=MGREY, align=PP_ALIGN.CENTER)
tb(s1, "3", Inches(10.55), Inches(3.75), Inches(2.2), Inches(1.1),
   size=72, bold=True, color=RED, align=PP_ALIGN.CENTER)
tb(s1, "/ 10", Inches(10.55), Inches(4.7), Inches(2.2), Inches(0.4),
   size=18, color=LGREY, align=PP_ALIGN.CENTER)
tag(s1, "HIGH PRIORITY", Inches(10.75), Inches(5.1), fill=RED, size=10)


# SLIDE 2 - AT A GLANCE
s2 = prs.slides.add_slide(BLANK)
bg(s2)
rect(s2, 0, 0, Inches(0.12), H, GOLD)
tb(s2, "AT A GLANCE", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GOLD)
tb(s2, "Boris Bronshteyn - Who He Is",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s2, Inches(1.35), color=GOLD)

stats = [
    ("Munich", "Bavaria, Germany HQ"),
    ("Bronshteyn\nGmbH", "Munich Family Office"),
    ("11", "Team Members"),
    ("3", "Service Lines: RE / M&A / Advisory"),
    ("Capital\nFactory", "Tech Mentor, Austin TX"),
]
cw = Inches(2.3); ch = Inches(2.2); gap = Inches(0.18)
for i, (num, label) in enumerate(stats):
    l = Inches(0.55) + i * (cw + gap)
    rect(s2, l, Inches(1.8), cw, ch, DARK_BLUE)
    tb(s2, num, l + Inches(0.1), Inches(1.95), cw - Inches(0.2), Inches(1.1),
       size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s2, label, l + Inches(0.1), Inches(3.05), cw - Inches(0.2), Inches(0.7),
       size=12, color=LGREY, align=PP_ALIGN.CENTER)

rect(s2, Inches(0.55), Inches(4.2), Inches(12.23), Inches(0.6), DARK_BLUE)
tb(s2, "Family Office | Investment Real Estate | M&A Advisory | Business Acquisitions  -  Leibnizstrasse 35, 80686 Munich  -  Amtsgericht Munich",
   Inches(0.65), Inches(4.28), Inches(12), Inches(0.45), size=12, color=LGREY, align=PP_ALIGN.CENTER)

tb(s2, "WHAT BORIS MANAGES", Inches(0.55), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=GOLD)
bullets(s2, ["Real estate mediation and investment advisory across Germany",
             "M&A consulting and company acquisition mandates",
             "Family office operations for HNW investor clients"],
        Inches(0.55), Inches(5.35), Inches(5.8), Inches(1.5), size=13, color=LGREY)

tb(s2, "WHO'S ON THE TEAM", Inches(6.8), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=GOLD)
bullets(s2, ["Boris Bronshteyn - CEO, deal lead and investor relationships",
             "Anastasia Gerasimova - Head of Real Estate",
             "Ksenia Kuligina - Key Account Manager"],
        Inches(6.8), Inches(5.35), Inches(5.8), Inches(1.5), size=13, color=LGREY)


# SLIDE 3 - STRENGTHS + WHY THE TIMING IS RIGHT
s3 = prs.slides.add_slide(BLANK)
bg(s3)
rect(s3, 0, 0, Inches(0.12), H, GREEN)
tb(s3, "WHAT'S WORKING", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s3, "Why This Works - and Why Now",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s3, Inches(1.35), color=GREEN)

strengths = [
    ("Bronshteyn GmbH Covers Three Mandate Types That Most Munich Advisories Don't",
     "Real estate mediation, M&A consulting, and business acquisition advisory under one operation means a client relationship that starts as a property deal can evolve into an M&A mandate without switching advisors. That breadth is a genuine differentiator in a Munich market where most operators are either pure-play RE or pure-play M&A - but only if prospects understand what Bronshteyn GmbH actually covers before the first call."),
    ("Capital Factory Mentorship Means Boris Already Evaluates Tech on Its Merits",
     "A Munich family office principal who mentors Austin startup founders thinks differently about technology than most European RE advisors. The Capital Factory connection signals that Boris has already filtered for what works in tech-enabled operations - which makes the AI automation and site rebuild conversation practical rather than philosophical."),
    ("Anastasia and Ksenia Provide Operational Depth Beyond a One-Principal Shop",
     "A Head of Real Estate and a Key Account Manager means the business has the capacity to handle increased deal flow if inbound is improved. The team already exists to close the mandates that better digital infrastructure would generate. This is not a situation where automation would outpace the team's ability to execute - the headcount is there."),
    ("WhiteWill Advisory Relationship Provides a Direct Line Into a AED 2.23B Operation",
     "Boris's advisory role at WhiteWill Dubai - AED 2.23B in 2025 transactions, 500+ employees, 6,000 partner agents - is a bridge relationship that most advisory firms never have. The same AI automation stack that benefits Bronshteyn GmbH scales directly to WhiteWill's 500 weekly inbound inquiries. That conversation becomes available once the Bronshteyn stack has proof of performance."),
    ("German Real Estate and M&A Advisory Is Still Largely Relationship and Network-Driven",
     "Most family offices and investment advisories in the Munich and Bavarian market operate with minimal digital infrastructure. They win on network and do not invest in online presence. A family office that publishes market commentary, ranks for relevant search terms, and presents professionally online has a structural first-mover advantage over operators who do not - and that advantage compounds as content builds."),
]
for i, (title, desc) in enumerate(strengths):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s3, Inches(0.55), t, Inches(0.06), Inches(0.7), GREEN)
    tb(s3, title, Inches(0.75), t, Inches(11.8), Inches(0.38), size=14, bold=True, color=WHITE)
    tb(s3, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# SLIDE 4 - THE THREE GAPS (distinct problem categories, not a repeat of slide 7)
s4 = prs.slides.add_slide(BLANK)
bg(s4)
rect(s4, 0, 0, Inches(0.12), H, RED)
tb(s4, "THE THREE GAPS", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=RED)
tb(s4, "One Placeholder Site. Three Categories of Lost Opportunity.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=28, bold=True, color=WHITE)
aline(s4, Inches(1.35), color=RED)

gap_cards = [
    ("GAP 1\nCREDIBILITY\nINFRASTRUCTURE", RED,
     "What's missing from bronshteyn.com",
     [
         "No service scope - visitors cannot assess mandate fit without calling",
         "No deal track record - zero evidence of outcomes achieved",
         "No team credentials - Anastasia and Ksenia are invisible to prospects",
         "Single-page site fails the pre-call research filter HNW clients apply",
     ]),
    ("GAP 2\nDEAL FLOW\nQUALIFICATION", GOLD,
     "How inbound currently reaches Boris",
     [
         "Email, WhatsApp, and Telegram all deliver unscreened to Boris directly",
         "No deal type, ticket size, or geography screening before his involvement",
         "Boris personally filters every inquiry - a direct cost on mandate execution time",
         "No structured entry point to route acquisition vs. advisory vs. co-investment",
     ]),
    ("GAP 3\nMARKET\nPRESENCE", ACCENT,
     "How Bronshteyn GmbH appears to searching prospects",
     [
         "Single-page site indexes zero content - cannot rank for any keyword",
         "Sellers researching 'family office Munich' or 'M&A Beratung Bayern' never find it",
         "No content or commentary - no reason for Google to surface the domain",
         "First contact for most prospects is still cold - no warm organic channel exists",
     ]),
]
cw3 = Inches(3.95); gap3i = Inches(0.18)
for i, (title, color, subtitle, pts) in enumerate(gap_cards):
    l = Inches(0.55) + i * (cw3 + gap3i)
    rect(s4, l, Inches(1.8), cw3, Inches(0.42), color)
    tb(s4, title, l + Inches(0.12), Inches(1.85), cw3 - Inches(0.24), Inches(0.36),
       size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    rect(s4, l, Inches(2.22), cw3, Inches(4.7), DARK_BLUE)
    tb(s4, subtitle, l + Inches(0.12), Inches(2.3), cw3 - Inches(0.24), Inches(0.35),
       size=10, italic=True, color=color)
    bullets(s4, pts, l + Inches(0.12), Inches(2.72), cw3 - Inches(0.24), Inches(3.9),
            size=12, color=LGREY, dot=color)

rect(s4, Inches(0.55), Inches(7.05), Inches(12.23), Inches(0.3), DARK_BLUE)
tb(s4, "Each gap is a separate problem with a separate fix. All three are solvable.",
   Inches(0.7), Inches(7.08), Inches(12), Inches(0.25), size=11, italic=True, color=GOLD)


# SLIDE 5 - WEBSITE DESIGN AUDIT (before/after visual - what to build, not why it's broken)
s5 = prs.slides.add_slide(BLANK)
bg(s5)
rect(s5, 0, 0, Inches(0.12), H, PURPLE)
tb(s5, "WEBSITE DESIGN AUDIT", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=PURPLE)
tb(s5, "From Contact Card to Family Office Authority Site.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.8), size=32, bold=True, color=WHITE)
aline(s5, Inches(1.45), color=PURPLE, w=Inches(1.5))

rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(3.6), DARK_BLUE)
rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(0.42), RED)
tb(s5, "BRONSHTEYN.COM TODAY", Inches(0.65), Inches(1.92), Inches(3.6), Inches(0.38),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Name + subtitle only: Family Office | Real estate",
    "Email, WhatsApp, Telegram - no other content",
    "Single page, no navigation, no depth",
    "No service descriptions or mandate scope",
    "No deal track record or transaction history",
    "No team page - Anastasia and Ksenia invisible",
    "No blog or content hub",
    "No SEO - single page cannot rank for any term",
], Inches(0.7), Inches(2.42), Inches(3.6), Inches(2.9), size=11, color=LGREY, dot=RED)

tb(s5, "->", Inches(4.45), Inches(3.6), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(3.6), DARK_BLUE)
rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(0.42), GOLD)
tb(s5, "WHAT A FAMILY OFFICE SITE DOES", Inches(5.1), Inches(1.92), Inches(3.4), Inches(0.38),
   size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Services scoped with mandate types and deal criteria",
    "Track record: deal count, asset class, geography covered",
    "Team page: Boris, Anastasia, Ksenia with credentials and roles",
    "Case studies or transaction summaries where disclosable",
    "Content hub: market commentary and deal flow insights",
    "SEO infrastructure for target keywords in German",
    "Structured inquiry form routing by deal type",
    "Design tone matching HNW investor expectations",
], Inches(5.1), Inches(2.42), Inches(3.4), Inches(2.9), size=11, color=LGREY, dot=GOLD)

tb(s5, "->", Inches(8.68), Inches(3.6), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(3.6), DARK_BLUE)
rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(0.42), GREEN)
tb(s5, "WHAT GENOS AI BUILDS", Inches(9.35), Inches(1.92), Inches(3.4), Inches(0.38),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Full rebuild: services, deal criteria, mandate types scoped",
    "Track record section where disclosable",
    "Team page with roles and credentials for all three",
    "AI chat: deal qualifier before Boris's direct involvement",
    "Content layer with Munich RE and M&A commentary",
    "SEO infrastructure: family office + M&A advisory Germany",
    "Structured inquiry routing by deal type",
    "Investor-grade design built for HNW audience",
], Inches(9.35), Inches(2.42), Inches(3.4), Inches(2.9), size=11, color=LGREY, dot=GREEN)

rect(s5, Inches(0.55), Inches(5.72), Inches(12.23), Inches(1.28), DARK_BLUE)
tb(s5, "DESIGN BENCHMARK:", Inches(0.7), Inches(5.78), Inches(2.5), Inches(0.38),
   size=11, bold=True, color=PURPLE)
stats2 = [
    ("0 pages", "currently indexable\nfor any keyword"),
    ("3 channels", "unstructured inbound\nall reaching Boris"),
    ("0 team", "credentials visible\nto prospects today"),
    ("1 site rebuild", "clears all three\ncredibility gaps"),
]
sw = Inches(2.7); sl = Inches(3.2)
for i, (num, label) in enumerate(stats2):
    xl = sl + i * sw
    tb(s5, num, xl, Inches(5.75), Inches(2.5), Inches(0.55),
       size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s5, label, xl, Inches(6.22), Inches(2.5), Inches(0.7),
       size=10, color=LGREY, align=PP_ALIGN.CENTER)


# SLIDE 6 - AI AUTOMATION STACK (automation only - site rebuild is slide 5)
s6 = prs.slides.add_slide(BLANK)
bg(s6)
rect(s6, 0, 0, Inches(0.12), H, ACCENT)
tb(s6, "AI AUTOMATION STACK", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s6, "Four Layers That Run the Business Side Boris Doesn't Have Time to Run Manually.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=26, bold=True, color=WHITE)
aline(s6, Inches(1.35), color=ACCENT)

opps = [
    ("AI Deal Flow\nInquiry Qualifier", GOLD,
     "Trained on Bronshteyn GmbH service scope, deal criteria, ticket size ranges, geographic focus, and mandate types. A seller or investor who messages at 11pm gets an immediate, structured response that screens deal fit before Boris sees the inquiry. Qualified inbound arrives pre-categorised by mandate type and deal size. Unqualified volume is filtered automatically."),
    ("Structured Inquiry\nRouting by Mandate", ACCENT,
     "One entry point replacing three undifferentiated channels. Routes by deal type: direct acquisition, co-investment, advisory mandate, M&A facilitation. Each route triggers a different follow-up sequence matched to that mandate. The seller exploring an exit gets different content than the investor evaluating a co-investment. Relevance from the first response."),
    ("Investor and Seller\nFollow-Up Sequences", GREEN,
     "HNW relationships and M&A mandates close across multiple touchpoints, not one. A first-contact prospect who doesn't immediately proceed enters a 60-90 day automated sequence: relevant market commentary, deal flow insights, soft re-engagement CTAs. Boris's direct involvement is reserved for mandate-ready conversations. Relationships stay warm without manual tracking."),
    ("Contact Intelligence\nand Pipeline Visibility", PURPLE,
     "A CRM layer that tracks where every prospect is in the relationship: first contact, engaged, follow-up sequence active, deal criteria confirmed, mandate-ready. Boris sees at a glance which relationships are warm, which need a personal touch, and which are in automated nurture. No mandate opportunity falls through the gap between manual notes and inbox management."),
]
cw2 = Inches(2.95); gap2 = Inches(0.16)
for i, (title, color, desc) in enumerate(opps):
    l2 = Inches(0.55) + i * (cw2 + gap2)
    rect(s6, l2, Inches(1.8), cw2, Inches(0.06), color)
    rect(s6, l2, Inches(1.86), cw2, Inches(4.7), DARK_BLUE)
    tb(s6, title, l2 + Inches(0.12), Inches(1.98), cw2 - Inches(0.24), Inches(0.65),
       size=13, bold=True, color=color)
    tb(s6, desc, l2 + Inches(0.12), Inches(2.65), cw2 - Inches(0.24), Inches(3.65),
       size=11, color=LGREY)

rect(s6, Inches(0.55), Inches(6.7), Inches(12.23), Inches(0.45), DARK_BLUE)
tb(s6, "A Munich family office that closes on credibility. The automation runs the inbound layer so Boris's time goes to deals that are already qualified.",
   Inches(0.7), Inches(6.75), Inches(12), Inches(0.35), size=11, italic=True, color=GOLD)


# SLIDE 7 - THE WHITEWILL ANGLE (unique to Boris - not present in any other deck)
s7 = prs.slides.add_slide(BLANK)
bg(s7)
rect(s7, 0, 0, Inches(0.12), H, ORANGE)
tb(s7, "THE WHITEWILL ANGLE", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ORANGE)
tb(s7, "One Advisory Relationship. A Second Conversation Worth Having.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=30, bold=True, color=WHITE)
aline(s7, Inches(1.35), color=ORANGE, w=Inches(1.8))

# WhiteWill stats bar
rect(s7, Inches(0.55), Inches(1.7), Inches(12.23), Inches(1.55), DARK_BLUE)
ww_stats = [
    ("AED 2.23B", "2025 transaction volume"),
    ("500+", "employees"),
    ("6,000", "partner agents"),
    ("500/week", "inbound inquiries"),
    ("4 offices", "Dubai / London / Abu Dhabi / Miami"),
]
ssw = Inches(2.35)
for i, (num, label) in enumerate(ww_stats):
    xl = Inches(0.7) + i * ssw
    tb(s7, num, xl, Inches(1.78), Inches(2.2), Inches(0.65),
       size=22, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    tb(s7, label, xl, Inches(2.38), Inches(2.2), Inches(0.5),
       size=10, color=LGREY, align=PP_ALIGN.CENTER)

# Two-column explanation
rect(s7, Inches(0.55), Inches(3.45), Inches(5.9), Inches(3.2), DARK_BLUE)
tb(s7, "WHY THE SAME STACK APPLIES", Inches(0.7), Inches(3.52), Inches(5.6), Inches(0.35),
   size=11, bold=True, color=ORANGE)
bullets(s7, [
    "500 weekly inbound inquiries from agents and investors cannot be handled manually - AI qualification is not optional at that volume",
    "6,000 partner agents need automated communication infrastructure to stay aligned on listing inventory, deal status, and priority leads",
    "Multi-city operations (Dubai, London, Abu Dhabi, Miami) require centralized inquiry routing with geography-aware follow-up sequences",
    "Investor and HNW buyer relationships at AED 2.23B deal volume close on the same multiple-touchpoint follow-up automation that applies to Bronshteyn GmbH",
], Inches(0.7), Inches(3.9), Inches(5.6), Inches(2.5), size=12, color=LGREY, dot=ORANGE)

rect(s7, Inches(6.65), Inches(3.45), Inches(6.13), Inches(3.2), DARK_BLUE)
tb(s7, "HOW THE CONVERSATION FLOWS", Inches(6.8), Inches(3.52), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=GOLD)
bullets(s7, [
    "Genos AI builds and deploys the automation stack for Bronshteyn GmbH first - Boris sees the performance directly",
    "Once the Bronshteyn stack has 60-90 days of live results, Boris has a concrete proof point to take to Oleg Torbosov at WhiteWill",
    "Boris's advisory relationship is the introduction - the performance data is the pitch",
    "If WhiteWill engages, Boris is the bridge that brought a working AI infrastructure to a AED 2B+ operation - that has its own value in the advisory relationship",
], Inches(6.8), Inches(3.9), Inches(5.8), Inches(2.5), size=12, color=LGREY, dot=GOLD)

rect(s7, Inches(0.55), Inches(6.8), Inches(12.23), Inches(0.38), DARK_BLUE)
tb(s7, "Bronshteyn GmbH is the proof of concept. WhiteWill is the scale opportunity.",
   Inches(0.7), Inches(6.85), Inches(12), Inches(0.28), size=11, italic=True, color=ORANGE)


# SLIDE 8 - PHASED IMPLEMENTATION PLAN (what gets built and when - not a re-list of features)
s8 = prs.slides.add_slide(BLANK)
bg(s8)
rect(s8, 0, 0, Inches(0.12), H, GREEN)
tb(s8, "IMPLEMENTATION PLAN", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s8, "What Gets Built - and When.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s8, Inches(1.35), color=GREEN)

phases = [
    ("PHASE 1", "Weeks 1-3", GOLD, [
        ("Site Architecture and Design", "Structure, service scope, team page, deal criteria, investor-grade design - all approved before dev begins. No surprises on delivery."),
        ("AI Chat Trained and Deployed", "Chat trained on Bronshteyn GmbH service scope, deal criteria, ticket size, and mandate types. Live on the rebuilt site from day one."),
        ("Structured Inquiry Form Live", "One entry point replacing email/WhatsApp/Telegram. Routes by deal type. Each route triggers the correct follow-up path."),
    ]),
    ("PHASE 2", "Weeks 4-6", ACCENT, [
        ("Site Live on bronshteyn.com", "Full rebuild deployed. Services, team, track record, AI chat, inquiry routing. SEO infrastructure in place from day one of launch."),
        ("Follow-Up Sequences Active", "Investor and seller sequences running for all first-contact prospects in the pipeline, including anyone Boris already has in progress."),
        ("CRM Pipeline View Live", "Contact intelligence layer active. Boris sees every prospect's relationship stage at a glance without manual tracking."),
    ]),
    ("PHASE 3", "Weeks 7-12", GREEN, [
        ("First Content Pieces Published", "Munich real estate market commentary and M&A advisory perspective pieces live. Content calendar set for quarterly publishing cadence."),
        ("SEO Indexing in Progress", "Target keywords for family office Munich, investment real estate Germany, M&A advisory Bavaria indexed and tracked for ranking progress."),
        ("WhiteWill Proof Point Ready", "60-90 days of live performance data from the Bronshteyn stack. Inquiry volume, pre-qualification rate, follow-up conversion - ready to present."),
    ]),
]
pw = Inches(3.9); gap4 = Inches(0.18)
for i, (phase, timing, color, items) in enumerate(phases):
    l3 = Inches(0.55) + i * (pw + gap4)
    rect(s8, l3, Inches(1.8), pw, Inches(0.52), color)
    tb(s8, phase, l3 + Inches(0.12), Inches(1.84), pw * 0.5 - Inches(0.12), Inches(0.42),
       size=13, bold=True, color=NAVY)
    tb(s8, timing, l3 + pw * 0.5, Inches(1.88), pw * 0.5 - Inches(0.15), Inches(0.38),
       size=11, color=NAVY, align=PP_ALIGN.RIGHT)
    rect(s8, l3, Inches(2.32), pw, Inches(4.6), DARK_BLUE)
    y_item = Inches(2.42)
    for title, desc in items:
        rect(s8, l3 + Inches(0.12), y_item + Inches(0.06), Inches(0.05), Inches(0.52), color)
        tb(s8, title, l3 + Inches(0.25), y_item, pw - Inches(0.38), Inches(0.32),
           size=12, bold=True, color=color)
        tb(s8, desc, l3 + Inches(0.25), y_item + Inches(0.3), pw - Inches(0.38), Inches(0.45),
           size=11, color=LGREY)
        y_item += Inches(0.95)


# SLIDE 9 - BEFORE vs AFTER - MEASURABLE OUTCOMES (numbers, not deliverables)
s9 = prs.slides.add_slide(BLANK)
bg(s9)
rect(s9, 0, 0, Inches(0.12), H, ACCENT)
tb(s9, "BEFORE vs AFTER", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s9, "What Actually Changes at Bronshteyn GmbH in 90 Days",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=30, bold=True, color=WHITE)
aline(s9, Inches(1.35))

comparisons = [
    ("Inbound Qualification",
     "Boris screens every inquiry personally across 3 channels. Time lost to unqualified volume.",
     "AI qualifier screens all first-contact inbound. Boris's inbox contains only pre-categorised, mandate-relevant inquiries."),
    ("Prospect Research",
     "HNW investors and M&A targets search for Boris Bronshteyn and find a name, email, and Telegram handle.",
     "They land on a family office authority site with service scope, team credentials, and a case for why Bronshteyn GmbH is the right call."),
    ("Follow-Up and Pipeline",
     "First contacts that don't immediately proceed are manually tracked - or quietly lost to inactivity.",
     "60-90 day automated sequences keep every prospect warm. Boris re-engages when they're mandate-ready, not when he remembers to follow up."),
    ("Organic Discovery",
     "bronshteyn.com cannot appear in any organic search. Sellers and investors searching the market before reaching out never find it.",
     "Content layer and SEO infrastructure build search presence for target keywords over 90 days. Warm inbound begins to compound."),
    ("WhiteWill Conversation",
     "No tech proof point exists to bring to Oleg Torbosov at WhiteWill. Advisory relationship is separate from business development.",
     "90 days of live performance data from the Bronshteyn stack. A working AI infrastructure to introduce to a AED 2B+ operation."),
]

lw = Inches(5.4)
rect(s9, Inches(0.55), Inches(1.65), lw, Inches(0.35), RED)
tb(s9, "TODAY", Inches(0.65), Inches(1.68), lw - Inches(0.2), Inches(0.28),
   size=11, bold=True, color=WHITE)
rect(s9, Inches(6.45), Inches(1.65), lw + Inches(0.43), Inches(0.35), GREEN)
tb(s9, "AFTER 90 DAYS", Inches(6.55), Inches(1.68), lw, Inches(0.28),
   size=11, bold=True, color=WHITE)

for i, (topic, before, after) in enumerate(comparisons):
    t = Inches(2.1) + i * Inches(1.02)
    rect(s9, Inches(0.55), t, Inches(12.23), Inches(0.22), DARK_BLUE)
    tb(s9, topic.upper(), Inches(0.65), t + Inches(0.02), Inches(12), Inches(0.2),
       size=9, bold=True, color=MGREY)
    tb(s9, before, Inches(0.65), t + Inches(0.24), lw - Inches(0.2), Inches(0.66),
       size=11, color=LGREY)
    tb(s9, after, Inches(6.55), t + Inches(0.24), lw, Inches(0.66),
       size=11, color=WHITE)


# SLIDE 10 - NEXT STEPS
s10 = prs.slides.add_slide(BLANK)
bg(s10)
rect(s10, 0, 0, Inches(0.12), H, GOLD)
rect(s10, Inches(7.5), 0, Inches(5.83), H, DARK_BLUE)
tb(s10, "NEXT STEPS", Inches(0.55), Inches(0.3), Inches(6.5), Inches(0.5),
   size=11, bold=True, color=GOLD)
tb(s10, "Let's Talk for\n15 Minutes.",
   Inches(0.55), Inches(0.8), Inches(6.8), Inches(2.0), size=48, bold=True, color=WHITE)
aline(s10, Inches(2.8), color=GOLD, w=Inches(2.5))
tb(s10, "No pitch deck required. Reply or book a call directly.\nWe'll walk through what fits Bronshteyn GmbH first.",
   Inches(0.55), Inches(3.05), Inches(6.5), Inches(0.8), size=14, color=LGREY)

steps = [
    "Reply to this message - we'll set up a time",
    "Walk through the audit findings in 15 minutes",
    "Agree on which gap to close first and in what order",
    "You decide what, if anything, makes sense to move on",
]
bullets(s10, steps, Inches(0.55), Inches(3.95), Inches(6.5), Inches(2.0), size=14, color=WHITE, dot=GOLD)

rect(s10, 0, Inches(6.5), Inches(7.5), Inches(1.0), DARK_BLUE)
tb(s10, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699  |  www.genosai.tech",
   Inches(0.55), Inches(6.6), Inches(6.8), Inches(0.5), size=12, color=MGREY)

tb(s10, "THREE GAPS.\nTHREE FIXES.", Inches(7.85), Inches(0.6), Inches(5.0), Inches(1.1),
   size=22, bold=True, color=GOLD)
aline(s10, Inches(1.65), color=GOLD)

summary_items = [
    ("Credibility Gap", "Full site rebuild: services, team, track record, investor-grade design"),
    ("Qualification Gap", "AI chat + structured inquiry routing by mandate type"),
    ("Pipeline Gap", "Investor and seller follow-up sequences across 60-90 day window"),
    ("Presence Gap", "Content layer + SEO for family office Munich and M&A advisory keywords"),
    ("Contact Intelligence", "CRM pipeline view - relationship stage visible at a glance"),
    ("WhiteWill Bridge", "90-day proof point ready to take to Oleg Torbosov"),
]
for i, (label, val) in enumerate(summary_items):
    t2 = Inches(1.85) + i * Inches(0.88)
    rect(s10, Inches(7.85), t2, Inches(1.5), Inches(0.72), NAVY)
    tb(s10, label, Inches(7.9), t2 + Inches(0.05), Inches(1.4), Inches(0.32),
       size=10, bold=True, color=GOLD)
    tb(s10, val, Inches(9.5), t2 + Inches(0.05), Inches(3.6), Inches(0.65),
       size=11, color=LGREY)


prs.save(ROOT / "output" / "decks" / "Bronshteyn_Audit_GenosAI.pptx")
print("Saved: Bronshteyn_Audit_GenosAI.pptx")
