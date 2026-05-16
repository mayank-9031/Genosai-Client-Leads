import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.parent / "utils"))
from ppt_theme import (
    NAVY, DARK_BLUE, ACCENT, GOLD, RED, GREEN, PURPLE, ORANGE, TEAL, CYAN,
    WHITE, LGREY, MGREY, W, H,
    bg, rect, tb, aline, bullets, tag,
    slide_header, slide_footer, score_badge,
    stat_cards, top_bar_card, phase_card, before_after, genos_logo_block,
)

ROOT = Path(__file__).parent.parent.parent
prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]

# Client-specific colors
AMBER     = RGBColor(0xFF, 0xA0, 0x00)

# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
bg(s1)
rect(s1, 0, 0, Inches(0.12), H, RED)
rect(s1, Inches(8.5), 0, Inches(4.83), Inches(2.2), DARK_BLUE)
tb(s1, "GENOS AI", Inches(9.2), Inches(0.4), Inches(3.5), Inches(0.8),
   size=28, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
tb(s1, "www.genosai.tech", Inches(9.2), Inches(1.05), Inches(3.5), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)

tb(s1, "RAJENDRA UNIVERSITY, BALANGIR  |  ODISHA", Inches(0.55), Inches(1.38),
   Inches(9), Inches(0.6), size=13, bold=False, color=RED)
tb(s1, "Website Audit Report\n— An Honest Assessment", Inches(0.55), Inches(1.88),
   Inches(9), Inches(1.8), size=40, bold=True, color=WHITE)
aline(s1, Inches(3.62), color=RED, w=Inches(3))
tb(s1, "Prepared exclusively for the Administration of Rajendra University",
   Inches(0.55), Inches(3.92), Inches(9), Inches(0.5), size=14, color=LGREY)

rect(s1, 0, Inches(6.5), W, Inches(1.0), DARK_BLUE)
tb(s1, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699",
   Inches(0.55), Inches(6.6), Inches(9), Inches(0.5), size=12, color=MGREY)
tb(s1, "May 2026  |  CONFIDENTIAL", Inches(9.5), Inches(6.6), Inches(3.3), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)

rect(s1, Inches(10.5), Inches(3.3), Inches(2.3), Inches(2.5), DARK_BLUE)
tb(s1, "AUDIT SCORE", Inches(10.55), Inches(3.4), Inches(2.2), Inches(0.5),
   size=11, bold=True, color=MGREY, align=PP_ALIGN.CENTER)
tb(s1, "2", Inches(10.55), Inches(3.75), Inches(2.2), Inches(1.1),
   size=72, bold=True, color=RED, align=PP_ALIGN.CENTER)
tb(s1, "/ 10", Inches(10.55), Inches(4.7), Inches(2.2), Inches(0.4),
   size=18, color=LGREY, align=PP_ALIGN.CENTER)
tag(s1, "CRITICAL — REBUILD NEEDED", Inches(10.2), Inches(5.1), fill=RED, size=9)


# ══════════════════════════════════════════════════════════
# SLIDE 2 — UNIVERSITY SNAPSHOT + THE REAL PROBLEM
# ══════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
bg(s2)
rect(s2, 0, 0, Inches(0.12), H, ACCENT)
tb(s2, "UNIVERSITY SNAPSHOT", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s2, "Rajendra University — Who They Are",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s2, Inches(1.35))

stats = [
    ("1944", "Established\n(Origin as\nIntermediate College)"),
    ("2015", "Upgraded to\nUniversity\nStatus"),
    ("UGC", "Recognised\nState Public\nUniversity"),
    ("Balangir", "Western\nOdisha —\nRegional Hub"),
    ("NAAC", "Accredited\n(Quality\nRecognition)"),
]
cw = Inches(2.3); ch = Inches(2.2); gap = Inches(0.18)
for i, (num, label) in enumerate(stats):
    l = Inches(0.55) + i * (cw + gap)
    rect(s2, l, Inches(1.8), cw, ch, DARK_BLUE)
    tb(s2, num, l + Inches(0.1), Inches(1.95), cw - Inches(0.2), Inches(0.85),
       size=26, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s2, label, l + Inches(0.1), Inches(2.82), cw - Inches(0.2), Inches(0.9),
       size=10.5, color=LGREY, align=PP_ALIGN.CENTER)

rect(s2, Inches(0.55), Inches(4.2), Inches(12.23), Inches(0.6), DARK_BLUE)
tb(s2, "Departments include Commerce, Political Science, Physics, Education, Statistics, Computer Science, Journalism & Mass Communication, and more",
   Inches(0.65), Inches(4.28), Inches(12), Inches(0.45), size=11, color=LGREY, align=PP_ALIGN.CENTER)

# The fatal problem box
rect(s2, Inches(0.55), Inches(5.0), Inches(12.23), Inches(1.72), RGBColor(0x2A, 0x05, 0x05))
tb(s2, "THE FUNDAMENTAL PROBLEM:", Inches(0.7), Inches(5.08), Inches(4), Inches(0.38),
   size=12, bold=True, color=RED)
tb(s2,
   "The site being audited (scius.tech/link/rajendrauniversity) is a DEMO template built by Scius Solution — not the real university "
   "website (rajendrauniversity.ac.in). It is hosted on a third-party subdomain. It contains placeholder phone numbers, fake alumni who are "
   "Scius Solution employees, dummy research scholars, broken navigation on every page, and blank sub-pages. A UGC-recognised state public "
   "university with a 80-year history does not have a functioning digital presence. That is what this audit addresses.",
   Inches(0.7), Inches(5.48), Inches(11.9), Inches(1.12),
   size=11.5, color=WHITE, wrap=True)


# ══════════════════════════════════════════════════════════
# SLIDE 3 — THE MOST EMBARRASSING ISSUES (CRITICAL ERRORS)
# ══════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
bg(s3)
rect(s3, 0, 0, Inches(0.12), H, RED)
tb(s3, "CRITICAL ERRORS — FOUND IN AUDIT", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=RED)
tb(s3, "Issues That Cannot Be Left As-Is",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s3, Inches(1.35), color=RED)

issues = [
    ("Header Shows an OMAN Phone Number", "CRITICAL",
     'The very first line on the site reads "Have any question? +968547856254". +968 is Oman\'s country code. This is a placeholder that was never replaced. A prospective student, parent, or faculty member calling this number reaches Oman. This is the most visible element of the entire website and it is factually wrong.'),
    ("'Notable Alumni' Are Scius Solution Employees", "CRITICAL",
     "The Notable Alumni section showcases two people: Kshirod Biswal (Developer at Scius Solution, 2018) and Abinash Choudhury (Frontend Developer at Scius Solution, 2020). These are employees of the website vendor — not university alumni. Placeholder data was shipped to production and presented as real graduates of the institution."),
    ("Every Navigation Link Shows the Homepage", "CRITICAL",
     "Clicking Administration, Academics, Admission & Fee, Examination, Faculty, Departments, Placements, International — every single one loads identical homepage content (TEXT_LEN: 6075 across all). Not one sub-section works. The navigation is decorative. There is effectively one page on the entire site."),
    ("9 Pages Render Completely Blank", "HIGH",
     "The following routes return near-zero content: /research (241 chars), /campus (146), /alumni (146), /news (146), /gallery (146), /hostel (146), /contact (611 — just footer). These routes exist in the navigation but deliver nothing. Visitors who click them see an empty shell."),
    ("Site Hosted on a Vendor Subdomain — Not the University Domain", "HIGH",
     "The URL is scius.tech/link/rajendrauniversity — not rajendrauniversity.ac.in. Hosting the official digital presence of a state public university on a third-party vendor's subdomain means the university has zero ownership, zero SEO control, and loses all content if the vendor relationship ends."),
]
for i, (title, pri, desc) in enumerate(issues):
    t = Inches(1.62) + i * Inches(1.02)
    pc = RED if pri == "CRITICAL" else AMBER
    rect(s3, Inches(0.55), t, Inches(0.06), Inches(0.72), pc)
    tb(s3, title, Inches(0.75), t, Inches(9.5), Inches(0.35), size=13, bold=True, color=WHITE)
    tag(s3, pri, Inches(10.75), t + Inches(0.04), fill=pc, size=9)
    tb(s3, desc, Inches(0.75), t + Inches(0.36), Inches(11.8), Inches(0.55), size=11, color=LGREY)


# ══════════════════════════════════════════════════════════
# SLIDE 4 — DESIGN HEAD-TO-HEAD vs IIT KANPUR
# ══════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
bg(s4)
rect(s4, 0, 0, Inches(0.12), H, PURPLE)
tb(s4, "DESIGN HEAD-TO-HEAD", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=PURPLE)
tb(s4, "Rajendra University vs IIT Kanpur — Element by Element",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=28, bold=True, color=WHITE)
aline(s4, Inches(1.35), color=PURPLE, w=Inches(1.5))

# Column headers
rect(s4, Inches(0.55), Inches(1.45), Inches(3.0), Inches(0.38), DARK_BLUE)
tb(s4, "DESIGN ELEMENT", Inches(0.55), Inches(1.47), Inches(3.0), Inches(0.34),
   size=10, bold=True, color=MGREY, align=PP_ALIGN.CENTER)
rect(s4, Inches(3.6), Inches(1.45), Inches(4.85), Inches(0.38), RED)
tb(s4, "RAJENDRA UNIVERSITY (Today)", Inches(3.6), Inches(1.47), Inches(4.85), Inches(0.34),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rect(s4, Inches(8.5), Inches(1.45), Inches(4.28), Inches(0.38), GREEN)
tb(s4, "IIT KANPUR (Benchmark)", Inches(8.5), Inches(1.47), Inches(4.28), Inches(0.34),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rows = [
    ("Header",
     "Wrong Oman phone (+968 547856254). Cosmetic +/- accessibility buttons. Broken Odia toggle.",
     "Persistent search bar. Working English/हिन्दी toggle. 'Plan Your Visit' link prominent."),
    ("Hero Section",
     "Generic welcome text. Truncated Chancellor/VC messages ending '..more' (link broken).",
     "Welcome statement + 5 institutional stats: 9500+ Students, 570 Faculty, ₹1540+ Cr funding."),
    ("Institutional Stats",
     "ZERO metrics visible anywhere. No student count, no funding figures, no ranking.",
     "Dedicated 'Key Statistics': 2000+ publications/yr, 1000+ patents filed, 150+ companies incubated."),
    ("Research Display",
     "Names + department only: 'Dr. Kunal Mishra — Commerce'. No titles, no journals.",
     "Real project showcase: 'Svan M2 Quadruped Robot' (Mech), 'VTOL UAVs' (Aero) with Learn More."),
    ("News Pattern",
     "Scrolling marquee ticker (2005-era). Carousel duplicates same 6 items including 2021 events.",
     "Static card list, latest-first. Real headlines: 672 job offers, ₹100 Cr alumni pledge, MoUs."),
    ("Navigation",
     "9 top-level menus — ALL load homepage. Navigation is decorative, not functional.",
     "8 top-level menus — every single one leads to a fully built-out section."),
    ("Department Discovery",
     "No department pages exist. Cannot find a single faculty profile or course syllabus.",
     "7 visual 'Explore' cards: Engineering, Humanities, Economics, Sciences, Management, Schools."),
    ("Admission Section",
     "Buried inside cluttered notice board. No dates calendar, no application process.",
     "Active admission notices: 'PhD/MTech/MDes/MS 2026-27', GATE-JAM, Online PGP — each clickable."),
    ("Visual Hierarchy",
     "Marquee + ticker + carousel all competing. No whitespace. Everything on one page.",
     "Clear section breaks. Generous whitespace. One focus per scroll-view. H1→H2→body rhythm."),
    ("Domain Ownership",
     "scius.tech/link/rajendrauniversity — hosted on the website vendor's subdomain.",
     "iitk.ac.in — institutional ownership, .ac.in trust, full SEO control."),
]

row_h = Inches(0.5)
for i, (elem, ru, iitk) in enumerate(rows):
    t = Inches(1.83) + i * row_h
    row_color = DARK_BLUE if i % 2 == 0 else RGBColor(0x16, 0x21, 0x55)
    rect(s4, Inches(0.55), t, Inches(12.23), row_h, row_color)
    tb(s4, elem, Inches(0.62), t + Inches(0.08), Inches(2.9), Inches(0.35),
       size=10.5, bold=True, color=PURPLE)
    tb(s4, ru, Inches(3.62), t + Inches(0.05), Inches(4.8), Inches(0.42),
       size=9.5, color=LGREY, wrap=True)
    tb(s4, iitk, Inches(8.52), t + Inches(0.05), Inches(4.22), Inches(0.42),
       size=9.5, color=LGREY, wrap=True)


# ══════════════════════════════════════════════════════════
# SLIDE 5 — VISUAL & UX DESIGN BREAKDOWN
# ══════════════════════════════════════════════════════════
s_design = prs.slides.add_slide(BLANK)
bg(s_design)
rect(s_design, 0, 0, Inches(0.12), H, PURPLE)
tb(s_design, "VISUAL & UX DESIGN BREAKDOWN", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=PURPLE)
tb(s_design, "Specific Design Failures, Element by Element",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=30, bold=True, color=WHITE)
aline(s_design, Inches(1.35), color=PURPLE)

flaws = [
    ("Marquee Ticker for 'Latest News'",
     "The homepage uses a horizontally scrolling marquee text ticker for news headlines — a pattern abandoned by serious websites around 2008. It is hard to read, accessibility-hostile, cannot be paused, and competes with the static notice board directly below it. IIT Kanpur uses a clean static card list — each headline gets its own row with proper typography."),
    ("Three Carousels Stacked Above the Fold",
     "Rajendra University stacks: (1) news ticker, (2) events carousel, (3) research scholars carousel — all auto-rotating on the same screen. The eye cannot settle. IIT Kanpur uses zero auto-rotating carousels — every piece of content is statically laid out and discoverable on scroll."),
    ("No Whitespace — Everything Is Squeezed",
     "Sections butt against each other with no breathing room. Cards have minimal padding. Headings sit on top of body text. The result is a wall of information with no visual hierarchy. IIT Kanpur uses generous whitespace between sections — each block gets room to be read."),
    ("Identical Card Treatment Everywhere",
     "Every card on the site looks the same: same border, same fill, same size, same shadow. There is no visual distinction between a Notice, an Event, a Research highlight, or an Alumnus. IIT Kanpur uses themed cards — research has imagery, news has dates, departments have icons."),
    ("Broken Typography Rhythm",
     "Headings, sub-headings, and body text are all visually similar weight and size. There is no clear H1 → H2 → H3 → body scale. Visitors cannot scan the page — they have to read line by line. IIT Kanpur uses a clear 4-level type scale with consistent weight and spacing."),
    ("No Photography — Just Stock-Feel Graphics",
     "There is no actual campus imagery on the homepage — no aerial shots, no building photographs, no faculty portraits, no student life. The 'Gallery' page returns 146 characters. IIT Kanpur leads with campus and research photography in every major section."),
    ("Accessibility Theatre",
     "The site shows a 'Skip To Content | Screen Reader Access +  -  Switch to ଓଡିଆ' bar at the top — but the +/- buttons do not work, screen reader markup is incomplete, and the Odia toggle is broken. This is performative accessibility, not real WCAG compliance."),
    ("Footer Is a Link Dump",
     "The footer contains 14+ external links (UGC, NAAC, RUSA, AISHE, etc.) but the internal navigation it should provide — site map, key pages, search — is missing. IIT Kanpur's footer is an organisational map: Board of Governors, Directorate, Statutes, Acts, Ordinances, Annual Reports, RTI — all linked."),
]

for i, (title, desc) in enumerate(flaws):
    col = i % 2; row = i // 2
    lx = Inches(0.55) + col * Inches(6.4)
    ty = Inches(1.55) + row * Inches(1.42)
    rect(s_design, lx, ty, Inches(6.1), Inches(1.3), DARK_BLUE)
    rect(s_design, lx, ty, Inches(0.06), Inches(1.3), PURPLE)
    tb(s_design, title, lx + Inches(0.16), ty + Inches(0.08), Inches(5.85), Inches(0.4),
       size=12, bold=True, color=PURPLE)
    tb(s_design, desc, lx + Inches(0.16), ty + Inches(0.48), Inches(5.85), Inches(0.8),
       size=9.5, color=LGREY, wrap=True)


# ══════════════════════════════════════════════════════════
# SLIDE 5 — CONTENT GAPS — PAGE BY PAGE
# ══════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
bg(s5)
rect(s5, 0, 0, Inches(0.12), H, AMBER)
tb(s5, "CONTENT GAP ANALYSIS", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=AMBER)
tb(s5, "Every Major Section — What Exists vs. What Must Exist",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=28, bold=True, color=WHITE)
aline(s5, Inches(1.35), color=AMBER)

sections = [
    ("HOMEPAGE", AMBER,
     "What exists: Welcome text, Chancellor/VC messages (cut off with '..more'), Notice Board, Events carousel (includes 2021 events), Research names-only list, 2 alumni cards (vendor employees).",
     "What's missing: Real photography, program highlights, NAAC score, NIRF ranking, student count, placement stats, campus life imagery, social media feeds, course search."),
    ("ACADEMICS / DEPARTMENTS", RED,
     "What exists: Nothing. Clicking Academics in the nav loads the homepage. There are no department pages, no course listings, no faculty pages, no syllabus downloads, no program-specific information anywhere on the site.",
     "What's missing: Every department needs its own page — faculty profiles with qualifications, courses offered with syllabi, research areas, labs and facilities, contact details."),
    ("ADMISSION & FEE", RED,
     "What exists: Nothing. The Admission & Fee nav item loads the homepage. No admission procedure, no eligibility criteria, no fee structure, no application form, no important dates, no scholarship information is accessible.",
     "What's missing: Online application portal, eligibility criteria per program, fee structure table, important dates calendar, documents required, reservation policy."),
    ("EXAMINATION", RED,
     "What exists: The Examination nav loads the homepage. However, the Notice Board on the homepage DOES show exam notices — but these are buried in a scrolling list with no search, no filter, no downloadable PDFs visible.",
     "What's missing: Dedicated exam portal with date-wise schedule, hall ticket download, result check, re-evaluation application, back paper form, examination fee payment."),
    ("RESEARCH", RED,
     "What exists: A section showing 10 faculty names + department labels. No paper titles, no journal names, no publication dates, no DOI links, no abstract, no impact factor. Functionally useless.",
     "What's missing: Searchable research database, faculty-wise publication lists with citations, ongoing projects, funded research, PhD scholars list, conference proceedings."),
    ("ALUMNI", RED,
     "What exists: Two cards — both Scius Solution employees listed as 'Notable Alumni'. No real alumni. No alumni directory, no batch-wise search, no achievement listings.",
     "What's missing: Alumni registration form, searchable directory, achievement gallery, alumni chapter contacts, mentorship programme link, social media community links."),
]

col_w = Inches(6.05)
for i, (name, color, exists, missing) in enumerate(sections):
    col = i % 2; row = i // 2
    lx = Inches(0.55) + col * Inches(6.4)
    ty = Inches(1.62) + row * Inches(1.82)
    rect(s5, lx, ty, col_w, Inches(1.72), DARK_BLUE)
    tag(s5, name, lx + Inches(0.1), ty + Inches(0.1), fill=color, tc=WHITE, size=9)
    tb(s5, "Has: " + exists[:110] + ("..." if len(exists) > 110 else ""),
       lx + Inches(0.12), ty + Inches(0.48), col_w - Inches(0.22), Inches(0.58),
       size=9.5, color=LGREY, wrap=True)
    tb(s5, "Needs: " + missing[:110] + ("..." if len(missing) > 110 else ""),
       lx + Inches(0.12), ty + Inches(1.1), col_w - Inches(0.22), Inches(0.55),
       size=9.5, color=AMBER, wrap=True)


# ══════════════════════════════════════════════════════════
# SLIDE 6 — WHAT FOREIGN UNIVERSITIES PUT ON THEIR WEBSITES
# ══════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
bg(s6)
rect(s6, 0, 0, Inches(0.12), H, TEAL)
tb(s6, "INTERNATIONAL BENCHMARK", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=TEAL)
tb(s6, "What Foreign Universities Put on Their Websites",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=30, bold=True, color=WHITE)
aline(s6, Inches(1.35), color=TEAL)

tb(s6, "Universities like University of Edinburgh, NUS Singapore, University of Melbourne and their peers treat their website as the university's primary interface with the world. Here is what they mandate:",
   Inches(0.55), Inches(1.5), Inches(12.23), Inches(0.52),
   size=12, color=LGREY, wrap=True)

categories = [
    ("ACADEMICS", TEAL, [
        "Every programme has its own page: entry requirements, learning outcomes, assessment methods, career pathways, fee schedule",
        "Faculty directory searchable by name, department, research interest — each with full bio, photo, publications, office hours",
        "Course catalogue with unit/module descriptions, credit hours, prerequisites — downloadable as PDF",
        "Academic calendar (term dates, exam periods, holiday schedule) updated annually",
    ]),
    ("ADMISSIONS", GOLD, [
        "Step-by-step application guide — what to submit, when, to whom, and what happens next",
        "Eligibility matrices by nationality, category, programme — no ambiguity",
        "Fee structure broken down: tuition, hostel, examination, miscellaneous — updated every year",
        "Scholarship and financial aid search tool with eligibility criteria and application deadlines",
    ]),
    ("RESEARCH", PURPLE, [
        "Research centres and labs listed with active projects, funding sources, and team members",
        "Faculty publications auto-synced from Scopus/ORCID — with paper titles, journal names, year, DOI",
        "PhD programme page: how to apply, supervisors available, current scholars, completion stats",
        "Industry partnerships, patents, and technology transfer disclosures",
    ]),
    ("STUDENT LIFE", ACCENT, [
        "Virtual campus tour — either 360° photography or video walkthrough of key buildings",
        "Hostel: room types, facilities, rules, fee, warden contact, application process",
        "Student clubs and societies directory — with contact person for each",
        "Sports, cultural events calendar with registration links — not just past events",
    ]),
    ("TRANSPARENCY", GREEN, [
        "NAAC/NIRF scores and accreditation documents — prominently displayed, not buried",
        "Annual reports, audit reports, RTI compliance disclosures — downloadable",
        "Grievance redressal portal with ticket tracking — not just an email address",
        "Anti-ragging committee, ICC (Internal Complaints Committee) contact — legally required",
    ]),
    ("TECHNICAL", RED, [
        "Mobile-first, WCAG 2.1 AA accessible — tested on actual assistive technology",
        "SSL, proper domain (not a vendor subdomain), fast CDN-backed loading",
        "CMS backend so admin staff can update notices, events, and results without a developer",
        "SEO: each department page ranks for its own name + location + courses",
    ]),
]
for i, (cat, color, items) in enumerate(categories):
    col = i % 3; row = i // 3
    lx = Inches(0.55) + col * Inches(4.28)
    ty = Inches(2.22) + row * Inches(2.35)
    rect(s6, lx, ty, Inches(4.0), Inches(2.18), DARK_BLUE)
    tag(s6, cat, lx + Inches(0.12), ty + Inches(0.1), fill=color, tc=WHITE, size=9)
    bullets(s6, items, lx + Inches(0.12), ty + Inches(0.5),
            Inches(3.8), Inches(1.6), size=9.5, color=LGREY, dot=color)


# ══════════════════════════════════════════════════════════
# SLIDE 7 — SPECIFIC PAGE-BY-PAGE FAILURES
# ══════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
bg(s7)
rect(s7, 0, 0, Inches(0.12), H, RED)
tb(s7, "PAGE-BY-PAGE FAILURE AUDIT", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=RED)
tb(s7, "Navigation Exists. Content Does Not.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s7, Inches(1.35), color=RED)

pages = [
    ("/home (Homepage)", "PARTIAL", AMBER,
     "Only working page. Contains: welcome text, VC/Chancellor messages (truncated), notice board, events (stale), research (names only), 2 fake alumni. Missing: programs, stats, campus photos."),
    ("/administration and all sub-items", "BROKEN", RED,
     "Loads homepage. No Registrar info, no Statutes, no Governance structure, no Finance committee details, no Senate/Syndicate composition — all expected on a public university site."),
    ("/academics and all sub-items", "BROKEN", RED,
     "Loads homepage. No departments, no courses, no faculty, no syllabus, no academic calendar. A student choosing this university has zero information about what they would study."),
    ("/admission-and-fee and all sub-items", "BROKEN", RED,
     "Loads homepage. Prospective students cannot find: eligibility, fee structure, application process, cut-offs, or important dates. This is the most critical journey on any university site."),
    ("/examination and all sub-items", "BROKEN", RED,
     "Loads homepage. Exam schedules, hall ticket downloads, result checking, re-evaluation forms — all missing. Students must rely on physical notice boards or WhatsApp groups."),
    ("/research and sub-items", "BLANK", RED,
     "Near-zero content (241 characters). No projects, no publications with actual titles, no PhD programme details, no research centres listed. The homepage research section is names-only."),
    ("/student-corner and all sub-items", "BROKEN", RED,
     "Loads homepage. NSS, NCC, YRC activities have zero dedicated pages. Student clubs, cultural events, sports — all absent. No student grievance portal despite being mandatory."),
    ("/campus-harmony and sub-items", "BLANK", RED,
     "Near-zero content. Campus infrastructure, sports facilities, labs, canteen, medical centre — none described. A student deciding where to spend 3–5 years cannot assess campus life."),
    ("/alumni and sub-items", "BLANK", RED,
     "Near-zero content (146 chars). Two 'notable alumni' are Scius Solution employees. No real graduates listed. No alumni network, no achievement gallery, no mentorship programme."),
    ("/information and sub-items", "BROKEN", RED,
     "Loads homepage. RTI disclosures, annual reports, audit reports, NAAC documentation, NIRF data — none available. These are legally required disclosures for a public institution."),
]

for i, (page, status, sc, desc) in enumerate(pages):
    t = Inches(1.58) + i * Inches(0.58)
    rect(s7, Inches(0.55), t, Inches(12.23), Inches(0.5), DARK_BLUE)
    tag(s7, status, Inches(0.62), t + Inches(0.1), fill=sc, tc=WHITE, size=8)
    tb(s7, page, Inches(1.92), t + Inches(0.06), Inches(2.5), Inches(0.35),
       size=11, bold=True, color=WHITE)
    tb(s7, desc, Inches(4.55), t + Inches(0.06), Inches(8.1), Inches(0.38),
       size=10, color=LGREY, wrap=True)


# ══════════════════════════════════════════════════════════
# SLIDE 8 — DATA INTEGRITY FAILURES
# ══════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
bg(s8)
rect(s8, 0, 0, Inches(0.12), H, AMBER)
tb(s8, "DATA INTEGRITY FAILURES", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=AMBER)
tb(s8, "Placeholder Data Published as Real Content",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=30, bold=True, color=WHITE)
aline(s8, Inches(1.35), color=AMBER)

issues = [
    ("Oman Phone Number in Site Header",
     "The header reads: 'Have any question? +968547856254'. Country code +968 = Oman. The university's actual number (06652 250 307) is in the footer. Anyone calling the header number reaches Oman. This is the first thing every visitor reads."),
    ("Vendor Employees Listed as University Alumni",
     "The 'Notable Alumni' section lists: Kshirod Biswal (Developer, Scius Solution, 2018) and Abinash Choudhury (Frontend Developer, Scius Solution, 2020). These are the web development team's employees, not graduates. This content was left in from the demo build."),
    ("Research Scholars Are Generic Placeholders",
     "'Scholars In Focus' shows: Ananya Sen (Statistics), Abhisek Das (Statistics), Arjun Nair (Statistics), Rohan Mehta (Statistics) — all with identical-sounding descriptions like 'Research scholar focusing on Bayesian statistics and data modeling.' These are dummy names cycling in a carousel."),
    ("Events From 2021 Shown as Current",
     "The Events section in 2026 still prominently features: 'One Day National Webinar — 03-02-2021', 'One Day National Seminar — 06-06-2021', 'First Foundation Day — 01-09-2021'. Stale 5-year-old events sit alongside 2026 entries because there is no content moderation system."),
    ("Chancellor/VC Messages Cut Off Mid-Sentence",
     "The Chancellor message ends with '...more' and the VC message ends with '...more'. Neither link works — both lead to the homepage. Visitors who want to read leadership messages are left with incomplete paragraphs. Standard university protocol requires full published statements."),
    ("Recent Publications Show No Publication Titles",
     "The 'Recent Publications' section lists: Dr. Kunal Mishra (Commerce), Miss Manisha Meher (Commerce), and 8 others — with only name and department. No paper title, no journal, no year, no DOI. This section communicates nothing about the university's research output."),
]

for i, (title, desc) in enumerate(issues):
    t = Inches(1.65) + i * Inches(0.97)
    rect(s8, Inches(0.55), t, Inches(12.23), Inches(0.85), DARK_BLUE)
    tb(s8, "!  " + title, Inches(0.7), t + Inches(0.08), Inches(5.5), Inches(0.35),
       size=13, bold=True, color=AMBER)
    tb(s8, desc, Inches(0.7), t + Inches(0.45), Inches(11.8), Inches(0.35),
       size=11, color=LGREY, wrap=True)


# ══════════════════════════════════════════════════════════
# SLIDE 9 — STUDENT-FACING FUNCTIONS — ALL ABSENT
# ══════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(BLANK)
bg(s9)
rect(s9, 0, 0, Inches(0.12), H, ACCENT)
tb(s9, "STUDENT-FACING DIGITAL FUNCTIONS", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s9, "What Students Need from a University Website",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=30, bold=True, color=WHITE)
aline(s9, Inches(1.35))

tb(s9, "The primary users of a university website are students — prospective, current, and alumni. Below is a checklist of what they need and what Rajendra University currently provides:",
   Inches(0.55), Inches(1.52), Inches(12.23), Inches(0.5), size=12, color=LGREY, wrap=True)

functions = [
    ("Online Admission Application", "Apply for programmes online, track application status", "ABSENT"),
    ("Fee Payment Portal", "Pay tuition, hostel, exam fees online — receipt download", "ABSENT"),
    ("Exam Schedule & Hall Tickets", "Download admit cards, check exam timetable by programme", "ABSENT"),
    ("Result / Marks Portal", "Check semester results, marksheets, revaluation apply", "ABSENT"),
    ("Course & Syllabus Download", "Download semester syllabus PDFs per programme", "ABSENT"),
    ("Faculty Directory", "Find faculty member contact, office hours, research area", "ABSENT"),
    ("Notice Board Search", "Filter notices by category, department, date range", "PARTIAL"),
    ("Library Catalogue (OPAC)", "Search books, journals, e-resources online", "ABSENT"),
    ("Grievance Portal", "Submit complaints, track resolution status online", "ABSENT"),
    ("Hostel Application & Info", "Apply for hostel, view room availability, rules, fees", "ABSENT"),
    ("Scholarship Information", "Eligibility, amount, application process per scheme", "ABSENT"),
    ("Placement Statistics", "Company-wise placements, packages, sector breakdown", "ABSENT"),
]

for i, (func, desc, status) in enumerate(functions):
    col = i % 2; row = i // 2
    lx = Inches(0.55) + col * Inches(6.4)
    ty = Inches(2.18) + row * Inches(0.62)
    sc = RED if status == "ABSENT" else AMBER
    rect(s9, lx, ty, Inches(6.1), Inches(0.52), DARK_BLUE)
    tag(s9, status, lx + Inches(0.1), ty + Inches(0.11), fill=sc, tc=WHITE, size=9)
    tb(s9, func, lx + Inches(1.3), ty + Inches(0.08), Inches(2.5), Inches(0.32),
       size=12, bold=True, color=WHITE)
    tb(s9, desc, lx + Inches(3.85), ty + Inches(0.1), Inches(2.2), Inches(0.32),
       size=10, color=MGREY, wrap=True)


# ══════════════════════════════════════════════════════════
# SLIDE 10 — WHAT A PROPER REBUILD LOOKS LIKE
# ══════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(BLANK)
bg(s10)
rect(s10, 0, 0, Inches(0.12), H, GREEN)
tb(s10, "WHAT A PROPER REBUILD LOOKS LIKE", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s10, "The Standard Rajendra University Deserves",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=30, bold=True, color=WHITE)
aline(s10, Inches(1.35), color=GREEN)

phases = [
    ("PHASE 1\nFoundation", GREEN, [
        "University-owned domain: rajendrauniversity.ac.in (not a vendor subdomain)",
        "CMS (WordPress/Drupal) so admin staff update content without a developer",
        "Correct all placeholder data: real phone, real alumni, real scholars",
        "Homepage redesign: campus photography, NAAC score, programme highlights",
        "Functional navigation with all 9 top-level sections properly built out",
    ]),
    ("PHASE 2\nCore Pages", GOLD, [
        "Department pages: all departments with faculty profiles, courses, research",
        "Admission hub: eligibility, fee structure, application process, dates",
        "Examination portal: schedule download, hall ticket, result link",
        "Research section: actual publications with titles, journals, DOIs",
        "News & Events: date-filtered, searchable, with full articles and registration",
    ]),
    ("PHASE 3\nStudent Portal", ACCENT, [
        "Student self-service: fee payment, marksheets, admit card, grievance form",
        "Library OPAC integration or link to catalogue",
        "Hostel application and information page",
        "Placement statistics with company names and packages",
        "Alumni directory with registration and batch-wise search",
    ]),
]

for i, (phase, color, items) in enumerate(phases):
    cl = Inches(0.55) + i * Inches(4.25)
    rect(s10, cl, Inches(1.62), Inches(4.0), Inches(5.48), DARK_BLUE)
    rect(s10, cl, Inches(1.62), Inches(4.0), Inches(0.5), color)
    tb(s10, phase, cl, Inches(1.62), Inches(4.0), Inches(0.5),
       size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    bullets(s10, items, cl + Inches(0.15), Inches(2.28),
            Inches(3.72), Inches(4.6), size=12, color=LGREY, dot=color)

rect(s10, Inches(0.55), Inches(7.1), Inches(12.23), Inches(0.28), DARK_BLUE)
tb(s10, "Timeline: Phase 1 — 4 weeks  |  Phase 2 — 6–8 weeks  |  Phase 3 — 4–6 weeks  |  Total: 14–18 weeks for a complete, functional university digital presence",
   Inches(0.7), Inches(7.12), Inches(12), Inches(0.22), size=10, color=MGREY, wrap=False)


# ══════════════════════════════════════════════════════════
# SLIDE 11 — NEXT STEPS
# ══════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(BLANK)
bg(s11)
rect(s11, 0, 0, Inches(0.12), H, GOLD)
tb(s11, "NEXT STEPS", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GOLD)
tb(s11, "From a Vendor Demo to a Real University Website",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=30, bold=True, color=WHITE)
aline(s11, Inches(1.35), color=GOLD)

steps = [
    ("1", "Discovery Call — 30 Minutes",
     "We walk the Administration team through every specific finding in this audit with screen recordings — so there is no ambiguity about what is broken and why. No pitch, no commitment."),
    ("2", "Design Brief & Sitemap Review",
     "We produce a full sitemap covering all required pages, their content structure, and user journeys for the 4 primary audiences: prospective students, current students, faculty, and administrative bodies."),
    ("3", "Phased Rebuild with Handover Training",
     "We build on a proper CMS with full documentation and staff training so the university team can update notices, results, events, and news independently without requiring the vendor for every change."),
]
for i, (num, title, desc) in enumerate(steps):
    t = Inches(2.0) + i * Inches(1.5)
    rect(s11, Inches(0.55), t, Inches(12.23), Inches(1.3), DARK_BLUE)
    rect(s11, Inches(0.55), t, Inches(1.0), Inches(1.3), GOLD)
    tb(s11, num, Inches(0.55), t + Inches(0.2), Inches(1.0), Inches(0.7),
       size=36, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    tb(s11, title, Inches(1.7), t + Inches(0.15), Inches(10), Inches(0.45),
       size=16, bold=True, color=WHITE)
    tb(s11, desc, Inches(1.7), t + Inches(0.6), Inches(10.5), Inches(0.55),
       size=12.5, color=LGREY)

rect(s11, 0, Inches(6.5), W, Inches(1.0), ACCENT)
tb(s11, "Rohan Malik  ·  CEO, Genos AI  ·  hello@genosai.tech  ·  +91 638-714-2699  ·  www.genosai.tech",
   Inches(0.55), Inches(6.6), Inches(12.23), Inches(0.55),
   size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ── save ──────────────────────────────────────────────────
OUT = "RajendraUniversity_Audit_GenosAI.pptx"
prs.save(OUT)
print(f"Saved: {OUT}  |  11 slides")
