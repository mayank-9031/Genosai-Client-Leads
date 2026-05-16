import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.parent / "utils"))
from ppt_theme import (
    NAVY, DARK_BLUE, ACCENT, GOLD, RED, GREEN, PURPLE, ORANGE, TEAL, CYAN,
    WHITE, LGREY, MGREY, W, H,
    bg, rect, tb, aline, bullets, tag,
    slide_header, slide_footer, score_badge,
    stat_cards, top_bar_card, phase_card, before_after, genos_logo_block,
)

ROOT = Path(__file__).parent.parent.parent
prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ══ SLIDE 1 — COVER ══
s1 = prs.slides.add_slide(BLANK)
bg(s1)
rect(s1, 0, 0, Inches(0.12), H, GOLD)
rect(s1, Inches(8.5), 0, Inches(4.83), Inches(2.2), DARK_BLUE)
tb(s1, "GENOS AI", Inches(9.2), Inches(0.4), Inches(3.5), Inches(0.8),
   size=28, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
tb(s1, "www.genosai.tech", Inches(9.2), Inches(1.05), Inches(3.5), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
tb(s1, "QUALITY ACCEPTANCE / NEXA OPTIMUM SOLUTIONS", Inches(0.55), Inches(1.6),
   Inches(9), Inches(0.6), size=13, bold=False, color=GOLD)
tb(s1, "Website & AI Growth\nAudit Report", Inches(0.55), Inches(2.1),
   Inches(9), Inches(1.8), size=48, bold=True, color=WHITE)
aline(s1, Inches(3.85), color=GOLD, w=Inches(3))
tb(s1, "Prepared exclusively for Ofer Alon, CEO & Founder",
   Inches(0.55), Inches(4.15), Inches(9), Inches(0.5), size=14, color=LGREY)
rect(s1, 0, Inches(6.5), W, Inches(1.0), DARK_BLUE)
tb(s1, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699",
   Inches(0.55), Inches(6.6), Inches(9), Inches(0.5), size=12, color=MGREY)
tb(s1, "May 2026  |  CONFIDENTIAL", Inches(9.5), Inches(6.6), Inches(3.3), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
rect(s1, Inches(10.5), Inches(3.3), Inches(2.3), Inches(2.5), DARK_BLUE)
tb(s1, "AUDIT SCORE", Inches(10.55), Inches(3.4), Inches(2.2), Inches(0.5),
   size=11, bold=True, color=MGREY, align=PP_ALIGN.CENTER)
tb(s1, "5", Inches(10.55), Inches(3.75), Inches(2.2), Inches(1.1),
   size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
tb(s1, "/ 10", Inches(10.55), Inches(4.7), Inches(2.2), Inches(0.4),
   size=18, color=LGREY, align=PP_ALIGN.CENTER)
tag(s1, "HIGH PRIORITY", Inches(10.75), Inches(5.1), fill=RED, size=10)


# ══ SLIDE 2 — AT A GLANCE ══
s2 = prs.slides.add_slide(BLANK)
bg(s2)
rect(s2, 0, 0, Inches(0.12), H, GOLD)
tb(s2, "AT A GLANCE", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GOLD)
tb(s2, "Ofer Alon — Who He Is",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s2, Inches(1.35), color=GOLD)

stats = [
    ("25+", "Years Auto Finance"),
    ("$1B+", "Funded Transactions"),
    ("180K+", "Consumers Served"),
    ("4 States", "CA, NV, TX, AZ"),
    ("A+", "BBB Rating"),
]
cw = Inches(2.3); ch = Inches(2.2); gap = Inches(0.18)
for i, (num, label) in enumerate(stats):
    l = Inches(0.55) + i * (cw + gap)
    rect(s2, l, Inches(1.8), cw, ch, DARK_BLUE)
    tb(s2, num, l + Inches(0.1), Inches(1.95), cw - Inches(0.2), Inches(1.1),
       size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s2, label, l + Inches(0.1), Inches(3.05), cw - Inches(0.2), Inches(0.7),
       size=12, color=LGREY, align=PP_ALIGN.CENTER)

rect(s2, Inches(0.55), Inches(4.2), Inches(12.23), Inches(0.6), DARK_BLUE)
tb(s2, "Van Nuys CA  ·  Founded 1998  ·  IDB Bank-Backed Facility  ·  AFSA Board Member  ·  1st Franklin Financial Board (April 2026)  ·  Nexa Optimum Solutions (2025)",
   Inches(0.65), Inches(4.28), Inches(12), Inches(0.45), size=12, color=LGREY, align=PP_ALIGN.CENTER)

tb(s2, "BACKGROUND THAT SETS HIM APART", Inches(0.55), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=GOLD)
bullets(s2, ["Harvard Business School, University of Bradford, Tel Aviv University CPA",
             "Former Ernst & Young manager — discipline of a Big Four firm",
             "Former tech CEO at Buildcom.com — not a typical auto finance background"],
        Inches(0.55), Inches(5.35), Inches(5.8), Inches(1.5), size=13, color=LGREY)

tb(s2, "WHAT HE BUILT", Inches(6.8), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=GOLD)
bullets(s2, ["Quality Acceptance: subprime auto finance, 1998 to present",
             "Nexa Optimum Solutions: collections-focused lender adaptation platform, 2025",
             "Auto Remarketing commentary author on the AI implementation gap in finance"],
        Inches(6.8), Inches(5.35), Inches(5.8), Inches(1.5), size=13, color=LGREY)


# ══ SLIDE 3 — WHAT'S WORKING ══
s3 = prs.slides.add_slide(BLANK)
bg(s3)
rect(s3, 0, 0, Inches(0.12), H, GREEN)
tb(s3, "WHAT'S WORKING", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s3, "Strengths to Build On",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s3, Inches(1.35), color=GREEN)

strengths = [
    ("25+ Years in Subprime — Where Others Exit",
     "Quality Acceptance has operated continuously since 1998 through multiple credit cycles, recessions, and regulatory shifts. Most lenders in non-prime auto finance either blow up or exit. That survival record across 25+ years and $1B+ in funded transactions is institutional credibility most competitors can't replicate."),
    ("Ofer's Tech-First Background Is Unusual for This Space",
     "Auto finance founders typically come from car dealerships or traditional lending. Ofer came from Ernst & Young, Harvard, and ran an early tech company (Buildcom.com). He writes publicly about AI and implementation gaps in financial services. The sophistication is there — the question is whether the infrastructure underneath matches it."),
    ("Nexa Optimum Solutions — He Sees Where Lending Is Going",
     "Founding Nexa in 2025 specifically to help lenders adapt to a changing market signals Ofer isn't waiting for the industry to modernize. He's building the infrastructure for it. That same instinct applied internally to Quality Acceptance's own dealer and consumer channels is a significant leverage point."),
    ("Institutional Backing and Industry Credibility",
     "IDB Bank-backed asset revolving credit facility gives Quality Acceptance a serious capital foundation. AFSA board membership and the April 2026 appointment to 1st Franklin Financial's board put Ofer in rooms most auto finance operators don't access. The operation behind the website is substantially more sophisticated than the website suggests."),
    ("Mobile App Already Deployed — Consumer Tech Infrastructure Exists",
     "Quality Acceptance has an iOS and Android app for consumer payment management. That's infrastructure most lenders this size don't have. The foundation for self-service and push notification automation is already built — the automation layer on top is what's missing."),
]
for i, (title, desc) in enumerate(strengths):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s3, Inches(0.55), t, Inches(0.06), Inches(0.7), GREEN)
    tb(s3, title, Inches(0.75), t, Inches(11.8), Inches(0.38), size=14, bold=True, color=WHITE)
    tb(s3, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# ══ SLIDE 4 — THE IMPLEMENTATION GAP ══
s4 = prs.slides.add_slide(BLANK)
bg(s4)
rect(s4, 0, 0, Inches(0.12), H, RED)
tb(s4, "THE IMPLEMENTATION GAP", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=RED)
tb(s4, "The Tech Is There. The Automation Layer Underneath Isn't.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=30, bold=True, color=WHITE)
aline(s4, Inches(1.35), color=RED)

issues = [
    ("NO DEALER ACQUISITION FUNNEL — JUST A PHONE NUMBER", "HIGH",
     "A dealer researching subprime financing partners lands on qualityacceptance.com and finds a 1-800 number. No online enrollment form, no CTA, no automated onboarding sequence. The conversion path from site visit to enrolled dealer is a phone call during business hours — which means every dealer who lands after hours either calls back or doesn't."),
    ("NO CONSUMER SELF-SERVICE OR AI CHAT", "HIGH",
     "Consumers with questions about their balance, payment options, or insurance status call 818-503-1322. That's a servicing team handling inquiries that AI chat could resolve at 2am on a Sunday without any human involvement. The mobile app exists but the website has no chat widget, no self-service portal link, no automated response path."),
    ("PAYMENT REMINDERS APPEAR STATIC, NOT BEHAVIORAL", "HIGH",
     "Standard payment reminder sequences send the same messages on the same schedule regardless of consumer behavior. Behavioral-trigger sequences respond to actual patterns — days since last payment, communication response history, payment method preference — and consistently outperform static schedules on early-stage delinquency reduction."),
    ("INSURANCE LAPSE MANAGEMENT APPEARS MANUAL", "MEDIUM",
     "Insurance lapse is both a compliance risk and a loss exposure. Automated alert sequences triggered immediately on lapse notification, with escalation at defined intervals before forced placement, reduce both exposure and the cost of managing it. Manual monitoring at this portfolio scale has gaps."),
    ("NO DEALER REACTIVATION AUTOMATION", "MEDIUM",
     "Dealers who stop submitting deals don't always leave — they get busy, get recruited by a competitor, or stop after one friction point. Automated sequences at 30, 60, and 90 days of dormancy, with content relevant to what they submitted before, recover a percentage of those relationships without requiring a manual call for every one."),
]
for i, (title, pri, desc) in enumerate(issues):
    t = Inches(1.7) + i * Inches(1.05)
    pc = RED if pri == "HIGH" else GOLD
    rect(s4, Inches(0.55), t, Inches(0.06), Inches(0.7), pc)
    tb(s4, title, Inches(0.75), t, Inches(9.2), Inches(0.38), size=12, bold=True, color=WHITE)
    tag(s4, pri, Inches(10.1), t + Inches(0.05), fill=pc, size=9)
    tb(s4, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# ══ SLIDE 5 — WEBSITE DESIGN AUDIT ══
s5 = prs.slides.add_slide(BLANK)
bg(s5)
rect(s5, 0, 0, Inches(0.12), H, PURPLE)
tb(s5, "WEBSITE DESIGN AUDIT", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=PURPLE)
tb(s5, "$1B+ Funded. 25 Years. The Site Doesn't Show It.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.8), size=32, bold=True, color=WHITE)
aline(s5, Inches(1.45), color=PURPLE, w=Inches(1.5))

rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(3.6), DARK_BLUE)
rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(0.42), RED)
tb(s5, "QUALITYACCEPTANCE.COM TODAY", Inches(0.65), Inches(1.92), Inches(3.6), Inches(0.38),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Dated design understates institutional scale",
    "Dealer page = phone number only, no enrollment funnel",
    "No AI chat or self-service widget for consumers",
    "No testimonials from dealers or borrowers",
    "No case studies or funded volume data visible",
    "1-800 number as primary conversion action",
    "Mobile app exists but site doesn't drive to it",
    "No credibility signals: A+ BBB, $1B+, AFSA board",
], Inches(0.7), Inches(2.42), Inches(3.6), Inches(2.9), size=11, color=LGREY, dot=RED)

tb(s5, "->", Inches(4.45), Inches(3.6), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(3.6), DARK_BLUE)
rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(0.42), GOLD)
tb(s5, "TOP AUTO FINANCE LENDER SITES DO THIS", Inches(5.1), Inches(1.92), Inches(3.4), Inches(0.38),
   size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Online dealer enrollment with instant confirmation",
    "Consumer AI chat for balance/payment/insurance 24/7",
    "Dealer testimonials and funded volume data",
    "Credibility signals: ratings, years, capital backing",
    "Behavioral payment reminder sequences, not static",
    "Insurance lapse alert automation",
    "Dealer dashboard/portal access from homepage",
    "Case studies by dealer type and state",
], Inches(5.1), Inches(2.42), Inches(3.4), Inches(2.9), size=11, color=LGREY, dot=GOLD)

tb(s5, "->", Inches(8.68), Inches(3.6), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(3.6), DARK_BLUE)
rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(0.42), GREEN)
tb(s5, "WHAT GENOS AI BUILDS", Inches(9.35), Inches(1.92), Inches(3.4), Inches(0.38),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Dealer acquisition funnel with online enrollment CTA",
    "Consumer AI chat: balance, payment, insurance 24/7",
    "Dealer and consumer testimonials section",
    "Credibility signals front and center ($1B+, A+, 25yr)",
    "Behavioral-trigger payment reminder sequences",
    "Insurance lapse alert and escalation automation",
    "Dealer reactivation sequences (30/60/90 day)",
    "Site rebuild reflecting institutional scale",
], Inches(9.35), Inches(2.42), Inches(3.4), Inches(2.9), size=11, color=LGREY, dot=GREEN)

rect(s5, Inches(0.55), Inches(5.72), Inches(12.23), Inches(1.28), DARK_BLUE)
tb(s5, "WHY IT MATTERS:", Inches(0.7), Inches(5.78), Inches(2.2), Inches(0.38),
   size=11, bold=True, color=PURPLE)
stats2 = [
    ("60-70%", "of inbound calls are\nroutine self-service questions"),
    ("30/60/90", "day dormancy sequences\nrecover 15-25% of lapsed dealers"),
    ("48 hrs", "insurance lapse window\nbefore meaningful loss exposure"),
    ("3x", "conversion lift with\nonline enrollment vs. phone-only"),
]
sw = Inches(2.7); sl = Inches(2.9)
for i, (num, label) in enumerate(stats2):
    xl = sl + i * sw
    tb(s5, num, xl, Inches(5.75), Inches(2.5), Inches(0.55),
       size=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s5, label, xl, Inches(6.22), Inches(2.5), Inches(0.7),
       size=10, color=LGREY, align=PP_ALIGN.CENTER)


# ══ SLIDE 6 — AI AUTOMATION OPPORTUNITIES ══
s6 = prs.slides.add_slide(BLANK)
bg(s6)
rect(s6, 0, 0, Inches(0.12), H, ACCENT)
tb(s6, "AI AUTOMATION OPPORTUNITIES", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s6, "Closing the Implementation Gap.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=34, bold=True, color=WHITE)
aline(s6, Inches(1.35), color=ACCENT)

opps = [
    ("Dealer Acquisition\nFunnel", GOLD,
     "Online enrollment form with automated onboarding sequence. New dealer submits online, gets instant confirmation, automated deal-submission guide, first check-in at day 7. Converts site visits into active dealer relationships without a phone call per prospect."),
    ("Consumer AI Chat\n24/7 Self-Service", ACCENT,
     "Chat widget trained on Quality Acceptance payment options, insurance requirements, balance inquiries, and account procedures. Handles 60-70% of inbound call volume at any hour without touching servicing staff. Routes complex issues to the team with context already captured."),
    ("Behavioral-Trigger\nPayment Sequences", GREEN,
     "Replace static reminder schedules with sequences that respond to consumer behavior: days since last payment, response to prior communications, payment method history. Early-stage delinquency reduction is measurable within the first 90 days."),
    ("Insurance Lapse\nAlert Automation", RED,
     "Triggered immediately on lapse notification: consumer alert sequence, dealer notification if applicable, escalation at defined intervals before forced placement. Reduces both loss exposure and the manual cost of tracking lapse status across a large portfolio."),
    ("Dealer Reactivation\nSequences", PURPLE,
     "Dealers who stop submitting deals at 30, 60, and 90 days receive automated outreach with content relevant to their prior submission history. Recovers a measurable percentage of dormant relationships without requiring a manual call for each one."),
]
cw2 = Inches(2.35); gap2 = Inches(0.14)
for i, (title, color, desc) in enumerate(opps):
    l2 = Inches(0.55) + i * (cw2 + gap2)
    rect(s6, l2, Inches(1.8), cw2, Inches(0.06), color)
    rect(s6, l2, Inches(1.86), cw2, Inches(4.6), DARK_BLUE)
    tb(s6, title, l2 + Inches(0.12), Inches(1.98), cw2 - Inches(0.24), Inches(0.65),
       size=13, bold=True, color=color)
    tb(s6, desc, l2 + Inches(0.12), Inches(2.65), cw2 - Inches(0.24), Inches(3.6),
       size=11, color=LGREY)

rect(s6, Inches(0.55), Inches(6.7), Inches(12.23), Inches(0.45), DARK_BLUE)
tb(s6, '"Technology is not failing; implementation is failing." — Ofer Alon, Auto Remarketing',
   Inches(0.7), Inches(6.75), Inches(12), Inches(0.35), size=11, italic=True, color=GOLD)


# ══ SLIDE 7 — WHAT THIS IS COSTING YOU ══
s7 = prs.slides.add_slide(BLANK)
bg(s7)
rect(s7, 0, 0, Inches(0.12), H, RED)
tb(s7, "WHAT THIS IS COSTING YOU", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=RED)
tb(s7, "The Operation Is Sophisticated. The Digital Layer Isn't.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s7, Inches(1.35), color=RED)

costs = [
    ("No dealer funnel = every new dealer relationship starts with a phone tag",
     "A dealer who finds Quality Acceptance after hours, fills out no form, gets no follow-up, and moves on to the next lender in their search results. The phone-only enrollment model means dealer acquisition is capped by when your team is available, not by how many dealers are looking."),
    ("No consumer self-service = your servicing team fields questions AI could handle",
     "60-70% of consumer inbound calls are routine: what's my balance, when is my payment due, what are my payment options, why did my insurance get flagged. Those calls don't require a human. They require a system that answers reliably at 2am. Every one your team takes is cost that doesn't have to exist."),
    ("Static payment reminders = higher early-stage delinquency than necessary",
     "The consumer who paid late twice and then recovered responds differently to reminders than the consumer who has never missed. Static schedules treat both the same. Behavioral-trigger sequences send the right message at the right moment based on actual account history, and the delinquency reduction shows up in the data."),
    ("Manual insurance lapse tracking = exposure window between lapse and action",
     "At portfolio scale, manual monitoring of insurance lapse status has gaps. Every day between lapse and detection is risk. Automated alert sequences triggered at the moment of lapse notification close that window and create a documented escalation trail."),
    ("No dealer reactivation = losing relationships that haven't actually ended",
     "A dealer who goes quiet at month 3 hasn't necessarily switched lenders. They may have had a bad experience with one deal, gotten busy, or simply not been contacted. An automated reactivation sequence at 30/60/90 days recovers the recoverable ones without requiring someone to manually review every dormant account."),
]
for i, (title, desc) in enumerate(costs):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s7, Inches(0.55), t, Inches(0.06), Inches(0.7), RED)
    tb(s7, title, Inches(0.75), t, Inches(11.8), Inches(0.38), size=14, bold=True, color=WHITE)
    tb(s7, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# ══ SLIDE 8 — WHAT WE'D BUILD ══
s8 = prs.slides.add_slide(BLANK)
bg(s8)
rect(s8, 0, 0, Inches(0.12), H, GREEN)
tb(s8, "WHAT WE'D BUILD", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s8, "Genos AI for Quality Acceptance",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s8, Inches(1.35), color=GREEN)

builds = [
    ("Dealer Acquisition Funnel + Automated Onboarding",
     "Online enrollment form on qualityacceptance.com with instant confirmation. Automated onboarding sequence: deal submission guide, program overview, first check-in at day 7, review request at day 30. Converts site visits into active dealer relationships without manual touchpoints for each prospect."),
    ("Consumer AI Chat — 24/7 Self-Service",
     "Chat widget trained on Quality Acceptance products: payment options, insurance requirements, balance inquiries, account status, app download guidance. Handles 60-70% of inbound volume at any hour. Routes complex issues to the servicing team with full conversation context captured."),
    ("Behavioral-Trigger Payment Reminder Sequences",
     "Replace static schedules with sequences built on account history. Consumer A (consistent payer, one late) receives different messaging than Consumer B (pattern of lates). Triggers respond to payment behavior in real time. Early-stage delinquency reduction is measurable within 90 days."),
    ("Insurance Lapse Alert Automation + Escalation",
     "Triggered at the moment of lapse notification. Consumer sequence: immediate alert, payment options for reinstating coverage, escalation at day 3 and day 7. Internal escalation trail for compliance documentation. Closes the gap between lapse detection and action."),
    ("Site Rebuild Reflecting Institutional Scale + Dealer Reactivation",
     "qualityacceptance.com rebuilt to reflect 25+ years, $1B+ funded, A+ BBB, IDB Bank backing. Dealer testimonials, case studies, funded volume data. Dealer reactivation sequences at 30/60/90 day dormancy marks, personalized by prior submission history and dealer type."),
]
for i, (title, desc) in enumerate(builds):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s8, Inches(0.55), t, Inches(0.06), Inches(0.7), GREEN)
    tb(s8, title, Inches(0.75), t, Inches(11.8), Inches(0.38), size=14, bold=True, color=WHITE)
    tb(s8, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# ══ SLIDE 9 — 90-DAY OUTCOMES ══
s9 = prs.slides.add_slide(BLANK)
bg(s9)
rect(s9, 0, 0, Inches(0.12), H, ACCENT)
tb(s9, "90-DAY OUTCOMES", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s9, "What Changes in 90 Days",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s9, Inches(1.35))

phases = [
    ("Month 1", GOLD, [
        "AI chat live on qualityacceptance.com - consumer self-service running 24/7",
        "Dealer acquisition funnel active with online enrollment and automated onboarding",
        "Insurance lapse alert automation deployed - gap between lapse and action closed",
        "Behavioral-trigger payment sequence framework built and first segments running",
    ]),
    ("Month 2", ACCENT, [
        "Behavioral payment sequences fully deployed across consumer segments",
        "Dealer reactivation sequences running for 30/60/90 day dormant accounts",
        "Site rebuild underway with credibility signals, dealer testimonials, case studies",
        "First AI chat deflection data: call volume reduction measurable",
    ]),
    ("Month 3", GREEN, [
        "Site live with institutional credibility: $1B+ funded, A+ BBB, dealer testimonials",
        "Early-stage delinquency change vs. prior static schedule measurable",
        "Dealer reactivation sequences recovering dormant relationships",
        "Full implementation gap closed - automation layer live across both audiences",
    ]),
]
pw = Inches(3.9); gap3 = Inches(0.18)
for i, (month, color, pts) in enumerate(phases):
    l3 = Inches(0.55) + i * (pw + gap3)
    rect(s9, l3, Inches(1.8), pw, Inches(0.42), color)
    tb(s9, month.upper(), l3 + Inches(0.1), Inches(1.84), pw - Inches(0.2), Inches(0.36),
       size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    rect(s9, l3, Inches(2.22), pw, Inches(4.7), DARK_BLUE)
    bullets(s9, pts, l3 + Inches(0.15), Inches(2.35), pw - Inches(0.3), Inches(4.4),
            size=12, color=LGREY, dot=color)


# ══ SLIDE 10 — NEXT STEPS ══
s10 = prs.slides.add_slide(BLANK)
bg(s10)
rect(s10, 0, 0, Inches(0.12), H, GOLD)
rect(s10, Inches(7.5), 0, Inches(5.83), H, DARK_BLUE)
tb(s10, "NEXT STEPS", Inches(0.55), Inches(0.3), Inches(6.5), Inches(0.5),
   size=11, bold=True, color=GOLD)
tb(s10, "Let's Talk for\n15 Minutes.",
   Inches(0.55), Inches(0.8), Inches(6.8), Inches(2.0), size=48, bold=True, color=WHITE)
aline(s10, Inches(2.8), color=GOLD, w=Inches(2.5))
tb(s10, "No pitch deck required. Reply or book a call directly.\nWe'll walk through what fits Quality Acceptance and Nexa.",
   Inches(0.55), Inches(3.05), Inches(6.5), Inches(0.8), size=14, color=LGREY)

steps = [
    "Reply to this email - we'll set up a time",
    "Book directly: www.genosai.tech/call",
    "Walk through the audit findings in 15 minutes",
    "You decide what, if anything, makes sense to move on",
]
bullets(s10, steps, Inches(0.55), Inches(3.95), Inches(6.5), Inches(2.0), size=14, color=WHITE, dot=GOLD)

rect(s10, 0, Inches(6.5), Inches(7.5), Inches(1.0), DARK_BLUE)
tb(s10, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699  |  www.genosai.tech",
   Inches(0.55), Inches(6.6), Inches(6.8), Inches(0.5), size=12, color=MGREY)

tb(s10, "WHAT WE BUILD FOR\nQUALITY ACCEPTANCE", Inches(7.85), Inches(0.6), Inches(5.0), Inches(0.9),
   size=18, bold=True, color=GOLD)
aline(s10, Inches(1.45), color=GOLD)

summary_items = [
    ("Dealer Funnel", "Online enrollment + automated onboarding sequences"),
    ("AI Chat", "24/7 consumer self-service - 60-70% call deflection"),
    ("Payment", "Behavioral-trigger sequences replacing static schedules"),
    ("Insurance", "Lapse alert automation - closes the detection-action gap"),
    ("Reactivation", "30/60/90 day sequences for dormant dealer relationships"),
    ("Site Rebuild", "Reflects $1B+ funded, 25 years, A+ BBB, IDB backing"),
]
for i, (label, val) in enumerate(summary_items):
    t2 = Inches(1.6) + i * Inches(0.88)
    rect(s10, Inches(7.85), t2, Inches(1.5), Inches(0.72), NAVY)
    tb(s10, label, Inches(7.9), t2 + Inches(0.05), Inches(1.4), Inches(0.32),
       size=10, bold=True, color=GOLD)
    tb(s10, val, Inches(9.5), t2 + Inches(0.05), Inches(3.6), Inches(0.65),
       size=11, color=LGREY)


prs.save(ROOT / "output" / "decks" / "QualityAcceptance_Audit_GenosAI.pptx")
print("Saved: QualityAcceptance_Audit_GenosAI.pptx")
