import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.parent / "utils"))
from ppt_theme import (
    NAVY, DARK_BLUE, ACCENT, GOLD, RED, GREEN, PURPLE, ORANGE, TEAL, CYAN,
    WHITE, LGREY, MGREY, W, H,
    bg, rect, tb, aline, bullets, tag,
    slide_header, slide_footer, score_badge,
    stat_cards, top_bar_card, phase_card, before_after, genos_logo_block,
)

ROOT = Path(__file__).parent.parent.parent
prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]

# SLIDE 1 - COVER
s1 = prs.slides.add_slide(BLANK)
bg(s1)
rect(s1, 0, 0, Inches(0.12), H, GOLD)
rect(s1, Inches(8.5), 0, Inches(4.83), Inches(2.2), DARK_BLUE)
tb(s1, "GENOS AI", Inches(9.2), Inches(0.4), Inches(3.5), Inches(0.8),
   size=28, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
tb(s1, "www.genosai.tech", Inches(9.2), Inches(1.05), Inches(3.5), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
tb(s1, "MMJ NORTH REAL ESTATE - CORRIMAL", Inches(0.55), Inches(1.6),
   Inches(9), Inches(0.6), size=13, bold=False, color=GOLD)
tb(s1, "Website & AI Growth\nAudit Report", Inches(0.55), Inches(2.1),
   Inches(9), Inches(1.8), size=48, bold=True, color=WHITE)
aline(s1, Inches(3.85), color=GOLD, w=Inches(3))
tb(s1, "Prepared exclusively for Greg Ellul, North Director / Licensee",
   Inches(0.55), Inches(4.15), Inches(9), Inches(0.5), size=14, color=LGREY)
rect(s1, 0, Inches(6.5), W, Inches(1.0), DARK_BLUE)
tb(s1, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699",
   Inches(0.55), Inches(6.6), Inches(9), Inches(0.5), size=12, color=MGREY)
tb(s1, "May 2026  |  CONFIDENTIAL", Inches(9.5), Inches(6.6), Inches(3.3), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
rect(s1, Inches(10.5), Inches(3.3), Inches(2.3), Inches(2.5), DARK_BLUE)
tb(s1, "AUDIT SCORE", Inches(10.55), Inches(3.4), Inches(2.2), Inches(0.5),
   size=11, bold=True, color=MGREY, align=PP_ALIGN.CENTER)
tb(s1, "5", Inches(10.55), Inches(3.75), Inches(2.2), Inches(1.1),
   size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
tb(s1, "/ 10", Inches(10.55), Inches(4.7), Inches(2.2), Inches(0.4),
   size=18, color=LGREY, align=PP_ALIGN.CENTER)
tag(s1, "HIGH PRIORITY", Inches(10.75), Inches(5.1), fill=RED, size=10)


# SLIDE 2 - AT A GLANCE
s2 = prs.slides.add_slide(BLANK)
bg(s2)
rect(s2, 0, 0, Inches(0.12), H, GOLD)
tb(s2, "AT A GLANCE", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GOLD)
tb(s2, "Greg Ellul - Who He Is",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s2, Inches(1.35), color=GOLD)

stats = [
    ("64", "Years of MMJ"),
    ("30+", "Years Greg in RE"),
    ("6", "REINSW Finalists 2025"),
    ("30+", "Agents at North Office"),
    ("11-50", "Office Headcount"),
]
cw = Inches(2.3); ch = Inches(2.2); gap = Inches(0.18)
for i, (num, label) in enumerate(stats):
    l = Inches(0.55) + i * (cw + gap)
    rect(s2, l, Inches(1.8), cw, ch, DARK_BLUE)
    tb(s2, num, l + Inches(0.1), Inches(1.95), cw - Inches(0.2), Inches(1.1),
       size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s2, label, l + Inches(0.1), Inches(3.05), cw - Inches(0.2), Inches(0.7),
       size=12, color=LGREY, align=PP_ALIGN.CENTER)

rect(s2, Inches(0.55), Inches(4.2), Inches(12.23), Inches(0.6), DARK_BLUE)
tb(s2, "Corrimal NSW  -  MMJ Founded 1960  -  Towradgi to Stanwell Park  -  John Hill Award (Tiana McGrath)  -  Residential + Commercial + PM + Business Sales",
   Inches(0.65), Inches(4.28), Inches(12), Inches(0.45), size=12, color=LGREY, align=PP_ALIGN.CENTER)

tb(s2, "BACKGROUND THAT SETS HIM APART", Inches(0.55), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=GOLD)
bullets(s2, ["30+ years in real estate since 1992 - across multiple Illawarra cycles",
             "Construction industry background - reads property fundamentals, not just listings",
             "Director and Licensee - runs MMJ North as both operator and rainmaker"],
        Inches(0.55), Inches(5.35), Inches(5.8), Inches(1.5), size=13, color=LGREY)

tb(s2, "WHAT HE OPERATES", Inches(6.8), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=GOLD)
bullets(s2, ["Coverage: Towradgi to Stanwell Park - whole Illawarra coastline north of Wollongong",
             "Multi-service: residential sales, commercial, property management, business sales",
             "11 active sale listings + 15 recently sold - active market position"],
        Inches(6.8), Inches(5.35), Inches(5.8), Inches(1.5), size=13, color=LGREY)


# SLIDE 3 - WHAT'S WORKING
s3 = prs.slides.add_slide(BLANK)
bg(s3)
rect(s3, 0, 0, Inches(0.12), H, GREEN)
tb(s3, "WHAT'S WORKING", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s3, "Strengths to Build On",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s3, Inches(1.35), color=GREEN)

strengths = [
    ("64-Year MMJ Track Record - Australian RE Longevity",
     "MMJ has operated since 1960. Most agencies that started then are gone, merged, or absorbed by national portals. MMJ's 64-year independence in Australian real estate is the kind of proof point a national chain can't manufacture and a startup can't buy. The current site doesn't lead with it - it should."),
    ("6 REINSW Awards Finalists in 2025 + John Hill Award",
     "Six MMJ team members named REINSW finalists in 2025. Tiana McGrath won the John Hill Award for outstanding contribution as chairperson on the residential property management chapter committee. That's the kind of industry recognition agencies build entire campaigns around. At MMJ today it appears once and isn't surfaced where it would close trust gaps pre-call."),
    ("Greg's 30+ Years + Construction Background Is Genuinely Unusual",
     "30 years across multiple Illawarra cycles is rare. The construction industry background is rarer still. Most agents read a listing. Greg reads the structure, the build, the renovation potential, the investment angle. That's a credibility differentiator with vendors and investors that the agent bio page mentions but doesn't fully exploit."),
    ("Geographic Authority - Towradgi to Stanwell Park",
     "MMJ North covers the entire Illawarra coastline north of Wollongong: Corrimal, Woonona, Tarrawanna, Bulli, Stanwell Park. Deep, narrow geographic focus is the one thing national portals like Domain and realestate.com.au can't replicate. Built into suburb-level office pages, that authority is the strongest organic moat MMJ has."),
    ("Multi-Service Operation Beyond Residential",
     "Most agencies in this size band do residential sales and PM and stop there. MMJ North runs commercial sales, property management, and business sales - a four-line operation. That's a referral and cross-sell engine that the website currently doesn't surface as one operation. Each service feels like its own silo on the site."),
]
for i, (title, desc) in enumerate(strengths):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s3, Inches(0.55), t, Inches(0.06), Inches(0.7), GREEN)
    tb(s3, title, Inches(0.75), t, Inches(11.8), Inches(0.38), size=14, bold=True, color=WHITE)
    tb(s3, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# SLIDE 4 - WHERE LEADS ARE LEAKING
s4 = prs.slides.add_slide(BLANK)
bg(s4)
rect(s4, 0, 0, Inches(0.12), H, RED)
tb(s4, "WHERE LEADS ARE LEAKING", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=RED)
tb(s4, "The Vendor Searching at 9pm Doesn't Become a Lead.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=28, bold=True, color=WHITE)
aline(s4, Inches(1.35), color=RED)

issues = [
    ("NO INSTANT PROPERTY VALUATION TOOL", "HIGH",
     "The single highest-converting lead magnet in residential real estate is missing. A vendor in Woonona thinking about listing in spring punches their address into Domain or realestate.com.au, gets an instant range, and the agency that captured that data is the agency they hear from. MMJ has the comparable sold data and the local expertise. The site doesn't ask for the lead."),
    ("OFFICE PAGE IS A PAGINATED AGENT ROSTER", "HIGH",
     "30+ agents across 6 paginated pages. No agent matching tool, no recent sold data by suburb, no neighborhood market data. A buyer trying to find the right agent for Tarrawanna versus Corrimal versus Bulli does the matching themselves. Most don't. They click the next agency that does the matching for them."),
    ("AWARDS AND CREDIBILITY NOT SURFACED ON PRIMARY PAGES", "HIGH",
     "REINSW finalist nominations, John Hill Award, 64 years of MMJ - these are the trust signals that close pre-call. They don't appear on the homepage above the fold, on office pages, or on agent bios in a structured way. The credibility exists. The pages where buyers and vendors decide whether to call don't show it."),
    ("NO AI CHAT FOR AFTER-HOURS QUESTIONS", "MEDIUM",
     "Vendor and buyer questions arrive when offices are closed. No appraisal availability, no inspection times, no rental application status, no maintenance request triage. Every after-hours visitor either calls back the next morning or moves on. AI chat trained on MMJ's process closes that gap without staffing the night shift."),
    ("NO PAST-CLIENT RE-ENGAGEMENT AUTOMATION", "MEDIUM",
     "The buyer who closed in 2018 is the appraisal request in 2026. Australians on average sell every 6-9 years. Without anniversary sequences, market reports, or re-engagement flows, the past-client list is dormant. Every settled sale is a future opportunity that goes to whoever is in the inbox at the right moment."),
]
for i, (title, pri, desc) in enumerate(issues):
    t = Inches(1.7) + i * Inches(1.05)
    pc = RED if pri == "HIGH" else GOLD
    rect(s4, Inches(0.55), t, Inches(0.06), Inches(0.7), pc)
    tb(s4, title, Inches(0.75), t, Inches(9.2), Inches(0.38), size=12, bold=True, color=WHITE)
    tag(s4, pri, Inches(10.1), t + Inches(0.05), fill=pc, size=9)
    tb(s4, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# SLIDE 5 - WEBSITE DESIGN AUDIT
s5 = prs.slides.add_slide(BLANK)
bg(s5)
rect(s5, 0, 0, Inches(0.12), H, PURPLE)
tb(s5, "WEBSITE DESIGN AUDIT", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=PURPLE)
tb(s5, "64 Years of MMJ. The Site Doesn't Lead With It.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.8), size=32, bold=True, color=WHITE)
aline(s5, Inches(1.45), color=PURPLE, w=Inches(1.5))

rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(3.6), DARK_BLUE)
rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(0.42), RED)
tb(s5, "MMJ NORTH OFFICE PAGE TODAY", Inches(0.65), Inches(1.92), Inches(3.6), Inches(0.38),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "No instant valuation tool anywhere",
    "30+ agents paginated 6 deep, no matching",
    "Office page = team roster, not suburb authority",
    "No recent sold data by suburb",
    "No AI chat / after-hours self-service",
    "No tenant portal for PM clients",
    "Awards/credibility not on primary pages",
    "Only 4 testimonials surfaced on homepage",
], Inches(0.7), Inches(2.42), Inches(3.6), Inches(2.9), size=11, color=LGREY, dot=RED)

tb(s5, "->", Inches(4.45), Inches(3.6), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(3.6), DARK_BLUE)
rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(0.42), GOLD)
tb(s5, "WHAT TOP AU AGENCY SITES DO", Inches(5.1), Inches(1.92), Inches(3.4), Inches(0.38),
   size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Instant valuation on homepage as primary CTA",
    "Office pages = suburb authorities, not rosters",
    "Agent matching by suburb and property type",
    "AI chat for vendor / buyer / tenant 24/7",
    "Tenant portal: applications, maintenance, rent",
    "Awards wall and verified review counts",
    "Past-client anniversary nurture sequences",
    "Vendor weekly suburb market reports",
], Inches(5.1), Inches(2.42), Inches(3.4), Inches(2.9), size=11, color=LGREY, dot=GOLD)

tb(s5, "->", Inches(8.68), Inches(3.6), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(3.6), DARK_BLUE)
rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(0.42), GREEN)
tb(s5, "WHAT GENOS AI BUILDS", Inches(9.35), Inches(1.92), Inches(3.4), Inches(0.38),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Instant valuation tool on homepage + office",
    "Office pages rebuilt as suburb authorities",
    "Agent matching by suburb and property type",
    "AI chat: appraisals, inspections, PM 24/7",
    "Tenant portal: applications, maintenance, rent",
    "Awards wall on every primary page",
    "Past-client anniversary sequences (1/3/5/7yr)",
    "Vendor weekly suburb-specific market reports",
], Inches(9.35), Inches(2.42), Inches(3.4), Inches(2.9), size=11, color=LGREY, dot=GREEN)

rect(s5, Inches(0.55), Inches(5.72), Inches(12.23), Inches(1.28), DARK_BLUE)
tb(s5, "WHY IT MATTERS:", Inches(0.7), Inches(5.78), Inches(2.2), Inches(0.38),
   size=11, bold=True, color=PURPLE)
stats2 = [
    ("3-5x", "vendor lead capture\nvs. plain contact form"),
    ("60-70%", "after-hours questions\nresolvable by AI chat"),
    ("6-9 yrs", "average AU resale window\npast-client = future vendor"),
    ("64 yrs", "credibility doing no work\nwhere buyers decide to call"),
]
sw = Inches(2.7); sl = Inches(2.9)
for i, (num, label) in enumerate(stats2):
    xl = sl + i * sw
    tb(s5, num, xl, Inches(5.75), Inches(2.5), Inches(0.55),
       size=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s5, label, xl, Inches(6.22), Inches(2.5), Inches(0.7),
       size=10, color=LGREY, align=PP_ALIGN.CENTER)


# SLIDE 6 - AI AUTOMATION OPPORTUNITIES
s6 = prs.slides.add_slide(BLANK)
bg(s6)
rect(s6, 0, 0, Inches(0.12), H, ACCENT)
tb(s6, "AI AUTOMATION OPPORTUNITIES", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s6, "Capturing the Vendor Searching at 9pm.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=34, bold=True, color=WHITE)
aline(s6, Inches(1.35), color=ACCENT)

opps = [
    ("Instant Valuation\n+ Suburb Routing", GOLD,
     "Vendor enters address, gets instant comparable-based range, optionally requests a refined appraisal. Lead routes automatically to the agent specialising in that suburb (Greg for commercial / development, suburb specialist for residential). The single highest-converting lead magnet in AU residential real estate, missing today."),
    ("AI Chat\n24/7 Self-Service", ACCENT,
     "Trained on MMJ's appraisal process, current listings, suburb market data, and PM workflows. Handles vendor, buyer, and tenant questions after hours. Books appraisals and inspections. Triages PM maintenance requests. Routes complex deals to the right agent with conversation context attached."),
    ("Vendor Nurture\n+ Market Reports", GREEN,
     "Weekly suburb-specific market reports automated to vendor prospects in the 6-12 month window before listing. Different content for someone who valued their Bulli home vs. Greg's Tarrawanna investment audience. Stays in front of vendors until they're ready, then captures the appraisal request when they are."),
    ("Past-Client\nRe-Engagement", RED,
     "Buyer from 2018, 2020, 2022 receives anniversary sequences with current market value of their home, suburb sales activity, and a soft appraisal CTA. Australians sell every 6-9 years on average. Past-client list activated as a recurring vendor pipeline rather than a dormant CRM table."),
    ("Tenant Portal\n+ PM Automation", PURPLE,
     "Online rental applications, maintenance request triage with photo upload, lease renewal sequences starting 90 days out, behavioural rent reminders. PM team handles exceptions, not routine work. Reduces tenant churn at 115+ studios... applies the same way to MMJ's PM book."),
]
cw2 = Inches(2.35); gap2 = Inches(0.14)
for i, (title, color, desc) in enumerate(opps):
    l2 = Inches(0.55) + i * (cw2 + gap2)
    rect(s6, l2, Inches(1.8), cw2, Inches(0.06), color)
    rect(s6, l2, Inches(1.86), cw2, Inches(4.6), DARK_BLUE)
    tb(s6, title, l2 + Inches(0.12), Inches(1.98), cw2 - Inches(0.24), Inches(0.65),
       size=13, bold=True, color=color)
    tb(s6, desc, l2 + Inches(0.12), Inches(2.65), cw2 - Inches(0.24), Inches(3.6),
       size=11, color=LGREY)

rect(s6, Inches(0.55), Inches(6.7), Inches(12.23), Inches(0.45), DARK_BLUE)
tb(s6, "MMJ - More Than Just Real Estate. The site should do more than list agents and listings.",
   Inches(0.7), Inches(6.75), Inches(12), Inches(0.35), size=11, italic=True, color=GOLD)


# SLIDE 7 - WHAT THIS IS COSTING YOU
s7 = prs.slides.add_slide(BLANK)
bg(s7)
rect(s7, 0, 0, Inches(0.12), H, RED)
tb(s7, "WHAT THIS IS COSTING YOU", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=RED)
tb(s7, "Domain and realestate.com.au Capture What You Don't.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=30, bold=True, color=WHITE)
aline(s7, Inches(1.35), color=RED)

costs = [
    ("No instant valuation = vendors hand their data to portals, not to MMJ",
     "Domain and realestate.com.au capture every Illawarra vendor punching in their address for an instant value range, and they monetise that data back to whoever pays for leads. MMJ has the same comparable sold data and 30+ years of geographic specialty. The site just doesn't ask. Every vendor who values online without filling out an MMJ form is a future listing that may not be MMJ's."),
    ("Office page as a roster = buyers self-match, often to the wrong agent",
     "30+ agents paginated 6 deep with no matching tool means a buyer who wants Tarrawanna gets the same experience as one who wants Stanwell Park. Some pick the wrong agent. Some pick a competitor's site that does match them. The agent-suburb mapping exists in your team's heads. The website needs to expose it."),
    ("Awards not on primary pages = trust signals doing no work pre-call",
     "Six REINSW finalists in 2025. John Hill Award. 64 years of MMJ. These are the credentials that close trust gaps before the first conversation. Buried below the fold or on a single about page, they reach almost no one. Surfaced on the homepage and every office page, they shorten every sales cycle."),
    ("No past-client re-engagement = the 2018 buyer becomes someone else's 2026 vendor",
     "Australians sell every 6-9 years on average. The buyer Greg or any MMJ agent settled in 2018 is starting to think about it now. Without an anniversary sequence with current home value and suburb activity, they don't hear from MMJ - they hear from the agency that does. Every settled sale is a future appraisal opportunity that needs to be programmed."),
    ("No AI chat or PM portal = staff time on questions automation handles",
     "Maintenance requests, balance and lease questions, inspection availability, application status. PM teams field these all day. Most are answered the same way every time. AI chat trained on MMJ's process and a tenant portal with self-service handles 60-70% of routine volume, and frees the team for actual exception management."),
]
for i, (title, desc) in enumerate(costs):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s7, Inches(0.55), t, Inches(0.06), Inches(0.7), RED)
    tb(s7, title, Inches(0.75), t, Inches(11.8), Inches(0.38), size=14, bold=True, color=WHITE)
    tb(s7, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# SLIDE 8 - WHAT WE'D BUILD
s8 = prs.slides.add_slide(BLANK)
bg(s8)
rect(s8, 0, 0, Inches(0.12), H, GREEN)
tb(s8, "WHAT WE'D BUILD", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s8, "Genos AI for MMJ North",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s8, Inches(1.35), color=GREEN)

builds = [
    ("Instant Property Valuation + Suburb-Routed Lead Capture",
     "Embedded on the homepage and the North Corrimal office page. Vendor enters address, gets instant range based on MMJ's sold data, optionally books a refined appraisal. Lead routes automatically to the agent specialising in that suburb with the calculator inputs already attached to the CRM record."),
    ("AI Chat - 24/7 Vendor / Buyer / Tenant Self-Service",
     "Trained on MMJ's appraisal process, current listings, suburb market data, PM workflows, and tenant procedures. Handles after-hours questions, books appraisals and inspections, triages maintenance requests with photo upload. Routes complex inquiries to the right agent with full conversation context."),
    ("Office Page Rebuild as Suburb Authority",
     "North Corrimal page rebuilt around the suburbs MMJ actually owns: Corrimal, Woonona, Tarrawanna, Bulli, Stanwell Park. Recent sold data, agent-suburb matching, neighborhood data, market reports. The page becomes the suburb authority instead of a paginated team roster, and ranks for suburb searches."),
    ("Past-Client Re-Engagement at 1, 3, 5, 7 Years",
     "Buyer settled in 2024 enters an automated 1-year anniversary sequence with current home value and suburb activity, then 3-year, 5-year, 7-year sequences with appraisal CTA. Past-client list activated as a recurring vendor pipeline. Every settled sale becomes a future opportunity programmed."),
    ("Tenant Portal + PM Automation + Vendor Nurture",
     "Online rental applications, maintenance triage, lease renewal sequences 90 days out, behavioural rent reminders. Plus weekly suburb-specific vendor market reports for the 6-12 month pre-listing window. Plus awards wall on homepage and office pages: 64 years, 6 REINSW finalists, John Hill Award."),
]
for i, (title, desc) in enumerate(builds):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s8, Inches(0.55), t, Inches(0.06), Inches(0.7), GREEN)
    tb(s8, title, Inches(0.75), t, Inches(11.8), Inches(0.38), size=14, bold=True, color=WHITE)
    tb(s8, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# SLIDE 9 - 90-DAY OUTCOMES
s9 = prs.slides.add_slide(BLANK)
bg(s9)
rect(s9, 0, 0, Inches(0.12), H, ACCENT)
tb(s9, "90-DAY OUTCOMES", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s9, "What Changes in 90 Days",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s9, Inches(1.35))

phases = [
    ("Month 1", GOLD, [
        "Instant valuation tool live on mmj.com.au and North Corrimal office page",
        "AI chat deployed for vendor, buyer, and tenant after-hours questions",
        "North Corrimal office page rebuilt with recent sold data and agent matching",
        "Awards wall on homepage and office page: 64 years, REINSW, John Hill Award",
    ]),
    ("Month 2", ACCENT, [
        "Vendor nurture sequences running with weekly suburb-specific market reports",
        "Tenant portal live - online applications, maintenance triage, lease renewals",
        "Past-client re-engagement sequences active for 1, 3, 5, 7-year anniversaries",
        "First measurable inbound appraisal requests from instant valuation tool",
    ]),
    ("Month 3", GREEN, [
        "Office pages acting as suburb authorities for Corrimal, Woonona, Tarrawanna, Bulli, Stanwell Park",
        "Buyer alerts segmented by suburb and property type running",
        "PM automation deflecting routine tenant questions, freeing team for exceptions",
        "Vendor pipeline pre-listing measurable, past-client appraisal requests landing",
    ]),
]
pw = Inches(3.9); gap3 = Inches(0.18)
for i, (month, color, pts) in enumerate(phases):
    l3 = Inches(0.55) + i * (pw + gap3)
    rect(s9, l3, Inches(1.8), pw, Inches(0.42), color)
    tb(s9, month.upper(), l3 + Inches(0.1), Inches(1.84), pw - Inches(0.2), Inches(0.36),
       size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    rect(s9, l3, Inches(2.22), pw, Inches(4.7), DARK_BLUE)
    bullets(s9, pts, l3 + Inches(0.15), Inches(2.35), pw - Inches(0.3), Inches(4.4),
            size=12, color=LGREY, dot=color)


# SLIDE 10 - NEXT STEPS
s10 = prs.slides.add_slide(BLANK)
bg(s10)
rect(s10, 0, 0, Inches(0.12), H, GOLD)
rect(s10, Inches(7.5), 0, Inches(5.83), H, DARK_BLUE)
tb(s10, "NEXT STEPS", Inches(0.55), Inches(0.3), Inches(6.5), Inches(0.5),
   size=11, bold=True, color=GOLD)
tb(s10, "Let's Talk for\n15 Minutes.",
   Inches(0.55), Inches(0.8), Inches(6.8), Inches(2.0), size=48, bold=True, color=WHITE)
aline(s10, Inches(2.8), color=GOLD, w=Inches(2.5))
tb(s10, "No pitch deck required. Reply or book a call directly.\nWe'll walk through what fits MMJ North Corrimal.",
   Inches(0.55), Inches(3.05), Inches(6.5), Inches(0.8), size=14, color=LGREY)

steps = [
    "Reply to this email - we'll set up a time",
    "Book directly: www.genosai.tech/call",
    "Walk through the audit findings in 15 minutes",
    "You decide what, if anything, makes sense to move on",
]
bullets(s10, steps, Inches(0.55), Inches(3.95), Inches(6.5), Inches(2.0), size=14, color=WHITE, dot=GOLD)

rect(s10, 0, Inches(6.5), Inches(7.5), Inches(1.0), DARK_BLUE)
tb(s10, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699  |  www.genosai.tech",
   Inches(0.55), Inches(6.6), Inches(6.8), Inches(0.5), size=12, color=MGREY)

tb(s10, "WHAT WE BUILD FOR\nMMJ NORTH", Inches(7.85), Inches(0.6), Inches(5.0), Inches(0.9),
   size=18, bold=True, color=GOLD)
aline(s10, Inches(1.45), color=GOLD)

summary_items = [
    ("Valuation", "Instant tool with suburb-routed agent follow-up"),
    ("AI Chat", "24/7 vendor / buyer / tenant self-service"),
    ("Office Pages", "Suburb authorities, not just team rosters"),
    ("Past Clients", "1/3/5/7 yr anniversary sequences = appraisal pipeline"),
    ("PM Portal", "Applications, maintenance, lease renewals automated"),
    ("Awards Wall", "64 yrs, REINSW, John Hill Award - on every primary page"),
]
for i, (label, val) in enumerate(summary_items):
    t2 = Inches(1.6) + i * Inches(0.88)
    rect(s10, Inches(7.85), t2, Inches(1.5), Inches(0.72), NAVY)
    tb(s10, label, Inches(7.9), t2 + Inches(0.05), Inches(1.4), Inches(0.32),
       size=10, bold=True, color=GOLD)
    tb(s10, val, Inches(9.5), t2 + Inches(0.05), Inches(3.6), Inches(0.65),
       size=11, color=LGREY)


prs.save(ROOT / "output" / "decks" / "MMJNorth_Audit_GenosAI.pptx")
print("Saved: MMJNorth_Audit_GenosAI.pptx")
