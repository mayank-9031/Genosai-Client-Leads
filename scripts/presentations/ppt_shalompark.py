import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.parent / "utils"))
from ppt_theme import (
    NAVY, DARK_BLUE, ACCENT, GOLD, RED, GREEN, PURPLE, ORANGE, TEAL, CYAN,
    WHITE, LGREY, MGREY, W, H,
    bg, rect, tb, aline, bullets, tag,
    slide_header, slide_footer, score_badge,
    stat_cards, top_bar_card, phase_card, before_after, genos_logo_block,
)

ROOT = Path(__file__).parent.parent.parent
prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]
# Client-specific colors
EMERALD = RGBColor(0x00, 0x92, 0x62)
DGREY   = RGBColor(0x2C, 0x38, 0x4A)
AMBER   = RGBColor(0xF3, 0x9C, 0x12)
WARM    = RGBColor(0xFD, 0xF0, 0xE2)
SLATE   = RGBColor(0x2E, 0x45, 0x5F)


# ── Local helpers ────────────────────────────────────────────────────────────

def header(slide, title, subtitle, accent=GREEN):
    rect(slide, 0, 0, Inches(0.12), H, accent)
    tb(slide, title.upper(), Inches(0.55 + 0.12), Inches(0.3),
       Inches(12), Inches(0.45), size=11, bold=True, color=accent)
    tb(slide, subtitle, Inches(0.55 + 0.12), Inches(0.72),
       Inches(12 - 0.12), Inches(0.65), size=22, bold=True, color=WHITE)
    aline(slide, Inches(1.38), color=accent, w=Inches(1.5))


def footer(slide,
           text='Genos AI  |  hello@genosai.tech  |  www.genosai.tech  |  +91 638-714-2699'):
    rect(slide, 0, Inches(7.05), W, Inches(0.45), DARK_BLUE)
    tb(slide, text, Inches(0.55), Inches(7.1), Inches(12.4), Inches(0.35),
       size=9, color=MGREY, align=PP_ALIGN.CENTER)


def card(slide, x, y, w, h, fill=DARK_BLUE):
    return rect(slide, x, y, w, h, fill)


def stat_card(slide, val, label, x, y, w=Inches(3.0), h=Inches(1.6), accent=GREEN):
    rect(slide, x, y, w, h, DARK_BLUE)
    rect(slide, x, y, w, Inches(0.05), accent)
    tb(slide, str(val), x + Inches(0.1), y + Inches(0.1),
       w - Inches(0.2), Inches(0.9), size=28, bold=True, color=accent,
       align=PP_ALIGN.CENTER)
    tb(slide, label, x + Inches(0.1), y + Inches(1.1),
       w - Inches(0.2), Inches(0.42), size=10, color=LGREY, align=PP_ALIGN.CENTER)


def bullet_card(slide, title, items, x, y, w, h,
                accent=GREEN, title_bg=None):
    if title_bg is None:
        title_bg = DARK_BLUE
    rect(slide, x, y, w, h, DARK_BLUE)
    rect(slide, x, y, w, Inches(0.42), title_bg)
    rect(slide, x, y, Inches(0.05), h, accent)
    tb(slide, title, x + Inches(0.1), y + Inches(0.05),
       w - Inches(0.15), Inches(0.35), size=11, bold=True, color=WHITE)
    bx = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.48),
                                  w - Inches(0.15), h - Inches(0.55))
    tf = bx.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r0 = p.add_run(); r0.text = '  '; r0.font.size = Pt(9)
        r1 = p.add_run(); r1.text = '● '; r1.font.size = Pt(9)
        r1.font.color.rgb = accent; r1.font.bold = True
        r2 = p.add_run(); r2.text = str(item)
        r2.font.size = Pt(10.5); r2.font.color.rgb = LGREY
        if i < len(items) - 1:
            p.space_after = Pt(5)




# SLIDE 1 – TITLE
# ════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
bg(s1, NAVY)
rect(s1, 0, 0, Inches(0.7), H, GOLD)
rect(s1, Inches(0.7), Inches(4.5), W - Inches(0.7), Inches(0.06), GOLD)

tb(s1, "SHALOM PARK ESTATE  |  IFT REALTY LTD",
   Inches(1.1), Inches(1.0), Inches(11), Inches(0.7),
   size=15, bold=True, color=GOLD)
tb(s1, "Web Development &\nAI Automation Audit",
   Inches(1.1), Inches(1.7), Inches(11), Inches(2.0),
   size=44, bold=True, color=WHITE)
tb(s1, "Strategic Review & Digital Growth Roadmap  |  Lekki-Epe, Lagos, Nigeria",
   Inches(1.1), Inches(3.72), Inches(10), Inches(0.5),
   size=16, color=GOLD, italic=True)
tb(s1, "Prepared by: Genos Apollo  |  May 2026",
   Inches(1.1), Inches(5.0), Inches(8), Inches(0.4),
   size=12, color=MGREY)
tb(s1, "CONFIDENTIAL",
   Inches(1.1), Inches(6.6), Inches(4), Inches(0.35),
   size=10, color=MGREY, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 2 – AGENDA
# ════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
bg(s2)
header(s2, "Agenda", "What we will cover today")
footer(s2)

agenda = [
    ("01", "Company Digital Overview",    "Current state of shalomparknigeria.com & digital footprint"),
    ("02", "Website Audit Findings",      "WordPress gaps: pricing, UX, SEO and compliance"),
    ("03", "Web Development Roadmap",     "Buyer portal, diaspora features & modern stack"),
    ("04", "AI Automation Opportunity",   "WhatsApp, investor comms & realtor programme AI"),
    ("05", "Benefits of Automation",      "ROI, lead velocity and diaspora sales advantage"),
    ("06", "Implementation Timeline",     "Phased rollout plan across 12 months"),
    ("07", "Investment & Next Steps",     "Budget guidance and immediate actions"),
]

for i, (num, title, desc) in enumerate(agenda):
    col = i % 2
    row = i // 2
    x = Inches(0.55) + col * Inches(6.4)
    y = Inches(1.55) + row * Inches(1.22)
    card(s2, x, y, Inches(6.1), Inches(1.05))
    rect(s2, x, y, Inches(0.7), Inches(1.05), NAVY)
    tb(s2, num, x, y + Inches(0.25), Inches(0.7), Inches(0.5),
       size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s2, title, x + Inches(0.8), y + Inches(0.1),
       Inches(5.0), Inches(0.4), size=13, bold=True, color=NAVY)
    tb(s2, desc, x + Inches(0.8), y + Inches(0.52),
       Inches(5.0), Inches(0.42), size=10.5, color=MGREY, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 3 – COMPANY DIGITAL OVERVIEW
# ════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
bg(s3)
header(s3, "Company Digital Overview", "Shalom Park Estate – current digital footprint")
footer(s3)

card(s3, Inches(0.4), Inches(1.4), Inches(5.8), Inches(5.5))
rect(s3, Inches(0.4), Inches(1.4), Inches(5.8), Inches(0.45), NAVY)
tb(s3, "Company Snapshot", Inches(0.55), Inches(1.44),
   Inches(5.5), Inches(0.38), size=13, bold=True, color=WHITE)

snapshot = [
    ("Developer",    "IFT Realty Ltd  (subsidiary of IFT Realty Dev. LLC, New York)"),
    ("CEO",          "Oluremi Oshikanlu"),
    ("Location",     "Abijo, Ibeju Lekki, Lagos State — Lekki-Epe axis"),
    ("Estate Size",  "25 acres — gated luxury community"),
    ("Units",        "2-bed condos · 4-bed · 5-bed semi/fully detached · serviced plots"),
    ("Legal",        "Lagos State Governor's Consent · Approved layout · Megacity Plan"),
    ("Website",      "www.shalomparknigeria.com  (WordPress CMS)"),
    ("Social",       "Instagram · Facebook · LinkedIn · Twitter/X · YouTube"),
]
ty = Inches(2.02)
for k, v in snapshot:
    tb(s3, k, Inches(0.6), ty, Inches(1.55), Inches(0.3),
       size=10.5, bold=True, color=NAVY)
    tb(s3, v, Inches(2.2), ty, Inches(3.8), Inches(0.3),
       size=10.5, color=MGREY)
    ty += Inches(0.52)

card(s3, Inches(6.5), Inches(1.4), Inches(6.5), Inches(5.5))
rect(s3, Inches(6.5), Inches(1.4), Inches(6.5), Inches(0.45), GOLD)
tb(s3, "Key Digital Health Indicators", Inches(6.65), Inches(1.44),
   Inches(6.2), Inches(0.38), size=13, bold=True, color=NAVY)

metrics = [
    ("Custom Domain",             "✓ Active",   "shalomparknigeria.com live",            GREEN),
    ("Mobile Page Speed",         "~48 / 100",  "WordPress + image-heavy — est. slow",   RED),
    ("SEO Score",                 "~60 / 100",  "Blog exists; diaspora keywords weak",   AMBER),
    ("Property Pricing on Site",  "None",       "No prices shown — kills conversions",   RED),
    ("Buyer / Investor Portal",   "None",       "No login or payment tracking system",   RED),
    ("Realtor Portal",            "Manual",     "Google Meet webinar — not scalable",    AMBER),
    ("WhatsApp Integration",      "Unknown",    "No visible click-to-chat on listings",  AMBER),
    ("NDPR Compliance",           "Partial",    "No visible cookie consent or policy",   AMBER),
]
ty2 = Inches(2.02)
for label, val, note, col in metrics:
    tb(s3, label, Inches(6.65), ty2, Inches(2.7), Inches(0.3),
       size=10.5, bold=True, color=NAVY)
    tb(s3, val,   Inches(9.4), ty2, Inches(1.25), Inches(0.3),
       size=10.5, bold=True, color=col, align=PP_ALIGN.CENTER)
    tb(s3, note,  Inches(10.7), ty2, Inches(2.1), Inches(0.3),
       size=9, color=MGREY, italic=True)
    ty2 += Inches(0.52)

# ════════════════════════════════════════════════════════════════
# SLIDE 4 – WEBSITE AUDIT FINDINGS
# ════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
bg(s4)
header(s4, "Website Audit Findings", "shalomparknigeria.com – gaps limiting sales and investor conversions")
footer(s4)

audit_cards = [
    ("Performance",
     ["WordPress with plugins — mobile load est. >4s", "Image-heavy property galleries not optimised",
      "No CDN confirmed for Nigerian edge delivery", "No lazy-loading detected on property photos"]),
    ("UX & Conversion",
     ["Zero pricing on all property pages (5 listings)", "Google Meet for realtors — low trust for luxury brand",
      "No virtual tour or 3D walkthrough embed", "No WhatsApp click-to-chat on property listings"]),
    ("SEO & Content",
     ["Blog exists but diaspora keyword gaps (US/UK)", "No LocalBusiness or RealEstateListing schema",
      "MREIF fund has no dedicated SEO landing page", "Weak backlink profile vs. Lekki competitors"]),
    ("Trust & Compliance",
     ["No NDPR cookie consent or privacy policy visible", "No buyer portal or subscription tracker",
      "Construction progress not visible to subscribers", "No USD pricing despite NY parent company focus"]),
]

for i, (title, bullets) in enumerate(audit_cards):
    x = Inches(0.4) + i * Inches(3.2)
    y = Inches(1.45)
    card(s4, x, y, Inches(3.05), Inches(5.3))
    rect(s4, x, y, Inches(3.05), Inches(0.45), NAVY)
    tb(s4, title, x + Inches(0.1), y + Inches(0.06),
       Inches(2.85), Inches(0.36), size=12, bold=True, color=WHITE)
    ty = y + Inches(0.62)
    for b in bullets:
        rect(s4, x + Inches(0.15), ty + Inches(0.07), Inches(0.1), Inches(0.1), GOLD)
        tb(s4, b, x + Inches(0.35), ty, Inches(2.6), Inches(0.55),
           size=10, color=NAVY)
        ty += Inches(0.72)

rect(s4, Inches(0.4), Inches(6.8), Inches(12.5), Inches(0.06), GOLD)
tb(s4, "Immediate Priority: Add pricing + WhatsApp CTAs + buyer portal — these three changes directly unlock stalled sales conversions",
   Inches(0.4), Inches(6.86), Inches(12.5), Inches(0.3),
   size=10.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 5 – WEB DEVELOPMENT ROADMAP
# ════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
bg(s5)
header(s5, "Web Development Roadmap", "From brochure website to a full buyer & investor platform")
footer(s5)

card(s5, Inches(0.4), Inches(1.4), Inches(4.5), Inches(5.6))
rect(s5, Inches(0.4), Inches(1.4), Inches(4.5), Inches(0.45), EMERALD)
tb(s5, "Recommended Tech Stack", Inches(0.55), Inches(1.44),
   Inches(4.2), Inches(0.38), size=13, bold=True, color=WHITE)

stack = [
    ("Frontend",     "Next.js 15 – SSR/SSG (replace WordPress)"),
    ("CMS",          "Sanity.io – listings, blog, MREIF content"),
    ("Hosting",      "Vercel + Cloudflare CDN (Nigeria + US edge)"),
    ("Payments",     "Paystack (NGN) + Stripe (USD/diaspora)"),
    ("Portal",       "Auth0 – buyer & realtor secure dashboards"),
    ("WhatsApp",     "WhatsApp Business API – enquiry & updates"),
    ("Analytics",    "GA4 + Hotjar — conversion funnel tracking"),
    ("Email CRM",    "Mailchimp / Brevo — diaspora nurture flows"),
]
ty = Inches(2.02)
for k, v in stack:
    tb(s5, k, Inches(0.6), ty, Inches(1.4), Inches(0.3),
       size=10, bold=True, color=EMERALD)
    tb(s5, v, Inches(2.05), ty, Inches(2.75), Inches(0.3),
       size=10, color=NAVY)
    ty += Inches(0.5)

card(s5, Inches(5.2), Inches(1.4), Inches(7.75), Inches(5.6))
rect(s5, Inches(5.2), Inches(1.4), Inches(7.75), Inches(0.45), NAVY)
tb(s5, "Key Improvements & Features",
   Inches(5.35), Inches(1.44), Inches(7.4), Inches(0.38),
   size=13, bold=True, color=WHITE)

improvements = [
    ("Transparent Pricing (NGN + USD)",
     "Display all unit prices in both currencies on every property page. Diaspora buyers transact in USD — hiding prices loses them immediately."),
    ("Buyer & Subscriber Portal",
     "Secure login: payment schedule tracker, construction milestone photos, document downloads, next-payment reminders — no more manual email chasing."),
    ("Realtor Partner Dashboard",
     "Replace Google Meet with a portal: realtor onboarding, lead assignment, commission tracker, marketing collateral downloads, and referral leaderboard."),
    ("MREIF Investment Landing Page",
     "Dedicated SEO-optimised page: fund structure, projected returns (NGN + USD), minimum investment, application form, and MREIF FAQs."),
    ("Virtual Tour & 3D Walkthrough Embed",
     "Matterport or YouTube 3D embed on each unit page. Diaspora buyers sign contracts without visiting — virtual tours are non-negotiable for remote sales."),
    ("NDPR Compliance & Cookie Consent",
     "Nigeria Data Protection Regulation compliant consent banner, privacy policy, data request form, and secure storage for buyer personal data."),
]
ty2 = Inches(1.98)
for title, desc in improvements:
    rect(s5, Inches(5.35), ty2 + Inches(0.1), Inches(0.08), Inches(0.25), GOLD)
    tb(s5, title, Inches(5.55), ty2, Inches(7.1), Inches(0.28),
       size=11, bold=True, color=NAVY)
    tb(s5, desc, Inches(5.55), ty2 + Inches(0.3), Inches(7.1), Inches(0.35),
       size=10, color=MGREY)
    ty2 += Inches(0.78)

# ════════════════════════════════════════════════════════════════
# SLIDE 6 – AI AUTOMATION OPPORTUNITY MAP
# ════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
bg(s6)
header(s6, "AI Automation Opportunity Map", "Where AI drives sales, trust and operational scale for Shalom Park")
footer(s6)

opps = [
    (NAVY,    "WhatsApp Lead & Sales Automation",
     ["AI chatbot handles property enquiries 24/7", "Auto-send unit brochures + pricing on request",
      "Qualify buyers: budget, timeline, unit preference", "Book site visits or virtual tours via WhatsApp"]),
    (GOLD,    "Buyer & Subscriber Communications",
     ["Monthly construction progress auto-updates (WA + email)", "Payment due-date reminders with Paystack link",
      "Personalised milestone messages (foundation, roofing)", "Post-purchase welcome sequence & handover pack"]),
    (EMERALD, "Realtor Programme Automation",
     ["AI-guided realtor onboarding (replace Google Meet)", "Auto-assign leads from web form to nearest realtor",
      "Commission calculation & payout notification", "Weekly leaderboard + performance nudges"]),
    (SLATE,   "Diaspora & Investor Marketing",
     ["AI-targeted email campaigns (US/UK Nigerian diaspora)", "Instagram Reels auto-captioned from site milestones",
      "MREIF fund: AI-generated quarterly investor reports", "Monitor Lekki competitors & alert on pricing shifts"]),
]

for i, (col, title, bullets) in enumerate(opps):
    c = i % 2
    r = i // 2
    x = Inches(0.4) + c * Inches(6.5)
    y = Inches(1.45) + r * Inches(2.85)
    card(s6, x, y, Inches(6.2), Inches(2.65))
    rect(s6, x, y, Inches(6.2), Inches(0.45), col)
    tb(s6, title, x + Inches(0.15), y + Inches(0.06),
       Inches(5.9), Inches(0.36), size=13, bold=True,
       color=NAVY if col == GOLD else WHITE)
    ty = y + Inches(0.62)
    for b in bullets:
        rect(s6, x + Inches(0.18), ty + Inches(0.08), Inches(0.1), Inches(0.1), col)
        tb(s6, b, x + Inches(0.38), ty, Inches(5.65), Inches(0.38),
           size=10.5, color=NAVY)
        ty += Inches(0.48)

# ════════════════════════════════════════════════════════════════
# SLIDE 7 – BENEFITS OF AI AUTOMATION
# ════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
bg(s7)
header(s7, "Benefits of AI Automation", "Measurable impact across Shalom Park's sales and operations")
footer(s7)

stats = [
    ("24/7",    "WhatsApp property\nenquiry coverage"),
    ("3×",      "Faster realtor\nonboarding & activation"),
    ("60%",     "Reduction in manual\nbuyer follow-up time"),
    ("40%",     "More diaspora leads\nconverted remotely"),
]
for i, (val, label) in enumerate(stats):
    stat_card(s7, val, label,
              Inches(0.4) + i * Inches(3.2), Inches(1.45),
              w=Inches(3.0), h=Inches(1.6))

benefits = [
    ("Sales Conversion",
     ["Pricing visible = buyers self-qualify 24/7", "WhatsApp bot books site visits while you sleep",
      "Virtual tour links close diaspora deals remotely", "Fewer unqualified enquiries wasting sales time"]),
    ("Buyer Trust & Retention",
     ["Construction updates build confidence in developer", "Payment reminders reduce subscription defaults",
      "Milestone photos shared automatically to buyers", "Post-sale care drives referrals and reviews"]),
    ("Realtor Programme Scale",
     ["100+ realtors managed via portal, not WhatsApp groups", "Auto-assigned leads = no dropped referrals",
      "Transparent commissions increase realtor loyalty", "Best performers identified and rewarded by data"]),
    ("Diaspora & Investor Edge",
     ["USD pricing + Stripe unlocks North America buyers", "AI investor reports build MREIF fund credibility",
      "Nigerian diaspora is the #1 Lekki buyer segment", "Automated comms close across time zones"]),
]

for i, (title, bullets) in enumerate(benefits):
    c = i % 2
    r = i // 2
    x = Inches(0.4) + c * Inches(6.5)
    y = Inches(3.2) + r * Inches(1.85)
    bullet_card(s7, title, bullets, x, y, Inches(6.2), Inches(1.65),
                title_bg=NAVY if r == 0 else EMERALD)

# ════════════════════════════════════════════════════════════════
# SLIDE 8 – ROI & BUSINESS CASE
# ════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
bg(s8)
header(s8, "ROI & Business Case", "Projected returns from web + AI investment (NGN / USD)")
footer(s8)

card(s8, Inches(0.4), Inches(1.4), Inches(5.8), Inches(5.6))
rect(s8, Inches(0.4), Inches(1.4), Inches(5.8), Inches(0.45), NAVY)
tb(s8, "Indicative Investment (Year 1 – USD)",
   Inches(0.55), Inches(1.44), Inches(5.5), Inches(0.38),
   size=13, bold=True, color=WHITE)

invest_rows = [
    ("Website Rebuild (Next.js + Sanity CMS)",    "$3,500 – $ 6,500"),
    ("Buyer & Realtor Portal (Auth0)",             "$2,000 – $ 4,000"),
    ("WhatsApp Business API + AI Bot",             "$  800 – $ 1,800"),
    ("Paystack + Stripe Payment Integration",      "$  600 – $ 1,200"),
    ("MREIF Landing Page + SEO Sprint",            "$  500 – $ 1,000"),
    ("NDPR Compliance & Cookie Consent",           "$  300 – $   600"),
    ("Ongoing monthly support (est.)",             "$  400 – $   700 / mo"),
    ("TOTAL YEAR 1 (once-off + 12 mo)",           "$12,500 – $23,900"),
]
ty = Inches(2.0)
for i, (item, cost) in enumerate(invest_rows):
    is_total = i == len(invest_rows) - 1
    if is_total:
        rect(s8, Inches(0.4), ty - Inches(0.05), Inches(5.8), Inches(0.06), GOLD)
        ty += Inches(0.12)
        tb(s8, item, Inches(0.6), ty, Inches(3.5), Inches(0.35),
           size=11, bold=True, color=NAVY)
        tb(s8, cost, Inches(4.05), ty, Inches(1.95), Inches(0.35),
           size=11, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)
    else:
        tb(s8, item, Inches(0.6), ty, Inches(3.5), Inches(0.3),
           size=10, color=NAVY)
        tb(s8, cost, Inches(4.05), ty, Inches(1.95), Inches(0.3),
           size=10, color=MGREY, align=PP_ALIGN.RIGHT)
    ty += Inches(0.5)

card(s8, Inches(6.5), Inches(1.4), Inches(6.5), Inches(5.6))
rect(s8, Inches(6.5), Inches(1.4), Inches(6.5), Inches(0.45), GREEN)
tb(s8, "Projected Returns & Value Unlocked",
   Inches(6.65), Inches(1.44), Inches(6.2), Inches(0.38),
   size=13, bold=True, color=WHITE)

returns = [
    ("1 extra 5-bed unit sold via diaspora SEO",   "$120K–$200K+ revenue (₦180M–₦300M)"),
    ("Realtor programme: 20 more active realtors", "+20 units referred / year potential"),
    ("Reduced buyer default via payment reminders", "5–10% fewer subscription cancellations"),
    ("Brochure → virtual tour → remote close",     "+30% diaspora conversion rate est."),
    ("Transparent MREIF page → fund inflows",      "₦50M–₦200M+ additional AUM potential"),
    ("Staff time saved on manual follow-ups",       "15–20 hrs/week reclaimed"),
    ("BREAKEVEN (fraction of one unit sale)",      "Full Year 1 investment recovered"),
]
ty2 = Inches(2.0)
for i, (item, val) in enumerate(returns):
    is_total = i == len(returns) - 1
    if is_total:
        rect(s8, Inches(6.5), ty2 - Inches(0.05), Inches(6.5), Inches(0.06), GOLD)
        ty2 += Inches(0.12)
        tb(s8, item, Inches(6.65), ty2, Inches(3.9), Inches(0.35),
           size=11, bold=True, color=NAVY)
        tb(s8, val, Inches(10.6), ty2, Inches(2.2), Inches(0.35),
           size=11, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)
    else:
        tb(s8, item, Inches(6.65), ty2, Inches(3.9), Inches(0.3),
           size=10, color=NAVY)
        tb(s8, val, Inches(10.6), ty2, Inches(2.2), Inches(0.3),
           size=10, color=MGREY, align=PP_ALIGN.RIGHT)
    ty2 += Inches(0.5)

tb(s8, "Commission on ONE additional 5-bedroom unit closes the entire Year 1 investment — every unit after that is pure upside",
   Inches(0.4), Inches(7.07), Inches(12.5), Inches(0.28),
   size=10.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 9 – IMPLEMENTATION TIMELINE
# ════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(BLANK)
bg(s9)
header(s9, "Implementation Timeline", "Phased 12-month delivery plan")
footer(s9)

phases = [
    ("Phase 1\nMon 1–2",  GOLD,     [
        "Add pricing to all property pages",
        "WhatsApp click-to-chat on all listings",
        "NDPR cookie consent + privacy policy",
        "Google Search Console + GA4 setup",
    ]),
    ("Phase 2\nMon 3–5",  NAVY,     [
        "Next.js + Sanity CMS rebuild launch",
        "Paystack (NGN) + Stripe (USD) live",
        "SEO: diaspora keywords + RealEstate schema",
        "MREIF dedicated landing page",
    ]),
    ("Phase 3\nMon 6–8",  EMERALD,  [
        "Buyer portal: payments + milestone tracker",
        "Realtor dashboard replaces Google Meet",
        "WhatsApp AI bot for full enquiry flow",
        "Virtual tour embeds on all unit pages",
    ]),
    ("Phase 4\nMon 9–12", GREEN,    [
        "AI construction progress auto-updates",
        "Diaspora email campaign automation",
        "MREIF investor quarterly report AI",
        "Full audit, review & optimise",
    ]),
]

for i, (label, col, items) in enumerate(phases):
    x = Inches(0.4) + i * Inches(3.2)
    y = Inches(1.45)
    card(s9, x, y, Inches(3.0), Inches(5.55))
    rect(s9, x, y, Inches(3.0), Inches(0.7), col)
    tb(s9, label, x, y + Inches(0.06), Inches(3.0), Inches(0.62),
       size=13, bold=True,
       color=NAVY if col == GOLD else WHITE,
       align=PP_ALIGN.CENTER)
    if i < 3:
        rect(s9, x + Inches(3.02), y + Inches(0.3), Inches(0.16), Inches(0.06), MGREY)
    ty = y + Inches(0.88)
    for item in items:
        rect(s9, x + Inches(0.18), ty + Inches(0.08), Inches(0.1), Inches(0.1), col)
        tb(s9, item, x + Inches(0.38), ty, Inches(2.5), Inches(0.42),
           size=10.5, color=NAVY)
        ty += Inches(0.9)

# ════════════════════════════════════════════════════════════════
# SLIDE 10 – WHY ACT NOW?
# ════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(BLANK)
bg(s10, NAVY)
rect(s10, 0, 0, Inches(0.5), H, GOLD)
footer(s10)

tb(s10, "Why Act — Now?",
   Inches(0.8), Inches(0.5), Inches(11), Inches(0.65),
   size=30, bold=True, color=GOLD)
tb(s10, "Ibeju-Lekki is Nigeria's fastest-appreciating corridor. The developers who own digital now will own the buyer relationship for a decade.",
   Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.55),
   size=14, color=WHITE, italic=True)

reasons = [
    ("Lekki-Epe is the hottest real estate corridor in West Africa",
     "Dangote Refinery, the proposed Lekki International Airport, and the Lekki Free Trade Zone are driving land appreciation at 15–30% annually. Buyers are researching now — Shalom Park must be findable and closeable digitally before competitors dominate."),
    ("Hidden pricing is costing sales today",
     "Every property page with no price forces a buyer to call or email. In 2026, most serious buyers — especially diaspora — move on immediately. Adding transparent NGN and USD pricing is the single highest-ROI change available at zero development cost."),
    ("The Nigerian diaspora market is $25B+ in annual remittances",
     "With IFT Realty's New York roots, Shalom Park is uniquely positioned to capture diaspora buyers. But without USD pricing, Stripe integration, or virtual tours, the diaspora channel is effectively closed — leaving the highest-value segment untapped."),
    ("The realtor programme is a flywheel — if it's automated",
     "Manual Google Meet onboarding cannot scale. A digital realtor portal with self-serve onboarding, lead assignment, and commission tracking can support 500+ active realtors — multiplying sales capacity without multiplying headcount."),
    ("Construction is underway — subscriber confidence needs feeding",
     "Buyers who have subscribed need consistent proof that their investment is safe and progressing. Automated WhatsApp milestone updates and portal photo galleries build the trust that generates referrals, top-up purchases, and 5-star reviews."),
]
ty = Inches(1.95)
for i, (title, desc) in enumerate(reasons):
    rect(s10, Inches(0.8), ty, Inches(0.55), Inches(0.55),
         GOLD if i % 2 == 0 else EMERALD)
    tb(s10, str(i + 1), Inches(0.8), ty, Inches(0.55), Inches(0.55),
       size=16, bold=True, color=NAVY if i % 2 == 0 else WHITE,
       align=PP_ALIGN.CENTER)
    tb(s10, title, Inches(1.5), ty, Inches(11), Inches(0.3),
       size=12, bold=True, color=WHITE)
    tb(s10, desc, Inches(1.5), ty + Inches(0.3), Inches(11.3), Inches(0.55),
       size=10.5, color=RGBColor(0xBB, 0xCC, 0xDD))
    ty += Inches(0.98)

# ════════════════════════════════════════════════════════════════
# SLIDE 11 – NEXT STEPS
# ════════════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(BLANK)
bg(s11)
header(s11, "Next Steps", "From audit to action — what happens now")
footer(s11)

steps = [
    ("This Week",  GOLD,    "Zero-Cost Quick Wins",
     "Add pricing to all property pages. Add WhatsApp click-to-chat on listings. These two changes alone will lift enquiry conversions immediately."),
    ("Week 2",     NAVY,    "Discovery Workshop",
     "30-min session with Genos Apollo. Review findings with Oluremi & team. Agree Phase 1 priorities and budget. Share brand assets and content."),
    ("Week 3–5",   EMERALD, "Sprint 1 Builds",
     "NDPR compliance, Paystack/Stripe integration, and SSR website rebuild start. WhatsApp Business API provisioned."),
    ("Month 2",    GREEN,   "Platform Launch",
     "New site live with pricing, USD option, and WhatsApp bot. Realtor portal v1 active. MREIF page indexed by Google."),
]

for i, (when, col, title, desc) in enumerate(steps):
    x = Inches(0.4) + i * Inches(3.2)
    y = Inches(1.5)
    card(s11, x, y, Inches(3.0), Inches(4.0))
    rect(s11, x, y, Inches(3.0), Inches(0.45), col)
    tb(s11, when, x, y + Inches(0.06), Inches(3.0), Inches(0.36),
       size=12, bold=True,
       color=NAVY if col == GOLD else WHITE,
       align=PP_ALIGN.CENTER)
    tb(s11, title, x + Inches(0.15), y + Inches(0.58),
       Inches(2.7), Inches(0.36), size=12, bold=True,
       color=WARM if col == GOLD else col)
    tb(s11, desc, x + Inches(0.15), y + Inches(1.0),
       Inches(2.7), Inches(2.8), size=10.5, color=NAVY)

card(s11, Inches(0.4), Inches(5.8), Inches(12.5), Inches(1.1), fill=NAVY)
tb(s11, "Ready to make Shalom Park Nigeria's most digitally advanced gated estate?",
   Inches(0.6), Inches(5.88), Inches(7.5), Inches(0.4),
   size=14, bold=True, color=GOLD)
tb(s11, "Contact Genos Apollo to schedule your discovery workshop.",
   Inches(0.6), Inches(6.28), Inches(7), Inches(0.38),
   size=11, color=WHITE)
tb(s11, "www.shalomparknigeria.com  |  info@shalomparknigeria.com  |  +234 703 546 1661",
   Inches(7.5), Inches(6.0), Inches(5.5), Inches(0.38),
   size=10, color=GOLD, italic=True, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════
# SLIDE 12 – THANK YOU / BACK COVER
# ════════════════════════════════════════════════════════════════
s12 = prs.slides.add_slide(BLANK)
bg(s12, NAVY)
rect(s12, 0, 0, Inches(0.7), H, GOLD)
rect(s12, Inches(0.7), Inches(3.5), W - Inches(0.7), Inches(0.06), GOLD)

tb(s12, "Thank You",
   Inches(1.1), Inches(2.0), Inches(11), Inches(1.2),
   size=52, bold=True, color=WHITE)
tb(s12, "Paradise on the Lekki-Epe Axis — built for the world to discover.",
   Inches(1.1), Inches(3.1), Inches(10), Inches(0.55),
   size=16, color=GOLD, italic=True)
tb(s12, "Shalom Park Estate  ·  IFT Realty Ltd  ·  Abijo, Ibeju Lekki, Lagos",
   Inches(1.1), Inches(4.0), Inches(9), Inches(0.4),
   size=12, color=RGBColor(0xAA, 0xBB, 0xCC))
tb(s12, "Prepared by Genos Apollo  |  May 2026  |  Confidential",
   Inches(1.1), Inches(6.7), Inches(10), Inches(0.35),
   size=10, color=MGREY)

# ── Save ────────────────────────────────────────────────────────
OUT = ROOT / "output" / "decks" / "ShalomPark_WebDev_AI_Automation_Audit_GenosAI.pptx"
prs.save(str(OUT))
print(f"Saved: {OUT}")
