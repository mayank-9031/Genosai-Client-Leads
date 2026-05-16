import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.parent / "utils"))
from ppt_theme import (
    NAVY, DARK_BLUE, ACCENT, GOLD, RED, GREEN, PURPLE, ORANGE, TEAL, CYAN,
    WHITE, LGREY, MGREY, W, H,
    bg, rect, tb, aline, bullets, tag,
    slide_header, slide_footer, score_badge,
    stat_cards, top_bar_card, phase_card, before_after, genos_logo_block,
)

ROOT = Path(__file__).parent.parent.parent
prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ══ SLIDE 1 — COVER ══
s1 = prs.slides.add_slide(BLANK)
bg(s1)
rect(s1, 0, 0, Inches(0.12), H, ACCENT)
rect(s1, Inches(8.5), 0, Inches(4.83), Inches(2.2), DARK_BLUE)
tb(s1, "GENOS AI", Inches(9.2), Inches(0.4), Inches(3.5), Inches(0.8),
   size=28, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
tb(s1, "www.genosai.tech", Inches(9.2), Inches(1.05), Inches(3.5), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
tb(s1, "ISELLVIRGINIA.COM / DIGIOVANNA REALTY GROUP", Inches(0.55), Inches(1.6),
   Inches(9), Inches(0.6), size=13, bold=False, color=ACCENT)
tb(s1, "Website & AI Growth\nAudit Report", Inches(0.55), Inches(2.1),
   Inches(9), Inches(1.8), size=48, bold=True, color=WHITE)
aline(s1, Inches(3.85), color=GOLD, w=Inches(3))
tb(s1, "Prepared exclusively for Rich DiGiovanna, Associate Broker & Team Leader",
   Inches(0.55), Inches(4.15), Inches(9), Inches(0.5), size=14, color=LGREY)
rect(s1, 0, Inches(6.5), W, Inches(1.0), DARK_BLUE)
tb(s1, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699",
   Inches(0.55), Inches(6.6), Inches(9), Inches(0.5), size=12, color=MGREY)
tb(s1, "May 2026  |  CONFIDENTIAL", Inches(9.5), Inches(6.6), Inches(3.3), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
rect(s1, Inches(10.5), Inches(3.3), Inches(2.3), Inches(2.5), DARK_BLUE)
tb(s1, "AUDIT SCORE", Inches(10.55), Inches(3.4), Inches(2.2), Inches(0.5),
   size=11, bold=True, color=MGREY, align=PP_ALIGN.CENTER)
tb(s1, "5", Inches(10.55), Inches(3.75), Inches(2.2), Inches(1.1),
   size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
tb(s1, "/ 10", Inches(10.55), Inches(4.7), Inches(2.2), Inches(0.4),
   size=18, color=LGREY, align=PP_ALIGN.CENTER)
tag(s1, "HIGH PRIORITY", Inches(10.75), Inches(5.1), fill=RED, size=10)


# ══ SLIDE 2 — AT A GLANCE ══
s2 = prs.slides.add_slide(BLANK)
bg(s2)
rect(s2, 0, 0, Inches(0.12), H, ACCENT)
tb(s2, "AT A GLANCE", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s2, "Rich DiGiovanna — Who He Is",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s2, Inches(1.35))

stats = [
    ("39\nYears", "In N. Virginia"),
    ("2,000+", "Families Helped"),
    ("$500M+", "Career Sales"),
    ("14 Years", "U.S. Coast Guard"),
    ("CRS / GRI\nABR", "Top 5% Certifications"),
]
cw = Inches(2.3); ch = Inches(2.2); gap = Inches(0.18)
for i, (num, label) in enumerate(stats):
    l = Inches(0.55) + i * (cw + gap)
    rect(s2, l, Inches(1.8), cw, ch, DARK_BLUE)
    tb(s2, num, l + Inches(0.1), Inches(1.95), cw - Inches(0.2), Inches(1.1),
       size=30, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s2, label, l + Inches(0.1), Inches(3.05), cw - Inches(0.2), Inches(0.7),
       size=12, color=LGREY, align=PP_ALIGN.CENTER)

rect(s2, Inches(0.55), Inches(4.2), Inches(12.23), Inches(0.6), DARK_BLUE)
tb(s2, "RE/MAX Allegiance  ·  Burke, VA & Woodbridge, VA  ·  YouTube: @DiGiovanna_VA  ·  Property Mgmt: Superior RE Services",
   Inches(0.65), Inches(4.28), Inches(12), Inches(0.45), size=12, color=LGREY, align=PP_ALIGN.CENTER)

tb(s2, "SERVICES", Inches(0.55), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=ACCENT)
bullets(s2, ["Residential buying & selling", "Property management (Superior RE Services)",
             "Mortgage (Intercoastal Mortgage partnership)"],
        Inches(0.55), Inches(5.35), Inches(5.8), Inches(1.5), size=13, color=LGREY)

tb(s2, "AREAS SERVED", Inches(6.8), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=ACCENT)
bullets(s2, ["Prince William County", "Fairfax County / Arlington County",
             "Burke, Woodbridge, Alexandria, Springfield"],
        Inches(6.8), Inches(5.35), Inches(5.8), Inches(1.5), size=13, color=LGREY)


# ══ SLIDE 3 — WHAT'S WORKING ══
s3 = prs.slides.add_slide(BLANK)
bg(s3)
rect(s3, 0, 0, Inches(0.12), H, GREEN)
tb(s3, "WHAT'S WORKING", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s3, "Strengths to Build On",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s3, Inches(1.35), color=GREEN)

strengths = [
    ("39-Year Track Record in One Market",
     "Deep local roots in Northern Virginia are a moat. 39 years in the same market means relationships, referrals, and trust that new agents can't replicate."),
    ("YouTube Channel Active",
     "Rich is posting video content - a signal that he understands modern marketing. Most agents his age don't. The infrastructure is there to turn that into serious lead generation."),
    ("Multiple Revenue Streams",
     "Brokerage, property management (Superior RE Services), and mortgage referrals through Intercoastal give the business diversified income that pure agents don't have."),
    ("IDX + CloudFlare + CRM in Place",
     "IDXBroker, LeadConnectorHQ, and CloudFlare CDN show a team that's already invested in digital infrastructure. These are good foundations to build AI automation on top of."),
    ("Military Background as Differentiator",
     "14 years in the Coast Guard is a real story. Precision, integrity, showing up when others don't. It's the kind of background that converts skeptical sellers into loyal clients."),
]
for i, (title, desc) in enumerate(strengths):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s3, Inches(0.55), t, Inches(0.06), Inches(0.7), GREEN)
    tb(s3, title, Inches(0.75), t, Inches(11.8), Inches(0.38), size=14, bold=True, color=WHITE)
    tb(s3, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# ══ SLIDE 4 — WHERE LEADS ARE LEAKING ══
s4 = prs.slides.add_slide(BLANK)
bg(s4)
rect(s4, 0, 0, Inches(0.12), H, RED)
tb(s4, "WHERE LEADS ARE LEAKING", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=RED)
tb(s4, "Critical Gaps Found in the Audit",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s4, Inches(1.35), color=RED)

issues = [
    ("TESTIMONIALS SECTION IS EMPTY", "HIGH",
     "There is a testimonials section on the site. It has no reviews in it. Sellers comparing agents will land here, see nothing, and move on. Every review from 2,000+ families should be here."),
    ("YOUTUBE CONTENT NOT ON THE HOMEPAGE", "HIGH",
     "Rich is posting to YouTube. None of it is embedded on the homepage. That video content builds trust faster than any text, and right now it's siloed to a platform Rich doesn't control."),
    ("NO LIVE CHAT OR AI BOT", "HIGH",
     "Zillow and Redfin run AI bots 24/7. When a buyer or seller visits iSellVirginia.com at 9pm on a Sunday and has a question, there's nothing there. That lead goes to whoever does respond."),
    ("TEMPLATE DESIGN LACKS PERSONALITY", "MEDIUM",
     "The site is functional but generic. Nothing about it reflects 39 years, a Coast Guard background, or 2,000 families. A competitor with a newer template looks just as credible on screen."),
    ("CTАС LACK URGENCY", "MEDIUM",
     "Multiple contact options with no clear hierarchy. No sticky header. No 'Get Your Free Home Valuation' or 'See What Your Home Is Worth Today' above the fold. Leads are leaking to friction."),
]
for i, (title, pri, desc) in enumerate(issues):
    t = Inches(1.7) + i * Inches(1.05)
    pc = RED if pri == "HIGH" else GOLD
    rect(s4, Inches(0.55), t, Inches(0.06), Inches(0.7), pc)
    tb(s4, title, Inches(0.75), t, Inches(9.2), Inches(0.38), size=13, bold=True, color=WHITE)
    tag(s4, pri, Inches(10.1), t + Inches(0.05), fill=pc, size=9)
    tb(s4, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# ══ SLIDE 5 — WEBSITE DESIGN AUDIT ══
s5 = prs.slides.add_slide(BLANK)
bg(s5)
rect(s5, 0, 0, Inches(0.12), H, PURPLE)
tb(s5, "WEBSITE DESIGN AUDIT", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=PURPLE)
tb(s5, "The Site Doesn't Match the Agent.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.8), size=34, bold=True, color=WHITE)
aline(s5, Inches(1.45), color=PURPLE, w=Inches(1.5))

# 3-column
rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(3.6), DARK_BLUE)
rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(0.42), RED)
tb(s5, "ISELLVIRGINIA TODAY", Inches(0.65), Inches(1.92), Inches(3.6), Inches(0.38),
   size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Generic template, no visual identity",
    "No video on homepage despite YouTube",
    "Empty testimonials section",
    "Centered CTAs, no sticky header",
    "Team page lacks personality/story",
    "Community pages are template clones",
    "No home valuation tool above the fold",
    "Coast Guard story buried in bio text",
], Inches(0.7), Inches(2.42), Inches(3.6), Inches(2.9), size=11, color=LGREY, dot=RED)

tb(s5, "->", Inches(4.45), Inches(3.6), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(3.6), DARK_BLUE)
rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(0.42), GOLD)
tb(s5, "TOP REALTORS DO THIS", Inches(5.1), Inches(1.92), Inches(3.4), Inches(0.38),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Video hero on homepage, autoplay muted",
    "Real reviews front and center (Zillow/Google)",
    "Agent story told with emotion, not a resume",
    "Home valuation CTA above the fold",
    "Sticky header with phone number always visible",
    "Neighborhood pages with real local content",
    "Mobile-first, sub-2s load time",
    "95+ Lighthouse mobile score",
], Inches(5.1), Inches(2.42), Inches(3.4), Inches(2.9), size=11, color=LGREY, dot=GOLD)

tb(s5, "->", Inches(8.68), Inches(3.6), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(3.6), DARK_BLUE)
rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(0.42), GREEN)
tb(s5, "WHAT GENOS AI BUILDS", Inches(9.35), Inches(1.92), Inches(3.4), Inches(0.38),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Premium site built around Rich's story",
    "YouTube feed embedded on homepage",
    "Real testimonials section, auto-populating",
    "Home valuation tool above the fold",
    "Sticky header with click-to-call",
    "Unique community pages for Burke/Woodbridge",
    "Mobile-first, sub-2s load",
    "AI chat on every page",
], Inches(9.35), Inches(2.42), Inches(3.4), Inches(2.9), size=11, color=LGREY, dot=GREEN)

rect(s5, Inches(0.55), Inches(5.72), Inches(12.23), Inches(1.28), DARK_BLUE)
tb(s5, "WHY IT MATTERS:", Inches(0.7), Inches(5.78), Inches(2.2), Inches(0.38),
   size=11, bold=True, color=PURPLE)
stats2 = [
    ("94%", "of first impressions\nare design-related"),
    ("75%", "of sellers judge agent\ncredibility by website"),
    ("3 sec", "before a visitor leaves\nor stays"),
    ("200%+", "more leads from a\nmodern agent site"),
]
for j, (num, lbl) in enumerate(stats2):
    lx = Inches(3.0) + j * Inches(2.45)
    tb(s5, num, lx, Inches(5.75), Inches(1.6), Inches(0.48),
       size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s5, lbl, lx, Inches(6.22), Inches(1.6), Inches(0.55),
       size=10, color=MGREY, align=PP_ALIGN.CENTER)


# ══ SLIDE 6 — AI AUTOMATION ══
s6 = prs.slides.add_slide(BLANK)
bg(s6)
rect(s6, 0, 0, Inches(0.12), H, ACCENT)
tb(s6, "AI AUTOMATION OPPORTUNITIES", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s6, "Seven AI Systems Built for iSellVirginia",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s6, Inches(1.35))

opps = [
    ("01", "AI Lead Follow-Up Sequences",
     "Every new inquiry triggers a personalized follow-up sequence automatically. No lead goes cold over a weekend or a holiday."),
    ("02", "Automated Seller Market Reports",
     "Monthly CMAs and market updates auto-sent to Rich's farm areas. Keeps iSellVirginia top of mind 60-90 days before someone decides to list."),
    ("03", "AI Chat (24/7 Lead Capture)",
     "Trained on Rich's services, pricing, and farm areas. Captures buyers and sellers at 11pm on a Sunday, qualifies them, and routes to Rich or the team."),
    ("04", "Post-Closing Review Automation",
     "Automated review requests sent after closing to Google, Zillow, and Realtor.com. Fills the empty testimonials section automatically over time."),
    ("05", "Property Mgmt Workflow Automation",
     "Superior Real Estate Services runs on manual processes. AI automates maintenance requests, rent reminders, lease renewals, and tenant communications."),
    ("06", "Listing Description Generator",
     "AI writes compelling, SEO-optimized listing descriptions from photos and data in minutes - freeing agent time for client work."),
    ("07", "Social Content Automation",
     "Auto-generate YouTube scripts, Instagram captions, and Facebook posts from listings and market data. Keeps all channels active without extra effort."),
]
col1 = opps[:4]; col2 = opps[4:]
for i, (num, title, desc) in enumerate(col1):
    t = Inches(1.65) + i * Inches(1.42)
    rect(s6, Inches(0.55), t, Inches(6.0), Inches(1.28), DARK_BLUE)
    tb(s6, num, Inches(0.65), t + Inches(0.1), Inches(0.7), Inches(0.55),
       size=20, bold=True, color=ACCENT)
    tb(s6, title, Inches(1.35), t + Inches(0.1), Inches(5.0), Inches(0.38),
       size=13, bold=True, color=WHITE)
    tb(s6, desc, Inches(0.65), t + Inches(0.52), Inches(5.7), Inches(0.65),
       size=11, color=LGREY)
for i, (num, title, desc) in enumerate(col2):
    t = Inches(1.65) + i * Inches(1.42)
    rect(s6, Inches(7.0), t, Inches(5.8), Inches(1.28), DARK_BLUE)
    tb(s6, num, Inches(7.1), t + Inches(0.1), Inches(0.7), Inches(0.55),
       size=20, bold=True, color=GOLD)
    tb(s6, title, Inches(7.8), t + Inches(0.1), Inches(4.8), Inches(0.38),
       size=13, bold=True, color=WHITE)
    tb(s6, desc, Inches(7.1), t + Inches(0.52), Inches(5.5), Inches(0.65),
       size=11, color=LGREY)


# ══ SLIDE 7 — WHAT THIS IS COSTING YOU ══
s7 = prs.slides.add_slide(BLANK)
bg(s7)
rect(s7, 0, 0, Inches(0.12), H, GOLD)
tb(s7, "WHAT THIS IS COSTING YOU", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GOLD)
tb(s7, "The Hidden Cost of the Status Quo",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s7, Inches(1.35), color=GOLD)

costs = [
    ("Cold Leads Every Weekend",
     "Real estate inquiries don't respect business hours. Every Friday evening to Monday morning, leads hitting iSellVirginia.com with no chat, no bot, and no follow-up automation are going to Zillow agents who respond in under 5 minutes."),
    ("Reviews Left on the Table",
     "Rich has helped 2,000+ families. His testimonials section is empty. Every seller comparing agents is making a decision with incomplete information because those reviews were never asked for systematically."),
    ("YouTube Views That Don't Convert",
     "Rich posts on YouTube but the site doesn't link back to it meaningfully. Video builds trust faster than any other format. That investment is currently disconnected from the lead pipeline."),
    ("Manual Property Management",
     "Superior Real Estate Services is a real business running on manual processes. Maintenance requests, reminders, renewals handled by hand. Every manual process is a margin leak and a scaling ceiling."),
]
for i, (title, desc) in enumerate(costs):
    t = Inches(1.7) + i * Inches(1.35)
    rect(s7, Inches(0.55), t, Inches(12.23), Inches(1.22), DARK_BLUE)
    tb(s7, "!  " + title, Inches(0.7), t + Inches(0.1), Inches(5), Inches(0.42),
       size=14, bold=True, color=GOLD)
    tb(s7, desc, Inches(0.7), t + Inches(0.5), Inches(11.8), Inches(0.6),
       size=12, color=LGREY)


# ══ SLIDE 8 — WHAT WE'D BUILD ══
s8 = prs.slides.add_slide(BLANK)
bg(s8)
rect(s8, 0, 0, Inches(0.12), H, ACCENT)
tb(s8, "WHAT WE'D BUILD", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s8, "The Genos AI Stack for iSellVirginia",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s8, Inches(1.35))

builds = [
    ("New Website Rebuild",
     "Premium design built around Rich's story: YouTube embedded, valuation CTA above fold, real testimonials, community pages for Burke and Woodbridge.",
     "Weeks 1-4"),
    ("AI Chat Widget",
     "Trained on Rich's listings, farm areas, and services. Captures and qualifies leads 24/7, routes urgent ones to Rich directly.",
     "Weeks 1-3"),
    ("Lead Follow-Up Engine",
     "30/60/90-day automated nurture sequences for every new lead. Personalized by buyer vs. seller, location, and price range.",
     "Weeks 2-5"),
    ("Seller Market Report Automation",
     "Monthly automated CMAs to Rich's farm area database. Keeps iSellVirginia top of mind before someone decides to list.",
     "Weeks 3-6"),
    ("Review Automation + Property Mgmt",
     "Post-closing review requests auto-sent to Google/Zillow. Superior RE Services workflows automated: maintenance, reminders, renewals.",
     "Weeks 4-8"),
]
for i, (title, desc, tl) in enumerate(builds):
    t = Inches(1.65) + i * Inches(1.12)
    rect(s8, Inches(0.55), t, Inches(10.8), Inches(1.0), DARK_BLUE)
    rect(s8, Inches(0.62), t + Inches(0.15), Inches(0.55), Inches(0.55), ACCENT)
    tb(s8, str(i + 1), Inches(0.62), t + Inches(0.12), Inches(0.55), Inches(0.55),
       size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s8, title, Inches(1.28), t + Inches(0.08), Inches(7.5), Inches(0.38),
       size=14, bold=True, color=WHITE)
    tb(s8, desc, Inches(1.28), t + Inches(0.48), Inches(7.5), Inches(0.42),
       size=11, color=LGREY)
    tag(s8, tl, Inches(11.5), t + Inches(0.3), fill=DARK_BLUE, tc=MGREY, size=10)


# ══ SLIDE 9 — 90-DAY OUTCOMES ══
s9 = prs.slides.add_slide(BLANK)
bg(s9)
rect(s9, 0, 0, Inches(0.12), H, ACCENT)
tb(s9, "90-DAY OUTCOMES", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s9, "What iSellVirginia Looks Like 90 Days From Now",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=30, bold=True, color=WHITE)
aline(s9, Inches(1.35))

months = [
    ("MONTH 1", ACCENT, [
        "New site live - YouTube embedded, valuation CTA, modern design",
        "AI chat capturing leads 24/7",
        "First automated market reports sent to farm areas",
        "Post-closing review requests active",
    ]),
    ("MONTH 2", GOLD, [
        "Lead nurture sequences running for all new inquiries",
        "Testimonials section filling automatically",
        "Superior RE Services maintenance workflows automated",
        "First SEO-optimized community pages ranking",
    ]),
    ("MONTH 3", GREEN, [
        "Full lead pipeline automated start to finish",
        "Review velocity 2-3x what it was before",
        "Property mgmt on autopilot, margin recovered",
        "YouTube content feeding directly into lead funnel",
    ]),
]
for i, (month, color, items) in enumerate(months):
    cl = Inches(0.55) + i * Inches(4.2)
    rect(s9, cl, Inches(1.6), Inches(3.9), Inches(5.4), DARK_BLUE)
    rect(s9, cl, Inches(1.6), Inches(3.9), Inches(0.5), color)
    tb(s9, month, cl, Inches(1.6), Inches(3.9), Inches(0.5),
       size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    bullets(s9, items, cl + Inches(0.15), Inches(2.25),
            Inches(3.6), Inches(4.4), size=13, color=LGREY, dot=color)


# ══ SLIDE 10 — NEXT STEPS ══
s10 = prs.slides.add_slide(BLANK)
bg(s10)
rect(s10, 0, 0, Inches(0.12), H, GOLD)
tb(s10, "NEXT STEPS", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GOLD)
tb(s10, "One Call to See What's Possible",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s10, Inches(1.35), color=GOLD)

steps = [
    ("1", "Reply or Book a 15-Min Call",
     "No pitch. No commitment. We walk you through the audit and answer any questions about what we found."),
    ("2", "Free Strategy Session",
     "We map out which AI automations and design changes would deliver the fastest ROI for iSellVirginia specifically."),
    ("3", "Custom Build Roadmap",
     "We deliver a bespoke implementation plan, timeline, and cost estimate tailored to Rich's team size and goals."),
]
for i, (num, title, desc) in enumerate(steps):
    t = Inches(2.0) + i * Inches(1.5)
    rect(s10, Inches(0.55), t, Inches(12.23), Inches(1.3), DARK_BLUE)
    rect(s10, Inches(0.55), t, Inches(1.0), Inches(1.3), GOLD)
    tb(s10, num, Inches(0.55), t + Inches(0.2), Inches(1.0), Inches(0.7),
       size=36, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    tb(s10, title, Inches(1.7), t + Inches(0.15), Inches(10), Inches(0.45),
       size=16, bold=True, color=WHITE)
    tb(s10, desc, Inches(1.7), t + Inches(0.6), Inches(10), Inches(0.5),
       size=13, color=LGREY)

rect(s10, 0, Inches(6.5), W, Inches(1.0), ACCENT)
tb(s10, "Rohan Malik  ·  CEO, Genos AI  ·  hello@genosai.tech  ·  +91 638-714-2699  ·  www.genosai.tech",
   Inches(0.55), Inches(6.6), Inches(12.23), Inches(0.55),
   size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

prs.save(ROOT / "output" / "decks" / "iSellVirginia_Audit_GenosAI.pptx")
print("Saved: iSellVirginia_Audit_GenosAI.pptx  |  10 slides")
