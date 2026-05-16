import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.parent / "utils"))
from ppt_theme import (
    NAVY, DARK_BLUE, ACCENT, GOLD, RED, GREEN, PURPLE, ORANGE, TEAL, CYAN,
    WHITE, LGREY, MGREY, W, H,
    bg, rect, tb, aline, bullets, tag,
    slide_header, slide_footer, score_badge,
    stat_cards, top_bar_card, phase_card, before_after, genos_logo_block,
)

ROOT = Path(__file__).parent.parent.parent
prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ── SLIDE 1 — COVER ──────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(BLANK)
bg(s1)
rect(s1, 0, 0, Inches(0.12), H, ACCENT)
genos_logo_block(s1)
tb(s1, "[CLIENT NAME]", Inches(0.55), Inches(1.6),
   Inches(9), Inches(0.6), size=13, color=ACCENT)
tb(s1, "AI Operations &\nProductivity Audit", Inches(0.55), Inches(2.05),
   Inches(9), Inches(1.8), size=44, bold=True, color=WHITE)
aline(s1, Inches(3.85), color=GOLD, w=Inches(3))
tb(s1, "Prepared exclusively for [Contact Name], [Title]",
   Inches(0.55), Inches(4.15), Inches(9), Inches(0.5), size=13, color=LGREY)
slide_footer(s1)
score_badge(s1, score=5, priority_text="HIGH PRIORITY", priority_color=RED)


# ── SLIDE 2 — EXECUTIVE SUMMARY ───────────────────────────────────────────────
s2 = prs.slides.add_slide(BLANK)
bg(s2)
slide_header(s2, "EXECUTIVE SUMMARY", "At a Glance — Current Reality & Opportunity",
             accent_color=CYAN, aline_w=Inches(1.5))

# Left card — Current Reality
rect(s2, Inches(0.55), Inches(1.6), Inches(3.9), Inches(4.6), DARK_BLUE)
rect(s2, Inches(0.55), Inches(1.6), Inches(3.9), Inches(0.06), RED)
tb(s2, "Current Operational Reality", Inches(0.65), Inches(1.72),
   Inches(3.7), Inches(0.4), size=12, bold=True, color=RED)
bullets(s2, [
    "90+ employees across 6 departments",
    "Heavy manual dependency in reporting & revisions",
    "Repetitive workflows consuming skilled hours",
    "Disconnected tools and manual approvals",
    "High dependency on human follow-ups",
], Inches(0.65), Inches(2.15), Inches(3.7), Inches(3.8), size=11, dot=RED)

# Middle card — Bottlenecks
rect(s2, Inches(4.65), Inches(1.6), Inches(3.9), Inches(4.6), DARK_BLUE)
rect(s2, Inches(4.65), Inches(1.6), Inches(3.9), Inches(0.06), GOLD)
tb(s2, "Current Bottlenecks", Inches(4.75), Inches(1.72),
   Inches(3.7), Inches(0.4), size=12, bold=True, color=GOLD)
bullets(s2, [
    "Delayed execution",
    "Manual reporting overhead",
    "Slower client turnaround",
    "Inconsistent workflow quality",
    "High operational fatigue",
    "Difficulty scaling output",
], Inches(4.75), Inches(2.15), Inches(3.7), Inches(3.8), size=11, dot=GOLD)

# Right card — AI Opportunity
rect(s2, Inches(8.75), Inches(1.6), Inches(4.13), Inches(4.6), DARK_BLUE)
rect(s2, Inches(8.75), Inches(1.6), Inches(4.13), Inches(0.06), GREEN)
tb(s2, "AI Transformation Opportunity", Inches(8.85), Inches(1.72),
   Inches(3.9), Inches(0.4), size=12, bold=True, color=GREEN)
metrics = [
    ("Task Speed", "+35–60%"),
    ("Reporting Time", "70% faster"),
    ("Manual Work", "−40%"),
    ("Efficiency", "+45%"),
    ("Productivity", "1.8×"),
    ("Annual Savings", "₹45L–₹1.2Cr"),
]
for i, (k, v) in enumerate(metrics):
    y = Inches(2.2) + i * Inches(0.62)
    rect(s2, Inches(8.85), y, Inches(3.9), Inches(0.52), NAVY)
    tb(s2, k, Inches(8.95), y + Inches(0.08), Inches(2.2), Inches(0.36),
       size=11, color=LGREY)
    tb(s2, v, Inches(11.1), y + Inches(0.08), Inches(1.55), Inches(0.36),
       size=11, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)

rect(s2, Inches(0.55), Inches(6.32), Inches(12.23), Inches(0.38), ACCENT)
tb(s2, 'KEY INSIGHT: "The goal is not replacement — the goal is MORE OUTPUT per employee."',
   Inches(0.65), Inches(6.37), Inches(12), Inches(0.3), size=12, bold=True,
   color=NAVY, align=PP_ALIGN.CENTER)


# ── SLIDE 3 — THE 5 OPERATIONAL GAPS ─────────────────────────────────────────
s3 = prs.slides.add_slide(BLANK)
bg(s3)
slide_header(s3, "OPERATIONAL DIAGNOSIS", "The 5 Gaps Limiting Scale",
             accent_color=RED, aline_w=Inches(1.2))

gap_colors = [RED, GOLD, ORANGE, PURPLE, CYAN]
gaps = [
    ("GAP 1\nREPETITIVE\nMANUAL WORK", RED,
     "Skilled employees on manual tasks",
     ["Hours lost to edits, data entry, formatting",
      "High salary cost on low-value tasks",
      "AI removes the backlog permanently"]),
    ("GAP 2\nKNOWLEDGE\nSILOS", GOLD,
     "Workflows live only in people's heads",
     ["Onboarding delays + inconsistent output",
      "Operational chaos when key staff leave",
      "AI codifies SOP into automated processes"]),
    ("GAP 3\nREPORTING\nIS SLOW", ORANGE,
     "Decision-making waits on manual reports",
     ["Managers wait hours/days for metrics",
      "Campaign, HR, and inventory data delayed",
      "AI dashboards deliver it in real time"]),
    ("GAP 4\nTOOL\nFRAGMENTATION", PURPLE,
     "Excel, WhatsApp, Email, Drive, approvals",
     ["Inefficiency and communication gaps",
      "No single source of truth for operations",
      "AI unifies data across all tools"]),
    ("GAP 5\nSCALING\nREQUIRES HIRING", CYAN,
     "Without AI: 150 hires for 2× growth",
     ["With AI: 90 employees → output of 150+",
      "Every new hire costs ₹6–18L/year",
      "AI closes the headcount gap"]),
]
cw = Inches(2.35); gap = Inches(0.12)
for i, (title, color, subtitle, pts) in enumerate(gaps):
    l = Inches(0.55) + i * (cw + gap)
    rect(s3, l, Inches(1.8), cw, Inches(0.52), color)
    tb(s3, title, l + Inches(0.1), Inches(1.84),
       cw - Inches(0.2), Inches(0.44),
       size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    rect(s3, l, Inches(2.32), cw, Inches(4.48), DARK_BLUE)
    tb(s3, subtitle, l + Inches(0.1), Inches(2.4),
       cw - Inches(0.2), Inches(0.35),
       size=10, italic=True, color=color)
    bullets(s3, pts, l + Inches(0.1), Inches(2.82),
            cw - Inches(0.2), Inches(3.8), size=11, dot=color)

rect(s3, Inches(0.55), Inches(6.9), Inches(12.23), Inches(0.36), DARK_BLUE)
tb(s3, "Every gap has a documented AI solution. None require replacing your team — only augmenting it.",
   Inches(0.7), Inches(6.95), Inches(12), Inches(0.26),
   size=11, italic=True, color=GOLD)


# ── SLIDE 4 — AI AUTOMATION STACK ────────────────────────────────────────────
s4 = prs.slides.add_slide(BLANK)
bg(s4)
slide_header(s4, "AI AUTOMATION STACK",
             "Four Layers That Remove Bottlenecks Across Every Department",
             accent_color=ACCENT, aline_w=Inches(1.8))

opps = [
    ("HR Automation", TEAL,
     "AI resume ranking, employee support chatbot (Hindi + English), automated onboarding flows, and live HR analytics dashboards. Reduces HR coordination by 50% and screening time by 80%."),
    ("Design & CAD\nAcceleration", GOLD,
     "Batch Photoshop automation for exports and resizing, AI-assisted design variations, CAD template acceleration, and smart asset management. Design throughput increases 60%."),
    ("Marketing &\nContent Engine", PURPLE,
     "AI blogs, ad copy, and captions at 3× current velocity. Automated PPC reporting, keyword clustering, social scheduling, and wasted ad spend detection. Reporting time cut 90%."),
    ("Analytics &\nInventory AI", GREEN,
     "Predictive inventory forecasting, automated daily/weekly reports, AI anomaly detection, and executive dashboards. Reporting speed 10× faster. Forecasting accuracy +35%."),
]
cw = Inches(2.95); gap = Inches(0.16)
for i, (title, color, desc) in enumerate(opps):
    l = Inches(0.55) + i * (cw + gap)
    top_bar_card(s4, l, Inches(1.8), cw, Inches(4.7), color, title, desc,
                 title_size=13, desc_size=11)

rect(s4, Inches(0.55), Inches(6.62), Inches(12.23), Inches(0.55), DARK_BLUE)
tb(s4, "Anurag's team builds the products. Genos AI handles the operational layer that scales without adding headcount.",
   Inches(0.7), Inches(6.68), Inches(12), Inches(0.4), size=11, italic=True, color=GOLD)


# ── SLIDE 5 — BEFORE vs AFTER ────────────────────────────────────────────────
s5 = prs.slides.add_slide(BLANK)
bg(s5)
slide_header(s5, "BEFORE vs AFTER",
             "What Actually Changes in 90 Days", accent_color=ACCENT)

before_after(s5, [
    ("Hotel Developer Inquiry",
     "Form submitted Sunday. Reply arrives Monday. Competitor already has site visit booked.",
     "WhatsApp AI responds in 60 seconds. 5 qualifying questions. Site visit booked before Monday."),
    ("AMC Renewal (490 Clients)",
     "Contracts expire. Manual calls or no contact. Each lapsed AMC is direct revenue lost.",
     "WhatsApp reminders at 30/7/1 days. Payment link in message. 490 clients managed automatically."),
    ("Reporting",
     "Managers wait hours for campaign, HR, or inventory data. Decisions are delayed.",
     "Live dashboards refresh in real time. Every department has data without asking for it."),
    ("Content Production",
     "1 article published last quarter. SCI invisible to organic search. Competitors rank instead.",
     "4–8 articles published monthly. SEO traction compounds. Inbound qualified leads start flowing."),
    ("Candidate Screening",
     "Consultants spend 60%+ of time on manual resume review. High-value time on low-value work.",
     "AI scoring pipeline: top candidates surfaced in minutes. Consultants focus on relationships."),
])


# ── SLIDE 6 — IMPLEMENTATION ROADMAP ─────────────────────────────────────────
s6 = prs.slides.add_slide(BLANK)
bg(s6)
slide_header(s6, "IMPLEMENTATION ROADMAP", "Your Path to AI-Powered Operations",
             accent_color=GREEN, aline_w=Inches(1.5))

pw = Inches(2.95); gap = Inches(0.16)
phase_card(s6, Inches(0.55), Inches(1.6), pw, Inches(4.75), GOLD,
           "PHASE 1", "Weeks 1–3", [
               ("Workflow Audit & SOP Mapping",
                "Department workflow mapping, SOP analysis, repetitive task identification, ROI priority ranking."),
               ("Quick-Win Automations Live",
                "HR chatbot, AI reporting dashboards, batch design automation deployed."),
           ])

phase_card(s6, Inches(0.55) + pw + gap, Inches(1.6), pw, Inches(4.75), ACCENT,
           "PHASE 2", "Weeks 4–8", [
               ("Department AI Copilots",
                "HR, Marketing, Analytics, and Dev copilots deployed and trained on company SOPs."),
               ("AI Content Engine Running",
                "4–8 articles/month. PPC reporting automated. Social scheduling active."),
           ])

phase_card(s6, Inches(0.55) + 2 * (pw + gap), Inches(1.6), pw, Inches(4.75), PURPLE,
           "PHASE 3", "Months 2–3", [
               ("Advanced AI Integrations",
                "Inventory forecasting, candidate outreach engine, executive movement intelligence."),
               ("Unified Dashboard Layer",
                "Leadership gets live operational visibility across all departments."),
           ])

phase_card(s6, Inches(0.55) + 3 * (pw + gap), Inches(1.6), pw, Inches(4.75), GREEN,
           "PHASE 4", "Ongoing", [
               ("Full AI Operations Layer",
                "All copilots live, workflow orchestration active, continuous optimisation cadence."),
               ("Scale Without Headcount",
                "90 employees delivering output of 140–160. AI closes the gap."),
           ])

rect(s6, Inches(0.55), Inches(6.45), Inches(12.23), Inches(0.48), ACCENT)
tb(s6, '"Companies that scale fastest will not have the most employees — but the best Human + AI systems."',
   Inches(0.7), Inches(6.5), Inches(12), Inches(0.4),
   size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ── SLIDE 7 — ROI & BUSINESS IMPACT ─────────────────────────────────────────
s7 = prs.slides.add_slide(BLANK)
bg(s7)
slide_header(s7, "BUSINESS IMPACT", "Projected ROI — 12 Months",
             accent_color=GOLD, aline_w=Inches(1.2))

# Impact stats
impact_stats = [
    ("+45%", "Overall\nproductivity"),
    ("10×", "Reporting\nspeed"),
    ("−40%", "Repetitive\nmanual work"),
    ("1.8×", "Employee\noutput"),
    ("₹45L–1.2Cr", "Annual\nsavings"),
]
stat_cards(s7, impact_stats, top=Inches(1.8), card_w=Inches(2.3), card_h=Inches(2.0),
           num_color=GOLD)

# Without vs With AI
rect(s7, Inches(0.55), Inches(4.1), Inches(5.8), Inches(2.0), DARK_BLUE)
rect(s7, Inches(0.55), Inches(4.1), Inches(5.8), Inches(0.06), RED)
tb(s7, "WITHOUT AI", Inches(0.65), Inches(4.22), Inches(5.6), Inches(0.38),
   size=12, bold=True, color=RED)
bullets(s7, [
    "25–40 additional hires required for 2× growth",
    "Yearly operational cost increase: ₹1.8Cr – ₹3.5Cr",
    "Reporting delays slow every strategic decision",
], Inches(0.65), Inches(4.62), Inches(5.5), Inches(1.38), size=12, dot=RED)

rect(s7, Inches(6.65), Inches(4.1), Inches(6.13), Inches(2.0), DARK_BLUE)
rect(s7, Inches(6.65), Inches(4.1), Inches(6.13), Inches(0.06), GREEN)
tb(s7, "WITH AI + AUTOMATION", Inches(6.75), Inches(4.22), Inches(5.9), Inches(0.38),
   size=12, bold=True, color=GREEN)
bullets(s7, [
    "Same 90 employees, output of 140–160",
    "Annual operational savings: ₹45L – ₹1.2Cr",
    "Every department has live data without asking",
], Inches(6.75), Inches(4.62), Inches(5.9), Inches(1.38), size=12, dot=GREEN)

rect(s7, Inches(0.55), Inches(6.22), Inches(12.23), Inches(0.48), DARK_BLUE)
tb(s7, "AI removes repetitive work, coordination bottlenecks & operational inefficiency — not leadership, creativity, or strategy.",
   Inches(0.7), Inches(6.28), Inches(12), Inches(0.38),
   size=11, color=LGREY, align=PP_ALIGN.CENTER)


# ── SLIDE 8 — NEXT STEPS ────────────────────────────────────────────────────
s8 = prs.slides.add_slide(BLANK)
bg(s8)
rect(s8, 0, 0, Inches(0.12), H, CYAN)
rect(s8, Inches(7.5), 0, Inches(5.83), H, DARK_BLUE)

tb(s8, "NEXT STEPS", Inches(0.55), Inches(0.3), Inches(6.5), Inches(0.5),
   size=11, bold=True, color=CYAN)
tb(s8, "Let's Talk for\n15 Minutes.",
   Inches(0.55), Inches(0.8), Inches(6.8), Inches(2.0), size=44, bold=True, color=WHITE)
aline(s8, Inches(2.8), color=CYAN, w=Inches(2.5))
tb(s8, "No pitch. No commitment.\nWe walk through the findings and show you exactly what we'd automate first.",
   Inches(0.55), Inches(3.05), Inches(6.5), Inches(0.8), size=13, color=LGREY)

steps = [
    "Reply or book a 15-min call: hello@genosai.tech",
    "Free strategy session: map highest-ROI automations for your team",
    "Custom AI roadmap: timeline, deliverables, cost estimate",
    "Phase 1 live in 2–3 weeks — quick wins before month end",
]
bullets(s8, steps, Inches(0.55), Inches(3.95), Inches(6.5), Inches(2.0),
        size=13, color=WHITE, dot=CYAN)

rect(s8, 0, Inches(6.5), Inches(7.5), Inches(1.0), DARK_BLUE)
tb(s8, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699  |  www.genosai.tech",
   Inches(0.55), Inches(6.6), Inches(6.8), Inches(0.5), size=11, color=MGREY)

tb(s8, "5 GAPS.\n5 FIXES.", Inches(7.85), Inches(0.6), Inches(5.0), Inches(1.1),
   size=22, bold=True, color=CYAN)
aline(s8, Inches(1.65), color=CYAN)

summary = [
    ("Manual Work", "AI copilots eliminate repetitive tasks across every dept"),
    ("Slow Reports", "Live dashboards — no more waiting for data"),
    ("Tool Chaos", "AI unifies your stack into a single operational layer"),
    ("Content Gap", "AI content engine: 4-8 articles/month, SEO compounds"),
    ("Scaling Cost", "90 employees → output of 140+ with AI augmentation"),
]
for i, (label, val) in enumerate(summary):
    t2 = Inches(1.85) + i * Inches(0.88)
    rect(s8, Inches(7.85), t2, Inches(1.5), Inches(0.72), NAVY)
    tb(s8, label, Inches(7.9), t2 + Inches(0.05), Inches(1.4), Inches(0.32),
       size=10, bold=True, color=CYAN)
    tb(s8, val, Inches(9.5), t2 + Inches(0.05), Inches(3.6), Inches(0.65),
       size=11, color=LGREY)


out = ROOT / "output" / "decks" / "GenosAI_Operations_Audit.pptx"
prs.save(out)
print(f"Saved: {out}")
