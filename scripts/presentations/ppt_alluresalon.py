import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.parent / "utils"))
from ppt_theme import (
    NAVY, DARK_BLUE, ACCENT, GOLD, RED, GREEN, PURPLE, ORANGE, TEAL, CYAN,
    WHITE, LGREY, MGREY, W, H,
    bg, rect, tb, aline, bullets, tag,
    slide_header, slide_footer, score_badge,
    stat_cards, top_bar_card, phase_card, before_after, genos_logo_block,
)

ROOT = Path(__file__).parent.parent.parent
prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]

# SLIDE 1 - COVER
s1 = prs.slides.add_slide(BLANK)
bg(s1)
rect(s1, 0, 0, Inches(0.12), H, GOLD)
rect(s1, Inches(8.5), 0, Inches(4.83), Inches(2.2), DARK_BLUE)
tb(s1, "GENOS AI", Inches(9.2), Inches(0.4), Inches(3.5), Inches(0.8),
   size=28, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
tb(s1, "www.genosai.tech", Inches(9.2), Inches(1.05), Inches(3.5), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
tb(s1, "ALLURE SALON SUITE CONSULTING & ALLURE SALON GROUP", Inches(0.55), Inches(1.6),
   Inches(9), Inches(0.6), size=13, bold=False, color=GOLD)
tb(s1, "Website & AI Growth\nAudit Report", Inches(0.55), Inches(2.1),
   Inches(9), Inches(1.8), size=48, bold=True, color=WHITE)
aline(s1, Inches(3.85), color=GOLD, w=Inches(3))
tb(s1, "Prepared exclusively for Karen A. Kaminski, Founder & CEO",
   Inches(0.55), Inches(4.15), Inches(9), Inches(0.5), size=14, color=LGREY)
rect(s1, 0, Inches(6.5), W, Inches(1.0), DARK_BLUE)
tb(s1, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699",
   Inches(0.55), Inches(6.6), Inches(9), Inches(0.5), size=12, color=MGREY)
tb(s1, "May 2026  |  CONFIDENTIAL", Inches(9.5), Inches(6.6), Inches(3.3), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
rect(s1, Inches(10.5), Inches(3.3), Inches(2.3), Inches(2.5), DARK_BLUE)
tb(s1, "AUDIT SCORE", Inches(10.55), Inches(3.4), Inches(2.2), Inches(0.5),
   size=11, bold=True, color=MGREY, align=PP_ALIGN.CENTER)
tb(s1, "5", Inches(10.55), Inches(3.75), Inches(2.2), Inches(1.1),
   size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
tb(s1, "/ 10", Inches(10.55), Inches(4.7), Inches(2.2), Inches(0.4),
   size=18, color=LGREY, align=PP_ALIGN.CENTER)
tag(s1, "HIGH PRIORITY", Inches(10.75), Inches(5.1), fill=RED, size=10)


# SLIDE 2 - AT A GLANCE
s2 = prs.slides.add_slide(BLANK)
bg(s2)
rect(s2, 0, 0, Inches(0.12), H, GOLD)
tb(s2, "AT A GLANCE", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GOLD)
tb(s2, "Karen A. Kaminski - Who She Is",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s2, Inches(1.35), color=GOLD)

stats = [
    ("5", "Owned Locations"),
    ("115+", "Salon Studios"),
    ("$150K", "Avg Revenue / Location"),
    ("30-40%", "Client Startup Savings"),
    ("US + CA", "Operating Footprint"),
]
cw = Inches(2.3); ch = Inches(2.2); gap = Inches(0.18)
for i, (num, label) in enumerate(stats):
    l = Inches(0.55) + i * (cw + gap)
    rect(s2, l, Inches(1.8), cw, ch, DARK_BLUE)
    tb(s2, num, l + Inches(0.1), Inches(1.95), cw - Inches(0.2), Inches(1.1),
       size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s2, label, l + Inches(0.1), Inches(3.05), cw - Inches(0.2), Inches(0.7),
       size=12, color=LGREY, align=PP_ALIGN.CENTER)

rect(s2, Inches(0.55), Inches(4.2), Inches(12.23), Inches(0.6), DARK_BLUE)
tb(s2, "Moosic, PA  -  Founded 2018  -  Modern Salon Top 10 Nominee 2024  -  Authority Magazine  -  IAMCEO Podcast  -  SF Commercial  -  Elite Property News",
   Inches(0.65), Inches(4.28), Inches(12), Inches(0.45), size=12, color=LGREY, align=PP_ALIGN.CENTER)

tb(s2, "BACKGROUND THAT SETS HER APART", Inches(0.55), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=GOLD)
bullets(s2, ["BA Advertising/Graphics/Fine Arts, Marywood University - design eye is real",
             "15+ years beauty industry: CosmoProf, East Coast Salon Services, Moroccan Oil",
             "Opened first independent salon suite in 2016 - hands-on operator, not just a consultant"],
        Inches(0.55), Inches(5.35), Inches(5.8), Inches(1.5), size=13, color=LGREY)

tb(s2, "WHAT SHE BUILT", Inches(6.8), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=GOLD)
bullets(s2, ["Allure Salon Group: 5 owned locations housing 115+ working studios",
             "Allure Salon Suite Consulting (2018): end-to-end build for investors across US + Canada",
             "Allure University in development - online education for salon suite operators"],
        Inches(6.8), Inches(5.35), Inches(5.8), Inches(1.5), size=13, color=LGREY)


# SLIDE 3 - WHAT'S WORKING
s3 = prs.slides.add_slide(BLANK)
bg(s3)
rect(s3, 0, 0, Inches(0.12), H, GREEN)
tb(s3, "WHAT'S WORKING", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s3, "Strengths to Build On",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s3, Inches(1.35), color=GREEN)

strengths = [
    ("The Operator-Consultant Combination Is Genuinely Rare",
     "Most salon suite consultants either come from a beauty background but never built a facility, or from real estate but never worked in a salon. Karen owns and operates 5 locations housing 115+ studios. She runs the same playbook she sells. That's the differentiator the homepage should be leading with, and the credibility that separates Allure from anyone else in the space."),
    ("Press Footprint Most Industry Niches Don't Get",
     "Modern Salon nominated the 2024 piece for Top 10 Article of the year. Authority Magazine, IAMCEO Podcast, Salon Scoop, SF Commercial Property Conversations, Elite Property News. That's the kind of media coverage most consultants would build their entire site around. Right now it sits on a single news page rather than serving as the trust anchor on every primary page."),
    ("Concrete Numbers Most Consultants Can't Quote",
     "$150K in average annual revenue per location. 30-40% savings on initial startup and construction. $100/sqft+ tenant improvement on 10-20 year leases. Investors looking at this category drown in vague pitches. Numbers like these belong in a calculator, not buried in interview transcripts."),
    ("Clear Long-Term Product Roadmap",
     "Allure University, Employee-to-Owner Program with bi-weekly training, M&A activity with other private salon suite operators, expansion across US and Canada. This is product thinking - not a typical consulting practice. The infrastructure to actually capture and onboard each of these audiences hasn't been built yet."),
    ("Brand Sensibility Backed by Design Training",
     "Marywood University BA in Advertising, Graphics, and Fine Arts. The visual instincts are real. The current site doesn't fully express them - which means there's a gap between what Allure could look like and what visitors actually see today. That's a closeable gap, not a structural weakness."),
]
for i, (title, desc) in enumerate(strengths):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s3, Inches(0.55), t, Inches(0.06), Inches(0.7), GREEN)
    tb(s3, title, Inches(0.75), t, Inches(11.8), Inches(0.38), size=14, bold=True, color=WHITE)
    tb(s3, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# SLIDE 4 - WHERE LEADS ARE LEAKING
s4 = prs.slides.add_slide(BLANK)
bg(s4)
rect(s4, 0, 0, Inches(0.12), H, RED)
tb(s4, "WHERE LEADS ARE LEAKING", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=RED)
tb(s4, "Two Audiences. Both End at the Same Phone Tag.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=30, bold=True, color=WHITE)
aline(s4, Inches(1.35), color=RED)

issues = [
    ("PARTNER WITH US PAGE = PHONE NUMBER + EMAIL", "HIGH",
     "An investor researching salon suite consultants on a Sunday night reads the page and sees two phone numbers and an info@ email address. No form, no calendar booking, no automated follow-up. The page asks the visitor to do all the work. Most don't, and the lead never enters the funnel."),
    ("NO INVESTOR ROI CALCULATOR DESPITE A CLEAR BENCHMARK", "HIGH",
     "Karen quotes $150K average annual revenue per location publicly. That number is the strongest investor magnet Allure has, and it lives in interview transcripts rather than on the site. A calculator that takes square footage and metro area as inputs and returns projected revenue is the kind of tool that turns site visits into qualified investor leads."),
    ("PRESS WINS DON'T SHOW UP ON THE HOMEPAGE", "HIGH",
     "Modern Salon Top 10 nominee. Authority Magazine. IAMCEO. SF Commercial. Elite Property News. None of it is on the homepage in a way that builds trust before the first call. Most consulting buyers do their research before they call. The trust signals that should be doing that work are sitting on a separate page."),
    ("ALLURE UNIVERSITY HAS NO INFRASTRUCTURE YET", "MEDIUM",
     "Karen has talked publicly about launching Allure University as an online education platform for salon suite operators. There's no landing page, no signup funnel, no drip course delivery, no completion tracking. The audience exists - the infrastructure to capture and onboard them doesn't."),
    ("TENANT RETENTION ACROSS 115+ STUDIOS IS MANUAL", "MEDIUM",
     "Karen says publicly: 'if your tenants are successful business owners, they stay long term.' That thesis is the whole game. But across 115+ studios in 5 locations, lease renewal sequences, employee-to-owner program nurture, and early-warning churn signals appear to be run manually. Automation here directly protects the recurring revenue base."),
]
for i, (title, pri, desc) in enumerate(issues):
    t = Inches(1.7) + i * Inches(1.05)
    pc = RED if pri == "HIGH" else GOLD
    rect(s4, Inches(0.55), t, Inches(0.06), Inches(0.7), pc)
    tb(s4, title, Inches(0.75), t, Inches(9.2), Inches(0.38), size=12, bold=True, color=WHITE)
    tag(s4, pri, Inches(10.1), t + Inches(0.05), fill=pc, size=9)
    tb(s4, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# SLIDE 5 - WEBSITE DESIGN AUDIT
s5 = prs.slides.add_slide(BLANK)
bg(s5)
rect(s5, 0, 0, Inches(0.12), H, PURPLE)
tb(s5, "WEBSITE DESIGN AUDIT", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=PURPLE)
tb(s5, "5 Locations. 115+ Studios. The Site Doesn't Show It.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.8), size=32, bold=True, color=WHITE)
aline(s5, Inches(1.45), color=PURPLE, w=Inches(1.5))

rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(3.6), DARK_BLUE)
rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(0.42), RED)
tb(s5, "ALLURESALONSUITECONSULTING.COM TODAY", Inches(0.65), Inches(1.92), Inches(3.6), Inches(0.38),
   size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Partner With Us = phone + email, no form",
    "No AI chat or after-hours self-service",
    "No investor ROI calculator on site",
    "Press wins not on homepage",
    "No investor / tenant audience separation",
    "No Allure University funnel",
    "No tenant portal for 115+ studios",
    "Two phone numbers, no unified routing",
], Inches(0.7), Inches(2.42), Inches(3.6), Inches(2.9), size=11, color=LGREY, dot=RED)

tb(s5, "->", Inches(4.45), Inches(3.6), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(3.6), DARK_BLUE)
rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(0.42), GOLD)
tb(s5, "WHAT TOP CONSULTING SITES DO", Inches(5.1), Inches(1.92), Inches(3.4), Inches(0.38),
   size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "ROI calculator with instant revenue projection",
    "AI chat 24/7 for investor + tenant questions",
    "Press wall on homepage building pre-call trust",
    "Separate investor / member / press tracks",
    "Calendar booking embedded on every CTA",
    "Case studies with location-level financials",
    "Online member portal with self-service",
    "Course platform with drip + completion tracking",
], Inches(5.1), Inches(2.42), Inches(3.4), Inches(2.9), size=11, color=LGREY, dot=GOLD)

tb(s5, "->", Inches(8.68), Inches(3.6), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(3.6), DARK_BLUE)
rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(0.42), GREEN)
tb(s5, "WHAT GENOS AI BUILDS", Inches(9.35), Inches(1.92), Inches(3.4), Inches(0.38),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Investor ROI calculator anchored on $150K",
    "AI chat: process, TI negotiation, pricing 24/7",
    "Press wall on homepage + auto-syndication",
    "Separated investor / tenant / press tracks",
    "Lead capture form + calendar booking",
    "Allure University funnel + drip courses",
    "Tenant portal for 115+ studios",
    "Site rebuild with case studies + financials",
], Inches(9.35), Inches(2.42), Inches(3.4), Inches(2.9), size=11, color=LGREY, dot=GREEN)

rect(s5, Inches(0.55), Inches(5.72), Inches(12.23), Inches(1.28), DARK_BLUE)
tb(s5, "WHY IT MATTERS:", Inches(0.7), Inches(5.78), Inches(2.2), Inches(0.38),
   size=11, bold=True, color=PURPLE)
stats2 = [
    ("3-5x", "investor lead capture\nwith ROI calculator vs. plain page"),
    ("60-70%", "of after-hours questions\nresolvable by AI chat"),
    ("90 days", "lease renewal sequence\nstart point for retention"),
    ("8 features", "Modern Salon, IAMCEO,\nAuth Mag - homepage uses 0"),
]
sw = Inches(2.7); sl = Inches(2.9)
for i, (num, label) in enumerate(stats2):
    xl = sl + i * sw
    tb(s5, num, xl, Inches(5.75), Inches(2.5), Inches(0.55),
       size=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s5, label, xl, Inches(6.22), Inches(2.5), Inches(0.7),
       size=10, color=LGREY, align=PP_ALIGN.CENTER)


# SLIDE 6 - AI AUTOMATION OPPORTUNITIES
s6 = prs.slides.add_slide(BLANK)
bg(s6)
rect(s6, 0, 0, Inches(0.12), H, ACCENT)
tb(s6, "AI AUTOMATION OPPORTUNITIES", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s6, "Turning Credibility Into Inbound.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=34, bold=True, color=WHITE)
aline(s6, Inches(1.35), color=ACCENT)

opps = [
    ("Investor ROI\nCalculator", GOLD,
     "Inputs: square footage, metro area, target investor budget. Outputs: projected annual revenue anchored on Allure's $150K-per-location benchmark, projected TI negotiation savings, payback timeline. Captures qualified investor leads in a way no plain contact form can. Replaces vague conversations with anchored numbers."),
    ("AI Chat\n24/7 Self-Service", ACCENT,
     "Trained on Allure's process, TI negotiation playbook, pricing approach, US vs. Canada operations, and Allure Salon Group tenant onboarding. Handles investor and beauty professional inquiries at any hour. Routes complex deals to Karen with full conversation context already captured. Replaces the after-hours dead end."),
    ("Allure University\nFunnel + LMS", GREEN,
     "Landing page, signup sequence, drip course delivery, completion tracking, certificate generation. The infrastructure Karen has talked about publicly but doesn't exist yet. Built once, runs as a permanent acquisition channel for both DIY operators and consulting prospects."),
    ("Tenant Retention\nAutomation", RED,
     "Across 115+ studios at Allure Salon Group: 90-day pre-renewal sequences, employee-to-owner program nurture, early-warning churn signals based on payment patterns and engagement, post-onboarding check-ins. Protects the recurring revenue base directly. Karen's tenant-success thesis applied at scale."),
    ("Press Syndication\n+ Trust Wall", PURPLE,
     "Homepage section featuring Modern Salon, Authority Magazine, IAMCEO, SF Commercial, Elite Property News - ordered for the audience the visitor matches. Automated syndication when the next feature lands: cross-posts to social, newsletter, partner channels in minutes instead of hours."),
]
cw2 = Inches(2.35); gap2 = Inches(0.14)
for i, (title, color, desc) in enumerate(opps):
    l2 = Inches(0.55) + i * (cw2 + gap2)
    rect(s6, l2, Inches(1.8), cw2, Inches(0.06), color)
    rect(s6, l2, Inches(1.86), cw2, Inches(4.6), DARK_BLUE)
    tb(s6, title, l2 + Inches(0.12), Inches(1.98), cw2 - Inches(0.24), Inches(0.65),
       size=13, bold=True, color=color)
    tb(s6, desc, l2 + Inches(0.12), Inches(2.65), cw2 - Inches(0.24), Inches(3.6),
       size=11, color=LGREY)

rect(s6, Inches(0.55), Inches(6.7), Inches(12.23), Inches(0.45), DARK_BLUE)
tb(s6, '"If your tenants are successful business owners, they stay long term." - Karen Kaminski',
   Inches(0.7), Inches(6.75), Inches(12), Inches(0.35), size=11, italic=True, color=GOLD)


# SLIDE 7 - WHAT THIS IS COSTING YOU
s7 = prs.slides.add_slide(BLANK)
bg(s7)
rect(s7, 0, 0, Inches(0.12), H, RED)
tb(s7, "WHAT THIS IS COSTING YOU", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=RED)
tb(s7, "The Brand Is Built. The Digital Layer Isn't Yet.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s7, Inches(1.35), color=RED)

costs = [
    ("No investor capture path = the Sunday night researcher never enters your funnel",
     "An investor with capital, looking at salon suites as a passive-income asset, lands on the Partner With Us page after hours. Phone number, email, no form, no calculator, no calendar. They don't call. They go back to Google and click the next consultant. Allure never sees the lead. The credibility doesn't matter if the front door doesn't open after 5pm."),
    ("No ROI calculator = your strongest number lives in interview transcripts",
     "$150K average annual revenue per location is the kind of anchor that turns interest into action. Right now it appears in podcast transcripts and Authority Magazine, not on the site where investors are evaluating. A calculator with that number as the anchor is the difference between a brochure site and a lead engine."),
    ("Press wins not on the homepage = trust signals doing no work pre-call",
     "Modern Salon Top 10 nominee. Authority Magazine. IAMCEO. SF Commercial. Elite Property News. Every one of these is the kind of credential that closes the trust gap before the first conversation. Buried on a separate page, they reach almost no one. Surfaced on the homepage, they shorten every sales cycle by removing the credibility question."),
    ("No Allure University infrastructure = an audience built and not captured",
     "Karen has been talking about Allure University for at least a year. The audience for it is being built every time she does a podcast or feature. Without a landing page, signup funnel, and drip course delivery, that audience is being introduced to the idea and not given anywhere to go next."),
    ("Manual tenant retention across 115+ studios = revenue at risk",
     "Lease renewal sequences, employee-to-owner program nurture, churn early-warning. At 115+ studios across 5 locations, manual handling has gaps. A tenant who goes quiet at month 14 of an 18-month lease should trigger a sequence, not wait for someone to notice. Karen's own thesis on tenant success points directly at this."),
]
for i, (title, desc) in enumerate(costs):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s7, Inches(0.55), t, Inches(0.06), Inches(0.7), RED)
    tb(s7, title, Inches(0.75), t, Inches(11.8), Inches(0.38), size=14, bold=True, color=WHITE)
    tb(s7, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# SLIDE 8 - WHAT WE'D BUILD
s8 = prs.slides.add_slide(BLANK)
bg(s8)
rect(s8, 0, 0, Inches(0.12), H, GREEN)
tb(s8, "WHAT WE'D BUILD", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s8, "Genos AI for Allure Salon",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s8, Inches(1.35), color=GREEN)

builds = [
    ("Investor ROI Calculator + Qualified Lead Capture",
     "Embedded on Partner With Us and the homepage. Square footage and metro area in, projected annual revenue (anchored on the $150K benchmark), TI negotiation savings, and payback timeline out. Calculator output triggers a lead capture form and an automated follow-up sequence routed to Karen with the prospect's calculator inputs already attached."),
    ("AI Chat - 24/7 Investor + Tenant Self-Service",
     "Chat widget trained on Allure's process, TI negotiation, pricing approach, US vs. Canada operations, and Allure Salon Group tenant onboarding. Routes complex investor inquiries to Karen with conversation context. Handles tenant questions at Allure Salon Group locations. Closes the after-hours gap on both audiences."),
    ("Allure University - Funnel, LMS, Drip Delivery",
     "Landing page, signup sequence, drip course delivery, completion tracking, certificate generation. The platform Karen has talked about, built. Connects directly to the consulting funnel: course completers become qualified consulting prospects with already-demonstrated commitment."),
    ("Tenant Retention Automation Across 115+ Studios",
     "Pre-renewal sequences starting 90 days before lease end. Employee-to-Owner program drip nurture. Churn early-warning signals based on payment timing, communication response, and engagement patterns. Post-onboarding check-ins at 30, 60, and 90 days. Protects the recurring revenue base where it lives."),
    ("Site Rebuild + Press Syndication + Trust Wall",
     "alluresalonsuiteconsulting.com rebuilt with separated investor and tenant tracks, press wall on homepage (Modern Salon, Authority Magazine, IAMCEO, SF Commercial, Elite Property News), case studies with location-level financials, calendar booking on every CTA. Press syndication automation when the next feature lands."),
]
for i, (title, desc) in enumerate(builds):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s8, Inches(0.55), t, Inches(0.06), Inches(0.7), GREEN)
    tb(s8, title, Inches(0.75), t, Inches(11.8), Inches(0.38), size=14, bold=True, color=WHITE)
    tb(s8, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# SLIDE 9 - 90-DAY OUTCOMES
s9 = prs.slides.add_slide(BLANK)
bg(s9)
rect(s9, 0, 0, Inches(0.12), H, ACCENT)
tb(s9, "90-DAY OUTCOMES", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s9, "What Changes in 90 Days",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s9, Inches(1.35))

phases = [
    ("Month 1", GOLD, [
        "AI chat live on alluresalonsuiteconsulting.com - 24/7 investor and tenant questions handled",
        "Investor ROI calculator deployed on Partner With Us, anchored on $150K benchmark",
        "Lead capture form replacing phone-and-email-only contact, calendar booking embedded",
        "Press wall added to homepage - Modern Salon, Authority Magazine, IAMCEO, SF Commercial",
    ]),
    ("Month 2", ACCENT, [
        "Allure University landing page live with first drip course active",
        "Tenant retention sequences running across Allure Salon Group's 115+ studios",
        "Employee-to-Owner program nurture sequences deployed",
        "First measurable inbound investor leads from ROI calculator submissions",
    ]),
    ("Month 3", GREEN, [
        "Site rebuild complete with separated investor / tenant tracks and case studies",
        "Press syndication automation in production - next feature spreads in minutes",
        "Cross-border (US / Canada) investor capture flows live with geo-segmentation",
        "Tenant churn early-warning signals running, first save data measurable",
    ]),
]
pw = Inches(3.9); gap3 = Inches(0.18)
for i, (month, color, pts) in enumerate(phases):
    l3 = Inches(0.55) + i * (pw + gap3)
    rect(s9, l3, Inches(1.8), pw, Inches(0.42), color)
    tb(s9, month.upper(), l3 + Inches(0.1), Inches(1.84), pw - Inches(0.2), Inches(0.36),
       size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    rect(s9, l3, Inches(2.22), pw, Inches(4.7), DARK_BLUE)
    bullets(s9, pts, l3 + Inches(0.15), Inches(2.35), pw - Inches(0.3), Inches(4.4),
            size=12, color=LGREY, dot=color)


# SLIDE 10 - NEXT STEPS
s10 = prs.slides.add_slide(BLANK)
bg(s10)
rect(s10, 0, 0, Inches(0.12), H, GOLD)
rect(s10, Inches(7.5), 0, Inches(5.83), H, DARK_BLUE)
tb(s10, "NEXT STEPS", Inches(0.55), Inches(0.3), Inches(6.5), Inches(0.5),
   size=11, bold=True, color=GOLD)
tb(s10, "Let's Talk for\n15 Minutes.",
   Inches(0.55), Inches(0.8), Inches(6.8), Inches(2.0), size=48, bold=True, color=WHITE)
aline(s10, Inches(2.8), color=GOLD, w=Inches(2.5))
tb(s10, "No pitch deck required. Reply or book a call directly.\nWe'll walk through what fits Allure Salon Suite Consulting and Allure Salon Group.",
   Inches(0.55), Inches(3.05), Inches(6.5), Inches(0.8), size=14, color=LGREY)

steps = [
    "Reply to this email - we'll set up a time",
    "Book directly: www.genosai.tech/call",
    "Walk through the audit findings in 15 minutes",
    "You decide what, if anything, makes sense to move on",
]
bullets(s10, steps, Inches(0.55), Inches(3.95), Inches(6.5), Inches(2.0), size=14, color=WHITE, dot=GOLD)

rect(s10, 0, Inches(6.5), Inches(7.5), Inches(1.0), DARK_BLUE)
tb(s10, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699  |  www.genosai.tech",
   Inches(0.55), Inches(6.6), Inches(6.8), Inches(0.5), size=12, color=MGREY)

tb(s10, "WHAT WE BUILD FOR\nALLURE SALON", Inches(7.85), Inches(0.6), Inches(5.0), Inches(0.9),
   size=18, bold=True, color=GOLD)
aline(s10, Inches(1.45), color=GOLD)

summary_items = [
    ("ROI Calculator", "Anchored on $150K - turns interest into qualified leads"),
    ("AI Chat", "24/7 investor + tenant self-service, full context routing"),
    ("Allure U", "Funnel, drip courses, completion tracking - built and live"),
    ("Retention", "Lease renewal + churn early-warning across 115+ studios"),
    ("Press Wall", "Homepage trust signals + automated syndication"),
    ("Site Rebuild", "Investor / tenant tracks, case studies, calendar booking"),
]
for i, (label, val) in enumerate(summary_items):
    t2 = Inches(1.6) + i * Inches(0.88)
    rect(s10, Inches(7.85), t2, Inches(1.5), Inches(0.72), NAVY)
    tb(s10, label, Inches(7.9), t2 + Inches(0.05), Inches(1.4), Inches(0.32),
       size=10, bold=True, color=GOLD)
    tb(s10, val, Inches(9.5), t2 + Inches(0.05), Inches(3.6), Inches(0.65),
       size=11, color=LGREY)


prs.save(ROOT / "output" / "decks" / "AllureSalon_Audit_GenosAI.pptx")
print("Saved: AllureSalon_Audit_GenosAI.pptx")
