import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.parent / "utils"))
from ppt_theme import (
    NAVY, DARK_BLUE, ACCENT, GOLD, RED, GREEN, PURPLE, ORANGE, TEAL, CYAN,
    WHITE, LGREY, MGREY, W, H,
    bg, rect, tb, aline, bullets, tag,
    slide_header, slide_footer, score_badge,
    stat_cards, top_bar_card, phase_card, before_after, genos_logo_block,
)

ROOT = Path(__file__).parent.parent.parent
prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]
# Client-specific colors
FOREST = RGBColor(0x1D, 0x6B, 0x3C)
DGREY  = RGBColor(0x2C, 0x38, 0x4A)
AMBER  = RGBColor(0xF3, 0x9C, 0x12)
SLATE  = RGBColor(0x2E, 0x45, 0x5F)


# ── Local helpers ────────────────────────────────────────────────────────────

def header(slide, title, subtitle, accent=GREEN):
    rect(slide, 0, 0, Inches(0.12), H, accent)
    tb(slide, title.upper(), Inches(0.55 + 0.12), Inches(0.3),
       Inches(12), Inches(0.45), size=11, bold=True, color=accent)
    tb(slide, subtitle, Inches(0.55 + 0.12), Inches(0.72),
       Inches(12 - 0.12), Inches(0.65), size=22, bold=True, color=WHITE)
    aline(slide, Inches(1.38), color=accent, w=Inches(1.5))


def footer(slide,
           text='Genos AI  |  hello@genosai.tech  |  www.genosai.tech  |  +91 638-714-2699'):
    rect(slide, 0, Inches(7.05), W, Inches(0.45), DARK_BLUE)
    tb(slide, text, Inches(0.55), Inches(7.1), Inches(12.4), Inches(0.35),
       size=9, color=MGREY, align=PP_ALIGN.CENTER)


def card(slide, x, y, w, h, fill=DARK_BLUE):
    return rect(slide, x, y, w, h, fill)


def stat_card(slide, val, label, x, y, w=Inches(3.0), h=Inches(1.6), accent=GREEN):
    rect(slide, x, y, w, h, DARK_BLUE)
    rect(slide, x, y, w, Inches(0.05), accent)
    tb(slide, str(val), x + Inches(0.1), y + Inches(0.1),
       w - Inches(0.2), Inches(0.9), size=28, bold=True, color=accent,
       align=PP_ALIGN.CENTER)
    tb(slide, label, x + Inches(0.1), y + Inches(1.1),
       w - Inches(0.2), Inches(0.42), size=10, color=LGREY, align=PP_ALIGN.CENTER)


def bullet_card(slide, title, items, x, y, w, h,
                accent=GREEN, title_bg=None):
    if title_bg is None:
        title_bg = DARK_BLUE
    rect(slide, x, y, w, h, DARK_BLUE)
    rect(slide, x, y, w, Inches(0.42), title_bg)
    rect(slide, x, y, Inches(0.05), h, accent)
    tb(slide, title, x + Inches(0.1), y + Inches(0.05),
       w - Inches(0.15), Inches(0.35), size=11, bold=True, color=WHITE)
    bx = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.48),
                                  w - Inches(0.15), h - Inches(0.55))
    tf = bx.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r0 = p.add_run(); r0.text = '  '; r0.font.size = Pt(9)
        r1 = p.add_run(); r1.text = '● '; r1.font.size = Pt(9)
        r1.font.color.rgb = accent; r1.font.bold = True
        r2 = p.add_run(); r2.text = str(item)
        r2.font.size = Pt(10.5); r2.font.color.rgb = LGREY
        if i < len(items) - 1:
            p.space_after = Pt(5)




# SLIDE 1 – TITLE
# ════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
bg(s1, FOREST)
rect(s1, 0, 0, Inches(0.7), H, ORANGE)
rect(s1, Inches(0.7), Inches(4.5), W - Inches(0.7), Inches(0.06), ORANGE)

tb(s1, "IDILI AFRICA",
   Inches(1.1), Inches(1.0), Inches(11), Inches(0.7),
   size=16, bold=True, color=ORANGE)
tb(s1, "Web Development &\nAI Automation Audit",
   Inches(1.1), Inches(1.7), Inches(11), Inches(2.0),
   size=44, bold=True, color=WHITE)
tb(s1, "Strategic Review & Digital Growth Roadmap  |  Cape Town, South Africa",
   Inches(1.1), Inches(3.72), Inches(10), Inches(0.5),
   size=16, color=ORANGE, italic=True)
tb(s1, "Prepared by: Genos Apollo  |  May 2026",
   Inches(1.1), Inches(5.0), Inches(8), Inches(0.4),
   size=12, color=MGREY)
tb(s1, "CONFIDENTIAL",
   Inches(1.1), Inches(6.6), Inches(4), Inches(0.35),
   size=10, color=MGREY, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 2 – AGENDA
# ════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
bg(s2)
header(s2, "Agenda", "What we will cover today")
footer(s2)

agenda = [
    ("01", "Company Digital Overview",    "Current state of idili.africa and digital footprint"),
    ("02", "Website Audit Findings",      "UX, SEO, platform gaps and compliance issues"),
    ("03", "Web Development Roadmap",     "Talent portal, CMS & modern platform stack"),
    ("04", "AI Automation Opportunity",   "Where AI accelerates talent matching & operations"),
    ("05", "Benefits of Automation",      "ROI, placement velocity and partner experience"),
    ("06", "Implementation Timeline",     "Phased rollout plan across 12 months"),
    ("07", "Investment & Next Steps",     "Budget guidance and immediate actions"),
]

for i, (num, title, desc) in enumerate(agenda):
    col = i % 2
    row = i // 2
    x = Inches(0.55) + col * Inches(6.4)
    y = Inches(1.55) + row * Inches(1.22)
    card(s2, x, y, Inches(6.1), Inches(1.05))
    rect(s2, x, y, Inches(0.7), Inches(1.05), FOREST)
    tb(s2, num, x, y + Inches(0.25), Inches(0.7), Inches(0.5),
       size=20, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    tb(s2, title, x + Inches(0.8), y + Inches(0.1),
       Inches(5.0), Inches(0.4), size=13, bold=True, color=FOREST)
    tb(s2, desc, x + Inches(0.8), y + Inches(0.52),
       Inches(5.0), Inches(0.42), size=10.5, color=MGREY, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 3 – COMPANY DIGITAL OVERVIEW
# ════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
bg(s3)
header(s3, "Company Digital Overview", "Idili Africa – current digital footprint")
footer(s3)

card(s3, Inches(0.4), Inches(1.4), Inches(5.8), Inches(5.5))
rect(s3, Inches(0.4), Inches(1.4), Inches(5.8), Inches(0.45), FOREST)
tb(s3, "Company Snapshot", Inches(0.55), Inches(1.44),
   Inches(5.5), Inches(0.38), size=13, bold=True, color=WHITE)

snapshot = [
    ("Mission",       '"Talent is global. Opportunity should be too."'),
    ("Industry",      "International Talent Mobility & Graduate Recruitment"),
    ("HQ",            "Innovation City Cape Town, Darters Road, 8001"),
    ("Focus",         "SA graduates (18–25) → Swedish / global companies"),
    ("Partners",      "Saab · Atlas Copco · Bitprop · REY17 · Oakvale Invest"),
    ("Affiliations",  "Sweden Consulate Cape Town · Made with Sweden initiative"),
    ("Extra Service", "Seed funding for African AI founders via Oakvale Invest"),
    ("Website",       "idili.africa  (custom domain ✓ | platform TBC)"),
]
ty = Inches(2.02)
for k, v in snapshot:
    tb(s3, k, Inches(0.6), ty, Inches(1.55), Inches(0.3),
       size=10.5, bold=True, color=FOREST)
    tb(s3, v, Inches(2.2), ty, Inches(3.8), Inches(0.3),
       size=10.5, color=MGREY)
    ty += Inches(0.52)

card(s3, Inches(6.5), Inches(1.4), Inches(6.5), Inches(5.5))
rect(s3, Inches(6.5), Inches(1.4), Inches(6.5), Inches(0.45), ORANGE)
tb(s3, "Key Digital Health Indicators", Inches(6.65), Inches(1.44),
   Inches(6.2), Inches(0.38), size=13, bold=True, color=WHITE)

metrics = [
    ("Custom Domain",              "✓ Active",    "idili.africa live",                     GREEN),
    ("Mobile Page Speed",          "~52 / 100",   "Est. — CMS/platform not confirmed",     AMBER),
    ("SEO Visibility",             "Low",         "Thin content, minimal keyword strategy", RED),
    ("Candidate Application Portal","None",       "No ATS or talent portal on site",        RED),
    ("Company Partner Portal",     "None",        "No dashboard for hiring partners",       RED),
    ("Social Media on Site",       "None",        "No LinkedIn / Instagram links visible",  RED),
    ("Privacy Policy / POPIA",     "None",        "No legal footer or compliance notice",   RED),
    ("Email / Lead Capture",       "Form only",   "No newsletter or talent pipeline CRM",   AMBER),
]
ty2 = Inches(2.02)
for label, val, note, col in metrics:
    tb(s3, label, Inches(6.65), ty2, Inches(2.7), Inches(0.3),
       size=10.5, bold=True, color=FOREST)
    tb(s3, val,   Inches(9.4), ty2, Inches(1.25), Inches(0.3),
       size=10.5, bold=True, color=col, align=PP_ALIGN.CENTER)
    tb(s3, note,  Inches(10.7), ty2, Inches(2.1), Inches(0.3),
       size=9, color=MGREY, italic=True)
    ty2 += Inches(0.52)

# ════════════════════════════════════════════════════════════════
# SLIDE 4 – WEBSITE AUDIT FINDINGS
# ════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
bg(s4)
header(s4, "Website Audit Findings", "idili.africa – gaps blocking talent pipeline growth")
footer(s4)

audit_cards = [
    ("Platform & Performance",
     ["CMS/platform not publicly identifiable", "Subpages return 404 (fragile URL structure)",
      "No CDN or edge caching confirmed", "Mobile performance likely under-optimised"]),
    ("UX & Platform Features",
     ["No applicant tracking or talent portal", "Opportunities link to external company pages",
      "No partner company dashboard or login", "No WhatsApp or live chat for candidate queries"]),
    ("SEO & Content",
     ["4 pages only – extremely thin content", "No blog, success stories or placement stats",
      "Missing schema (JobPosting, Organization)", "No keyword targeting for graduate job queries"]),
    ("Trust & Compliance",
     ["No privacy policy or POPIA notice", "No cookie consent mechanism on site",
      "No testimonials from placed graduates", "No social media links – isolated digital presence"]),
]

for i, (title, bullets) in enumerate(audit_cards):
    x = Inches(0.4) + i * Inches(3.2)
    y = Inches(1.45)
    card(s4, x, y, Inches(3.05), Inches(5.3))
    rect(s4, x, y, Inches(3.05), Inches(0.45), FOREST)
    tb(s4, title, x + Inches(0.1), y + Inches(0.06),
       Inches(2.85), Inches(0.36), size=12, bold=True, color=WHITE)
    ty = y + Inches(0.62)
    for b in bullets:
        rect(s4, x + Inches(0.15), ty + Inches(0.07), Inches(0.1), Inches(0.1), ORANGE)
        tb(s4, b, x + Inches(0.35), ty, Inches(2.6), Inches(0.55),
           size=10, color=FOREST)
        ty += Inches(0.72)

rect(s4, Inches(0.4), Inches(6.8), Inches(12.5), Inches(0.06), ORANGE)
tb(s4, "Priority: Build a proper talent portal + POPIA compliance + SEO content strategy to scale the candidate pipeline",
   Inches(0.4), Inches(6.86), Inches(12.5), Inches(0.3),
   size=10.5, bold=True, color=FOREST, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 5 – WEB DEVELOPMENT ROADMAP
# ════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
bg(s5)
header(s5, "Web Development Roadmap", "Transforming idili.africa into a scalable talent platform")
footer(s5)

card(s5, Inches(0.4), Inches(1.4), Inches(4.5), Inches(5.6))
rect(s5, Inches(0.4), Inches(1.4), Inches(4.5), Inches(0.45), TEAL)
tb(s5, "Recommended Tech Stack", Inches(0.55), Inches(1.44),
   Inches(4.2), Inches(0.38), size=13, bold=True, color=WHITE)

stack = [
    ("Frontend",    "Next.js 15 – SSR/SSG for SEO-rich talent pages"),
    ("CMS",         "Sanity.io – manage opportunities, stories, content"),
    ("ATS",         "Greenhouse (free tier) or custom Sanity portal"),
    ("Hosting",     "Vercel + Cloudflare CDN (SA + EU edge nodes)"),
    ("Auth",        "Auth0 – candidate & partner company portals"),
    ("Email CRM",   "Brevo (free tier) – candidate nurture sequences"),
    ("Analytics",   "GA4 + Hotjar – funnel & application drop-off"),
    ("WhatsApp",    "WhatsApp Business API – candidate comms"),
]
ty = Inches(2.02)
for k, v in stack:
    tb(s5, k, Inches(0.6), ty, Inches(1.4), Inches(0.3),
       size=10, bold=True, color=TEAL)
    tb(s5, v, Inches(2.05), ty, Inches(2.75), Inches(0.3),
       size=10, color=FOREST)
    ty += Inches(0.5)

card(s5, Inches(5.2), Inches(1.4), Inches(7.75), Inches(5.6))
rect(s5, Inches(5.2), Inches(1.4), Inches(7.75), Inches(0.45), FOREST)
tb(s5, "Key Improvements & Features",
   Inches(5.35), Inches(1.44), Inches(7.4), Inches(0.38),
   size=13, bold=True, color=WHITE)

improvements = [
    ("Candidate Talent Portal",
     "Graduates create profiles, upload CVs, track application status, receive assessment invites, and access placement resources in one secure dashboard."),
    ("Partner Company Dashboard",
     "Saab, Atlas Copco and future partners log in to view shortlisted candidates, schedule interviews, and track pipeline status — no manual email chains."),
    ("SEO-Optimised Opportunity Listings",
     "Each role becomes a Google-indexable page with JobPosting schema, salary range, skills required, and an apply-now CTA. Target 'graduate jobs South Africa Sweden'."),
    ("Success Stories & Impact Hub",
     "Placed graduates featured by name, company and journey. Tracks cumulative stats: X graduates placed, X countries, X partner companies — building social proof."),
    ("POPIA Compliance & Privacy Framework",
     "Full cookie consent, privacy policy, data processing agreements with partners, and candidate data rights management in line with POPIA 2021."),
    ("WhatsApp Application Journey",
     "Candidates apply or enquire via WhatsApp. Auto-guided flow: eligibility check → document upload → assessment booking → status updates."),
]
ty2 = Inches(1.98)
for title, desc in improvements:
    rect(s5, Inches(5.35), ty2 + Inches(0.1), Inches(0.08), Inches(0.25), ORANGE)
    tb(s5, title, Inches(5.55), ty2, Inches(7.1), Inches(0.28),
       size=11, bold=True, color=FOREST)
    tb(s5, desc, Inches(5.55), ty2 + Inches(0.3), Inches(7.1), Inches(0.35),
       size=10, color=MGREY)
    ty2 += Inches(0.78)

# ════════════════════════════════════════════════════════════════
# SLIDE 6 – AI AUTOMATION OPPORTUNITY MAP
# ════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
bg(s6)
header(s6, "AI Automation Opportunity Map", "Where AI accelerates Idili Africa's talent mission")
footer(s6)

opps = [
    (FOREST,  "AI Talent Screening & Matching",
     ["AI CV parser + skills extraction on upload", "Auto-score candidates against role requirements",
      "Psychometric pre-screening bot (before human review)", "AI-match candidates to best-fit partner companies"]),
    (ORANGE,  "Candidate Communications",
     ["WhatsApp bot: 24/7 application status updates", "Auto-send assessment links & interview schedules",
      "Personalised placement journey nudges (email/WA)", "Post-placement check-ins & feedback collection"]),
    (TEAL,    "Partner & Employer Experience",
     ["AI-generated shortlist summaries for hiring managers", "Auto-notify partners when new candidates qualify",
      "Weekly pipeline reports auto-emailed to partners", "Interview debrief templates auto-generated"]),
    (SLATE,   "Content, SEO & Growth",
     ["AI-draft graduate success stories from a template", "Auto-post placements to LinkedIn + Instagram",
      "Blog content: 'Working in Sweden as a SA engineer'", "Monitor SA university graduation calendars (AI digest)"]),
]

for i, (col, title, bullets) in enumerate(opps):
    c = i % 2
    r = i // 2
    x = Inches(0.4) + c * Inches(6.5)
    y = Inches(1.45) + r * Inches(2.85)
    card(s6, x, y, Inches(6.2), Inches(2.65))
    rect(s6, x, y, Inches(6.2), Inches(0.45), col)
    tb(s6, title, x + Inches(0.15), y + Inches(0.06),
       Inches(5.9), Inches(0.36), size=13, bold=True,
       color=FOREST if col == ORANGE else WHITE)
    ty = y + Inches(0.62)
    for b in bullets:
        rect(s6, x + Inches(0.18), ty + Inches(0.08), Inches(0.1), Inches(0.1), col)
        tb(s6, b, x + Inches(0.38), ty, Inches(5.65), Inches(0.38),
           size=10.5, color=FOREST)
        ty += Inches(0.48)

# ════════════════════════════════════════════════════════════════
# SLIDE 7 – BENEFITS OF AI AUTOMATION
# ════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
bg(s7)
header(s7, "Benefits of AI Automation", "Measurable impact across Idili Africa's talent pipeline")
footer(s7)

stats = [
    ("5×",      "More candidates\nscreened per week"),
    ("70%",     "Reduction in\nmanual admin time"),
    ("24/7",    "WhatsApp candidate\nsupport coverage"),
    ("3×",      "Faster shortlist\ndelivery to partners"),
]
for i, (val, label) in enumerate(stats):
    stat_card(s7, val, label,
              Inches(0.4) + i * Inches(3.2), Inches(1.45),
              w=Inches(3.0), h=Inches(1.6))

benefits = [
    ("Talent Pipeline Scale",
     ["AI pre-screens CVs — human time on best fits only", "Handle 10× volume without adding headcount",
      "Auto-qualify candidates before assessment day", "Continuous intake, not just cohort-based"]),
    ("Partner Experience",
     ["Partners get shortlists, not raw applicant dumps", "Real-time portal beats weekly email updates",
      "AI summaries reduce partner review time by 60%", "Stronger relationships → more partner companies"]),
    ("Candidate Experience",
     ["WhatsApp guidance from apply to placement", "No candidate left wondering about their status",
      "Personalised comms in English (+ local languages)", "Better experience → higher referral rate"]),
    ("Organisational Impact",
     ["Data on every cohort: conversion rates, placement speed", "Identify top-performing universities & faculties",
      "Build case for grant funding with AI-tracked outcomes", "Scale to 5 countries without proportional cost"]),
]

for i, (title, bullets) in enumerate(benefits):
    c = i % 2
    r = i // 2
    x = Inches(0.4) + c * Inches(6.5)
    y = Inches(3.2) + r * Inches(1.85)
    bullet_card(s7, title, bullets, x, y, Inches(6.2), Inches(1.65),
                title_bg=FOREST if r == 0 else TEAL)

# ════════════════════════════════════════════════════════════════
# SLIDE 8 – ROI & BUSINESS CASE
# ════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
bg(s8)
header(s8, "ROI & Business Case", "Projected returns from platform + AI investment (USD)")
footer(s8)

card(s8, Inches(0.4), Inches(1.4), Inches(5.8), Inches(5.6))
rect(s8, Inches(0.4), Inches(1.4), Inches(5.8), Inches(0.45), FOREST)
tb(s8, "Indicative Investment (Year 1 – USD)",
   Inches(0.55), Inches(1.44), Inches(5.5), Inches(0.38),
   size=13, bold=True, color=WHITE)

invest_rows = [
    ("Website Rebuild (Next.js + Sanity CMS)",   "$2,500 – $ 5,000"),
    ("Talent Portal & Partner Dashboard",         "$2,000 – $ 4,000"),
    ("AI Screening & WhatsApp Bot Setup",         "$1,200 – $ 2,500"),
    ("SEO Content Sprint & JobPosting Schema",    "$  600 – $ 1,200"),
    ("POPIA Compliance & Privacy Framework",      "$  400 – $   800"),
    ("CRM / Email Automation (Brevo)",            "$  300 – $   600"),
    ("Ongoing monthly support (est.)",            "$  300 – $   600 / mo"),
    ("TOTAL YEAR 1 (once-off + 12 mo)",          "$10,600 – $21,300"),
]
ty = Inches(2.0)
for i, (item, cost) in enumerate(invest_rows):
    is_total = i == len(invest_rows) - 1
    if is_total:
        rect(s8, Inches(0.4), ty - Inches(0.05), Inches(5.8), Inches(0.06), ORANGE)
        ty += Inches(0.12)
        tb(s8, item, Inches(0.6), ty, Inches(3.5), Inches(0.35),
           size=11, bold=True, color=FOREST)
        tb(s8, cost, Inches(4.05), ty, Inches(1.95), Inches(0.35),
           size=11, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)
    else:
        tb(s8, item, Inches(0.6), ty, Inches(3.5), Inches(0.3),
           size=10, color=FOREST)
        tb(s8, cost, Inches(4.05), ty, Inches(1.95), Inches(0.3),
           size=10, color=MGREY, align=PP_ALIGN.RIGHT)
    ty += Inches(0.5)

card(s8, Inches(6.5), Inches(1.4), Inches(6.5), Inches(5.6))
rect(s8, Inches(6.5), Inches(1.4), Inches(6.5), Inches(0.45), GREEN)
tb(s8, "Projected Returns & Value Unlocked",
   Inches(6.65), Inches(1.44), Inches(6.2), Inches(0.38),
   size=13, bold=True, color=WHITE)

returns = [
    ("Placement fee per candidate (est.)",         "$1,500 – $4,000 per hire"),
    ("10× more candidates screened / month",       "Unlock higher volume → more placements"),
    ("Reduced manual screening time saved",        "15–20 hrs/week reclaimed"),
    ("New partner companies attracted by portal",  "2–4 new partners / year est."),
    ("Grant & impact funding credibility",         "Data-tracked outcomes unlock NGO/govt funding"),
    ("Seed funding pipeline (Oakvale / AI arm)",   "Platform data qualifies founder cohort quality"),
    ("BREAKEVEN (3–7 additional placements)",      "Full Year 1 investment recovered"),
]
ty2 = Inches(2.0)
for i, (item, val) in enumerate(returns):
    is_total = i == len(returns) - 1
    if is_total:
        rect(s8, Inches(6.5), ty2 - Inches(0.05), Inches(6.5), Inches(0.06), ORANGE)
        ty2 += Inches(0.12)
        tb(s8, item, Inches(6.65), ty2, Inches(3.9), Inches(0.35),
           size=11, bold=True, color=FOREST)
        tb(s8, val, Inches(10.6), ty2, Inches(2.2), Inches(0.35),
           size=11, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)
    else:
        tb(s8, item, Inches(6.65), ty2, Inches(3.9), Inches(0.3),
           size=10, color=FOREST)
        tb(s8, val, Inches(10.6), ty2, Inches(2.2), Inches(0.3),
           size=10, color=MGREY, align=PP_ALIGN.RIGHT)
    ty2 += Inches(0.5)

tb(s8, "3–7 additional placements per year at $1.5K–$4K each recovers the full investment — everything beyond is pure growth",
   Inches(0.4), Inches(7.07), Inches(12.5), Inches(0.28),
   size=10.5, bold=True, color=FOREST, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 9 – IMPLEMENTATION TIMELINE
# ════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(BLANK)
bg(s9)
header(s9, "Implementation Timeline", "Phased 12-month delivery plan")
footer(s9)

phases = [
    ("Phase 1\nMon 1–2",  ORANGE,   [
        "POPIA privacy policy & cookie consent",
        "Social media links added to site",
        "WhatsApp Business API connected",
        "Google Search Console + GA4 setup",
    ]),
    ("Phase 2\nMon 3–5",  FOREST,   [
        "Next.js + Sanity CMS rebuild launch",
        "SEO opportunity pages (JobPosting schema)",
        "Candidate talent portal v1 live",
        "Email CRM pipeline (Brevo) active",
    ]),
    ("Phase 3\nMon 6–8",  TEAL,     [
        "Partner company dashboard live",
        "AI CV screening & match scoring",
        "Success stories hub + placement stats",
        "LinkedIn + Instagram auto-posting",
    ]),
    ("Phase 4\nMon 9–12", GREEN,    [
        "WhatsApp full application journey bot",
        "AI psychometric pre-screening flow",
        "University partnership landing pages",
        "Full audit, review & optimise",
    ]),
]

for i, (label, col, items) in enumerate(phases):
    x = Inches(0.4) + i * Inches(3.2)
    y = Inches(1.45)
    card(s9, x, y, Inches(3.0), Inches(5.55))
    rect(s9, x, y, Inches(3.0), Inches(0.7), col)
    tb(s9, label, x, y + Inches(0.06), Inches(3.0), Inches(0.62),
       size=13, bold=True,
       color=FOREST if col == ORANGE else WHITE,
       align=PP_ALIGN.CENTER)
    if i < 3:
        rect(s9, x + Inches(3.02), y + Inches(0.3), Inches(0.16), Inches(0.06), MGREY)
    ty = y + Inches(0.88)
    for item in items:
        rect(s9, x + Inches(0.18), ty + Inches(0.08), Inches(0.1), Inches(0.1), col)
        tb(s9, item, x + Inches(0.38), ty, Inches(2.5), Inches(0.42),
           size=10.5, color=FOREST)
        ty += Inches(0.9)

# ════════════════════════════════════════════════════════════════
# SLIDE 10 – WHY ACT NOW?
# ════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(BLANK)
bg(s10, FOREST)
rect(s10, 0, 0, Inches(0.5), H, ORANGE)
footer(s10)

tb(s10, "Why Act — Now?",
   Inches(0.8), Inches(0.5), Inches(11), Inches(0.65),
   size=30, bold=True, color=ORANGE)
tb(s10, "Africa produces 12M+ graduates annually. The pipeline is vast. The platform to harness it professionally is not yet built.",
   Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.55),
   size=14, color=WHITE, italic=True)

reasons = [
    ("The talent pipeline is real — the infrastructure isn't",
     "Idili has the partnerships (Saab, Atlas Copco, Sweden Consulate) and the mission. What's missing is the platform to scale from dozens of candidates to thousands. Every month without a portal is a month of manually managed applications that could be automated."),
    ("European skills shortages are intensifying",
     "Sweden and broader Europe face critical shortages in engineering, software, and technology through 2030. The window to establish Idili as the go-to Africa-to-Europe talent bridge is open now — and competitors will emerge."),
    ("No social media links on the site means no trust loop",
     "Graduates and companies research before engaging. A site with no LinkedIn link, no success stories, and no privacy policy signals 'unfinished' — even if the mission and partnerships are exceptional. This is a fast, cheap fix."),
    ("POPIA compliance is non-negotiable for a data-rich platform",
     "Idili handles personal data: CVs, psychometric results, interview notes. Operating without a privacy policy or POPIA notice is a legal exposure risk. Fixing this also unlocks trust with European partner companies who have GDPR obligations."),
    ("AI screening turns a boutique operation into a scalable one",
     "Today, candidate screening requires human hours. AI can pre-qualify 10× the volume before a single recruiter reviews a profile. That efficiency multiplier is the difference between placing 50 graduates per year and placing 500."),
]
ty = Inches(1.95)
for i, (title, desc) in enumerate(reasons):
    rect(s10, Inches(0.8), ty, Inches(0.55), Inches(0.55),
         ORANGE if i % 2 == 0 else TEAL)
    tb(s10, str(i + 1), Inches(0.8), ty, Inches(0.55), Inches(0.55),
       size=16, bold=True, color=FOREST if i % 2 == 0 else WHITE,
       align=PP_ALIGN.CENTER)
    tb(s10, title, Inches(1.5), ty, Inches(11), Inches(0.3),
       size=12, bold=True, color=WHITE)
    tb(s10, desc, Inches(1.5), ty + Inches(0.3), Inches(11.3), Inches(0.55),
       size=10.5, color=RGBColor(0xBB, 0xCC, 0xDD))
    ty += Inches(0.98)

# ════════════════════════════════════════════════════════════════
# SLIDE 11 – NEXT STEPS
# ════════════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(BLANK)
bg(s11)
header(s11, "Next Steps", "From audit to action — what happens now")
footer(s11)

steps = [
    ("This Week",  ORANGE,  "Quick Wins (48 hrs)",
     "Add LinkedIn & social links to site. Add a basic privacy policy. Connect WhatsApp Business number. Zero cost, immediate trust uplift."),
    ("Week 2",     FOREST,  "Discovery Workshop",
     "30-min session with Genos Apollo: review audit, confirm platform choice, share partner requirements and candidate journey map."),
    ("Week 3–5",   TEAL,    "Platform Sprint Begins",
     "Next.js rebuild starts. Candidate portal scoped. POPIA compliance drafted. Google Search Console connected. SEO keyword plan agreed."),
    ("Month 2",    GREEN,   "Portal & SEO Live",
     "New site live. Candidates apply via portal. First SEO-optimised opportunity pages indexed. WhatsApp intake bot active."),
]

for i, (when, col, title, desc) in enumerate(steps):
    x = Inches(0.4) + i * Inches(3.2)
    y = Inches(1.5)
    card(s11, x, y, Inches(3.0), Inches(4.0))
    rect(s11, x, y, Inches(3.0), Inches(0.45), col)
    tb(s11, when, x, y + Inches(0.06), Inches(3.0), Inches(0.36),
       size=12, bold=True,
       color=FOREST if col == ORANGE else WHITE,
       align=PP_ALIGN.CENTER)
    tb(s11, title, x + Inches(0.15), y + Inches(0.58),
       Inches(2.7), Inches(0.36), size=12, bold=True,
       color=ORANGE if col != ORANGE else FOREST)
    tb(s11, desc, x + Inches(0.15), y + Inches(1.0),
       Inches(2.7), Inches(2.8), size=10.5, color=FOREST)

card(s11, Inches(0.4), Inches(5.8), Inches(12.5), Inches(1.1), fill=FOREST)
tb(s11, "Ready to scale Africa's most compelling talent bridge?",
   Inches(0.6), Inches(5.88), Inches(7.5), Inches(0.4),
   size=14, bold=True, color=ORANGE)
tb(s11, "Contact Genos Apollo to schedule your discovery workshop.",
   Inches(0.6), Inches(6.28), Inches(7), Inches(0.38),
   size=11, color=WHITE)
tb(s11, "www.idili.africa  |  Innovation City Cape Town",
   Inches(8.2), Inches(6.0), Inches(4.8), Inches(0.38),
   size=10, color=ORANGE, italic=True, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════
# SLIDE 12 – THANK YOU / BACK COVER
# ════════════════════════════════════════════════════════════════
s12 = prs.slides.add_slide(BLANK)
bg(s12, FOREST)
rect(s12, 0, 0, Inches(0.7), H, ORANGE)
rect(s12, Inches(0.7), Inches(3.5), W - Inches(0.7), Inches(0.06), ORANGE)

tb(s12, "Thank You",
   Inches(1.1), Inches(2.0), Inches(11), Inches(1.2),
   size=52, bold=True, color=WHITE)
tb(s12, "Talent is global. Opportunity should be too.",
   Inches(1.1), Inches(3.1), Inches(10), Inches(0.55),
   size=16, color=ORANGE, italic=True)
tb(s12, "Idili Africa  ·  www.idili.africa  ·  Cape Town, South Africa",
   Inches(1.1), Inches(4.0), Inches(9), Inches(0.4),
   size=12, color=RGBColor(0xAA, 0xCC, 0xBB))
tb(s12, "Prepared by Genos Apollo  |  May 2026  |  Confidential",
   Inches(1.1), Inches(6.7), Inches(10), Inches(0.35),
   size=10, color=MGREY)

# ── Save ────────────────────────────────────────────────────────
OUT = ROOT / "output" / "decks" / "IdiliAfrica_WebDev_AI_Automation_Audit_GenosAI.pptx"
prs.save(str(OUT))
print(f"Saved: {OUT}")
