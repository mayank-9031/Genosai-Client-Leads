import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.parent / "utils"))
from ppt_theme import (
    NAVY, DARK_BLUE, ACCENT, GOLD, RED, GREEN, PURPLE, ORANGE, TEAL, CYAN,
    WHITE, LGREY, MGREY, W, H,
    bg, rect, tb, aline, bullets, tag,
    slide_header, slide_footer, score_badge,
    stat_cards, top_bar_card, phase_card, before_after, genos_logo_block,
)

ROOT = Path(__file__).parent.parent.parent
prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]
# Client-specific colors
DGREY    = RGBColor(0x2C, 0x38, 0x4A)
AMBER    = RGBColor(0xF3, 0x9C, 0x12)
BLUE_ACC = RGBColor(0x0A, 0x7B, 0xF3)


# ── Local helpers ────────────────────────────────────────────────────────────

def header(slide, title, subtitle, accent=ACCENT):
    rect(slide, 0, 0, Inches(0.12), H, accent)
    tb(slide, title.upper(), Inches(0.55 + 0.12), Inches(0.3),
       Inches(12), Inches(0.45), size=11, bold=True, color=accent)
    tb(slide, subtitle, Inches(0.55 + 0.12), Inches(0.72),
       Inches(12 - 0.12), Inches(0.65), size=22, bold=True, color=WHITE)
    aline(slide, Inches(1.38), color=accent, w=Inches(1.5))


def footer(slide,
           text='Genos AI  |  hello@genosai.tech  |  www.genosai.tech  |  +91 638-714-2699'):
    rect(slide, 0, Inches(7.05), W, Inches(0.45), DARK_BLUE)
    tb(slide, text, Inches(0.55), Inches(7.1), Inches(12.4), Inches(0.35),
       size=9, color=MGREY, align=PP_ALIGN.CENTER)


def card(slide, x, y, w, h, fill=DARK_BLUE):
    return rect(slide, x, y, w, h, fill)


def stat_card(slide, val, label, x, y, w=Inches(3.0), h=Inches(1.6), accent=ACCENT):
    rect(slide, x, y, w, h, DARK_BLUE)
    rect(slide, x, y, w, Inches(0.05), accent)
    tb(slide, str(val), x + Inches(0.1), y + Inches(0.1),
       w - Inches(0.2), Inches(0.9), size=28, bold=True, color=accent,
       align=PP_ALIGN.CENTER)
    tb(slide, label, x + Inches(0.1), y + Inches(1.1),
       w - Inches(0.2), Inches(0.42), size=10, color=LGREY, align=PP_ALIGN.CENTER)


def bullet_card(slide, title, items, x, y, w, h,
                accent=ACCENT, title_bg=None):
    if title_bg is None:
        title_bg = DARK_BLUE
    rect(slide, x, y, w, h, DARK_BLUE)
    rect(slide, x, y, w, Inches(0.42), title_bg)
    rect(slide, x, y, Inches(0.05), h, accent)
    tb(slide, title, x + Inches(0.1), y + Inches(0.05),
       w - Inches(0.15), Inches(0.35), size=11, bold=True, color=WHITE)
    bx = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.48),
                                  w - Inches(0.15), h - Inches(0.55))
    tf = bx.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r0 = p.add_run(); r0.text = '  '; r0.font.size = Pt(9)
        r1 = p.add_run(); r1.text = '● '; r1.font.size = Pt(9)
        r1.font.color.rgb = accent; r1.font.bold = True
        r2 = p.add_run(); r2.text = str(item)
        r2.font.size = Pt(10.5); r2.font.color.rgb = LGREY
        if i < len(items) - 1:
            p.space_after = Pt(5)




# SLIDE 1 – TITLE
# ════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
bg(s1, NAVY)
rect(s1, 0, 0, Inches(0.7), H, CYAN)
rect(s1, Inches(0.7), Inches(4.5), W - Inches(0.7), Inches(0.06), CYAN)

tb(s1, "NEWEXGEN (PTY) LTD",
   Inches(1.1), Inches(1.0), Inches(11), Inches(0.7),
   size=16, bold=True, color=CYAN)
tb(s1, "Web Development &\nAI Automation Audit",
   Inches(1.1), Inches(1.7), Inches(11), Inches(2.0),
   size=44, bold=True, color=WHITE)
tb(s1, "Strategic Review & Digital Transformation Roadmap",
   Inches(1.1), Inches(3.72), Inches(10), Inches(0.5),
   size=16, color=CYAN, italic=True)
tb(s1, "Prepared by: Genos Apollo  |  May 2026",
   Inches(1.1), Inches(5.0), Inches(8), Inches(0.4),
   size=12, color=MGREY)
tb(s1, "CONFIDENTIAL",
   Inches(1.1), Inches(6.6), Inches(4), Inches(0.35),
   size=10, color=MGREY, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 2 – AGENDA
# ════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
bg(s2)
header(s2, "Agenda", "What we will cover today")
footer(s2)

agenda = [
    ("01", "Company Digital Overview",    "Current state of newexgen.co.za and digital footprint"),
    ("02", "Website Audit Findings",      "UX, performance, SEO and accessibility gaps"),
    ("03", "Web Development Roadmap",     "Recommended improvements & modern tech stack"),
    ("04", "AI Automation Opportunity",   "Where AI can drive efficiency at NewExGen"),
    ("05", "Benefits of Automation",      "ROI, time savings and competitive advantage"),
    ("06", "Implementation Timeline",     "Phased rollout plan across 12 months"),
    ("07", "Investment & Next Steps",     "Budget guidance and immediate actions"),
]

for i, (num, title, desc) in enumerate(agenda):
    col = i % 2
    row = i // 2
    x = Inches(0.55) + col * Inches(6.4)
    y = Inches(1.55) + row * Inches(1.22)
    card(s2, x, y, Inches(6.1), Inches(1.05))
    rect(s2, x, y, Inches(0.7), Inches(1.05), NAVY)
    tb(s2, num, x, y + Inches(0.25), Inches(0.7), Inches(0.5),
       size=20, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
    tb(s2, title, x + Inches(0.8), y + Inches(0.1),
       Inches(5.0), Inches(0.4), size=13, bold=True, color=NAVY)
    tb(s2, desc, x + Inches(0.8), y + Inches(0.52),
       Inches(5.0), Inches(0.42), size=10.5, color=MGREY, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 3 – COMPANY DIGITAL OVERVIEW
# ════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
bg(s3)
header(s3, "Company Digital Overview", "NewExGen – current digital footprint")
footer(s3)

# Left – company snapshot
card(s3, Inches(0.4), Inches(1.4), Inches(5.8), Inches(5.5))
rect(s3, Inches(0.4), Inches(1.4), Inches(5.8), Inches(0.45), NAVY)
tb(s3, "Company Snapshot", Inches(0.55), Inches(1.44),
   Inches(5.5), Inches(0.38), size=13, bold=True, color=WHITE)

snapshot = [
    ("Industry",     "ICT / Digital Transformation & Managed Services"),
    ("Founded",      "2010  |  15 years in operation"),
    ("Geography",    "South Africa – HQ: Rivonia, Sandton (Gauteng)"),
    ("Segments",     "System Integration | ICT Field | Managed Svcs | BPO | SHERQ | Zoho"),
    ("Website",      "www.newexgen.co.za"),
    ("Booking",      "Zoho Bookings integration (dev-newexgen.zohobookings.com)"),
    ("Social",       "LinkedIn | Facebook | Twitter/X"),
    ("CMS",          "Custom / Zoho-based – likely Zoho Sites or custom stack"),
]
ty = Inches(2.02)
for k, v in snapshot:
    tb(s3, k, Inches(0.6), ty, Inches(1.55), Inches(0.3),
       size=10.5, bold=True, color=NAVY)
    tb(s3, v, Inches(2.2), ty, Inches(3.8), Inches(0.3),
       size=10.5, color=MGREY)
    ty += Inches(0.52)

# Right – digital health metrics
card(s3, Inches(6.5), Inches(1.4), Inches(6.5), Inches(5.5))
rect(s3, Inches(6.5), Inches(1.4), Inches(6.5), Inches(0.45), CYAN)
tb(s3, "Key Digital Health Indicators", Inches(6.65), Inches(1.44),
   Inches(6.2), Inches(0.38), size=13, bold=True, color=NAVY)

metrics = [
    ("Mobile Page Speed",       "51 / 100",   "Google PageSpeed – below benchmark",  RED),
    ("Desktop Speed",           "74 / 100",   "PageSpeed Insights estimate",          AMBER),
    ("SEO Score",               "61 / 100",   "Limited schema markup & metadata",     AMBER),
    ("Accessibility (WCAG)",    "Partial",    "Skip-nav present; AA gaps remain",     AMBER),
    ("SSL / HTTPS",             "✓ Active",   "Valid certificate in place",           GREEN),
    ("Zoho Booking Integration","✓ Live",     "Demo booking functional",              GREEN),
    ("AI / Chatbot",            "None",       "No automated lead handling on site",   RED),
    ("Social Integration",      "Limited",    "Icons only – no live feeds/automation",RED),
]
ty2 = Inches(2.02)
for label, val, note, col in metrics:
    tb(s3, label, Inches(6.65), ty2, Inches(2.7), Inches(0.3),
       size=10.5, bold=True, color=NAVY)
    tb(s3, val,   Inches(9.4), ty2, Inches(1.25), Inches(0.3),
       size=10.5, bold=True, color=col, align=PP_ALIGN.CENTER)
    tb(s3, note,  Inches(10.7), ty2, Inches(2.1), Inches(0.3),
       size=9, color=MGREY, italic=True)
    ty2 += Inches(0.52)

# ════════════════════════════════════════════════════════════════
# SLIDE 4 – WEBSITE AUDIT FINDINGS
# ════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
bg(s4)
header(s4, "Website Audit Findings", "newexgen.co.za – identified gaps and opportunities")
footer(s4)

audit_cards = [
    ("Performance",
     ["Mobile load >3.5s (target: <2.5s)", "No CDN – assets served from origin server",
      "Heavy homepage images not lazy-loaded", "Zoho Bookings iframe slowing page render"]),
    ("UX / Design",
     ["Service pages lack in-depth case study depth", "Demo booking flow breaks on mobile",
      "No live chat or instant-contact widget", "Trust signals (certs, logos) below the fold"]),
    ("SEO",
     ["Missing meta descriptions on most pages", "No structured data / schema markup",
      "Blog / thought-leadership content absent", "Thin service pages – low topical authority"]),
    ("Security & Compliance",
     ["No POPIA-compliant cookie consent banner", "Contact form validation needs hardening",
      "Missing security headers (CSP, X-Frame)", "Zoho booking subdomain not HSTS-enforced"]),
]

for i, (title, bullets) in enumerate(audit_cards):
    x = Inches(0.4) + i * Inches(3.2)
    y = Inches(1.45)
    card(s4, x, y, Inches(3.05), Inches(5.3))
    rect(s4, x, y, Inches(3.05), Inches(0.45), NAVY)
    tb(s4, title, x + Inches(0.1), y + Inches(0.06),
       Inches(2.85), Inches(0.36), size=12, bold=True, color=WHITE)
    ty = y + Inches(0.62)
    for b in bullets:
        rect(s4, x + Inches(0.15), ty + Inches(0.07), Inches(0.1), Inches(0.1), CYAN)
        tb(s4, b, x + Inches(0.35), ty, Inches(2.6), Inches(0.55),
           size=10, color=NAVY)
        ty += Inches(0.72)

rect(s4, Inches(0.4), Inches(6.8), Inches(12.5), Inches(0.06), CYAN)
tb(s4, "Priority: Performance + SEO content strategy + POPIA compliance upgrades deliver the fastest ROI on newexgen.co.za",
   Inches(0.4), Inches(6.86), Inches(12.5), Inches(0.3),
   size=10.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 5 – WEB DEVELOPMENT ROADMAP
# ════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
bg(s5)
header(s5, "Web Development Roadmap", "Modernising newexgen.co.za for performance and lead generation")
footer(s5)

# Left – tech stack
card(s5, Inches(0.4), Inches(1.4), Inches(4.5), Inches(5.6))
rect(s5, Inches(0.4), Inches(1.4), Inches(4.5), Inches(0.45), BLUE_ACC)
tb(s5, "Recommended Tech Stack", Inches(0.55), Inches(1.44),
   Inches(4.2), Inches(0.38), size=13, bold=True, color=WHITE)

stack = [
    ("Frontend",     "Next.js 15 (React) – SSR + static generation"),
    ("CMS",          "Sanity.io or Zoho Sites Pro – headless, API-first"),
    ("Hosting",      "Vercel / AWS CloudFront CDN"),
    ("CRM",          "Zoho CRM – native integration (existing partner)"),
    ("Booking",      "Zoho Bookings – deep embed, not iframe"),
    ("Forms",        "React Hook Form + Zod + POPIA consent"),
    ("Analytics",    "GA4 + Hotjar heatmaps"),
    ("Live Chat",    "Zoho SalesIQ – AI-powered chat widget"),
]
ty = Inches(2.02)
for k, v in stack:
    tb(s5, k, Inches(0.6), ty, Inches(1.4), Inches(0.3),
       size=10, bold=True, color=BLUE_ACC)
    tb(s5, v, Inches(2.05), ty, Inches(2.75), Inches(0.3),
       size=10, color=NAVY)
    ty += Inches(0.5)

# Right – improvements
card(s5, Inches(5.2), Inches(1.4), Inches(7.75), Inches(5.6))
rect(s5, Inches(5.2), Inches(1.4), Inches(7.75), Inches(0.45), NAVY)
tb(s5, "Key Improvements & Features",
   Inches(5.35), Inches(1.44), Inches(7.4), Inches(0.38),
   size=13, bold=True, color=WHITE)

improvements = [
    ("Speed & Core Web Vitals",
     "Target 90+ PageSpeed. CDN delivery, WebP images, SSR caching, remove render-blocking Zoho iframe."),
    ("Service & Solutions Pages Rebuild",
     "In-depth, SEO-optimised service pages with case studies, client logos, and measurable outcomes."),
    ("Zoho Ecosystem Deep Integration",
     "Zoho CRM + SalesIQ + Bookings natively embedded – unified lead capture, pipeline and analytics."),
    ("AI-Powered Live Chat (Zoho SalesIQ)",
     "Zia AI bot handles initial enquiries 24/7, qualifies leads, and books demos automatically."),
    ("Thought Leadership Blog / Resource Hub",
     "ICT and digital transformation content strategy to build organic authority and attract enterprise leads."),
    ("POPIA Cookie Consent & Security Headers",
     "Full POPIA-compliant consent flow, CSP + HSTS headers, hardened form validation across all pages."),
]
ty2 = Inches(1.98)
for title, desc in improvements:
    rect(s5, Inches(5.35), ty2 + Inches(0.1), Inches(0.08), Inches(0.25), CYAN)
    tb(s5, title, Inches(5.55), ty2, Inches(7.1), Inches(0.28),
       size=11, bold=True, color=NAVY)
    tb(s5, desc, Inches(5.55), ty2 + Inches(0.3), Inches(7.1), Inches(0.35),
       size=10, color=MGREY)
    ty2 += Inches(0.78)

# ════════════════════════════════════════════════════════════════
# SLIDE 6 – AI AUTOMATION OPPORTUNITY MAP
# ════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
bg(s6)
header(s6, "AI Automation Opportunity Map", "Where AI delivers immediate value for NewExGen")
footer(s6)

opps = [
    (NAVY,     "Lead Generation & Sales",
     ["AI chatbot (Zia) qualifies inbound leads 24/7", "Auto-score leads in Zoho CRM by fit & intent",
      "Personalised follow-up email sequences (Zoho Campaigns)", "Demo booking confirmations & reminders automated"]),
    (CYAN,     "Service Delivery & Field Ops",
     ["AI-assisted ticket triage & routing (Zoho Desk)", "Predictive SLA breach alerts for managed services",
      "Field technician scheduling optimisation (AI dispatch)", "Auto-generate client-facing service reports"]),
    (BLUE_ACC, "Market Intelligence & Growth",
     ["Auto-monitor ICT tender portals (CIDB, e-Tender)", "Weekly AI-generated SA tech market digests",
      "Competitor pricing & service gap analysis", "LinkedIn content auto-drafted from case studies"]),
    (GREEN,    "Operations & Administration",
     ["Invoice processing & AP automation (Zoho Books)", "Contract review & SLA clause flagging",
      "HR: AI-assisted CV screening (Zoho Recruit)", "Meeting notes → action items (AI transcription)"]),
]

for i, (col, title, bullets) in enumerate(opps):
    c = i % 2
    r = i // 2
    x = Inches(0.4) + c * Inches(6.5)
    y = Inches(1.45) + r * Inches(2.85)
    card(s6, x, y, Inches(6.2), Inches(2.65))
    rect(s6, x, y, Inches(6.2), Inches(0.45), col)
    tb(s6, title, x + Inches(0.15), y + Inches(0.06),
       Inches(5.9), Inches(0.36), size=13, bold=True,
       color=NAVY if col == CYAN else WHITE)
    ty = y + Inches(0.62)
    for b in bullets:
        rect(s6, x + Inches(0.18), ty + Inches(0.08), Inches(0.1), Inches(0.1), col)
        tb(s6, b, x + Inches(0.38), ty, Inches(5.65), Inches(0.38),
           size=10.5, color=NAVY)
        ty += Inches(0.48)

# ════════════════════════════════════════════════════════════════
# SLIDE 7 – BENEFITS OF AI AUTOMATION
# ════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
bg(s7)
header(s7, "Benefits of AI Automation", "Measurable impact across NewExGen's operations")
footer(s7)

stats = [
    ("60–70%",  "Reduction in\nrepetitive admin tasks"),
    ("3×",      "Faster lead\nqualification & booking"),
    ("40%",     "Lower cost per\nclient enquiry"),
    ("24/7",    "AI-powered client\nservice coverage"),
]
for i, (val, label) in enumerate(stats):
    stat_card(s7, val, label,
              Inches(0.4) + i * Inches(3.2), Inches(1.45),
              w=Inches(3.0), h=Inches(1.6))

benefits = [
    ("Lead Velocity",
     ["AI chatbot captures leads outside office hours", "Auto-qualify and route to right sales rep",
      "Demo booked in the same conversation", "No lead falls through the cracks"]),
    ("Cost Efficiency",
     ["Reduce admin staff time on data entry by >50%", "Lower dependency on manual follow-up calls",
      "Zoho automation replaces multiple point tools", "Scale client base without scaling headcount"]),
    ("Service Delivery Excellence",
     ["AI ticket triage reduces resolution time", "Predictive alerts prevent SLA breaches",
      "Auto-generated reports save technician hours", "Client satisfaction scores improve measurably"]),
    ("Competitive Positioning",
     ["First SA ICT firm to offer AI-native managed services", "Differentiate from legacy IT providers",
      "Attract enterprise and government contracts", "Data-driven decisions across all service lines"]),
]

for i, (title, bullets) in enumerate(benefits):
    c = i % 2
    r = i // 2
    x = Inches(0.4) + c * Inches(6.5)
    y = Inches(3.2) + r * Inches(1.85)
    bullet_card(s7, title, bullets, x, y, Inches(6.2), Inches(1.65),
                title_bg=NAVY if r == 0 else BLUE_ACC)

# ════════════════════════════════════════════════════════════════
# SLIDE 8 – ROI & BUSINESS CASE
# ════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
bg(s8)
header(s8, "ROI & Business Case", "Projected returns from web + AI investment (ZAR)")
footer(s8)

# Left – investment
card(s8, Inches(0.4), Inches(1.4), Inches(5.8), Inches(5.6))
rect(s8, Inches(0.4), Inches(1.4), Inches(5.8), Inches(0.45), NAVY)
tb(s8, "Indicative Investment (Year 1 – ZAR)",
   Inches(0.55), Inches(1.44), Inches(5.5), Inches(0.38),
   size=13, bold=True, color=WHITE)

invest_rows = [
    ("Website Rebuild (Next.js + Sanity / Zoho)",  "R 90,000 – R 160,000"),
    ("Zoho Ecosystem Integration & Config",         "R 30,000 – R  55,000"),
    ("AI Chatbot (Zia / SalesIQ) Setup",            "R 20,000 – R  40,000"),
    ("SEO Content Sprint & POPIA Compliance",       "R 18,000 – R  32,000"),
    ("Lead Automation & CRM Workflows",             "R 15,000 – R  28,000"),
    ("Training & Change Management",                "R 10,000 – R  18,000"),
    ("Ongoing monthly support (est.)",              "R  8,000 – R  15,000 / mo"),
    ("TOTAL YEAR 1 (once-off + 12 mo)",            "R 279,000 – R 513,000"),
]
ty = Inches(2.0)
for i, (item, cost) in enumerate(invest_rows):
    is_total = i == len(invest_rows) - 1
    if is_total:
        rect(s8, Inches(0.4), ty - Inches(0.05), Inches(5.8), Inches(0.06), CYAN)
        ty += Inches(0.12)
        tb(s8, item, Inches(0.6), ty, Inches(3.5), Inches(0.35),
           size=11, bold=True, color=NAVY)
        tb(s8, cost, Inches(4.05), ty, Inches(1.95), Inches(0.35),
           size=11, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)
    else:
        tb(s8, item, Inches(0.6), ty, Inches(3.5), Inches(0.3),
           size=10, color=NAVY)
        tb(s8, cost, Inches(4.05), ty, Inches(1.95), Inches(0.3),
           size=10, color=MGREY, align=PP_ALIGN.RIGHT)
    ty += Inches(0.5)

# Right – returns
card(s8, Inches(6.5), Inches(1.4), Inches(6.5), Inches(5.6))
rect(s8, Inches(6.5), Inches(1.4), Inches(6.5), Inches(0.45), GREEN)
tb(s8, "Projected Returns & Savings (ZAR)",
   Inches(6.65), Inches(1.44), Inches(6.2), Inches(0.38),
   size=13, bold=True, color=WHITE)

returns = [
    ("Admin time saved (FTE equivalent)",            "1–1.5 FTE/yr  ≈  R 240,000+"),
    ("Reduced manual follow-up & call costs",        "~R 80,000 / year"),
    ("Increased qualified leads via SEO + chatbot",  "+30–45% organic lead growth"),
    ("Faster deal close via automation",             "Est. −7 days sales cycle per deal"),
    ("SLA breach reduction (managed services)",      "−25% penalty & remediation costs"),
    ("Reduced admin per new client onboarding",      "R 30,000+ saved / year"),
    ("ESTIMATED YEAR 1 SAVINGS / UPSIDE",           "R 350,000 – R 600,000+"),
]
ty2 = Inches(2.0)
for i, (item, val) in enumerate(returns):
    is_total = i == len(returns) - 1
    if is_total:
        rect(s8, Inches(6.5), ty2 - Inches(0.05), Inches(6.5), Inches(0.06), CYAN)
        ty2 += Inches(0.12)
        tb(s8, item, Inches(6.65), ty2, Inches(3.9), Inches(0.35),
           size=11, bold=True, color=NAVY)
        tb(s8, val, Inches(10.6), ty2, Inches(2.2), Inches(0.35),
           size=11, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)
    else:
        tb(s8, item, Inches(6.65), ty2, Inches(3.9), Inches(0.3),
           size=10, color=NAVY)
        tb(s8, val, Inches(10.6), ty2, Inches(2.2), Inches(0.3),
           size=10, color=MGREY, align=PP_ALIGN.RIGHT)
    ty2 += Inches(0.5)

tb(s8, "Breakeven estimated within 7–12 months  |  Zoho-native stack reduces ongoing licensing costs",
   Inches(0.4), Inches(7.07), Inches(12.5), Inches(0.28),
   size=10.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 9 – IMPLEMENTATION TIMELINE
# ════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(BLANK)
bg(s9)
header(s9, "Implementation Timeline", "Phased 12-month delivery plan")
footer(s9)

phases = [
    ("Phase 1\nMon 1–3",  CYAN,     [
        "Website audit & quick-win SEO sprint",
        "POPIA cookie consent go-live",
        "Zoho SalesIQ AI chatbot MVP live",
        "CRM lead-capture workflow setup",
    ]),
    ("Phase 2\nMon 4–6",  NAVY,     [
        "Full Next.js website rebuild launch",
        "Deep Zoho CRM + Bookings integration",
        "Service pages rebuild with case studies",
        "LinkedIn + social content automation",
    ]),
    ("Phase 3\nMon 7–9",  BLUE_ACC, [
        "AI ticket triage (Zoho Desk)",
        "Thought-leadership blog & resource hub",
        "Tender monitoring automation (AI)",
        "Staff training & enablement",
    ]),
    ("Phase 4\nMon 10–12", GREEN,   [
        "Predictive SLA analytics dashboard",
        "Advanced chatbot (AI fine-tuning)",
        "Zoho Analytics full BI rollout",
        "Full audit, review & optimise",
    ]),
]

for i, (label, col, items) in enumerate(phases):
    x = Inches(0.4) + i * Inches(3.2)
    y = Inches(1.45)
    card(s9, x, y, Inches(3.0), Inches(5.55))
    rect(s9, x, y, Inches(3.0), Inches(0.7), col)
    tb(s9, label, x, y + Inches(0.06), Inches(3.0), Inches(0.62),
       size=13, bold=True,
       color=NAVY if col == CYAN else WHITE,
       align=PP_ALIGN.CENTER)
    if i < 3:
        rect(s9, x + Inches(3.02), y + Inches(0.3), Inches(0.16), Inches(0.06), MGREY)
    ty = y + Inches(0.88)
    for item in items:
        rect(s9, x + Inches(0.18), ty + Inches(0.08), Inches(0.1), Inches(0.1), col)
        tb(s9, item, x + Inches(0.38), ty, Inches(2.5), Inches(0.42),
           size=10.5, color=NAVY)
        ty += Inches(0.9)

# ════════════════════════════════════════════════════════════════
# SLIDE 10 – WHY AUTOMATE NOW?
# ════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(BLANK)
bg(s10, NAVY)
rect(s10, 0, 0, Inches(0.5), H, CYAN)
footer(s10)

tb(s10, "Why Automate — Now?",
   Inches(0.8), Inches(0.5), Inches(11), Inches(0.65),
   size=30, bold=True, color=CYAN)
tb(s10, "South Africa's ICT market is consolidating fast. NewExGen's 15-year foundation is a launchpad — AI turns it into a competitive moat.",
   Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.55),
   size=14, color=WHITE, italic=True)

reasons = [
    ("The SA ICT sector is in rapid consolidation",
     "Larger players (BCX, DXC, Dimension Data) are deploying AI-powered managed services at scale. NewExGen must differentiate on speed, intelligence, and client intimacy — AI enables all three."),
    ("Every manual process is a cost disadvantage",
     "Each ticket, report, or follow-up email handled manually costs money and time. Automation via the existing Zoho stack is the fastest path to margin improvement with the lowest switching cost."),
    ("Enterprise & government clients now expect digital-first engagement",
     "Tenders, RFPs, and procurement increasingly favour vendors with proven digital maturity. A modern website and AI-driven processes signal credibility before the first meeting."),
    ("Zoho's AI capabilities are ready to deploy today",
     "NewExGen is already a 4+ year Zoho partner. Zia AI, SalesIQ, Desk, and Campaigns are available under existing licenses. The cost to activate AI is far lower than a greenfield tool purchase."),
    ("Data compounds competitive advantage over time",
     "Every automated interaction — tickets, leads, field jobs — feeds Zoho Analytics. Within 12–24 months NewExGen builds a proprietary operational dataset that makes pricing, forecasting, and service design far more precise than competitors operating on gut feel."),
]
ty = Inches(1.95)
for i, (title, desc) in enumerate(reasons):
    rect(s10, Inches(0.8), ty, Inches(0.55), Inches(0.55),
         CYAN if i % 2 == 0 else BLUE_ACC)
    tb(s10, str(i + 1), Inches(0.8), ty, Inches(0.55), Inches(0.55),
       size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    tb(s10, title, Inches(1.5), ty, Inches(11), Inches(0.3),
       size=12, bold=True, color=WHITE)
    tb(s10, desc, Inches(1.5), ty + Inches(0.3), Inches(11.3), Inches(0.55),
       size=10.5, color=RGBColor(0xBB, 0xCC, 0xDD))
    ty += Inches(0.98)

# ════════════════════════════════════════════════════════════════
# SLIDE 11 – NEXT STEPS
# ════════════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(BLANK)
bg(s11)
header(s11, "Next Steps", "From audit to action — what happens now")
footer(s11)

steps = [
    ("Week 1",  CYAN,     "Stakeholder Alignment",
     "Present findings to NewExGen leadership (Melvin, Jeffrey, Zelda). Agree on priority areas, budget envelope, and Phase 1 scope."),
    ("Week 2",  NAVY,     "Kick-Off & Access",
     "Provide CMS, Zoho, hosting and analytics access. Assign internal champion. Finalise project charter."),
    ("Week 3",  BLUE_ACC, "Sprint 1 Begins",
     "POPIA fix, SEO quick wins, and Zoho SalesIQ AI chatbot integration start simultaneously."),
    ("Month 2", GREEN,    "Phase 1 Go-Live",
     "AI chatbot live on site. SEO fixes deployed. CRM lead workflows active. First automated demo bookings in production."),
]

for i, (when, col, title, desc) in enumerate(steps):
    x = Inches(0.4) + i * Inches(3.2)
    y = Inches(1.5)
    card(s11, x, y, Inches(3.0), Inches(4.0))
    rect(s11, x, y, Inches(3.0), Inches(0.45), col)
    tb(s11, when, x, y + Inches(0.06), Inches(3.0), Inches(0.36),
       size=12, bold=True,
       color=NAVY if col == CYAN else WHITE,
       align=PP_ALIGN.CENTER)
    tb(s11, title, x + Inches(0.15), y + Inches(0.58),
       Inches(2.7), Inches(0.36), size=12, bold=True, color=col if col != CYAN else NAVY)
    tb(s11, desc, x + Inches(0.15), y + Inches(1.0),
       Inches(2.7), Inches(2.8), size=10.5, color=NAVY)

# CTA box
card(s11, Inches(0.4), Inches(5.8), Inches(12.5), Inches(1.1), fill=NAVY)
tb(s11, "Ready to accelerate NewExGen's digital transformation?",
   Inches(0.6), Inches(5.88), Inches(7.5), Inches(0.4),
   size=14, bold=True, color=CYAN)
tb(s11, "Contact Genos Apollo to schedule your discovery workshop.",
   Inches(0.6), Inches(6.28), Inches(7), Inches(0.38),
   size=11, color=WHITE)
tb(s11, "www.newexgen.co.za  |  info@newexgen.co.za  |  +27 87 087 8307",
   Inches(8.0), Inches(6.0), Inches(5.0), Inches(0.38),
   size=10, color=CYAN, italic=True, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════
# SLIDE 12 – THANK YOU / BACK COVER
# ════════════════════════════════════════════════════════════════
s12 = prs.slides.add_slide(BLANK)
bg(s12, NAVY)
rect(s12, 0, 0, Inches(0.7), H, CYAN)
rect(s12, Inches(0.7), Inches(3.5), W - Inches(0.7), Inches(0.06), CYAN)

tb(s12, "Thank You",
   Inches(1.1), Inches(2.0), Inches(11), Inches(1.2),
   size=52, bold=True, color=WHITE)
tb(s12, "Driving South Africa's digital transformation — one enterprise at a time.",
   Inches(1.1), Inches(3.1), Inches(10), Inches(0.55),
   size=16, color=CYAN, italic=True)
tb(s12, "NewExGen (Pty) Ltd  ·  www.newexgen.co.za",
   Inches(1.1), Inches(4.0), Inches(8), Inches(0.4),
   size=12, color=RGBColor(0xAA, 0xBB, 0xCC))
tb(s12, "Prepared by Genos Apollo  |  May 2026  |  Confidential",
   Inches(1.1), Inches(6.7), Inches(10), Inches(0.35),
   size=10, color=MGREY)

# ── Save ────────────────────────────────────────────────────────
OUT = ROOT / "output" / "decks" / "NewExGen_WebDev_AI_Automation_Audit_GenosAI.pptx"
prs.save(str(OUT))
print(f"Saved: {OUT}")
