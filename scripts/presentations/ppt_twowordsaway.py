import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.parent / "utils"))
from ppt_theme import (
    NAVY, DARK_BLUE, ACCENT, GOLD, RED, GREEN, PURPLE, ORANGE, TEAL, CYAN,
    WHITE, LGREY, MGREY, W, H,
    bg, rect, tb, aline, bullets, tag,
    slide_header, slide_footer, score_badge,
    stat_cards, top_bar_card, phase_card, before_after, genos_logo_block,
)

ROOT = Path(__file__).parent.parent.parent
prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]

# Client-specific colors
AMBER     = RGBColor(0xFF, 0xA0, 0x00)
WA_GREEN  = RGBColor(0x25, 0xD3, 0x66)

# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
bg(s1)
rect(s1, 0, 0, Inches(0.12), H, ACCENT)
rect(s1, Inches(8.5), 0, Inches(4.83), Inches(2.2), DARK_BLUE)
tb(s1, "GENOS AI", Inches(9.2), Inches(0.4), Inches(3.5), Inches(0.8),
   size=28, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
tb(s1, "www.genosai.tech", Inches(9.2), Inches(1.05), Inches(3.5), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)

# Shark Tank badge
rect(s1, Inches(0.55), Inches(0.8), Inches(3.0), Inches(0.38), GOLD)
tb(s1, "  AS SEEN ON SHARK TANK INDIA", Inches(0.55), Inches(0.82), Inches(3.0), Inches(0.35),
   size=10, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

tb(s1, "TWOWORDSAWAY.COM", Inches(0.55), Inches(1.38),
   Inches(9), Inches(0.6), size=13, bold=False, color=ACCENT)
tb(s1, "Website & AI Growth\nAudit Report", Inches(0.55), Inches(1.88),
   Inches(9), Inches(1.8), size=46, bold=True, color=WHITE)
aline(s1, Inches(3.68), color=GOLD, w=Inches(3))
tb(s1, "Prepared exclusively for the TwoWordsAway founding team",
   Inches(0.55), Inches(3.98), Inches(9), Inches(0.5), size=14, color=LGREY)

rect(s1, 0, Inches(6.5), W, Inches(1.0), DARK_BLUE)
tb(s1, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699",
   Inches(0.55), Inches(6.6), Inches(9), Inches(0.5), size=12, color=MGREY)
tb(s1, "May 2026  |  CONFIDENTIAL", Inches(9.5), Inches(6.6), Inches(3.3), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)

rect(s1, Inches(10.5), Inches(3.3), Inches(2.3), Inches(2.5), DARK_BLUE)
tb(s1, "AUDIT SCORE", Inches(10.55), Inches(3.4), Inches(2.2), Inches(0.5),
   size=11, bold=True, color=MGREY, align=PP_ALIGN.CENTER)
tb(s1, "7", Inches(10.55), Inches(3.75), Inches(2.2), Inches(1.1),
   size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
tb(s1, "/ 10", Inches(10.55), Inches(4.7), Inches(2.2), Inches(0.4),
   size=18, color=LGREY, align=PP_ALIGN.CENTER)
tag(s1, "HIGH PRIORITY", Inches(10.75), Inches(5.1), fill=RED, size=10)


# ══════════════════════════════════════════════════════════
# SLIDE 2 — AT A GLANCE
# ══════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
bg(s2)
rect(s2, 0, 0, Inches(0.12), H, ACCENT)
tb(s2, "AT A GLANCE", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s2, "Two Words Away — What They've Built",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s2, Inches(1.35))

stats = [
    ("120+", "Professionals\nat Tables"),
    ("60+", "Tables\nHosted"),
    ("6", "Cities\nActive"),
    ("97%", "Post-Event\nSatisfaction"),
    ("4", "Event Types\nDinner / Coffee\nDrinks / Lunch"),
]
cw = Inches(2.3); ch = Inches(2.2); gap = Inches(0.18)
for i, (num, label) in enumerate(stats):
    l = Inches(0.55) + i * (cw + gap)
    rect(s2, l, Inches(1.8), cw, ch, DARK_BLUE)
    tb(s2, num, l + Inches(0.1), Inches(1.95), cw - Inches(0.2), Inches(1.1),
       size=30, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s2, label, l + Inches(0.1), Inches(3.05), cw - Inches(0.2), Inches(0.7),
       size=11, color=LGREY, align=PP_ALIGN.CENTER)

rect(s2, Inches(0.55), Inches(4.2), Inches(12.23), Inches(0.6), DARK_BLUE)
tb(s2, "Curated networking dinners & coffees  ·  Bangalore · Mumbai · Delhi · Gurgaon + 2 cities  ·  Free signup, per-seat fee per table  ·  As seen on Shark Tank India",
   Inches(0.65), Inches(4.28), Inches(12), Inches(0.45), size=11, color=LGREY, align=PP_ALIGN.CENTER)

tb(s2, "WHAT THEY DO", Inches(0.55), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=ACCENT)
bullets(s2, [
    "Curate small-group dinners (4–8 people) for founders, investors, engineers, and growth professionals",
    "Handle venue, curation, and invites — attendees just show up",
    "Free to sign up; paid per-seat model for confirmed tables",
], Inches(0.55), Inches(5.35), Inches(5.8), Inches(1.5), size=12, color=LGREY)

tb(s2, "TARGET USERS", Inches(6.95), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=ACCENT)
bullets(s2, [
    "Founders — fundraising, co-founders, warm intros to VCs",
    "Investors & Angels — early access to high-intent founders",
    "Operators, Product, Growth, Engineers, Ambitious Employees",
], Inches(6.95), Inches(5.35), Inches(5.8), Inches(1.5), size=12, color=LGREY)


# ══════════════════════════════════════════════════════════
# SLIDE 3 — WHAT'S WORKING
# ══════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
bg(s3)
rect(s3, 0, 0, Inches(0.12), H, GREEN)
tb(s3, "WHAT'S WORKING", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s3, "Strengths to Build On",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s3, Inches(1.35), color=GREEN)

strengths = [
    ("Shark Tank India — On-Site and Prominent",
     "The Shark Tank India credibility badge is front-and-center on the homepage. For a bootstrapped product in India, this is a moat. New visitors immediately see third-party validation before reading a single line of copy. Most early-stage startups spend years trying to earn this kind of trust — TWA already has it."),
    ("Razor-Sharp Value Proposition",
     '"We handle the people, table and venue — you just show up." That\'s the entire product in one sentence. It\'s clear, memorable, and directly eliminates the two biggest objections to networking events: effort and awkwardness. The homepage passes the 5-second test.'),
    ("Live Event Inventory with Seat Scarcity",
     "The /events page shows 8 upcoming tables across Bangalore, Mumbai, Delhi, and Gurgaon — with real seat counts like '2 seats left' and '1 seat left'. This is not vaporware. Real events, real scarcity, real urgency. It's the strongest conversion signal on the whole site."),
    ("6 Persona-Specific Value Sections",
     "The homepage breaks down benefits for Founders, Investors, Operators, Product/Growth, Engineers, and Ambitious Employees — each with 3 tailored bullet points. Visitors self-identify. This dramatically improves conversion rates compared to generic 'meet professionals' messaging."),
    ("Frictionless Entry — Free Signup",
     "No credit card, no commitment at signup. The free model eliminates the #1 reason people don't try new networking products: financial risk. Once someone is in the system, the conversion from 'registered' to 'paid seat' is a much easier ask."),
]
for i, (title, desc) in enumerate(strengths):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s3, Inches(0.55), t, Inches(0.06), Inches(0.72), GREEN)
    tb(s3, title, Inches(0.75), t, Inches(11.8), Inches(0.38), size=14, bold=True, color=WHITE)
    tb(s3, desc, Inches(0.75), t + Inches(0.38), Inches(11.8), Inches(0.58), size=11.5, color=LGREY)


# ══════════════════════════════════════════════════════════
# SLIDE 4 — WHERE IT'S BROKEN
# ══════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
bg(s4)
rect(s4, 0, 0, Inches(0.12), H, RED)
tb(s4, "WHERE IT'S BROKEN", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=RED)
tb(s4, "6 Issues Found in the Audit",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s4, Inches(1.35), color=RED)

issues = [
    ("/about — 404 Page Not Found", "HIGH",
     "The footer prominently links to 'About' — clicking it returns a 404. Visitors looking for the founding story, team background, and company mission hit a dead end. This page is critical for enterprise clients and investors doing due diligence."),
    ("/contact-us — 404 Page Not Found", "HIGH",
     "'Contact Us' appears in both the main nav and the footer. Both routes return 404. There is no way to reach the team through the website other than the two email addresses buried in the footer. Any serious B2B or press inquiry dies here."),
    ("/blog — 404 Page Not Found", "HIGH",
     "The blog is listed in the footer nav — a high-value SEO and content marketing channel. The route doesn't exist. Every 'networking tips', 'how to meet investors', 'Shark Tank experience' article that should be driving organic search traffic is simply not there."),
    ("/feedback and /support — Both 404", "MEDIUM",
     "Both Feedback and Support are linked in the footer nav. Both throw 404 errors. Users who want to report an issue, suggest an improvement, or get help before booking a table have no self-serve path. This erodes trust post-signup."),
    ("Pricing Page Has Wrong Product Copy", "HIGH",
     "The /pricing page still describes an older concept: 'Two Words Away is free to use. Say when you\'re free, get matched, make plans.' This contradicts the homepage which clearly states attendees pay a per-seat fee. A visitor who checks pricing before booking reads two completely different products."),
    ("Events Require Login to See Venue & Guest List", "MEDIUM",
     "Logged-out visitors can see event cards but cannot see the venue or who's attending until they sign up. While this drives registration, it also means curious visitors can't assess table quality before committing — a silent conversion friction point."),
]
for i, (title, pri, desc) in enumerate(issues):
    t = Inches(1.62) + i * Inches(0.97)
    pc = RED if pri == "HIGH" else AMBER
    rect(s4, Inches(0.55), t, Inches(0.06), Inches(0.68), pc)
    tb(s4, title, Inches(0.75), t, Inches(9.8), Inches(0.35), size=13, bold=True, color=WHITE)
    tag(s4, pri, Inches(10.85), t + Inches(0.04), fill=pc, size=9)
    tb(s4, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.52), size=11, color=LGREY)


# ══════════════════════════════════════════════════════════
# SLIDE 5 — WEBSITE DESIGN AUDIT (3-col comparison)
# ══════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
bg(s5)
rect(s5, 0, 0, Inches(0.12), H, PURPLE)
tb(s5, "WEBSITE DESIGN AUDIT", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=PURPLE)
tb(s5, "The Site Works. It Could Convert 3x Better.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.8), size=32, bold=True, color=WHITE)
aline(s5, Inches(1.45), color=PURPLE, w=Inches(1.5))

rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(3.8), DARK_BLUE)
rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(0.42), RED)
tb(s5, "TWOWORDSAWAY TODAY", Inches(0.65), Inches(1.92), Inches(3.6), Inches(0.38),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "/about, /contact, /blog all return 404",
    "Pricing page shows old product — wrong copy",
    "No team/founder story on the site",
    "No testimonial depth — only 3 short quotes",
    "No FAQ expansion — 4 questions, all collapsed",
    "Events are auth-gated — no venue preview",
    "No blog / content marketing engine",
    "No sticky CTA or floating 'Book a Seat' button",
], Inches(0.7), Inches(2.42), Inches(3.6), Inches(3.1), size=11, color=LGREY, dot=RED)

tb(s5, "→", Inches(4.45), Inches(3.7), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(3.8), DARK_BLUE)
rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(0.42), GOLD)
tb(s5, "TOP EVENT PLATFORMS DO THIS", Inches(5.1), Inches(1.92), Inches(3.4), Inches(0.38),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "About page: founder story + mission + team photos",
    "Contact page: form + WhatsApp + email + response SLA",
    "Rich blog driving SEO (networking tips, founder guides)",
    "Testimonials: 20+ with names, roles, company logos",
    "FAQ: 10+ questions, fully expanded, schema markup",
    "Event previews: show venue type before login",
    "Sticky 'Sign Up Free' header always visible",
    "Social proof ticker: '47 people joined this week'",
], Inches(5.1), Inches(2.42), Inches(3.4), Inches(3.1), size=11, color=LGREY, dot=GOLD)

tb(s5, "→", Inches(8.68), Inches(3.7), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(3.8), DARK_BLUE)
rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(0.42), GREEN)
tb(s5, "WHAT GENOS AI BUILDS", Inches(9.35), Inches(1.92), Inches(3.4), Inches(0.38),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Fix all 5 broken pages — about, contact, blog, feedback, support",
    "Rewrite /pricing to match actual paid-seat model",
    "Founder story page: Shark Tank journey, mission, team",
    "30+ testimonials section, auto-populated via WhatsApp",
    "Blog: SEO articles on networking, hiring, India startups",
    "Sticky 'Book a Table' CTA on every page",
    "WhatsApp chat widget for instant attendee queries",
    "Social proof: live seat counter + recent signups ticker",
], Inches(9.35), Inches(2.42), Inches(3.4), Inches(3.1), size=11, color=LGREY, dot=GREEN)

rect(s5, Inches(0.55), Inches(5.88), Inches(12.23), Inches(1.12), DARK_BLUE)
tb(s5, "WHY IT MATTERS:", Inches(0.7), Inches(5.95), Inches(2.2), Inches(0.38),
   size=11, bold=True, color=PURPLE)
stats2 = [
    ("88%", "of users won't return\nafter a bad experience"),
    ("3 sec", "to decide stay or leave\n— first impression is design"),
    ("200%+", "more conversions from\na trust-optimised site"),
    ("5 broken\npages", "found in this audit\ncausing silent drop-off"),
]
for j, (num, lbl) in enumerate(stats2):
    lx = Inches(3.0) + j * Inches(2.45)
    tb(s5, num, lx, Inches(5.92), Inches(1.6), Inches(0.58),
       size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s5, lbl, lx, Inches(6.48), Inches(1.6), Inches(0.5),
       size=10, color=MGREY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
# SLIDE 6 — WHAT THIS IS COSTING YOU
# ══════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
bg(s6)
rect(s6, 0, 0, Inches(0.12), H, GOLD)
tb(s6, "WHAT THIS IS COSTING YOU", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GOLD)
tb(s6, "The Hidden Cost of the Status Quo",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s6, Inches(1.35), color=GOLD)

costs = [
    ("Every Signup Gets Silence After the Thank-You Page",
     "Someone discovers Two Words Away, creates a free account, and then... nothing happens. No WhatsApp message. No 'here's a table in your city'. No nudge when a relevant seat opens. The conversion window from signup to first booking is highest in the first 48 hours — and right now that window closes with zero automated touchpoints."),
    ("8 Live Tables, Zero Automated Reminders",
     "There are 8 upcoming tables across 4 cities. Every confirmed attendee is relying on their own calendar to show up. There's no 48-hour reminder, no day-of venue reveal, no 'here's who you'll meet tonight' message. No-show rates for unreminded events average 20-30%. That's lost revenue and degraded table quality every single event."),
    ("Post-Event: The Warm Moment Is Wasted",
     "The highest-intent moment in the entire attendee journey is the 24 hours after a table. People are energised, connections are fresh, and they're most likely to book again and refer friends. Right now, there is no post-event WhatsApp follow-up, no connection list, no feedback ask, and no 'book your next table' nudge. That moment evaporates every time."),
    ("5 Broken Pages Silently Killing Credibility",
     "Every time a VC, journalist, enterprise HR team, or serious professional clicks 'About' and lands on a 404, the brand takes a credibility hit. It signals either a site that's been abandoned or a team that doesn't test their own product. First impressions with decision-makers don't get a second chance."),
]
for i, (title, desc) in enumerate(costs):
    t = Inches(1.65) + i * Inches(1.38)
    rect(s6, Inches(0.55), t, Inches(12.23), Inches(1.25), DARK_BLUE)
    tb(s6, "!  " + title, Inches(0.7), t + Inches(0.1), Inches(8), Inches(0.42),
       size=14, bold=True, color=GOLD)
    tb(s6, desc, Inches(0.7), t + Inches(0.52), Inches(11.8), Inches(0.62),
       size=11.5, color=LGREY)


# ══════════════════════════════════════════════════════════
# SLIDE 7 — WHATSAPP AUTOMATION OPPORTUNITY
# ══════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
bg(s7)
rect(s7, 0, 0, Inches(0.12), H, WA_GREEN)
tb(s7, "WHATSAPP AUTOMATION OPPORTUNITY", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=WA_GREEN)
tb(s7, "Seven AI Flows Built for TwoWordsAway",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s7, Inches(1.35), color=WA_GREEN)

opps = [
    ("01", "Signup Nurture Sequence",
     "The moment someone signs up, a WhatsApp message fires: welcome, profile confirmation, and the next available table in their city with a direct booking link. No cold inbox. No waiting. The highest-intent window captured automatically."),
    ("02", "Seat Availability Alerts",
     "When a new table opens in a city, all registered users matched to that city + persona type get an alert: event type, date, seats available, and one-tap booking. Fills seats faster than any email campaign. Urgency is built-in when only 2 seats remain."),
    ("03", "48-Hour & Day-Of Event Reminders",
     "Two messages per confirmed booking: a 48h reminder with event details and a day-of message with the venue address + a preview of who they'll be meeting. Reduces no-shows by 25-40%. Table quality improves. Attendee satisfaction goes up."),
    ("04", "Post-Event Follow-Up & Connection List",
     "Within 2 hours of a table ending, every attendee gets a WhatsApp message with the names and details of who they met, a prompt to connect, and a 'book your next table' CTA. The warm moment is captured. Repeat bookings increase significantly."),
    ("05", "Feedback Collection (NPS via WhatsApp)",
     "Post-event feedback request via WhatsApp — a simple 1-5 rating plus one open question. Response rates on WhatsApp are 5x higher than email surveys. The data feeds directly into curation quality: better tables, better testimonials, better product."),
    ("06", "No-Show Recovery & Waitlist Management",
     "If an attendee doesn't show up, an automated message checks in the next morning and offers them the next available table in their city. If a table is full, users join a WhatsApp waitlist and get notified the moment a seat opens — no manual coordination needed."),
    ("07", "Re-engagement for Lapsed Signups",
     "Users who signed up but never booked, or who haven't attended in 60+ days, get a personalised re-engagement message: 'A Founders Dinner is happening in Delhi this Friday — 3 seats left. Want in?' Converts dormant users back into revenue with zero manual effort."),
]
col1 = opps[:4]; col2 = opps[4:]
for i, (num, title, desc) in enumerate(col1):
    t = Inches(1.65) + i * Inches(1.43)
    rect(s7, Inches(0.55), t, Inches(6.0), Inches(1.3), DARK_BLUE)
    tb(s7, num, Inches(0.65), t + Inches(0.1), Inches(0.7), Inches(0.55),
       size=20, bold=True, color=WA_GREEN)
    tb(s7, title, Inches(1.35), t + Inches(0.1), Inches(5.0), Inches(0.38),
       size=13, bold=True, color=WHITE)
    tb(s7, desc, Inches(0.65), t + Inches(0.55), Inches(5.7), Inches(0.65),
       size=11, color=LGREY)
for i, (num, title, desc) in enumerate(col2):
    t = Inches(1.65) + i * Inches(1.43)
    rect(s7, Inches(7.0), t, Inches(5.8), Inches(1.3), DARK_BLUE)
    tb(s7, num, Inches(7.1), t + Inches(0.1), Inches(0.7), Inches(0.55),
       size=20, bold=True, color=GOLD)
    tb(s7, title, Inches(7.8), t + Inches(0.1), Inches(4.8), Inches(0.38),
       size=13, bold=True, color=WHITE)
    tb(s7, desc, Inches(7.1), t + Inches(0.55), Inches(5.5), Inches(0.65),
       size=11, color=LGREY)


# ══════════════════════════════════════════════════════════
# SLIDE 8 — ATTENDEE LIFECYCLE MAP
# ══════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
bg(s8)
rect(s8, 0, 0, Inches(0.12), H, ACCENT)
tb(s8, "THE ATTENDEE LIFECYCLE", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s8, "What Happens Today vs. What Should Happen",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s8, Inches(1.35))

stages = [
    ("DISCOVERS SITE", "Lands on homepage\nShark Tank badge → trust\nClean value prop"),
    ("SIGNS UP FREE", "Creates account\nNo credit card\nNo welcome WA message"),
    ("BROWSES EVENTS", "Sees event cards\nSeat counts visible\nVenue hidden — must login"),
    ("BOOKS A TABLE", "Per-seat payment\nNo confirmation WA\nNo calendar link sent"),
    ("ATTENDS DINNER", "Meets 4–8 people\nNo day-of reminder\nVenue sent manually?"),
    ("POST-EVENT", "Goes home\nNo follow-up\nNo connection list\nNo feedback ask"),
]

for i, (stage, notes) in enumerate(stages):
    lx = Inches(0.55) + i * Inches(2.15)
    # stage box
    rect(s8, lx, Inches(1.72), Inches(1.95), Inches(0.42), ACCENT)
    tb(s8, stage, lx, Inches(1.75), Inches(1.95), Inches(0.38),
       size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    rect(s8, lx, Inches(2.18), Inches(1.95), Inches(2.0), DARK_BLUE)
    tb(s8, notes, lx + Inches(0.1), Inches(2.25), Inches(1.78), Inches(1.85),
       size=10.5, color=LGREY, wrap=True)
    if i < 5:
        tb(s8, "→", lx + Inches(2.0), Inches(2.4), Inches(0.2), Inches(0.5),
           size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# Gap labels row
tb(s8, "AUTOMATION GAPS (what fires at each stage with Genos AI):", Inches(0.55), Inches(4.42),
   Inches(12), Inches(0.38), size=11, bold=True, color=WA_GREEN)

gaps = [
    "No action",
    "WA welcome\n+ city table alert",
    "WA seat alert\nwhen match opens",
    "WA confirmation\n+ calendar link",
    "48h reminder\n+ day-of venue\n+ who you'll meet",
    "Connection list\n+ NPS ask\n+ rebook CTA",
]
for i, gap in enumerate(gaps):
    lx = Inches(0.55) + i * Inches(2.15)
    c = RED if gap == "No action" else WA_GREEN
    rect(s8, lx, Inches(4.85), Inches(1.95), Inches(1.85), DARK_BLUE)
    tb(s8, gap, lx + Inches(0.08), Inches(4.95), Inches(1.82), Inches(1.72),
       size=10.5, color=c, wrap=True)


# ══════════════════════════════════════════════════════════
# SLIDE 9 — WHAT WE'D BUILD
# ══════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(BLANK)
bg(s9)
rect(s9, 0, 0, Inches(0.12), H, ACCENT)
tb(s9, "WHAT WE'D BUILD", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s9, "The Genos AI Stack for TwoWordsAway",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s9, Inches(1.35))

builds = [
    ("Fix All Broken Pages + Rewrite Pricing",
     "Build /about (founder story, Shark Tank journey, team), /contact-us (form + WhatsApp + email), /blog (SEO-first content engine), /feedback and /support pages. Rewrite /pricing to accurately describe the per-seat model with a clear value comparison.",
     "Week 1–2"),
    ("WhatsApp Automation Engine — Full Lifecycle",
     "7 automated flows covering every stage from signup to re-engagement: signup nurture, seat alerts, 48h/day-of reminders, post-event follow-up with connection list, NPS collection, no-show recovery, and lapsed-user reactivation. All personalised by city + persona.",
     "Week 2–4"),
    ("Testimonials Engine + Social Proof",
     "Automated post-event WhatsApp request asking attendees to share a quote. Quotes feed into a live testimonials section on the homepage. Over 60 tables hosted, that's 60+ testimonials waiting to be captured. Adds massive trust signal for new visitors.",
     "Week 3–5"),
    ("SEO Blog — Networking & Startup Content",
     "Launch blog with 5 cornerstone articles targeting high-intent India keywords: 'how to meet investors in Bangalore', 'startup networking events Mumbai', 'find a co-founder in Delhi'. Each article links back to the events page with a table-booking CTA.",
     "Week 3–6"),
    ("Conversion Optimisation — CTA + Scarcity Layer",
     "Add sticky 'Book a Table' header, real-time seat counter on homepage, social proof ticker ('14 people joined this week'), and a 'Founders Dinner in your city' personalised CTA block based on user location. Estimated lift: 40-60% more free signups.",
     "Week 4–6"),
]
for i, (title, desc, tl) in enumerate(builds):
    t = Inches(1.62) + i * Inches(1.1)
    rect(s9, Inches(0.55), t, Inches(10.8), Inches(1.0), DARK_BLUE)
    rect(s9, Inches(0.62), t + Inches(0.2), Inches(0.52), Inches(0.52), ACCENT)
    tb(s9, str(i + 1), Inches(0.62), t + Inches(0.18), Inches(0.52), Inches(0.52),
       size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s9, title, Inches(1.28), t + Inches(0.08), Inches(7.5), Inches(0.38),
       size=14, bold=True, color=WHITE)
    tb(s9, desc, Inches(1.28), t + Inches(0.48), Inches(8.2), Inches(0.42),
       size=11, color=LGREY)
    tag(s9, tl, Inches(11.5), t + Inches(0.3), fill=DARK_BLUE, tc=MGREY, size=10)


# ══════════════════════════════════════════════════════════
# SLIDE 10 — NEXT STEPS
# ══════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(BLANK)
bg(s10)
rect(s10, 0, 0, Inches(0.12), H, GOLD)
tb(s10, "NEXT STEPS", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GOLD)
tb(s10, "One Call to See What's Possible",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s10, Inches(1.35), color=GOLD)

steps = [
    ("1", "Reply or Book a 15-Min Call",
     "No pitch. No commitment. We walk you through exactly what we found in the audit and what a fixed, automated version of the site looks like for TwoWordsAway specifically."),
    ("2", "Free Automation Demo",
     "We run a live WhatsApp agent demo using your actual event flow — signup, seat alert, and post-event follow-up — so you can see exactly what your attendees would experience."),
    ("3", "Custom Build Roadmap",
     "We deliver a bespoke implementation plan: which pages to fix first, which automation flows to launch, in what order, and a clear timeline from kickoff to first live WhatsApp message."),
]
for i, (num, title, desc) in enumerate(steps):
    t = Inches(2.0) + i * Inches(1.5)
    rect(s10, Inches(0.55), t, Inches(12.23), Inches(1.3), DARK_BLUE)
    rect(s10, Inches(0.55), t, Inches(1.0), Inches(1.3), GOLD)
    tb(s10, num, Inches(0.55), t + Inches(0.2), Inches(1.0), Inches(0.7),
       size=36, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    tb(s10, title, Inches(1.7), t + Inches(0.15), Inches(10), Inches(0.45),
       size=16, bold=True, color=WHITE)
    tb(s10, desc, Inches(1.7), t + Inches(0.6), Inches(10.5), Inches(0.55),
       size=12.5, color=LGREY)

rect(s10, 0, Inches(6.5), W, Inches(1.0), ACCENT)
tb(s10, "Rohan Malik  ·  CEO, Genos AI  ·  hello@genosai.tech  ·  +91 638-714-2699  ·  www.genosai.tech",
   Inches(0.55), Inches(6.6), Inches(12.23), Inches(0.55),
   size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ── save ──────────────────────────────────────────────────
OUT = "TwoWordsAway_Audit_GenosAI.pptx"
prs.save(OUT)
print(f"Saved: {OUT}  |  10 slides")
