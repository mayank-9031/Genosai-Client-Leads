import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.parent / "utils"))
from ppt_theme import (
    NAVY, DARK_BLUE, ACCENT, GOLD, RED, GREEN, PURPLE, ORANGE, TEAL, CYAN,
    WHITE, LGREY, MGREY, W, H,
    bg, rect, tb, aline, bullets, tag,
    slide_header, slide_footer, score_badge,
    stat_cards, top_bar_card, phase_card, before_after, genos_logo_block,
)

ROOT = Path(__file__).parent.parent.parent
prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ══ SLIDE 1 — COVER ══
s1 = prs.slides.add_slide(BLANK)
bg(s1)
rect(s1, 0, 0, Inches(0.12), H, RED)
rect(s1, Inches(8.5), 0, Inches(4.83), Inches(2.2), DARK_BLUE)
tb(s1, "GENOS AI", Inches(9.2), Inches(0.4), Inches(3.5), Inches(0.8),
   size=28, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
tb(s1, "www.genosai.tech", Inches(9.2), Inches(1.05), Inches(3.5), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
tb(s1, "THE FIRST RESPONDERS / DAVE POLITZER", Inches(0.55), Inches(1.6),
   Inches(9), Inches(0.6), size=13, bold=False, color=RED)
tb(s1, "Website & AI Growth\nAudit Report", Inches(0.55), Inches(2.1),
   Inches(9), Inches(1.8), size=48, bold=True, color=WHITE)
aline(s1, Inches(3.85), color=GOLD, w=Inches(3))
tb(s1, "Prepared exclusively for Dave Politzer, REALTOR / Founder",
   Inches(0.55), Inches(4.15), Inches(9), Inches(0.5), size=14, color=LGREY)
rect(s1, 0, Inches(6.5), W, Inches(1.0), DARK_BLUE)
tb(s1, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699",
   Inches(0.55), Inches(6.6), Inches(9), Inches(0.5), size=12, color=MGREY)
tb(s1, "May 2026  |  CONFIDENTIAL", Inches(9.5), Inches(6.6), Inches(3.3), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
rect(s1, Inches(10.5), Inches(3.3), Inches(2.3), Inches(2.5), DARK_BLUE)
tb(s1, "AUDIT SCORE", Inches(10.55), Inches(3.4), Inches(2.2), Inches(0.5),
   size=11, bold=True, color=MGREY, align=PP_ALIGN.CENTER)
tb(s1, "5", Inches(10.55), Inches(3.75), Inches(2.2), Inches(1.1),
   size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
tb(s1, "/ 10", Inches(10.55), Inches(4.7), Inches(2.2), Inches(0.4),
   size=18, color=LGREY, align=PP_ALIGN.CENTER)
tag(s1, "HIGH PRIORITY", Inches(10.75), Inches(5.1), fill=RED, size=10)


# ══ SLIDE 2 — AT A GLANCE ══
s2 = prs.slides.add_slide(BLANK)
bg(s2)
rect(s2, 0, 0, Inches(0.12), H, RED)
tb(s2, "AT A GLANCE", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=RED)
tb(s2, "Dave Politzer — Who He Is",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s2, Inches(1.35), color=RED)

stats = [
    ("21 Years", "Oakland PD, Sgt & Detective"),
    ("1984", "Started Real Estate"),
    ("1999", "Founded Interloan.com"),
    ("40+", "Years RE + Mortgage"),
    ("5.0", "Stars — 14 Reviews"),
]
cw = Inches(2.3); ch = Inches(2.2); gap = Inches(0.18)
for i, (num, label) in enumerate(stats):
    l = Inches(0.55) + i * (cw + gap)
    rect(s2, l, Inches(1.8), cw, ch, DARK_BLUE)
    tb(s2, num, l + Inches(0.1), Inches(1.95), cw - Inches(0.2), Inches(1.1),
       size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s2, label, l + Inches(0.1), Inches(3.05), cw - Inches(0.2), Inches(0.7),
       size=12, color=LGREY, align=PP_ALIGN.CENTER)

rect(s2, Inches(0.55), Inches(4.2), Inches(12.23), Inches(0.6), DARK_BLUE)
tb(s2, "Danville & East Bay CA  ·  Compass  ·  DRE# 00866434  ·  Specialties: Buyer, Listing, Relocation, New Construction  ·  thefirstresponders.org",
   Inches(0.65), Inches(4.28), Inches(12), Inches(0.45), size=12, color=LGREY, align=PP_ALIGN.CENTER)

tb(s2, "WHY FIRST RESPONDERS TRUST DAVE", Inches(0.55), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=RED)
bullets(s2, ["21 years Oakland PD — the credential no other agent can fake",
             "Understands the financial pressures of public service salaries",
             "Knows what first responders actually need in a home (commute, schedule, lifestyle)"],
        Inches(0.55), Inches(5.35), Inches(5.8), Inches(1.5), size=13, color=LGREY)

tb(s2, "TRANSACTION RANGE", Inches(6.8), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=RED)
bullets(s2, ["Sales range: $1.3M to $3.1M",
             "East Bay focus: Danville, San Ramon, Pleasanton, Alameda County",
             "Mortgage broker license — handles full transaction from search to close"],
        Inches(6.8), Inches(5.35), Inches(5.8), Inches(1.5), size=13, color=LGREY)


# ══ SLIDE 3 — WHAT'S WORKING ══
s3 = prs.slides.add_slide(BLANK)
bg(s3)
rect(s3, 0, 0, Inches(0.12), H, GREEN)
tb(s3, "WHAT'S WORKING", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s3, "Strengths to Build On",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s3, Inches(1.35), color=GREEN)

strengths = [
    ("Irreplaceable First Responder Credibility",
     "21 years in Oakland PD, retiring as a Sergeant and Detective, is a credential no competitor can manufacture. When a police officer, firefighter, or paramedic meets Dave, they're talking to someone who did the job. That trust converts."),
    ("Niche Positioning with a Dedicated Domain",
     "thefirstresponders.org gives Dave a standalone identity beyond just being a Compass agent. A dedicated program site signals commitment to the community, not just a checkbox on a Compass profile that says 'first responders welcome.'"),
    ("Early Tech Adopter Track Record",
     "Dave built Interloan.com in 1999 - one of the first online multi-lender mortgage platforms. He sold it to a public company in 2000. That instinct to adopt tech early is exactly what makes AI automation a natural next step, not a foreign concept."),
    ("Dual License Advantage",
     "Most agents are just agents. Dave handles both real estate and mortgage brokering, which means he controls more of the transaction, can offer better guidance on first responder loan programs, and builds deeper client relationships."),
    ("5.0 Stars Across 14 Reviews",
     "Perfect rating from 14 clients signals a high-satisfaction, repeat and referral business. The first responder community is tight-knit - one satisfied police officer or firefighter who refers within their department is worth five cold leads."),
]
for i, (title, desc) in enumerate(strengths):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s3, Inches(0.55), t, Inches(0.06), Inches(0.7), GREEN)
    tb(s3, title, Inches(0.75), t, Inches(11.8), Inches(0.38), size=14, bold=True, color=WHITE)
    tb(s3, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# ══ SLIDE 4 — WHERE LEADS ARE LEAKING ══
s4 = prs.slides.add_slide(BLANK)
bg(s4)
rect(s4, 0, 0, Inches(0.12), H, RED)
tb(s4, "WHERE LEADS ARE LEAKING", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=RED)
tb(s4, "The Niche Is Perfect. The Digital Infrastructure Isn't.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=30, bold=True, color=WHITE)
aline(s4, Inches(1.35), color=RED)

issues = [
    ("NO AI CHAT FOR AFTER-HOURS FIRST RESPONDER BUYERS", "HIGH",
     "First responders work nights, weekends, 48-hour shifts. The buyer who's ready is browsing thefirstresponders.org at 2am after a night shift, not at 10am. With no chat widget, that visitor either fills a contact form and waits, or goes to Zillow. There's no real-time engagement path built for the schedule Dave's clients actually keep."),
    ("NO AUTOMATED FOLLOW-UP FOR INQUIRIES", "HIGH",
     "A first responder who fills out a contact form on a Friday night won't hear back until Monday. They may be back on shift by then. Automated follow-up within 30 minutes of an inquiry - even just an acknowledgment with next steps - keeps that lead warm through the weekend."),
    ("NO VIDEO TESTIMONIALS OR AUTOMATED REVIEW COLLECTION", "HIGH",
     "The site has three first responder testimonials - Dave S. from Pleasanton PD, V. Bullock from Oakland PD, R. Harmon from the DA's office. That's real social proof. But static text testimonials don't compound. With video testimonials and post-closing review automation that asks clients to mention their profession and department, the peer proof builds with every transaction."),
    ("NO AUTOMATED MARKET REPORTS TO FIRST RESPONDER DATABASE", "MEDIUM",
     "Dave has served police, fire, and EMS buyers across the East Bay for 40 years. That's a database. Automated monthly market reports sent to that list - customized for their areas and price ranges - keep Dave top of mind for the moment a first responder decides it's time to buy or sell."),
    ("PROGRAM BENEFITS NOT CLEARLY EXPLAINED ON SITE", "MEDIUM",
     "First responder home buying programs (CalHFA, HUD Good Neighbor Next Door, local grants) are complex. Most buyers don't know what they qualify for. Dedicated landing pages by profession - police, fire, EMS, dispatch - that clearly explain the benefits and have a calculator would convert research visits into appointments."),
]
for i, (title, pri, desc) in enumerate(issues):
    t = Inches(1.7) + i * Inches(1.05)
    pc = RED if pri == "HIGH" else GOLD
    rect(s4, Inches(0.55), t, Inches(0.06), Inches(0.7), pc)
    tb(s4, title, Inches(0.75), t, Inches(9.2), Inches(0.38), size=12, bold=True, color=WHITE)
    tag(s4, pri, Inches(10.1), t + Inches(0.05), fill=pc, size=9)
    tb(s4, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# ══ SLIDE 5 — WEBSITE DESIGN AUDIT ══
s5 = prs.slides.add_slide(BLANK)
bg(s5)
rect(s5, 0, 0, Inches(0.12), H, PURPLE)
tb(s5, "WEBSITE DESIGN AUDIT", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=PURPLE)
tb(s5, "A Niche This Strong Deserves a Site That Matches.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.8), size=32, bold=True, color=WHITE)
aline(s5, Inches(1.45), color=PURPLE, w=Inches(1.5))

rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(3.6), DARK_BLUE)
rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(0.42), RED)
tb(s5, "THEFIRSTRESPONDERS.ORG TODAY", Inches(0.65), Inches(1.92), Inches(3.6), Inches(0.38),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Good content: benefits, stats, 3 FR testimonials",
    "Static contact form - no automated follow-up",
    "No AI chat for 2am shift-worker browsers",
    "No video testimonials or review automation",
    "No department-specific pages (Police/Fire/EMS)",
    "No program benefit calculator by profession",
    "No market report pipeline to FR database",
    "No department referral landing pages",
], Inches(0.7), Inches(2.42), Inches(3.6), Inches(2.9), size=11, color=LGREY, dot=RED)

tb(s5, "->", Inches(4.45), Inches(3.6), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(3.6), DARK_BLUE)
rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(0.42), GOLD)
tb(s5, "TOP NICHE REAL ESTATE SITES DO THIS", Inches(5.1), Inches(1.92), Inches(3.4), Inches(0.38),
   size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Homes for Heroes: profession-specific pages with calculators",
    "Testimonials from buyers in the same profession",
    "Clear program benefit comparison by profession",
    "AI chat or chatbot for after-hours inquiry capture",
    "Automated email sequences post-contact",
    "Market reports segmented by buyer type",
    "Video testimonials from satisfied clients",
    "Community pages for local first responder areas",
], Inches(5.1), Inches(2.42), Inches(3.4), Inches(2.9), size=11, color=LGREY, dot=GOLD)

tb(s5, "->", Inches(8.68), Inches(3.6), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(3.6), DARK_BLUE)
rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(0.42), GREEN)
tb(s5, "WHAT GENOS AI BUILDS", Inches(9.35), Inches(1.92), Inches(3.4), Inches(0.38),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Dedicated pages: Police / Fire / EMS / Dispatch",
    "First responder program benefit calculator",
    "Peer testimonials from fellow officers and firefighters",
    "AI chat trained on first responder programs 24/7",
    "Automated follow-up within 30 min of any contact form",
    "Market reports auto-sent to Dave's FR database monthly",
    "Video testimonials section auto-populating post-close",
    "Department referral landing pages for OPD, LAFD, SFFD",
], Inches(9.35), Inches(2.42), Inches(3.4), Inches(2.9), size=11, color=LGREY, dot=GREEN)

rect(s5, Inches(0.55), Inches(5.72), Inches(12.23), Inches(1.28), DARK_BLUE)
tb(s5, "WHY IT MATTERS:", Inches(0.7), Inches(5.78), Inches(2.2), Inches(0.38),
   size=11, bold=True, color=PURPLE)
stats2 = [
    ("67%", "of home searches start\nbetween 8pm and 8am"),
    ("5x", "referral rate within\ntight-knit peer communities"),
    ("30 min", "response time = 10x\nmore likely to convert"),
    ("200%+", "lift in conversions with\nniche-specific landing pages"),
]
sw = Inches(2.7); sl = Inches(2.9)
for i, (num, label) in enumerate(stats2):
    xl = sl + i * sw
    tb(s5, num, xl, Inches(5.75), Inches(2.5), Inches(0.55),
       size=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s5, label, xl, Inches(6.22), Inches(2.5), Inches(0.7),
       size=10, color=LGREY, align=PP_ALIGN.CENTER)


# ══ SLIDE 6 — AI AUTOMATION OPPORTUNITIES ══
s6 = prs.slides.add_slide(BLANK)
bg(s6)
rect(s6, 0, 0, Inches(0.12), H, ACCENT)
tb(s6, "AI AUTOMATION OPPORTUNITIES", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s6, "Built for a Buyer Who Works When You're Asleep.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=34, bold=True, color=WHITE)
aline(s6, Inches(1.35), color=ACCENT)

opps = [
    ("AI Chat — 24/7\nFirst Responder Intake", RED,
     "A chat widget trained on all first responder home buying programs. CalHFA, Good Neighbor Next Door, local East Bay grants, VA loans for vets who are also first responders. Captures the 2am browser. Routes qualified leads to Dave by morning."),
    ("Lead Follow-Up\nSequences", ACCENT,
     "Contact form submitted at midnight on a Saturday? Automated acknowledgment within 30 minutes, follow-up sequence every 48 hours until Dave can personally connect. Built around shift schedules - not a 9-5 email cadence."),
    ("Market Report\nAutomation", GREEN,
     "Monthly automated reports to Dave's first responder database. Segmented by location and price range - Oakland buyer sees Alameda County data, San Ramon buyer sees Contra Costa. Keeps Dave top of mind before the first responder decides it's time to move."),
    ("Peer Review\nAutomation", PURPLE,
     "30 days after closing, automated review request goes to the client. Specifically asks them to mention their profession in the review - because a review from 'Officer Martinez from OPD' converts 5x better with other first responders than a review from 'a satisfied buyer.'"),
    ("Department\nReferral Sequences", GOLD,
     "When a first responder closes with Dave, an automated sequence goes out asking for referrals within their department. Pre-written message they can forward. Targets the exact community where Dave has the strongest trust advantage and where referrals cluster."),
]
cw2 = Inches(2.35); gap2 = Inches(0.14)
for i, (title, color, desc) in enumerate(opps):
    l2 = Inches(0.55) + i * (cw2 + gap2)
    rect(s6, l2, Inches(1.8), cw2, Inches(0.06), color)
    rect(s6, l2, Inches(1.86), cw2, Inches(4.6), DARK_BLUE)
    tb(s6, title, l2 + Inches(0.12), Inches(1.98), cw2 - Inches(0.24), Inches(0.65),
       size=13, bold=True, color=color)
    tb(s6, desc, l2 + Inches(0.12), Inches(2.65), cw2 - Inches(0.24), Inches(3.6),
       size=11, color=LGREY)

rect(s6, Inches(0.55), Inches(6.7), Inches(12.23), Inches(0.45), DARK_BLUE)
tb(s6, "Dave built Interloan.com in 1999 before any agent was online. The instinct hasn't changed - it's just AI now instead of the internet.",
   Inches(0.7), Inches(6.75), Inches(12), Inches(0.35), size=11, italic=True, color=GOLD)


# ══ SLIDE 7 — WHAT THIS IS COSTING YOU ══
s7 = prs.slides.add_slide(BLANK)
bg(s7)
rect(s7, 0, 0, Inches(0.12), H, RED)
tb(s7, "WHAT THIS IS COSTING YOU", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=RED)
tb(s7, "The Niche Is There. The Leads Are Leaking.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s7, Inches(1.35), color=RED)

costs = [
    ("No after-hours chat = every 2am first responder browser goes somewhere else",
     "Most first responders don't browse homes at noon. They browse after shifts, on nights and weekends. Without AI chat on thefirstresponders.org, every visitor who has a question at midnight gets silence. That lead goes to Zillow, Redfin, or a competitor who has a bot running."),
    ("No automated follow-up = leads going cold on 48-hour shifts",
     "A paramedic who fills out a contact form on a Sunday morning and goes on shift won't be reachable until Tuesday. If Dave doesn't respond until Monday and they're unavailable until Tuesday, the lead is cold by Wednesday. Automated follow-up keeps the conversation alive regardless of shift schedule."),
    ("Text testimonials don't compound — video and review automation do",
     "Dave has three genuine first responder testimonials on the site. That's more than most. But static text from 2024 doesn't grow. With automated post-closing review requests that ask clients to mention their profession and department, a new profession-specific testimonial gets added after every transaction. Video converts 3x better than text in trust-based niches."),
    ("No market reports = out of sight before buyers are ready",
     "A first responder thinking about buying in 6 months is not actively searching right now. They're forming an opinion about which agent to call when the time comes. Monthly market reports keep Dave's name in front of that decision. Without them, whoever sends the most relevant content wins."),
    ("No program landing pages = confused visitors who don't convert",
     "First responder home buying programs are genuinely confusing. Different programs for police vs. fire vs. EMS. Federal programs vs. CalHFA vs. local grants. A visitor who can't figure out what they qualify for will not schedule an appointment - they'll keep researching. Clear pages by profession with calculators turn that research into a call."),
]
for i, (title, desc) in enumerate(costs):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s7, Inches(0.55), t, Inches(0.06), Inches(0.7), RED)
    tb(s7, title, Inches(0.75), t, Inches(11.8), Inches(0.38), size=14, bold=True, color=WHITE)
    tb(s7, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# ══ SLIDE 8 — WHAT WE'D BUILD ══
s8 = prs.slides.add_slide(BLANK)
bg(s8)
rect(s8, 0, 0, Inches(0.12), H, GREEN)
tb(s8, "WHAT WE'D BUILD", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s8, "Genos AI for The First Responders",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s8, Inches(1.35), color=GREEN)

builds = [
    ("thefirstresponders.org Full Rebuild",
     "Dedicated landing pages by profession: Police, Fire, EMS, Dispatch. Each explains the specific programs available for that role, with a calculator. Peer testimonial sections with names and departments visible. Dave's OPD background front and center, not buried in a bio."),
    ("AI Chat — 24/7 First Responder Intake",
     "Widget trained on all first responder home buying programs - CalHFA, Good Neighbor Next Door, East Bay grants, VA overlaps for vets. Captures inquiries at 2am on a Tuesday. Routes the conversation to Dave by morning with context already gathered."),
    ("Lead Follow-Up & Shift-Aware Sequences",
     "When a first responder submits a contact form, automated sequence activates. Acknowledgment within 30 minutes. Follow-up cadence adjusted for the schedule Dave's clients keep - not a corporate 9-to-5 email sequence. Keeps the lead warm through a 48-hour shift."),
    ("Automated Market Reports + Peer Review Automation",
     "Monthly market reports segmented by buyer location and price range, auto-sent to Dave's first responder database. Post-closing, automated review requests that specifically prompt clients to mention their profession and department - the testimonials that convert most in this niche."),
    ("Department Referral Landing Pages",
     "Dedicated pages for specific departments - OPD, SFFD, Contra Costa EMS. When a satisfied client refers within their department, the referral link goes to a page that speaks directly to that department's specific programs. Warm transfer, not a cold website visit."),
]
for i, (title, desc) in enumerate(builds):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s8, Inches(0.55), t, Inches(0.06), Inches(0.7), GREEN)
    tb(s8, title, Inches(0.75), t, Inches(11.8), Inches(0.38), size=14, bold=True, color=WHITE)
    tb(s8, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# ══ SLIDE 9 — 90-DAY OUTCOMES ══
s9 = prs.slides.add_slide(BLANK)
bg(s9)
rect(s9, 0, 0, Inches(0.12), H, ACCENT)
tb(s9, "90-DAY OUTCOMES", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s9, "What Changes in 90 Days",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s9, Inches(1.35))

phases = [
    ("Month 1", GOLD, [
        "AI chat live on thefirstresponders.org - capturing 2am first responder browsers",
        "Automated lead follow-up sequences active - no inquiry goes cold over a weekend shift",
        "Market report pipeline set up - first automated send to Dave's FR database",
        "Site rebuild begins - profession-specific pages and program calculators in design",
    ]),
    ("Month 2", ACCENT, [
        "thefirstresponders.org live with profession-specific pages (Police/Fire/EMS/Dispatch)",
        "Peer review automation running - post-closing prompts generating profession-identified reviews",
        "Department referral pages live for OPD, SFFD, Contra Costa EMS",
        "Second market report send - tracking opens and click-throughs by segment",
    ]),
    ("Month 3", GREEN, [
        "AI chat conversion data: after-hours inquiries captured vs. prior period",
        "Review velocity measurably up with profession-specific peer testimonials accumulating",
        "Referral sequences driving qualified leads from within first responder departments",
        "Market reports building top-of-mind presence with the 6-month-out buyer pool",
    ]),
]
pw = Inches(3.9); gap3 = Inches(0.18)
for i, (month, color, pts) in enumerate(phases):
    l3 = Inches(0.55) + i * (pw + gap3)
    rect(s9, l3, Inches(1.8), pw, Inches(0.42), color)
    tb(s9, month.upper(), l3 + Inches(0.1), Inches(1.84), pw - Inches(0.2), Inches(0.36),
       size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    rect(s9, l3, Inches(2.22), pw, Inches(4.7), DARK_BLUE)
    bullets(s9, pts, l3 + Inches(0.15), Inches(2.35), pw - Inches(0.3), Inches(4.4),
            size=12, color=LGREY, dot=color)


# ══ SLIDE 10 — NEXT STEPS ══
s10 = prs.slides.add_slide(BLANK)
bg(s10)
rect(s10, 0, 0, Inches(0.12), H, GOLD)
rect(s10, Inches(7.5), 0, Inches(5.83), H, DARK_BLUE)
tb(s10, "NEXT STEPS", Inches(0.55), Inches(0.3), Inches(6.5), Inches(0.5),
   size=11, bold=True, color=GOLD)
tb(s10, "Let's Talk for\n15 Minutes.",
   Inches(0.55), Inches(0.8), Inches(6.8), Inches(2.0), size=48, bold=True, color=WHITE)
aline(s10, Inches(2.8), color=GOLD, w=Inches(2.5))
tb(s10, "No pitch deck required. Reply or book a call directly.\nWe'll walk through what fits The First Responders program.",
   Inches(0.55), Inches(3.05), Inches(6.5), Inches(0.8), size=14, color=LGREY)

steps = [
    "Reply to this email - we'll set up a time",
    "Book directly: www.genosai.tech/call",
    "Walk through the audit findings in 15 minutes",
    "You decide what, if anything, makes sense to move on",
]
bullets(s10, steps, Inches(0.55), Inches(3.95), Inches(6.5), Inches(2.0), size=14, color=WHITE, dot=GOLD)

rect(s10, 0, Inches(6.5), Inches(7.5), Inches(1.0), DARK_BLUE)
tb(s10, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699  |  www.genosai.tech",
   Inches(0.55), Inches(6.6), Inches(6.8), Inches(0.5), size=12, color=MGREY)

tb(s10, "WHAT WE BUILD FOR\nFIRST RESPONDERS", Inches(7.85), Inches(0.6), Inches(5.0), Inches(0.9),
   size=18, bold=True, color=GOLD)
aline(s10, Inches(1.45), color=GOLD)

summary_items = [
    ("Site Rebuild", "Profession pages: Police / Fire / EMS / Dispatch + calculators"),
    ("AI Chat", "24/7 first responder program intake - catches the 2am browser"),
    ("Follow-Up", "Shift-aware sequences - no lead cold over a 48-hour shift"),
    ("Market Reports", "Monthly auto-send to Dave's first responder database"),
    ("Peer Reviews", "Post-closing automation prompts profession-identified testimonials"),
    ("Referrals", "Department referral pages for OPD, SFFD, Contra Costa EMS"),
]
for i, (label, val) in enumerate(summary_items):
    t2 = Inches(1.6) + i * Inches(0.88)
    rect(s10, Inches(7.85), t2, Inches(1.5), Inches(0.72), NAVY)
    tb(s10, label, Inches(7.9), t2 + Inches(0.05), Inches(1.4), Inches(0.32),
       size=10, bold=True, color=GOLD)
    tb(s10, val, Inches(9.5), t2 + Inches(0.05), Inches(3.6), Inches(0.65),
       size=11, color=LGREY)


prs.save(ROOT / "output" / "decks" / "FirstResponders_Audit_GenosAI.pptx")
print("Saved: FirstResponders_Audit_GenosAI.pptx")
