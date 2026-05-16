import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.parent / "utils"))
from ppt_theme import (
    NAVY, DARK_BLUE, ACCENT, GOLD, RED, GREEN, PURPLE, ORANGE, TEAL, CYAN,
    WHITE, LGREY, MGREY, W, H,
    bg, rect, tb, aline, bullets, tag,
    slide_header, slide_footer, score_badge,
    stat_cards, top_bar_card, phase_card, before_after, genos_logo_block,
)

ROOT = Path(__file__).parent.parent.parent
prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── SLIDE 1 — COVER ──────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(BLANK)
bg(s1)
rect(s1, 0, 0, Inches(0.12), H, GOLD)
rect(s1, Inches(8.5), 0, Inches(4.83), Inches(2.2), DARK_BLUE)
tb(s1, "GENOS AI", Inches(9.2), Inches(0.4), Inches(3.5), Inches(0.8),
   size=28, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)
tb(s1, "www.genosai.tech", Inches(9.2), Inches(1.05), Inches(3.5), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
tb(s1, "FIRSTPATHWAY PARTNERS — MILWAUKEE, WI", Inches(0.55), Inches(1.6),
   Inches(9), Inches(0.6), size=13, color=GOLD)
tb(s1, "AI Automation & Digital\nPresence Audit", Inches(0.55), Inches(2.05),
   Inches(9), Inches(1.8), size=44, bold=True, color=WHITE)
aline(s1, Inches(3.85), color=GOLD, w=Inches(3))
tb(s1, "Prepared exclusively for Asif Chhipa — America's Trusted EB-5 Immigration Investment Specialists",
   Inches(0.55), Inches(4.15), Inches(9), Inches(0.5), size=13, color=LGREY)
rect(s1, 0, Inches(6.5), W, Inches(1.0), DARK_BLUE)
tb(s1, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699",
   Inches(0.55), Inches(6.6), Inches(9), Inches(0.5), size=12, color=MGREY)
tb(s1, "May 2026  |  CONFIDENTIAL", Inches(9.5), Inches(6.6), Inches(3.3), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
rect(s1, Inches(10.5), Inches(3.3), Inches(2.3), Inches(2.5), DARK_BLUE)
tb(s1, "AUDIT SCORE", Inches(10.55), Inches(3.4), Inches(2.2), Inches(0.5),
   size=11, bold=True, color=MGREY, align=PP_ALIGN.CENTER)
tb(s1, "4", Inches(10.55), Inches(3.75), Inches(2.2), Inches(1.1),
   size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
tb(s1, "/ 10", Inches(10.55), Inches(4.7), Inches(2.2), Inches(0.4),
   size=18, color=LGREY, align=PP_ALIGN.CENTER)
tag(s1, "HIGH PRIORITY", Inches(10.75), Inches(5.1), fill=RED, tc=WHITE, size=10)


# ── SLIDE 2 — AT A GLANCE ────────────────────────────────────────────────────
s2 = prs.slides.add_slide(BLANK)
bg(s2)
rect(s2, 0, 0, Inches(0.12), H, GOLD)
tb(s2, "AT A GLANCE", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GOLD)
tb(s2, "FirstPathway Partners — EB-5 Immigration Specialists",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=30, bold=True, color=WHITE)
aline(s2, Inches(1.35), color=GOLD)

stats = [
    ("$800K+", "Minimum investor\ncapital required"),
    ("6 Countries", "India, UAE, China,\nDenmark, Venezuela, Hungary"),
    ("34+ yrs", "Immigration law\nexperience (Wilka Toppins)"),
    ("2 Active", "Projects: Ameyalli\n& Kaia, Utah"),
    ("I-829", "Approval track\nrecord — full lifecycle"),
]
cw = Inches(2.3); ch = Inches(2.2); gap = Inches(0.18)
for i, (num, label) in enumerate(stats):
    l = Inches(0.55) + i * (cw + gap)
    rect(s2, l, Inches(1.8), cw, ch, DARK_BLUE)
    tb(s2, num, l + Inches(0.1), Inches(1.95), cw - Inches(0.2), Inches(1.1),
       size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s2, label, l + Inches(0.1), Inches(3.05), cw - Inches(0.2), Inches(0.7),
       size=11, color=LGREY, align=PP_ALIGN.CENTER)

rect(s2, Inches(0.55), Inches(4.2), Inches(12.23), Inches(0.6), DARK_BLUE)
tb(s2, "Milwaukee, WI  |  info@firstpathway.com  |  +1 (414) 431-0742  |  firstpathway.com  |  WeChat · Weibo · LinkedIn · Facebook",
   Inches(0.65), Inches(4.28), Inches(12), Inches(0.45), size=11, color=LGREY, align=PP_ALIGN.CENTER)

tb(s2, "SERVICES OFFERED", Inches(0.55), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=GOLD)
bullets(s2, [
    "EB-5 visa application processing and USCIS filing guidance",
    "Investment project placement (Ameyalli, Kaia — Utah)",
    "I-526E and I-829 documentation preparation and submission",
    "Legal support via 34-year immigration law specialist",
    "Multilingual consultation: English, Mandarin, regional support"],
        Inches(0.55), Inches(5.35), Inches(5.8), Inches(1.6), size=12, color=LGREY)

tb(s2, "INVESTOR ORIGINS SERVED", Inches(6.8), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=GOLD)
bullets(s2, [
    "India — fastest-growing EB-5 applicant country globally",
    "UAE — high-net-worth Gulf investors seeking US residency",
    "China — historically #1 EB-5 market, WeChat-first audience",
    "Europe: Denmark, Hungary and other European nationals",
    "Latin America: Venezuela and South American investors"],
        Inches(6.8), Inches(5.35), Inches(5.8), Inches(1.6), size=12, color=LGREY)


# ── SLIDE 3 — WHAT GENOS AI DOES ─────────────────────────────────────────────
s3 = prs.slides.add_slide(BLANK)
bg(s3)
rect(s3, 0, 0, Inches(0.12), H, ACCENT)
tb(s3, "ABOUT GENOS AI", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s3, "AI Automation for Immigration, Investment & Professional Services.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=28, bold=True, color=WHITE)
aline(s3, Inches(1.35), color=ACCENT)

tb(s3, "What We Build", Inches(0.55), Inches(1.6), Inches(12.5), Inches(0.4),
   size=14, bold=True, color=ACCENT)

services = [
    ("Multilingual AI\nChatbot Agents", TEAL,
     "AI agents that speak English, Hindi, Mandarin, and Arabic — live on your website, WhatsApp, and WeChat. Qualify investors by net worth, investment timeline, and country. The most effective entry channel for Indian and Chinese HNI families exploring EB-5."),
    ("Client Journey\nAutomation", GOLD,
     "Automated milestone communications across the 3-5 year EB-5 process: I-526E filed, I-526E approved, DS-260 reminder, visa interview prep, I-829 filing reminder, congratulations on approval. No investor goes silent for months. Every stage is acknowledged."),
    ("HNI Lead\nQualification Engine", PURPLE,
     "Score, qualify, and route inbound investor inquiries in real time. Collect key signals: net worth threshold, investment timeline, country, visa history. Route serious investors to Asif immediately. Filter early-stage inquiries into a nurture sequence."),
    ("Content & SEO\nOutreach Engine", ORANGE,
     "LinkedIn thought leadership, blog content targeting 'EB-5 visa India', 'US residency investment program', 'EB-5 investor 2025' — automated and published weekly. Email sequences in Hindi, Mandarin, and Arabic for investor databases. Webinar registration and follow-up automation."),
]
cw2 = Inches(2.95); gap2 = Inches(0.16)
for i, (title, color, desc) in enumerate(services):
    l2 = Inches(0.55) + i * (cw2 + gap2)
    rect(s3, l2, Inches(2.1), cw2, Inches(0.06), color)
    rect(s3, l2, Inches(2.16), cw2, Inches(3.5), DARK_BLUE)
    tb(s3, title, l2 + Inches(0.12), Inches(2.28), cw2 - Inches(0.24), Inches(0.6),
       size=13, bold=True, color=color)
    tb(s3, desc, l2 + Inches(0.12), Inches(2.9), cw2 - Inches(0.24), Inches(2.6),
       size=11, color=LGREY)

rect(s3, Inches(0.55), Inches(5.85), Inches(12.23), Inches(1.3), DARK_BLUE)
tb(s3, "WHY GENOS AI FOR FIRSTPATHWAY", Inches(0.7), Inches(5.92), Inches(5.5), Inches(0.35),
   size=11, bold=True, color=ACCENT)
bullets(s3, [
    "Immigration expertise: built for document-heavy, multi-year client journeys",
    "Multilingual by default: Hindi, Mandarin, Arabic — not an add-on",
    "High-value client focus: every system tuned for $800K+ decision makers",
    "ROI-clear: one additional qualified investor per quarter justifies the entire program"],
        Inches(0.7), Inches(6.28), Inches(5.8), Inches(0.8), size=11, color=LGREY, dot=ACCENT)

tb(s3, "PROVEN DELIVERY", Inches(7.0), Inches(5.92), Inches(5.5), Inches(0.35),
   size=11, bold=True, color=ACCENT)
bullets(s3, [
    "23+ clients across USA, UK, UAE, India, Australia — 98% retention",
    "50+ production AI systems shipped — not demos, not pilots",
    "4-12 week delivery: production-grade from day one",
    "Integrates with your existing CRM, email, and communication stack"],
        Inches(7.0), Inches(6.28), Inches(5.8), Inches(0.8), size=11, color=LGREY, dot=ACCENT)


# ── SLIDE 4 — THE THREE GAPS ─────────────────────────────────────────────────
s4 = prs.slides.add_slide(BLANK)
bg(s4)
rect(s4, 0, 0, Inches(0.12), H, RED)
tb(s4, "THE THREE GAPS", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=RED)
tb(s4, "34 Years of Immigration Experience. Investors Still Falling Through the Cracks.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=26, bold=True, color=WHITE)
aline(s4, Inches(1.35), color=RED)

gap_cards = [
    ("GAP 1\nNO MULTILINGUAL\nDIGITAL CHANNEL", RED,
     "Serving 6 countries — in English only",
     [
         "India is now the #1 EB-5 applicant country — their HNI investors prefer Hindi WhatsApp conversations, not English contact forms",
         "Chinese investors — your historical top market — use WeChat exclusively; the website has no WeChat bot or Mandarin chat support",
         "UAE investors browse in Arabic; a generic 'fill the form' experience loses them to competitors with Arabic-first digital presence",
         "FirstPathway has WeChat and Weibo accounts but no automated AI agent behind those channels to qualify investors 24/7",
     ]),
    ("GAP 2\nMANUAL CLIENT\nJOURNEY (3-5 YRS)", ORANGE,
     "No automated touchpoints across a multi-year lifecycle",
     [
         "The EB-5 process spans 3-5 years: I-526E → DS-260 → visa → conditional green card → I-829 → permanent residency",
         "Investors waiting for USCIS updates go silent for months with no proactive communication from the firm",
         "No automated milestone emails/WhatsApp when key stages are reached or documents are due",
         "Manual follow-up for every renewal reminder, document collection, and status update — not scalable at volume",
     ]),
    ("GAP 3\nWEAK HNI LEAD\nGENERATION", GOLD,
     "No digital funnel for the world's fastest-growing EB-5 market",
     [
         "Google searches for 'EB-5 visa India 2025', 'US residency investment program', 'green card for investors' return no FirstPathway content",
         "No lead capture tool on the website — Indian and UAE visitors read the page and leave with no follow-up mechanism",
         "No LinkedIn content engine reaching Indian HNI families, wealth managers, or immigration attorneys who refer clients",
         "One qualified EB-5 investor represents $800K+ in investment and years of recurring advisory relationship",
     ]),
]
cw3 = Inches(3.95); gap3i = Inches(0.18)
for i, (title, color, subtitle, pts) in enumerate(gap_cards):
    l = Inches(0.55) + i * (cw3 + gap3i)
    rect(s4, l, Inches(1.8), cw3, Inches(0.5), color)
    tb(s4, title, l + Inches(0.12), Inches(1.85), cw3 - Inches(0.24), Inches(0.42),
       size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    rect(s4, l, Inches(2.3), cw3, Inches(4.6), DARK_BLUE)
    tb(s4, subtitle, l + Inches(0.12), Inches(2.38), cw3 - Inches(0.24), Inches(0.35),
       size=10, italic=True, color=color)
    bullets(s4, pts, l + Inches(0.12), Inches(2.78), cw3 - Inches(0.24), Inches(3.9),
            size=11, color=LGREY, dot=color)

rect(s4, Inches(0.55), Inches(7.05), Inches(12.23), Inches(0.3), DARK_BLUE)
tb(s4, "FirstPathway has the legal expertise, the track record, and the projects. What it lacks is a digital system that finds, qualifies, and retains investors at the speed the market now demands.",
   Inches(0.7), Inches(7.08), Inches(12), Inches(0.25), size=11, italic=True, color=GOLD)


# ── SLIDE 5 — WEBSITE AUDIT ──────────────────────────────────────────────────
s5 = prs.slides.add_slide(BLANK)
bg(s5)
rect(s5, 0, 0, Inches(0.12), H, PURPLE)
tb(s5, "WEBSITE & DIGITAL AUDIT", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=PURPLE)
tb(s5, "Strong Legal Credentials. Underperforming Digital Presence.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.8), size=28, bold=True, color=WHITE)
aline(s5, Inches(1.45), color=PURPLE, w=Inches(1.5))

rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(3.6), DARK_BLUE)
rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(0.42), RED)
tb(s5, "FIRSTPATHWAY.COM TODAY", Inches(0.65), Inches(1.92), Inches(3.6), Inches(0.38),
   size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "English-only website — no multilingual content or chatbot",
    "Static contact form — no automated investor intake workflow",
    "No WhatsApp or WeChat AI bot behind social accounts",
    "No investor qualification tool or ROI calculator",
    "No automated status update system for active clients",
    "No SEO content targeting EB-5 investor search terms",
    "No LinkedIn content engine for HNI investor outreach",
    "No webinar scheduling or automated consultation booking",
], Inches(0.7), Inches(2.42), Inches(3.6), Inches(2.9), size=11, color=LGREY, dot=RED)

tb(s5, "->", Inches(4.45), Inches(3.6), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(3.6), DARK_BLUE)
rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(0.42), GOLD)
tb(s5, "WHAT TOP EB-5 FIRMS DO", Inches(5.1), Inches(1.92), Inches(3.4), Inches(0.38),
   size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Hindi / Mandarin / Arabic chatbot on homepage",
    "WeChat AI bot qualifying Chinese investors 24/7",
    "Automated investor intake with net worth qualification",
    "Milestone email sequences throughout the 3-5 year journey",
    "EB-5 ROI calculator capturing serious investor contact",
    "SEO content ranking for 'EB-5 visa India' terms",
    "LinkedIn articles reaching wealth managers and attorneys",
    "Webinar series with automatic registration and follow-up",
], Inches(5.1), Inches(2.42), Inches(3.4), Inches(2.9), size=11, color=LGREY, dot=GOLD)

tb(s5, "->", Inches(8.68), Inches(3.6), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(3.6), DARK_BLUE)
rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(0.42), GREEN)
tb(s5, "WHAT GENOS AI BUILDS", Inches(9.35), Inches(1.92), Inches(3.4), Inches(0.38),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Multilingual chatbot: Hindi + Mandarin + Arabic + English",
    "WeChat AI agent: qualify Chinese investors on their platform",
    "Automated investor intake: 8-question qualification flow",
    "Client journey emails: I-526E → DS-260 → I-829 milestones",
    "EB-5 ROI / residency benefit calculator with lead capture",
    "SEO content: 12+ articles targeting EB-5 India/UAE/China",
    "LinkedIn content engine: 3 posts/week for HNI audience",
    "Webinar automation: registration, reminders, replay email",
], Inches(9.35), Inches(2.42), Inches(3.4), Inches(2.9), size=11, color=LGREY, dot=GREEN)

rect(s5, Inches(0.55), Inches(5.72), Inches(12.23), Inches(1.28), DARK_BLUE)
tb(s5, "DIGITAL BENCHMARKS:", Inches(0.7), Inches(5.78), Inches(2.8), Inches(0.38),
   size=11, bold=True, color=PURPLE)
bstats = [
    ("#1 Market", "India: fastest-growing\nEB-5 applicant country"),
    ("$800K+", "Revenue per investor —\njustifies AI at any cost"),
    ("3-5 Years", "Client lifecycle with\nzero automated touchpoints"),
    ("70%+", "HNI investors research\nonline before contacting"),
]
sw = Inches(2.55); sl = Inches(3.65)
for i, (num, label) in enumerate(bstats):
    xl = sl + i * sw
    tb(s5, num, xl, Inches(5.75), Inches(2.4), Inches(0.55),
       size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s5, label, xl, Inches(6.22), Inches(2.4), Inches(0.7),
       size=10, color=LGREY, align=PP_ALIGN.CENTER)


# ── SLIDE 6 — AI AUTOMATION STACK ────────────────────────────────────────────
s6 = prs.slides.add_slide(BLANK)
bg(s6)
rect(s6, 0, 0, Inches(0.12), H, ACCENT)
tb(s6, "AI AUTOMATION STACK", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s6, "Four Systems Built for an EB-5 Immigration Firm.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=28, bold=True, color=WHITE)
aline(s6, Inches(1.35), color=ACCENT)

opps = [
    ("Multilingual\nInvestor AI Agent", TEAL,
     "An AI agent deployed on your website, WhatsApp, and WeChat that communicates fluently in English, Hindi, Mandarin, and Arabic. It asks 8 qualifying questions: investment budget, timeline, current country, visa history, family composition, preferred project type, urgency, and contact preference. Serious investors ($800K+ capacity, clear timeline) are routed to Asif with a full brief. Early-stage inquiries enter a nurture sequence. Runs 24/7 — no staff required."),
    ("Client Journey\nAutomation", GOLD,
     "A milestone-based communication system that tracks every active client through the EB-5 lifecycle and sends proactive updates at each stage. I-526E filed → confirmation and expected timeline. I-526E approved → celebration + DS-260 action steps. DS-260 submitted → NVC processing update. Visa interview scheduled → prep guide. I-829 window approaching → automated reminder and document checklist. Congratulations on permanent residency. Every investor stays informed without a single manual email."),
    ("HNI Lead\nQualification Engine", PURPLE,
     "A real-time lead scoring and routing system that evaluates every inbound inquiry against 6 criteria: net worth signal, investment timeline, country eligibility, visa history, seriousness indicators, and source quality. High-scoring leads are pushed to Asif's calendar as priority. Mid-tier leads enter a 30-day education sequence. Low-tier leads get automated resources. FirstPathway's pipeline is always sorted by value — the team focuses on the $800K conversations, not the 'just curious' traffic."),
    ("Content & SEO\nOutreach Engine", ORANGE,
     "A content production and distribution system that publishes 3 LinkedIn posts per week, 2 blog articles per month targeting Indian and UAE EB-5 search terms, and a monthly webinar in Hindi or Mandarin. Topics include: EB-5 visa India 2025 timeline, how Indian HNIs qualify for US permanent residency, Ameyalli and Kaia project deep dives, EB-5 vs. other US residency routes. Email sequences in native languages warm the audience before Asif ever gets on a call."),
]
cw2 = Inches(2.95); gap2 = Inches(0.16)
for i, (title, color, desc) in enumerate(opps):
    l2 = Inches(0.55) + i * (cw2 + gap2)
    rect(s6, l2, Inches(1.8), cw2, Inches(0.06), color)
    rect(s6, l2, Inches(1.86), cw2, Inches(4.7), DARK_BLUE)
    tb(s6, title, l2 + Inches(0.12), Inches(1.98), cw2 - Inches(0.24), Inches(0.65),
       size=13, bold=True, color=color)
    tb(s6, desc, l2 + Inches(0.12), Inches(2.65), cw2 - Inches(0.24), Inches(3.65),
       size=11, color=LGREY)

rect(s6, Inches(0.55), Inches(6.7), Inches(12.23), Inches(0.45), DARK_BLUE)
tb(s6, "Asif and Wilka handle the legal expertise and investor relationships. Genos AI handles every digital touchpoint before the first call and throughout the 5-year journey.",
   Inches(0.7), Inches(6.75), Inches(12), Inches(0.35), size=11, italic=True, color=GOLD)


# ── SLIDE 7 — THE MULTILINGUAL IMPERATIVE ────────────────────────────────────
s7 = prs.slides.add_slide(BLANK)
bg(s7)
rect(s7, 0, 0, Inches(0.12), H, TEAL)
tb(s7, "THE MULTILINGUAL IMPERATIVE", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=TEAL)
tb(s7, "India, UAE, and China Are the EB-5 Market. None of Them Speak English First.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=26, bold=True, color=WHITE)
aline(s7, Inches(1.35), color=TEAL, w=Inches(2.5))

rect(s7, Inches(0.55), Inches(1.65), Inches(12.23), Inches(1.55), DARK_BLUE)
market_stats = [
    ("#1 Market", "India: largest new EB-5\napplicant nation (2024-25)"),
    ("500M+", "WhatsApp users in India —\ntheir preferred channel"),
    ("1.4B", "WeChat monthly active users —\nhow Chinese investors communicate"),
    ("Arabic", "Official language of UAE,\n#2 Gulf EB-5 source country"),
    ("$800K", "Investment per client — worth\nbuilding in their language"),
]
ssw = Inches(2.35)
for i, (num, label) in enumerate(market_stats):
    xl = Inches(0.7) + i * ssw
    tb(s7, num, xl, Inches(1.73), Inches(2.2), Inches(0.65),
       size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    tb(s7, label, xl, Inches(2.33), Inches(2.2), Inches(0.5),
       size=10, color=LGREY, align=PP_ALIGN.CENTER)

rect(s7, Inches(0.55), Inches(3.4), Inches(5.9), Inches(3.25), DARK_BLUE)
tb(s7, "HOW AN INDIAN HNI INVESTOR INQUIRES TODAY", Inches(0.7), Inches(3.47), Inches(5.6), Inches(0.35),
   size=11, bold=True, color=RED)
bullets(s7, [
    "Indian family researching EB-5 for children's US education pathway",
    "Finds firstpathway.com — reads the English page, no Hindi option",
    "Wants to ask a quick question — sees only a contact form",
    "Fills the form at 10pm IST (midnight Milwaukee time)",
    "Simultaneously WhatsApps 3 other EB-5 consultancies",
    "Competitor with Hindi WhatsApp bot responds in 4 minutes",
    "FirstPathway's reply arrives 2 business days later — investor already on a call with rival",
], Inches(0.7), Inches(3.85), Inches(5.6), Inches(2.55), size=11, color=LGREY, dot=RED)

rect(s7, Inches(6.65), Inches(3.4), Inches(6.13), Inches(3.25), DARK_BLUE)
tb(s7, "WITH GENOS AI MULTILINGUAL AGENT", Inches(6.8), Inches(3.47), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=GREEN)
bullets(s7, [
    "Same Indian family lands on firstpathway.com at 10pm IST",
    "Hindi chatbot greets them: 'Namaste! EB-5 ke baare mein jaankari chahiye?'",
    "Agent runs 8-question qualification in Hindi — 5 minutes",
    "Family qualifies: $2M net worth, 2-year timeline, children US-education goal",
    "Agent sends: project overview in Hindi + Asif's calendar link",
    "Asif wakes up to: qualified brief, contact details, preferred call slot",
    "FirstPathway books the consultation before the competitor has been Googled",
], Inches(6.8), Inches(3.85), Inches(5.8), Inches(2.55), size=11, color=LGREY, dot=GREEN)

rect(s7, Inches(0.55), Inches(6.8), Inches(12.23), Inches(0.38), DARK_BLUE)
tb(s7, "The first credible response in the investor's language wins the consultation. The consultation wins the $800K+ client.",
   Inches(0.7), Inches(6.85), Inches(12), Inches(0.28), size=11, italic=True, color=TEAL)


# ── SLIDE 8 — IMPLEMENTATION PLAN ────────────────────────────────────────────
s8 = prs.slides.add_slide(BLANK)
bg(s8)
rect(s8, 0, 0, Inches(0.12), H, GREEN)
tb(s8, "IMPLEMENTATION PLAN", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s8, "Multilingual First. Client Journey Second. HNI Outreach Third.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=28, bold=True, color=WHITE)
aline(s8, Inches(1.35), color=GREEN)

phases = [
    ("PHASE 1", "Weeks 1-2", TEAL, [
        ("Multilingual Website Chatbot", "AI agent deployed on firstpathway.com responding in English, Hindi, Mandarin, and Arabic. 8-question investor qualification flow. Serious leads routed to Asif's calendar automatically."),
        ("WhatsApp Business AI Agent", "WhatsApp AI bot for Indian and Arabic-speaking investors. Qualifies leads, answers FAQs about EB-5 process and Ameyalli/Kaia projects, books consultations in under 60 seconds."),
        ("WeChat AI Agent Setup", "WeChat Official Account AI bot for Chinese investor audience. Matches the qualification flow in Mandarin. Existing Weibo account connected for content amplification."),
    ]),
    ("PHASE 2", "Weeks 3-6", GOLD, [
        ("Client Journey Automation", "Milestone email and WhatsApp sequences for all active clients: I-526E filing, approval, DS-260, visa processing, I-829 approach. No investor waits in silence for months."),
        ("HNI Lead Scoring Engine", "Real-time scoring of every inbound inquiry. High-value leads ($800K+ signals, clear timeline) pushed to Asif as priority. Mid-tier enters 30-day nurture. CRM integration."),
        ("EB-5 ROI / Benefit Calculator", "Interactive tool on website: investor enters country, family size, investment goal — gets a personalised US residency benefit summary. Lead captured before results shown."),
    ]),
    ("PHASE 3", "Months 2-3", GREEN, [
        ("SEO Content Engine", "12 blog articles targeting: 'EB-5 visa India 2025', 'US green card for investors', 'EB-5 Ameyalli Utah', 'how to get US permanent residency from UAE'. Published monthly."),
        ("LinkedIn Thought Leadership", "3 posts per week on EB-5 updates, project news, investor success stories, and US residency tips. Targeting Indian and UAE HNI audience and referring attorneys/wealth managers."),
        ("Hindi/Mandarin Webinar Series", "Monthly investor webinar in Hindi or Mandarin with Asif. Automated registration, reminder, replay email. Positions FirstPathway as the go-to firm for non-English EB-5 investors."),
    ]),
]
pw = Inches(3.9); gap4 = Inches(0.18)
for i, (phase, timing, color, items) in enumerate(phases):
    l3 = Inches(0.55) + i * (pw + gap4)
    rect(s8, l3, Inches(1.8), pw, Inches(0.52), color)
    tb(s8, phase, l3 + Inches(0.12), Inches(1.84), pw * 0.5 - Inches(0.12), Inches(0.42),
       size=13, bold=True, color=NAVY)
    tb(s8, timing, l3 + pw * 0.45, Inches(1.88), pw * 0.55 - Inches(0.15), Inches(0.38),
       size=11, color=NAVY, align=PP_ALIGN.RIGHT)
    rect(s8, l3, Inches(2.32), pw, Inches(4.6), DARK_BLUE)
    y_item = Inches(2.42)
    for title, desc in items:
        rect(s8, l3 + Inches(0.12), y_item + Inches(0.06), Inches(0.05), Inches(0.52), color)
        tb(s8, title, l3 + Inches(0.25), y_item, pw - Inches(0.38), Inches(0.32),
           size=12, bold=True, color=color)
        tb(s8, desc, l3 + Inches(0.25), y_item + Inches(0.3), pw - Inches(0.38), Inches(0.45),
           size=11, color=LGREY)
        y_item += Inches(0.95)


# ── SLIDE 9 — BEFORE vs AFTER ────────────────────────────────────────────────
s9 = prs.slides.add_slide(BLANK)
bg(s9)
rect(s9, 0, 0, Inches(0.12), H, ACCENT)
tb(s9, "BEFORE vs AFTER", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s9, "What Actually Changes for FirstPathway in 90 Days",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=28, bold=True, color=WHITE)
aline(s9, Inches(1.35))

comparisons = [
    ("Indian HNI Investor Inquiry",
     "Investor finds firstpathway.com at 10pm IST. Sees a contact form. Fills it. Simultaneously WhatsApps 3 other EB-5 firms. Competitor with Hindi bot responds in 4 minutes. FirstPathway's reply arrives 2 days later — investor is already on a competitor's call.",
     "Hindi WhatsApp AI agent responds within 60 seconds at any hour. 8-question qualification in Hindi. Asif receives a full investor brief by morning. Consultation booked before the investor has heard back from anyone else."),
    ("Chinese Investor via WeChat",
     "Chinese investor finds FirstPathway on Weibo. Clicks to website — English only. No WeChat bot. Sends a message on WeChat, waits days for a reply. Historical #1 EB-5 market, now effectively unreachable.",
     "WeChat AI agent in Mandarin. Qualifies investor: budget, timeline, family size. Answers Ameyalli and Kaia questions. Routes serious investor to Asif with full Mandarin brief. WeChat becomes an active lead channel, not a dormant account."),
    ("Active Client Communication (3-5 Year Journey)",
     "Client files I-526E. Waits. Months pass with no proactive update. Calls the office asking for status. Team provides manual update. Client remains anxious throughout a 3-5 year process with no structured communication.",
     "Automated milestone messages: I-526E confirmed, expected timeline, DS-260 due date reminder, visa interview prep, I-829 filing approach. Every client knows where they stand at every stage — without a single manual email from the team."),
    ("HNI Lead Qualification",
     "Every inbound inquiry is treated the same — Asif manually reviews each email and responds personally. Early-stage 'just curious' inquiries consume the same time as $2M-capacity investors ready to commit.",
     "AI scores every inquiry in real time. $800K+ investors with clear timelines go straight to Asif's calendar as priority. Mid-tier leads enter a 30-day education sequence. Team time focused entirely on conversations that convert."),
    ("Digital Visibility",
     "Typing 'EB-5 visa India 2025' or 'US green card investor program' on Google returns no FirstPathway content. Indian and UAE investors searching online never encounter the firm. Traffic is entirely dependent on referrals.",
     "12 SEO articles published over 90 days targeting exact India, UAE, China EB-5 search terms. LinkedIn posts 3x/week reaching HNI audience. Monthly Hindi/Mandarin webinar. FirstPathway becomes the visible expert for non-English EB-5 investors."),
]

lw = Inches(5.4)
rect(s9, Inches(0.55), Inches(1.65), lw, Inches(0.35), RED)
tb(s9, "TODAY", Inches(0.65), Inches(1.68), lw - Inches(0.2), Inches(0.28),
   size=11, bold=True, color=WHITE)
rect(s9, Inches(6.45), Inches(1.65), lw + Inches(0.43), Inches(0.35), GREEN)
tb(s9, "AFTER 90 DAYS", Inches(6.55), Inches(1.68), lw, Inches(0.28),
   size=11, bold=True, color=WHITE)

for i, (topic, before, after) in enumerate(comparisons):
    t = Inches(2.1) + i * Inches(1.02)
    rect(s9, Inches(0.55), t, Inches(12.23), Inches(0.22), DARK_BLUE)
    tb(s9, topic.upper(), Inches(0.65), t + Inches(0.02), Inches(12), Inches(0.2),
       size=9, bold=True, color=MGREY)
    tb(s9, before, Inches(0.65), t + Inches(0.24), lw - Inches(0.2), Inches(0.66),
       size=11, color=LGREY)
    tb(s9, after, Inches(6.55), t + Inches(0.24), lw, Inches(0.66),
       size=11, color=WHITE)


# ── SLIDE 10 — NEXT STEPS ────────────────────────────────────────────────────
s10 = prs.slides.add_slide(BLANK)
bg(s10)
rect(s10, 0, 0, Inches(0.12), H, TEAL)
rect(s10, Inches(7.5), 0, Inches(5.83), H, DARK_BLUE)
tb(s10, "NEXT STEPS", Inches(0.55), Inches(0.3), Inches(6.5), Inches(0.5),
   size=11, bold=True, color=TEAL)
tb(s10, "Let's Talk for\n15 Minutes.",
   Inches(0.55), Inches(0.8), Inches(6.8), Inches(2.0), size=48, bold=True, color=WHITE)
aline(s10, Inches(2.8), color=TEAL, w=Inches(2.5))
tb(s10, "No pitch needed. Call or WhatsApp.\nWe'll start with the multilingual investor agent — live in 2 weeks.",
   Inches(0.55), Inches(3.05), Inches(6.5), Inches(0.8), size=14, color=LGREY)

steps = [
    "Call or WhatsApp Rohan: +91 638-714-2699",
    "Week 1-2: Hindi + Mandarin + Arabic investor agent live on website + WhatsApp + WeChat",
    "Week 3-6: Client journey automation + HNI lead scoring + ROI calculator",
    "Month 2-3: SEO content engine + LinkedIn posts + Hindi/Mandarin webinar series",
]
bullets(s10, steps, Inches(0.55), Inches(3.95), Inches(6.5), Inches(2.0), size=14, color=WHITE, dot=TEAL)

rect(s10, 0, Inches(6.5), Inches(7.5), Inches(1.0), DARK_BLUE)
tb(s10, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699  |  www.genosai.tech",
   Inches(0.55), Inches(6.6), Inches(6.8), Inches(0.5), size=11, color=MGREY)

tb(s10, "THREE GAPS.\nTHREE FIXES.", Inches(7.85), Inches(0.6), Inches(5.0), Inches(1.1),
   size=22, bold=True, color=TEAL)
aline(s10, Inches(1.65), color=TEAL)

summary_items = [
    ("No Hindi Bot", "Multilingual AI agent: Hindi + Mandarin + Arabic, 24/7"),
    ("No WeChat AI", "WeChat bot qualifying Chinese investors in Mandarin"),
    ("Manual Journey", "Automated 3-5 yr milestone comms: I-526E to I-829"),
    ("No Lead Score", "HNI scoring: $800K+ investors routed to Asif instantly"),
    ("No SEO", "12 articles: EB-5 India, UAE, China search terms"),
    ("No Webinars", "Hindi/Mandarin webinar series with auto registration"),
]
for i, (label, val) in enumerate(summary_items):
    t2 = Inches(1.85) + i * Inches(0.88)
    rect(s10, Inches(7.85), t2, Inches(1.5), Inches(0.72), NAVY)
    tb(s10, label, Inches(7.9), t2 + Inches(0.05), Inches(1.4), Inches(0.32),
       size=10, bold=True, color=TEAL)
    tb(s10, val, Inches(9.5), t2 + Inches(0.05), Inches(3.6), Inches(0.65),
       size=11, color=LGREY)


prs.save(ROOT / "output" / "decks" / "FirstPathway_Audit_GenosAI.pptx")
print("Saved: FirstPathway_Audit_GenosAI.pptx")
