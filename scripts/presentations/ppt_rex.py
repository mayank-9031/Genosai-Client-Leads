import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.parent / "utils"))
from ppt_theme import (
    NAVY, DARK_BLUE, ACCENT, GOLD, RED, GREEN, PURPLE, ORANGE, TEAL, CYAN,
    WHITE, LGREY, MGREY, W, H,
    bg, rect, tb, aline, bullets, tag,
    slide_header, slide_footer, score_badge,
    stat_cards, top_bar_card, phase_card, before_after, genos_logo_block,
)

ROOT = Path(__file__).parent.parent.parent
prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ══ SLIDE 1 — COVER ══
s1 = prs.slides.add_slide(BLANK)
bg(s1)
rect(s1, 0, 0, Inches(0.12), H, ACCENT)
rect(s1, Inches(8.5), 0, Inches(4.83), Inches(2.2), DARK_BLUE)
tb(s1, "GENOS AI", Inches(9.2), Inches(0.4), Inches(3.5), Inches(0.8),
   size=28, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
tb(s1, "www.genosai.tech", Inches(9.2), Inches(1.05), Inches(3.5), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
tb(s1, "REX COMPANIES / REX.COM", Inches(0.55), Inches(1.6),
   Inches(9), Inches(0.6), size=13, bold=False, color=ACCENT)
tb(s1, "Website & AI Growth\nAudit Report", Inches(0.55), Inches(2.1),
   Inches(9), Inches(1.8), size=48, bold=True, color=WHITE)
aline(s1, Inches(3.85), color=GOLD, w=Inches(3))
tb(s1, "Prepared exclusively for Peter Rex, Founder & CEO",
   Inches(0.55), Inches(4.15), Inches(9), Inches(0.5), size=14, color=LGREY)
rect(s1, 0, Inches(6.5), W, Inches(1.0), DARK_BLUE)
tb(s1, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699",
   Inches(0.55), Inches(6.6), Inches(9), Inches(0.5), size=12, color=MGREY)
tb(s1, "May 2026  |  CONFIDENTIAL", Inches(9.5), Inches(6.6), Inches(3.3), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
rect(s1, Inches(10.5), Inches(3.3), Inches(2.3), Inches(2.5), DARK_BLUE)
tb(s1, "AUDIT SCORE", Inches(10.55), Inches(3.4), Inches(2.2), Inches(0.5),
   size=11, bold=True, color=MGREY, align=PP_ALIGN.CENTER)
tb(s1, "5", Inches(10.55), Inches(3.75), Inches(2.2), Inches(1.1),
   size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
tb(s1, "/ 10", Inches(10.55), Inches(4.7), Inches(2.2), Inches(0.4),
   size=18, color=LGREY, align=PP_ALIGN.CENTER)
tag(s1, "HIGH PRIORITY", Inches(10.75), Inches(5.1), fill=RED, size=10)


# ══ SLIDE 2 — AT A GLANCE ══
s2 = prs.slides.add_slide(BLANK)
bg(s2)
rect(s2, 0, 0, Inches(0.12), H, ACCENT)
tb(s2, "AT A GLANCE", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s2, "Peter Rex — Who He Is",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s2, Inches(1.35))

stats = [
    ("19,000+", "Apartments Managed"),
    ("$5B+", "Assets Under Management"),
    ("10", "AI Tech Companies"),
    ("9", "Patents Pending"),
    ("AI R&D\n2017", "Early Mover in PropTech"),
]
cw = Inches(2.3); ch = Inches(2.2); gap = Inches(0.18)
for i, (num, label) in enumerate(stats):
    l = Inches(0.55) + i * (cw + gap)
    rect(s2, l, Inches(1.8), cw, ch, DARK_BLUE)
    tb(s2, num, l + Inches(0.1), Inches(1.95), cw - Inches(0.2), Inches(1.1),
       size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s2, label, l + Inches(0.1), Inches(3.05), cw - Inches(0.2), Inches(0.7),
       size=12, color=LGREY, align=PP_ALIGN.CENTER)

rect(s2, Inches(0.55), Inches(4.2), Inches(12.23), Inches(0.6), DARK_BLUE)
tb(s2, "Georgetown BA (Philosophy/Govt)  ·  Harvard Law JD  ·  CPA  ·  Austin TX  ·  WSJ Op-Ed 2020  ·  600+ Teammates  ·  19 Metros  ·  6 States",
   Inches(0.65), Inches(4.28), Inches(12), Inches(0.45), size=12, color=LGREY, align=PP_ALIGN.CENTER)

tb(s2, "REX TECH COMPANIES", Inches(0.55), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=ACCENT)
bullets(s2, ["OwnProp (fractional tokenized real estate)",
             "GetDone (maintenance AI/ML)",
             "PayUp, SmartFuse, PropData, InsurePro + 4 more"],
        Inches(0.55), Inches(5.35), Inches(5.8), Inches(1.5), size=13, color=LGREY)

tb(s2, "REAL ESTATE PORTFOLIO", Inches(6.8), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=ACCENT)
bullets(s2, ["19,000 apartments — Sun Belt + Mountain West",
             "Hotel Ella + Arena Hall (Austin)",
             "Office space — Dallas and Austin"],
        Inches(6.8), Inches(5.35), Inches(5.8), Inches(1.5), size=13, color=LGREY)


# ══ SLIDE 3 — WHAT'S WORKING ══
s3 = prs.slides.add_slide(BLANK)
bg(s3)
rect(s3, 0, 0, Inches(0.12), H, GREEN)
tb(s3, "WHAT'S WORKING", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s3, "Strengths to Build On",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s3, Inches(1.35), color=GREEN)

strengths = [
    ("Genuine AI Infrastructure Since 2017",
     "Peter started AI R&D in 2017 - before most real estate firms had even heard of machine learning. Nine pending patents and NLP/ML solutions powering active products means Rex isn't talking about AI, they're shipping it."),
    ("10 Revenue-Generating Tech Companies",
     "OwnProp, GetDone, PayUp, SmartFuse, PropData, InsurePro - all mobile-first, API-driven, AI-leveraged, all generating revenue. This is a tech holding company, not a landlord with a website."),
    ("Vertically Integrated at Scale",
     "Investment, asset management, property management, construction, and tech under one roof across 19 metros. Most operators at this scale have fragmented their stack. Rex owns it end to end."),
    ("OwnProp First-Mover in Fractional Tokenized RE",
     "Fractional real estate is still early. OwnProp has a genuine first-mover advantage in a category where the incumbents (Fundrise, RealtyMogul) are not tech-native. The product infrastructure is ahead of the marketing funnel."),
    ("Premium Brand Design - rex.com",
     "The website is clean, minimalist, and premium. It doesn't look like most real estate companies. That's a deliberate identity choice that works. The gap is functional depth, not visual design."),
]
for i, (title, desc) in enumerate(strengths):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s3, Inches(0.55), t, Inches(0.06), Inches(0.7), GREEN)
    tb(s3, title, Inches(0.75), t, Inches(11.8), Inches(0.38), size=14, bold=True, color=WHITE)
    tb(s3, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# ══ SLIDE 4 — WHERE LEADS ARE LEAKING ══
s4 = prs.slides.add_slide(BLANK)
bg(s4)
rect(s4, 0, 0, Inches(0.12), H, RED)
tb(s4, "REVENUE LAYER GAPS", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=RED)
tb(s4, "Tech Depth on Product Side. Not on the Revenue Side.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=30, bold=True, color=WHITE)
aline(s4, Inches(1.35), color=RED)

issues = [
    ("INVESTMENT PAGE IS A BRAND DECK, NOT A FUNNEL", "HIGH",
     "Rex.com/real-estate has no performance data, no fund structure, no team bios, no deal case studies. An accredited investor researching whether to wire money finds nothing concrete to act on. Greystar, Aimco, and Starwood all publish track records. Rex doesn't."),
    ("OWNPROP LACKS AN INVESTOR ACQUISITION FUNNEL", "HIGH",
     "Fractional real estate has a longer consideration cycle than most products - investors typically research for weeks before committing. OwnProp needs education sequences, re-engagement flows, and content that moves prospects from awareness to first investment."),
    ("NO LEASING AUTOMATION AT PORTFOLIO SCALE", "HIGH",
     "19,000 apartments across 19 metros means thousands of prospect touches monthly. Without automated follow-up, tour scheduling, post-tour nurtures, and renewal sequences, leads go cold between Friday PM and Monday AM. That's vacancy cost at scale."),
    ("NO LIVE CHAT ON REX.COM", "MEDIUM",
     "No chat widget anywhere on rex.com. Leasing prospects and potential investors who visit in the evening or on weekends have no real-time engagement path. GetDone handles maintenance. The front door doesn't have the same coverage."),
    ("RESIDENT RETENTION NOT VISIBLE IN DIGITAL LAYER", "MEDIUM",
     "With 19,000 units, even a 1% improvement in lease renewal rates is significant. Automated renewal sequences starting 90 days out, personalized by unit type, tenure, and market conditions, are table stakes at this portfolio scale."),
]
for i, (title, pri, desc) in enumerate(issues):
    t = Inches(1.7) + i * Inches(1.05)
    pc = RED if pri == "HIGH" else GOLD
    rect(s4, Inches(0.55), t, Inches(0.06), Inches(0.7), pc)
    tb(s4, title, Inches(0.75), t, Inches(9.2), Inches(0.38), size=13, bold=True, color=WHITE)
    tag(s4, pri, Inches(10.1), t + Inches(0.05), fill=pc, size=9)
    tb(s4, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# ══ SLIDE 5 — WEBSITE DESIGN AUDIT ══
s5 = prs.slides.add_slide(BLANK)
bg(s5)
rect(s5, 0, 0, Inches(0.12), H, PURPLE)
tb(s5, "WEBSITE DESIGN AUDIT", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=PURPLE)
tb(s5, "Premium Brand. Investor Funnel Missing.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.8), size=34, bold=True, color=WHITE)
aline(s5, Inches(1.45), color=PURPLE, w=Inches(1.5))

rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(3.6), DARK_BLUE)
rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(0.42), RED)
tb(s5, "REX.COM TODAY", Inches(0.65), Inches(1.92), Inches(3.6), Inches(0.38),
   size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Brand deck, not investor conversion funnel",
    "No performance data or returns history",
    "No fund structure or investment minimums",
    "No team bios beyond founder",
    "No deal case studies or portfolio highlights",
    "No live chat on any page",
    "OwnProp lacks investor education content",
    "No testimonials from LPs or residents",
], Inches(0.7), Inches(2.42), Inches(3.6), Inches(2.9), size=11, color=LGREY, dot=RED)

tb(s5, "->", Inches(4.45), Inches(3.6), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(3.6), DARK_BLUE)
rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(0.42), GOLD)
tb(s5, "TOP MULTIFAMILY OPERATORS DO THIS", Inches(5.1), Inches(1.92), Inches(3.4), Inches(0.38),
   size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Greystar, Aimco publish full track record pages",
    "Starwood, Blackstone show fund structure + team",
    "LP testimonials and deal spotlights front and center",
    "OwnProp competitors (Fundrise) have deep education funnels",
    "Property pages with virtual tours and automated scheduling",
    "AI chat for investor and leasing inquiries 24/7",
    "Resident portal login prominent for current tenants",
    "Performance dashboards visible to prospective LPs",
], Inches(5.1), Inches(2.42), Inches(3.4), Inches(2.9), size=11, color=LGREY, dot=GOLD)

tb(s5, "->", Inches(8.68), Inches(3.6), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(3.6), DARK_BLUE)
rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(0.42), GREEN)
tb(s5, "WHAT GENOS AI BUILDS", Inches(9.35), Inches(1.92), Inches(3.4), Inches(0.38),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Investment pages: returns data, fund structure, team",
    "OwnProp: full investor education funnel + re-engagement",
    "Deal case studies: deal-by-deal portfolio highlights",
    "LP testimonials section with real partner quotes",
    "AI chat trained on Rex investment criteria",
    "Property pages: virtual tours, automated tour booking",
    "Resident portal integration and retention sequences",
    "Peter's thought leadership auto-syndicating to all channels",
], Inches(9.35), Inches(2.42), Inches(3.4), Inches(2.9), size=11, color=LGREY, dot=GREEN)

rect(s5, Inches(0.55), Inches(5.72), Inches(12.23), Inches(1.28), DARK_BLUE)
tb(s5, "WHY IT MATTERS:", Inches(0.7), Inches(5.78), Inches(2.2), Inches(0.38),
   size=11, bold=True, color=PURPLE)
stats2 = [
    ("94%", "of first impressions\nare design-related"),
    ("67%", "of institutional LPs screen\nfirms via website before call"),
    ("3 sec", "before a visitor decides\nto engage or leave"),
    ("200%+", "lift in conversions with\npersonalized investor funnel"),
]
sw = Inches(2.7); sl = Inches(2.9)
for i, (num, label) in enumerate(stats2):
    xl = sl + i * sw
    tb(s5, num, xl, Inches(5.75), Inches(2.5), Inches(0.55),
       size=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s5, label, xl, Inches(6.22), Inches(2.5), Inches(0.7),
       size=10, color=LGREY, align=PP_ALIGN.CENTER)


# ══ SLIDE 6 — AI AUTOMATION OPPORTUNITIES ══
s6 = prs.slides.add_slide(BLANK)
bg(s6)
rect(s6, 0, 0, Inches(0.12), H, ACCENT)
tb(s6, "AI AUTOMATION OPPORTUNITIES", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s6, "The Revenue Layer. Not the Product Layer.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=34, bold=True, color=WHITE)
aline(s6, Inches(1.35), color=ACCENT)
tb(s6, "Rex already builds AI products. This is about the commercial automation layer on top of the portfolio.",
   Inches(0.55), Inches(1.3), Inches(12), Inches(0.38), size=12, italic=True, color=MGREY)

opps = [
    ("Leasing Funnel\nAutomation", GOLD,
     "19,000 units across 19 metros = thousands of prospect touches monthly. Automated follow-up, tour scheduling and confirmation, post-tour nudges, and application reminders. Leads that go cold over weekends get recovered automatically."),
    ("OwnProp Investor\nAcquisition", ACCENT,
     "Fractional real estate has a 3-6 week consideration cycle. Prospect education sequences, comparison content vs. competitors, re-engagement for people who visited OwnProp but didn't invest. Build the funnel the product deserves."),
    ("Resident Retention\nSequences", GREEN,
     "90-day pre-renewal outreach, personalized by unit type, tenure, and local market conditions. At 19,000 units, a 1% improvement in renewal rates is hundreds of avoided vacancy cycles worth tens of thousands each."),
    ("Investment Page\nInquiry Funnel", PURPLE,
     "Returns track record, fund structure, deal spotlights, LP testimonials. Convert brand impressions into accredited investor inquiries. Automated follow-up for contact form submissions and investor meeting scheduling."),
    ("Peter's Thought\nLeadership Pipeline", RED,
     "One X post or article triggers an automated pipeline: email to investor list, LinkedIn native post, press release draft, newsletter section. Peter creates once. AI distributes everywhere. Builds LP trust at scale."),
]
cw2 = Inches(2.35); gap2 = Inches(0.14)
for i, (title, color, desc) in enumerate(opps):
    l2 = Inches(0.55) + i * (cw2 + gap2)
    rect(s6, l2, Inches(1.8), cw2, Inches(0.06), color)
    rect(s6, l2, Inches(1.86), cw2, Inches(4.6), DARK_BLUE)
    tb(s6, title, l2 + Inches(0.12), Inches(1.98), cw2 - Inches(0.24), Inches(0.65),
       size=13, bold=True, color=color)
    tb(s6, desc, l2 + Inches(0.12), Inches(2.65), cw2 - Inches(0.24), Inches(3.6),
       size=11, color=LGREY)

rect(s6, Inches(0.55), Inches(6.7), Inches(12.23), Inches(0.45), DARK_BLUE)
tb(s6, "\"Building product for real estate tech is very hard. There are multiple stakeholders. We love hard things at Rex.\" — Peter Rex, X  |  The revenue automation layer is the other hard thing.",
   Inches(0.7), Inches(6.75), Inches(12), Inches(0.35), size=11, italic=True, color=GOLD)


# ══ SLIDE 7 — WHAT THIS IS COSTING YOU ══
s7 = prs.slides.add_slide(BLANK)
bg(s7)
rect(s7, 0, 0, Inches(0.12), H, RED)
tb(s7, "WHAT THIS IS COSTING YOU", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=RED)
tb(s7, "Nine Patents. Revenue Layer Running Manual.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s7, Inches(1.35), color=RED)

costs = [
    ("No investor funnel = capital sitting on the table",
     "An accredited investor who researches rex.com/real-estate and finds no performance data, no team, and no fund structure doesn't move forward. That's not a product problem. It's a funnel problem. Rex's actual track record would convert many of them."),
    ("OwnProp with no nurture = first-mover advantage being given away",
     "Fundrise and RealtyMogul have deep email nurture sequences, investor education content, and re-engagement flows. OwnProp has the better product and the first-mover position. Without the marketing funnel, the incumbent wins on attention."),
    ("19,000 units without leasing automation = weekend vacancy cycles",
     "Every prospect who tours on a Thursday or Friday and doesn't hear back until Monday is a potential vacancy. Multiply that across 19 metros and the cost is not small. Automated follow-up within 2 hours of a tour request is standard at Greystar. Rex is not there yet."),
    ("Resident retention running on inertia",
     "Renewals that happen because the tenant didn't get around to moving are not the same as renewals earned through proactive communication. Personalized retention sequences reduce vacancy and give Rex data on which residents are flight risks."),
    ("Peter's thought leadership not being amplified",
     "Peter posts on X, writes press pieces, goes on podcasts. None of it is systematically converted into investor outreach or LP relationship building. One syndication pipeline would turn every public post into an investor touchpoint."),
]
for i, (title, desc) in enumerate(costs):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s7, Inches(0.55), t, Inches(0.06), Inches(0.7), RED)
    tb(s7, title, Inches(0.75), t, Inches(11.8), Inches(0.38), size=14, bold=True, color=WHITE)
    tb(s7, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# ══ SLIDE 8 — WHAT WE'D BUILD ══
s8 = prs.slides.add_slide(BLANK)
bg(s8)
rect(s8, 0, 0, Inches(0.12), H, GREEN)
tb(s8, "WHAT WE'D BUILD", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s8, "Genos AI for REX",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s8, Inches(1.35), color=GREEN)

builds = [
    ("Rex Investment Pages Rebuild",
     "Returns track record, fund structure, deal case studies, team bios. Built for accredited investor conversion, not brand admiration. A prospect who visits this page should have enough to request a meeting without emailing for a deck."),
    ("OwnProp Investor Acquisition Funnel",
     "Education sequences for fractional RE first-timers, comparison content for prospects evaluating Fundrise or RealtyMogul, re-engagement flows for people who visited but didn't invest. Turns OwnProp's product lead into a conversion lead."),
    ("Leasing Automation Layer",
     "Automated prospect follow-up within 2 hours of inquiry, tour scheduling and confirmation, post-tour nurture sequences, application reminders. Deployed across all 19 metros. Weekend coverage built in - no lead goes cold between Friday PM and Monday AM."),
    ("Resident Retention Sequences",
     "90-day pre-renewal outreach by unit type, tenure, and local market. Automated satisfaction check-ins at 6 and 12 months. Flight-risk flagging for residents showing early exit signals. At 19,000 units, this has measurable P&L impact."),
    ("Peter Thought Leadership Pipeline",
     "One post on X or one press piece triggers automated distribution: investor email list, LinkedIn native, newsletter section, press release. Peter creates. AI handles reach. Builds LP trust at scale without adding marketing headcount."),
]
for i, (title, desc) in enumerate(builds):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s8, Inches(0.55), t, Inches(0.06), Inches(0.7), GREEN)
    tb(s8, title, Inches(0.75), t, Inches(11.8), Inches(0.38), size=14, bold=True, color=WHITE)
    tb(s8, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# ══ SLIDE 9 — 90-DAY OUTCOMES ══
s9 = prs.slides.add_slide(BLANK)
bg(s9)
rect(s9, 0, 0, Inches(0.12), H, ACCENT)
tb(s9, "90-DAY OUTCOMES", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s9, "What Changes in 90 Days",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s9, Inches(1.35))

phases = [
    ("Month 1", GOLD, [
        "AI chat live on rex.com for investor and leasing inquiries",
        "Leasing automation deployed - follow-up, tour scheduling, post-tour nudge sequences",
        "OwnProp nurture sequences active - education content flowing to prospects",
        "Investment pages rebuild begins - returns data, team, fund structure in design",
    ]),
    ("Month 2", ACCENT, [
        "Rex investment pages live - track record, team bios, deal spotlights visible to LPs",
        "Resident retention sequences running across full portfolio",
        "OwnProp first campaign performance data: open rates, click-to-invest conversion",
        "Peter thought leadership pipeline active - every post now auto-syndicating",
    ]),
    ("Month 3", GREEN, [
        "Leasing lead recovery from weekend gaps measurable in per-metro data",
        "OwnProp investor conversion rate up vs. organic baseline",
        "Renewal rate improvement visible in first cohorts of retention sequences",
        "LP inquiry volume from investment pages vs. prior period tracked",
    ]),
]
pw = Inches(3.9); gap3 = Inches(0.18)
for i, (month, color, pts) in enumerate(phases):
    l3 = Inches(0.55) + i * (pw + gap3)
    rect(s9, l3, Inches(1.8), pw, Inches(0.42), color)
    tb(s9, month.upper(), l3 + Inches(0.1), Inches(1.84), pw - Inches(0.2), Inches(0.36),
       size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    rect(s9, l3, Inches(2.22), pw, Inches(4.7), DARK_BLUE)
    bullets(s9, pts, l3 + Inches(0.15), Inches(2.35), pw - Inches(0.3), Inches(4.4),
            size=12, color=LGREY, dot=color)


# ══ SLIDE 10 — NEXT STEPS ══
s10 = prs.slides.add_slide(BLANK)
bg(s10)
rect(s10, 0, 0, Inches(0.12), H, GOLD)
rect(s10, Inches(7.5), 0, Inches(5.83), H, DARK_BLUE)
tb(s10, "NEXT STEPS", Inches(0.55), Inches(0.3), Inches(6.5), Inches(0.5),
   size=11, bold=True, color=GOLD)
tb(s10, "Let's Talk for\n15 Minutes.",
   Inches(0.55), Inches(0.8), Inches(6.8), Inches(2.0), size=48, bold=True, color=WHITE)
aline(s10, Inches(2.8), color=GOLD, w=Inches(2.5))
tb(s10, "No pitch deck needed. Reply to the email or book directly.\nWe'll compare notes on the revenue-layer AI stuff.",
   Inches(0.55), Inches(3.05), Inches(6.5), Inches(0.8), size=14, color=LGREY)

steps = [
    "Reply to this email - we'll set up a time",
    "Book directly: www.genosai.tech/call",
    "We walk through the audit findings in 15 minutes",
    "You decide what, if anything, makes sense",
]
bullets(s10, steps, Inches(0.55), Inches(3.95), Inches(6.5), Inches(2.0), size=14, color=WHITE, dot=GOLD)

rect(s10, 0, Inches(6.5), Inches(7.5), Inches(1.0), DARK_BLUE)
tb(s10, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699  |  www.genosai.tech",
   Inches(0.55), Inches(6.6), Inches(6.8), Inches(0.5), size=12, color=MGREY)

tb(s10, "WHAT WE BUILD\nFOR REX", Inches(7.85), Inches(0.6), Inches(5.0), Inches(0.9),
   size=20, bold=True, color=GOLD)
aline(s10, Inches(1.45), color=GOLD)

summary_items = [
    ("Investment Pages", "Track record, team, fund structure - built to convert LPs"),
    ("OwnProp Funnel", "Education sequences + re-engagement for fractional RE investors"),
    ("Leasing Automation", "Follow-up, tour scheduling, post-tour nurture - 19 metros"),
    ("Resident Retention", "90-day pre-renewal sequences, personalized by tenure/unit"),
    ("AI Chat", "24/7 investor and leasing inquiry widget on rex.com"),
    ("Thought Leadership", "Peter's posts auto-syndicated to email, LinkedIn, press"),
]
for i, (label, val) in enumerate(summary_items):
    t2 = Inches(1.6) + i * Inches(0.88)
    rect(s10, Inches(7.85), t2, Inches(1.5), Inches(0.72), NAVY)
    tb(s10, label, Inches(7.9), t2 + Inches(0.05), Inches(1.4), Inches(0.32),
       size=10, bold=True, color=GOLD)
    tb(s10, val, Inches(9.5), t2 + Inches(0.05), Inches(3.6), Inches(0.65),
       size=11, color=LGREY)


prs.save(ROOT / "output" / "decks" / "REX_Audit_GenosAI.pptx")
print("Saved: REX_Audit_GenosAI.pptx")
