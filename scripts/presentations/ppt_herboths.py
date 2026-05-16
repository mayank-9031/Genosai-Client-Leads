import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.parent / "utils"))
from ppt_theme import (
    NAVY, DARK_BLUE, ACCENT, GOLD, RED, GREEN, PURPLE, ORANGE, TEAL, CYAN,
    WHITE, LGREY, MGREY, W, H,
    bg, rect, tb, aline, bullets, tag,
    slide_header, slide_footer, score_badge,
    stat_cards, top_bar_card, phase_card, before_after, genos_logo_block,
)

ROOT = Path(__file__).parent.parent.parent
prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]

# Client-specific colors
DGREY     = RGBColor(0x3A, 0x4A, 0x5A)
AMBER     = RGBColor(0xFF, 0xA5, 0x00)


def score_bar(slide, label, score, x, y, w):
    """Horizontal score bar: label + filled bar out of 10 + score number."""
    bar_h = Inches(0.32)
    fill_color = RED if score <= 3 else (AMBER if score <= 6 else GREEN)
    # background track
    rect(slide, x, y + Inches(0.42), w, Inches(0.18), DARK_BLUE)
    # filled portion
    fill_w = w * (score / 10)
    rect(slide, x, y + Inches(0.42), fill_w, Inches(0.18), fill_color)
    # label
    tb(slide, label, x, y, w - Inches(0.55), bar_h,
       size=11, color=LGREY)
    # score number
    tb(slide, str(score), x + w - Inches(0.5), y, Inches(0.48), bar_h,
       size=12, bold=True, color=fill_color, align=PP_ALIGN.RIGHT)


# ── SLIDE 1 — COVER ──────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(BLANK)
bg(s1)
rect(s1, 0, 0, Inches(0.12), H, GREEN)
rect(s1, Inches(8.5), 0, Inches(4.83), Inches(2.2), DARK_BLUE)
tb(s1, "GENOS AI", Inches(9.2), Inches(0.4), Inches(3.5), Inches(0.8),
   size=28, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)
tb(s1, "www.genosai.tech", Inches(9.2), Inches(1.05), Inches(3.5), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
tb(s1, "HERBOTHS PROPERTY — WINDHOEK, NAMIBIA", Inches(0.55), Inches(1.6),
   Inches(9), Inches(0.6), size=13, color=GREEN)
tb(s1, "Website Audit\nReport", Inches(0.55), Inches(2.05),
   Inches(9), Inches(1.8), size=48, bold=True, color=WHITE)
aline(s1, Inches(3.85), color=GREEN, w=Inches(3))
tb(s1, "Prepared exclusively for Lizette da Fonseca — Woodlands @ Herboths Lifestyle Estate",
   Inches(0.55), Inches(4.15), Inches(9), Inches(0.5), size=13, color=LGREY)
rect(s1, 0, Inches(6.5), W, Inches(1.0), DARK_BLUE)
tb(s1, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699",
   Inches(0.55), Inches(6.6), Inches(9), Inches(0.5), size=12, color=MGREY)
tb(s1, "May 2026  |  CONFIDENTIAL", Inches(9.5), Inches(6.6), Inches(3.3), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
# Score box
rect(s1, Inches(10.5), Inches(3.1), Inches(2.3), Inches(2.7), DARK_BLUE)
tb(s1, "OVERALL SCORE", Inches(10.55), Inches(3.2), Inches(2.2), Inches(0.4),
   size=10, bold=True, color=MGREY, align=PP_ALIGN.CENTER)
tb(s1, "3.4", Inches(10.55), Inches(3.55), Inches(2.2), Inches(1.1),
   size=66, bold=True, color=RED, align=PP_ALIGN.CENTER)
tb(s1, "/ 10", Inches(10.55), Inches(4.5), Inches(2.2), Inches(0.38),
   size=18, color=LGREY, align=PP_ALIGN.CENTER)
tag(s1, "CRITICAL PRIORITY", Inches(10.6), Inches(5.0), fill=RED, tc=WHITE, size=9)


# ── SLIDE 2 — AT A GLANCE ────────────────────────────────────────────────────
s2 = prs.slides.add_slide(BLANK)
bg(s2)
rect(s2, 0, 0, Inches(0.12), H, GREEN)
tb(s2, "AT A GLANCE", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s2, "Herboths Property — Woodlands @ Herboths",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=30, bold=True, color=WHITE)
aline(s2, Inches(1.35), color=GREEN)

stats = [
    ("31 Plots",  "Freehold erven\n3–8 hectares each"),
    ("160 ha",    "Total estate,\nprivate woodland"),
    ("20 min",    "From Windhoek CBD\n21km east"),
    ("Sold Out",  "Herboths Blick —\npredecessor project"),
    ("Off-Grid",  "Solar + independent\nwater & sewerage"),
]
cw = Inches(2.3); gap = Inches(0.18)
for i, (num, label) in enumerate(stats):
    l = Inches(0.55) + i * (cw + gap)
    rect(s2, l, Inches(1.8), cw, Inches(2.2), DARK_BLUE)
    tb(s2, num, l + Inches(0.1), Inches(1.95), cw - Inches(0.2), Inches(1.05),
       size=20, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    tb(s2, label, l + Inches(0.1), Inches(3.0), cw - Inches(0.2), Inches(0.75),
       size=11, color=LGREY, align=PP_ALIGN.CENTER)

rect(s2, Inches(0.55), Inches(4.2), Inches(12.23), Inches(0.55), DARK_BLUE)
tb(s2, "Windhoek, Namibia  |  Lizette da Fonseca: +264 81 142 8785  |  herboths.com  |  LinkedIn · Facebook",
   Inches(0.65), Inches(4.3), Inches(12), Inches(0.42), size=11, color=LGREY, align=PP_ALIGN.CENTER)

tb(s2, "ESTATE HIGHLIGHTS", Inches(0.55), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=GREEN)
bullets(s2, [
    "31 freehold plots — 3 to 8 hectares, large serviced erven",
    "Solar power generation — fully independent from load-shedding grid",
    "Independent water and sewerage — no reliance on municipal supply",
    "Natural bush corridors supporting free-roaming wildlife",
    "HOA management with architectural controls for value protection"],
        Inches(0.55), Inches(5.35), Inches(5.8), Inches(1.6), size=12, color=LGREY)

tb(s2, "BUYER PROFILE", Inches(6.8), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=GREEN)
bullets(s2, [
    "Windhoek professionals seeking space, privacy and bush lifestyle",
    "Families wanting off-grid security against load-shedding and water cuts",
    "Investors in proven freehold land with HOA value protection",
    "Expats and retirees relocating to Namibia's capital region",
    "Nature-oriented buyers valuing wildlife and low-density living"],
        Inches(6.8), Inches(5.35), Inches(5.8), Inches(1.6), size=12, color=LGREY)


# ── SLIDE 3 — AUDIT SCORECARD ────────────────────────────────────────────────
s3 = prs.slides.add_slide(BLANK)
bg(s3)
rect(s3, 0, 0, Inches(0.12), H, RED)
tb(s3, "WEBSITE AUDIT SCORECARD", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=RED)
tb(s3, "herboths.com — Scored Across 8 Dimensions",
   Inches(0.55), Inches(0.7), Inches(10), Inches(0.7), size=30, bold=True, color=WHITE)
aline(s3, Inches(1.35), color=RED, w=Inches(2.0))

# Overall score badge (right side)
rect(s3, Inches(10.8), Inches(0.55), Inches(2.1), Inches(2.3), DARK_BLUE)
tb(s3, "OVERALL", Inches(10.85), Inches(0.65), Inches(2.0), Inches(0.35),
   size=10, bold=True, color=MGREY, align=PP_ALIGN.CENTER)
tb(s3, "3.4", Inches(10.85), Inches(0.95), Inches(2.0), Inches(0.95),
   size=56, bold=True, color=RED, align=PP_ALIGN.CENTER)
tb(s3, "/ 10", Inches(10.85), Inches(1.78), Inches(2.0), Inches(0.32),
   size=15, color=LGREY, align=PP_ALIGN.CENTER)

# Score bars — left column
categories_left = [
    ("Security & SSL",         1),
    ("SEO & Discoverability",  3),
    ("Content Quality",        4),
    ("Mobile Experience",      5),
]
categories_right = [
    ("Page Speed",             5),
    ("UX & Design",            4),
    ("Conversion & Lead Cap.", 2),
    ("Social & Digital Pres.", 1),
]

bar_w = Inches(4.8)
for i, (label, score) in enumerate(categories_left):
    t = Inches(1.7) + i * Inches(1.2)
    score_bar(s3, label, score, Inches(0.55), t, bar_w)

for i, (label, score) in enumerate(categories_right):
    t = Inches(1.7) + i * Inches(1.2)
    score_bar(s3, label, score, Inches(6.55), t, bar_w)

# Legend
rect(s3, Inches(0.55), Inches(6.85), Inches(12.23), Inches(0.38), DARK_BLUE)
for lx, lcolor, ltxt in [
    (Inches(0.8),  RED,   "0–3  Critical"),
    (Inches(3.2),  AMBER, "4–6  Needs Work"),
    (Inches(5.6),  GREEN, "7–10 Healthy"),
]:
    rect(s3, lx, Inches(6.93), Inches(0.18), Inches(0.18), lcolor)
    tb(s3, ltxt, lx + Inches(0.22), Inches(6.9), Inches(1.8), Inches(0.25),
       size=10, color=MGREY)

tb(s3, "Security and Conversion are the two lowest scores — both are immediate revenue risks.",
   Inches(7.2), Inches(6.88), Inches(5.8), Inches(0.3), size=10, italic=True, color=AMBER)


# ── SLIDE 4 — CRITICAL FINDINGS ──────────────────────────────────────────────
s4 = prs.slides.add_slide(BLANK)
bg(s4)
rect(s4, 0, 0, Inches(0.12), H, RED)
tb(s4, "CRITICAL FINDINGS", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=RED)
tb(s4, "10 Issues Found on herboths.com — Prioritised by Impact",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.65), size=28, bold=True, color=WHITE)
aline(s4, Inches(1.3), color=RED, w=Inches(2.0))

issues = [
    (RED,   "CRITICAL", "SSL Certificate Invalid",
     "Site shows a browser security warning on HTTPS. Every buyer who types herboths.com sees an 'unsafe site' alert. Many leave without ever reading about the estate."),
    (RED,   "CRITICAL", "No WhatsApp Contact Widget",
     "WhatsApp is the primary business channel in Namibia (92% penetration). No widget on the site means buyers who want to message immediately have no frictionless path."),
    (RED,   "CRITICAL", "No Lead Capture or Inquiry Automation",
     "The only contact mechanism appears to be a basic form or phone number. No automated response, no confirmation, no follow-up. Inquiries that arrive after hours go cold."),
    (AMBER, "HIGH",     "No Plot Availability Map or Status",
     "31 plots is the core product. Buyers can't see which plots are available, their orientation, size, or proximity to wildlife corridors without calling. This friction kills conversions."),
    (AMBER, "HIGH",     "Thin Content — No FAQs, No Blog, No Detail Pages",
     "Premium buyers do research before calling. There is no FAQ answering solar setup, HOA fees, building rules, or wildlife. No blog. No long-form content building trust before the call."),
    (AMBER, "HIGH",     "No SEO — Invisible on Google",
     "Searches for 'lifestyle estate Windhoek', 'freehold plots Namibia', or 'off-grid property near Windhoek' return no Herboths content. Buyers researching online never find the estate."),
    (AMBER, "HIGH",     "No Virtual Tour or Visual Estate Map",
     "160 hectares of woodland, wildlife corridors, and 31 plots — none of it is visually communicated. Buyers can't orient themselves to the estate before visiting."),
    (GOLD,  "MEDIUM",   "Social Media Disconnected from Website",
     "LinkedIn and Facebook exist but have 29 and minimal followers respectively. No feed embedded on site, no share prompts, no content driving traffic back to the website."),
    (GOLD,  "MEDIUM",   "No Testimonials from Herboths Blick Owners",
     "The predecessor project sold out. Those buyers are the most credible social proof available. None of their stories appear on the Woodlands site."),
    (GOLD,  "MEDIUM",   "No Mobile-Optimised Inquiry Flow",
     "Most Namibian buyers browse on mobile. The contact path on mobile needs to be a single tap to WhatsApp or call — not a form that requires typing a long message on a small keyboard."),
]

col_w = Inches(5.85)
left_issues  = issues[:5]
right_issues = issues[5:]

for col_idx, col_issues in enumerate([left_issues, right_issues]):
    lx = Inches(0.55) + col_idx * (col_w + Inches(0.3))
    for row_idx, (color, priority, title, desc) in enumerate(col_issues):
        ty = Inches(1.6) + row_idx * Inches(1.12)
        rect(s4, lx, ty, col_w, Inches(1.05), DARK_BLUE)
        rect(s4, lx, ty, Inches(0.07), Inches(1.05), color)
        # priority tag
        tag_w = Inches(len(priority) * 0.1 + 0.55)
        p_tag = s4.shapes.add_shape(9, lx + Inches(0.15), ty + Inches(0.08),
                                     tag_w, Inches(0.24))
        p_tag.fill.solid(); p_tag.fill.fore_color.rgb = color
        p_tag.line.fill.background()
        p_tf = p_tag.text_frame; p_p = p_tf.paragraphs[0]
        p_p.alignment = PP_ALIGN.CENTER
        p_r = p_p.add_run(); p_r.text = priority
        p_r.font.size = Pt(8); p_r.font.bold = True
        p_r.font.color.rgb = NAVY if color == GOLD else WHITE
        tb(s4, title, lx + Inches(0.15), ty + Inches(0.34),
           col_w - Inches(0.25), Inches(0.28), size=12, bold=True, color=WHITE)
        tb(s4, desc, lx + Inches(0.15), ty + Inches(0.62),
           col_w - Inches(0.25), Inches(0.4), size=10, color=LGREY, wrap=True)


# ── SLIDE 5 — SECURITY & SSL DEEP DIVE ───────────────────────────────────────
s5 = prs.slides.add_slide(BLANK)
bg(s5)
rect(s5, 0, 0, Inches(0.12), H, RED)
tb(s5, "CRITICAL ISSUE #1 — SECURITY", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=RED)
tb(s5, "Every Buyer Who Visits herboths.com Sees a Security Warning.",
   Inches(0.55), Inches(0.7), Inches(11), Inches(0.75), size=28, bold=True, color=WHITE)
aline(s5, Inches(1.4), color=RED, w=Inches(2.5))

# Score badge
rect(s5, Inches(11.8), Inches(0.55), Inches(1.3), Inches(1.3), DARK_BLUE)
tb(s5, "1/10", Inches(11.82), Inches(0.72), Inches(1.26), Inches(0.72),
   size=28, bold=True, color=RED, align=PP_ALIGN.CENTER)

# What the buyer sees
rect(s5, Inches(0.55), Inches(1.65), Inches(5.7), Inches(4.9), DARK_BLUE)
rect(s5, Inches(0.55), Inches(1.65), Inches(5.7), Inches(0.42), RED)
tb(s5, "WHAT A BUYER SEES ON CHROME / SAFARI", Inches(0.65), Inches(1.68),
   Inches(5.5), Inches(0.38), size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Simulated browser warning
rect(s5, Inches(0.75), Inches(2.18), Inches(5.3), Inches(3.9), RGBColor(0x20, 0x10, 0x10))
tb(s5, "Your connection is not private", Inches(0.85), Inches(2.35),
   Inches(5.1), Inches(0.45), size=14, bold=True, color=RED, align=PP_ALIGN.CENTER)
tb(s5, "Attackers might be trying to steal your information\nfrom herboths.com",
   Inches(0.85), Inches(2.82), Inches(5.1), Inches(0.55),
   size=11, color=LGREY, align=PP_ALIGN.CENTER)
tb(s5, "NET::ERR_CERT_COMMON_NAME_INVALID",
   Inches(0.85), Inches(3.42), Inches(5.1), Inches(0.35),
   size=10, color=MGREY, align=PP_ALIGN.CENTER)
rect(s5, Inches(1.3), Inches(3.85), Inches(2.0), Inches(0.45), DARK_BLUE)
tb(s5, "Back to safety", Inches(1.3), Inches(3.87), Inches(2.0), Inches(0.38),
   size=11, color=ACCENT, align=PP_ALIGN.CENTER)
rect(s5, Inches(3.4), Inches(3.85), Inches(2.2), Inches(0.45), RGBColor(0x30, 0x20, 0x20))
tb(s5, "Advanced (unsafe)", Inches(3.4), Inches(3.87), Inches(2.2), Inches(0.38),
   size=10, color=MGREY, align=PP_ALIGN.CENTER)
tb(s5, "Most buyers click 'Back to safety' and never return.",
   Inches(0.75), Inches(4.42), Inches(5.3), Inches(0.28),
   size=10, italic=True, color=RED, align=PP_ALIGN.CENTER)

# Impact breakdown
rect(s5, Inches(6.55), Inches(1.65), Inches(6.33), Inches(4.9), DARK_BLUE)
rect(s5, Inches(6.55), Inches(1.65), Inches(6.33), Inches(0.42), AMBER)
tb(s5, "WHY THIS IS A REVENUE PROBLEM", Inches(6.65), Inches(1.68),
   Inches(6.1), Inches(0.38), size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
impact_pts = [
    ("Buyer trust is destroyed before they read a word",
     "The security warning appears before any estate content. The product never gets a chance to speak — the browser already told them this site is dangerous."),
    ("Google penalises insecure sites in search rankings",
     "HTTPS is a Google ranking signal. A site without a valid SSL certificate ranks lower than competitors with identical content. Herboths is invisible on Google in part because of this."),
    ("WhatsApp link shares trigger the same warning",
     "When a buyer shares the herboths.com link on WhatsApp, the preview shows an 'unsafe' flag. Word-of-mouth referrals are being undermined by the SSL error."),
    ("Fix cost: under NAD 500/year. Impact: immediate",
     "An SSL certificate costs less than a tank of petrol. It can be installed in under an hour. This is the single highest-ROI fix available — and it should happen before anything else."),
]
for i, (heading, detail) in enumerate(impact_pts):
    ty = Inches(2.18) + i * Inches(1.12)
    rect(s5, Inches(6.65), ty + Inches(0.04), Inches(0.06), Inches(0.9), AMBER)
    tb(s5, heading, Inches(6.8), ty, Inches(5.9), Inches(0.32),
       size=12, bold=True, color=AMBER)
    tb(s5, detail, Inches(6.8), ty + Inches(0.32), Inches(5.9), Inches(0.55),
       size=11, color=LGREY, wrap=True)

rect(s5, Inches(0.55), Inches(6.72), Inches(12.23), Inches(0.45), DARK_BLUE)
tb(s5, "Priority: Fix the SSL certificate before any other work begins. Every other improvement is undermined while the security warning is live.",
   Inches(0.7), Inches(6.78), Inches(12), Inches(0.32), size=11, italic=True, color=RED)


# ── SLIDE 6 — CONVERSION & LEAD CAPTURE ──────────────────────────────────────
s6 = prs.slides.add_slide(BLANK)
bg(s6)
rect(s6, 0, 0, Inches(0.12), H, AMBER)
tb(s6, "CRITICAL ISSUE #2 — CONVERSION & LEAD CAPTURE", Inches(0.55), Inches(0.3),
   Inches(12), Inches(0.5), size=11, bold=True, color=AMBER)
tb(s6, "Buyers Arrive, Read, and Leave. Nothing Catches Them.",
   Inches(0.55), Inches(0.7), Inches(11), Inches(0.75), size=28, bold=True, color=WHITE)
aline(s6, Inches(1.4), color=AMBER, w=Inches(2.5))

rect(s6, Inches(11.8), Inches(0.55), Inches(1.3), Inches(1.3), DARK_BLUE)
tb(s6, "2/10", Inches(11.82), Inches(0.72), Inches(1.26), Inches(0.72),
   size=28, bold=True, color=RED, align=PP_ALIGN.CENTER)

conv_issues = [
    ("No WhatsApp Chat Widget", RED,
     "WhatsApp is how Namibians communicate with businesses. A buyer who wants to ask 'Is plot 7 still available?' has to call or fill a form — neither is as fast as WhatsApp. "
     "Competitors with a WhatsApp button on their site respond before the buyer has finished reading yours."),
    ("No Automated Inquiry Response", RED,
     "When a buyer submits a contact form at 7pm on a Friday, what happens? Nothing automated. No confirmation email. No acknowledgement. "
     "The buyer doesn't know if their message was received. Many re-inquire with a competitor who confirms instantly."),
    ("No Plot Availability Status", AMBER,
     "31 plots is the core inventory. Buyers can't see which plots are available, their size, orientation, or proximity to wildlife corridors without making a call. "
     "This single friction point is likely responsible for the majority of lost conversions — serious buyers won't call before they've done online research."),
    ("No CTA on Key Pages", AMBER,
     "Estate pages should end with a clear, single call to action: 'WhatsApp us to book a site visit' or 'See available plots.' "
     "Buyers who are interested but unsure of the next step leave rather than guess. The site doesn't tell them what to do next."),
]

cw = Inches(5.85); gap = Inches(0.25)
for i, (title, color, desc) in enumerate(conv_issues):
    col = i % 2; row = i // 2
    lx = Inches(0.55) + col * (cw + gap)
    ty = Inches(1.65) + row * Inches(2.45)
    rect(s6, lx, ty, cw, Inches(2.32), DARK_BLUE)
    rect(s6, lx, ty, cw, Inches(0.06), color)
    tb(s6, title, lx + Inches(0.15), ty + Inches(0.14),
       cw - Inches(0.3), Inches(0.35), size=13, bold=True, color=color)
    tb(s6, desc, lx + Inches(0.15), ty + Inches(0.52),
       cw - Inches(0.3), Inches(1.65), size=11, color=LGREY, wrap=True)

rect(s6, Inches(0.55), Inches(6.72), Inches(12.23), Inches(0.45), DARK_BLUE)
tb(s6, "A buyer who can't easily ask a question on their preferred channel will ask a competitor. With only 31 plots, every unconverted inquiry is a meaningful loss.",
   Inches(0.7), Inches(6.78), Inches(12), Inches(0.32), size=11, italic=True, color=AMBER)


# ── SLIDE 7 — SEO & CONTENT ───────────────────────────────────────────────────
s7 = prs.slides.add_slide(BLANK)
bg(s7)
rect(s7, 0, 0, Inches(0.12), H, PURPLE)
tb(s7, "ISSUE #3 — SEO & CONTENT", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=PURPLE)
tb(s7, "Buyers Searching for This Estate Can't Find It.",
   Inches(0.55), Inches(0.7), Inches(11), Inches(0.75), size=28, bold=True, color=WHITE)
aline(s7, Inches(1.4), color=PURPLE, w=Inches(2.0))

rect(s7, Inches(11.8), Inches(0.55), Inches(1.3), Inches(1.3), DARK_BLUE)
tb(s7, "3/10", Inches(11.82), Inches(0.72), Inches(1.26), Inches(0.72),
   size=28, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

# Left: keyword gaps
rect(s7, Inches(0.55), Inches(1.65), Inches(5.7), Inches(5.0), DARK_BLUE)
rect(s7, Inches(0.55), Inches(1.65), Inches(5.7), Inches(0.42), RED)
tb(s7, "SEARCH TERMS HERBOTHS IS MISSING", Inches(0.65), Inches(1.68),
   Inches(5.5), Inches(0.38), size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

kw_rows = [
    ("lifestyle estate Windhoek",         "High", "No content"),
    ("freehold plots near Windhoek",      "High", "No content"),
    ("off-grid property Namibia",         "High", "No content"),
    ("solar powered estate Namibia",      "Med",  "No content"),
    ("Woodlands Herboths review",         "Med",  "No page"),
    ("HOA estate Windhoek",               "Med",  "No content"),
    ("wildlife estate Namibia",           "Med",  "No content"),
    ("buy plot Windhoek",                 "High", "No content"),
    ("property investment Namibia 2025",  "Med",  "No content"),
]
col_heads = ["Keyword", "Volume", "Status"]
col_xs = [Inches(0.7), Inches(4.05), Inches(4.85)]
col_ws = [Inches(3.25), Inches(0.7), Inches(1.15)]
for ci, (head, cx, cw_) in enumerate(zip(col_heads, col_xs, col_ws)):
    tb(s7, head, cx, Inches(2.14), cw_, Inches(0.25),
       size=9, bold=True, color=MGREY)
for ri, (kw, vol, status) in enumerate(kw_rows):
    ty = Inches(2.42) + ri * Inches(0.28)
    row_fill_color = DARK_BLUE if ri % 2 == 0 else RGBColor(0x14, 0x22, 0x52)
    rect(s7, Inches(0.65), ty - Inches(0.02), Inches(5.3), Inches(0.27), row_fill_color)
    tb(s7, kw,     Inches(0.7),  ty, Inches(3.25), Inches(0.25), size=10, color=LGREY)
    vol_color = RED if vol == "High" else AMBER
    tb(s7, vol,    Inches(4.05), ty, Inches(0.7),  Inches(0.25), size=10, color=vol_color, bold=True)
    tb(s7, status, Inches(4.85), ty, Inches(1.15), Inches(0.25), size=10, color=RED)

# Right: content gaps
rect(s7, Inches(6.55), Inches(1.65), Inches(6.33), Inches(5.0), DARK_BLUE)
rect(s7, Inches(6.55), Inches(1.65), Inches(6.33), Inches(0.42), PURPLE)
tb(s7, "CONTENT GAPS IDENTIFIED", Inches(6.65), Inches(1.68),
   Inches(6.1), Inches(0.38), size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
content_gaps = [
    ("No FAQ Page",
     "Solar setup, HOA fees, building restrictions, plot transfer process, wildlife policy — "
     "buyers research all of this before calling. No answers on the site means buyers call competitors who do answer."),
    ("No Individual Plot Pages",
     "Each of the 31 plots should have its own page: size, orientation, proximity to wildlife corridor, views, price range. "
     "Buyers shortlist plots online. Without detail pages, they shortlist someone else's estate instead."),
    ("No Blog or Educational Content",
     "Articles like 'Why solar independence matters in Namibia' or 'What to know before buying a lifestyle estate' "
     "attract buyers who are still deciding. Currently, Herboths captures zero of this early-stage research traffic."),
    ("No Herboths Blick Success Content",
     "The predecessor project sold out — that's the most powerful proof point available. "
     "No testimonials, no case studies, no 'what Blick owners say' page. "
     "Social proof from real owners converts hesitant buyers faster than any marketing copy."),
]
for i, (heading, detail) in enumerate(content_gaps):
    ty = Inches(2.18) + i * Inches(1.1)
    rect(s7, Inches(6.65), ty + Inches(0.04), Inches(0.06), Inches(0.9), PURPLE)
    tb(s7, heading, Inches(6.8), ty, Inches(5.9), Inches(0.32),
       size=12, bold=True, color=PURPLE)
    tb(s7, detail, Inches(6.8), ty + Inches(0.32), Inches(5.9), Inches(0.6),
       size=11, color=LGREY, wrap=True)

rect(s7, Inches(0.55), Inches(6.72), Inches(12.23), Inches(0.45), DARK_BLUE)
tb(s7, "A buyer who searches 'lifestyle estate Windhoek' and lands on a competitor's FAQ page is not coming back to Herboths. Content is the pipeline that doesn't require ad spend.",
   Inches(0.7), Inches(6.78), Inches(12), Inches(0.32), size=11, italic=True, color=PURPLE)


# ── SLIDE 8 — COMPETITOR COMPARISON ──────────────────────────────────────────
s8 = prs.slides.add_slide(BLANK)
bg(s8)
rect(s8, 0, 0, Inches(0.12), H, TEAL)
tb(s8, "COMPETITOR WEBSITE COMPARISON", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=TEAL)
tb(s8, "What Strong Lifestyle Estate Websites Do — vs. Herboths Today",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.65), size=28, bold=True, color=WHITE)
aline(s8, Inches(1.3), color=TEAL, w=Inches(2.0))

features = [
    "SSL — secure site on all browsers",
    "WhatsApp chat button (mobile tap-to-message)",
    "Interactive plot map with availability status",
    "Individual plot detail pages with specs",
    "FAQ page: HOA, solar, building rules, transfer",
    "Photo gallery: wildlife, estate lifestyle, sunsets",
    "Testimonials from existing/previous buyers",
    "Virtual estate walkthrough or video tour",
    "Blog or news section with 5+ articles",
    "SEO titles & meta descriptions on all pages",
    "Google Maps integration with estate pin",
    "Social media feed embedded on homepage",
    "Mobile-first layout (single-tap call / WhatsApp)",
    "Email newsletter sign-up for plot updates",
]

# Header row
rect(s8, Inches(0.55), Inches(1.58), Inches(5.8), Inches(0.38), DARK_BLUE)
tb(s8, "FEATURE", Inches(0.65), Inches(1.61), Inches(5.6), Inches(0.32),
   size=10, bold=True, color=MGREY)
rect(s8, Inches(6.45), Inches(1.58), Inches(2.9), Inches(0.38), RED)
tb(s8, "HERBOTHS TODAY", Inches(6.45), Inches(1.61), Inches(2.9), Inches(0.32),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rect(s8, Inches(9.45), Inches(1.58), Inches(3.43), Inches(0.38), GREEN)
tb(s8, "TOP ESTATE STANDARD", Inches(9.45), Inches(1.61), Inches(3.43), Inches(0.32),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

herboths_status = [
    ("No", RED), ("No", RED), ("No", RED), ("No", RED),
    ("No", RED), ("Partial", AMBER), ("No", RED), ("No", RED),
    ("No", RED), ("No", RED), ("No", RED), ("No", RED),
    ("Partial", AMBER), ("No", RED),
]
row_h = Inches(0.355)
for i, (feature, (h_status, h_color)) in enumerate(zip(features, herboths_status)):
    ty = Inches(1.98) + i * row_h
    alt = DARK_BLUE if i % 2 == 0 else RGBColor(0x14, 0x22, 0x52)
    rect(s8, Inches(0.55), ty, Inches(5.8), row_h - Inches(0.02), alt)
    tb(s8, feature, Inches(0.68), ty + Inches(0.04),
       Inches(5.5), row_h - Inches(0.06), size=11, color=LGREY)
    rect(s8, Inches(6.45), ty, Inches(2.9), row_h - Inches(0.02), alt)
    tb(s8, h_status, Inches(6.45), ty + Inches(0.04), Inches(2.9), row_h - Inches(0.06),
       size=11, bold=True, color=h_color, align=PP_ALIGN.CENTER)
    rect(s8, Inches(9.45), ty, Inches(3.43), row_h - Inches(0.02), alt)
    tb(s8, "Yes", Inches(9.45), ty + Inches(0.04), Inches(3.43), row_h - Inches(0.06),
       size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)


# ── SLIDE 9 — RECOMMENDED FIXES ──────────────────────────────────────────────
s9 = prs.slides.add_slide(BLANK)
bg(s9)
rect(s9, 0, 0, Inches(0.12), H, GREEN)
tb(s9, "RECOMMENDED FIXES", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s9, "Priority-Ordered Action Plan for herboths.com",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.65), size=28, bold=True, color=WHITE)
aline(s9, Inches(1.3), color=GREEN, w=Inches(2.0))

fix_phases = [
    ("WEEK 1\nImmediate Fixes", RED, [
        ("Fix SSL Certificate",
         "Resolve the HTTPS certificate error. Site must load securely on all browsers. "
         "This is the prerequisite for everything else — no other fix matters while the security warning is live."),
        ("Add WhatsApp Chat Widget",
         "Single tap-to-chat button on every page. Mobile-first. "
         "Links directly to Lizette's WhatsApp number with a pre-filled opening message. "
         "Zero friction for the buyer's preferred contact channel."),
        ("Add Automated Inquiry Response",
         "Basic auto-reply to every contact form submission: confirmation of receipt, "
         "expected response time, and Lizette's WhatsApp number as an alternative. "
         "No inquiry goes unacknowledged."),
    ]),
    ("WEEKS 2–5\nCore Website Rebuild", AMBER, [
        ("Interactive Plot Map",
         "Visual map of the 160-hectare estate showing all 31 plots — "
         "available (green), reserved (amber), sold (grey). "
         "Buyers can click a plot to see size, orientation, price range, and wildlife proximity."),
        ("Individual Plot Detail Pages",
         "One page per plot: area in hectares, orientation, distance to wildlife corridor, "
         "views, topography notes, price range, and a direct WhatsApp CTA. "
         "Serious buyers shortlist online — give them the information to shortlist Herboths."),
        ("FAQ Page",
         "Answer the 15 most common buyer questions: HOA fees and rules, solar system specs, "
         "water and sewerage setup, building restrictions, transfer timeline, "
         "wildlife management, security, and proximity to schools and amenities."),
        ("Testimonials from Herboths Blick Owners",
         "Contact 5–8 Blick owners for short quotes or video testimonials. "
         "Their sold-out project is the strongest possible proof that the Herboths concept works."),
    ]),
    ("MONTH 2–3\nSEO & Content Engine", GREEN, [
        ("Blog Content — 8 Articles",
         "Publish articles targeting: 'lifestyle estate Windhoek 2025', 'off-grid living Namibia', "
         "'freehold plot investment Namibia', 'solar estate near Windhoek'. "
         "Each article draws in early-stage buyers who are still researching options."),
        ("Social Media Content Engine",
         "3 posts per week on LinkedIn and Facebook: wildlife sightings, solar independence stories, "
         "plot availability updates, buyer testimonials, and investment angle pieces. "
         "Grow from 29 followers to an audience of Windhoek property buyers."),
        ("AI Lead Nurture (Enhancement Layer)",
         "Once the website is sound: add WhatsApp AI agent for 24/7 inquiry handling + "
         "a 6-touch automated nurture sequence for warm leads on a 3–6 month decision cycle. "
         "This layer sits on top of a working website — not instead of one."),
    ]),
]

pw = Inches(3.9); gap4 = Inches(0.18)
for i, (phase_title, color, items) in enumerate(fix_phases):
    lx = Inches(0.55) + i * (pw + gap4)
    rect(s9, lx, Inches(1.58), pw, Inches(0.55), color)
    tb(s9, phase_title, lx + Inches(0.12), Inches(1.62),
       pw - Inches(0.24), Inches(0.48), size=12, bold=True, color=NAVY)
    rect(s9, lx, Inches(2.13), pw, Inches(5.0), DARK_BLUE)
    y_item = Inches(2.22)
    for title, desc in items:
        rect(s9, lx + Inches(0.12), y_item + Inches(0.05),
             Inches(0.05), Inches(0.55), color)
        tb(s9, title, lx + Inches(0.25), y_item,
           pw - Inches(0.38), Inches(0.32), size=12, bold=True, color=color)
        tb(s9, desc, lx + Inches(0.25), y_item + Inches(0.32),
           pw - Inches(0.38), Inches(0.52), size=10, color=LGREY, wrap=True)
        y_item += Inches(0.98) if len(items) == 4 else Inches(1.55)


# ── SLIDE 10 — NEXT STEPS ────────────────────────────────────────────────────
s10 = prs.slides.add_slide(BLANK)
bg(s10)
rect(s10, 0, 0, Inches(0.12), H, TEAL)
rect(s10, Inches(7.5), 0, Inches(5.83), H, DARK_BLUE)
tb(s10, "NEXT STEPS", Inches(0.55), Inches(0.3), Inches(6.5), Inches(0.5),
   size=11, bold=True, color=TEAL)
tb(s10, "Let's Talk for\n15 Minutes.",
   Inches(0.55), Inches(0.8), Inches(6.8), Inches(2.0), size=48, bold=True, color=WHITE)
aline(s10, Inches(2.8), color=TEAL, w=Inches(2.5))
tb(s10, "Start with the SSL fix and WhatsApp button — Week 1.\nEverything else builds on that foundation.",
   Inches(0.55), Inches(3.05), Inches(6.5), Inches(0.75), size=14, color=LGREY)

steps = [
    "Call or WhatsApp Rohan: +91 638-714-2699",
    "Week 1: SSL fix + WhatsApp button + automated inquiry response",
    "Weeks 2-5: Plot map, detail pages, FAQ, Blick testimonials",
    "Month 2-3: SEO content, social engine, AI nurture layer",
]
bullets(s10, steps, Inches(0.55), Inches(3.88), Inches(6.5), Inches(2.1),
        size=14, color=WHITE, dot=TEAL)

rect(s10, 0, Inches(6.5), Inches(7.5), Inches(1.0), DARK_BLUE)
tb(s10, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699  |  www.genosai.tech",
   Inches(0.55), Inches(6.6), Inches(6.8), Inches(0.5), size=11, color=MGREY)

tb(s10, "AUDIT SUMMARY", Inches(7.85), Inches(0.6), Inches(5.0), Inches(0.4),
   size=14, bold=True, color=TEAL)
aline(s10, Inches(1.0), color=TEAL, w=Inches(1.5))

summary_items = [
    ("Security",   "1/10", RED,   "SSL cert invalid — security warning on every visit"),
    ("Conversion", "2/10", RED,   "No WhatsApp, no plot map, no automated response"),
    ("SEO",        "3/10", RED,   "Zero keyword content — invisible on Google"),
    ("Content",    "4/10", AMBER, "Thin pages — no FAQ, no plot detail, no testimonials"),
    ("UX/Design",  "4/10", AMBER, "Basic — no visual estate map, weak CTAs"),
    ("Social",     "1/10", RED,   "29 followers — no content engine behind LinkedIn/FB"),
]
for i, (label, score, sc, detail) in enumerate(summary_items):
    ty = Inches(1.2) + i * Inches(0.88)
    rect(s10, Inches(7.85), ty, Inches(1.1), Inches(0.75), DARK_BLUE)
    tb(s10, score, Inches(7.88), ty + Inches(0.1), Inches(1.04), Inches(0.5),
       size=20, bold=True, color=sc, align=PP_ALIGN.CENTER)
    tb(s10, label, Inches(9.05), ty + Inches(0.04), Inches(1.2), Inches(0.28),
       size=11, bold=True, color=TEAL)
    tb(s10, detail, Inches(9.05), ty + Inches(0.35), Inches(3.95), Inches(0.38),
       size=10, color=LGREY, wrap=True)


out_path = ROOT / "output" / "decks" / "Herboths_Audit_GenosAI.pptx"
alt_path = ROOT / "output" / "decks" / "Herboths_Audit_GenosAI_v2.pptx"
try:
    prs.save(out_path)
    print(f"Saved: {out_path.name}")
except PermissionError:
    prs.save(alt_path)
    print(f"Saved as {alt_path.name} (close the original file in PowerPoint to overwrite it)")
