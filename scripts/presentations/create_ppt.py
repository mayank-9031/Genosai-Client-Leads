import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.parent / "utils"))
from ppt_theme import (
    NAVY, DARK_BLUE, ACCENT, GOLD, RED, GREEN, PURPLE, ORANGE, TEAL, CYAN,
    WHITE, LGREY, MGREY, W, H,
    bg, rect, tb, aline, bullets, tag,
    slide_header, slide_footer, score_badge,
    stat_cards, top_bar_card, phase_card, before_after, genos_logo_block,
)

ROOT = Path(__file__).parent.parent.parent
prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ── SLIDE 1 — COVER ──────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(BLANK)
bg(s1)
rect(s1, 0, 0, Inches(0.12), H, ACCENT)
genos_logo_block(s1)
tb(s1, "SPECIALTY CONSULTANTS INC.", Inches(0.55), Inches(1.6),
   Inches(9), Inches(0.6), size=13, color=ACCENT)
tb(s1, "Website & AI Growth\nAudit Report", Inches(0.55), Inches(2.05),
   Inches(9), Inches(1.8), size=44, bold=True, color=WHITE)
aline(s1, Inches(3.85), color=GOLD, w=Inches(3))
tb(s1, "Prepared exclusively for Thomas Williams, COO & EVP",
   Inches(0.55), Inches(4.15), Inches(9), Inches(0.5), size=13, color=LGREY)
slide_footer(s1)
score_badge(s1, score=5, priority_text="HIGH PRIORITY", priority_color=RED)


# ── SLIDE 2 — AT A GLANCE ────────────────────────────────────────────────────
s2 = prs.slides.add_slide(BLANK)
bg(s2)
slide_header(s2, "AT A GLANCE", "Specialty Consultants Inc. — Who They Are",
             accent_color=CYAN, aline_w=Inches(1.5))

stat_cards(s2, [
    ("5\nDecades", "In Business"),
    ("7,500+", "Searches\nCompleted"),
    ("Forbes", "Best Exec\nSearch Firms"),
    ("30+", "Blue-Chip\nClients"),
    ("SCI.ON", "Site-Level\nStaffing Product"),
], num_color=CYAN)

rect(s2, Inches(0.55), Inches(4.3), Inches(12.23), Inches(0.6), DARK_BLUE)
tb(s2, "KEY CLIENTS:  Prologis  ·  Brookfield  ·  Related Companies  ·  Truist  ·  GID  ·  CenterPoint",
   Inches(0.65), Inches(4.38), Inches(12), Inches(0.45), size=12, color=LGREY, align=PP_ALIGN.CENTER)

tb(s2, "SERVICE AREAS", Inches(0.55), Inches(5.1), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=CYAN)
bullets(s2, ["Investment Management & Finance", "Development & Construction",
             "Property Operations & Management"],
        Inches(0.55), Inches(5.45), Inches(5.8), Inches(1.5), size=12, dot=CYAN)

tb(s2, "ASSET CLASSES", Inches(6.8), Inches(5.1), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=CYAN)
bullets(s2, ["Multifamily & Affordable Housing", "Industrial / Logistics",
             "Seniors Housing & Healthcare"],
        Inches(6.8), Inches(5.45), Inches(5.8), Inches(1.5), size=12, dot=CYAN)


# ── SLIDE 3 — WHAT'S WORKING ─────────────────────────────────────────────────
s3 = prs.slides.add_slide(BLANK)
bg(s3)
slide_header(s3, "STRENGTHS", "What's Working — Strengths to Build On",
             accent_color=GREEN, aline_w=Inches(1.5))

strengths = [
    ("Forbes Recognition",
     "Named one of 'America's Best Executive Search Firms' — a credibility signal few competitors can match."),
    ("Blue-Chip Client Roster",
     "Working with Prologis, Brookfield, Related Companies, Truist, and 30+ others provides undeniable social proof."),
    ("5-Decade Track Record",
     "7,500+ completed executive searches is a compounding moat. Longevity builds trust in a relationship-driven business."),
    ("SCI.ON Initiative",
     "Launching a dedicated site-level staffing technology product shows forward-thinking leadership beyond headhunting."),
    ("Clear Service Structure",
     "The website segments Investment Management, Development, Construction, Management, and Finance — easy to self-qualify."),
]
for i, (title, desc) in enumerate(strengths):
    t = Inches(1.7) + i * Inches(1.02)
    rect(s3, Inches(0.55), t, Inches(0.06), Inches(0.72), GREEN)
    tb(s3, title, Inches(0.75), t, Inches(4), Inches(0.38),
       size=13, bold=True, color=WHITE)
    tb(s3, desc, Inches(0.75), t + Inches(0.36), Inches(11.8), Inches(0.55),
       size=12, color=LGREY)


# ── SLIDE 4 — WHERE LEADS ARE LEAKING ────────────────────────────────────────
s4 = prs.slides.add_slide(BLANK)
bg(s4)
slide_header(s4, "AUDIT FINDINGS", "Critical Gaps Identified in the Audit",
             accent_color=RED, aline_w=Inches(1.2))

issues = [
    ("NO LIVE CHAT OR AI BOT", "HIGH",
     "Zero real-time engagement for visitors. Candidates and clients landing outside business hours leave with no way to connect — permanently lost to a competitor."),
    ("ONLY 1 BLOG POST", "HIGH",
     "Real estate hiring managers search Google monthly for executive search resources. With 1 article, SCI is invisible. Competitors rank instead."),
    ("NO CLIENT TESTIMONIALS", "MEDIUM",
     "Stats (7,500 searches) are good. Quotes are better. No client voices reduces conversion for first-time visitors who can't independently validate the track record."),
    ("WEAK CTAs THROUGHOUT", "MEDIUM",
     "'Schedule a Call' buried in the nav is the only conversion path. No lead magnets, gated content, or email capture. The site is a brochure, not a funnel."),
    ("SCI.ON VALUE PROP MISSING", "MEDIUM",
     "A product page with only a tagline and 'Learn More' fails to sell SCI.ON to site-level staffing clients. This is a missed revenue stream."),
]
for i, (title, priority, desc) in enumerate(issues):
    t = Inches(1.7) + i * Inches(1.02)
    p_color = RED if priority == "HIGH" else GOLD
    rect(s4, Inches(0.55), t, Inches(0.06), Inches(0.72), p_color)
    tb(s4, title, Inches(0.75), t, Inches(9), Inches(0.38),
       size=13, bold=True, color=WHITE)
    tag(s4, priority, Inches(10.1), t + Inches(0.05), fill=p_color, size=9)
    tb(s4, desc, Inches(0.75), t + Inches(0.36), Inches(11.8), Inches(0.55),
       size=12, color=LGREY)


# ── SLIDE 5 — WEBSITE DESIGN AUDIT ───────────────────────────────────────────
s5 = prs.slides.add_slide(BLANK)
bg(s5)
slide_header(s5, "WEBSITE DESIGN AUDIT",
             "Your Website Looks Like 2015. Your Clients Expect 2025.",
             accent_color=PURPLE, aline_w=Inches(1.5), headline_size=26)

# 3-col comparison
rect(s5, Inches(0.55), Inches(2.1), Inches(3.8), Inches(3.5), DARK_BLUE)
rect(s5, Inches(0.55), Inches(2.1), Inches(3.8), Inches(0.42), RED)
tb(s5, "SCI TODAY — WHAT WE SEE", Inches(0.65), Inches(2.13), Inches(3.6), Inches(0.38),
   size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Dense text, no visual breathing room",
    "Generic blue-on-white — zero brand personality",
    "No hero video or dynamic content",
    "2015-style corporate brochure typography",
    "CTAs buried — no conversion hierarchy",
    "No sticky nav or UX micro-interactions",
    "Team page: no photos, bios, or LinkedIn",
    "SCI.ON: tagline only, no landing funnel",
], Inches(0.7), Inches(2.62), Inches(3.6), Inches(2.8), size=11, dot=RED)

tb(s5, "→", Inches(4.43), Inches(3.55), Inches(0.5), Inches(0.55),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(5.0), Inches(2.1), Inches(3.6), Inches(3.5), DARK_BLUE)
rect(s5, Inches(5.0), Inches(2.1), Inches(3.6), Inches(0.42), GOLD)
tb(s5, "INDUSTRY LEADERS — WHAT THEY DO", Inches(5.1), Inches(2.13), Inches(3.4), Inches(0.38),
   size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "CBRE: full-screen video hero, dark premium palette",
    "JLL: animated stats, bold type, card case studies",
    "Newmark: clean white space, strong hierarchy",
    "Cushman & Wakefield: motion design, modular UX",
    "Lincoln Property: fast, mobile-first, clear CTAs",
    "Top search firms: AI chat on every page",
    "Modern leaders: gated thought leadership content",
    "Best-in-class: 95+ mobile Lighthouse scores",
], Inches(5.1), Inches(2.62), Inches(3.4), Inches(2.8), size=11, dot=GOLD)

tb(s5, "→", Inches(8.68), Inches(3.55), Inches(0.5), Inches(0.55),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(9.25), Inches(2.1), Inches(3.6), Inches(3.5), DARK_BLUE)
rect(s5, Inches(9.25), Inches(2.1), Inches(3.6), Inches(0.42), GREEN)
tb(s5, "WHAT GENOS AI BUILDS FOR SCI", Inches(9.35), Inches(2.13), Inches(3.4), Inches(0.38),
   size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Dark/premium design with bold brand identity",
    "Full-screen hero with animated stat counters",
    "Sector case study cards (Multifamily, Industrial…)",
    "SCI.ON dedicated landing page + funnel",
    "Team page: pro headshots, bios, LinkedIn links",
    "AI chat trained on SCI services — always on",
    "Mobile-first, sub-2s load, SEO architecture",
    "Lead capture: forms, calendar booking, CRM sync",
], Inches(9.35), Inches(2.62), Inches(3.4), Inches(2.8), size=11, dot=GREEN)

# Bottom stat bar
rect(s5, Inches(0.55), Inches(5.82), Inches(12.23), Inches(1.28), DARK_BLUE)
tb(s5, "WHY THIS MATTERS:", Inches(0.7), Inches(5.88), Inches(2.5), Inches(0.38),
   size=11, bold=True, color=PURPLE)
ux_stats = [
    ("94%", "of first impressions\nare design-related"),
    ("75%", "of visitors judge\ncredibility by design"),
    ("3 sec", "before a visitor\nbounces — or stays"),
    ("200%+", "more leads from\na modern redesign"),
]
for j, (num, label) in enumerate(ux_stats):
    lx = Inches(3.0) + j * Inches(2.45)
    tb(s5, num, lx, Inches(5.85), Inches(1.6), Inches(0.48),
       size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s5, label, lx, Inches(6.32), Inches(1.6), Inches(0.55),
       size=10, color=MGREY, align=PP_ALIGN.CENTER)


# ── SLIDE 6 — AI AUTOMATION OPPORTUNITIES ────────────────────────────────────
s6 = prs.slides.add_slide(BLANK)
bg(s6)
slide_header(s6, "AI AUTOMATION OPPORTUNITIES", "Six AI Systems We'd Build for SCI",
             accent_color=ACCENT, aline_w=Inches(1.2))

opps = [
    ("AI Candidate\nPre-Screening", TEAL,
     "Automated resume parsing and scoring against role requirements — cuts consultant time-to-shortlist by 30–50%."),
    ("24/7 AI\nChat Widget", CYAN,
     "Conversational AI trained on SCI's services captures and qualifies candidate and client inquiries around the clock."),
    ("SCI.ON AI\nMatching Engine", GOLD,
     "Layer AI matching intelligence onto SCI.ON to score site-level candidates against open roles in real time."),
    ("Automated Candidate\nOutreach", PURPLE,
     "Multi-touch personalized outreach sequences that nurture passive candidates until they're ready to engage — at scale."),
    ("Executive Movement\nIntelligence", ORANGE,
     "AI monitors LinkedIn and news for signals (promotions, company events) to surface warm, timely outreach opportunities."),
    ("Automated Client\nReporting", GREEN,
     "Auto-generated weekly search progress updates sent to clients — saving 3–5 hours per consultant per week."),
]
cw = Inches(2.05); gap = Inches(0.13)
for i, (title, color, desc) in enumerate(opps):
    l = Inches(0.55) + i * (cw + gap)
    top_bar_card(s6, l, Inches(1.8), cw, Inches(4.7), color, title, desc,
                 title_size=12, desc_size=10)

rect(s6, Inches(0.55), Inches(6.62), Inches(12.23), Inches(0.55), DARK_BLUE)
tb(s6, "Every system plugs into SCI's existing ATS and workflow — no replacement, only augmentation.",
   Inches(0.7), Inches(6.68), Inches(12), Inches(0.4), size=11, italic=True, color=GOLD)


# ── SLIDE 7 — WHAT THIS IS COSTING YOU ───────────────────────────────────────
s7 = prs.slides.add_slide(BLANK)
bg(s7)
slide_header(s7, "THE COST OF INACTION", "The Hidden Cost of the Status Quo",
             accent_color=GOLD, aline_w=Inches(1.5))

costs = [
    ("Manual Screening",
     "Consultants spend 60%+ of their time on manual resume review. At SCI's billing rates, that's significant revenue-generating time going to administrative work instead of client relationships."),
    ("After-Hours Blind Spots",
     "No chat or AI bot means every website visitor outside 9–5 leaves without capturing their interest. In executive search, first contact timing is everything — best candidates are passive and busy."),
    ("SEO Invisibility",
     "Real estate executives and HR leaders Google 'real estate executive search firm' thousands of times monthly. With 1 blog post, SCI doesn't appear. Competitors with content strategies capture those leads."),
    ("Lost SCI.ON Revenue",
     "An under-explained product page leaves site-level staffing clients on the table. SCI.ON has the potential to be a significant ARR business — it needs a proper digital funnel."),
]
for i, (title, desc) in enumerate(costs):
    t = Inches(1.7) + i * Inches(1.28)
    rect(s7, Inches(0.55), t, Inches(12.23), Inches(1.12), DARK_BLUE)
    rect(s7, Inches(0.55), t, Inches(0.06), Inches(1.12), GOLD)
    tb(s7, title, Inches(0.75), t + Inches(0.1), Inches(4), Inches(0.42),
       size=13, bold=True, color=GOLD)
    tb(s7, desc, Inches(0.75), t + Inches(0.52), Inches(11.8), Inches(0.55),
       size=12, color=LGREY)


# ── SLIDE 8 — WHAT WE'D BUILD ─────────────────────────────────────────────────
s8 = prs.slides.add_slide(BLANK)
bg(s8)
slide_header(s8, "WHAT WE'D BUILD", "The Genos AI Stack for SCI",
             accent_color=ACCENT, aline_w=Inches(1.2))

builds = [
    ("AI Candidate Screening Pipeline",
     "ATS-integrated AI scores inbound resumes, flags top matches, and de-prioritizes poor fits — automatically.",
     "Weeks 1–4"),
    ("AI Chat Widget (Website)",
     "Conversational assistant trained on SCI's services, sectors, and FAQs. Captures leads 24/7 and routes to the right consultant.",
     "Weeks 1–3"),
    ("SCI.ON AI Matching Layer",
     "AI ranking engine for site-level candidates — matches open roles to profiles using skills, location, and availability signals.",
     "Weeks 4–8"),
    ("Automated Candidate Outreach Engine",
     "Personalized multi-touch outreach sequences (email + LinkedIn) that nurture passive executive candidates through a warm pipeline.",
     "Weeks 3–6"),
    ("SEO Content Engine",
     "4–8 expert articles per month targeting real estate hiring managers on Google — built from SCI's expertise and case studies.",
     "Ongoing"),
]
for i, (title, desc, timeline) in enumerate(builds):
    t = Inches(1.65) + i * Inches(1.0)
    rect(s8, Inches(0.55), t, Inches(11.6), Inches(0.88), DARK_BLUE)
    rect(s8, Inches(0.55), t, Inches(0.06), Inches(0.88), ACCENT)
    rect(s8, Inches(0.68), t + Inches(0.18), Inches(0.48), Inches(0.48), ACCENT)
    tb(s8, str(i + 1), Inches(0.68), t + Inches(0.14), Inches(0.48), Inches(0.48),
       size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s8, title, Inches(1.28), t + Inches(0.08), Inches(7.5), Inches(0.36),
       size=13, bold=True, color=WHITE)
    tb(s8, desc, Inches(1.28), t + Inches(0.46), Inches(7.5), Inches(0.38),
       size=11, color=LGREY)
    tag(s8, timeline, Inches(12.25), t + Inches(0.28), fill=DARK_BLUE, tc=MGREY, size=10)


# ── SLIDE 9 — 90-DAY OUTCOMES ────────────────────────────────────────────────
s9 = prs.slides.add_slide(BLANK)
bg(s9)
slide_header(s9, "90-DAY OUTCOMES", "What SCI Looks Like 90 Days from Now",
             accent_color=CYAN, aline_w=Inches(1.2))

phase_card(s9, Inches(0.55), Inches(1.55), Inches(3.9), Inches(5.05), ACCENT,
           "MONTH 1", "Days 1–30", [
               ("AI chat live on specialtyconsultants.com",
                "Capturing candidate & client inquiries 24/7. First 3 SEO articles published and indexing."),
               ("Automated outreach sequences drafted",
                "Sequences tested and deployed for passive executive candidates."),
               ("SCI.ON landing page live",
                "Dedicated funnel with value prop, case studies, and contact gate."),
           ])

phase_card(s9, Inches(4.65), Inches(1.55), Inches(3.9), Inches(5.05), GOLD,
           "MONTH 2", "Days 31–60", [
               ("AI screening pipeline in ATS",
                "30% reduction in time-to-shortlist per mandate. 4–6 articles live and ranking."),
               ("SCI.ON AI matching prototype tested",
                "Internal validation against active site-level mandates."),
               ("Content velocity ramping",
                "6–8 articles published; first ranking signals from Google."),
           ])

phase_card(s9, Inches(8.75), Inches(1.55), Inches(3.9), Inches(5.05), GREEN,
           "MONTH 3", "Days 61–90", [
               ("SCI.ON AI matching live for clients",
                "Executive movement intelligence alerts active. Timely, warm outreach at scale."),
               ("Automated client reporting active",
                "Saving 3–5 hrs/week per consultant. Full funnel: search → chat → call."),
               ("10+ articles indexing, SEO traction",
                "Inbound qualified leads from organic search beginning to flow."),
           ])


# ── SLIDE 10 — NEXT STEPS ────────────────────────────────────────────────────
s10 = prs.slides.add_slide(BLANK)
bg(s10)
rect(s10, 0, 0, Inches(0.12), H, CYAN)
rect(s10, Inches(7.5), 0, Inches(5.83), H, DARK_BLUE)

tb(s10, "NEXT STEPS", Inches(0.55), Inches(0.3), Inches(6.5), Inches(0.5),
   size=11, bold=True, color=CYAN)
tb(s10, "One Call to See\nWhat's Possible.",
   Inches(0.55), Inches(0.8), Inches(6.8), Inches(2.0), size=44, bold=True, color=WHITE)
aline(s10, Inches(2.8), color=CYAN, w=Inches(2.5))
tb(s10, "No pitch. No commitment.\nWe walk you through the audit findings and answer any questions.",
   Inches(0.55), Inches(3.05), Inches(6.5), Inches(0.8), size=14, color=LGREY)

steps = [
    "Reply or book a 15-min call with Rohan: hello@genosai.tech",
    "Free strategy session: map the highest-ROI automations for SCI",
    "Custom AI roadmap: timeline, deliverables, cost estimate",
    "Month 1: AI chat live + screening pipeline + SCI.ON funnel",
]
bullets(s10, steps, Inches(0.55), Inches(3.95), Inches(6.5), Inches(2.0),
        size=13, color=WHITE, dot=CYAN)

rect(s10, 0, Inches(6.5), Inches(7.5), Inches(1.0), DARK_BLUE)
tb(s10, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699  |  www.genosai.tech",
   Inches(0.55), Inches(6.6), Inches(6.8), Inches(0.5), size=11, color=MGREY)

tb(s10, "THREE GAPS.\nTHREE FIXES.", Inches(7.85), Inches(0.6), Inches(5.0), Inches(1.1),
   size=22, bold=True, color=CYAN)
aline(s10, Inches(1.65), color=CYAN)

summary = [
    ("No Live Chat", "AI chat widget — captures leads 24/7, trained on SCI services"),
    ("1 Blog Post", "SEO content engine — 4-8 articles/month targeting real estate searches"),
    ("No Testimonials", "Client quote capture — automated post-search review requests"),
    ("Weak CTAs", "Full funnel redesign — lead magnets, calendar booking, CRM sync"),
    ("SCI.ON Gap", "Dedicated landing page + AI matching layer + client demo flow"),
]
for i, (label, val) in enumerate(summary):
    t2 = Inches(1.85) + i * Inches(0.88)
    rect(s10, Inches(7.85), t2, Inches(1.5), Inches(0.72), NAVY)
    tb(s10, label, Inches(7.9), t2 + Inches(0.05), Inches(1.4), Inches(0.32),
       size=10, bold=True, color=CYAN)
    tb(s10, val, Inches(9.5), t2 + Inches(0.05), Inches(3.6), Inches(0.65),
       size=11, color=LGREY)


prs.save(ROOT / "output" / "decks" / "SCI_Audit_GenosAI.pptx")
print("Saved: SCI_Audit_GenosAI.pptx")
