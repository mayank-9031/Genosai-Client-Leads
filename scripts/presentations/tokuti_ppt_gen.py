from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Tokuti / Web3 Palette  (dark mode · electric blue · purple) ─
VOID        = RGBColor(0x08, 0x0A, 0x14)   # near-black bg
DEEP_NAVY   = RGBColor(0x0D, 0x12, 0x2A)   # dark card
ELECTRIC    = RGBColor(0x00, 0xC2, 0xFF)   # electric cyan/blue
PURPLE      = RGBColor(0x7B, 0x2F, 0xFF)   # vivid purple
VIOLET      = RGBColor(0xA0, 0x60, 0xFF)   # soft violet
NEON_GREEN  = RGBColor(0x00, 0xE5, 0x9E)   # neon green accent
GOLD        = RGBColor(0xF5, 0xC5, 0x42)   # gold
IVORY       = RGBColor(0xF4, 0xF6, 0xFB)   # light slide bg
SILVER      = RGBColor(0x8A, 0x90, 0xAA)   # silver text
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DANGER      = RGBColor(0xFF, 0x45, 0x45)   # red alert
CAUTION     = RGBColor(0xFF, 0xA5, 0x30)   # amber caution
SUCCESS     = RGBColor(0x00, 0xE5, 0x9E)   # green success

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

# ── Helpers ─────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.line.fill.background(); s.line.width = 0
    s.fill.solid(); s.fill.fore_color.rgb = fill
    return s

def add_textbox(slide, text, x, y, w, h,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(font_size); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return txb

def bg(slide):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, IVORY)

def dark_bg(slide):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, VOID)

def header(slide, title, sub=None, dark=False):
    hbg = VOID if not dark else DEEP_NAVY
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.25), hbg)
    add_rect(slide, 0, Inches(1.25), Inches(0.38), Inches(6.25), ELECTRIC)
    add_textbox(slide, title, Inches(0.58), Inches(0.16),
                Inches(10.5), Inches(0.64), 27, True, WHITE)
    if sub:
        add_textbox(slide, sub, Inches(0.58), Inches(0.74),
                    Inches(10.5), Inches(0.42), 12, False, ELECTRIC, italic=True)

def footer(slide, txt="Tokuti  |  Web Development & Automation Proposal  |  Genos Apollo  |  2026"):
    add_rect(slide, 0, Inches(7.2), SLIDE_W, Inches(0.3), VOID)
    add_textbox(slide, txt, Inches(0.3), Inches(7.2),
                Inches(12.7), Inches(0.3), 9, False, SILVER, PP_ALIGN.CENTER)

def card(slide, x, y, w, h, fill=WHITE, border_col=None, dark=False):
    cf = RGBColor(0x14, 0x1C, 0x38) if dark else fill
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = cf
    if border_col:
        s.line.color.rgb = border_col; s.line.width = Pt(0.8)
    else:
        s.line.fill.background(); s.line.width = 0
    return s

def glow_bar(slide, x, y, w, h=Inches(0.052), col=ELECTRIC):
    add_rect(slide, x, y, w, h, col)

def stat_card_dark(slide, val, lbl, x, y, w=Inches(2.8), h=Inches(1.52)):
    card(slide, x, y, w, h, dark=True)
    glow_bar(slide, x, y, w, Inches(0.045), ELECTRIC)
    add_textbox(slide, val, x, y + Inches(0.18), w, Inches(0.72),
                30, True, ELECTRIC, PP_ALIGN.CENTER)
    add_textbox(slide, lbl, x, y + Inches(0.9), w, Inches(0.52),
                11, False, SILVER, PP_ALIGN.CENTER)

def bullet_card(slide, title, bullets, x, y, w, h,
                hdr_bg=VOID, hdr_col=ELECTRIC, dark=False):
    card(slide, x, y, w, h, dark=dark, border_col=ELECTRIC)
    add_rect(slide, x, y, w, Inches(0.44), hdr_bg)
    add_textbox(slide, title, x + Inches(0.14), y + Inches(0.06),
                w - Inches(0.28), Inches(0.36), 12, True, hdr_col)
    ty = y + Inches(0.58)
    tc = WHITE if dark else DEEP_NAVY
    for b in bullets:
        add_textbox(slide, f"◆  {b}", x + Inches(0.15), ty,
                    w - Inches(0.3), Inches(0.36), 10.5, False, tc)
        ty += Inches(0.34)

# ════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE  (dark cinematic)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_bg(s)
add_rect(s, 0, 0, Inches(0.72), SLIDE_H, ELECTRIC)
add_rect(s, Inches(0.72), Inches(4.72), SLIDE_W - Inches(0.72), Inches(0.055), ELECTRIC)

# Subtle grid lines
for i in range(1, 8):
    add_rect(s, Inches(0.72), Inches(i * 0.9), SLIDE_W, Inches(0.008),
             RGBColor(0x14, 0x20, 0x40))

add_textbox(s, "TOKUTI  ·  WEB3 REAL ESTATE TOKENIZATION",
            Inches(1.1), Inches(0.82), Inches(11.5), Inches(0.55),
            12, True, ELECTRIC)
add_textbox(s, "Website Services &\nAI Automation Proposal",
            Inches(1.1), Inches(1.52), Inches(11.2), Inches(2.15),
            44, True, WHITE)
add_textbox(s, "Why Webflow is your starting point — not your destination",
            Inches(1.1), Inches(3.78), Inches(11.2), Inches(0.55),
            15, False, VIOLET, italic=True)
add_textbox(s, "Prepared for: Mark Mariampillai, CEO & Founder — Tokuti.io",
            Inches(1.1), Inches(5.05), Inches(10), Inches(0.4),
            12, False, SILVER)
add_textbox(s, "By: Genos Apollo  |  May 2026",
            Inches(1.1), Inches(5.52), Inches(8), Inches(0.35),
            11, False, RGBColor(0x55, 0x66, 0x88))
add_textbox(s, "CONFIDENTIAL",
            Inches(1.1), Inches(6.68), Inches(4), Inches(0.32),
            10, False, SILVER, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Agenda", "What we cover in this proposal"); footer(s)

items = [
    ("01", "About Tokuti",               "Your vision, platform and where you are building today"),
    ("02", "The Webflow Ceiling",         "Why Webflow works now — and where it breaks down"),
    ("03", "Our Web Development Pitch",  "What a professional build gives Tokuti instead"),
    ("04", "AI Automation Opportunities","Where AI accelerates your platform, deals & growth"),
    ("05", "Benefits of Automation",     "Speed, scale and investor confidence"),
    ("06", "Implementation Plan",        "How we build it, phase by phase"),
    ("07", "Investment & Next Steps",    "Transparent pricing and how to get started"),
]
for i, (num, title, desc) in enumerate(items):
    col, row = i % 2, i // 2
    x = Inches(0.55) + col * Inches(6.4)
    y = Inches(1.55) + row * Inches(1.22)
    card(s, x, y, Inches(6.1), Inches(1.05), border_col=RGBColor(0xCC, 0xD8, 0xF0))
    add_rect(s, x, y, Inches(0.7), Inches(1.05), VOID)
    glow_bar(s, x, y + Inches(1.05) - Inches(0.052), Inches(6.1))
    add_textbox(s, num, x, y + Inches(0.25), Inches(0.7), Inches(0.5),
                20, True, ELECTRIC, PP_ALIGN.CENTER)
    add_textbox(s, title, x + Inches(0.82), y + Inches(0.1),
                Inches(5.1), Inches(0.4), 13, True, VOID)
    add_textbox(s, desc,  x + Inches(0.82), y + Inches(0.56),
                Inches(5.1), Inches(0.38), 10, False, SILVER, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 3 — ABOUT TOKUTI
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "About Tokuti", "Understanding the vision before we build the platform"); footer(s)

# Left – company snapshot
card(s, Inches(0.4), Inches(1.4), Inches(5.8), Inches(5.58),
     border_col=RGBColor(0xCC, 0xD8, 0xF0))
add_rect(s, Inches(0.4), Inches(1.4), Inches(5.8), Inches(0.45), VOID)
add_textbox(s, "Tokuti Snapshot", Inches(0.56), Inches(1.46),
            Inches(5.5), Inches(0.35), 13, True, ELECTRIC)
snap = [
    ("Founder",     "Mark Mariampillai, CEO — 25 yrs in finance (Credit Suisse, HSBC)"),
    ("Type",        "Web3 / PropTech Startup"),
    ("Core Product","Tokuti Auction — blockchain real estate auction platform"),
    ("Target",      "High-net-worth investors & crypto-native buyers globally"),
    ("Technology",  "Blockchain / smart contracts / tokenized real estate (RWA)"),
    ("Linked",      "GATENet — GATE token & staking ecosystem"),
    ("Market",      "Global RWA tokenization ($4T+ by 2030 per Deloitte)"),
    ("Current Site","Building on Webflow — smart MVP starting point"),
]
ty = Inches(2.02)
for k, v in snap:
    add_textbox(s, k, Inches(0.6), ty, Inches(1.55), Inches(0.28), 10, True, VOID)
    add_textbox(s, v, Inches(2.22), ty, Inches(3.82), Inches(0.28), 10, False, SILVER)
    ty += Inches(0.53)

# Right – market context
card(s, Inches(6.5), Inches(1.4), Inches(6.5), Inches(5.58),
     border_col=RGBColor(0xCC, 0xD8, 0xF0))
add_rect(s, Inches(6.5), Inches(1.4), Inches(6.5), Inches(0.45), ELECTRIC)
add_textbox(s, "The Market Opportunity", Inches(6.66), Inches(1.46),
            Inches(6.2), Inches(0.35), 13, True, VOID)
stats = [
    ("$4 Trillion+",   "Projected tokenized RWA market by 2030 (Deloitte)"),
    ("50%+ CAGR",      "Growth rate of tokenized real estate assets globally"),
    ("$10B+ (2025)",   "Tokenized real estate assets already on-chain today"),
    ("$1.4T (2026)",   "Market forecast — explosive near-term expansion"),
    ("1st-Mover",      "RWA tokenization still early — no dominant auction platform"),
]
ty = Inches(2.05)
for val, lbl in stats:
    glow_bar(s, Inches(6.66), ty + Inches(0.28), Inches(5.5), Inches(0.035), ELECTRIC)
    add_textbox(s, val, Inches(6.66), ty, Inches(2.2), Inches(0.32),
                16, True, ELECTRIC)
    add_textbox(s, lbl, Inches(8.95), ty + Inches(0.02), Inches(3.9), Inches(0.28),
                10, False, SILVER)
    ty += Inches(0.75)

add_textbox(s,
            "Why it matters: The platform Tokuti builds today must be ready to handle institutional-grade traffic, compliance and scale within 12–18 months.",
            Inches(6.66), Inches(5.72), Inches(6.1), Inches(0.52),
            10, False, VIOLET, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 4 — THE WEBFLOW CEILING
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "The Webflow Ceiling",
       "Webflow is a great start — here's exactly where it stops working for Tokuti")
footer(s)

add_textbox(s, "Webflow is perfect for validating your brand and generating early interest. "
               "But a blockchain auction platform has needs Webflow cannot meet.",
            Inches(0.55), Inches(1.35), Inches(12.3), Inches(0.42),
            11, False, SILVER, italic=True)

limits = [
    (DANGER,  "No Blockchain / Web3 Integration",
     "Webflow cannot connect to smart contracts, wallet authentication (MetaMask/WalletConnect), or on-chain auction logic. You'll need custom code for every crypto interaction — making Webflow a liability, not an asset.",
     "Critical — Tokuti's core product requires native Web3"),
    (DANGER,  "No Real-Time Auction Engine",
     "Live bidding, countdown timers synced to blockchain state, dynamic price updates — all require WebSocket connections and server-side logic that Webflow's static hosting cannot support.",
     "Critical — auction integrity depends on this"),
    (CAUTION, "Investor-Grade Security & Compliance",
     "KYC/AML flows, accredited investor verification, jurisdictional access controls, and audit trails need backend infrastructure. Webflow's form handling is not built for regulated financial transactions.",
     "High — HNW investors expect institutional security"),
    (CAUTION, "Scalability for High-Traffic Launch Events",
     "A major tokenized property auction could spike traffic 100× overnight. Webflow's CDN is optimised for marketing sites, not financial platforms under simultaneous global load.",
     "High — one bad launch = reputational damage"),
    (RGBColor(0x44, 0x88, 0xCC), "Ownership & Portability",
     "Your Webflow site lives on Webflow's infrastructure. For an investor-facing platform, you need full code ownership, custom deployment pipelines, and the ability to migrate or self-host.",
     "Medium — long-term strategic risk"),
]
ty = Inches(1.92)
for i, (col, title, desc, severity) in enumerate(limits):
    card(s, Inches(0.4), ty, Inches(12.5), Inches(0.92),
         border_col=RGBColor(0xCC, 0xD8, 0xF0))
    add_rect(s, Inches(0.4), ty, Inches(0.22), Inches(0.92), col)
    add_textbox(s, title, Inches(0.72), ty + Inches(0.08),
                Inches(7.0), Inches(0.32), 12, True, VOID)
    add_textbox(s, desc,  Inches(0.72), ty + Inches(0.44),
                Inches(9.2), Inches(0.4), 9.5, False, SILVER)
    add_textbox(s, severity, Inches(9.98), ty + Inches(0.28),
                Inches(2.8), Inches(0.32), 10, True, col, PP_ALIGN.RIGHT)
    ty += Inches(1.0)

# ════════════════════════════════════════════════════════════════
# SLIDE 5 — OUR WEB DEVELOPMENT PITCH
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Our Web Development Pitch",
       "What Genos Apollo builds for Tokuti — purpose-built for Web3 auctions")
footer(s)

# Tech stack
card(s, Inches(0.4), Inches(1.4), Inches(4.65), Inches(5.62),
     border_col=RGBColor(0xCC, 0xD8, 0xF0))
add_rect(s, Inches(0.4), Inches(1.4), Inches(4.65), Inches(0.45), VOID)
add_textbox(s, "Purpose-Built Tech Stack", Inches(0.56), Inches(1.46),
            Inches(4.4), Inches(0.35), 13, True, ELECTRIC)
stack = [
    ("Frontend",     "Next.js 15 — fast, SEO-ready, server components"),
    ("Web3",         "wagmi + viem — wallet connect, smart contract reads"),
    ("Wallet Auth",  "RainbowKit — MetaMask, WalletConnect, Coinbase"),
    ("Smart Ctrcts", "Solidity / EVM — auction logic on-chain"),
    ("Real-Time",    "WebSockets — live bids, countdowns, price feeds"),
    ("CMS",          "Sanity.io — property listings, legal docs, news"),
    ("KYC / AML",    "Sumsub / Onfido — accredited investor verification"),
    ("Payments",     "Stripe + crypto gateway (USDC / ETH / BTC)"),
    ("Auth",         "Auth0 + wallet-signature login"),
    ("Analytics",    "GA4 + Mixpanel + on-chain event tracking"),
    ("Hosting",      "Vercel + Cloudflare — global low-latency"),
]
ty = Inches(2.0)
for k, v in stack:
    add_textbox(s, k, Inches(0.6), ty, Inches(1.45), Inches(0.27), 10, True, ELECTRIC)
    add_textbox(s, v, Inches(2.1), ty, Inches(2.85), Inches(0.27), 10, False, VOID)
    ty += Inches(0.44)

# Key deliverables
card(s, Inches(5.22), Inches(1.4), Inches(7.73), Inches(5.62),
     border_col=RGBColor(0xCC, 0xD8, 0xF0))
add_rect(s, Inches(5.22), Inches(1.4), Inches(7.73), Inches(0.45), VOID)
add_textbox(s, "Key Platform Deliverables",
            Inches(5.38), Inches(1.46), Inches(7.4), Inches(0.35), 13, True, ELECTRIC)
deliverables = [
    ("Live Blockchain Auction Engine",
     "Real-time bidding synced to smart contracts. Countdown clocks, bid history, reserve price logic, automatic settlement on-chain at auction close."),
    ("Wallet-Native Investor Experience",
     "MetaMask / WalletConnect login. Portfolio dashboard showing owned tokens, auction history, upcoming drops, and live P&L."),
    ("KYC / AML Onboarding Flow",
     "Integrated Sumsub identity verification. Accredited investor status checks. Jurisdictional access controls (block non-compliant regions automatically)."),
    ("Property Listing & Tokenisation Interface",
     "Rich property pages: 3D tours, legal docs, tokenomics breakdown, projected yields, smart contract address with Etherscan link."),
    ("Admin Dashboard for Tokuti Team",
     "Create auctions, set reserve prices, approve bidders, monitor bids in real-time, trigger settlement, generate compliance reports."),
    ("Mobile-Responsive, SEO-Optimised Marketing Site",
     "Sanity-powered public site: how it works, listed properties, press, investor FAQs — fast, indexed, and brand-perfect."),
    ("Webflow → Custom Migration Path",
     "We migrate your existing Webflow content cleanly. Zero downtime. Your brand and copy preserved; the engine replaced underneath."),
]
ty = Inches(1.98)
for title, desc in deliverables:
    add_rect(s, Inches(5.38), ty + Inches(0.12), Inches(0.08), Inches(0.22), ELECTRIC)
    add_textbox(s, title, Inches(5.58), ty, Inches(7.1), Inches(0.28), 11, True, VOID)
    add_textbox(s, desc,  Inches(5.58), ty + Inches(0.3), Inches(7.1), Inches(0.38), 10, False, SILVER)
    ty += Inches(0.7)

# ════════════════════════════════════════════════════════════════
# SLIDE 6 — WHY CHOOSE GENOS APOLLO  (differentiators)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_bg(s)
header(s, "Why Genos Apollo", "What makes us the right partner for a Web3 startup pitch")
footer(s)

add_textbox(s, "We don't build generic websites. We build platforms that close deals, build investor confidence, and scale.",
            Inches(0.55), Inches(1.35), Inches(12.3), Inches(0.42),
            11.5, False, SILVER, italic=True)

diff = [
    (ELECTRIC, "Web3-Native Development Team",
     "We work with wagmi, viem, Solidity and EVM daily. No learning curve — no 'we'll figure out the blockchain part later.' Your auction logic is our core competency."),
    (PURPLE,   "Startup-Speed Delivery",
     "We run 2-week sprints with daily standups. MVP auction engine in 6–8 weeks. Full platform in 16–20 weeks. You're racing to market — we build at founder pace."),
    (NEON_GREEN,"Investor-Grade UX & Trust Design",
     "HNW investors and crypto whales judge a platform in 4 seconds. We design for the trust signals that matter: security badges, audit links, KYC flow clarity, smart contract transparency."),
    (GOLD,      "Compliance Architecture from Day 1",
     "KYC/AML, accredited investor checks, GDPR/POPIA data handling and jurisdictional access controls built in — not bolted on. Avoid costly re-architecture when regulators come knocking."),
    (VIOLET,    "Full Code Ownership — Always",
     "You own 100% of the codebase, hosted on your infrastructure. No vendor lock-in. No Webflow monthly plan holding your platform hostage. Raise, pivot, or self-host — your choice."),
]
ty = Inches(1.88)
for i, (col, title, desc) in enumerate(diff):
    card(s, Inches(0.4), ty, Inches(12.5), Inches(0.9), dark=True, border_col=col)
    glow_bar(s, Inches(0.4), ty, Inches(0.2), Inches(0.9), col)
    add_textbox(s, title, Inches(0.75), ty + Inches(0.08),
                Inches(5.0), Inches(0.32), 12, True, col)
    add_textbox(s, desc,  Inches(0.75), ty + Inches(0.44),
                Inches(11.5), Inches(0.38), 10, False, RGBColor(0xBB, 0xC4, 0xD8))
    ty += Inches(0.98)

# ════════════════════════════════════════════════════════════════
# SLIDE 7 — AI AUTOMATION OPPORTUNITIES
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "AI Automation Opportunities",
       "Where AI gives Tokuti an unfair advantage over manual processes")
footer(s)

opps = [
    (VOID, ELECTRIC, "Investor Acquisition & Onboarding", [
        "AI chatbot: answers RWA, auction, KYC FAQs 24/7",
        "Lead scoring — rank investors by wallet size & activity",
        "Auto-personalise outreach by investor type (crypto / HNW / TradFi)",
        "AI-drafted welcome sequences for newly verified investors",
        "Accreditation check automation — minutes not days",
    ]),
    (PURPLE, WHITE, "Auction Intelligence & Pricing", [
        "AI property valuation model — cross-reference comparables",
        "Dynamic reserve price recommendations per property",
        "Bidder behaviour analysis — predict auction close price",
        "Fraud / shill bidding detection (on-chain pattern analysis)",
        "Post-auction AI report — market insights for sellers",
    ]),
    (DEEP_NAVY, ELECTRIC, "Compliance & Risk Automation", [
        "Auto-trigger KYC re-verification at transaction thresholds",
        "Jurisdictional compliance check per bidder wallet address",
        "AML transaction monitoring — flag suspicious patterns",
        "Automated FICA/GDPR data retention and deletion flows",
        "Regulatory change alerts — update access rules automatically",
    ]),
    (RGBColor(0x05, 0x3A, 0x20), NEON_GREEN, "Operations & Growth", [
        "AI generates property listing copy from spec documents",
        "Social media auto-posts for new auctions (LinkedIn, X, Telegram)",
        "Smart contract deployment automation (reduce dev bottleneck)",
        "Investor newsletter — AI-curated weekly RWA market brief",
        "Meeting notes → action items (AI transcription for BD calls)",
    ]),
]
for i, (hdr_bg, hdr_txt, title, bullets) in enumerate(opps):
    c, r = i % 2, i // 2
    x = Inches(0.4)  + c * Inches(6.5)
    y = Inches(1.45) + r * Inches(2.88)
    card(s, x, y, Inches(6.2), Inches(2.72),
         border_col=RGBColor(0xCC, 0xD8, 0xF0))
    add_rect(s, x, y, Inches(6.2), Inches(0.45), hdr_bg)
    add_textbox(s, title, x + Inches(0.15), y + Inches(0.07),
                Inches(5.9), Inches(0.35), 12, True, hdr_txt)
    ty = y + Inches(0.58)
    for b in bullets:
        add_rect(s, x + Inches(0.18), ty + Inches(0.07),
                 Inches(0.1), Inches(0.1), ELECTRIC)
        add_textbox(s, b, x + Inches(0.38), ty, Inches(5.65), Inches(0.38),
                    10.5, False, VOID)
        ty += Inches(0.43)

# ════════════════════════════════════════════════════════════════
# SLIDE 8 — BENEFITS OF AUTOMATION  (hero dark)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Benefits of Automation",
       "Why automation is not optional for a Web3 auction platform")
footer(s)

for i, (val, lbl) in enumerate([
    ("10×",   "Faster investor\nonboarding vs. manual KYC"),
    ("24/7",  "AI handles enquiries\nwhile you sleep"),
    ("60%+",  "Reduction in ops\nstaff time on repetitive tasks"),
    ("Zero",  "Compliance gaps with\nautomated regulatory checks"),
]):
    stat_card_dark(s, val, lbl, Inches(0.4) + i * Inches(3.2), Inches(1.45),
                   w=Inches(3.0), h=Inches(1.65))

benefits = [
    ("Trust & Investor Confidence", [
        "Instant KYC + accredited verification = professional signal",
        "AI chatbot answers questions without 48hr email delays",
        "Automated smart contract audit trail = full transparency",
        "Compliance automation proves you take regulation seriously",
    ]),
    ("Speed to Scale", [
        "10 auctions or 10,000 — AI handles the workload either way",
        "Onboard 500 investors in a week without extra headcount",
        "Automated listing copy means new properties go live in hours",
        "Smart contract deployment automation reduces dev lag",
    ]),
    ("Revenue & Monetisation", [
        "AI lead scoring focuses BD effort on highest-value investors",
        "Personalised emails outperform generic blasts by 3–5×",
        "Automated follow-up captures deals that manual outreach misses",
        "AI market reports become a premium investor subscription product",
    ]),
    ("Competitive Moat", [
        "First RWA auction platform with AI-native investor experience",
        "Compliance automation = market access competitors can't unlock",
        "On-chain data compounding builds irreplaceable platform intelligence",
        "AI + blockchain = defensible, hard-to-replicate product layer",
    ]),
]
for i, (title, bullets) in enumerate(benefits):
    c, r = i % 2, i // 2
    x = Inches(0.4)  + c * Inches(6.5)
    y = Inches(3.28) + r * Inches(1.88)
    bullet_card(s, title, bullets, x, y, Inches(6.2), Inches(1.72),
                hdr_bg=VOID, hdr_col=ELECTRIC)

# ════════════════════════════════════════════════════════════════
# SLIDE 9 — IMPLEMENTATION PLAN
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Implementation Plan",
       "From Webflow prototype to production-grade Web3 auction platform")
footer(s)

phases = [
    ("Phase 1\nWk 1–6",   ELECTRIC,    VOID, [
        "Webflow → Next.js migration",
        "Wallet auth (RainbowKit) integration",
        "KYC/AML flow (Sumsub) live",
        "AI chatbot MVP on site",
        "Sanity CMS + property pages",
    ]),
    ("Phase 2\nWk 7–12",  PURPLE,      WHITE, [
        "Live auction engine (smart contract)",
        "Real-time bidding UI (WebSockets)",
        "Investor dashboard & portfolio",
        "Admin auction management panel",
        "Crypto + fiat payment gateway",
    ]),
    ("Phase 3\nWk 13–18", NEON_GREEN,  VOID, [
        "AI lead scoring + investor CRM",
        "AI listing copy generator",
        "Social auto-posting pipeline",
        "Compliance monitoring automation",
        "Fraud / shill bid detection AI",
    ]),
    ("Phase 4\nWk 19–24", GOLD,        VOID, [
        "AI property valuation model",
        "Investor newsletter automation",
        "Smart contract auto-deployment",
        "Market intelligence AI reports",
        "Scale testing & security audit",
    ]),
]
for i, (label, col, tc, items) in enumerate(phases):
    x = Inches(0.4) + i * Inches(3.22)
    y = Inches(1.45)
    card(s, x, y, Inches(3.05), Inches(5.62), border_col=col)
    add_rect(s, x, y, Inches(3.05), Inches(0.72), col)
    add_textbox(s, label, x, y + Inches(0.06), Inches(3.05), Inches(0.62),
                13, True, tc, PP_ALIGN.CENTER)
    if i < 3:
        add_rect(s, x + Inches(3.07), y + Inches(0.3),
                 Inches(0.14), Inches(0.06), SILVER)
    ty = y + Inches(0.9)
    for item in items:
        add_rect(s, x + Inches(0.18), ty + Inches(0.08),
                 Inches(0.1), Inches(0.1), col)
        add_textbox(s, item, x + Inches(0.38), ty, Inches(2.55), Inches(0.42),
                    10.5, False, VOID)
        ty += Inches(0.84)

add_textbox(s, "MVP auction engine ready in ~12 weeks  |  Full platform + AI suite live in ~24 weeks",
            Inches(0.4), Inches(7.07), Inches(12.5), Inches(0.28),
            10.5, True, VOID, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 10 — INVESTMENT / PRICING
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_bg(s)
header(s, "Investment & Pricing",
       "Transparent, milestone-based — built for a startup's runway")
footer(s)

# Tiers side by side
tiers = [
    ("STARTER",  ELECTRIC, [
        ("Webflow → Next.js migration",   "✓"),
        ("Wallet auth (RainbowKit)",       "✓"),
        ("KYC/AML onboarding (Sumsub)",   "✓"),
        ("Sanity CMS + property pages",   "✓"),
        ("AI chatbot (FAQ bot)",          "✓"),
        ("Auction engine (smart ctrct)",  "✗"),
        ("Live bidding UI (WebSocket)",   "✗"),
        ("Investor dashboard",            "✗"),
        ("AI lead scoring + CRM",         "✗"),
        ("AI listing copy generator",     "✗"),
    ], "$12 000 – $18 000", "~6 weeks"),
    ("GROWTH",   PURPLE, [
        ("Webflow → Next.js migration",   "✓"),
        ("Wallet auth (RainbowKit)",       "✓"),
        ("KYC/AML onboarding (Sumsub)",   "✓"),
        ("Sanity CMS + property pages",   "✓"),
        ("AI chatbot (FAQ bot)",          "✓"),
        ("Auction engine (smart ctrct)",  "✓"),
        ("Live bidding UI (WebSocket)",   "✓"),
        ("Investor dashboard",            "✓"),
        ("AI lead scoring + CRM",         "✗"),
        ("AI listing copy generator",     "✗"),
    ], "$28 000 – $40 000", "~12 weeks"),
    ("FULL PLATFORM", NEON_GREEN, [
        ("Webflow → Next.js migration",   "✓"),
        ("Wallet auth (RainbowKit)",       "✓"),
        ("KYC/AML onboarding (Sumsub)",   "✓"),
        ("Sanity CMS + property pages",   "✓"),
        ("AI chatbot (FAQ bot)",          "✓"),
        ("Auction engine (smart ctrct)",  "✓"),
        ("Live bidding UI (WebSocket)",   "✓"),
        ("Investor dashboard",            "✓"),
        ("AI lead scoring + CRM",         "✓"),
        ("AI listing copy generator",     "✓"),
    ], "$55 000 – $80 000", "~24 weeks"),
]
for i, (tier, col, features, price, timeline) in enumerate(tiers):
    x = Inches(0.5) + i * Inches(4.28)
    y = Inches(1.45)
    card(s, x, y, Inches(4.05), Inches(5.62), dark=True, border_col=col)
    add_rect(s, x, y, Inches(4.05), Inches(0.5), col)
    tc = VOID if col == NEON_GREEN else WHITE
    add_textbox(s, tier, x, y + Inches(0.07), Inches(4.05), Inches(0.38),
                14, True, tc, PP_ALIGN.CENTER)
    ty = y + Inches(0.62)
    for feat, included in features:
        ic = col if included == "✓" else SILVER
        add_textbox(s, f"{included}  {feat}", x + Inches(0.18), ty,
                    Inches(3.7), Inches(0.28), 9.5,
                    included == "✓", ic)
        ty += Inches(0.36)
    glow_bar(s, x, ty + Inches(0.08), Inches(4.05), Inches(0.04), col)
    add_textbox(s, price, x, ty + Inches(0.18), Inches(4.05), Inches(0.4),
                17, True, col, PP_ALIGN.CENTER)
    add_textbox(s, timeline, x, ty + Inches(0.56), Inches(4.05), Inches(0.28),
                10, False, SILVER, PP_ALIGN.CENTER, italic=True)

add_textbox(s,
            "All tiers: milestone billing  ·  full code ownership  ·  monthly retainer available post-launch  ·  no vendor lock-in",
            Inches(0.4), Inches(7.07), Inches(12.5), Inches(0.28),
            9.5, False, SILVER, PP_ALIGN.CENTER, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 11 — NEXT STEPS
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Next Steps", "How we start — fast, low-friction, no obligation")
footer(s)

steps = [
    ("Day 1",   ELECTRIC, "Discovery Call",
     "30-minute call with Mark. We review the current Webflow build, audit the smart contract architecture, and align on Phase 1 scope and timeline."),
    ("Day 3",   PURPLE,   "Technical Scoping",
     "We deliver a detailed technical spec: chain selection, wallet auth approach, KYC provider, API architecture, and a refined project plan."),
    ("Day 7",   NEON_GREEN,"Proposal Accepted",
     "Sign engagement letter. Milestone 1 deposit triggered. Sprint 1 begins. Your Webflow site stays live — zero downtime during migration."),
    ("Week 6",  GOLD,     "MVP Live",
     "Starter tier: Next.js site with wallet login, KYC flow, AI chatbot and property pages deployed. Ready for your first investor beta audience."),
]
for i, (when, col, title, desc) in enumerate(steps):
    x = Inches(0.4) + i * Inches(3.22)
    y = Inches(1.5)
    card(s, x, y, Inches(3.05), Inches(4.1), border_col=col)
    add_rect(s, x, y, Inches(3.05), Inches(0.45), col)
    tc = VOID if col in (ELECTRIC, NEON_GREEN, GOLD) else WHITE
    add_textbox(s, when, x, y + Inches(0.06), Inches(3.05), Inches(0.35),
                12, True, tc, PP_ALIGN.CENTER)
    add_textbox(s, title, x + Inches(0.15), y + Inches(0.6),
                Inches(2.75), Inches(0.36), 12, True, col)
    add_textbox(s, desc, x + Inches(0.15), y + Inches(1.05),
                Inches(2.75), Inches(2.9), 10.5, False, SILVER)

card(s, Inches(0.4), Inches(5.82), Inches(12.5), Inches(1.1),
     fill=VOID, border_col=ELECTRIC)
glow_bar(s, Inches(0.4), Inches(5.82), Inches(12.5), Inches(0.055), ELECTRIC)
add_textbox(s, "Ready to build the platform Tokuti deserves — not just a Webflow placeholder?",
            Inches(0.65), Inches(5.92), Inches(8.5), Inches(0.4),
            14, True, WHITE)
add_textbox(s, "Book your discovery call — we move at startup speed.",
            Inches(0.65), Inches(6.32), Inches(8.5), Inches(0.38),
            11, False, ELECTRIC)
add_textbox(s, "tokuti.io  •  mark@tokuti.io",
            Inches(9.2), Inches(6.1), Inches(3.8), Inches(0.35),
            10, False, ELECTRIC, PP_ALIGN.RIGHT, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 12 — BACK COVER
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_bg(s)
add_rect(s, 0, 0, Inches(0.72), SLIDE_H, ELECTRIC)
add_rect(s, Inches(0.72), Inches(3.55),
         SLIDE_W - Inches(0.72), Inches(0.055), ELECTRIC)

add_textbox(s, "Let's Build\nTokuti.",
            Inches(1.1), Inches(1.65), Inches(11), Inches(1.8),
            52, True, WHITE)
add_textbox(s, "The future of real estate investment is on-chain. The platform that captures it should be built right — from day one.",
            Inches(1.1), Inches(3.22), Inches(10.5), Inches(0.6),
            14, False, ELECTRIC, italic=True)
add_textbox(s, "Genos Apollo  ·  Web3 & AI Development  ·  2026",
            Inches(1.1), Inches(4.1), Inches(9), Inches(0.4),
            12, False, RGBColor(0x55, 0x66, 0x88))
add_textbox(s, "Prepared for Mark Mariampillai, CEO & Founder — Tokuti.io  |  Confidential",
            Inches(1.1), Inches(6.7), Inches(10), Inches(0.35),
            10, False, SILVER, italic=True)

# ── Save ────────────────────────────────────────────────────────
out = r"c:\Users\deevansho\Desktop\GenosApollp clients\Tokuti_WebDev_Automation_Proposal.pptx"
prs.save(out)
print(f"Saved: {out}")
