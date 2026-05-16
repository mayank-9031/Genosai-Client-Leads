import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.parent / "utils"))
from ppt_theme import (
    NAVY, DARK_BLUE, ACCENT, GOLD, RED, GREEN, PURPLE, ORANGE, TEAL, CYAN,
    WHITE, LGREY, MGREY, W, H,
    bg, rect, tb, aline, bullets, tag,
    slide_header, slide_footer, score_badge,
    stat_cards, top_bar_card, phase_card, before_after, genos_logo_block,
)

ROOT = Path(__file__).parent.parent.parent
prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]

# SLIDE 1 - COVER
s1 = prs.slides.add_slide(BLANK)
bg(s1)
rect(s1, 0, 0, Inches(0.12), H, GOLD)
rect(s1, Inches(8.5), 0, Inches(4.83), Inches(2.2), DARK_BLUE)
tb(s1, "GENOS AI", Inches(9.2), Inches(0.4), Inches(3.5), Inches(0.8),
   size=28, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
tb(s1, "www.genosai.tech", Inches(9.2), Inches(1.05), Inches(3.5), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
tb(s1, "COSTAR REAL ESTATE MANAGER", Inches(0.55), Inches(1.6),
   Inches(9), Inches(0.6), size=13, bold=False, color=GOLD)
tb(s1, "Website & AI Growth\nAudit Report", Inches(0.55), Inches(2.1),
   Inches(9), Inches(1.8), size=48, bold=True, color=WHITE)
aline(s1, Inches(3.85), color=GOLD, w=Inches(3))
tb(s1, "Prepared exclusively for Mark McDonald, President",
   Inches(0.55), Inches(4.15), Inches(9), Inches(0.5), size=14, color=LGREY)
rect(s1, 0, Inches(6.5), W, Inches(1.0), DARK_BLUE)
tb(s1, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699",
   Inches(0.55), Inches(6.6), Inches(9), Inches(0.5), size=12, color=MGREY)
tb(s1, "May 2026  |  CONFIDENTIAL", Inches(9.5), Inches(6.6), Inches(3.3), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
rect(s1, Inches(10.5), Inches(3.3), Inches(2.3), Inches(2.5), DARK_BLUE)
tb(s1, "AUDIT SCORE", Inches(10.55), Inches(3.4), Inches(2.2), Inches(0.5),
   size=11, bold=True, color=MGREY, align=PP_ALIGN.CENTER)
tb(s1, "5", Inches(10.55), Inches(3.75), Inches(2.2), Inches(1.1),
   size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
tb(s1, "/ 10", Inches(10.55), Inches(4.7), Inches(2.2), Inches(0.4),
   size=18, color=LGREY, align=PP_ALIGN.CENTER)
tag(s1, "HIGH PRIORITY", Inches(10.75), Inches(5.1), fill=RED, size=10)


# SLIDE 2 - AT A GLANCE
s2 = prs.slides.add_slide(BLANK)
bg(s2)
rect(s2, 0, 0, Inches(0.12), H, GOLD)
tb(s2, "AT A GLANCE", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GOLD)
tb(s2, "Mark McDonald - Who He Is",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s2, Inches(1.35), color=GOLD)

stats = [
    ("15 yrs", "Tenure at CoStar / Virtual Premise"),
    ("Fortune\n500", "Customer Portfolio"),
    ("ASC 842\nIFRS 16", "Compliance Standards Served"),
    ("2", "Standards: Lease Acctg + Admin"),
    ("Georgia\nTech+MIT", "Engineering + Sloan MBA"),
]
cw = Inches(2.3); ch = Inches(2.2); gap = Inches(0.18)
for i, (num, label) in enumerate(stats):
    l = Inches(0.55) + i * (cw + gap)
    rect(s2, l, Inches(1.8), cw, ch, DARK_BLUE)
    tb(s2, num, l + Inches(0.1), Inches(1.95), cw - Inches(0.2), Inches(1.1),
       size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s2, label, l + Inches(0.1), Inches(3.05), cw - Inches(0.2), Inches(0.7),
       size=12, color=LGREY, align=PP_ALIGN.CENTER)

rect(s2, Inches(0.55), Inches(4.2), Inches(12.23), Inches(0.6), DARK_BLUE)
tb(s2, "President, CoStar Real Estate Manager  -  Formerly Virtual Premise  -  Also oversees Visual Lease  -  Speaker: CoreNet Global, NRTA  -  Atlanta, GA",
   Inches(0.65), Inches(4.28), Inches(12), Inches(0.45), size=12, color=LGREY, align=PP_ALIGN.CENTER)

tb(s2, "MARK'S STRATEGIC FOCUS", Inches(0.55), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=GOLD)
bullets(s2, ["Profitable growth of CoStar Real Estate Manager customer base",
             "Product leadership in lease administration and lease accounting markets",
             "Guiding Fortune 500 companies through digital real estate transformation"],
        Inches(0.55), Inches(5.35), Inches(5.8), Inches(1.5), size=13, color=LGREY)

tb(s2, "WHAT THE PLATFORM COVERS", Inches(6.8), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=GOLD)
bullets(s2, ["Lease accounting (ASC 842, IFRS 16, FRS 102 compliance)",
             "Lease administration with integrated CoStar market data and analytics",
             "Transaction management, ESG, Visual Lease capabilities"],
        Inches(6.8), Inches(5.35), Inches(5.8), Inches(1.5), size=13, color=LGREY)


# SLIDE 3 - WHAT'S WORKING
s3 = prs.slides.add_slide(BLANK)
bg(s3)
rect(s3, 0, 0, Inches(0.12), H, GREEN)
tb(s3, "WHAT'S WORKING", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s3, "Strengths to Build On",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s3, Inches(1.35), color=GREEN)

strengths = [
    ("The CoStar Data Moat Is a Genuine Differentiator",
     "No other lease management platform ships with CoStar's proprietary commercial market data and analytics integrated at the platform level. That's not a feature competitors can replicate by Q3. It's an institutional advantage that changes the ROI story for any corporate real estate team evaluating the market - and the website largely buries it below the fold."),
    ("15 Years of Institutional Depth From Startup Through Enterprise Scale",
     "Mark joined Virtual Premise during the startup phase, led hyper-growth field sales, built strategy and customer success, and now runs the platform as President inside CoStar Group. That tenure is rare in SaaS - it means the implementation playbook, the customer escalation patterns, and the compliance edge cases are baked into the product, not assembled from consulting engagements."),
    ("Visual Lease Acquisition Expands the Addressable Market",
     "Adding Visual Lease deepens accounting compliance specialisation and widens the buyer pool to mid-market teams that need a standalone lease accounting answer rather than a full portfolio management platform. The acquisition is a meaningful growth move and the site could surface what it means for buyers more clearly."),
    ("Compliance-Driven Demand Creates Real Buying Windows",
     "ASC 842 and IFRS 16 mandates create finite deadline pressure that most SaaS categories don't get. Companies coming into scope, changing auditors, or hitting contract renewal periods have a specific urgency that makes the buying decision time-bound. That demand window is structural and predictable - it just needs a website that captures it."),
    ("Thought Leadership Presence Across Publishing and Speaking",
     "Contributing to Accounting Today and Propmodo, speaking at CoreNet Global and NRTA, running the Lease Alert podcast - the content engine exists. The challenge is that the website doesn't convert that authority into pipeline. Prospects arrive from thought leadership and hit a form-gated demo with no intermediate engagement step."),
]
for i, (title, desc) in enumerate(strengths):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s3, Inches(0.55), t, Inches(0.06), Inches(0.7), GREEN)
    tb(s3, title, Inches(0.75), t, Inches(11.8), Inches(0.38), size=14, bold=True, color=WHITE)
    tb(s3, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# SLIDE 4 - WHERE LEADS ARE LEAKING
s4 = prs.slides.add_slide(BLANK)
bg(s4)
rect(s4, 0, 0, Inches(0.12), H, RED)
tb(s4, "WHERE LEADS ARE LEAKING", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=RED)
tb(s4, "The Compliance Buyer at 11pm Hits a Form and Moves On.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=28, bold=True, color=WHITE)
aline(s4, Inches(1.35), color=RED)

issues = [
    ("NO PRICING PAGE - DEMO REQUIRED TO LEARN ANY COST INFORMATION", "HIGH",
     "A VP of Corporate Real Estate evaluating lease management platforms for a 200-location portfolio can't answer the CFO's budget question without a demo call. LeaseQuery shows pricing tiers on site. LeaseAccelerator has ROI framing. CoStar Real Estate Manager asks the prospect to fill out a form and wait. In competitive shortlisting, the vendor that removes friction wins the first meeting."),
    ("NO ROI CALCULATOR DESPITE SAVINGS BEING A CORE VALUE PROPOSITION", "HIGH",
     "The platform's core pitch is compliance savings, audit cycle reduction, and lease overpayment prevention. A compliance buyer has a specific portfolio with a measurable current cost - whether that's spreadsheet hours, auditor fees, or overpaid rent. An ROI calculator that takes portfolio size and current workflow and returns a savings estimate pre-qualifies the lead by deal size and hands the prospect a number to take to the CFO before the first call."),
    ("DEMO FULLY FORM-GATED WITH NO INTERACTIVE PRODUCT TOUR", "HIGH",
     "A corporate real estate director doing platform research at 11pm can read about unified lease accounting but can't see it working. No product walkthrough, no sandbox environment, no embedded video of the lease accounting module pulling CoStar market data in real time. Buyers who can't evaluate the product before a call are lower-quality meetings - and many don't book at all."),
    ("NO AI CHAT FOR 24/7 PROSPECT OR CUSTOMER QUESTIONS", "MEDIUM",
     "Compliance questions arrive outside business hours. A buyer researching ASC 842 requirements in January before a February audit has questions at 9pm that the site cannot answer. Current customers with implementation questions wait for support response. AI chat trained on the platform, compliance standards, and integration landscape captures the prospect before they move to the next search result."),
    ("NO ROLE-SPECIFIC LANDING PAGES FOR DIFFERENT BUYING PERSONAS", "HIGH",
     "Accounting teams buy for ASC 842 audit close and financial reporting speed. Real estate teams buy for portfolio visibility and transaction workflow. Retail tenants buy for lease optimisation and renewal strategy. One homepage does not convert all three personas. Each needs a different ROI story, a different urgency frame, and a different entry point - none of which exist today on costarmanager.com."),
]
for i, (title, pri, desc) in enumerate(issues):
    t = Inches(1.7) + i * Inches(1.05)
    pc = RED if pri == "HIGH" else GOLD
    rect(s4, Inches(0.55), t, Inches(0.06), Inches(0.7), pc)
    tb(s4, title, Inches(0.75), t, Inches(9.2), Inches(0.38), size=12, bold=True, color=WHITE)
    tag(s4, pri, Inches(10.1), t + Inches(0.05), fill=pc, size=9)
    tb(s4, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# SLIDE 5 - WEBSITE DESIGN AUDIT
s5 = prs.slides.add_slide(BLANK)
bg(s5)
rect(s5, 0, 0, Inches(0.12), H, PURPLE)
tb(s5, "WEBSITE DESIGN AUDIT", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=PURPLE)
tb(s5, "15 Years of Depth. The Site Doesn't Show the ROI.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.8), size=32, bold=True, color=WHITE)
aline(s5, Inches(1.45), color=PURPLE, w=Inches(1.5))

rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(3.6), DARK_BLUE)
rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(0.42), RED)
tb(s5, "COSTARMANAGER.COM TODAY", Inches(0.65), Inches(1.92), Inches(3.6), Inches(0.38),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "No pricing page - form gate to discover cost",
    "No ROI calculator or savings estimator",
    "Demo fully gated, no interactive product tour",
    "No AI chat / 24/7 prospect self-service",
    "No role-specific pages (accounting vs. RE vs. retail)",
    "Testimonials have no quantified results",
    "No analyst recognition or customer count surfaced",
    "No compliance nurture or re-engagement sequences",
], Inches(0.7), Inches(2.42), Inches(3.6), Inches(2.9), size=11, color=LGREY, dot=RED)

tb(s5, "->", Inches(4.45), Inches(3.6), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(3.6), DARK_BLUE)
rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(0.42), GOLD)
tb(s5, "WHAT TOP B2B SaaS SITES DO", Inches(5.1), Inches(1.92), Inches(3.4), Inches(0.38),
   size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Pricing page or ROI frame visible on homepage",
    "Interactive product tour without demo gate",
    "Role-specific landing pages by buying persona",
    "AI chat for 24/7 prospect and customer routing",
    "Quantified case studies with before/after metrics",
    "Customer count or trusted-by statement on homepage",
    "Analyst recognition (Gartner, Forrester) surfaced",
    "Compliance deadline nurture sequences active",
], Inches(5.1), Inches(2.42), Inches(3.4), Inches(2.9), size=11, color=LGREY, dot=GOLD)

tb(s5, "->", Inches(8.68), Inches(3.6), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(3.6), DARK_BLUE)
rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(0.42), GREEN)
tb(s5, "WHAT GENOS AI BUILDS", Inches(9.35), Inches(1.92), Inches(3.4), Inches(0.38),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "ROI calculator as homepage CTA by portfolio size",
    "Interactive product tour: lease accounting + CoStar data",
    "Role-specific pages: accounting, RE team, retail tenant",
    "AI chat trained on platform, ASC 842, integrations",
    "Quantified case studies: before/after metrics per client",
    "Compliance nurture sequences by ASC 842 stage",
    "Re-engagement for 2022-2023 adopters at renewal window",
    "Awards and data moat surfaced as homepage differentiators",
], Inches(9.35), Inches(2.42), Inches(3.4), Inches(2.9), size=11, color=LGREY, dot=GREEN)

rect(s5, Inches(0.55), Inches(5.72), Inches(12.23), Inches(1.28), DARK_BLUE)
tb(s5, "WHY IT MATTERS:", Inches(0.7), Inches(5.78), Inches(2.2), Inches(0.38),
   size=11, bold=True, color=PURPLE)
stats2 = [
    ("3-4x", "demo booking rate\nwith interactive product tour"),
    ("60-70%", "prospect questions\nanswerable by AI chat"),
    ("3 personas", "different ROI stories\none homepage loses all three"),
    ("2022-23", "ASC 842 adopters now\nat contract renewal window"),
]
sw = Inches(2.7); sl = Inches(2.9)
for i, (num, label) in enumerate(stats2):
    xl = sl + i * sw
    tb(s5, num, xl, Inches(5.75), Inches(2.5), Inches(0.55),
       size=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s5, label, xl, Inches(6.22), Inches(2.5), Inches(0.7),
       size=10, color=LGREY, align=PP_ALIGN.CENTER)


# SLIDE 6 - AI AUTOMATION OPPORTUNITIES
s6 = prs.slides.add_slide(BLANK)
bg(s6)
rect(s6, 0, 0, Inches(0.12), H, ACCENT)
tb(s6, "AI AUTOMATION OPPORTUNITIES", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s6, "Turning Compliance Pressure Into Inbound Pipeline.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s6, Inches(1.35), color=ACCENT)

opps = [
    ("ROI Calculator\n+ Lead Capture", GOLD,
     "Prospect inputs portfolio size, lease count, and current compliance workflow. Calculator returns estimated savings vs. manual process and ASC 842 audit cycle advantage. Captures lead with deal-size pre-qualification before any sales conversation starts. CFO-ready number that removes the pricing friction without exposing list rates."),
    ("AI Chat\n24/7 Self-Service", ACCENT,
     "Trained on CoStar Real Estate Manager's platform, ASC 842/IFRS 16 requirements, integration landscape, and implementation timeline. Handles prospect questions at 11pm, books demo calls, routes by accounting team vs. real estate team vs. retail tenant persona, handles current customer questions without waiting for support queue."),
    ("Interactive\nProduct Tour", GREEN,
     "Lease accounting module walkthrough, CoStar market data layer in action, Visual Lease capabilities demonstrated - all without a form gate. A compliance buyer researching at night can evaluate the platform before committing to a 45-minute sales call. Higher-intent demos, shorter sales cycles, fewer no-show meetings."),
    ("Compliance\nNurture Sequences", RED,
     "Prospects by ASC 842 compliance stage: research phase, evaluation, adoption, post-adoption optimisation. Each stage gets a different content sequence - regulatory guides, implementation timelines, customer case studies with audit cycle metrics. Stays in front of deadline-pressured buyers through their entire evaluation window."),
    ("Re-Engagement:\n2022-23 Renewal\nCohort", PURPLE,
     "Companies that adopted compliance tools in 2022 and 2023 are approaching contract renewal. Automated sequences targeting that cohort - Visual Lease acquisition value, new CoStar data capabilities, competitive positioning - captures the re-evaluation window before LeaseQuery or a competitor makes the case first."),
]
cw2 = Inches(2.35); gap2 = Inches(0.14)
for i, (title, color, desc) in enumerate(opps):
    l2 = Inches(0.55) + i * (cw2 + gap2)
    rect(s6, l2, Inches(1.8), cw2, Inches(0.06), color)
    rect(s6, l2, Inches(1.86), cw2, Inches(4.6), DARK_BLUE)
    tb(s6, title, l2 + Inches(0.12), Inches(1.98), cw2 - Inches(0.24), Inches(0.65),
       size=13, bold=True, color=color)
    tb(s6, desc, l2 + Inches(0.12), Inches(2.65), cw2 - Inches(0.24), Inches(3.6),
       size=11, color=LGREY)

rect(s6, Inches(0.55), Inches(6.7), Inches(12.23), Inches(0.45), DARK_BLUE)
tb(s6, "15 years of implementation depth + CoStar data moat + Fortune 500 references. The website should close the trust gap before the first call.",
   Inches(0.7), Inches(6.75), Inches(12), Inches(0.35), size=11, italic=True, color=GOLD)


# SLIDE 7 - WHAT THIS IS COSTING YOU
s7 = prs.slides.add_slide(BLANK)
bg(s7)
rect(s7, 0, 0, Inches(0.12), H, RED)
tb(s7, "WHAT THIS IS COSTING YOU", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=RED)
tb(s7, "LeaseQuery Shows Pricing. CoStar Real Estate Manager Shows a Form.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=28, bold=True, color=WHITE)
aline(s7, Inches(1.35), color=RED)

costs = [
    ("Form-gated demos mean competitors win the shortlist before the first call",
     "A VP of Corporate Real Estate shortlisting lease management platforms in January gives 20 minutes to each vendor site before deciding who to call. The vendor with an ROI calculator and a pricing frame gets the call. The vendor with a contact form gets deprioritised. The CoStar data moat and Fortune 500 track record don't get a chance to land - the prospect moved on before reading them."),
    ("No pricing visibility means the CFO question stays unanswered until a call",
     "Corporate buying teams need to answer a budget question before they can book a demo. If the answer requires a discovery call, the evaluation stalls or the CFO approves a competitor whose pricing was visible on site. A pricing page with portfolio-size framing doesn't require publishing list rates - it gives the prospect enough to build a business case and justify the conversation internally."),
    ("No product tour means evaluators can't assess the CoStar data integration before investing in a meeting",
     "The CoStar market data integration is the feature no competitor can replicate. A corporate real estate director who sees it working - lease administration data mapped against live CoStar market comps in a 3-minute product tour - is a different quality of meeting than one who read a bullet point about it on a homepage. The tour converts the data moat into a tangible advantage before the first call."),
    ("One homepage trying to convert three different buying personas loses all three",
     "An accounting team controller focused on ASC 842 audit close, a VP Real Estate optimising a 500-lease portfolio, and a retail tenant CFO negotiating renewals have different urgencies, different objections, and different ROI definitions. A single homepage positioned at the generic middle converts none of them at the rate that persona-specific pages would."),
    ("Testimonials without metrics don't build the trust they should",
     "The Hertz reference and Allegion reference are credible names. John Grotto's quote about doubling down on CoStar is a positive signal. But a compliance buyer asking their CFO for approval needs a before-and-after: how many hours did month-end close take before, how many after. What was the audit risk exposure, and what is it now. Every named Fortune 500 reference that lacks quantified outcomes is a missed conversion opportunity."),
]
for i, (title, desc) in enumerate(costs):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s7, Inches(0.55), t, Inches(0.06), Inches(0.7), RED)
    tb(s7, title, Inches(0.75), t, Inches(11.8), Inches(0.38), size=14, bold=True, color=WHITE)
    tb(s7, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# SLIDE 8 - WHAT WE'D BUILD
s8 = prs.slides.add_slide(BLANK)
bg(s8)
rect(s8, 0, 0, Inches(0.12), H, GREEN)
tb(s8, "WHAT WE'D BUILD", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s8, "Genos AI for CoStar Real Estate Manager",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s8, Inches(1.35), color=GREEN)

builds = [
    ("ROI Calculator as Homepage CTA - Savings Before the Demo",
     "Prospect inputs portfolio size, lease count, and current compliance workflow. Calculator returns estimated savings vs. manual process, ASC 842 audit cycle advantage, and estimated time to compliance. Lead captured with deal-size pre-qualification. CFO-ready number generated before a single sales call. Every LeaseQuery evaluation that starts on costarmanager.com ends differently."),
    ("AI Chat - 24/7 Prospect and Customer Self-Service",
     "Trained on CoStar Real Estate Manager's platform features, ASC 842/IFRS 16/FRS 102 requirements, integration landscape, Visual Lease capabilities, and implementation timeline. Routes by accounting team, real estate team, and retail tenant persona. Books demo calls after hours. Handles current customer implementation questions without adding to the support queue."),
    ("Interactive Product Tour - CoStar Data Moat Made Visible",
     "Lease accounting module demonstration, CoStar market data layer live on a sample portfolio, Visual Lease compliance workflow shown end-to-end. No form gate. A compliance buyer evaluating at 11pm sees the platform working before committing to a meeting. Higher-intent demos, shorter sales cycles, fewer no-shows from buyers who couldn't pre-evaluate."),
    ("Role-Specific Landing Pages by Buying Persona",
     "Accounting team page: ASC 842 audit close speed, financial reporting accuracy, audit trail automation. Real estate team page: portfolio visibility, transaction management, CoStar market data for renewal negotiation. Retail tenant page: lease optimisation, renewal strategy, overpayment prevention. Each persona gets the ROI story that maps to their specific problem."),
    ("Quantified Case Studies + Compliance Nurture + Renewal Re-Engagement",
     "Every Fortune 500 reference gets a before-and-after metric (month-end close hours, audit risk exposure, lease savings). Compliance nurture sequences routing by ASC 842 stage keep CoStar Real Estate Manager in front of deadline-pressured buyers through their full evaluation window. Re-engagement sequences for 2022-2023 adopters hit the renewal decision before competitors do."),
]
for i, (title, desc) in enumerate(builds):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s8, Inches(0.55), t, Inches(0.06), Inches(0.7), GREEN)
    tb(s8, title, Inches(0.75), t, Inches(11.8), Inches(0.38), size=14, bold=True, color=WHITE)
    tb(s8, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# SLIDE 9 - 90-DAY OUTCOMES
s9 = prs.slides.add_slide(BLANK)
bg(s9)
rect(s9, 0, 0, Inches(0.12), H, ACCENT)
tb(s9, "90-DAY OUTCOMES", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s9, "What Changes in 90 Days",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s9, Inches(1.35))

phases = [
    ("Month 1", GOLD, [
        "ROI calculator live on costarmanager.com homepage and pricing entry point",
        "AI chat deployed for 24/7 prospect routing, demo booking, and persona identification",
        "Interactive product tour launched: lease accounting + CoStar market data walkthrough",
        "Site rebuild begins with role-specific landing pages for accounting, RE team, retail tenant",
    ]),
    ("Month 2", ACCENT, [
        "Accounting team, real estate team, and retail tenant landing pages live with persona ROI stories",
        "Quantified case studies published for Hertz, Allegion, and top Fortune 500 references",
        "Compliance nurture sequences active by ASC 842 stage: research, evaluation, adoption",
        "First measurable inbound demo requests from ROI calculator before sales call",
    ]),
    ("Month 3", GREEN, [
        "Re-engagement sequences running for 2022-2023 ASC 842 adopters at contract renewal window",
        "Buyer alert sequences for companies newly coming into ASC 842/IFRS 16 compliance scope",
        "CoStar data moat and Visual Lease acquisition surfaced as homepage differentiators",
        "Measurable shortlist capture from compliance-deadline buyers who previously bounced on form",
    ]),
]
pw = Inches(3.9); gap3 = Inches(0.18)
for i, (month, color, pts) in enumerate(phases):
    l3 = Inches(0.55) + i * (pw + gap3)
    rect(s9, l3, Inches(1.8), pw, Inches(0.42), color)
    tb(s9, month.upper(), l3 + Inches(0.1), Inches(1.84), pw - Inches(0.2), Inches(0.36),
       size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    rect(s9, l3, Inches(2.22), pw, Inches(4.7), DARK_BLUE)
    bullets(s9, pts, l3 + Inches(0.15), Inches(2.35), pw - Inches(0.3), Inches(4.4),
            size=12, color=LGREY, dot=color)


# SLIDE 10 - NEXT STEPS
s10 = prs.slides.add_slide(BLANK)
bg(s10)
rect(s10, 0, 0, Inches(0.12), H, GOLD)
rect(s10, Inches(7.5), 0, Inches(5.83), H, DARK_BLUE)
tb(s10, "NEXT STEPS", Inches(0.55), Inches(0.3), Inches(6.5), Inches(0.5),
   size=11, bold=True, color=GOLD)
tb(s10, "Let's Talk for\n15 Minutes.",
   Inches(0.55), Inches(0.8), Inches(6.8), Inches(2.0), size=48, bold=True, color=WHITE)
aline(s10, Inches(2.8), color=GOLD, w=Inches(2.5))
tb(s10, "No pitch deck required. Reply or book a call directly.\nWe'll walk through what fits CoStar Real Estate Manager.",
   Inches(0.55), Inches(3.05), Inches(6.5), Inches(0.8), size=14, color=LGREY)

steps = [
    "Reply to this email - we'll set up a time",
    "Book directly: www.genosai.tech/call",
    "Walk through the audit findings in 15 minutes",
    "You decide what, if anything, makes sense to move on",
]
bullets(s10, steps, Inches(0.55), Inches(3.95), Inches(6.5), Inches(2.0), size=14, color=WHITE, dot=GOLD)

rect(s10, 0, Inches(6.5), Inches(7.5), Inches(1.0), DARK_BLUE)
tb(s10, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699  |  www.genosai.tech",
   Inches(0.55), Inches(6.6), Inches(6.8), Inches(0.5), size=12, color=MGREY)

tb(s10, "WHAT WE BUILD FOR\nCOSTAR REAL ESTATE MGR", Inches(7.85), Inches(0.6), Inches(5.0), Inches(0.9),
   size=16, bold=True, color=GOLD)
aline(s10, Inches(1.45), color=GOLD)

summary_items = [
    ("ROI Calc", "Portfolio inputs return savings + ASC 842 audit advantage"),
    ("AI Chat", "24/7 prospect/customer routing and demo booking"),
    ("Product Tour", "Lease accounting + CoStar data moat - no form gate"),
    ("Role Pages", "Accounting / RE team / retail tenant persona pages"),
    ("Case Studies", "Before/after metrics on every Fortune 500 reference"),
    ("Nurture", "Compliance stage sequences + 2022-23 renewal re-engagement"),
]
for i, (label, val) in enumerate(summary_items):
    t2 = Inches(1.6) + i * Inches(0.88)
    rect(s10, Inches(7.85), t2, Inches(1.5), Inches(0.72), NAVY)
    tb(s10, label, Inches(7.9), t2 + Inches(0.05), Inches(1.4), Inches(0.32),
       size=10, bold=True, color=GOLD)
    tb(s10, val, Inches(9.5), t2 + Inches(0.05), Inches(3.6), Inches(0.65),
       size=11, color=LGREY)


prs.save(ROOT / "output" / "decks" / "CoStarManager_Audit_GenosAI.pptx")
print("Saved: CoStarManager_Audit_GenosAI.pptx")
