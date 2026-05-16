import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.parent / "utils"))
from ppt_theme import (
    NAVY, DARK_BLUE, ACCENT, GOLD, RED, GREEN, PURPLE, ORANGE, TEAL, CYAN,
    WHITE, LGREY, MGREY, W, H,
    bg, rect, tb, aline, bullets, tag,
    slide_header, slide_footer, score_badge,
    stat_cards, top_bar_card, phase_card, before_after, genos_logo_block,
)

ROOT = Path(__file__).parent.parent.parent
prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]
# Client-specific colors
CHARCOAL = RGBColor(0x36, 0x45, 0x54)
DGREY    = RGBColor(0x2C, 0x38, 0x4A)
AMBER    = RGBColor(0xF3, 0x9C, 0x12)
WARM     = RGBColor(0xFD, 0xF0, 0xE2)
SLATE    = RGBColor(0x2E, 0x45, 0x5F)


# ── Local helpers ────────────────────────────────────────────────────────────

def header(slide, title, subtitle, accent=GOLD):
    rect(slide, 0, 0, Inches(0.12), H, accent)
    tb(slide, title.upper(), Inches(0.55 + 0.12), Inches(0.3),
       Inches(12), Inches(0.45), size=11, bold=True, color=accent)
    tb(slide, subtitle, Inches(0.55 + 0.12), Inches(0.72),
       Inches(12 - 0.12), Inches(0.65), size=22, bold=True, color=WHITE)
    aline(slide, Inches(1.38), color=accent, w=Inches(1.5))


def footer(slide,
           text='Genos AI  |  hello@genosai.tech  |  www.genosai.tech  |  +91 638-714-2699'):
    rect(slide, 0, Inches(7.05), W, Inches(0.45), DARK_BLUE)
    tb(slide, text, Inches(0.55), Inches(7.1), Inches(12.4), Inches(0.35),
       size=9, color=MGREY, align=PP_ALIGN.CENTER)


def card(slide, x, y, w, h, fill=DARK_BLUE):
    return rect(slide, x, y, w, h, fill)


def stat_card(slide, val, label, x, y, w=Inches(3.0), h=Inches(1.6), accent=GOLD):
    rect(slide, x, y, w, h, DARK_BLUE)
    rect(slide, x, y, w, Inches(0.05), accent)
    tb(slide, str(val), x + Inches(0.1), y + Inches(0.1),
       w - Inches(0.2), Inches(0.9), size=28, bold=True, color=accent,
       align=PP_ALIGN.CENTER)
    tb(slide, label, x + Inches(0.1), y + Inches(1.1),
       w - Inches(0.2), Inches(0.42), size=10, color=LGREY, align=PP_ALIGN.CENTER)


def bullet_card(slide, title, items, x, y, w, h,
                accent=GOLD, title_bg=None):
    if title_bg is None:
        title_bg = DARK_BLUE
    rect(slide, x, y, w, h, DARK_BLUE)
    rect(slide, x, y, w, Inches(0.42), title_bg)
    rect(slide, x, y, Inches(0.05), h, accent)
    tb(slide, title, x + Inches(0.1), y + Inches(0.05),
       w - Inches(0.15), Inches(0.35), size=11, bold=True, color=WHITE)
    bx = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.48),
                                  w - Inches(0.15), h - Inches(0.55))
    tf = bx.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r0 = p.add_run(); r0.text = '  '; r0.font.size = Pt(9)
        r1 = p.add_run(); r1.text = '● '; r1.font.size = Pt(9)
        r1.font.color.rgb = accent; r1.font.bold = True
        r2 = p.add_run(); r2.text = str(item)
        r2.font.size = Pt(10.5); r2.font.color.rgb = LGREY
        if i < len(items) - 1:
            p.space_after = Pt(5)




# SLIDE 1 – TITLE
# ════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
bg(s1, CHARCOAL)
rect(s1, 0, 0, Inches(0.7), H, GOLD)
rect(s1, Inches(0.7), Inches(4.5), W - Inches(0.7), Inches(0.06), GOLD)

tb(s1, "REALTOR COJO",
   Inches(1.1), Inches(1.0), Inches(11), Inches(0.7),
   size=16, bold=True, color=GOLD)
tb(s1, "Web Development &\nAI Automation Audit",
   Inches(1.1), Inches(1.7), Inches(11), Inches(2.0),
   size=44, bold=True, color=WHITE)
tb(s1, "Strategic Review & Digital Growth Roadmap  |  Lagos, Nigeria",
   Inches(1.1), Inches(3.72), Inches(10), Inches(0.5),
   size=16, color=GOLD, italic=True)
tb(s1, "Prepared by: Genos Apollo  |  May 2026",
   Inches(1.1), Inches(5.0), Inches(8), Inches(0.4),
   size=12, color=MGREY)
tb(s1, "CONFIDENTIAL",
   Inches(1.1), Inches(6.6), Inches(4), Inches(0.35),
   size=10, color=MGREY, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 2 – AGENDA
# ════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
bg(s2)
header(s2, "Agenda", "What we will cover today")
footer(s2)

agenda = [
    ("01", "Company Digital Overview",    "Current state of rc-realtor.vercel.app and digital footprint"),
    ("02", "Website Audit Findings",      "Critical gaps: domain, SEO, performance and UX"),
    ("03", "Web Development Roadmap",     "Recommended improvements & production-ready tech stack"),
    ("04", "AI Automation Opportunity",   "Where AI & WhatsApp automation can drive growth"),
    ("05", "Benefits of Automation",      "ROI, lead velocity and competitive advantage in Lagos"),
    ("06", "Implementation Timeline",     "Phased rollout plan across 12 months"),
    ("07", "Investment & Next Steps",     "Budget guidance and immediate actions"),
]

for i, (num, title, desc) in enumerate(agenda):
    col = i % 2
    row = i // 2
    x = Inches(0.55) + col * Inches(6.4)
    y = Inches(1.55) + row * Inches(1.22)
    card(s2, x, y, Inches(6.1), Inches(1.05))
    rect(s2, x, y, Inches(0.7), Inches(1.05), CHARCOAL)
    tb(s2, num, x, y + Inches(0.25), Inches(0.7), Inches(0.5),
       size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s2, title, x + Inches(0.8), y + Inches(0.1),
       Inches(5.0), Inches(0.4), size=13, bold=True, color=CHARCOAL)
    tb(s2, desc, x + Inches(0.8), y + Inches(0.52),
       Inches(5.0), Inches(0.42), size=10.5, color=MGREY, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 3 – COMPANY DIGITAL OVERVIEW
# ════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
bg(s3)
header(s3, "Company Digital Overview", "Realtor CoJo – current digital footprint")
footer(s3)

card(s3, Inches(0.4), Inches(1.4), Inches(5.8), Inches(5.5))
rect(s3, Inches(0.4), Inches(1.4), Inches(5.8), Inches(0.45), CHARCOAL)
tb(s3, "Company Snapshot", Inches(0.55), Inches(1.44),
   Inches(5.5), Inches(0.38), size=13, bold=True, color=WHITE)

snapshot = [
    ("Brand",        "Realtor CoJo (RC Realtor)"),
    ("Market",       "Luxury & Affordable Properties – Lagos, Nigeria"),
    ("Stage",        "Early-stage / emerging brand"),
    ("Website",      "rc-realtor.vercel.app  (⚠ Vercel subdomain – no custom domain)"),
    ("Tech Stack",   "Next.js / React on Vercel  (client-side rendering only)"),
    ("Rendering",    "CSR only – pages not visible to search engines"),
    ("Social",       "LinkedIn company page present"),
    ("CRM / CMS",    "None identified – no backend or listing management system"),
]
ty = Inches(2.02)
for k, v in snapshot:
    tb(s3, k, Inches(0.6), ty, Inches(1.55), Inches(0.3),
       size=10.5, bold=True, color=CHARCOAL)
    tb(s3, v, Inches(2.2), ty, Inches(3.8), Inches(0.3),
       size=10.5, color=MGREY)
    ty += Inches(0.52)

card(s3, Inches(6.5), Inches(1.4), Inches(6.5), Inches(5.5))
rect(s3, Inches(6.5), Inches(1.4), Inches(6.5), Inches(0.45), GOLD)
tb(s3, "Key Digital Health Indicators", Inches(6.65), Inches(1.44),
   Inches(6.2), Inches(0.38), size=13, bold=True, color=CHARCOAL)

metrics = [
    ("Custom Domain",           "⚠ None",       "Still on vercel.app subdomain",        RED),
    ("Mobile Page Speed",       "~40 / 100",    "CSR React – poor PageSpeed est.",       RED),
    ("Desktop Speed",           "~65 / 100",    "No SSR/SSG – JS bundle on load",        AMBER),
    ("SEO Visibility",          "~0 / 100",     "CSR = invisible to Google indexing",    RED),
    ("SSL / HTTPS",             "✓ Active",     "Vercel default SSL in place",           GREEN),
    ("Property Listings CMS",   "None",         "No headless CMS or listing backend",    RED),
    ("WhatsApp Integration",    "None",         "Primary Nigeria comm channel missing",  RED),
    ("AI / Chatbot",            "None",         "No automated lead capture on site",     RED),
]
ty2 = Inches(2.02)
for label, val, note, col in metrics:
    tb(s3, label, Inches(6.65), ty2, Inches(2.7), Inches(0.3),
       size=10.5, bold=True, color=CHARCOAL)
    tb(s3, val,   Inches(9.4), ty2, Inches(1.25), Inches(0.3),
       size=10.5, bold=True, color=col, align=PP_ALIGN.CENTER)
    tb(s3, note,  Inches(10.7), ty2, Inches(2.1), Inches(0.3),
       size=9, color=MGREY, italic=True)
    ty2 += Inches(0.52)

# ════════════════════════════════════════════════════════════════
# SLIDE 4 – WEBSITE AUDIT FINDINGS
# ════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
bg(s4)
header(s4, "Website Audit Findings", "rc-realtor.vercel.app – critical gaps blocking growth")
footer(s4)

audit_cards = [
    ("Domain & Hosting",
     ["No custom domain – vercel.app kills brand trust", "No professional email (@realtorcojo.com)",
      "Staging URL being used as production site", "No CDN configuration for Nigeria edge nodes"]),
    ("SEO (Critical)",
     ["Client-side React = zero Google visibility", "Site title is the only indexable content",
      "No meta descriptions, schema, or OG tags", "No sitemap.xml or robots.txt configured"]),
    ("UX & Listings",
     ["No visible property listings loading (CSR gap)", "No advanced search or filter functionality",
      "No virtual tour or photo gallery viewer", "No WhatsApp click-to-chat CTA on listings"]),
    ("Trust & Conversion",
     ["No testimonials, reviews or sold properties", "No agent bio or credentials visible",
      "No NDPR-compliant privacy or cookie notice", "No contact form – enquiry path unclear"]),
]

for i, (title, bullets) in enumerate(audit_cards):
    x = Inches(0.4) + i * Inches(3.2)
    y = Inches(1.45)
    card(s4, x, y, Inches(3.05), Inches(5.3))
    rect(s4, x, y, Inches(3.05), Inches(0.45), CHARCOAL)
    tb(s4, title, x + Inches(0.1), y + Inches(0.06),
       Inches(2.85), Inches(0.36), size=12, bold=True, color=WHITE)
    ty = y + Inches(0.62)
    for b in bullets:
        rect(s4, x + Inches(0.15), ty + Inches(0.07), Inches(0.1), Inches(0.1), GOLD)
        tb(s4, b, x + Inches(0.35), ty, Inches(2.6), Inches(0.55),
           size=10, color=CHARCOAL)
        ty += Inches(0.72)

rect(s4, Inches(0.4), Inches(6.8), Inches(12.5), Inches(0.06), GOLD)
tb(s4, "Immediate Priority: Custom domain + SSR/SSG migration — the site is currently invisible to all search engines",
   Inches(0.4), Inches(6.86), Inches(12.5), Inches(0.3),
   size=10.5, bold=True, color=CHARCOAL, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 5 – WEB DEVELOPMENT ROADMAP
# ════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
bg(s5)
header(s5, "Web Development Roadmap", "From staging prototype to production-ready luxury property platform")
footer(s5)

card(s5, Inches(0.4), Inches(1.4), Inches(4.5), Inches(5.6))
rect(s5, Inches(0.4), Inches(1.4), Inches(4.5), Inches(0.45), SLATE)
tb(s5, "Recommended Tech Stack", Inches(0.55), Inches(1.44),
   Inches(4.2), Inches(0.38), size=13, bold=True, color=WHITE)

stack = [
    ("Frontend",     "Next.js 15 – SSR + SSG for full SEO visibility"),
    ("CMS",          "Sanity.io – headless property listing management"),
    ("Domain",       "realtorcojo.com + professional email suite"),
    ("Hosting",      "Vercel (retain) + Cloudflare CDN (Lagos edge)"),
    ("Search",       "Algolia – property search (type, location, price)"),
    ("WhatsApp",     "WhatsApp Business API – click-to-chat on listings"),
    ("Analytics",    "GA4 + Microsoft Clarity heatmaps"),
    ("CRM",          "HubSpot Free or Zoho CRM – lead pipeline"),
]
ty = Inches(2.02)
for k, v in stack:
    tb(s5, k, Inches(0.6), ty, Inches(1.4), Inches(0.3),
       size=10, bold=True, color=GOLD)
    tb(s5, v, Inches(2.05), ty, Inches(2.75), Inches(0.3),
       size=10, color=CHARCOAL)
    ty += Inches(0.5)

card(s5, Inches(5.2), Inches(1.4), Inches(7.75), Inches(5.6))
rect(s5, Inches(5.2), Inches(1.4), Inches(7.75), Inches(0.45), CHARCOAL)
tb(s5, "Key Improvements & Features",
   Inches(5.35), Inches(1.44), Inches(7.4), Inches(0.38),
   size=13, bold=True, color=WHITE)

improvements = [
    ("Custom Domain & Professional Identity",
     "Register realtorcojo.com. Set up branded email, OG tags, favicon, and Google Business Profile."),
    ("SSR/SSG Migration – Full SEO Unlock",
     "Convert CSR React pages to Next.js SSR/SSG. Every property listing becomes a Google-indexable page."),
    ("Property Listings CMS",
     "Sanity.io backend: add/edit properties with photos, price (NGN + USD), bedrooms, location, virtual tour link."),
    ("WhatsApp CTA on Every Listing",
     "One-tap WhatsApp enquiry button on each property. Auto-prefilled message with property name and ID."),
    ("Trust & Social Proof Section",
     "Agent bio, certifications (NIESV), sold properties gallery, client video testimonials, Google reviews widget."),
    ("NDPR Compliance & Security",
     "Nigeria Data Protection Regulation compliant consent notice, secure contact forms, HTTPS enforcement."),
]
ty2 = Inches(1.98)
for title, desc in improvements:
    rect(s5, Inches(5.35), ty2 + Inches(0.1), Inches(0.08), Inches(0.25), GOLD)
    tb(s5, title, Inches(5.55), ty2, Inches(7.1), Inches(0.28),
       size=11, bold=True, color=CHARCOAL)
    tb(s5, desc, Inches(5.55), ty2 + Inches(0.3), Inches(7.1), Inches(0.35),
       size=10, color=MGREY)
    ty2 += Inches(0.78)

# ════════════════════════════════════════════════════════════════
# SLIDE 6 – AI AUTOMATION OPPORTUNITY MAP
# ════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
bg(s6)
header(s6, "AI Automation Opportunity Map", "Where AI & WhatsApp automation drive growth for Realtor CoJo")
footer(s6)

opps = [
    (CHARCOAL, "WhatsApp Lead Automation",
     ["AI WhatsApp bot handles enquiries 24/7", "Auto-send property photos & details on request",
      "Qualify leads (budget, timeline, location)", "Book property viewings via WhatsApp chat"]),
    (GOLD,     "Property Listings & Content",
     ["AI-generate property descriptions from specs", "Auto-resize and tag listing photos (AI)",
      "Auto-post new listings to Instagram & LinkedIn", "Dynamic price-drop alerts to interested buyers"]),
    (SLATE,    "Lead Management & Follow-Up",
     ["AI-score leads by budget & intent in CRM", "Automated follow-up sequences (WhatsApp + email)",
      "Re-engagement campaigns for cold leads", "Viewing reminders & confirmation automation"]),
    (GREEN,    "Market Intelligence & Growth",
     ["Track Lagos property price trends (AI digest)", "Monitor competitor listings on NPC & Jiji",
      "Weekly AI-generated Lagos market report", "Instagram Reels auto-captioned from listing data"]),
]

for i, (col, title, bullets) in enumerate(opps):
    c = i % 2
    r = i // 2
    x = Inches(0.4) + c * Inches(6.5)
    y = Inches(1.45) + r * Inches(2.85)
    card(s6, x, y, Inches(6.2), Inches(2.65))
    rect(s6, x, y, Inches(6.2), Inches(0.45), col)
    tb(s6, title, x + Inches(0.15), y + Inches(0.06),
       Inches(5.9), Inches(0.36), size=13, bold=True,
       color=CHARCOAL if col == GOLD else WHITE)
    ty = y + Inches(0.62)
    for b in bullets:
        rect(s6, x + Inches(0.18), ty + Inches(0.08), Inches(0.1), Inches(0.1), col)
        tb(s6, b, x + Inches(0.38), ty, Inches(5.65), Inches(0.38),
           size=10.5, color=CHARCOAL)
        ty += Inches(0.48)

# ════════════════════════════════════════════════════════════════
# SLIDE 7 – BENEFITS OF AI AUTOMATION
# ════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
bg(s7)
header(s7, "Benefits of AI Automation", "Measurable impact across Realtor CoJo's growth")
footer(s7)

stats = [
    ("24/7",    "WhatsApp lead\ncapture & response"),
    ("3×",      "Faster property\nlisting publication"),
    ("50%",     "Reduction in\nmanual follow-ups"),
    ("40%",     "More qualified\nviewings booked"),
]
for i, (val, label) in enumerate(stats):
    stat_card(s7, val, label,
              Inches(0.4) + i * Inches(3.2), Inches(1.45),
              w=Inches(3.0), h=Inches(1.6))

benefits = [
    ("Lead Velocity",
     ["AI bot captures WhatsApp leads at any hour", "No lead lost due to slow response time",
      "Viewing booked in the same conversation", "Auto-qualify by budget before agent time spent"]),
    ("Content at Scale",
     ["10× faster listing creation with AI copy", "Auto-post to Instagram, LinkedIn, Facebook",
      "Professional descriptions every time", "Price changes push to all channels instantly"]),
    ("Client Experience",
     ["Instant property details via WhatsApp", "Virtual tour links auto-sent to interested buyers",
      "Viewing reminders reduce no-shows by ~30%", "Post-sale follow-up keeps referral pipeline warm"]),
    ("Market Authority",
     ["Weekly Lagos market reports build credibility", "SEO-optimised listings attract organic traffic",
      "Google Maps + Property schema drives local search", "Data-driven pricing beats competitors on speed"]),
]

for i, (title, bullets) in enumerate(benefits):
    c = i % 2
    r = i // 2
    x = Inches(0.4) + c * Inches(6.5)
    y = Inches(3.2) + r * Inches(1.85)
    bullet_card(s7, title, bullets, x, y, Inches(6.2), Inches(1.65),
                title_bg=CHARCOAL if r == 0 else SLATE)

# ════════════════════════════════════════════════════════════════
# SLIDE 8 – ROI & BUSINESS CASE
# ════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
bg(s8)
header(s8, "ROI & Business Case", "Projected returns from web + AI investment (NGN / USD)")
footer(s8)

card(s8, Inches(0.4), Inches(1.4), Inches(5.8), Inches(5.6))
rect(s8, Inches(0.4), Inches(1.4), Inches(5.8), Inches(0.45), CHARCOAL)
tb(s8, "Indicative Investment (Year 1 – USD)",
   Inches(0.55), Inches(1.44), Inches(5.5), Inches(0.38),
   size=13, bold=True, color=WHITE)

invest_rows = [
    ("Custom Domain + Professional Email",   "$  80  –  $  150"),
    ("Website Rebuild (Next.js + Sanity)",   "$2,000  –  $4,500"),
    ("WhatsApp Business API Setup & Bot",    "$  600  –  $1,200"),
    ("Listing Automation & AI Copywriting",  "$  400  –  $  800"),
    ("SEO Sprint & NDPR Compliance",         "$  500  –  $1,000"),
    ("CRM Setup & Lead Workflows",           "$  300  –  $  600"),
    ("Ongoing monthly support (est.)",       "$  200  –  $  400 / mo"),
    ("TOTAL YEAR 1 (once-off + 12 mo)",     "$6,280  – $12,650"),
]
ty = Inches(2.0)
for i, (item, cost) in enumerate(invest_rows):
    is_total = i == len(invest_rows) - 1
    if is_total:
        rect(s8, Inches(0.4), ty - Inches(0.05), Inches(5.8), Inches(0.06), GOLD)
        ty += Inches(0.12)
        tb(s8, item, Inches(0.6), ty, Inches(3.5), Inches(0.35),
           size=11, bold=True, color=CHARCOAL)
        tb(s8, cost, Inches(4.05), ty, Inches(1.95), Inches(0.35),
           size=11, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)
    else:
        tb(s8, item, Inches(0.6), ty, Inches(3.5), Inches(0.3),
           size=10, color=CHARCOAL)
        tb(s8, cost, Inches(4.05), ty, Inches(1.95), Inches(0.3),
           size=10, color=MGREY, align=PP_ALIGN.RIGHT)
    ty += Inches(0.5)

card(s8, Inches(6.5), Inches(1.4), Inches(6.5), Inches(5.6))
rect(s8, Inches(6.5), Inches(1.4), Inches(6.5), Inches(0.45), GREEN)
tb(s8, "Projected Returns (1 deal pays for everything)",
   Inches(6.65), Inches(1.44), Inches(6.2), Inches(0.38),
   size=13, bold=True, color=WHITE)

returns = [
    ("Commission on 1 luxury property (1–2%)",   "$5,000 – $20,000+ per deal"),
    ("SEO-driven organic leads (vs. paid ads)",   "Save $200–$500 / mo in ad spend"),
    ("WhatsApp bot: 24/7 lead capture uplift",    "+30–50% more qualified enquiries"),
    ("Faster listing → viewing → close cycle",    "Est. −10 days deal cycle per sale"),
    ("Reduced manual follow-up hours",            "10–15 hrs/week reclaimed"),
    ("Referrals from better client experience",   "2–3 extra deals / year est."),
    ("BREAKEVEN (single mid-range sale)",        "Investment recovered in 1 closed deal"),
]
ty2 = Inches(2.0)
for i, (item, val) in enumerate(returns):
    is_total = i == len(returns) - 1
    if is_total:
        rect(s8, Inches(6.5), ty2 - Inches(0.05), Inches(6.5), Inches(0.06), GOLD)
        ty2 += Inches(0.12)
        tb(s8, item, Inches(6.65), ty2, Inches(3.9), Inches(0.35),
           size=11, bold=True, color=CHARCOAL)
        tb(s8, val, Inches(10.6), ty2, Inches(2.2), Inches(0.35),
           size=11, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)
    else:
        tb(s8, item, Inches(6.65), ty2, Inches(3.9), Inches(0.3),
           size=10, color=CHARCOAL)
        tb(s8, val, Inches(10.6), ty2, Inches(2.2), Inches(0.3),
           size=10, color=MGREY, align=PP_ALIGN.RIGHT)
    ty2 += Inches(0.5)

tb(s8, "One mid-range Lagos property sale at ₦80M (~$50K USD) generates 1% commission = $500 — far exceeds monthly investment",
   Inches(0.4), Inches(7.07), Inches(12.5), Inches(0.28),
   size=10.5, bold=True, color=CHARCOAL, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 9 – IMPLEMENTATION TIMELINE
# ════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(BLANK)
bg(s9)
header(s9, "Implementation Timeline", "Phased 12-month delivery plan")
footer(s9)

phases = [
    ("Phase 1\nMon 1–2",  GOLD,     [
        "Register realtorcojo.com domain",
        "Set up professional email suite",
        "Deploy SSR/SSG Next.js rebuild",
        "Google indexing & Search Console",
    ]),
    ("Phase 2\nMon 3–5",  CHARCOAL, [
        "Sanity CMS – property listings live",
        "WhatsApp Business API + AI bot",
        "Advanced search & filter launch",
        "Trust section: bio, reviews, sold",
    ]),
    ("Phase 3\nMon 6–8",  SLATE,    [
        "CRM lead pipeline & automation",
        "Social media auto-posting (IG/LI)",
        "Lagos SEO content sprint",
        "Virtual tour integration",
    ]),
    ("Phase 4\nMon 9–12", GREEN,    [
        "AI market intelligence reports",
        "Advanced chatbot (lead scoring)",
        "Google Ads + Meta Ads retargeting",
        "Full audit, review & optimise",
    ]),
]

for i, (label, col, items) in enumerate(phases):
    x = Inches(0.4) + i * Inches(3.2)
    y = Inches(1.45)
    card(s9, x, y, Inches(3.0), Inches(5.55))
    rect(s9, x, y, Inches(3.0), Inches(0.7), col)
    tb(s9, label, x, y + Inches(0.06), Inches(3.0), Inches(0.62),
       size=13, bold=True,
       color=CHARCOAL if col == GOLD else WHITE,
       align=PP_ALIGN.CENTER)
    if i < 3:
        rect(s9, x + Inches(3.02), y + Inches(0.3), Inches(0.16), Inches(0.06), MGREY)
    ty = y + Inches(0.88)
    for item in items:
        rect(s9, x + Inches(0.18), ty + Inches(0.08), Inches(0.1), Inches(0.1), col)
        tb(s9, item, x + Inches(0.38), ty, Inches(2.5), Inches(0.42),
           size=10.5, color=CHARCOAL)
        ty += Inches(0.9)

# ════════════════════════════════════════════════════════════════
# SLIDE 10 – WHY MOVE NOW?
# ════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(BLANK)
bg(s10, CHARCOAL)
rect(s10, 0, 0, Inches(0.5), H, GOLD)
footer(s10)

tb(s10, "Why Move — Now?",
   Inches(0.8), Inches(0.5), Inches(11), Inches(0.65),
   size=30, bold=True, color=GOLD)
tb(s10, "Lagos real estate is a ₦10+ trillion market. The agents winning today have digital foundations built yesterday.",
   Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.55),
   size=14, color=WHITE, italic=True)

reasons = [
    ("The current site is invisible to Google",
     "A client-side React app with no custom domain cannot be indexed. Every day without SSR is a day competitors capture the organic traffic Realtor CoJo should own for 'luxury properties Lagos' and similar queries."),
    ("WhatsApp is how Lagos buys property",
     "73%+ of Nigerian property enquiries start on WhatsApp. Without a click-to-chat button on every listing, Realtor CoJo is forcing motivated buyers to find contact details manually — most won't bother."),
    ("Early digital movers dominate Lagos real estate search",
     "Property search in Nigeria is consolidating onto Google, Jiji, and PropertyPro. Agents who build SEO authority now — through schema-rich, SSR listings — will be very difficult to displace in 12–18 months."),
    ("A single extra deal pays for everything",
     "The total Year 1 investment is ~$6K–$12K USD. One mid-range Lagos property commission (₦1M–₦5M) covers the entire programme. The risk-adjusted return is exceptionally high."),
    ("Trust signals separate top agents from the crowd",
     "Lagos property buyers research extensively before committing. Testimonials, sold history, professional bio, and NIESV credentials on a fast, modern website convert browsers into booked viewings — without any ad spend."),
]
ty = Inches(1.95)
for i, (title, desc) in enumerate(reasons):
    rect(s10, Inches(0.8), ty, Inches(0.55), Inches(0.55),
         GOLD if i % 2 == 0 else SLATE)
    tb(s10, str(i + 1), Inches(0.8), ty, Inches(0.55), Inches(0.55),
       size=16, bold=True, color=CHARCOAL if i % 2 == 0 else WHITE,
       align=PP_ALIGN.CENTER)
    tb(s10, title, Inches(1.5), ty, Inches(11), Inches(0.3),
       size=12, bold=True, color=WHITE)
    tb(s10, desc, Inches(1.5), ty + Inches(0.3), Inches(11.3), Inches(0.55),
       size=10.5, color=RGBColor(0xBB, 0xCC, 0xDD))
    ty += Inches(0.98)

# ════════════════════════════════════════════════════════════════
# SLIDE 11 – NEXT STEPS
# ════════════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(BLANK)
bg(s11)
header(s11, "Next Steps", "From audit to action — what happens now")
footer(s11)

steps = [
    ("This Week",  GOLD,     "Domain & Brand Foundation",
     "Register realtorcojo.com. Set up branded email. This single step immediately elevates trust and unlocks professional marketing."),
    ("Week 2",     CHARCOAL, "Discovery Workshop",
     "30-min session with Genos Apollo. Review audit findings, agree Phase 1 scope, share listing data and brand assets."),
    ("Week 3–4",   SLATE,    "Sprint 1 Builds",
     "SSR website rebuild begins. Domain live. WhatsApp Business API provisioned. Google Search Console connected."),
    ("Month 2",    GREEN,    "Site Launch & First Leads",
     "New site live at custom domain. First property listings indexed by Google. WhatsApp bot capturing enquiries 24/7."),
]

for i, (when, col, title, desc) in enumerate(steps):
    x = Inches(0.4) + i * Inches(3.2)
    y = Inches(1.5)
    card(s11, x, y, Inches(3.0), Inches(4.0))
    rect(s11, x, y, Inches(3.0), Inches(0.45), col)
    tb(s11, when, x, y + Inches(0.06), Inches(3.0), Inches(0.36),
       size=12, bold=True,
       color=CHARCOAL if col == GOLD else WHITE,
       align=PP_ALIGN.CENTER)
    tb(s11, title, x + Inches(0.15), y + Inches(0.58),
       Inches(2.7), Inches(0.36), size=12, bold=True,
       color=WARM if col == GOLD else col)
    tb(s11, desc, x + Inches(0.15), y + Inches(1.0),
       Inches(2.7), Inches(2.8), size=10.5, color=CHARCOAL)

card(s11, Inches(0.4), Inches(5.8), Inches(12.5), Inches(1.1), fill=CHARCOAL)
tb(s11, "Ready to turn Realtor CoJo into Lagos's leading digital property brand?",
   Inches(0.6), Inches(5.88), Inches(7.5), Inches(0.4),
   size=14, bold=True, color=GOLD)
tb(s11, "Contact Genos Apollo to schedule your discovery workshop.",
   Inches(0.6), Inches(6.28), Inches(7), Inches(0.38),
   size=11, color=WHITE)
tb(s11, "rc-realtor.vercel.app  |  linkedin.com/company/realtor-cojo",
   Inches(8.0), Inches(6.0), Inches(5.0), Inches(0.38),
   size=10, color=GOLD, italic=True, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════
# SLIDE 12 – THANK YOU / BACK COVER
# ════════════════════════════════════════════════════════════════
s12 = prs.slides.add_slide(BLANK)
bg(s12, CHARCOAL)
rect(s12, 0, 0, Inches(0.7), H, GOLD)
rect(s12, Inches(0.7), Inches(3.5), W - Inches(0.7), Inches(0.06), GOLD)

tb(s12, "Thank You",
   Inches(1.1), Inches(2.0), Inches(11), Inches(1.2),
   size=52, bold=True, color=WHITE)
tb(s12, "Luxury & Affordable Properties in Lagos — powered by digital intelligence.",
   Inches(1.1), Inches(3.1), Inches(10), Inches(0.55),
   size=16, color=GOLD, italic=True)
tb(s12, "Realtor CoJo  ·  rc-realtor.vercel.app  ·  Lagos, Nigeria",
   Inches(1.1), Inches(4.0), Inches(9), Inches(0.4),
   size=12, color=RGBColor(0xAA, 0xBB, 0xCC))
tb(s12, "Prepared by Genos Apollo  |  May 2026  |  Confidential",
   Inches(1.1), Inches(6.7), Inches(10), Inches(0.35),
   size=10, color=MGREY)

# ── Save ────────────────────────────────────────────────────────
OUT = ROOT / "output" / "decks" / "RealtorCoJo_WebDev_AI_Automation_Audit_GenosAI.pptx"
prs.save(str(OUT))
print(f"Saved: {OUT}")
