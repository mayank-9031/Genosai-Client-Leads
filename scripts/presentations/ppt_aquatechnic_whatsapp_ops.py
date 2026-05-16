import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.parent / "utils"))
from ppt_theme import (
    NAVY, DARK_BLUE, ACCENT, GOLD, RED, GREEN, PURPLE, ORANGE, TEAL, CYAN,
    WHITE, LGREY, MGREY, W, H,
    bg, rect, tb, aline, bullets, tag,
    slide_header, slide_footer, score_badge,
    stat_cards, top_bar_card, phase_card, before_after, genos_logo_block,
)

ROOT = Path(__file__).parent.parent.parent
prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]

# Client-specific colors
WA_GREEN  = RGBColor(0x25, 0xD3, 0x66)   # WhatsApp brand green


def wa_bubble(slide, x, y, w, text, sender="", incoming=True):
    """WhatsApp-style chat bubble. incoming=True → left (system), False → right (user)."""
    h = Inches(max(0.55, 0.35 + len(text) / 120))
    fill = RGBColor(0x1E, 0x3A, 0x2C) if incoming else RGBColor(0x0A, 0x1A, 0x0E)
    shape = slide.shapes.add_shape(9, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = WA_GREEN if incoming else RGBColor(0x12, 0x3A, 0x1C)
    from pptx.util import Pt as _Pt
    shape.line.width = _Pt(0.5)
    tf = shape.text_frame; tf.word_wrap = True
    if sender:
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.LEFT
        r0 = p0.add_run(); r0.text = sender
        r0.font.size = _Pt(8); r0.font.bold = True
        r0.font.color.rgb = WA_GREEN if incoming else CYAN
        p1 = tf.add_paragraph()
    else:
        p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    r1 = p1.add_run(); r1.text = text
    r1.font.size = _Pt(9.5); r1.font.color.rgb = WHITE


# ── SLIDE 1 — COVER ──────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(BLANK)
bg(s1)

# diagonal accent strip
rect(s1, Inches(9.8), 0, Inches(3.53), H, DARK_BLUE)
rect(s1, Inches(9.6), 0, Inches(0.12), H, CYAN)

# WA badge
rect(s1, Inches(10.1), Inches(1.6), Inches(2.5), Inches(2.5), WA_GREEN)
tb(s1, "WA", Inches(10.1), Inches(1.65), Inches(2.5), Inches(2.4),
   size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# tagline on right
tb(s1, "Powered by", Inches(9.85), Inches(4.3), Inches(3.2), Inches(0.4),
   size=11, color=LGREY, align=PP_ALIGN.CENTER)
tb(s1, "WhatsApp", Inches(9.85), Inches(4.7), Inches(3.2), Inches(0.5),
   size=20, bold=True, color=WA_GREEN, align=PP_ALIGN.CENTER)

# left content
tb(s1, "AQUA TECHNIC", Inches(0.55), Inches(1.3), Inches(9.0), Inches(0.5),
   size=13, bold=True, color=CYAN)
aline(s1, Inches(1.0), color=CYAN, w=Inches(2.0))
tb(s1, "WhatsApp Operations System", Inches(0.55), Inches(1.15), Inches(9.0), Inches(1.1),
   size=40, bold=True, color=WHITE)
tb(s1, "Centralized Field Management · Site Tracking · Team Coordination · Client Updates",
   Inches(0.55), Inches(2.35), Inches(9.0), Inches(0.6),
   size=14, color=LGREY)

aline(s1, Inches(3.15), color=CYAN, w=Inches(4.5))

# module pills
for i, (label, col) in enumerate([
    ("Daily Check-in", WA_GREEN),
    ("Site Updates", CYAN),
    ("Task Assignment", ACCENT),
    ("Manager Dashboard", GOLD),
    ("Client Automation", ORANGE),
    ("Maintenance Mgmt", PURPLE),
]):
    tag(s1, label, Inches(0.55 + i * 2.18), Inches(3.45), fill=col, size=10)

tb(s1, "All 6 modules · All through WhatsApp your team already uses",
   Inches(0.55), Inches(3.95), Inches(9.0), Inches(0.45),
   size=12, color=MGREY)

# footer
rect(s1, 0, Inches(6.85), W, Inches(0.65), RGBColor(0x06, 0x09, 0x1F))
tb(s1, "Prepared by  Genos AI  |  hello@genosai.tech  |  www.genosai.tech  |  +91 638-714-2699",
   Inches(0.4), Inches(6.9), Inches(12.0), Inches(0.45),
   size=10, color=MGREY, align=PP_ALIGN.CENTER)


# ── SLIDE 2 — THE PROBLEM ────────────────────────────────────────────────────
s2 = prs.slides.add_slide(BLANK)
bg(s2)

# red accent bar
rect(s2, 0, 0, Inches(0.12), H, RED)

tb(s2, "THE CURRENT REALITY", Inches(0.55), Inches(0.35), Inches(12.0), Inches(0.4),
   size=11, bold=True, color=RED)
aline(s2, Inches(0.85), color=RED, w=Inches(2.5))
tb(s2, "How Aqua Technic runs operations today",
   Inches(0.55), Inches(0.5), Inches(12.0), Inches(0.7),
   size=28, bold=True, color=WHITE)

# two columns
# left — what you manage
rect(s2, Inches(0.5), Inches(1.55), Inches(5.9), Inches(4.7), DARK_BLUE)
tb(s2, "What you manage daily", Inches(0.7), Inches(1.65), Inches(5.5), Inches(0.45),
   size=13, bold=True, color=CYAN)
bullets(s2, [
    "Multiple active construction sites",
    "Supervisors + field workers across locations",
    "Maintenance teams on AMC routes",
    "490+ pool clients expecting updates",
    "Ongoing projects at different stages",
    "Waterproofing, electrical, renovation crews",
], Inches(0.7), Inches(2.15), Inches(5.5), Inches(3.8), size=12, dot=CYAN)

# right — the chaos
rect(s2, Inches(6.7), Inches(1.55), Inches(6.1), Inches(4.7), RGBColor(0x1A, 0x04, 0x04))
tb(s2, "What actually happens", Inches(6.9), Inches(1.65), Inches(5.7), Inches(0.45),
   size=13, bold=True, color=RED)
bullets(s2, [
    "Manual WhatsApp messages to chase attendance",
    "No standard format — each supervisor reports differently",
    "Manager calls workers to find who is at which site",
    "Site photos sent to personal chats, never organized",
    "Delayed or missed maintenance visit reminders",
    "Clients call to ask for updates that exist somewhere in chats",
    "Completed tasks forgotten — no record or follow-up",
    "No single view of all sites, all workers, all tasks",
], Inches(6.9), Inches(2.15), Inches(5.7), Inches(3.8), size=12, dot=RED)

# bottom banner
rect(s2, 0, Inches(6.4), W, Inches(0.65), RGBColor(0x1A, 0x04, 0x04))
tb(s2, "The team already uses WhatsApp. The problem is: there is no system behind it.",
   Inches(0.4), Inches(6.45), Inches(12.5), Inches(0.5),
   size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rect(s2, 0, Inches(7.05), W, Inches(0.45), RGBColor(0x06, 0x09, 0x1F))
tb(s2, "Genos AI  |  hello@genosai.tech  |  www.genosai.tech",
   Inches(0.4), Inches(7.08), Inches(12.0), Inches(0.35),
   size=9, color=MGREY, align=PP_ALIGN.CENTER)


# ── SLIDE 3 — THE SOLUTION OVERVIEW ─────────────────────────────────────────
s3 = prs.slides.add_slide(BLANK)
bg(s3)
rect(s3, 0, 0, Inches(0.12), H, CYAN)

tb(s3, "THE SOLUTION", Inches(0.55), Inches(0.35), Inches(12.0), Inches(0.4),
   size=11, bold=True, color=CYAN)
aline(s3, Inches(0.85), color=CYAN, w=Inches(2.0))
tb(s3, "One centralized system. Entirely through WhatsApp.",
   Inches(0.55), Inches(0.5), Inches(9.5), Inches(0.75),
   size=28, bold=True, color=WHITE)
tb(s3, "No new apps. No training. Your team keeps using WhatsApp — the system runs behind it.",
   Inches(0.55), Inches(1.3), Inches(9.5), Inches(0.45),
   size=13, color=LGREY)

# 6 module cards — 3×2 grid
modules = [
    ("01", "Daily Employee\nCheck-in", "Automated morning attendance + site assignment sent to every worker", WA_GREEN),
    ("02", "Site Update\nCollection", "Scheduled reminders to supervisors — photos/videos/notes auto-stored", CYAN),
    ("03", "Task Assignment\nSystem", "Assign workers to sites, set deadlines, track completion via WhatsApp", ACCENT),
    ("04", "Automated\nNudges", "Auto follow-up if attendance missing, report late, or task incomplete", ORANGE),
    ("05", "Manager\nDashboard", "Real-time view: who checked in, site status, pending tasks, delays", GOLD),
    ("06", "Client Update &\nMaintenance", "Progress photos + AMC reminders sent to clients automatically", PURPLE),
]

positions = [
    (Inches(0.45), Inches(1.95)),
    (Inches(4.6),  Inches(1.95)),
    (Inches(8.75), Inches(1.95)),
    (Inches(0.45), Inches(4.3)),
    (Inches(4.6),  Inches(4.3)),
    (Inches(8.75), Inches(4.3)),
]
cw = Inches(3.9); ch = Inches(2.1)

for (num, title, desc, col), (cx, cy) in zip(modules, positions):
    rect(s3, cx, cy, cw, ch, DARK_BLUE)
    rect(s3, cx, cy, Inches(0.07), ch, col)
    tb(s3, num, cx + Inches(0.18), cy + Inches(0.1), Inches(0.6), Inches(0.45),
       size=13, bold=True, color=col)
    tb(s3, title, cx + Inches(0.18), cy + Inches(0.35), cw - Inches(0.3), Inches(0.7),
       size=14, bold=True, color=WHITE)
    tb(s3, desc, cx + Inches(0.18), cy + Inches(1.0), cw - Inches(0.3), Inches(0.95),
       size=10, color=LGREY)

rect(s3, 0, Inches(6.6), W, Inches(0.9), RGBColor(0x06, 0x09, 0x1F))
tb(s3, "Genos AI  |  hello@genosai.tech  |  www.genosai.tech",
   Inches(0.4), Inches(6.65), Inches(12.0), Inches(0.4),
   size=9, color=MGREY, align=PP_ALIGN.CENTER)


# ── SLIDE 4 — DAILY EMPLOYEE CHECK-IN ────────────────────────────────────────
s4 = prs.slides.add_slide(BLANK)
bg(s4)
rect(s4, 0, 0, Inches(0.12), H, WA_GREEN)

tb(s4, "MODULE 01", Inches(0.55), Inches(0.3), Inches(5.0), Inches(0.35),
   size=10, bold=True, color=WA_GREEN)
tb(s4, "Daily Employee Check-in", Inches(0.55), Inches(0.55), Inches(7.5), Inches(0.75),
   size=30, bold=True, color=WHITE)
aline(s4, Inches(1.4), color=WA_GREEN, w=Inches(3.5))
tb(s4, "Every morning — automated. Zero manual effort from management.",
   Inches(0.55), Inches(1.5), Inches(7.5), Inches(0.4),
   size=13, color=LGREY)

# what happens
rect(s4, Inches(0.45), Inches(2.05), Inches(5.4), Inches(4.65), DARK_BLUE)
tb(s4, "How it works", Inches(0.65), Inches(2.2), Inches(5.0), Inches(0.4),
   size=13, bold=True, color=WA_GREEN)
bullets(s4, [
    "6:30 AM — system sends WhatsApp to all workers",
    "Worker replies to mark attendance (1 tap)",
    "System assigns site + sends reporting instructions",
    "Worker uploads morning site photo",
    "Live location optional — confirms worker is on site",
    "Late / non-response → manager gets alert",
    "All responses logged automatically — no manual tracking",
], Inches(0.65), Inches(2.65), Inches(5.0), Inches(3.8), size=11, dot=WA_GREEN)

# chat bubbles simulation
rect(s4, Inches(6.15), Inches(1.85), Inches(6.75), Inches(5.2), RGBColor(0x0B, 0x14, 0x1A))
tb(s4, "WhatsApp", Inches(6.15), Inches(1.87), Inches(6.75), Inches(0.35),
   size=10, bold=True, color=WA_GREEN, align=PP_ALIGN.CENTER)

wa_bubble(s4, Inches(6.35), Inches(2.3), Inches(5.8),
          "Good morning Ramesh! Please mark your attendance for today.\nSend: 1 = Present  2 = Leave",
          sender="Genos System", incoming=True)
wa_bubble(s4, Inches(7.5), Inches(3.45), Inches(4.5),
          "1",
          sender="Ramesh (Worker)", incoming=False)
wa_bubble(s4, Inches(6.35), Inches(4.25), Inches(5.8),
          "Attendance marked. Your site today: Jubilee Hills Project.\nPlease upload 1 photo when you arrive.",
          sender="Genos System", incoming=True)
wa_bubble(s4, Inches(7.5), Inches(5.4), Inches(4.5),
          "[Photo uploaded] + Live Location",
          sender="Ramesh (Worker)", incoming=False)

tag(s4, "LIVE 6:30 AM", Inches(6.35), Inches(6.7), fill=WA_GREEN, size=9)
tag(s4, "AUTO-LOGGED", Inches(8.1), Inches(6.7), fill=CYAN, size=9)
tag(s4, "NO MANUAL WORK", Inches(9.85), Inches(6.7), fill=DARK_BLUE, size=9)

rect(s4, 0, Inches(7.05), W, Inches(0.45), RGBColor(0x06, 0x09, 0x1F))
tb(s4, "Genos AI  |  hello@genosai.tech  |  www.genosai.tech",
   Inches(0.4), Inches(7.08), Inches(12.0), Inches(0.35),
   size=9, color=MGREY, align=PP_ALIGN.CENTER)


# ── SLIDE 5 — SITE UPDATES + TASK ASSIGNMENT ─────────────────────────────────
s5 = prs.slides.add_slide(BLANK)
bg(s5)
rect(s5, 0, 0, Inches(0.12), H, CYAN)

tb(s5, "MODULES 02 + 03", Inches(0.55), Inches(0.3), Inches(6.0), Inches(0.35),
   size=10, bold=True, color=CYAN)
tb(s5, "Site Updates + Task Assignment",
   Inches(0.55), Inches(0.55), Inches(12.0), Inches(0.75),
   size=28, bold=True, color=WHITE)
aline(s5, Inches(1.4), color=CYAN, w=Inches(4.0))

# left half — site updates
rect(s5, Inches(0.4), Inches(1.65), Inches(6.1), Inches(5.2), DARK_BLUE)
rect(s5, Inches(0.4), Inches(1.65), Inches(6.1), Inches(0.42), RGBColor(0x00, 0x6A, 0x5C))
tb(s5, "  Site Update Collection", Inches(0.4), Inches(1.65), Inches(6.1), Inches(0.42),
   size=13, bold=True, color=WHITE)

bullets(s5, [
    "Supervisor gets WhatsApp reminder at scheduled time",
    "End-of-day update: photos, progress notes, issues",
    "Multiple media types: photos, videos, text notes",
    "Auto-stored to dashboard + Google Sheets",
    "Tagged by site name and date — fully searchable",
    "No update submitted? — auto-escalation to manager",
], Inches(0.6), Inches(2.2), Inches(5.7), Inches(2.8), size=11, dot=CYAN)

tb(s5, "Example reminder:", Inches(0.6), Inches(4.9), Inches(5.7), Inches(0.3),
   size=10, bold=True, color=LGREY)

# mini bubble
s_box = s5.shapes.add_shape(9, Inches(0.6), Inches(5.15), Inches(5.5), Inches(1.1))
s_box.fill.solid(); s_box.fill.fore_color.rgb = WHITE
s_box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
tf = s_box.text_frame; tf.word_wrap = True
tf.margin_left = Pt(8); tf.margin_top = Pt(5)
p0 = tf.paragraphs[0]
r0 = p0.add_run(); r0.text = "Genos System  "
r0.font.size = Pt(8); r0.font.bold = True; r0.font.color.rgb = WA_GREEN
p1 = tf.add_paragraph()
r1 = p1.add_run()
r1.text = "Hi Ravi. Please send end-of-day update for Banjara Hills site.\nSend: photos + status (On track / Delayed)."
r1.font.size = Pt(10); r1.font.color.rgb = RGBColor(0x11, 0x11, 0x11)

# right half — task assignment
rect(s5, Inches(6.75), Inches(1.65), Inches(6.1), Inches(5.2), DARK_BLUE)
rect(s5, Inches(6.75), Inches(1.65), Inches(6.1), Inches(0.42), RGBColor(0x1A, 0x2A, 0x55))
tb(s5, "  Task Assignment System", Inches(6.75), Inches(1.65), Inches(6.1), Inches(0.42),
   size=13, bold=True, color=WHITE)

bullets(s5, [
    "Management sends task via dashboard or WhatsApp",
    "Worker receives directly: task, site, deadline",
    "Task types: waterproofing, electrical, maintenance, repair",
    "Worker acknowledges — logged with timestamp",
    "Progress updates tied to the specific task",
    "Completion marked by worker — manager notified",
    "Incomplete past deadline? Auto-escalation alert",
], Inches(6.95), Inches(2.2), Inches(5.7), Inches(3.2), size=11, dot=ACCENT)

tb(s5, "Example task message:", Inches(6.95), Inches(5.05), Inches(5.7), Inches(0.3),
   size=10, bold=True, color=LGREY)

t_box = s5.shapes.add_shape(9, Inches(6.95), Inches(5.3), Inches(5.5), Inches(1.1))
t_box.fill.solid(); t_box.fill.fore_color.rgb = WHITE
t_box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
tf2 = t_box.text_frame; tf2.word_wrap = True
tf2.margin_left = Pt(8); tf2.margin_top = Pt(5)
p2 = tf2.paragraphs[0]
r2 = p2.add_run(); r2.text = "Genos System  "
r2.font.size = Pt(8); r2.font.bold = True; r2.font.color.rgb = WA_GREEN
p3 = tf2.add_paragraph()
r3 = p3.add_run()
r3.text = "Suresh, you have been assigned waterproofing work at\nJubilee Hills site. Deadline: tomorrow 5 PM. Reply 1 to confirm."
r3.font.size = Pt(10); r3.font.color.rgb = RGBColor(0x11, 0x11, 0x11)

rect(s5, 0, Inches(7.05), W, Inches(0.45), RGBColor(0x06, 0x09, 0x1F))
tb(s5, "Genos AI  |  hello@genosai.tech  |  www.genosai.tech",
   Inches(0.4), Inches(7.08), Inches(12.0), Inches(0.35),
   size=9, color=MGREY, align=PP_ALIGN.CENTER)


# ── SLIDE 6 — AUTOMATED NUDGES ───────────────────────────────────────────────
s6 = prs.slides.add_slide(BLANK)
bg(s6)
rect(s6, 0, 0, Inches(0.12), H, ORANGE)

tb(s6, "MODULE 04", Inches(0.55), Inches(0.3), Inches(5.0), Inches(0.35),
   size=10, bold=True, color=ORANGE)
tb(s6, "Automated Nudges & Follow-ups",
   Inches(0.55), Inches(0.55), Inches(9.0), Inches(0.75),
   size=30, bold=True, color=WHITE)
aline(s6, Inches(1.4), color=ORANGE, w=Inches(3.5))
tb(s6, "The manager chase is the most exhausting part of running a multi-site operation. This eliminates it.",
   Inches(0.55), Inches(1.5), Inches(9.5), Inches(0.45),
   size=12, color=LGREY)

# trigger + action cards
triggers = [
    ("Attendance Missing", "8:00 AM — worker hasn't responded", "Reminder sent. If no response by 9 AM, manager alerted.", RED),
    ("Site Report Late", "Supervisor hasn't submitted by 6 PM", "Auto follow-up sent. Manager flagged if still missing by 7 PM.", ORANGE),
    ("Task Incomplete", "Deadline passed, no completion mark", "Escalation message sent to worker. Manager CC'd automatically.", GOLD),
    ("Maintenance Due", "AMC visit scheduled for tomorrow", "Team reminder sent + client notified. Reschedule flow if cancelled.", CYAN),
    ("Client Update Pending", "Project milestone reached", "Auto-trigger: photos + update sent to client via WhatsApp.", WA_GREEN),
    ("Payment Due (AMC)", "30 / 7 / 1 day before renewal", "WhatsApp reminder with payment link sent to client automatically.", PURPLE),
]

for i, (trigger, condition, action, col) in enumerate(triggers):
    row = i // 3; col_x = i % 3
    cx = Inches(0.45 + col_x * 4.28); cy = Inches(2.1 + row * 2.35)
    cw2 = Inches(4.0); ch2 = Inches(2.1)
    rect(s6, cx, cy, cw2, ch2, DARK_BLUE)
    rect(s6, cx, cy, Inches(0.07), ch2, col)
    tb(s6, trigger, cx + Inches(0.18), cy + Inches(0.1), cw2 - Inches(0.25), Inches(0.35),
       size=12, bold=True, color=col)
    tb(s6, "When: " + condition, cx + Inches(0.18), cy + Inches(0.48), cw2 - Inches(0.25), Inches(0.45),
       size=10, color=LGREY, italic=True)
    tb(s6, action, cx + Inches(0.18), cy + Inches(0.98), cw2 - Inches(0.25), Inches(0.95),
       size=10, color=WHITE)

rect(s6, 0, Inches(7.05), W, Inches(0.45), RGBColor(0x06, 0x09, 0x1F))
tb(s6, "Genos AI  |  hello@genosai.tech  |  www.genosai.tech",
   Inches(0.4), Inches(7.08), Inches(12.0), Inches(0.35),
   size=9, color=MGREY, align=PP_ALIGN.CENTER)


# ── SLIDE 7 — MANAGER DASHBOARD ──────────────────────────────────────────────
s7 = prs.slides.add_slide(BLANK)
bg(s7)
rect(s7, 0, 0, Inches(0.12), H, GOLD)

tb(s7, "MODULE 05", Inches(0.55), Inches(0.3), Inches(5.0), Inches(0.35),
   size=10, bold=True, color=GOLD)
tb(s7, "Manager Dashboard",
   Inches(0.55), Inches(0.55), Inches(9.0), Inches(0.75),
   size=30, bold=True, color=WHITE)
aline(s7, Inches(1.4), color=GOLD, w=Inches(2.5))
tb(s7, "One screen. Everything happening across all sites, all workers, all tasks — right now.",
   Inches(0.55), Inches(1.5), Inches(9.5), Inches(0.45),
   size=13, color=LGREY)

# dashboard mockup panel
rect(s7, Inches(0.45), Inches(2.05), Inches(8.1), Inches(4.85), DARK_BLUE)
tb(s7, "Today's Operations — 12 May 2026",
   Inches(0.65), Inches(2.15), Inches(7.7), Inches(0.4),
   size=13, bold=True, color=WHITE)
rect(s7, Inches(0.45), Inches(2.05), Inches(8.1), Inches(0.42), RGBColor(0x10, 0x25, 0x5A))

# stat boxes
stats = [
    ("18 / 22", "Workers Checked In", WA_GREEN),
    ("4 of 6", "Sites Reporting", CYAN),
    ("7", "Tasks Active", GOLD),
    ("2", "Alerts Pending", RED),
]
for i, (val, label, col) in enumerate(stats):
    sx = Inches(0.65 + i * 1.95); sy = Inches(2.65)
    rect(s7, sx, sy, Inches(1.75), Inches(1.0), RGBColor(0x14, 0x20, 0x50))
    tb(s7, val, sx, sy + Inches(0.05), Inches(1.75), Inches(0.55),
       size=22, bold=True, color=col, align=PP_ALIGN.CENTER)
    tb(s7, label, sx, sy + Inches(0.6), Inches(1.75), Inches(0.35),
       size=9, color=LGREY, align=PP_ALIGN.CENTER)

# site list
sites = [
    ("Jubilee Hills - Construction", "Ramesh, Suresh, Ali", "On Track", WA_GREEN),
    ("Banjara Hills - Renovation",   "Ravi, Kumar",         "Update Pending", ORANGE),
    ("Gachibowli - Maintenance",     "Mohan",               "Completed", CYAN),
    ("Kondapur - New Pool",          "Team A (4 workers)",  "On Track", WA_GREEN),
]
tb(s7, "Site  |  Workers  |  Status", Inches(0.65), Inches(3.82), Inches(7.7), Inches(0.3),
   size=9, bold=True, color=MGREY)
for i, (site, workers, status, scol) in enumerate(sites):
    sy2 = Inches(4.15 + i * 0.62)
    if i % 2 == 0:
        rect(s7, Inches(0.65), sy2, Inches(7.7), Inches(0.58), RGBColor(0x12, 0x1E, 0x48))
    tb(s7, site, Inches(0.75), sy2 + Inches(0.1), Inches(3.5), Inches(0.4), size=11, color=WHITE)
    tb(s7, workers, Inches(4.35), sy2 + Inches(0.1), Inches(2.2), Inches(0.4), size=10, color=LGREY)
    tag(s7, status, Inches(6.65), sy2 + Inches(0.12), fill=scol if scol != ORANGE else ORANGE, size=8)

# right side — key features list
rect(s7, Inches(8.75), Inches(2.05), Inches(4.13), Inches(4.85), DARK_BLUE)
rect(s7, Inches(8.75), Inches(2.05), Inches(0.07), Inches(4.85), GOLD)
tb(s7, "What the dashboard shows", Inches(8.95), Inches(2.2), Inches(3.75), Inches(0.45),
   size=13, bold=True, color=GOLD)
bullets(s7, [
    "Real-time attendance status",
    "Which workers are at which site",
    "Site progress photos (latest upload)",
    "Pending / overdue task list",
    "Alerts: missing reports, late workers",
    "AMC schedule for the week",
    "Client updates sent vs pending",
    "Full activity log — searchable by date",
], Inches(8.95), Inches(2.7), Inches(3.8), Inches(4.0), size=11, dot=GOLD)

rect(s7, 0, Inches(7.05), W, Inches(0.45), RGBColor(0x06, 0x09, 0x1F))
tb(s7, "Genos AI  |  hello@genosai.tech  |  www.genosai.tech",
   Inches(0.4), Inches(7.08), Inches(12.0), Inches(0.35),
   size=9, color=MGREY, align=PP_ALIGN.CENTER)


# ── SLIDE 8 — CLIENT UPDATE AUTOMATION ───────────────────────────────────────
s8 = prs.slides.add_slide(BLANK)
bg(s8)
rect(s8, 0, 0, Inches(0.12), H, PURPLE)

tb(s8, "MODULE 06A", Inches(0.55), Inches(0.3), Inches(5.0), Inches(0.35),
   size=10, bold=True, color=PURPLE)
tb(s8, "Client Update Automation",
   Inches(0.55), Inches(0.55), Inches(9.0), Inches(0.75),
   size=30, bold=True, color=WHITE)
aline(s8, Inches(1.4), color=PURPLE, w=Inches(3.0))
tb(s8, "Clients stop calling when they already know what is happening.",
   Inches(0.55), Inches(1.5), Inches(9.5), Inches(0.45),
   size=13, color=LGREY)

# before / after
rect(s8, Inches(0.45), Inches(2.05), Inches(5.9), Inches(4.6), RGBColor(0x1A, 0x04, 0x04))
rect(s8, Inches(0.45), Inches(2.05), Inches(5.9), Inches(0.4), RGBColor(0x55, 0x10, 0x10))
tb(s8, "  Before — Client experience today", Inches(0.45), Inches(2.05), Inches(5.9), Inches(0.4),
   size=12, bold=True, color=WHITE)
bullets(s8, [
    "Client calls to ask: 'What is the status of my pool?'",
    "Manager must stop to give update verbally",
    "No photographic progress evidence sent",
    "Maintenance visit date? Client doesn't know until last minute",
    "AMC renewal? Client forgets — chases company",
    "Post-project: no follow-up, no review request",
    "Premium project, below-average client communication",
], Inches(0.65), Inches(2.6), Inches(5.5), Inches(3.8), size=11, dot=RED)

rect(s8, Inches(6.7), Inches(2.05), Inches(6.1), Inches(4.6), RGBColor(0x04, 0x14, 0x0C))
rect(s8, Inches(6.7), Inches(2.05), Inches(6.1), Inches(0.4), RGBColor(0x10, 0x45, 0x28))
tb(s8, "  After — Automated client updates", Inches(6.7), Inches(2.05), Inches(6.1), Inches(0.4),
   size=12, bold=True, color=WHITE)
bullets(s8, [
    "Weekly WhatsApp update: progress photos + brief note",
    "Milestone reached? Auto message + photos sent immediately",
    "Construction complete: completion album sent via WhatsApp",
    "Maintenance visit: reminder 2 days before, confirmation day-of",
    "AMC renewal: 30/7/1 day reminders with payment link",
    "Post-project: review request 1 week after handover",
    "All automated — no manager intervention needed",
], Inches(6.9), Inches(2.6), Inches(5.7), Inches(3.8), size=11, dot=WA_GREEN)

rect(s8, 0, Inches(6.75), W, Inches(0.3), PURPLE)
tb(s8, "Every Aqua Technic client — hotel or homeowner — feels like they have a dedicated project manager.",
   Inches(0.4), Inches(6.78), Inches(12.5), Inches(0.28),
   size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rect(s8, 0, Inches(7.05), W, Inches(0.45), RGBColor(0x06, 0x09, 0x1F))
tb(s8, "Genos AI  |  hello@genosai.tech  |  www.genosai.tech",
   Inches(0.4), Inches(7.08), Inches(12.0), Inches(0.35),
   size=9, color=MGREY, align=PP_ALIGN.CENTER)


# ── SLIDE 9 — MAINTENANCE MANAGEMENT ─────────────────────────────────────────
s9 = prs.slides.add_slide(BLANK)
bg(s9)
rect(s9, 0, 0, Inches(0.12), H, CYAN)

tb(s9, "MODULE 06B", Inches(0.55), Inches(0.3), Inches(5.0), Inches(0.35),
   size=10, bold=True, color=CYAN)
tb(s9, "Maintenance Management System",
   Inches(0.55), Inches(0.55), Inches(9.0), Inches(0.75),
   size=30, bold=True, color=WHITE)
aline(s9, Inches(1.4), color=CYAN, w=Inches(3.5))
tb(s9, "490 AMC clients. Each needs scheduled visits, reminders, and service records. This handles all of it.",
   Inches(0.55), Inches(1.5), Inches(9.5), Inches(0.45),
   size=13, color=LGREY)

# left — features
rect(s9, Inches(0.45), Inches(2.05), Inches(6.2), Inches(4.85), DARK_BLUE)
tb(s9, "What the system manages", Inches(0.65), Inches(2.2), Inches(5.8), Inches(0.4),
   size=13, bold=True, color=CYAN)
bullets(s9, [
    "All 490 AMC clients stored with service schedule",
    "Automated visit reminders to maintenance team",
    "Client notified 48h before visit — no surprise calls",
    "Post-visit report submitted by technician via WhatsApp",
    "Full service history per client — searchable anytime",
    "Parts used, chemicals added — logged per visit",
    "Rescheduling flow: client replies with alternate date",
    "AMC renewal: 30/7/1-day WhatsApp reminders + payment link",
    "New AMC signup: triggered automatically after construction handover",
], Inches(0.65), Inches(2.7), Inches(5.8), Inches(4.0), size=11, dot=CYAN)

# right — value cards
right_items = [
    ("490", "AMC clients on automated schedule", CYAN),
    ("0", "Manager calls needed for routine visit reminders", WA_GREEN),
    ("100%", "Service history logged — no paper records needed", GOLD),
    ("3x", "Faster AMC renewal conversion with automated follow-up", ORANGE),
]
rect(s9, Inches(6.9), Inches(2.05), Inches(5.95), Inches(4.85), DARK_BLUE)
for i, (val, label, col) in enumerate(right_items):
    cy2 = Inches(2.2 + i * 1.15)
    rect(s9, Inches(7.05), cy2, Inches(5.65), Inches(0.95), RGBColor(0x0C, 0x15, 0x38))
    tb(s9, val, Inches(7.15), cy2 + Inches(0.05), Inches(1.2), Inches(0.85),
       size=28, bold=True, color=col, align=PP_ALIGN.CENTER)
    tb(s9, label, Inches(8.45), cy2 + Inches(0.2), Inches(4.1), Inches(0.55),
       size=11, color=WHITE)

rect(s9, 0, Inches(7.05), W, Inches(0.45), RGBColor(0x06, 0x09, 0x1F))
tb(s9, "Genos AI  |  hello@genosai.tech  |  www.genosai.tech",
   Inches(0.4), Inches(7.08), Inches(12.0), Inches(0.35),
   size=9, color=MGREY, align=PP_ALIGN.CENTER)


# ── SLIDE 10 — IMPLEMENTATION + CTA ──────────────────────────────────────────
s10 = prs.slides.add_slide(BLANK)
bg(s10)
rect(s10, 0, 0, Inches(0.12), H, WA_GREEN)

tb(s10, "NEXT STEPS", Inches(0.55), Inches(0.3), Inches(5.0), Inches(0.35),
   size=10, bold=True, color=WA_GREEN)
tb(s10, "Implementation Timeline",
   Inches(0.55), Inches(0.55), Inches(9.0), Inches(0.75),
   size=30, bold=True, color=WHITE)
aline(s10, Inches(1.4), color=WA_GREEN, w=Inches(3.0))
tb(s10, "Fully operational in 3 weeks. Your team keeps using WhatsApp — the system activates around them.",
   Inches(0.55), Inches(1.5), Inches(9.5), Inches(0.45),
   size=13, color=LGREY)

# timeline blocks
weeks = [
    ("Week 1", "Core Operations Live",
     ["Daily worker check-in automation active",
      "Site update collection for all active sites",
      "Manager dashboard up with real-time data",
      "Attendance + report alert system running"],
     WA_GREEN),
    ("Week 2", "Task + Client System",
     ["Task assignment and tracking live",
      "Automated nudge system for all triggers",
      "Client update automation active",
      "Construction project update flows running"],
     CYAN),
    ("Week 3", "Maintenance + Full Go-Live",
     ["490 AMC clients loaded into system",
      "Maintenance schedule automation live",
      "AMC renewal reminder flows active",
      "Full system live — team trained in 1 session"],
     GOLD),
]

for i, (wk, title, items, col) in enumerate(weeks):
    cx2 = Inches(0.45 + i * 4.28); cy2 = Inches(2.1)
    cw3 = Inches(4.0); ch3 = Inches(4.0)
    rect(s10, cx2, cy2, cw3, ch3, DARK_BLUE)
    rect(s10, cx2, cy2, cw3, Inches(0.45), col)
    tb(s10, f"  {wk}  —  {title}", cx2, cy2, cw3, Inches(0.45),
       size=11, bold=True, color=NAVY)
    bullets(s10, items, cx2 + Inches(0.15), cy2 + Inches(0.55),
            cw3 - Inches(0.2), Inches(3.3), size=11, dot=col)

# bottom CTA
rect(s10, 0, Inches(6.2), W, Inches(0.85), DARK_BLUE)
tb(s10, "Ready to build this?",
   Inches(0.55), Inches(6.28), Inches(4.5), Inches(0.55),
   size=20, bold=True, color=WHITE)
tb(s10, "Call / WhatsApp: +91 638-714-2699   |   hello@genosai.tech   |   www.genosai.tech",
   Inches(0.55), Inches(6.68), Inches(10.0), Inches(0.35),
   size=12, color=LGREY)

tag(s10, "Rohan Malik, CEO — Genos AI", Inches(9.5), Inches(6.32), fill=WA_GREEN, size=11)

rect(s10, 0, Inches(7.05), W, Inches(0.45), RGBColor(0x06, 0x09, 0x1F))
tb(s10, "Genos AI  |  hello@genosai.tech  |  www.genosai.tech",
   Inches(0.4), Inches(7.08), Inches(12.0), Inches(0.35),
   size=9, color=MGREY, align=PP_ALIGN.CENTER)


# ── SAVE ─────────────────────────────────────────────────────────────────────
out = "AquaTechnic_WhatsApp_Ops_GenosAI.pptx"
prs.save(out)
print(f"Saved: {out}")
