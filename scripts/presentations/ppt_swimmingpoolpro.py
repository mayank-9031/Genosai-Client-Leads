import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.parent / "utils"))
from ppt_theme import (
    NAVY, DARK_BLUE, ACCENT, GOLD, RED, GREEN, PURPLE, ORANGE, TEAL, CYAN,
    WHITE, LGREY, MGREY, W, H,
    bg, rect, tb, aline, bullets, tag,
    slide_header, slide_footer, score_badge,
    stat_cards, top_bar_card, phase_card, before_after, genos_logo_block,
)

ROOT = Path(__file__).parent.parent.parent
prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]

# Client-specific colors
WA_GREEN  = RGBColor(0x25, 0xD3, 0x66)


def footer(slide, text="Genos AI  |  hello@genosai.tech  |  www.genosai.tech  |  +91 638-714-2699"):
    rect(slide, 0, Inches(7.05), W, Inches(0.45), DARK_BLUE)
    tb(slide, text, Inches(0.55), Inches(7.1), Inches(12.4), Inches(0.35),
       size=9, color=MGREY, align=PP_ALIGN.CENTER)


# ── SLIDE 1 — COVER ──────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(BLANK)
bg(s1)
rect(s1, Inches(9.8), 0, Inches(3.53), H, DARK_BLUE)
rect(s1, Inches(9.6), 0, Inches(0.12), H, CYAN)

# right panel accent
rect(s1, Inches(10.05), Inches(1.5), Inches(2.8), Inches(1.0), RGBColor(0x1A, 0x2E, 0x60))
tb(s1, "SCORE", Inches(10.05), Inches(1.53), Inches(2.8), Inches(0.35),
   size=10, bold=True, color=MGREY, align=PP_ALIGN.CENTER)
tb(s1, "6 / 10", Inches(10.05), Inches(1.82), Inches(2.8), Inches(0.6),
   size=26, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

rect(s1, Inches(10.05), Inches(2.7), Inches(2.8), Inches(0.75), RGBColor(0x2A, 0x08, 0x08))
tb(s1, "PRIORITY", Inches(10.05), Inches(2.73), Inches(2.8), Inches(0.3),
   size=9, bold=True, color=MGREY, align=PP_ALIGN.CENTER)
tb(s1, "HIGH", Inches(10.05), Inches(2.98), Inches(2.8), Inches(0.38),
   size=18, bold=True, color=RED, align=PP_ALIGN.CENTER)

rect(s1, Inches(10.05), Inches(3.6), Inches(2.8), Inches(0.75), RGBColor(0x06, 0x09, 0x1F))
tb(s1, "CRITICAL BUGS", Inches(10.05), Inches(3.63), Inches(2.8), Inches(0.3),
   size=9, bold=True, color=MGREY, align=PP_ALIGN.CENTER)
tb(s1, "4 FOUND", Inches(10.05), Inches(3.88), Inches(2.8), Inches(0.38),
   size=18, bold=True, color=RED, align=PP_ALIGN.CENTER)

# left content
tb(s1, "SWIMMINGPOOLPRO.IN", Inches(0.55), Inches(1.1), Inches(9.0), Inches(0.45),
   size=12, bold=True, color=CYAN)
aline(s1, Inches(1.65), color=CYAN, w=Inches(2.5))
tb(s1, "Website & AI Automation\nAudit", Inches(0.55), Inches(1.25), Inches(9.0), Inches(1.4),
   size=42, bold=True, color=WHITE)
tb(s1, "Aqua Technic  |  Hyderabad, Telangana  |  15+ Years  |  200+ Pools Built",
   Inches(0.55), Inches(2.75), Inches(9.0), Inches(0.45),
   size=13, color=LGREY)

aline(s1, Inches(3.3), color=CYAN, w=Inches(4.0))

# issue pills
issues = [
    ("Contact Page CRASH", RED),
    ("Pricing Gated", ORANGE),
    ("Nav Links 404", RED),
    ("Brand Inconsistency", ORANGE),
    ("No AI Qualification", GOLD),
    ("Static WhatsApp Widget", ACCENT),
]
for i, (label, col) in enumerate(issues):
    row = i // 3; xi = i % 3
    tag(s1, label, Inches(0.55 + xi * 3.0), Inches(3.55 + row * 0.45), fill=col, size=10)

tb(s1, "New React SPA launched Jan 2026. Modern stack. Critical pages broken at launch.",
   Inches(0.55), Inches(4.6), Inches(9.0), Inches(0.4),
   size=12, color=MGREY)

footer(s1)


# ── SLIDE 2 — COMPANY SNAPSHOT ───────────────────────────────────────────────
s2 = prs.slides.add_slide(BLANK)
bg(s2)
rect(s2, 0, 0, Inches(0.12), H, CYAN)

tb(s2, "COMPANY SNAPSHOT", Inches(0.55), Inches(0.3), Inches(12.0), Inches(0.4),
   size=11, bold=True, color=CYAN)
aline(s2, Inches(0.82), color=CYAN, w=Inches(2.5))
tb(s2, "What Aqua Technic / swimmingpoolpro.in is today",
   Inches(0.55), Inches(0.48), Inches(9.0), Inches(0.65),
   size=26, bold=True, color=WHITE)

# left — company data
rect(s2, Inches(0.45), Inches(1.35), Inches(5.9), Inches(5.4), DARK_BLUE)
tb(s2, "Company Profile", Inches(0.65), Inches(1.48), Inches(5.5), Inches(0.4),
   size=13, bold=True, color=CYAN)
data = [
    ("Brand",     "Aqua Technic (also uses 'PoolCraft' in blog)"),
    ("Website",   "swimmingpoolpro.in — new React SPA, Jan 2026"),
    ("Location",  "303 Surbhai Plaza, Lakdikapul, Hyderabad - 500004"),
    ("Phone 1",   "+91 93913 55379"),
    ("Phone 2",   "+91 93911 11962  (24/7 emergency)"),
    ("Office",    "040-23220724"),
    ("Email",     "info@swimmingpoolpro.in"),
    ("Hours",     "Mon–Fri 10:30am–5pm"),
    ("Track",     "15+ Years | 200+ Pools | Licensed & Insured"),
    ("Coverage",  "Telangana, AP, MP, Chhattisgarh, WB, Karnataka, TN, Maharashtra, Kerala, Delhi NCR"),
]
for i, (k, v) in enumerate(data):
    cy2 = Inches(1.95 + i * 0.45)
    tb(s2, k + ":", Inches(0.65), cy2, Inches(1.4), Inches(0.4),
       size=10, bold=True, color=MGREY)
    tb(s2, v, Inches(2.1), cy2, Inches(3.9), Inches(0.4),
       size=10, color=WHITE)

# right — services
rect(s2, Inches(6.65), Inches(1.35), Inches(6.2), Inches(2.55), DARK_BLUE)
tb(s2, "Services Offered", Inches(6.85), Inches(1.48), Inches(5.8), Inches(0.4),
   size=13, bold=True, color=CYAN)
bullets(s2, [
    "Infinity Edge, Architectural, Luxury, Therapeutic Pools",
    "Natural & Tropical Designs",
    "Prefabricated FRP / Fiberglass Pools",
    "Outdoor Living & Landscaping Integration",
    "Water Features — waterfalls, deck jets, fire & water",
    "Airport / Transit Hub commercial pools",
    "Equipment, Automation, Maintenance & Renovation",
], Inches(6.85), Inches(1.95), Inches(5.8), Inches(1.8), size=11, dot=CYAN)

# right — tech stack
rect(s2, Inches(6.65), Inches(4.1), Inches(6.2), Inches(2.65), DARK_BLUE)
tb(s2, "Website Tech Stack", Inches(6.85), Inches(4.23), Inches(5.8), Inches(0.4),
   size=13, bold=True, color=ACCENT)
bullets(s2, [
    "React SPA (Vite build — assets/vendor-DHGJfO9S.js)",
    "Client-side routing (React Router)",
    "WhatsApp chat widget present (static link)",
    "Newsletter subscription form",
    "Blog: 9 posts, good SEO content",
    "Google Business Profile linked in footer",
    "Financing calculator page",
], Inches(6.85), Inches(4.7), Inches(5.8), Inches(1.9), size=11, dot=ACCENT)

footer(s2)


# ── SLIDE 3 — WHAT'S WORKING ─────────────────────────────────────────────────
s3 = prs.slides.add_slide(BLANK)
bg(s3)
rect(s3, 0, 0, Inches(0.12), H, GREEN)

tb(s3, "WHAT'S WORKING", Inches(0.55), Inches(0.3), Inches(12.0), Inches(0.4),
   size=11, bold=True, color=GREEN)
aline(s3, Inches(0.82), color=GREEN, w=Inches(2.5))
tb(s3, "Strong foundations — modernised from aquatechnic.co.in",
   Inches(0.55), Inches(0.48), Inches(10.0), Inches(0.65),
   size=26, bold=True, color=WHITE)

goods = [
    ("WhatsApp Chat Widget", WA_GREEN,
     "Widget present on every page — 'Need help with your pool project? Chat with us on WhatsApp.' Correct for India market. Missing: AI qualification behind it."),
    ("Active Blog — 9 Posts", GREEN,
     "Strong SEO content: 'FRP Fibreglass Pools Hyderabad 2026', 'Swimming Pool Installation Cost Hyderabad 2026', maintenance guides, buyer guides. Active since Dec 2025."),
    ("Financing Page", ACCENT,
     "Detailed financing calculator: Standard (8.5% APR), Home Equity, 0% promo. Example monthly payments shown in INR. Approval process explained. Strong trust builder."),
    ("Multiple CTAs", GOLD,
     "Request Free Consultation, Get a Free Quote, Book Consultation — present throughout. Right density for high-ticket product."),
    ("Pool Type Range", CYAN,
     "7 pool types clearly categorised: Infinity, Architectural, FRP, Therapeutic, Natural, Luxury, Outdoor Living. Helps prospects self-identify."),
    ("Geographic Coverage", PURPLE,
     "Explicit coverage map: Telangana (33 districts), AP, MP, Chhattisgarh, WB, and major metro cities. Builds confidence in large-scale operations."),
]

for i, (title, col, desc) in enumerate(goods):
    row = i // 2; xi = i % 2
    cx = Inches(0.45 + xi * 6.43); cy = Inches(1.4 + row * 1.85)
    cw = Inches(6.15); ch = Inches(1.7)
    rect(s3, cx, cy, cw, ch, DARK_BLUE)
    rect(s3, cx, cy, Inches(0.07), ch, col)
    tag(s3, "GOOD", cx + Inches(0.18), cy + Inches(0.1), fill=col, size=8)
    tb(s3, title, cx + Inches(1.1), cy + Inches(0.1), cw - Inches(1.2), Inches(0.38),
       size=13, bold=True, color=WHITE)
    tb(s3, desc, cx + Inches(0.18), cy + Inches(0.55), cw - Inches(0.3), Inches(1.0),
       size=10, color=LGREY)

footer(s3)


# ── SLIDE 4 — CRITICAL BUGS ──────────────────────────────────────────────────
s4 = prs.slides.add_slide(BLANK)
bg(s4)
rect(s4, 0, 0, Inches(0.12), H, RED)

tb(s4, "CRITICAL BUGS — FIX IMMEDIATELY", Inches(0.55), Inches(0.3), Inches(12.0), Inches(0.4),
   size=11, bold=True, color=RED)
aline(s4, Inches(0.82), color=RED, w=Inches(3.5))
tb(s4, "4 broken pages losing prospects right now",
   Inches(0.55), Inches(0.48), Inches(10.0), Inches(0.65),
   size=26, bold=True, color=WHITE)

# Bug 1 — contact page crash
rect(s4, Inches(0.45), Inches(1.35), Inches(12.43), Inches(1.65), RGBColor(0x1A, 0x04, 0x04))
rect(s4, Inches(0.45), Inches(1.35), Inches(0.08), Inches(1.65), RED)
tb(s4, "BUG 1 — CRITICAL", Inches(0.65), Inches(1.42), Inches(3.0), Inches(0.35),
   size=10, bold=True, color=RED)
tb(s4, "Contact Page completely crashes — white screen + JS error",
   Inches(0.65), Inches(1.72), Inches(7.0), Inches(0.4),
   size=14, bold=True, color=WHITE)
tb(s4, 'Error: "ReferenceError: ArrowRight is not defined" in contact-B3CnmZCx.js:319',
   Inches(0.65), Inches(2.1), Inches(10.0), Inches(0.35),
   size=10, color=RGBColor(0xFF, 0xAA, 0xAA), italic=True)
tb(s4, "Impact: Every visitor who clicks Contact sees a blank crash page. Zero conversions from this page.",
   Inches(0.65), Inches(2.42), Inches(11.5), Inches(0.35),
   size=11, color=LGREY)
tag(s4, "URL: swimmingpoolpro.in/contact", Inches(8.5), Inches(1.72), fill=RGBColor(0x40, 0x08, 0x08), size=9)

# Bug 2 + 3 — 404 nav links
rect(s4, Inches(0.45), Inches(3.15), Inches(6.0), Inches(1.55), RGBColor(0x1A, 0x08, 0x04))
rect(s4, Inches(0.45), Inches(3.15), Inches(0.08), Inches(1.55), ORANGE)
tb(s4, "BUG 2 — HIGH", Inches(0.65), Inches(3.22), Inches(3.0), Inches(0.35),
   size=10, bold=True, color=ORANGE)
tb(s4, "/commercial → 404 Not Found",
   Inches(0.65), Inches(3.52), Inches(5.5), Inches(0.4),
   size=14, bold=True, color=WHITE)
tb(s4, "Nav link 'Commercial' prominently in header. Clicks go to 404. Hotel/resort clients lost immediately.",
   Inches(0.65), Inches(3.93), Inches(5.5), Inches(0.6),
   size=10, color=LGREY)

rect(s4, Inches(6.75), Inches(3.15), Inches(6.13), Inches(1.55), RGBColor(0x1A, 0x08, 0x04))
rect(s4, Inches(6.75), Inches(3.15), Inches(0.08), Inches(1.55), ORANGE)
tb(s4, "BUG 3 — HIGH", Inches(6.95), Inches(3.22), Inches(3.0), Inches(0.35),
   size=10, bold=True, color=ORANGE)
tb(s4, "/residential → 404 Not Found",
   Inches(6.95), Inches(3.52), Inches(5.7), Inches(0.4),
   size=14, bold=True, color=WHITE)
tb(s4, "Nav link 'Residential' also 404. Homeowner prospects click and bounce with no content.",
   Inches(6.95), Inches(3.93), Inches(5.7), Inches(0.6),
   size=10, color=LGREY)

# Bug 4 — brand inconsistency
rect(s4, Inches(0.45), Inches(4.85), Inches(12.43), Inches(1.6), RGBColor(0x12, 0x10, 0x04))
rect(s4, Inches(0.45), Inches(4.85), Inches(0.08), Inches(1.6), GOLD)
tb(s4, "BUG 4 — CREDIBILITY DAMAGE", Inches(0.65), Inches(4.92), Inches(4.0), Inches(0.35),
   size=10, bold=True, color=GOLD)
tb(s4, "Blog uses two brand names — 'Aqua Technic Team' and 'PoolCraft Team' as authors",
   Inches(0.65), Inches(5.22), Inches(10.0), Inches(0.4),
   size=14, bold=True, color=WHITE)
tb(s4, "Jan/Feb 2026 blog posts signed 'PoolCraft Team'. Dec 2025 posts signed 'Aqua Technic Team'. "
       "Prospect researching the company sees two different brand names — immediate trust damage.",
   Inches(0.65), Inches(5.65), Inches(11.8), Inches(0.65),
   size=11, color=LGREY)

footer(s4)


# ── SLIDE 5 — PRICING GATE PROBLEM ───────────────────────────────────────────
s5 = prs.slides.add_slide(BLANK)
bg(s5)
rect(s5, 0, 0, Inches(0.12), H, ORANGE)

tb(s5, "CONVERSION KILLER", Inches(0.55), Inches(0.3), Inches(12.0), Inches(0.4),
   size=11, bold=True, color=ORANGE)
aline(s5, Inches(0.82), color=ORANGE, w=Inches(2.5))
tb(s5, "Pricing hidden behind a signup wall",
   Inches(0.55), Inches(0.48), Inches(10.0), Inches(0.65),
   size=28, bold=True, color=WHITE)

# what the visitor sees
rect(s5, Inches(0.45), Inches(1.35), Inches(6.0), Inches(5.4), DARK_BLUE)
tb(s5, "What a prospect sees at /pricing", Inches(0.65), Inches(1.48), Inches(5.6), Inches(0.4),
   size=13, bold=True, color=ORANGE)

# simulate the signup gate
gate_box = s5.shapes.add_shape(9, Inches(0.65), Inches(2.0), Inches(5.6), Inches(3.6))
gate_box.fill.solid(); gate_box.fill.fore_color.rgb = RGBColor(0x14, 0x20, 0x50)
gate_box.line.color.rgb = ORANGE
tf = gate_box.text_frame; tf.word_wrap = True
tf.margin_left = Pt(15); tf.margin_top = Pt(12)
lines2 = [
    ("Create Account", 16, True, ORANGE),
    ("Sign up to view our pricing and packages", 11, False, WHITE),
    ("", 8, False, WHITE),
    ("First Name ___________", 10, False, LGREY),
    ("Last Name ___________", 10, False, LGREY),
    ("Email _______________", 10, False, LGREY),
    ("Phone (Optional) _____", 10, False, LGREY),
    ("Password ____________", 10, False, LGREY),
    ("", 8, False, WHITE),
    ("[ Create Account ]", 12, True, WHITE),
    ("Already have an account? Sign in", 9, False, MGREY),
]
for j, (txt, sz, bld, col2) in enumerate(lines2):
    p2 = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
    r2 = p2.add_run(); r2.text = txt
    r2.font.size = Pt(sz); r2.font.bold = bld; r2.font.color.rgb = col2
    p2.space_after = Pt(3)

# right — why this kills conversions
rect(s5, Inches(6.75), Inches(1.35), Inches(6.13), Inches(5.4), DARK_BLUE)
tb(s5, "Why this is a serious conversion problem", Inches(6.95), Inches(1.48), Inches(5.75), Inches(0.4),
   size=13, bold=True, color=ORANGE)
bullets(s5, [
    "Pool buyer's #1 question: 'How much does this cost?' — they hit a login wall",
    "Competitor with transparent pricing wins the comparison immediately",
    "Account creation = commitment signal — most visitors won't create an account just to see prices",
    "Pricing is already on homepage (₹40L / ₹65L / ₹1Cr examples in financing section) — inconsistency creates distrust",
    "B2C luxury buyer expects transparency, not a sales funnel gate",
    "Pricing is mentioned in blog posts as 'from ₹38L' — so it's not secret anyway",
    "Sign Up gate creates perception that company is hiding something",
], Inches(6.95), Inches(2.0), Inches(5.75), Inches(3.5), size=11, dot=ORANGE)

rect(s5, 0, Inches(6.85), W, Inches(0.2), ORANGE)
tb(s5, "Fix: Replace signup gate with transparent pricing tiers OR an AI cost estimator that captures email AFTER showing estimate.",
   Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.35),
   size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

footer(s5)


# ── SLIDE 6 — WHATSAPP: CURRENT VS POTENTIAL ─────────────────────────────────
s6 = prs.slides.add_slide(BLANK)
bg(s6)
rect(s6, 0, 0, Inches(0.12), H, WA_GREEN)

tb(s6, "WHATSAPP — PRESENT BUT PASSIVE", Inches(0.55), Inches(0.3), Inches(12.0), Inches(0.4),
   size=11, bold=True, color=WA_GREEN)
aline(s6, Inches(0.82), color=WA_GREEN, w=Inches(3.0))
tb(s6, "Widget exists. AI qualification doesn't.",
   Inches(0.55), Inches(0.48), Inches(10.0), Inches(0.65),
   size=28, bold=True, color=WHITE)

# current state
rect(s6, Inches(0.45), Inches(1.35), Inches(5.9), Inches(5.15), RGBColor(0x1A, 0x04, 0x04))
rect(s6, Inches(0.45), Inches(1.35), Inches(5.9), Inches(0.42), RGBColor(0x40, 0x10, 0x10))
tb(s6, "  Current — Static WhatsApp Button", Inches(0.45), Inches(1.35), Inches(5.9), Inches(0.42),
   size=12, bold=True, color=WHITE)
bullets(s6, [
    "Button present: 'Need help with your pool project?'",
    "Opens WhatsApp with blank chat to company number",
    "Prospect has to type their own question from scratch",
    "Response depends on who reads WhatsApp and when",
    "No qualification — no idea if prospect is homeowner, hotel, budget range",
    "No automated response if team is offline or busy",
    "No follow-up if prospect doesn't reply after first message",
    "All conversation data lost — no CRM logging",
], Inches(0.65), Inches(1.9), Inches(5.5), Inches(4.4), size=11, dot=RED)

# potential state
rect(s6, Inches(6.65), Inches(1.35), Inches(6.23), Inches(5.15), RGBColor(0x04, 0x14, 0x0C))
rect(s6, Inches(6.65), Inches(1.35), Inches(6.23), Inches(0.42), RGBColor(0x10, 0x40, 0x25))
tb(s6, "  With Genos AI — Intelligent WhatsApp Agent", Inches(6.65), Inches(1.35), Inches(6.23), Inches(0.42),
   size=12, bold=True, color=WHITE)
bullets(s6, [
    "Visitor clicks — automated greeting + qualification starts",
    "'Are you looking for a residential or commercial pool?'",
    "'What is your approximate budget? Under ₹40L / ₹40–80L / ₹80L+'",
    "'Which city? We cover Hyderabad, AP, Telangana and 10+ states'",
    "'When are you looking to build? 3 months / 6 months / exploring'",
    "Agent answers pricing, pool type, and FRP vs concrete questions 24/7",
    "Hot leads (clear budget + timeline) → auto-books site visit call",
    "All lead data → CRM automatically — name, project type, budget, timeline",
    "Weekend + night coverage — no missed inquiries",
], Inches(6.85), Inches(1.9), Inches(5.85), Inches(4.4), size=11, dot=WA_GREEN)

footer(s6)


# ── SLIDE 7 — MISSING AUTOMATION OPPORTUNITIES ───────────────────────────────
s7 = prs.slides.add_slide(BLANK)
bg(s7)
rect(s7, 0, 0, Inches(0.12), H, GOLD)

tb(s7, "AUTOMATION GAPS", Inches(0.55), Inches(0.3), Inches(12.0), Inches(0.4),
   size=11, bold=True, color=GOLD)
aline(s7, Inches(0.82), color=GOLD, w=Inches(2.2))
tb(s7, "Revenue leaking at 5 points right now",
   Inches(0.55), Inches(0.48), Inches(10.0), Inches(0.65),
   size=28, bold=True, color=WHITE)

gaps = [
    ("No Lead Qualification",
     "Consultation requests captured with zero context — no pool type, no budget, no timeline. Sales team goes in blind.",
     "AI chatbot qualifies before booking: pool type, budget, location, timeline → pre-qualified lead card in CRM",
     GOLD),
    ("No Follow-up Sequence",
     "Visitor submits quote request → goes silent → no automated Day 1/3/7 follow-up. Forgotten lead.",
     "Automated WhatsApp/email nurture: Day 1 'Here are 3 pools similar to what you described', Day 3 financing info, Day 7 soft close",
     CYAN),
    ("No AI Cost Estimator",
     "'How much does a pool cost?' is the #1 question. Answer is locked behind a signup form.",
     "Free instant estimator: pool type → size → finish → location → ballpark range shown. Email captured after estimate.",
     ACCENT),
    ("No Testimonials / Social Proof",
     "200+ pools built. Homepage shows portfolio images only — no client names, no reviews, no star ratings.",
     "Post-project WhatsApp review request sequence. Best reviews auto-embedded on homepage.",
     GREEN),
    ("No AMC / Maintenance Follow-up",
     "Post-construction client relationship drops off. No automated maintenance reminders.",
     "AMC onboarding: post-handover WhatsApp sequence offers annual maintenance package with WhatsApp scheduling.",
     ORANGE),
]

for i, (title, problem, solution, col) in enumerate(gaps):
    row = i // 2 if i < 4 else 2
    xi  = i % 2 if i < 4 else 0.5
    cx  = Inches(0.45 + xi * 6.43)
    cy  = Inches(1.35 + row * 1.9)
    if i == 4:
        cx = Inches(3.62); cy = Inches(5.2)
    cw = Inches(6.15) if i < 4 else Inches(6.15)
    ch = Inches(1.72)
    rect(s7, cx, cy, cw, ch, DARK_BLUE)
    rect(s7, cx, cy, Inches(0.07), ch, col)
    tb(s7, title, cx + Inches(0.18), cy + Inches(0.08), cw - Inches(0.25), Inches(0.35),
       size=12, bold=True, color=col)
    tb(s7, "Gap: " + problem, cx + Inches(0.18), cy + Inches(0.45), cw - Inches(0.25), Inches(0.55),
       size=10, color=LGREY, italic=True)
    tb(s7, "Fix: " + solution, cx + Inches(0.18), cy + Inches(1.02), cw - Inches(0.25), Inches(0.6),
       size=10, color=WHITE)

footer(s7)


# ── SLIDE 8 — ABOUT GENOS AI ─────────────────────────────────────────────────
s8 = prs.slides.add_slide(BLANK)
bg(s8)
rect(s8, 0, 0, Inches(0.12), H, ACCENT)

tb(s8, "ABOUT GENOS AI", Inches(0.55), Inches(0.3), Inches(12.0), Inches(0.4),
   size=11, bold=True, color=ACCENT)
aline(s8, Inches(0.82), color=ACCENT, w=Inches(2.0))
tb(s8, "We build AI automation for construction & service businesses",
   Inches(0.55), Inches(0.48), Inches(10.0), Inches(0.65),
   size=26, bold=True, color=WHITE)

# 4 service cards
services = [
    ("WhatsApp AI Agent", WA_GREEN,
     "Qualifies leads 24/7, answers pricing and product questions, books site visits automatically — responds in minutes, not hours."),
    ("Lead Qualification\n& Nurture Flows", ACCENT,
     "Multi-step sequences that warm cold prospects: Day 1/3/7 follow-ups, portfolio sharing, financing prompts — all automated."),
    ("AI Cost Estimator", GOLD,
     "Instant ballpark quote tool that captures email before showing estimate. Converts the #1 question into a qualified lead."),
    ("AMC & Maintenance\nAutomation", CYAN,
     "Post-construction client retention: AMC onboarding, renewal reminders (30/7/1 day), service scheduling via WhatsApp."),
]

for i, (title, col, desc) in enumerate(services):
    cx = Inches(0.45 + i * 3.2); cy = Inches(1.55)
    cw = Inches(3.0); ch = Inches(3.4)
    rect(s8, cx, cy, cw, ch, DARK_BLUE)
    rect(s8, cx, cy, cw, Inches(0.08), col)
    tb(s8, f"0{i+1}", cx + Inches(0.18), cy + Inches(0.18), cw - Inches(0.3), Inches(0.45),
       size=22, bold=True, color=col)
    tb(s8, title, cx + Inches(0.18), cy + Inches(0.68), cw - Inches(0.3), Inches(0.7),
       size=14, bold=True, color=WHITE)
    tb(s8, desc, cx + Inches(0.18), cy + Inches(1.48), cw - Inches(0.3), Inches(1.75),
       size=11, color=LGREY)

# bottom proof strip
rect(s8, 0, Inches(5.15), W, Inches(1.5), DARK_BLUE)
proof = [
    ("2 weeks", "Typical go-live\nfrom contract", ACCENT),
    ("24/7", "Automated\nresponse coverage", WA_GREEN),
    ("3-5x", "Lead conversion\nimprovement", GOLD),
    ("India-first", "WhatsApp-native\narchitecture", CYAN),
]
for i, (val, label, col) in enumerate(proof):
    px = Inches(1.65 + i * 2.55)
    tb(s8, val, px, Inches(5.3), Inches(2.3), Inches(0.6),
       size=26, bold=True, color=col, align=PP_ALIGN.CENTER)
    tb(s8, label, px, Inches(5.85), Inches(2.3), Inches(0.6),
       size=10, color=LGREY, align=PP_ALIGN.CENTER)

tb(s8, "Rohan Malik, CEO  |  hello@genosai.tech  |  www.genosai.tech  |  +91 638-714-2699",
   Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.35),
   size=11, color=MGREY, align=PP_ALIGN.CENTER)
footer(s8)


# ── SLIDE 9 — PRIORITY FIX ROADMAP ───────────────────────────────────────────
s9 = prs.slides.add_slide(BLANK)
bg(s9)
rect(s9, 0, 0, Inches(0.12), H, ACCENT)

tb(s9, "IMPLEMENTATION ROADMAP", Inches(0.55), Inches(0.3), Inches(12.0), Inches(0.4),
   size=11, bold=True, color=ACCENT)
aline(s9, Inches(0.82), color=ACCENT, w=Inches(3.0))
tb(s9, "What to fix and in what order",
   Inches(0.55), Inches(0.48), Inches(9.0), Inches(0.65),
   size=26, bold=True, color=WHITE)

phases = [
    ("WEEK 1\nImmediate Fixes", RED,
     ["Fix /contact page JS crash (import ArrowRight from lucide-react or swap icon)",
      "Build /commercial and /residential pages (currently 404 — nav is live)",
      "Remove signup gate from /pricing — show transparent tiers or redirect to estimator",
      "Fix blog brand names — standardise all author fields to 'Aqua Technic Team'"]),
    ("WEEK 2\nAI Lead System", WA_GREEN,
     ["Deploy Genos AI WhatsApp agent trained on pool types, pricing, FRP vs concrete, service areas",
      "Qualification flow: pool type → budget → location → timeline → site visit booking",
      "Day 1/3/7 automated follow-up for consultation form submissions",
      "Free pool cost estimator live on homepage with email gate before showing estimate"]),
    ("MONTH 2\nGrowth Layer", ACCENT,
     ["Customer review collection: WhatsApp request at project completion",
      "Best reviews embedded on homepage and service pages",
      "AMC onboarding sequence for all post-construction handovers",
      "Blog SEO: target 'swimming pool cost Hyderabad', 'FRP pool price India', 'pool builder near me'"]),
]

for i, (title, col, items) in enumerate(phases):
    cx = Inches(0.45 + i * 4.28); cy = Inches(1.35)
    cw = Inches(4.0); ch = Inches(5.5)
    rect(s9, cx, cy, cw, ch, DARK_BLUE)
    rect(s9, cx, cy, cw, Inches(0.55), col)
    tb(s9, title, cx, cy, cw, Inches(0.55),
       size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    bullets(s9, items, cx + Inches(0.15), cy + Inches(0.65),
            cw - Inches(0.2), Inches(4.7), size=11, dot=col)

footer(s9)


# ── SLIDE 10 — NEXT STEPS / CTA ──────────────────────────────────────────────
s10 = prs.slides.add_slide(BLANK)
bg(s10)

rect(s10, 0, 0, Inches(0.12), H, CYAN)
rect(s10, Inches(9.8), 0, Inches(3.53), H, DARK_BLUE)
rect(s10, Inches(9.6), 0, Inches(0.12), H, CYAN)

tb(s10, "NEXT STEPS", Inches(0.55), Inches(0.3), Inches(9.0), Inches(0.4),
   size=11, bold=True, color=CYAN)
aline(s10, Inches(0.82), color=CYAN, w=Inches(1.8))
tb(s10, "Swimmingpoolpro.in has\na solid foundation.",
   Inches(0.55), Inches(0.55), Inches(9.0), Inches(1.3),
   size=34, bold=True, color=WHITE)

tb(s10, "Modern design, active blog, WhatsApp widget, financing page — all real strengths. "
        "But 4 broken pages are losing prospects today, pricing friction is cutting conversions, "
        "and the WhatsApp widget is passive when it could be qualifying and booking automatically.",
   Inches(0.55), Inches(1.95), Inches(9.0), Inches(1.2),
   size=13, color=LGREY)

aline(s10, Inches(3.25), color=CYAN, w=Inches(4.0))

# summary metrics
tb(s10, "What we fix together:", Inches(0.55), Inches(3.4), Inches(9.0), Inches(0.4),
   size=13, bold=True, color=WHITE)

fixes = [
    ("1", "Contact crash, 404 nav links, pricing gate, brand confusion — Week 1"),
    ("2", "WhatsApp AI agent qualifying leads 24/7 — Week 2"),
    ("3", "Cost estimator + automated follow-up sequence — Week 2"),
    ("4", "Review automation + AMC sequence — Month 2"),
]
for i, (num, text) in enumerate(fixes):
    cy2 = Inches(3.9 + i * 0.52)
    rect(s10, Inches(0.55), cy2 + Inches(0.05), Inches(0.38), Inches(0.38), CYAN)
    tb(s10, num, Inches(0.55), cy2 + Inches(0.02), Inches(0.38), Inches(0.38),
       size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    tb(s10, text, Inches(1.05), cy2 + Inches(0.05), Inches(8.4), Inches(0.4),
       size=12, color=WHITE)

# right panel CTA
tb(s10, "Let's talk.", Inches(9.85), Inches(1.5), Inches(3.2), Inches(0.65),
   size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s10, "15 minutes to walk through\nthe audit and show you\nhow the WhatsApp agent works.",
   Inches(9.85), Inches(2.2), Inches(3.2), Inches(1.2),
   size=12, color=LGREY, align=PP_ALIGN.CENTER)

aline(s10, Inches(3.55), color=CYAN, w=Inches(2.0))

for label, val, col2 in [
    ("Call / WhatsApp", "+91 638-714-2699", WA_GREEN),
    ("Email", "hello@genosai.tech", ACCENT),
    ("Web", "www.genosai.tech", CYAN),
]:
    pass

tb(s10, "+91 638-714-2699", Inches(9.85), Inches(3.8), Inches(3.2), Inches(0.45),
   size=15, bold=True, color=WA_GREEN, align=PP_ALIGN.CENTER)
tb(s10, "hello@genosai.tech", Inches(9.85), Inches(4.3), Inches(3.2), Inches(0.4),
   size=12, color=ACCENT, align=PP_ALIGN.CENTER)
tb(s10, "www.genosai.tech", Inches(9.85), Inches(4.75), Inches(3.2), Inches(0.4),
   size=12, color=CYAN, align=PP_ALIGN.CENTER)

rect(s10, Inches(9.85), Inches(5.4), Inches(3.2), Inches(0.6), CYAN)
tb(s10, "Rohan Malik — CEO, Genos AI", Inches(9.85), Inches(5.42), Inches(3.2), Inches(0.55),
   size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

footer(s10)


# ── SAVE ─────────────────────────────────────────────────────────────────────
out = "SwimmingPoolPro_Audit_GenosAI.pptx"
prs.save(out)
print(f"Saved: {out}")
