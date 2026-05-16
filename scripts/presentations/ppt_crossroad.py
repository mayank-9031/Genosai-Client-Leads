import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.parent / "utils"))
from ppt_theme import (
    NAVY, DARK_BLUE, ACCENT, GOLD, RED, GREEN, PURPLE, ORANGE, TEAL, CYAN,
    WHITE, LGREY, MGREY, W, H,
    bg, rect, tb, aline, bullets, tag,
    slide_header, slide_footer, score_badge,
    stat_cards, top_bar_card, phase_card, before_after, genos_logo_block,
)

ROOT = Path(__file__).parent.parent.parent
prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]

# SLIDE 1 - COVER
s1 = prs.slides.add_slide(BLANK)
bg(s1)
rect(s1, 0, 0, Inches(0.12), H, TEAL)
rect(s1, Inches(8.5), 0, Inches(4.83), Inches(2.2), DARK_BLUE)
tb(s1, "GENOS AI", Inches(9.2), Inches(0.4), Inches(3.5), Inches(0.8),
   size=28, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
tb(s1, "www.genosai.tech", Inches(9.2), Inches(1.05), Inches(3.5), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
tb(s1, "CROSSROAD / ENERINVEST PROPTECH", Inches(0.55), Inches(1.6),
   Inches(9), Inches(0.6), size=13, bold=False, color=TEAL)
tb(s1, "Website & Sales Funnel\nAudit Report", Inches(0.55), Inches(2.05),
   Inches(9), Inches(1.8), size=44, bold=True, color=WHITE)
aline(s1, Inches(3.85), color=TEAL, w=Inches(3))
tb(s1, "Prepared exclusively for Mauro Terrinoni, CFA - CEO & Co-Founder",
   Inches(0.55), Inches(4.15), Inches(9), Inches(0.5), size=14, color=LGREY)
rect(s1, 0, Inches(6.5), W, Inches(1.0), DARK_BLUE)
tb(s1, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699",
   Inches(0.55), Inches(6.6), Inches(9), Inches(0.5), size=12, color=MGREY)
tb(s1, "May 2026  |  CONFIDENTIAL", Inches(9.5), Inches(6.6), Inches(3.3), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
rect(s1, Inches(10.5), Inches(3.3), Inches(2.3), Inches(2.5), DARK_BLUE)
tb(s1, "AUDIT SCORE", Inches(10.55), Inches(3.4), Inches(2.2), Inches(0.5),
   size=11, bold=True, color=MGREY, align=PP_ALIGN.CENTER)
tb(s1, "5", Inches(10.55), Inches(3.75), Inches(2.2), Inches(1.1),
   size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
tb(s1, "/ 10", Inches(10.55), Inches(4.7), Inches(2.2), Inches(0.4),
   size=18, color=LGREY, align=PP_ALIGN.CENTER)
tag(s1, "HIGH PRIORITY", Inches(10.75), Inches(5.1), fill=RED, size=10)


# SLIDE 2 - AT A GLANCE
s2 = prs.slides.add_slide(BLANK)
bg(s2)
rect(s2, 0, 0, Inches(0.12), H, TEAL)
tb(s2, "AT A GLANCE", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=TEAL)
tb(s2, "Mauro Terrinoni and Crossroad",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=30, bold=True, color=WHITE)
aline(s2, Inches(1.35), color=TEAL)

stats = [
    ("50 yrs", "Combined founder\nexperience (Mauro + Alex)"),
    ("CFA +\nMBA", "Cornell Johnson +\nJPMorgan/Citi/Fortress"),
    ("8 yrs", "PropTech product\ninnovation"),
    ("2020", "Crossroad founded\nLondon, UK"),
    ("2 pubs", "'PropTech can do for you'\n+ 'Digital era' article"),
]
cw = Inches(2.3); ch = Inches(2.2); gap = Inches(0.18)
for i, (num, label) in enumerate(stats):
    l = Inches(0.55) + i * (cw + gap)
    rect(s2, l, Inches(1.8), cw, ch, DARK_BLUE)
    tb(s2, num, l + Inches(0.1), Inches(1.95), cw - Inches(0.2), Inches(1.1),
       size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    tb(s2, label, l + Inches(0.1), Inches(3.05), cw - Inches(0.2), Inches(0.7),
       size=11, color=LGREY, align=PP_ALIGN.CENTER)

rect(s2, Inches(0.55), Inches(4.2), Inches(12.23), Inches(0.6), DARK_BLUE)
tb(s2, "36-37 Albert Embankment, London SE1 7TL  |  info@enerinvest.com  |  2-10 employees  |  PropTech / B2B SaaS",
   Inches(0.65), Inches(4.28), Inches(12), Inches(0.45), size=11, color=LGREY, align=PP_ALIGN.CENTER)

tb(s2, "MAURO'S BACKGROUND", Inches(0.55), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=TEAL)
bullets(s2, [
    "CFA charter holder (Oct 2004) + MBA Cornell Johnson",
    "Career: JPMorgan, Citi, Fortress, ACP Capital - IB + PE + distressed credit",
    "Executed trades in the hundreds of millions across European markets",
    "Published: 'Property Management enters the digital era' (Nov 2020)"],
        Inches(0.55), Inches(5.35), Inches(5.8), Inches(1.5), size=12, color=LGREY)

tb(s2, "CROSSROAD PLATFORM", Inches(6.8), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=TEAL)
bullets(s2, [
    "Asset, people, data, document, project management",
    "CRM for buyers, tenants, borrowers",
    "Branded property marketing website with built-in CRM",
    "Collaborative workspace for internal + external partners"],
        Inches(6.8), Inches(5.35), Inches(5.8), Inches(1.5), size=12, color=LGREY)


# SLIDE 3 - STRENGTHS
s3 = prs.slides.add_slide(BLANK)
bg(s3)
rect(s3, 0, 0, Inches(0.12), H, GREEN)
tb(s3, "WHAT'S WORKING", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s3, "Strong Institutional Credentials, a Real Product, and a Gap That Is the Pitch.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.8), size=22, bold=True, color=WHITE)
aline(s3, Inches(1.45), color=GREEN)

strengths = [
    ("Mauro Has Already Made the Case for Digitisation - In Writing",
     "Two published pieces arguing that property management must go digital give Mauro credibility with his target market and make the concept sell unnecessary. He is not asking property managers to believe in digital transformation - he is giving them a platform for it. The gap is applying the same argument to Crossroad's own customer acquisition."),
    ("CFA + JPMorgan/Citi/Fortress Background Means Rigorous ROI Standards",
     "Mauro evaluates investment arguments with institutional discipline. The pitch to him must be specific, data-led, and framed in terms of pipeline velocity and conversion rate improvement - not general AI enthusiasm. A CFA who ran distressed credit at Fortress will dismiss vague AI claims and respond to specific funnel metrics."),
    ("The Product Addresses a Real, Documented Pain",
     "Fragmented spreadsheets, disconnected service providers, manual document exchange, and siloed reporting are genuine operational problems for property portfolio managers. Crossroad's positioning around a single synoptic view of all assets, projects, and teams is correct. The product logic is sound. The gap is getting qualified evaluators into the funnel efficiently."),
    ("50 Years of Combined Experience Executing Large European Real Estate Trades",
     "Mauro and Alex have transacted hundreds of millions across European markets. They understand the operational complexity they are solving. That depth makes the product credible to institutional property managers who are the most valuable potential clients - and the most demanding evaluators during a sales process."),
    ("The Irony Is the Pitch - and It Is Not Hostile",
     "Pointing out that a PropTech company has a manual sales funnel is not a criticism of the product. It is a specific, accurate observation that only someone who looked carefully would make. Mauro published 'Property Management enters the digital era.' The argument applies to his own customer acquisition. He will recognise this immediately - and it opens a conversation rather than closing one."),
]
for i, (title, desc) in enumerate(strengths):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s3, Inches(0.55), t, Inches(0.06), Inches(0.7), GREEN)
    tb(s3, title, Inches(0.75), t, Inches(11.8), Inches(0.38), size=13, bold=True, color=WHITE)
    tb(s3, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.6), size=11, color=LGREY)


# SLIDE 4 - THE THREE GAPS
s4 = prs.slides.add_slide(BLANK)
bg(s4)
rect(s4, 0, 0, Inches(0.12), H, RED)
tb(s4, "THE THREE GAPS", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=RED)
tb(s4, "'Go Digital with Crossroad Today.' CTA Leads to an Email Inbox.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=24, bold=True, color=WHITE)
aline(s4, Inches(1.35), color=RED)

gap_cards = [
    ("GAP 1\nMANUAL\nSALES FUNNEL", RED,
     "The mechanism for digital transformation is an email",
     [
         "Homepage CTA: 'Go Digital with Crossroad Today' -> resolves to info@enerinvest.com",
         "No demo booking automation - prospect interest captured with zero qualifying context",
         "No AI chatbot to answer product questions in real time or qualify portfolio size/type",
         "No post-visit nurture - non-converting visitors disappear permanently",
     ]),
    ("GAP 2\nNO SOCIAL PROOF\nOR SEO CONTENT", GOLD,
     "8 years in market, 2 published pieces - none of it on the site",
     [
         "No case studies, no client logos, no outcome metrics despite 8 years operating",
         "Mauro's two LinkedIn articles not indexed as site content - SEO value entirely lost",
         "No pricing page: lengthens B2B sales cycles by forcing 1:1 qualification for every prospect",
         "No ROI calculator for property managers to self-qualify before requesting a demo",
     ]),
    ("GAP 3\nFRICTION IN THE\nEVALUATOR PATH", ACCENT,
     "No self-serve option for inbound evaluators",
     [
         "No free trial or freemium tier - every evaluator requires a full sales conversation",
         "A property manager comparing 5 SaaS options will choose the one they can try first",
         "No product tour, no explainer video, no interactive demo - evaluate by email only",
         "Institutional clients (Mauro's target) will evaluate rigorously - the evaluation process must match their standard",
     ]),
]
cw3 = Inches(3.95); gap3i = Inches(0.18)
for i, (title, color, subtitle, pts) in enumerate(gap_cards):
    l = Inches(0.55) + i * (cw3 + gap3i)
    rect(s4, l, Inches(1.8), cw3, Inches(0.5), color)
    tb(s4, title, l + Inches(0.12), Inches(1.85), cw3 - Inches(0.24), Inches(0.42),
       size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    rect(s4, l, Inches(2.3), cw3, Inches(4.6), DARK_BLUE)
    tb(s4, subtitle, l + Inches(0.12), Inches(2.38), cw3 - Inches(0.24), Inches(0.35),
       size=10, italic=True, color=color)
    bullets(s4, pts, l + Inches(0.12), Inches(2.78), cw3 - Inches(0.24), Inches(3.9),
            size=11, color=LGREY, dot=color)

rect(s4, Inches(0.55), Inches(7.05), Inches(12.23), Inches(0.3), DARK_BLUE)
tb(s4, "The product digitises property management. The sales funnel is still operating like it's 2010.",
   Inches(0.7), Inches(7.08), Inches(12), Inches(0.25), size=11, italic=True, color=GOLD)


# SLIDE 5 - WEBSITE DESIGN AUDIT
s5 = prs.slides.add_slide(BLANK)
bg(s5)
rect(s5, 0, 0, Inches(0.12), H, PURPLE)
tb(s5, "WEBSITE DESIGN AUDIT", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=PURPLE)
tb(s5, "Well Designed. Manually Operated.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.8), size=36, bold=True, color=WHITE)
aline(s5, Inches(1.45), color=PURPLE, w=Inches(1.5))

rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(3.6), DARK_BLUE)
rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(0.42), RED)
tb(s5, "CROSSROAD.APP TODAY", Inches(0.65), Inches(1.92), Inches(3.6), Inches(0.38),
   size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "CTA 'Go Digital Today' -> email inbox",
    "No AI chatbot or demo booking automation",
    "No qualifying intake (portfolio size, type, pain)",
    "No post-visit nurture sequence",
    "No pricing page",
    "No case studies or client logos",
    "LinkedIn articles not on site as blog",
    "No free trial or self-serve evaluation path",
], Inches(0.7), Inches(2.42), Inches(3.6), Inches(2.9), size=11, color=LGREY, dot=RED)

tb(s5, "->", Inches(4.45), Inches(3.6), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(3.6), DARK_BLUE)
rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(0.42), GOLD)
tb(s5, "WHAT B2B SAAS LEADERS DO", Inches(5.1), Inches(1.92), Inches(3.4), Inches(0.38),
   size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "AI chatbot qualifies + books demo automatically",
    "Automated demo booking with pre-qualification",
    "Pricing page: transparent tiers reduce friction",
    "Post-visit Day 1/3/7 nurture sequence",
    "Case studies with portfolio size + ROI metrics",
    "SEO blog indexed from thought leadership content",
    "Free trial or interactive demo for self-serve eval",
    "ROI calculator as lead magnet with email gate",
], Inches(5.1), Inches(2.42), Inches(3.4), Inches(2.9), size=11, color=LGREY, dot=GOLD)

tb(s5, "->", Inches(8.68), Inches(3.6), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(3.6), DARK_BLUE)
rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(0.42), GREEN)
tb(s5, "WHAT GENOS AI BUILDS", Inches(9.35), Inches(1.92), Inches(3.4), Inches(0.38),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "AI chatbot: qualifies by portfolio size, type, pain",
    "Automated demo booking replaces email CTA",
    "Post-visit Day 1/3/7 nurture sequence",
    "Pricing page with transparent starter tier",
    "Mauro's LinkedIn articles indexed as blog",
    "ROI calculator: hours saved per portfolio manager",
    "Case study template + first 2 client stories",
    "European SEO content: UK, Italian, Spanish markets",
], Inches(9.35), Inches(2.42), Inches(3.4), Inches(2.9), size=11, color=LGREY, dot=GREEN)

rect(s5, Inches(0.55), Inches(5.72), Inches(12.23), Inches(1.28), DARK_BLUE)
tb(s5, "SALES FUNNEL BENCHMARK:", Inches(0.7), Inches(5.78), Inches(2.4), Inches(0.38),
   size=11, bold=True, color=PURPLE)
stats2 = [
    ("0 sec", "Demo booking time\nwith AI chatbot"),
    ("3-5x", "More demos booked\nvs. email CTA"),
    ("Day 7", "Last automated\nnurture touchpoint"),
    ("2 pubs", "LinkedIn articles\nnot yet on site as SEO"),
]
sw = Inches(2.7); sl = Inches(3.2)
for i, (num, label) in enumerate(stats2):
    xl = sl + i * sw
    tb(s5, num, xl, Inches(5.75), Inches(2.5), Inches(0.55),
       size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s5, label, xl, Inches(6.22), Inches(2.5), Inches(0.7),
       size=10, color=LGREY, align=PP_ALIGN.CENTER)


# SLIDE 6 - AI AUTOMATION STACK
s6 = prs.slides.add_slide(BLANK)
bg(s6)
rect(s6, 0, 0, Inches(0.12), H, ACCENT)
tb(s6, "SALES AUTOMATION STACK", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s6, "Four Layers That Fill the Demo Pipeline Without Manual Work.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=26, bold=True, color=WHITE)
aline(s6, Inches(1.35), color=ACCENT)

opps = [
    ("AI Chatbot +\nDemo Booking", TEAL,
     "Replaces the email CTA. A visitor lands on crossroad.app, has questions about asset management scope or CRM functionality. An AI chatbot answers in real time, qualifies by portfolio size, property type (commercial, residential, mixed), geography, and current tooling. Qualified prospects book a demo directly. Non-qualified visitors are routed to content. Mauro and Alex only receive pre-qualified demo requests with full context attached."),
    ("Post-Visit\nNurture Sequence", GOLD,
     "A property manager who visits crossroad.app, reads the product pages, and does not immediately book is not saying no - they are comparing. A Day 1 follow-up email summarising the platform, Day 3 featuring a use case, and Day 7 with an ROI question re-engages them at the decision moment. For B2B SaaS where the sales cycle is weeks not minutes, the nurture sequence is where most deals are actually won."),
    ("Content &\nSEO Engine", GREEN,
     "Mauro published two articles that make the case for PropTech digitisation. Neither is indexed on crossroad.app. Importing them as blog posts, adding state- and market-specific property management guides (UK, Italian, Spanish markets - matching Mauro's language capabilities and European deal history), and building an ROI calculator targets property managers actively researching solutions - the highest-intent inbound traffic possible."),
    ("Social Proof\nInfrastructure", PURPLE,
     "8 years of PropTech work with 50 years of combined experience and European institutional trade history. None of it is surfaced as case studies, client logos, or outcome metrics on crossroad.app. A property manager evaluating 5 platforms will disproportionately weight the one with documented client results. A case study template capturing portfolio size, previous tooling, implementation time, and outcomes converts the existing track record into a conversion asset."),
]
cw2 = Inches(2.95); gap2 = Inches(0.16)
for i, (title, color, desc) in enumerate(opps):
    l2 = Inches(0.55) + i * (cw2 + gap2)
    rect(s6, l2, Inches(1.8), cw2, Inches(0.06), color)
    rect(s6, l2, Inches(1.86), cw2, Inches(4.7), DARK_BLUE)
    tb(s6, title, l2 + Inches(0.12), Inches(1.98), cw2 - Inches(0.24), Inches(0.65),
       size=13, bold=True, color=color)
    tb(s6, desc, l2 + Inches(0.12), Inches(2.65), cw2 - Inches(0.24), Inches(3.65),
       size=11, color=LGREY)

rect(s6, Inches(0.55), Inches(6.7), Inches(12.23), Inches(0.45), DARK_BLUE)
tb(s6, "Mauro's job is building and selling Crossroad. The automation layer handles prospect qualification, nurturing, and re-engagement.",
   Inches(0.7), Inches(6.75), Inches(12), Inches(0.35), size=11, italic=True, color=GOLD)


# SLIDE 7 - THE IRONY SLIDE
s7 = prs.slides.add_slide(BLANK)
bg(s7)
rect(s7, 0, 0, Inches(0.12), H, ORANGE)
tb(s7, "THE CORE ARGUMENT", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ORANGE)
tb(s7, "What Mauro Wrote vs. How Crossroad Acquires Customers.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=28, bold=True, color=WHITE)
aline(s7, Inches(1.35), color=ORANGE, w=Inches(2.5))

rect(s7, Inches(0.55), Inches(1.65), Inches(12.23), Inches(1.35), DARK_BLUE)
tb(s7, "\"Property management has traditionally been plagued by inefficiencies, fragmented data, and manual processes. Technology is the answer.\"",
   Inches(0.75), Inches(1.73), Inches(11.8), Inches(1.1),
   size=16, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
tb(s7, "- Mauro Terrinoni, 'Property Management enters the digital era', November 2020",
   Inches(0.75), Inches(2.75), Inches(11.8), Inches(0.3),
   size=11, color=MGREY, align=PP_ALIGN.CENTER)

rect(s7, Inches(0.55), Inches(3.2), Inches(5.9), Inches(3.45), DARK_BLUE)
tb(s7, "CROSSROAD'S PITCH TO PROPERTY MANAGERS", Inches(0.7), Inches(3.27), Inches(5.6), Inches(0.35),
   size=11, bold=True, color=ORANGE)
bullets(s7, [
    "Centralise fragmented asset data into one dashboard",
    "Replace manual document exchange with digital workflows",
    "Eliminate disconnected spreadsheets across team members",
    "Automate project and task tracking across service providers",
    "Build a CRM that connects buyers, tenants, and borrowers",
    "Go digital with your property operations today",
], Inches(0.7), Inches(3.65), Inches(5.6), Inches(2.6), size=12, color=LGREY, dot=ORANGE)

rect(s7, Inches(6.65), Inches(3.2), Inches(6.13), Inches(3.45), DARK_BLUE)
tb(s7, "HOW CROSSROAD ACQUIRES ITS OWN CUSTOMERS", Inches(6.8), Inches(3.27), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=RED)
bullets(s7, [
    "Prospect reads homepage, wants a demo",
    "CTA: 'Go Digital with Crossroad Today'",
    "Clicks CTA -> opens email client -> writes to info@enerinvest.com",
    "Waits for a human to reply and schedule a call",
    "No chatbot, no qualifier, no nurture if they don't follow through",
    "No pricing, no trial, no self-serve evaluation path",
], Inches(6.8), Inches(3.65), Inches(5.8), Inches(2.6), size=12, color=LGREY, dot=RED)

rect(s7, Inches(0.55), Inches(6.8), Inches(12.23), Inches(0.38), DARK_BLUE)
tb(s7, "The same inefficiency Crossroad eliminates for property managers is present in Crossroad's own sales process. Genos AI closes that gap.",
   Inches(0.7), Inches(6.85), Inches(12), Inches(0.28), size=11, italic=True, color=ORANGE)


# SLIDE 8 - IMPLEMENTATION PLAN
s8 = prs.slides.add_slide(BLANK)
bg(s8)
rect(s8, 0, 0, Inches(0.12), H, GREEN)
tb(s8, "IMPLEMENTATION PLAN", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s8, "Funnel Automation First. Content Second. Social Proof Third.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=28, bold=True, color=WHITE)
aline(s8, Inches(1.35), color=GREEN)

phases = [
    ("PHASE 1", "Weeks 1-2", GOLD, [
        ("AI Chatbot + Demo Booking", "Chatbot trained on Crossroad features, qualifying by portfolio size, property type, geography, current tools. Automated demo booking replaces email CTA. Mauro receives pre-qualified leads only."),
        ("Post-Visit Nurture Live", "Day 1/3/7 email sequence for non-converting visitors. B2B SaaS pipeline is won at touchpoint 3-5, not on first visit."),
        ("LinkedIn Articles as Blog", "Mauro's two published pieces imported and indexed as site blog. Immediate SEO value from existing content - zero additional writing required."),
    ]),
    ("PHASE 2", "Weeks 3-5", ACCENT, [
        ("Pricing Page", "Transparent starter tier (even 'contact for enterprise' with a clear entry point). Reduces friction for inbound evaluators comparing multiple platforms."),
        ("ROI Calculator", "Hours saved per property manager per month based on portfolio size and current tooling. Email gate before results. High-intent lead capture from self-qualifying property managers."),
        ("Case Study Template", "2-3 client stories with portfolio size, previous tooling, implementation time, outcomes. Institutional evaluators weight documented results heavily."),
    ]),
    ("PHASE 3", "Months 2-3", GREEN, [
        ("European SEO Content", "Property management guides for UK, Italian, Spanish markets - matching Mauro's language coverage and European deal geography. 10-15 articles targeting portfolio manager search terms."),
        ("Chatbot Data Analysis", "First 60 days of qualifying data refines ICP. What portfolio size converts? What property type? Which pain point closes fastest?"),
        ("Full Pipeline Reporting", "Demo booking rate, nurture sequence conversion, SEO traffic and ranking. CFA-level ROI reporting on the automation investment."),
    ]),
]
pw = Inches(3.9); gap4 = Inches(0.18)
for i, (phase, timing, color, items) in enumerate(phases):
    l3 = Inches(0.55) + i * (pw + gap4)
    rect(s8, l3, Inches(1.8), pw, Inches(0.52), color)
    tb(s8, phase, l3 + Inches(0.12), Inches(1.84), pw * 0.5 - Inches(0.12), Inches(0.42),
       size=13, bold=True, color=NAVY)
    tb(s8, timing, l3 + pw * 0.45, Inches(1.88), pw * 0.55 - Inches(0.15), Inches(0.38),
       size=11, color=NAVY, align=PP_ALIGN.RIGHT)
    rect(s8, l3, Inches(2.32), pw, Inches(4.6), DARK_BLUE)
    y_item = Inches(2.42)
    for title, desc in items:
        rect(s8, l3 + Inches(0.12), y_item + Inches(0.06), Inches(0.05), Inches(0.52), color)
        tb(s8, title, l3 + Inches(0.25), y_item, pw - Inches(0.38), Inches(0.32),
           size=12, bold=True, color=color)
        tb(s8, desc, l3 + Inches(0.25), y_item + Inches(0.3), pw - Inches(0.38), Inches(0.45),
           size=11, color=LGREY)
        y_item += Inches(0.95)


# SLIDE 9 - BEFORE vs AFTER
s9 = prs.slides.add_slide(BLANK)
bg(s9)
rect(s9, 0, 0, Inches(0.12), H, ACCENT)
tb(s9, "BEFORE vs AFTER", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s9, "What Actually Changes for Crossroad in 90 Days",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=28, bold=True, color=WHITE)
aline(s9, Inches(1.35))

comparisons = [
    ("Demo Request Process",
     "Prospect clicks 'Go Digital with Crossroad Today', opens email client, writes to info@, waits for a human reply, and schedules a call with no context captured.",
     "AI chatbot qualifies in 5 questions. Portfolio size, property type, geography, current tools, main pain. Pre-qualified demo booked automatically. Mauro's first call is substantive."),
    ("Non-Converting Visitors",
     "A property manager visits crossroad.app, reads the product pages, does not immediately email. No follow-up. They compare 4 other platforms and choose the one that stayed in touch.",
     "Day 1/3/7 nurture sequence re-engages. Day 3 use case matches their property type. Day 7 ROI question arrives at the decision moment. B2B deals close in the nurture, not the first visit."),
    ("Thought Leadership SEO",
     "Mauro published 'Property Management enters the digital era' and 'What PropTech can do for you' on LinkedIn. Neither is indexed as site content. SEO value = zero.",
     "Both articles live as indexed blog posts on crossroad.app. Property managers searching 'property management digital transformation' find Mauro's writing before a competitor's."),
    ("Evaluator Social Proof",
     "8 years in market with 50 years combined founder experience and institutional trade history. Zero case studies or client outcomes on the site. Evaluators have no reference points.",
     "2-3 client case studies with portfolio size, previous tooling, implementation time, and outcome metrics. ROI calculator for self-qualification. Institutional evaluators have evidence."),
    ("European Market SEO",
     "Mauro speaks English, Italian, and Spanish fluently and has executed deals across European markets. Zero content targeting UK, Italian, or Spanish property management search terms.",
     "10-15 articles targeting property management keywords in UK and European markets. Mauro's language capability and deal geography become content distribution channels."),
]

lw = Inches(5.4)
rect(s9, Inches(0.55), Inches(1.65), lw, Inches(0.35), RED)
tb(s9, "TODAY", Inches(0.65), Inches(1.68), lw - Inches(0.2), Inches(0.28),
   size=11, bold=True, color=WHITE)
rect(s9, Inches(6.45), Inches(1.65), lw + Inches(0.43), Inches(0.35), GREEN)
tb(s9, "AFTER 90 DAYS", Inches(6.55), Inches(1.68), lw, Inches(0.28),
   size=11, bold=True, color=WHITE)

for i, (topic, before, after) in enumerate(comparisons):
    t = Inches(2.1) + i * Inches(1.02)
    rect(s9, Inches(0.55), t, Inches(12.23), Inches(0.22), DARK_BLUE)
    tb(s9, topic.upper(), Inches(0.65), t + Inches(0.02), Inches(12), Inches(0.2),
       size=9, bold=True, color=MGREY)
    tb(s9, before, Inches(0.65), t + Inches(0.24), lw - Inches(0.2), Inches(0.66),
       size=11, color=LGREY)
    tb(s9, after, Inches(6.55), t + Inches(0.24), lw, Inches(0.66),
       size=11, color=WHITE)


# SLIDE 10 - NEXT STEPS
s10 = prs.slides.add_slide(BLANK)
bg(s10)
rect(s10, 0, 0, Inches(0.12), H, TEAL)
rect(s10, Inches(7.5), 0, Inches(5.83), H, DARK_BLUE)
tb(s10, "NEXT STEPS", Inches(0.55), Inches(0.3), Inches(6.5), Inches(0.5),
   size=11, bold=True, color=TEAL)
tb(s10, "Let's Talk for\n15 Minutes.",
   Inches(0.55), Inches(0.8), Inches(6.8), Inches(2.0), size=48, bold=True, color=WHITE)
aline(s10, Inches(2.8), color=TEAL, w=Inches(2.5))
tb(s10, "No pitch deck needed. Reply or book a call.\nFocused on demo automation and nurture sequence first.",
   Inches(0.55), Inches(3.05), Inches(6.5), Inches(0.8), size=14, color=LGREY)

steps = [
    "Reply to this email - we'll set up a time",
    "Start: AI chatbot + automated demo booking (2 weeks)",
    "Add post-visit nurture + LinkedIn articles as blog",
    "Months 2-3: ROI calculator, case studies, EU SEO",
]
bullets(s10, steps, Inches(0.55), Inches(3.95), Inches(6.5), Inches(2.0), size=14, color=WHITE, dot=TEAL)

rect(s10, 0, Inches(6.5), Inches(7.5), Inches(1.0), DARK_BLUE)
tb(s10, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699  |  www.genosai.tech",
   Inches(0.55), Inches(6.6), Inches(6.8), Inches(0.5), size=11, color=MGREY)

tb(s10, "THREE GAPS.\nTHREE FIXES.", Inches(7.85), Inches(0.6), Inches(5.0), Inches(1.1),
   size=22, bold=True, color=TEAL)
aline(s10, Inches(1.65), color=TEAL)

summary_items = [
    ("Manual Funnel", "AI chatbot + automated demo booking replaces email CTA"),
    ("No Nurture", "Day 1/3/7 sequence: B2B deals close at touchpoint 3-5"),
    ("Buried Content", "LinkedIn articles indexed as SEO blog (zero new writing)"),
    ("No Social Proof", "Case studies + ROI calculator for institutional evaluators"),
    ("No Pricing", "Transparent starter tier reduces B2B sales cycle length"),
    ("Zero EU SEO", "UK, Italian, Spanish property management content engine"),
]
for i, (label, val) in enumerate(summary_items):
    t2 = Inches(1.85) + i * Inches(0.88)
    rect(s10, Inches(7.85), t2, Inches(1.5), Inches(0.72), NAVY)
    tb(s10, label, Inches(7.9), t2 + Inches(0.05), Inches(1.4), Inches(0.32),
       size=10, bold=True, color=TEAL)
    tb(s10, val, Inches(9.5), t2 + Inches(0.05), Inches(3.6), Inches(0.65),
       size=11, color=LGREY)


prs.save(ROOT / "output" / "decks" / "Crossroad_Audit_GenosAI.pptx")
print("Saved: Crossroad_Audit_GenosAI.pptx")
