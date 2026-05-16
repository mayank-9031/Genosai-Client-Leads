import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.parent / "utils"))
from ppt_theme import (
    NAVY, DARK_BLUE, ACCENT, GOLD, RED, GREEN, PURPLE, ORANGE, TEAL, CYAN,
    WHITE, LGREY, MGREY, W, H,
    bg, rect, tb, aline, bullets, tag,
    slide_header, slide_footer, score_badge,
    stat_cards, top_bar_card, phase_card, before_after, genos_logo_block,
)

ROOT = Path(__file__).parent.parent.parent
prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]
# Client-specific colors
COBALT = RGBColor(0x1A, 0x4F, 0x8A)
DGREY  = RGBColor(0x2C, 0x38, 0x4A)
AMBER  = RGBColor(0xF3, 0x9C, 0x12)
SLATE  = RGBColor(0x2E, 0x45, 0x5F)


# ── Local helpers ────────────────────────────────────────────────────────────

def header(slide, title, subtitle, accent=GOLD):
    rect(slide, 0, 0, Inches(0.7), H, accent)
    tb(slide, title.upper(), Inches(0.55 + 0.7), Inches(0.3),
       Inches(12), Inches(0.45), size=11, bold=True, color=accent)
    tb(slide, subtitle, Inches(0.55 + 0.7), Inches(0.72),
       Inches(12 - 0.7), Inches(0.65), size=22, bold=True, color=WHITE)
    aline(slide, Inches(1.38), color=accent, w=Inches(1.5))


def footer(slide,
           text='Genos AI  |  hello@genosai.tech  |  www.genosai.tech  |  +91 638-714-2699'):
    rect(slide, 0, Inches(7.05), W, Inches(0.45), DARK_BLUE)
    tb(slide, text, Inches(0.55), Inches(7.1), Inches(12.4), Inches(0.35),
       size=9, color=MGREY, align=PP_ALIGN.CENTER)


def card(slide, x, y, w, h, fill=DARK_BLUE):
    return rect(slide, x, y, w, h, fill)


def stat_card(slide, val, label, x, y, w=Inches(3.0), h=Inches(1.6), accent=GOLD):
    rect(slide, x, y, w, h, DARK_BLUE)
    rect(slide, x, y, w, Inches(0.05), accent)
    tb(slide, str(val), x + Inches(0.1), y + Inches(0.1),
       w - Inches(0.2), Inches(0.9), size=28, bold=True, color=accent,
       align=PP_ALIGN.CENTER)
    tb(slide, label, x + Inches(0.1), y + Inches(1.1),
       w - Inches(0.2), Inches(0.42), size=10, color=LGREY, align=PP_ALIGN.CENTER)


def bullet_card(slide, title, items, x, y, w, h,
                accent=GOLD, title_bg=None):
    if title_bg is None:
        title_bg = DARK_BLUE
    rect(slide, x, y, w, h, DARK_BLUE)
    rect(slide, x, y, w, Inches(0.42), title_bg)
    rect(slide, x, y, Inches(0.05), h, accent)
    tb(slide, title, x + Inches(0.1), y + Inches(0.05),
       w - Inches(0.15), Inches(0.35), size=11, bold=True, color=WHITE)
    bx = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.48),
                                  w - Inches(0.15), h - Inches(0.55))
    tf = bx.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r0 = p.add_run(); r0.text = '  '; r0.font.size = Pt(9)
        r1 = p.add_run(); r1.text = '● '; r1.font.size = Pt(9)
        r1.font.color.rgb = accent; r1.font.bold = True
        r2 = p.add_run(); r2.text = str(item)
        r2.font.size = Pt(10.5); r2.font.color.rgb = LGREY
        if i < len(items) - 1:
            p.space_after = Pt(5)




# SLIDE 1 – TITLE
# ════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
bg(s1, NAVY)
rect(s1, 0, 0, Inches(0.7), H, GOLD)
rect(s1, Inches(0.7), Inches(4.5), W - Inches(0.7), Inches(0.06), GOLD)

tb(s1, "FORTRESS REAL ESTATE INVESTMENTS LIMITED  |  JSE: FFA & FFB",
   Inches(1.1), Inches(1.0), Inches(11.5), Inches(0.6),
   size=13, bold=True, color=GOLD)
tb(s1, "Web Development &\nAI Automation Audit",
   Inches(1.1), Inches(1.65), Inches(11), Inches(2.0),
   size=44, bold=True, color=WHITE)
tb(s1, "Digital Strategy Review & Innovation Roadmap  |  Johannesburg, South Africa",
   Inches(1.1), Inches(3.68), Inches(10), Inches(0.5),
   size=16, color=GOLD, italic=True)
tb(s1, "Prepared by: Genos Apollo  |  May 2026",
   Inches(1.1), Inches(5.0), Inches(8), Inches(0.4),
   size=12, color=MGREY)
tb(s1, "CONFIDENTIAL",
   Inches(1.1), Inches(6.6), Inches(4), Inches(0.35),
   size=10, color=MGREY, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 2 – AGENDA
# ════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
bg(s2)
header(s2, "Agenda", "What we will cover today")
footer(s2)

agenda = [
    ("01", "Company Digital Overview",    "Fortress's current digital footprint vs. REIT peers"),
    ("02", "Website Audit Findings",      "IR, leasing UX, SEO and platform performance gaps"),
    ("03", "Web Development Roadmap",     "Tenant portal, IR hub & modern platform stack"),
    ("04", "AI Automation Opportunity",   "Where AI drives leasing, IR and operations at scale"),
    ("05", "Benefits of Automation",      "ROI at R29.6B market cap scale"),
    ("06", "Implementation Timeline",     "Phased 12-month delivery plan"),
    ("07", "Investment & Next Steps",     "Budget guidance and immediate actions"),
]

for i, (num, title, desc) in enumerate(agenda):
    col = i % 2
    row = i // 2
    x = Inches(0.55) + col * Inches(6.4)
    y = Inches(1.55) + row * Inches(1.22)
    card(s2, x, y, Inches(6.1), Inches(1.05))
    rect(s2, x, y, Inches(0.7), Inches(1.05), NAVY)
    tb(s2, num, x, y + Inches(0.25), Inches(0.7), Inches(0.5),
       size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s2, title, x + Inches(0.8), y + Inches(0.1),
       Inches(5.0), Inches(0.4), size=13, bold=True, color=NAVY)
    tb(s2, desc, x + Inches(0.8), y + Inches(0.52),
       Inches(5.0), Inches(0.42), size=10.5, color=MGREY, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 3 – COMPANY DIGITAL OVERVIEW
# ════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
bg(s3)
header(s3, "Company Digital Overview", "Fortress Real Estate Investments – current digital footprint")
footer(s3)

card(s3, Inches(0.4), Inches(1.4), Inches(5.8), Inches(5.5))
rect(s3, Inches(0.4), Inches(1.4), Inches(5.8), Inches(0.45), NAVY)
tb(s3, "Company Snapshot", Inches(0.55), Inches(1.44),
   Inches(5.5), Inches(0.38), size=13, bold=True, color=WHITE)

snapshot = [
    ("Structure",     "Internally managed REIT  |  3rd largest in South Africa"),
    ("JSE Listing",   "FFA & FFB  |  Market cap: R29.6 billion"),
    ("CEO",           "Steven Brown"),
    ("HQ",            "Johannesburg, Gauteng, South Africa"),
    ("Portfolio",     "2.9M sqm GLA  |  4.0% vacancy  |  52 retail centres"),
    ("Segments",      "Logistics (SA + CEE)  |  Convenience Retail  |  NEPI Rockcastle 23.9%"),
    ("Dividend",      "7.36% yield  |  86.29 cps FFB H1 2025"),
    ("Website",       "fortressfund.co.za  (JS-rendered — SEO risk)"),
]
ty = Inches(2.02)
for k, v in snapshot:
    tb(s3, k, Inches(0.6), ty, Inches(1.55), Inches(0.3),
       size=10.5, bold=True, color=NAVY)
    tb(s3, v, Inches(2.2), ty, Inches(3.8), Inches(0.3),
       size=10.5, color=MGREY)
    ty += Inches(0.52)

card(s3, Inches(6.5), Inches(1.4), Inches(6.5), Inches(5.5))
rect(s3, Inches(6.5), Inches(1.4), Inches(6.5), Inches(0.45), GOLD)
tb(s3, "Key Digital Health Indicators", Inches(6.65), Inches(1.44),
   Inches(6.2), Inches(0.38), size=13, bold=True, color=NAVY)

metrics = [
    ("Custom Domain",               "✓ Active",   "fortressfund.co.za",                    GREEN),
    ("JS Rendering (SEO Risk)",     "⚠ SPA",      "JS-only = content invisible to crawlers",RED),
    ("Mobile Page Speed",           "~50 / 100",  "Heavy JS bundle — est. below benchmark", RED),
    ("Investor Relations Hub",      "Basic",      "SENS/reports exist; UX lags peers",      AMBER),
    ("Logistics Leasing Platform",  "Limited",    "No live availability or enquiry portal",  RED),
    ("AI / Automation Visible",     "None",       "No chatbot, no self-serve IR tools",      RED),
    ("ESG / Sustainability Section","Partial",    "Content exists; not interactively scored",AMBER),
    ("Social Media Integration",    "Limited",    "LinkedIn + Facebook present; no live feed",AMBER),
]
ty2 = Inches(2.02)
for label, val, note, col in metrics:
    tb(s3, label, Inches(6.65), ty2, Inches(2.7), Inches(0.3),
       size=10.5, bold=True, color=NAVY)
    tb(s3, val,   Inches(9.4), ty2, Inches(1.25), Inches(0.3),
       size=10.5, bold=True, color=col, align=PP_ALIGN.CENTER)
    tb(s3, note,  Inches(10.7), ty2, Inches(2.1), Inches(0.3),
       size=9, color=MGREY, italic=True)
    ty2 += Inches(0.52)

# ════════════════════════════════════════════════════════════════
# SLIDE 4 – WEBSITE AUDIT FINDINGS
# ════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
bg(s4)
header(s4, "Website Audit Findings", "fortressfund.co.za – gaps vs. global REIT digital benchmarks")
footer(s4)

audit_cards = [
    ("Performance & Tech",
     ["JS-rendered SPA — Google cannot index content", "Heavy bundle slows mobile load below 50/100",
      "No server-side rendering for critical IR pages", "Risk: institutional ESG screeners can't parse data"]),
    ("Investor Relations UX",
     ["SENS announcements buried — not auto-aggregated", "No live share price widget (FFA & FFB) on homepage",
      "Annual / interim reports require deep navigation", "No investor FAQ chatbot or self-serve IR tools"]),
    ("Leasing & Tenant Acquisition",
     ["No searchable logistics or retail availability database", "Prospective tenants cannot filter by GLA, location, spec",
      "No enquiry-to-lease pipeline or self-serve portal", "CEE portfolio barely visible to European prospects"]),
    ("ESG & Compliance",
     ["ESG dashboard not machine-readable for fund screeners", "POPIA cookie consent not prominently displayed",
      "Sustainability KPIs not dynamically updated", "No TCFD or carbon reporting interactive module"]),
]

for i, (title, bullets) in enumerate(audit_cards):
    x = Inches(0.4) + i * Inches(3.2)
    y = Inches(1.45)
    card(s4, x, y, Inches(3.05), Inches(5.3))
    rect(s4, x, y, Inches(3.05), Inches(0.45), NAVY)
    tb(s4, title, x + Inches(0.1), y + Inches(0.06),
       Inches(2.85), Inches(0.36), size=12, bold=True, color=WHITE)
    ty = y + Inches(0.62)
    for b in bullets:
        rect(s4, x + Inches(0.15), ty + Inches(0.07), Inches(0.1), Inches(0.1), GOLD)
        tb(s4, b, x + Inches(0.35), ty, Inches(2.6), Inches(0.55),
           size=10, color=NAVY)
        ty += Inches(0.72)

rect(s4, Inches(0.4), Inches(6.8), Inches(12.5), Inches(0.06), GOLD)
tb(s4, "Critical: SPA rendering blocks search engine indexing — Growthpoint and Redefine both serve SSR IR pages; Fortress is losing organic institutional traffic",
   Inches(0.4), Inches(6.86), Inches(12.5), Inches(0.3),
   size=10.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 5 – WEB DEVELOPMENT ROADMAP
# ════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
bg(s5)
header(s5, "Web Development Roadmap", "Elevating fortressfund.co.za to global REIT digital standards")
footer(s5)

card(s5, Inches(0.4), Inches(1.4), Inches(4.5), Inches(5.6))
rect(s5, Inches(0.4), Inches(1.4), Inches(4.5), Inches(0.45), COBALT)
tb(s5, "Recommended Tech Stack", Inches(0.55), Inches(1.44),
   Inches(4.2), Inches(0.38), size=13, bold=True, color=WHITE)

stack = [
    ("Frontend",     "Next.js 15 – SSR/SSG (full SEO & IR indexing)"),
    ("CMS",          "Sanity.io – IR docs, properties, ESG content"),
    ("Hosting",      "Vercel + Cloudflare CDN (SA + EU / CEE edge)"),
    ("Search",       "Algolia – logistics & retail availability search"),
    ("IR Data",      "JSE SENS API feed + live share price widget"),
    ("Portal",       "Auth0 – institutional investor & tenant portal"),
    ("Analytics",    "GA4 + Hotjar — IR funnel & leasing journey"),
    ("AI / Chat",    "Claude API – IR chatbot + leasing qualification"),
]
ty = Inches(2.02)
for k, v in stack:
    tb(s5, k, Inches(0.6), ty, Inches(1.4), Inches(0.3),
       size=10, bold=True, color=COBALT)
    tb(s5, v, Inches(2.05), ty, Inches(2.75), Inches(0.3),
       size=10, color=NAVY)
    ty += Inches(0.5)

card(s5, Inches(5.2), Inches(1.4), Inches(7.75), Inches(5.6))
rect(s5, Inches(5.2), Inches(1.4), Inches(7.75), Inches(0.45), NAVY)
tb(s5, "Key Improvements & Features",
   Inches(5.35), Inches(1.44), Inches(7.4), Inches(0.38),
   size=13, bold=True, color=WHITE)

improvements = [
    ("SSR Migration — Full IR & SEO Visibility",
     "Convert SPA to Next.js SSR/SSG. Every SENS page, property listing, and ESG report becomes fully indexed. Institutional investors and analysts find Fortress in organic search."),
    ("Live Investor Relations Hub",
     "Auto-aggregated SENS feed, live FFA/FFB share price ticker, interactive dividend history, downloadable annual and interim reports — all accessible within two clicks from homepage."),
    ("Logistics & Retail Availability Portal",
     "Searchable database of available GLA by province, spec, size and expected occupation. Prospective tenants enquire via portal — no cold-calling facilities teams."),
    ("ESG / Sustainability Interactive Dashboard",
     "Machine-readable TCFD-aligned ESG KPIs, live carbon intensity data, and downloadable ESG reports. Meets institutional fund screener requirements for JSE SRI and GRESB."),
    ("CEE Portfolio International Showcase",
     "Dedicated English-language section for Central & Eastern European logistics assets — targeting EU-based PE firms, logistics operators and sovereign wealth funds."),
    ("POPIA & Cookie Consent Compliance",
     "Full POPIA-compliant consent flow, privacy policy, data processing agreements and security headers — protecting Fortress from regulatory risk as a JSE-listed entity."),
]
ty2 = Inches(1.98)
for title, desc in improvements:
    rect(s5, Inches(5.35), ty2 + Inches(0.1), Inches(0.08), Inches(0.25), GOLD)
    tb(s5, title, Inches(5.55), ty2, Inches(7.1), Inches(0.28),
       size=11, bold=True, color=NAVY)
    tb(s5, desc, Inches(5.55), ty2 + Inches(0.3), Inches(7.1), Inches(0.35),
       size=10, color=MGREY)
    ty2 += Inches(0.78)

# ════════════════════════════════════════════════════════════════
# SLIDE 6 – AI AUTOMATION OPPORTUNITY MAP
# ════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
bg(s6)
header(s6, "AI Automation Opportunity Map", "Where AI delivers institutional-scale value for Fortress")
footer(s6)

opps = [
    (NAVY,    "Investor Relations Automation",
     ["AI chatbot answers IR queries 24/7 (dividends, SENS, NAV)", "Auto-summarise and publish SENS announcements",
      "Personalised IR email sequences by investor profile", "Earnings call transcript → key metrics auto-extracted"]),
    (GOLD,    "Leasing & Tenant Acquisition",
     ["AI qualifies inbound leasing enquiries by size & spec fit", "Auto-match available GLA to tenant requirements",
      "Heads of lease drafting assistant (AI-powered)", "Vacancy alert emails auto-sent to registered tenant prospects"]),
    (COBALT,  "Portfolio & Asset Management",
     ["AI-generated monthly portfolio performance summaries", "Predictive lease expiry & renewal risk flagging",
      "Automated valuation variance reports per asset", "CEE market intelligence digest (EU logistics trends)"]),
    (TEAL,    "ESG & Reporting Automation",
     ["Auto-compile ESG KPIs from property management systems", "AI-draft TCFD and GRESB submissions from internal data",
      "Carbon intensity trend alerts by portfolio segment", "Board ESG pack auto-generated monthly"]),
]

for i, (col, title, bullets) in enumerate(opps):
    c = i % 2
    r = i // 2
    x = Inches(0.4) + c * Inches(6.5)
    y = Inches(1.45) + r * Inches(2.85)
    card(s6, x, y, Inches(6.2), Inches(2.65))
    rect(s6, x, y, Inches(6.2), Inches(0.45), col)
    tb(s6, title, x + Inches(0.15), y + Inches(0.06),
       Inches(5.9), Inches(0.36), size=13, bold=True,
       color=NAVY if col == GOLD else WHITE)
    ty = y + Inches(0.62)
    for b in bullets:
        rect(s6, x + Inches(0.18), ty + Inches(0.08), Inches(0.1), Inches(0.1), col)
        tb(s6, b, x + Inches(0.38), ty, Inches(5.65), Inches(0.38),
           size=10.5, color=NAVY)
        ty += Inches(0.48)

# ════════════════════════════════════════════════════════════════
# SLIDE 7 – BENEFITS OF AI AUTOMATION
# ════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
bg(s7)
header(s7, "Benefits of AI Automation", "Institutional-scale impact across Fortress's operations")
footer(s7)

stats = [
    ("24/7",    "AI-powered IR &\nleasing enquiry coverage"),
    ("50%",     "Faster SENS &\nreport publication cycle"),
    ("3×",      "More tenant leads\nqualified per month"),
    ("R29.6B",  "Market cap scale —\nevery basis point counts"),
]
for i, (val, label) in enumerate(stats):
    stat_card(s7, val, label,
              Inches(0.4) + i * Inches(3.2), Inches(1.45),
              w=Inches(3.0), h=Inches(1.6))

benefits = [
    ("Investor Relations",
     ["Analysts find Fortress via Google (SSR unlocks this)", "IR chatbot answers dividend queries instantly",
      "SENS auto-summarised = media picks up faster", "Institutional buyers engage deeper with live data"]),
    ("Leasing & Occupancy",
     ["Available GLA visible 24/7 to logistics operators", "AI pre-qualifies tenants before BD team engages",
      "Vacancy periods shortened — direct NAV impact", "CEE assets marketed to EU capital via digital portal"]),
    ("ESG & Governance",
     ["Machine-readable ESG data satisfies fund mandates", "GRESB and JSE SRI submissions automated",
      "Board packs prepared in hours, not days", "Regulatory risk reduced across POPIA and GDPR (CEE)"]),
    ("Competitive Positioning",
     ["Match Growthpoint + Redefine's digital IR standard", "AI-native REIT attracts next-gen institutional allocators",
      "Global logistics operators find Fortress digitally", "Data flywheel: AI improves with every lease interaction"]),
]

for i, (title, bullets) in enumerate(benefits):
    c = i % 2
    r = i // 2
    x = Inches(0.4) + c * Inches(6.5)
    y = Inches(3.2) + r * Inches(1.85)
    bullet_card(s7, title, bullets, x, y, Inches(6.2), Inches(1.65),
                title_bg=NAVY if r == 0 else COBALT)

# ════════════════════════════════════════════════════════════════
# SLIDE 8 – ROI & BUSINESS CASE
# ════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
bg(s8)
header(s8, "ROI & Business Case", "Projected returns from digital + AI investment (ZAR)")
footer(s8)

card(s8, Inches(0.4), Inches(1.4), Inches(5.8), Inches(5.6))
rect(s8, Inches(0.4), Inches(1.4), Inches(5.8), Inches(0.45), NAVY)
tb(s8, "Indicative Investment (Year 1 – ZAR)",
   Inches(0.55), Inches(1.44), Inches(5.5), Inches(0.38),
   size=13, bold=True, color=WHITE)

invest_rows = [
    ("SSR Website Rebuild (Next.js + Sanity)",   "R 600,000 – R 1,200,000"),
    ("Live IR Hub + SENS API Integration",        "R 350,000 – R   700,000"),
    ("Logistics & Retail Leasing Portal",         "R 400,000 – R   800,000"),
    ("AI IR Chatbot (Claude API)",                "R 150,000 – R   300,000"),
    ("ESG Dashboard + TCFD Module",               "R 250,000 – R   500,000"),
    ("POPIA Compliance & Security Stack",         "R  80,000 – R   160,000"),
    ("Ongoing monthly support (est.)",            "R  80,000 – R   150,000 / mo"),
    ("TOTAL YEAR 1 (once-off + 12 mo)",          "R 2,790,000 – R 5,460,000"),
]
ty = Inches(2.0)
for i, (item, cost) in enumerate(invest_rows):
    is_total = i == len(invest_rows) - 1
    if is_total:
        rect(s8, Inches(0.4), ty - Inches(0.05), Inches(5.8), Inches(0.06), GOLD)
        ty += Inches(0.12)
        tb(s8, item, Inches(0.6), ty, Inches(3.5), Inches(0.35),
           size=11, bold=True, color=NAVY)
        tb(s8, cost, Inches(4.05), ty, Inches(1.95), Inches(0.35),
           size=11, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)
    else:
        tb(s8, item, Inches(0.6), ty, Inches(3.5), Inches(0.3),
           size=10, color=NAVY)
        tb(s8, cost, Inches(4.05), ty, Inches(1.95), Inches(0.3),
           size=10, color=MGREY, align=PP_ALIGN.RIGHT)
    ty += Inches(0.5)

card(s8, Inches(6.5), Inches(1.4), Inches(6.5), Inches(5.6))
rect(s8, Inches(6.5), Inches(1.4), Inches(6.5), Inches(0.45), GREEN)
tb(s8, "Projected Returns & Value Unlocked",
   Inches(6.65), Inches(1.44), Inches(6.2), Inches(0.38),
   size=13, bold=True, color=WHITE)

returns = [
    ("1 new logistics tenant (5,000 sqm @ R85/sqm)",  "R 5.1M+ annual rental income"),
    ("Vacancy reduction: 0.5% of 2.9M sqm GLA",       "R 12M–R 25M+ annual rental upside"),
    ("IR: reduced analyst roadshow admin cost",        "R 1.5M–R 3M / year est."),
    ("ESG rating uplift → lower cost of capital",      "10–20 bps saving on R29.6B portfolio"),
    ("AI lease admin time saving (10 staff × 20%)",    "R 2M–R 4M annual salary equivalent"),
    ("CEE digital presence → EU institutional inflows","Long-term AUM growth at fund level"),
    ("TOTAL YEAR 1 VALUE (conservative)",              "R 20M – R 45M+ identifiable upside"),
]
ty2 = Inches(2.0)
for i, (item, val) in enumerate(returns):
    is_total = i == len(returns) - 1
    if is_total:
        rect(s8, Inches(6.5), ty2 - Inches(0.05), Inches(6.5), Inches(0.06), GOLD)
        ty2 += Inches(0.12)
        tb(s8, item, Inches(6.65), ty2, Inches(3.9), Inches(0.35),
           size=11, bold=True, color=NAVY)
        tb(s8, val, Inches(10.6), ty2, Inches(2.2), Inches(0.35),
           size=11, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)
    else:
        tb(s8, item, Inches(6.65), ty2, Inches(3.9), Inches(0.3),
           size=10, color=NAVY)
        tb(s8, val, Inches(10.6), ty2, Inches(2.2), Inches(0.3),
           size=10, color=MGREY, align=PP_ALIGN.RIGHT)
    ty2 += Inches(0.5)

tb(s8, "Investment of R2.8M–R5.5M against R20M–R45M+ identifiable return = 7–16× ROI — and that excludes long-term ESG and AUM compounding",
   Inches(0.4), Inches(7.07), Inches(12.5), Inches(0.28),
   size=10.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 9 – IMPLEMENTATION TIMELINE
# ════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(BLANK)
bg(s9)
header(s9, "Implementation Timeline", "Phased 12-month delivery plan")
footer(s9)

phases = [
    ("Phase 1\nMon 1–3",  GOLD,     [
        "POPIA compliance & cookie consent",
        "SSR migration — unlock Google indexing",
        "Live FFA/FFB share price on homepage",
        "SENS auto-aggregation feed live",
    ]),
    ("Phase 2\nMon 4–6",  NAVY,     [
        "Full Next.js + Sanity CMS rebuild",
        "Logistics & retail availability portal",
        "IR Hub: reports, dividends, results centre",
        "AI IR chatbot (Claude API) MVP",
    ]),
    ("Phase 3\nMon 7–9",  COBALT,   [
        "ESG interactive dashboard + TCFD module",
        "CEE portfolio international showcase",
        "Tenant enquiry-to-pipeline portal",
        "AI leasing qualification workflows",
    ]),
    ("Phase 4\nMon 9–12", GREEN,    [
        "GRESB & JSE SRI auto-reporting",
        "AI portfolio performance summaries",
        "Institutional investor portal (Auth0)",
        "Full audit, optimise & benchmarking",
    ]),
]

for i, (label, col, items) in enumerate(phases):
    x = Inches(0.4) + i * Inches(3.2)
    y = Inches(1.45)
    card(s9, x, y, Inches(3.0), Inches(5.55))
    rect(s9, x, y, Inches(3.0), Inches(0.7), col)
    tb(s9, label, x, y + Inches(0.06), Inches(3.0), Inches(0.62),
       size=13, bold=True,
       color=NAVY if col == GOLD else WHITE,
       align=PP_ALIGN.CENTER)
    if i < 3:
        rect(s9, x + Inches(3.02), y + Inches(0.3), Inches(0.16), Inches(0.06), MGREY)
    ty = y + Inches(0.88)
    for item in items:
        rect(s9, x + Inches(0.18), ty + Inches(0.08), Inches(0.1), Inches(0.1), col)
        tb(s9, item, x + Inches(0.38), ty, Inches(2.5), Inches(0.42),
           size=10.5, color=NAVY)
        ty += Inches(0.9)

# ════════════════════════════════════════════════════════════════
# SLIDE 10 – WHY ACT NOW?
# ════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(BLANK)
bg(s10, NAVY)
rect(s10, 0, 0, Inches(0.5), H, GOLD)
footer(s10)

tb(s10, "Why Act — Now?",
   Inches(0.8), Inches(0.5), Inches(11), Inches(0.65),
   size=30, bold=True, color=GOLD)
tb(s10, "At R29.6 billion, every digital inefficiency compounds. The gap between Fortress and global REIT digital benchmarks is measurable in basis points.",
   Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.55),
   size=14, color=WHITE, italic=True)

reasons = [
    ("The SPA rendering issue is leaking institutional traffic today",
     "Fortress's JavaScript-only site is invisible to Google search crawlers for key queries: 'logistics property South Africa', 'SA REIT dividend yield', 'CEE logistics investment'. Growthpoint and Redefine serve fully-indexed SSR IR pages. Fixing this in Phase 1 recaptures that traffic immediately."),
    ("SA REITs returned 5.9% in April 2026 — investor interest is surging",
     "With rate cuts anticipated through 2026, institutional capital is rotating back into REITs. Fortress needs a world-class digital IR experience to capture these inflows before they commit to better-presented peers on the JSE."),
    ("Logistics vacancy is a direct NAV lever — digital leasing scales it",
     "At 2.9M sqm GLA and 4.0% vacancy, Fortress has ~116,000 sqm of unlet space. A searchable digital availability portal, combined with AI lead qualification, directly compresses that vacancy figure — each 0.5% improvement is worth R12M–R25M in annual rental income."),
    ("ESG mandates are now a capital allocation prerequisite",
     "Institutional fund managers operating under SFDR, PRI, and JSE SRI mandates require machine-readable ESG data. An interactive TCFD dashboard and GRESB-aligned reporting module directly influences Fortress's inclusion in ESG-screened portfolios — which represent a growing share of global institutional AUM."),
    ("The CEE portfolio is under-marketed to European capital",
     "Fortress holds a 23.9% stake in NEPI Rockcastle and has direct CEE logistics exposure. Without a dedicated English-language digital showcase targeting EU private equity, logistics operators, and sovereign wealth funds, this premium European positioning generates no inbound deal flow."),
]
ty = Inches(1.95)
for i, (title, desc) in enumerate(reasons):
    rect(s10, Inches(0.8), ty, Inches(0.55), Inches(0.55),
         GOLD if i % 2 == 0 else COBALT)
    tb(s10, str(i + 1), Inches(0.8), ty, Inches(0.55), Inches(0.55),
       size=16, bold=True, color=NAVY if i % 2 == 0 else WHITE,
       align=PP_ALIGN.CENTER)
    tb(s10, title, Inches(1.5), ty, Inches(11), Inches(0.3),
       size=12, bold=True, color=WHITE)
    tb(s10, desc, Inches(1.5), ty + Inches(0.3), Inches(11.3), Inches(0.55),
       size=10.5, color=RGBColor(0xBB, 0xCC, 0xDD))
    ty += Inches(0.98)

# ════════════════════════════════════════════════════════════════
# SLIDE 11 – NEXT STEPS
# ════════════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(BLANK)
bg(s11)
header(s11, "Next Steps", "From audit to action — what happens now")
footer(s11)

steps = [
    ("Week 1",  GOLD,    "Technical Discovery",
     "Genos Apollo conducts full technical audit: SSR gap analysis, SENS feed mapping, GLA data structure, existing CMS review and IR content inventory."),
    ("Week 2",  NAVY,    "Stakeholder Alignment",
     "Present findings to Steven Brown and digital / IR team. Agree Phase 1 scope: POPIA fix, SSR migration priority and live share price integration."),
    ("Week 3",  COBALT,  "Phase 1 Sprint Begins",
     "POPIA consent live. SSR migration starts (no design change required). SENS feed integrated. Google re-indexing begins within days of deployment."),
    ("Month 2", GREEN,   "IR Hub & Portal Live",
     "Fully indexed IR hub live. Logistics availability portal in testing. AI IR chatbot MVP deployed. Fortress visible in Google for key REIT queries."),
]

for i, (when, col, title, desc) in enumerate(steps):
    x = Inches(0.4) + i * Inches(3.2)
    y = Inches(1.5)
    card(s11, x, y, Inches(3.0), Inches(4.0))
    rect(s11, x, y, Inches(3.0), Inches(0.45), col)
    tb(s11, when, x, y + Inches(0.06), Inches(3.0), Inches(0.36),
       size=12, bold=True,
       color=NAVY if col == GOLD else WHITE,
       align=PP_ALIGN.CENTER)
    tb(s11, title, x + Inches(0.15), y + Inches(0.58),
       Inches(2.7), Inches(0.36), size=12, bold=True,
       color=COBALT if col == GOLD else col)
    tb(s11, desc, x + Inches(0.15), y + Inches(1.0),
       Inches(2.7), Inches(2.8), size=10.5, color=NAVY)

card(s11, Inches(0.4), Inches(5.8), Inches(12.5), Inches(1.1), fill=NAVY)
tb(s11, "Ready to bring Fortress's digital presence in line with its R29.6B market position?",
   Inches(0.6), Inches(5.88), Inches(7.8), Inches(0.4),
   size=14, bold=True, color=GOLD)
tb(s11, "Contact Genos Apollo to schedule your technical discovery session.",
   Inches(0.6), Inches(6.28), Inches(7), Inches(0.38),
   size=11, color=WHITE)
tb(s11, "fortressfund.co.za  |  JSE: FFA & FFB  |  Johannesburg",
   Inches(8.3), Inches(6.0), Inches(4.7), Inches(0.38),
   size=10, color=GOLD, italic=True, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════
# SLIDE 12 – THANK YOU / BACK COVER
# ════════════════════════════════════════════════════════════════
s12 = prs.slides.add_slide(BLANK)
bg(s12, NAVY)
rect(s12, 0, 0, Inches(0.7), H, GOLD)
rect(s12, Inches(0.7), Inches(3.5), W - Inches(0.7), Inches(0.06), GOLD)

tb(s12, "Thank You",
   Inches(1.1), Inches(2.0), Inches(11), Inches(1.2),
   size=52, bold=True, color=WHITE)
tb(s12, "Premium logistics and retail real estate — built for the digital investor era.",
   Inches(1.1), Inches(3.1), Inches(10), Inches(0.55),
   size=16, color=GOLD, italic=True)
tb(s12, "Fortress Real Estate Investments Limited  ·  JSE: FFA & FFB  ·  fortressfund.co.za",
   Inches(1.1), Inches(4.0), Inches(10), Inches(0.4),
   size=12, color=RGBColor(0xAA, 0xBB, 0xCC))
tb(s12, "Prepared by Genos Apollo  |  May 2026  |  Confidential",
   Inches(1.1), Inches(6.7), Inches(10), Inches(0.35),
   size=10, color=MGREY)

# ── Save ────────────────────────────────────────────────────────
OUT = ROOT / "output" / "decks" / "Fortress_WebDev_AI_Automation_Audit_GenosAI.pptx"
prs.save(str(OUT))
print(f"Saved: {OUT}")
