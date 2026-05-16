import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.parent / "utils"))
from ppt_theme import (
    NAVY, DARK_BLUE, ACCENT, GOLD, RED, GREEN, PURPLE, ORANGE, TEAL, CYAN,
    WHITE, LGREY, MGREY, W, H,
    bg, rect, tb, aline, bullets, tag,
    slide_header, slide_footer, score_badge,
    stat_cards, top_bar_card, phase_card, before_after, genos_logo_block,
)

ROOT = Path(__file__).parent.parent.parent
prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]

# SLIDE 1 - COVER
s1 = prs.slides.add_slide(BLANK)
bg(s1)
rect(s1, 0, 0, Inches(0.12), H, PURPLE)
rect(s1, Inches(8.5), 0, Inches(4.83), Inches(2.2), DARK_BLUE)
tb(s1, "GENOS AI", Inches(9.2), Inches(0.4), Inches(3.5), Inches(0.8),
   size=28, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
tb(s1, "www.genosai.tech", Inches(9.2), Inches(1.05), Inches(3.5), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
tb(s1, "NOTARY WORKFORCE", Inches(0.55), Inches(1.6),
   Inches(9), Inches(0.6), size=13, bold=False, color=PURPLE)
tb(s1, "Website & AI Growth\nAudit Report", Inches(0.55), Inches(2.05),
   Inches(9), Inches(1.8), size=48, bold=True, color=WHITE)
aline(s1, Inches(3.85), color=PURPLE, w=Inches(3))
tb(s1, "Prepared exclusively for June James, CNTDA, CNSA - Founder",
   Inches(0.55), Inches(4.15), Inches(9), Inches(0.5), size=14, color=LGREY)
rect(s1, 0, Inches(6.5), W, Inches(1.0), DARK_BLUE)
tb(s1, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699",
   Inches(0.55), Inches(6.6), Inches(9), Inches(0.5), size=12, color=MGREY)
tb(s1, "May 2026  |  CONFIDENTIAL", Inches(9.5), Inches(6.6), Inches(3.3), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
rect(s1, Inches(10.5), Inches(3.3), Inches(2.3), Inches(2.5), DARK_BLUE)
tb(s1, "AUDIT SCORE", Inches(10.55), Inches(3.4), Inches(2.2), Inches(0.5),
   size=11, bold=True, color=MGREY, align=PP_ALIGN.CENTER)
tb(s1, "5", Inches(10.55), Inches(3.75), Inches(2.2), Inches(1.1),
   size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
tb(s1, "/ 10", Inches(10.55), Inches(4.7), Inches(2.2), Inches(0.4),
   size=18, color=LGREY, align=PP_ALIGN.CENTER)
tag(s1, "HIGH PRIORITY", Inches(10.75), Inches(5.1), fill=RED, size=10)


# SLIDE 2 - AT A GLANCE
s2 = prs.slides.add_slide(BLANK)
bg(s2)
rect(s2, 0, 0, Inches(0.12), H, PURPLE)
tb(s2, "AT A GLANCE", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=PURPLE)
tb(s2, "June James and Notary Workforce",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=30, bold=True, color=WHITE)
aline(s2, Inches(1.35), color=PURPLE)

stats = [
    ("845+", "Notarized trusts\nacross all 50 states"),
    ("20+ yrs", "Notary + real estate\n+ mortgage consulting"),
    ("2002", "First commissioned\nat Wachovia/Wells Fargo"),
    ("CNTDA\nCNSA", "Certified Notary\nTrust Delivery Agent"),
    ("50 states", "Coverage - B2C\n+ B2B attorney clients"),
]
cw = Inches(2.3); ch = Inches(2.2); gap = Inches(0.18)
for i, (num, label) in enumerate(stats):
    l = Inches(0.55) + i * (cw + gap)
    rect(s2, l, Inches(1.8), cw, ch, DARK_BLUE)
    tb(s2, num, l + Inches(0.1), Inches(1.95), cw - Inches(0.2), Inches(1.1),
       size=18, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    tb(s2, label, l + Inches(0.1), Inches(3.05), cw - Inches(0.2), Inches(0.7),
       size=11, color=LGREY, align=PP_ALIGN.CENTER)

rect(s2, Inches(0.55), Inches(4.2), Inches(12.23), Inches(0.6), DARK_BLUE)
tb(s2, "Legacy Curator-in-Chief  |  Manassas, Virginia  |  WordPress site  |  Also Managing Director at The Signing Workforce LLC",
   Inches(0.65), Inches(4.28), Inches(12), Inches(0.45), size=11, color=LGREY, align=PP_ALIGN.CENTER)

tb(s2, "JUNE'S BACKGROUND", Inches(0.55), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=PURPLE)
bullets(s2, [
    "Notary since 2002 - Wachovia/Wells Fargo Financial Specialist",
    "Certified Notary Trust Delivery Agent (CNTDA) + CNSA credentials",
    "20+ years: notary, real estate, mortgage consulting",
    "B2B pivot: pitching attorneys to outsource trust delivery admin"],
        Inches(0.55), Inches(5.35), Inches(5.8), Inches(1.5), size=12, color=LGREY)

tb(s2, "NOTARY WORKFORCE SERVICES", Inches(6.8), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=PURPLE)
bullets(s2, [
    "Free 25-min estate planning webinar (ClickMeeting)",
    "Free 15-min consultation with Legacy Curators",
    "Document planning with partner attorney advocates",
    "White-Glove Trust Delivery: mobile, after-hours, weekends"],
        Inches(6.8), Inches(5.35), Inches(5.8), Inches(1.5), size=12, color=LGREY)


# SLIDE 3 - STRENGTHS
s3 = prs.slides.add_slide(BLANK)
bg(s3)
rect(s3, 0, 0, Inches(0.12), H, GREEN)
tb(s3, "WHAT'S WORKING", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s3, "845 Trusts. A Sharp B2B Pivot. And a Webinar Funnel That Just Needs a Back End.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.8), size=22, bold=True, color=WHITE)
aline(s3, Inches(1.45), color=GREEN)

strengths = [
    ("Free Webinar Top-of-Funnel Is the Right Move for Estate Planning",
     "Estate planning is a high-consideration purchase. Prospects need education before they need a pitch. The free 25-minute webinar on ClickMeeting is the correct lead magnet for this category. The gap is not the funnel entry - it is everything that happens after someone watches and does not immediately book. A post-webinar sequence turns the webinar investment into an asset that compounds."),
    ("845 Notarized Trusts Across 50 States = Proof That Scaling Is Possible",
     "Most notary operations are local. Operating across all 50 states with 845 completed trusts and 31 years of combined experience is a genuine differentiator. This is the headline social proof that should be above the fold on the homepage. Right now it is buried or invisible. The operational track record is there - the marketing of it is not."),
    ("B2B Attorney Pivot Is a High-Value Channel If the Infrastructure Matches",
     "June's LinkedIn posts are pitching estate planning attorneys directly: outsource your mobile trust delivery, client onboarding, and signing ceremonies to Notary Workforce. The angle is strong - attorneys want billable hours, not admin. But the website has no attorney landing page, no attorney-specific intake, and no qualifying form. The pitch is happening on LinkedIn and landing on a site built for consumers."),
    ("White-Glove Trust Delivery Is Premium and Differentiated",
     "After-hours, weekends, mobile delivery to clients' homes or offices. Most competitors do not offer this. It is the right service for high-value estate planning clients who will not rearrange their schedules for paperwork. The product is differentiated. The website does not describe it with enough specificity or urgency to convert the prospects who would value it most."),
    ("Probate Calculator Exists as a Lead Tool - Just Needs a Gate",
     "An interactive probate cost calculator is already live on the site. This is a genuine lead magnet - people searching 'how much does probate cost in Virginia' are high-intent estate planning prospects. The calculator currently shows results without requiring an email. One change - gate the results behind an email capture - turns passive traffic into a list-building machine."),
]
for i, (title, desc) in enumerate(strengths):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s3, Inches(0.55), t, Inches(0.06), Inches(0.7), GREEN)
    tb(s3, title, Inches(0.75), t, Inches(11.8), Inches(0.38), size=13, bold=True, color=WHITE)
    tb(s3, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.6), size=11, color=LGREY)


# SLIDE 4 - THE THREE GAPS
s4 = prs.slides.add_slide(BLANK)
bg(s4)
rect(s4, 0, 0, Inches(0.12), H, RED)
tb(s4, "THE THREE GAPS", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=RED)
tb(s4, "The Webinar Works. The 80% Who Don't Book Immediately Are Gone.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=26, bold=True, color=WHITE)
aline(s4, Inches(1.35), color=RED)

gap_cards = [
    ("GAP 1\nPOST-WEBINAR\nSILENCE", RED,
     "Where 80% of registrants go",
     [
         "Prospect watches 25-min webinar, does not immediately book 15-min call - nothing follows up",
         "No automated Day 1, Day 3, Day 7 email sequence to re-engage warm registrants",
         "No AI chatbot to answer questions they have at 10pm after watching",
         "Industry average: 3-5 touchpoints needed before estate planning purchase decision",
     ]),
    ("GAP 2\nB2B PITCH WITH\nNO B2B WEBSITE", GOLD,
     "Attorney pipeline built on LinkedIn, not the site",
     [
         "LinkedIn posts pitch attorneys to outsource trust delivery - strong angle, real pain point",
         "No dedicated attorney landing page - attorney lands on consumer homepage",
         "No qualifying intake form: state coverage needed, trust volume, current process, pain points",
         "No pricing or service scope for attorney partnerships vs. consumer clients",
     ]),
    ("GAP 3\nSOCIAL PROOF\nBURIED", ACCENT,
     "845 trusts invisible. Calculator leaks leads.",
     [
         "845 notarized trusts across 50 states - not visible on homepage above the fold",
         "Probate calculator shows results without email capture - high-intent traffic leaves with no contact",
         "No testimonials or case studies from attorney partners or consumer clients",
         "No state coverage map showing 50-state reach despite this being the key differentiator",
     ]),
]
cw3 = Inches(3.95); gap3i = Inches(0.18)
for i, (title, color, subtitle, pts) in enumerate(gap_cards):
    l = Inches(0.55) + i * (cw3 + gap3i)
    rect(s4, l, Inches(1.8), cw3, Inches(0.5), color)
    tb(s4, title, l + Inches(0.12), Inches(1.85), cw3 - Inches(0.24), Inches(0.42),
       size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    rect(s4, l, Inches(2.3), cw3, Inches(4.6), DARK_BLUE)
    tb(s4, subtitle, l + Inches(0.12), Inches(2.38), cw3 - Inches(0.24), Inches(0.35),
       size=10, italic=True, color=color)
    bullets(s4, pts, l + Inches(0.12), Inches(2.78), cw3 - Inches(0.24), Inches(3.9),
            size=11, color=LGREY, dot=color)

rect(s4, Inches(0.55), Inches(7.05), Inches(12.23), Inches(0.3), DARK_BLUE)
tb(s4, "The funnel entry works. The re-engagement layer, the B2B infrastructure, and the social proof deployment are what's missing.",
   Inches(0.7), Inches(7.08), Inches(12), Inches(0.25), size=11, italic=True, color=GOLD)


# SLIDE 5 - WEBSITE DESIGN AUDIT
s5 = prs.slides.add_slide(BLANK)
bg(s5)
rect(s5, 0, 0, Inches(0.12), H, PURPLE)
tb(s5, "WEBSITE DESIGN AUDIT", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=PURPLE)
tb(s5, "Good Entry. No Back End.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.8), size=36, bold=True, color=WHITE)
aline(s5, Inches(1.45), color=PURPLE, w=Inches(1.5))

rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(3.6), DARK_BLUE)
rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(0.42), RED)
tb(s5, "NOTARYWORKFORCE.COM TODAY", Inches(0.65), Inches(1.92), Inches(3.6), Inches(0.38),
   size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "No post-webinar email sequence",
    "No AI chatbot or live chat - 10pm questions get nothing",
    "845 trusts not visible on homepage",
    "Probate calculator shows results without email gate",
    "No attorney landing page or B2B intake form",
    "No scheduling automation for mobile delivery",
    "No SEO content by state despite 50-state coverage",
    "No pricing page or package transparency",
], Inches(0.7), Inches(2.42), Inches(3.6), Inches(2.9), size=11, color=LGREY, dot=RED)

tb(s5, "->", Inches(4.45), Inches(3.6), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(3.6), DARK_BLUE)
rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(0.42), GOLD)
tb(s5, "WHAT A CONVERTING SERVICE SITE DOES", Inches(5.1), Inches(1.92), Inches(3.4), Inches(0.38),
   size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Post-webinar nurture: Day 1/3/7/14 email sequence",
    "AI chatbot pre-qualifies and books 24/7",
    "Social proof above fold with trust stats + testimonials",
    "Calculator gated behind email capture",
    "Dedicated B2B page with attorney-specific intake",
    "AI scheduling for service appointments",
    "SEO content targeting state + service keywords",
    "Transparent packages for B2C and B2B tiers",
], Inches(5.1), Inches(2.42), Inches(3.4), Inches(2.9), size=11, color=LGREY, dot=GOLD)

tb(s5, "->", Inches(8.68), Inches(3.6), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(3.6), DARK_BLUE)
rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(0.42), GREEN)
tb(s5, "WHAT GENOS AI BUILDS", Inches(9.35), Inches(1.92), Inches(3.4), Inches(0.38),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Post-webinar nurture sequence (Day 1/3/7/14)",
    "AI chatbot: estate planning FAQ + pre-qualification",
    "845 trusts + testimonials above fold",
    "Probate calculator: email gate before results",
    "Attorney B2B landing page + qualifying intake form",
    "AI scheduling for mobile trust delivery",
    "State-specific SEO content engine (50 states)",
    "B2C + B2B pricing transparency pages",
], Inches(9.35), Inches(2.42), Inches(3.4), Inches(2.9), size=11, color=LGREY, dot=GREEN)

rect(s5, Inches(0.55), Inches(5.72), Inches(12.23), Inches(1.28), DARK_BLUE)
tb(s5, "CONVERSION LEAKS:", Inches(0.7), Inches(5.78), Inches(2.0), Inches(0.38),
   size=11, bold=True, color=PURPLE)
stats2 = [
    ("80%", "Webinar registrants lost\nwith no follow-up"),
    ("0", "Attorney-specific pages\ndespite B2B LinkedIn pivot"),
    ("845", "Trusts invisible\non homepage"),
    ("50 states", "Coverage with zero\nstate SEO content"),
]
sw = Inches(2.7); sl = Inches(3.0)
for i, (num, label) in enumerate(stats2):
    xl = sl + i * sw * 0.97
    tb(s5, num, xl, Inches(5.75), Inches(2.5), Inches(0.55),
       size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s5, label, xl, Inches(6.22), Inches(2.5), Inches(0.7),
       size=10, color=LGREY, align=PP_ALIGN.CENTER)


# SLIDE 6 - AI AUTOMATION STACK
s6 = prs.slides.add_slide(BLANK)
bg(s6)
rect(s6, 0, 0, Inches(0.12), H, ACCENT)
tb(s6, "AI AUTOMATION STACK", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s6, "Four Layers That Convert the Pipeline June Is Already Building.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=26, bold=True, color=WHITE)
aline(s6, Inches(1.35), color=ACCENT)

opps = [
    ("Post-Webinar\nNurture Sequence", TEAL,
     "The webinar already runs. The 80% of registrants who watch but don't immediately book a 15-minute call are not saying no - they are thinking about it. A Day 1 recap email, Day 3 FAQ email, Day 7 social proof email, and Day 14 urgency email re-engages them at the moment they are ready. Most estate planning purchases happen in the second or third touchpoint, not the first."),
    ("AI Chatbot for\nPre-Qualification", GOLD,
     "Someone watches the webinar at 10pm and has questions. Or they land on the site from a Google search and want to know if this is right for them. An AI chatbot trained on June's content, the service details, and estate planning FAQ answers in real time, pre-qualifies B2C and B2B visitors, and either books a 15-minute call or routes attorney prospects to the B2B intake. Runs 24/7 without June's time."),
    ("Attorney B2B\nInfrastructure", GREEN,
     "June's LinkedIn post is generating attorney interest. An attorney who clicks through lands on a consumer homepage. A dedicated attorney landing page, a qualifying intake form (states needed, monthly trust volume, current signing process, main pain point), and a tailored follow-up sequence convert that LinkedIn traffic into a qualified pipeline. The B2B pitch is already being made - it just needs a place to land."),
    ("AI Scheduling for\nMobile Trust Delivery", PURPLE,
     "Mobile trust delivery means logistics: client location, availability window, notary routing, document readiness check, post-signing follow-up. Each step currently involves manual communication. An AI scheduling system that handles intake, confirms appointment windows, sends reminders, and automates post-signing document delivery removes the administrative layer that scales poorly as client volume grows."),
]
cw2 = Inches(2.95); gap2 = Inches(0.16)
for i, (title, color, desc) in enumerate(opps):
    l2 = Inches(0.55) + i * (cw2 + gap2)
    rect(s6, l2, Inches(1.8), cw2, Inches(0.06), color)
    rect(s6, l2, Inches(1.86), cw2, Inches(4.7), DARK_BLUE)
    tb(s6, title, l2 + Inches(0.12), Inches(1.98), cw2 - Inches(0.24), Inches(0.65),
       size=13, bold=True, color=color)
    tb(s6, desc, l2 + Inches(0.12), Inches(2.65), cw2 - Inches(0.24), Inches(3.65),
       size=11, color=LGREY)

rect(s6, Inches(0.55), Inches(6.7), Inches(12.23), Inches(0.45), DARK_BLUE)
tb(s6, "June's job is delivering trust signings. The automation layer handles everything between 'watched the webinar' and 'ready to sign.'",
   Inches(0.7), Inches(6.75), Inches(12), Inches(0.35), size=11, italic=True, color=GOLD)


# SLIDE 7 - THE WEBINAR FUNNEL MAP
s7 = prs.slides.add_slide(BLANK)
bg(s7)
rect(s7, 0, 0, Inches(0.12), H, ORANGE)
tb(s7, "THE WEBINAR FUNNEL", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ORANGE)
tb(s7, "What Happens After Someone Watches - Today vs. After Genos AI.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=28, bold=True, color=WHITE)
aline(s7, Inches(1.35), color=ORANGE, w=Inches(2.5))

rect(s7, Inches(0.55), Inches(1.65), Inches(12.23), Inches(1.55), DARK_BLUE)
bm_stats = [
    ("80%", "Webinar registrants who\ndon't book immediately"),
    ("3-5x", "Touchpoints needed before\nestate planning decision"),
    ("Day 14", "Last automated re-engagement\nbefore lead goes cold"),
    ("24/7", "Chatbot coverage for\npost-webinar questions"),
    ("100%", "Attorney leads with\nB2B intake captured"),
]
ssw = Inches(2.35)
for i, (num, label) in enumerate(bm_stats):
    xl = Inches(0.7) + i * ssw
    tb(s7, num, xl, Inches(1.73), Inches(2.2), Inches(0.65),
       size=20, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    tb(s7, label, xl, Inches(2.33), Inches(2.2), Inches(0.5),
       size=10, color=LGREY, align=PP_ALIGN.CENTER)

rect(s7, Inches(0.55), Inches(3.4), Inches(5.9), Inches(3.25), DARK_BLUE)
tb(s7, "TODAY: WHAT HAPPENS AFTER THE WEBINAR", Inches(0.7), Inches(3.47), Inches(5.6), Inches(0.35),
   size=11, bold=True, color=RED)
bullets(s7, [
    "Prospect watches 25-min webinar",
    "Prompted to book a 15-minute call",
    "Does not book immediately -> nothing happens",
    "Has a question at 10pm -> no chatbot, no answer",
    "Visits again next week -> no retargeting sequence",
    "Searches for estate planning 3 weeks later -> finds a competitor",
    "Attorney sees LinkedIn post, visits homepage -> finds a consumer site",
], Inches(0.7), Inches(3.85), Inches(5.6), Inches(2.55), size=11, color=LGREY, dot=RED)

rect(s7, Inches(6.65), Inches(3.4), Inches(6.13), Inches(3.25), DARK_BLUE)
tb(s7, "AFTER GENOS AI: THE SAME PROSPECT", Inches(6.8), Inches(3.47), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=GREEN)
bullets(s7, [
    "Watches webinar -> Day 1 email: recap + FAQ answers",
    "Day 3 email: social proof (845 trusts, testimonials)",
    "Has a question at 10pm -> AI chatbot answers + qualifies",
    "Day 7 email: White-Glove Trust Delivery explained",
    "Day 14 email: probate cost urgency + booking CTA",
    "Books call at Day 14 -> arrives pre-qualified with context",
    "Attorney sees LinkedIn post -> dedicated attorney page + intake form",
], Inches(6.8), Inches(3.85), Inches(5.8), Inches(2.55), size=11, color=LGREY, dot=GREEN)

rect(s7, Inches(0.55), Inches(6.8), Inches(12.23), Inches(0.38), DARK_BLUE)
tb(s7, "The webinar is the asset. The sequence is what turns the asset into revenue.",
   Inches(0.7), Inches(6.85), Inches(12), Inches(0.28), size=11, italic=True, color=ORANGE)


# SLIDE 8 - IMPLEMENTATION PLAN
s8 = prs.slides.add_slide(BLANK)
bg(s8)
rect(s8, 0, 0, Inches(0.12), H, GREEN)
tb(s8, "IMPLEMENTATION PLAN", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s8, "Back End First. B2B Second. SEO Third.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=30, bold=True, color=WHITE)
aline(s8, Inches(1.35), color=GREEN)

phases = [
    ("PHASE 1", "Weeks 1-3", GOLD, [
        ("Post-Webinar Sequence Live", "Day 1/3/7/14 email sequence connected to ClickMeeting registration. Every future webinar registrant enters the nurture automatically. Historic registrants back-filled."),
        ("AI Chatbot Deployed", "Trained on webinar content, service details, estate planning FAQ. Pre-qualifies B2C and B2B visitors 24/7. Books calls or routes attorney prospects to intake."),
        ("Probate Calculator Gated", "Email capture before results shown. Immediate list-building from existing calculator traffic."),
    ]),
    ("PHASE 2", "Weeks 4-6", ACCENT, [
        ("Attorney B2B Page Live", "Dedicated landing page: mobile trust delivery outsourcing for estate planning attorneys. Qualifying intake form routing to June with full context."),
        ("Social Proof Above Fold", "845 trusts + testimonials surfaced on homepage. State coverage map or counter added for 50-state differentiation."),
        ("AI Scheduling for Delivery", "Mobile trust delivery appointment automation: intake, window confirmation, reminders, post-signing document delivery."),
    ]),
    ("PHASE 3", "Months 2-3", GREEN, [
        ("SEO Content Engine", "State-specific estate planning guides: Virginia, Maryland, DC, Florida, Texas first. Probate cost by state, trust delivery by state - high-intent keyword targeting."),
        ("Attorney Pipeline Data", "Intake form data analyzed to refine B2B targeting. Attorney conversion sequence built based on first 30 days of qualified intake."),
        ("Pricing + Package Pages", "Transparent B2C and B2B tier pages. White-Glove Trust Delivery positioned as premium tier. Attorney partnership pricing and scope documented."),
    ]),
]
pw = Inches(3.9); gap4 = Inches(0.18)
for i, (phase, timing, color, items) in enumerate(phases):
    l3 = Inches(0.55) + i * (pw + gap4)
    rect(s8, l3, Inches(1.8), pw, Inches(0.52), color)
    tb(s8, phase, l3 + Inches(0.12), Inches(1.84), pw * 0.5 - Inches(0.12), Inches(0.42),
       size=13, bold=True, color=NAVY)
    tb(s8, timing, l3 + pw * 0.45, Inches(1.88), pw * 0.55 - Inches(0.15), Inches(0.38),
       size=11, color=NAVY, align=PP_ALIGN.RIGHT)
    rect(s8, l3, Inches(2.32), pw, Inches(4.6), DARK_BLUE)
    y_item = Inches(2.42)
    for title, desc in items:
        rect(s8, l3 + Inches(0.12), y_item + Inches(0.06), Inches(0.05), Inches(0.52), color)
        tb(s8, title, l3 + Inches(0.25), y_item, pw - Inches(0.38), Inches(0.32),
           size=12, bold=True, color=color)
        tb(s8, desc, l3 + Inches(0.25), y_item + Inches(0.3), pw - Inches(0.38), Inches(0.45),
           size=11, color=LGREY)
        y_item += Inches(0.95)


# SLIDE 9 - BEFORE vs AFTER
s9 = prs.slides.add_slide(BLANK)
bg(s9)
rect(s9, 0, 0, Inches(0.12), H, ACCENT)
tb(s9, "BEFORE vs AFTER", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s9, "What Actually Changes for Notary Workforce in 90 Days",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=28, bold=True, color=WHITE)
aline(s9, Inches(1.35))

comparisons = [
    ("Post-Webinar Pipeline",
     "80% of webinar registrants who don't immediately book a call are permanently lost. The webinar spend generates a one-shot chance at conversion.",
     "Day 1/3/7/14 sequence re-engages every registrant. Estate planning conversions happen at touchpoint 3-5. The webinar becomes a compounding asset."),
    ("After-Hours Questions",
     "Prospect has estate planning questions at 10pm after watching the webinar. No chatbot, no reply. They Google a competitor.",
     "AI chatbot answers immediately, pre-qualifies, books a call or routes to the right intake. Zero late-night traffic exits without capture."),
    ("Attorney B2B Pipeline",
     "June pitches attorneys on LinkedIn. Interested attorneys visit notaryworkforce.com and find a consumer homepage. No qualifying form, no context captured.",
     "Dedicated attorney landing page. Qualifying intake captures state coverage, volume, current process. June receives only pre-qualified attorney prospects."),
    ("Probate Calculator",
     "High-intent visitors (searching 'probate cost Virginia') use the calculator and see results. No email captured. They leave with the answer and no relationship.",
     "Email gate before results. High-intent calculator traffic builds the list. Follow-up sequence addresses probate concern and introduces Notary Workforce services."),
    ("Mobile Delivery Logistics",
     "Scheduling mobile trust delivery involves manual back-and-forth for location, timing, document readiness, and post-signing follow-up as volume grows.",
     "AI scheduling handles intake, window confirmation, reminders, and post-signing document delivery. June's time goes to the signing itself, not the coordination."),
]

lw = Inches(5.4)
rect(s9, Inches(0.55), Inches(1.65), lw, Inches(0.35), RED)
tb(s9, "TODAY", Inches(0.65), Inches(1.68), lw - Inches(0.2), Inches(0.28),
   size=11, bold=True, color=WHITE)
rect(s9, Inches(6.45), Inches(1.65), lw + Inches(0.43), Inches(0.35), GREEN)
tb(s9, "AFTER 90 DAYS", Inches(6.55), Inches(1.68), lw, Inches(0.28),
   size=11, bold=True, color=WHITE)

for i, (topic, before, after) in enumerate(comparisons):
    t = Inches(2.1) + i * Inches(1.02)
    rect(s9, Inches(0.55), t, Inches(12.23), Inches(0.22), DARK_BLUE)
    tb(s9, topic.upper(), Inches(0.65), t + Inches(0.02), Inches(12), Inches(0.2),
       size=9, bold=True, color=MGREY)
    tb(s9, before, Inches(0.65), t + Inches(0.24), lw - Inches(0.2), Inches(0.66),
       size=11, color=LGREY)
    tb(s9, after, Inches(6.55), t + Inches(0.24), lw, Inches(0.66),
       size=11, color=WHITE)


# SLIDE 10 - NEXT STEPS
s10 = prs.slides.add_slide(BLANK)
bg(s10)
rect(s10, 0, 0, Inches(0.12), H, PURPLE)
rect(s10, Inches(7.5), 0, Inches(5.83), H, DARK_BLUE)
tb(s10, "NEXT STEPS", Inches(0.55), Inches(0.3), Inches(6.5), Inches(0.5),
   size=11, bold=True, color=PURPLE)
tb(s10, "Let's Talk for\n15 Minutes.",
   Inches(0.55), Inches(0.8), Inches(6.8), Inches(2.0), size=48, bold=True, color=WHITE)
aline(s10, Inches(2.8), color=PURPLE, w=Inches(2.5))
tb(s10, "No pitch deck needed. Reply or book a call.\nWe'll focus on the post-webinar sequence first.",
   Inches(0.55), Inches(3.05), Inches(6.5), Inches(0.8), size=14, color=LGREY)

steps = [
    "Reply to this email - we'll set up a time",
    "Start with post-webinar sequence (fastest revenue impact)",
    "Add AI chatbot + attorney B2B page in weeks 2-4",
    "SEO + scheduling automation in months 2-3",
]
bullets(s10, steps, Inches(0.55), Inches(3.95), Inches(6.5), Inches(2.0), size=14, color=WHITE, dot=PURPLE)

rect(s10, 0, Inches(6.5), Inches(7.5), Inches(1.0), DARK_BLUE)
tb(s10, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699  |  www.genosai.tech",
   Inches(0.55), Inches(6.6), Inches(6.8), Inches(0.5), size=11, color=MGREY)

tb(s10, "THREE GAPS.\nTHREE FIXES.", Inches(7.85), Inches(0.6), Inches(5.0), Inches(1.1),
   size=22, bold=True, color=PURPLE)
aline(s10, Inches(1.65), color=PURPLE)

summary_items = [
    ("Post-Webinar Silence", "Day 1/3/7/14 nurture sequence captures the 80%"),
    ("B2B No Infrastructure", "Attorney landing page + qualifying intake form"),
    ("Buried Social Proof", "845 trusts above fold + calculator email gate"),
    ("No Scheduling AI", "Mobile delivery automation removes manual coord."),
    ("Zero SEO by State", "State-specific content engine: all 50 states"),
    ("Late-Night Gap", "AI chatbot: estate planning FAQ + pre-qual 24/7"),
]
for i, (label, val) in enumerate(summary_items):
    t2 = Inches(1.85) + i * Inches(0.88)
    rect(s10, Inches(7.85), t2, Inches(1.5), Inches(0.72), NAVY)
    tb(s10, label, Inches(7.9), t2 + Inches(0.05), Inches(1.4), Inches(0.32),
       size=10, bold=True, color=PURPLE)
    tb(s10, val, Inches(9.5), t2 + Inches(0.05), Inches(3.6), Inches(0.65),
       size=11, color=LGREY)


prs.save(ROOT / "output" / "decks" / "NotaryWorkforce_Audit_GenosAI.pptx")
print("Saved: NotaryWorkforce_Audit_GenosAI.pptx")
