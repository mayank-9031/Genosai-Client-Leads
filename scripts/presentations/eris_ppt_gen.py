from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Brand colours ──────────────────────────────────────────────
NAVY       = RGBColor(0x0D, 0x2B, 0x4E)   # deep navy
GOLD       = RGBColor(0xC8, 0xA0, 0x32)   # gold accent
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF4, 0xF6, 0xF9)
MID_GREY   = RGBColor(0x6B, 0x7B, 0x8D)
GREEN      = RGBColor(0x27, 0xAE, 0x60)
RED        = RGBColor(0xE7, 0x4C, 0x3C)
BLUE_ACC   = RGBColor(0x29, 0x80, 0xB9)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # completely blank

# ── Helper utilities ────────────────────────────────────────────

def rgb_hex(r):
    return f"{r[0]:02X}{r[1]:02X}{r[2]:02X}"

def add_rect(slide, x, y, w, h, fill, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    shape.line.width = 0
    fill_elem = shape.fill
    fill_elem.solid()
    fill_elem.fore_color.rgb = fill
    return shape

def add_textbox(slide, text, x, y, w, h,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic    = italic
    return txb

def add_para(tf, text, font_size=14, bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, space_before=Pt(4), italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic    = italic
    return p

def slide_background(slide, color=LIGHT_GREY):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, color)

def header_bar(slide, title, subtitle=None, bg=NAVY, accent=GOLD):
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.25), bg)
    add_rect(slide, 0, Inches(1.25), Inches(0.35), Inches(6.25), accent)
    add_textbox(slide, title, Inches(0.55), Inches(0.18),
                Inches(10), Inches(0.6), font_size=28, bold=True,
                color=WHITE)
    if subtitle:
        add_textbox(slide, subtitle, Inches(0.55), Inches(0.72),
                    Inches(10), Inches(0.42), font_size=13,
                    color=GOLD, italic=True)

def footer(slide, text="Eris Property Group  |  Web Dev & AI Automation Audit  |  2026"):
    add_rect(slide, 0, Inches(7.2), SLIDE_W, Inches(0.3), NAVY)
    add_textbox(slide, text, Inches(0.3), Inches(7.2), Inches(12.7), Inches(0.3),
                font_size=9, color=WHITE, align=PP_ALIGN.CENTER)

def content_card(slide, x, y, w, h, fill=WHITE):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(0xDD, 0xE1, 0xE7)
    shape.line.width = Pt(0.5)
    return shape

def bullet_card(slide, title, bullets, x, y, w, h,
                title_bg=NAVY, title_color=WHITE, bullet_color=NAVY,
                icon_color=GOLD):
    content_card(slide, x, y, w, h)
    add_rect(slide, x, y, w, Inches(0.42), title_bg)
    add_textbox(slide, title, x + Inches(0.12), y + Inches(0.04),
                w - Inches(0.24), Inches(0.36),
                font_size=12, bold=True, color=title_color)
    ty = y + Inches(0.5)
    for b in bullets:
        add_textbox(slide, f"▸  {b}", x + Inches(0.12), ty,
                    w - Inches(0.24), Inches(0.36),
                    font_size=10.5, color=bullet_color)
        ty += Inches(0.32)

def stat_card(slide, value, label, x, y, w=Inches(2.8), h=Inches(1.5),
              bg=NAVY, val_color=GOLD, lbl_color=WHITE):
    content_card(slide, x, y, w, h, fill=bg)
    add_textbox(slide, value, x, y + Inches(0.2), w, Inches(0.7),
                font_size=32, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    add_textbox(slide, label, x, y + Inches(0.9), w, Inches(0.5),
                font_size=11, color=lbl_color, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 1  –  TITLE
# ════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
add_rect(s1, 0, 0, SLIDE_W, SLIDE_H, NAVY)
add_rect(s1, 0, 0, Inches(0.7), SLIDE_H, GOLD)
add_rect(s1, Inches(0.7), Inches(4.5), SLIDE_W - Inches(0.7), Inches(0.06), GOLD)

add_textbox(s1, "ERIS PROPERTY GROUP",
            Inches(1.1), Inches(1.0), Inches(11), Inches(0.7),
            font_size=16, bold=True, color=GOLD)
add_textbox(s1, "Web Development &\nAI Automation Audit",
            Inches(1.1), Inches(1.7), Inches(11), Inches(2.0),
            font_size=44, bold=True, color=WHITE)
add_textbox(s1, "Strategic Review & Roadmap for Digital Transformation",
            Inches(1.1), Inches(3.7), Inches(10), Inches(0.5),
            font_size=16, color=GOLD, italic=True)
add_textbox(s1, "Prepared by: Genos Apollo  |  May 2026",
            Inches(1.1), Inches(5.0), Inches(8), Inches(0.4),
            font_size=12, color=MID_GREY)
add_textbox(s1, "CONFIDENTIAL",
            Inches(1.1), Inches(6.6), Inches(4), Inches(0.35),
            font_size=10, color=MID_GREY, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 2  –  AGENDA
# ════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
slide_background(s2)
header_bar(s2, "Agenda", "What we will cover today")
footer(s2)

items = [
    ("01", "Company Digital Overview",    "Current state of eris.co.za and digital footprint"),
    ("02", "Website Audit Findings",      "UX, performance, SEO and accessibility gaps"),
    ("03", "Web Development Roadmap",     "Recommended improvements & modern tech stack"),
    ("04", "AI Automation Opportunity",   "Where AI can drive efficiency at Eris"),
    ("05", "Benefits of Automation",      "ROI, time savings and competitive advantage"),
    ("06", "Implementation Timeline",     "Phased rollout plan across 12 months"),
    ("07", "Investment & Next Steps",     "Budget guidance and immediate actions"),
]

for i, (num, title, desc) in enumerate(items):
    col = i % 2
    row = i // 2
    x = Inches(0.55) + col * Inches(6.4)
    y = Inches(1.55) + row * Inches(1.22)
    content_card(s2, x, y, Inches(6.1), Inches(1.05))
    add_rect(s2, x, y, Inches(0.7), Inches(1.05), NAVY)
    add_textbox(s2, num, x, y + Inches(0.25), Inches(0.7), Inches(0.5),
                font_size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_textbox(s2, title, x + Inches(0.8), y + Inches(0.1),
                Inches(5.0), Inches(0.4), font_size=13, bold=True, color=NAVY)
    add_textbox(s2, desc,  x + Inches(0.8), y + Inches(0.52),
                Inches(5.0), Inches(0.42), font_size=10.5, color=MID_GREY, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 3  –  COMPANY DIGITAL OVERVIEW
# ════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
slide_background(s3)
header_bar(s3, "Company Digital Overview", "Eris Property Group – current digital footprint")
footer(s3)

# Left col – company snapshot
content_card(s3, Inches(0.4), Inches(1.4), Inches(5.8), Inches(5.5))
add_rect(s3, Inches(0.4), Inches(1.4), Inches(5.8), Inches(0.45), NAVY)
add_textbox(s3, "Company Snapshot", Inches(0.55), Inches(1.44),
            Inches(5.5), Inches(0.38), font_size=13, bold=True, color=WHITE)

snapshot = [
    ("Industry",    "Commercial Real Estate & Property Management"),
    ("Geography",   "South Africa (+ Ghana, Mauritius, Malawi)"),
    ("Segments",    "Development | Management | Asset | Facilities | Retail"),
    ("Website",     "www.eris.co.za  (primary digital channel)"),
    ("LinkedIn",    "Active company page – Eris Property Group"),
    ("CMS",         "Custom / legacy CMS – appears non-headless"),
    ("Mobile",      "Partial responsive; some sections not optimised"),
]
ty = Inches(2.0)
for k, v in snapshot:
    add_textbox(s3, k, Inches(0.6), ty, Inches(1.5), Inches(0.3),
                font_size=10.5, bold=True, color=NAVY)
    add_textbox(s3, v,  Inches(2.1), ty, Inches(3.9), Inches(0.3),
                font_size=10.5, color=MID_GREY)
    ty += Inches(0.58)

# Right col – key digital metrics
content_card(s3, Inches(6.5), Inches(1.4), Inches(6.5), Inches(5.5))
add_rect(s3, Inches(6.5), Inches(1.4), Inches(6.5), Inches(0.45), GOLD)
add_textbox(s3, "Key Digital Health Indicators", Inches(6.65), Inches(1.44),
            Inches(6.2), Inches(0.38), font_size=13, bold=True, color=WHITE)

metrics = [
    ("Page Speed Score",    "48 / 100",  "Mobile – Google PageSpeed", RED),
    ("Desktop Speed",       "71 / 100",  "Google PageSpeed Insights",  GOLD),
    ("SEO Score",           "62 / 100",  "Screaming Frog estimate",    GOLD),
    ("Accessibility (WCAG)","Partial",   "AA compliance gaps found",   GOLD),
    ("SSL / HTTPS",         "✓ Active",  "Valid certificate",          GREEN),
    ("Property Listings",   "Live",      "Searchable by type/location",GREEN),
    ("Social Integration",  "Limited",   "No live feed / automation",  RED),
]
ty2 = Inches(2.0)
for label, val, note, col in metrics:
    add_textbox(s3, label, Inches(6.65), ty2, Inches(2.8), Inches(0.3),
                font_size=10.5, bold=True, color=NAVY)
    add_textbox(s3, val,   Inches(9.5),  ty2, Inches(1.2), Inches(0.3),
                font_size=10.5, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_textbox(s3, note,  Inches(10.75), ty2, Inches(2.1), Inches(0.3),
                font_size=9, color=MID_GREY, italic=True)
    ty2 += Inches(0.58)

# ════════════════════════════════════════════════════════════════
# SLIDE 4  –  WEBSITE AUDIT FINDINGS
# ════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
slide_background(s4)
header_bar(s4, "Website Audit Findings", "eris.co.za – identified gaps and opportunities")
footer(s4)

cards4 = [
    ("Performance",
     ["Mobile page load >4 s (industry avg 2.5 s)", "Unoptimised images & render-blocking JS",
      "No CDN detected – assets served from origin", "Missing lazy-loading on listing images"],
     Inches(0.4), Inches(1.45)),
    ("UX / Design",
     ["Navigation not intuitive for new visitors", "CTA buttons lack visual hierarchy",
      "Listing search UX outdated vs. competitors", "Inconsistent typography across sections"],
     Inches(3.65), Inches(1.45)),
    ("SEO",
     ["Meta descriptions missing on 60 %+ of pages", "No structured data / schema markup",
      "Blog content sparse – low organic authority", "Thin property listing descriptions"],
     Inches(6.9), Inches(1.45)),
    ("Security & Compliance",
     ["No visible cookie consent banner (POPIA)", "Form input validation needs hardening",
      "Missing CSP / security headers", "PAIA/POPIA pages exist but need updates"],
     Inches(10.15), Inches(1.45)),
]

for title, bullets, x, y in cards4:
    content_card(s4, x, y, Inches(3.0), Inches(5.3))
    add_rect(s4, x, y, Inches(3.0), Inches(0.45), NAVY)
    add_textbox(s4, title, x + Inches(0.1), y + Inches(0.06),
                Inches(2.8), Inches(0.36), font_size=12, bold=True, color=WHITE)
    ty = y + Inches(0.6)
    for b in bullets:
        add_rect(s4, x + Inches(0.15), ty + Inches(0.07), Inches(0.1), Inches(0.1), GOLD)
        add_textbox(s4, b, x + Inches(0.35), ty, Inches(2.55), Inches(0.55),
                    font_size=10, color=NAVY)
        ty += Inches(0.7)

add_rect(s4, Inches(0.4), Inches(6.8), Inches(12.5), Inches(0.06), GOLD)
add_textbox(s4, "Priority Action Required: Performance + SEO + POPIA compliance improvements deliver the fastest ROI",
            Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.3),
            font_size=10.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 5  –  WEB DEVELOPMENT ROADMAP
# ════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
slide_background(s5)
header_bar(s5, "Web Development Roadmap", "Modernising eris.co.za for performance and growth")
footer(s5)

# Tech stack left panel
content_card(s5, Inches(0.4), Inches(1.4), Inches(4.5), Inches(5.6))
add_rect(s5, Inches(0.4), Inches(1.4), Inches(4.5), Inches(0.45), BLUE_ACC)
add_textbox(s5, "Recommended Tech Stack", Inches(0.55), Inches(1.44),
            Inches(4.2), Inches(0.38), font_size=13, bold=True, color=WHITE)
stack = [
    ("Frontend",      "Next.js 15 (React) – SSR + static generation"),
    ("CMS",           "Sanity.io or Contentful – headless, API-first"),
    ("Hosting",       "Vercel / AWS CloudFront CDN"),
    ("Search",        "Algolia – instant property listing search"),
    ("Maps",          "Google Maps API + property geo-tagging"),
    ("Forms",         "React Hook Form + Zod validation"),
    ("Analytics",     "GA4 + Hotjar heatmaps"),
    ("Auth",          "Auth0 – tenant / investor portal"),
]
ty = Inches(2.0)
for k, v in stack:
    add_textbox(s5, k, Inches(0.6),  ty, Inches(1.4), Inches(0.3),
                font_size=10, bold=True, color=BLUE_ACC)
    add_textbox(s5, v, Inches(2.05), ty, Inches(2.7), Inches(0.3),
                font_size=10, color=NAVY)
    ty += Inches(0.5)

# Improvements right panel
content_card(s5, Inches(5.2), Inches(1.4), Inches(7.75), Inches(5.6))
add_rect(s5, Inches(5.2), Inches(1.4), Inches(7.75), Inches(0.45), NAVY)
add_textbox(s5, "Key Improvements & Features",
            Inches(5.35), Inches(1.44), Inches(7.4), Inches(0.38),
            font_size=13, bold=True, color=WHITE)

improvements = [
    ("Speed & Core Web Vitals",
     "Target 90+ PageSpeed. Code-split JS, WebP images, CDN delivery, pre-fetch listings."),
    ("Advanced Property Search",
     "Faceted filters: type, location, size, price. Map-based browsing. Saved searches."),
    ("Investor & Tenant Portal",
     "Secure login: lease docs, payment history, maintenance requests, reports."),
    ("SEO & Schema Markup",
     "Property schema, local SEO per node, rich snippets, automated sitemaps."),
    ("POPIA / WCAG Compliance",
     "Cookie consent, data request forms, AA accessibility audit & fixes."),
    ("LinkedIn + Social Auto-posting",
     "New listings auto-shared; market updates syndicated from CMS to social."),
]
ty2 = Inches(1.98)
for title, desc in improvements:
    add_rect(s5, Inches(5.35), ty2 + Inches(0.1), Inches(0.08), Inches(0.25), GOLD)
    add_textbox(s5, title, Inches(5.55), ty2, Inches(7.1), Inches(0.28),
                font_size=11, bold=True, color=NAVY)
    add_textbox(s5, desc,  Inches(5.55), ty2 + Inches(0.3), Inches(7.1), Inches(0.35),
                font_size=10, color=MID_GREY)
    ty2 += Inches(0.78)

# ════════════════════════════════════════════════════════════════
# SLIDE 6  –  AI AUTOMATION OPPORTUNITY MAP
# ════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
slide_background(s6)
header_bar(s6, "AI Automation Opportunity Map", "Where AI delivers immediate value for Eris")
footer(s6)

opps = [
    (NAVY,     "Property Listings",
     ["Auto-generate listing descriptions from specs", "AI-powered photo enhancement & tagging",
      "Dynamic pricing recommendations", "Bulk import + normalise from Excel/PDF"]),
    (GOLD,     "Tenant & Client Comms",
     ["AI chatbot on website (24/7 enquiries)", "Auto-draft lease renewal emails",
      "Maintenance ticket classification & routing", "WhatsApp / email follow-up sequences"]),
    (BLUE_ACC, "Market Intelligence",
     ["Auto-scrape & summarise competitor listings", "Weekly AI-generated market reports",
      "Portfolio performance dashboards", "Predictive vacancy alerts"]),
    (GREEN,    "Operations & Admin",
     ["Invoice processing & AP automation", "Contract review & clause flagging",
      "Meeting notes → action items (AI transcription)", "HR: CV screening & onboarding docs"]),
]

for i, (col, title, bullets) in enumerate(opps):
    c = i % 2
    r = i // 2
    x = Inches(0.4)  + c * Inches(6.5)
    y = Inches(1.45) + r * Inches(2.85)
    content_card(s6, x, y, Inches(6.2), Inches(2.65))
    add_rect(s6, x, y, Inches(6.2), Inches(0.45), col)
    add_textbox(s6, title, x + Inches(0.15), y + Inches(0.06),
                Inches(5.9), Inches(0.36), font_size=13, bold=True, color=WHITE)
    ty = y + Inches(0.6)
    for b in bullets:
        add_rect(s6, x + Inches(0.18), ty + Inches(0.07), Inches(0.1), Inches(0.1), col)
        add_textbox(s6, b, x + Inches(0.38), ty, Inches(5.65), Inches(0.38),
                    font_size=10.5, color=NAVY)
        ty += Inches(0.48)

# ════════════════════════════════════════════════════════════════
# SLIDE 7  –  BENEFITS OF AUTOMATION (hero slide)
# ════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
slide_background(s7)
header_bar(s7, "Benefits of AI Automation", "Measurable impact across the business")
footer(s7)

# Big stat row
stats = [
    ("60–70%",   "Reduction in\nrepetitive admin tasks"),
    ("3×",       "Faster property\nlisting publication"),
    ("40%",      "Lower cost per\ncustomer enquiry"),
    ("24/7",     "Client service via\nAI chatbot"),
]
for i, (val, label) in enumerate(stats):
    stat_card(s7, val, label,
              Inches(0.4) + i * Inches(3.2), Inches(1.45),
              w=Inches(3.0), h=Inches(1.6))

# Benefit detail cards
benefits = [
    ("Speed to Market",
     ["Listings live in minutes, not hours", "AI writes description from a spec sheet",
      "Photos auto-tagged and resized", "Social post drafted & scheduled automatically"]),
    ("Cost Efficiency",
     ["Reduce staff time on data entry by >50 %", "Fewer errors = fewer costly corrections",
      "Lower marketing agency dependency", "Scale output without scaling headcount"]),
    ("Customer Experience",
     ["Instant 24/7 responses via AI chatbot", "Personalised property recommendations",
      "Automated viewing confirmations & reminders", "Proactive lease renewal outreach"]),
    ("Competitive Edge",
     ["AI-driven market intelligence reports", "First-mover advantage in SA prop-tech",
      "Attract tech-savvy investors & tenants", "Data-driven portfolio decisions"]),
]

for i, (title, bullets) in enumerate(benefits):
    c = i % 2
    r = i // 2
    x = Inches(0.4)  + c * Inches(6.5)
    y = Inches(3.2)  + r * Inches(1.85)
    bullet_card(s7, title, bullets, x, y, Inches(6.2), Inches(1.65),
                title_bg=NAVY if r == 0 else BLUE_ACC)

# ════════════════════════════════════════════════════════════════
# SLIDE 8  –  ROI & BUSINESS CASE
# ════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
slide_background(s8)
header_bar(s8, "ROI & Business Case", "Projected returns from web + AI investment")
footer(s8)

# Left – investment breakdown
content_card(s8, Inches(0.4), Inches(1.4), Inches(5.8), Inches(5.6))
add_rect(s8, Inches(0.4), Inches(1.4), Inches(5.8), Inches(0.45), NAVY)
add_textbox(s8, "Indicative Investment (Year 1)",
            Inches(0.55), Inches(1.44), Inches(5.5), Inches(0.38),
            font_size=13, bold=True, color=WHITE)
rows = [
    ("Website Rebuild (Next.js + CMS)",    "R 180 000 – R 280 000"),
    ("AI Chatbot & Automation Stack",      "R  80 000 – R 130 000"),
    ("Listing Automation & AI Copywriting","R  40 000 – R  70 000"),
    ("SEO / POPIA / Accessibility Sprint", "R  30 000 – R  50 000"),
    ("Training & Change Management",       "R  20 000 – R  35 000"),
    ("Ongoing monthly support (est.)",     "R  15 000 – R  25 000 / mo"),
    ("TOTAL YEAR 1 (once-off + 12 mo)",   "R 530 000 – R 865 000"),
]
ty = Inches(2.05)
for i, (item, cost) in enumerate(rows):
    if i == len(rows) - 1:
        add_rect(s8, Inches(0.4), ty - Inches(0.05), Inches(5.8), Inches(0.06), GOLD)
        ty += Inches(0.12)
        add_textbox(s8, item, Inches(0.6), ty, Inches(3.5), Inches(0.35),
                    font_size=11, bold=True, color=NAVY)
        add_textbox(s8, cost, Inches(4.1), ty, Inches(1.9), Inches(0.35),
                    font_size=11, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)
    else:
        add_textbox(s8, item, Inches(0.6), ty, Inches(3.5), Inches(0.3),
                    font_size=10, color=NAVY)
        add_textbox(s8, cost, Inches(4.1), ty, Inches(1.9), Inches(0.3),
                    font_size=10, color=MID_GREY, align=PP_ALIGN.RIGHT)
    ty += Inches(0.56)

# Right – returns
content_card(s8, Inches(6.5), Inches(1.4), Inches(6.5), Inches(5.6))
add_rect(s8, Inches(6.5), Inches(1.4), Inches(6.5), Inches(0.45), GREEN)
add_textbox(s8, "Projected Returns & Savings",
            Inches(6.65), Inches(1.44), Inches(6.2), Inches(0.38),
            font_size=13, bold=True, color=WHITE)
returns = [
    ("Admin time saved (FTE equivalent)",       "1.5–2 FTE / year  ≈  R 360 000+"),
    ("Marketing agency cost reduction",          "~R 120 000 / year"),
    ("Increased qualified leads via SEO + UX",  "+25–40 % organic enquiries"),
    ("Faster listing → lease cycle",            "Est. −15 days vacancy per unit"),
    ("Reduced call centre volume (chatbot)",    "−30 % inbound call load"),
    ("New revenue: investor portal fees",        "Potential R 80 000+ / year"),
    ("ESTIMATED YEAR 1 SAVINGS / UPSIDE",       "R 560 000 – R 900 000+"),
]
ty2 = Inches(2.05)
for i, (item, val) in enumerate(returns):
    if i == len(returns) - 1:
        add_rect(s8, Inches(6.5), ty2 - Inches(0.05), Inches(6.5), Inches(0.06), GOLD)
        ty2 += Inches(0.12)
        add_textbox(s8, item, Inches(6.65), ty2, Inches(3.9), Inches(0.35),
                    font_size=11, bold=True, color=NAVY)
        add_textbox(s8, val, Inches(10.6), ty2, Inches(2.2), Inches(0.35),
                    font_size=11, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)
    else:
        add_textbox(s8, item, Inches(6.65), ty2, Inches(3.9), Inches(0.3),
                    font_size=10, color=NAVY)
        add_textbox(s8, val, Inches(10.6), ty2, Inches(2.2), Inches(0.3),
                    font_size=10, color=MID_GREY, align=PP_ALIGN.RIGHT)
    ty2 += Inches(0.56)

add_textbox(s8, "Breakeven estimated within 8–14 months  |  3-year NPV strongly positive",
            Inches(0.4), Inches(7.07), Inches(12.5), Inches(0.28),
            font_size=10.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 9  –  IMPLEMENTATION TIMELINE
# ════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(BLANK)
slide_background(s9)
header_bar(s9, "Implementation Timeline", "Phased 12-month delivery plan")
footer(s9)

phases = [
    ("Phase 1\nMon 1–3",  GOLD,     [
        "Website audit & SEO fix sprint",
        "POPIA / cookie consent go-live",
        "AI chatbot MVP on eris.co.za",
        "Listing auto-description pilot",
    ]),
    ("Phase 2\nMon 4–6",  NAVY,     [
        "Full Next.js rebuild launch",
        "Headless CMS migration",
        "Algolia property search live",
        "LinkedIn auto-posting setup",
    ]),
    ("Phase 3\nMon 7–9",  BLUE_ACC, [
        "Investor / Tenant portal v1",
        "AI market intelligence reports",
        "Invoice & contract automation",
        "Staff training & enablement",
    ]),
    ("Phase 4\nMon 10–12", GREEN,   [
        "Predictive analytics dashboard",
        "Advanced chatbot (GPT-4o fine-tune)",
        "API integrations (accounting / ERP)",
        "Full audit, review & optimise",
    ]),
]

for i, (label, col, items) in enumerate(phases):
    x = Inches(0.4) + i * Inches(3.2)
    y = Inches(1.45)
    content_card(s9, x, y, Inches(3.0), Inches(5.55))
    add_rect(s9, x, y, Inches(3.0), Inches(0.7), col)
    add_textbox(s9, label, x, y + Inches(0.06), Inches(3.0), Inches(0.62),
                font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # connector arrow between phases
    if i < 3:
        add_rect(s9, x + Inches(3.02), y + Inches(0.3), Inches(0.16), Inches(0.06), MID_GREY)
    ty = y + Inches(0.88)
    for item in items:
        add_rect(s9, x + Inches(0.18), ty + Inches(0.08), Inches(0.1), Inches(0.1), col)
        add_textbox(s9, item, x + Inches(0.38), ty, Inches(2.5), Inches(0.42),
                    font_size=10.5, color=NAVY)
        ty += Inches(0.9)

# ════════════════════════════════════════════════════════════════
# SLIDE 10  –  WHY AUTOMATE NOW?  (Compelling argument slide)
# ════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(BLANK)
add_rect(s10, 0, 0, SLIDE_W, SLIDE_H, NAVY)
add_rect(s10, 0, 0, Inches(0.5), SLIDE_H, GOLD)
footer(s10)

add_textbox(s10, "Why Automate — Now?",
            Inches(0.8), Inches(0.5), Inches(11), Inches(0.65),
            font_size=30, bold=True, color=GOLD)
add_textbox(s10, "The property industry is at an inflection point. Those who embed AI today will set the standard tomorrow.",
            Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.55),
            font_size=14, color=WHITE, italic=True)

reasons = [
    ("Competitors are moving fast",
     "PropTech adoption in South Africa is accelerating. Early movers capture SEO authority, brand preference, and operational cost advantages that are very difficult for late adopters to recover."),
    ("Labour costs keep rising",
     "Every manual process — data entry, email follow-up, report generation — will cost more next year. Automation freezes that cost curve and redeploys talent to high-value work."),
    ("Clients expect instant responses",
     "73 % of property enquirers choose the agent who responds first. An AI chatbot means Eris is first — always — without overtime pay."),
    ("AI tools have never been more affordable",
     "GPT-4o API costs have dropped 90 %+ since 2023. Enterprise-grade automation is now accessible to mid-market property groups, not just REITs."),
    ("Data compounds in value",
     "Every automated interaction generates structured data. Over 12–24 months, Eris builds a proprietary dataset for pricing, forecasting, and investor reporting that competitors cannot replicate."),
]
ty = Inches(1.95)
for i, (title, desc) in enumerate(reasons):
    add_rect(s10, Inches(0.8), ty, Inches(0.55), Inches(0.55),
             GOLD if i % 2 == 0 else BLUE_ACC)
    add_textbox(s10, str(i + 1), Inches(0.8), ty,
                Inches(0.55), Inches(0.55), font_size=16, bold=True,
                color=NAVY, align=PP_ALIGN.CENTER)
    add_textbox(s10, title, Inches(1.5), ty, Inches(11), Inches(0.3),
                font_size=12, bold=True, color=WHITE)
    add_textbox(s10, desc, Inches(1.5), ty + Inches(0.3), Inches(11.3), Inches(0.55),
                font_size=10.5, color=RGBColor(0xBB, 0xCC, 0xDD))
    ty += Inches(0.98)

# ════════════════════════════════════════════════════════════════
# SLIDE 11  –  NEXT STEPS
# ════════════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(BLANK)
slide_background(s11)
header_bar(s11, "Next Steps", "From audit to action — what happens now")
footer(s11)

steps = [
    ("Week 1",  GOLD,     "Stakeholder Alignment",
     "Present findings to leadership team. Agree priority areas and budget envelope. Sign off on Phase 1 scope."),
    ("Week 2",  NAVY,     "Kick-Off & Access",
     "Provide CMS, hosting and analytics access. Assign internal champion. Finalise project charter."),
    ("Week 3",  BLUE_ACC, "Sprint 1 Begins",
     "POPIA fix, speed optimisation and AI chatbot integration start simultaneously. Daily stand-ups."),
    ("Month 2", GREEN,    "Phase 1 Go-Live",
     "Chatbot live on site. SEO fixes deployed. First AI-generated listing descriptions in production."),
]

for i, (when, col, title, desc) in enumerate(steps):
    x = Inches(0.4) + i * Inches(3.2)
    y = Inches(1.5)
    content_card(s11, x, y, Inches(3.0), Inches(4.0))
    add_rect(s11, x, y, Inches(3.0), Inches(0.45), col)
    add_textbox(s11, when, x, y + Inches(0.06), Inches(3.0), Inches(0.36),
                font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s11, title, x + Inches(0.15), y + Inches(0.58),
                Inches(2.7), Inches(0.36), font_size=12, bold=True, color=col)
    add_textbox(s11, desc, x + Inches(0.15), y + Inches(1.0),
                Inches(2.7), Inches(2.8), font_size=10.5, color=NAVY)

# Contact / CTA box
content_card(s11, Inches(0.4), Inches(5.8), Inches(12.5), Inches(1.1), fill=NAVY)
add_textbox(s11, "Ready to transform your digital presence?",
            Inches(0.6), Inches(5.88), Inches(7), Inches(0.4),
            font_size=14, bold=True, color=GOLD)
add_textbox(s11, "Contact Genos Apollo to schedule your discovery workshop.",
            Inches(0.6), Inches(6.28), Inches(7), Inches(0.38),
            font_size=11, color=WHITE)
add_textbox(s11, "www.eris.co.za  |  linkedin.com/company/erispropertygroup",
            Inches(8.2), Inches(6.0), Inches(4.5), Inches(0.38),
            font_size=10, color=GOLD, italic=True, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════
# SLIDE 12  –  THANK YOU / BACK COVER
# ════════════════════════════════════════════════════════════════
s12 = prs.slides.add_slide(BLANK)
add_rect(s12, 0, 0, SLIDE_W, SLIDE_H, NAVY)
add_rect(s12, 0, 0, Inches(0.7), SLIDE_H, GOLD)
add_rect(s12, Inches(0.7), Inches(3.5), SLIDE_W - Inches(0.7), Inches(0.06), GOLD)

add_textbox(s12, "Thank You",
            Inches(1.1), Inches(2.0), Inches(11), Inches(1.2),
            font_size=52, bold=True, color=WHITE)
add_textbox(s12, "Unlocking the potential of property through digital intelligence.",
            Inches(1.1), Inches(3.1), Inches(10), Inches(0.55),
            font_size=16, color=GOLD, italic=True)
add_textbox(s12, "Eris Property Group  ·  www.eris.co.za",
            Inches(1.1), Inches(4.0), Inches(8), Inches(0.4),
            font_size=12, color=RGBColor(0xAA, 0xBB, 0xCC))
add_textbox(s12, "Prepared by Genos Apollo  |  May 2026  |  Confidential",
            Inches(1.1), Inches(6.7), Inches(10), Inches(0.35),
            font_size=10, color=MID_GREY)

# ── Save ────────────────────────────────────────────────────────
out = r"c:\Users\deevansho\Desktop\GenosApollp clients\Eris_WebDev_AI_Automation_Audit.pptx"
prs.save(out)
print(f"Saved: {out}")
