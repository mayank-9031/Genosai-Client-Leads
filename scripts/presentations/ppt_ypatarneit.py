import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.parent / "utils"))
from ppt_theme import (
    NAVY, DARK_BLUE, ACCENT, GOLD, RED, GREEN, PURPLE, ORANGE, TEAL, CYAN,
    WHITE, LGREY, MGREY, W, H,
    bg, rect, tb, aline, bullets, tag,
    slide_header, slide_footer, score_badge,
    stat_cards, top_bar_card, phase_card, before_after, genos_logo_block,
)

ROOT = Path(__file__).parent.parent.parent
prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]

# SLIDE 1 - COVER
s1 = prs.slides.add_slide(BLANK)
bg(s1)
rect(s1, 0, 0, Inches(0.12), H, GOLD)
rect(s1, Inches(8.5), 0, Inches(4.83), Inches(2.2), DARK_BLUE)
tb(s1, "GENOS AI", Inches(9.2), Inches(0.4), Inches(3.5), Inches(0.8),
   size=28, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
tb(s1, "www.genosai.tech", Inches(9.2), Inches(1.05), Inches(3.5), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
tb(s1, "YPA TARNEIT - TRUGANINA - WILLIAMS LANDING", Inches(0.55), Inches(1.6),
   Inches(9), Inches(0.6), size=13, bold=False, color=GOLD)
tb(s1, "Website & AI Growth\nAudit Report", Inches(0.55), Inches(2.1),
   Inches(9), Inches(1.8), size=48, bold=True, color=WHITE)
aline(s1, Inches(3.85), color=GOLD, w=Inches(3))
tb(s1, "Prepared exclusively for Jerry Jacob, Director / Partner",
   Inches(0.55), Inches(4.15), Inches(9), Inches(0.5), size=14, color=LGREY)
rect(s1, 0, Inches(6.5), W, Inches(1.0), DARK_BLUE)
tb(s1, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699",
   Inches(0.55), Inches(6.6), Inches(9), Inches(0.5), size=12, color=MGREY)
tb(s1, "May 2026  |  CONFIDENTIAL", Inches(9.5), Inches(6.6), Inches(3.3), Inches(0.5),
   size=12, color=MGREY, align=PP_ALIGN.RIGHT)
rect(s1, Inches(10.5), Inches(3.3), Inches(2.3), Inches(2.5), DARK_BLUE)
tb(s1, "AUDIT SCORE", Inches(10.55), Inches(3.4), Inches(2.2), Inches(0.5),
   size=11, bold=True, color=MGREY, align=PP_ALIGN.CENTER)
tb(s1, "5", Inches(10.55), Inches(3.75), Inches(2.2), Inches(1.1),
   size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
tb(s1, "/ 10", Inches(10.55), Inches(4.7), Inches(2.2), Inches(0.4),
   size=18, color=LGREY, align=PP_ALIGN.CENTER)
tag(s1, "HIGH PRIORITY", Inches(10.75), Inches(5.1), fill=RED, size=10)


# SLIDE 2 - AT A GLANCE
s2 = prs.slides.add_slide(BLANK)
bg(s2)
rect(s2, 0, 0, Inches(0.12), H, GOLD)
tb(s2, "AT A GLANCE", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GOLD)
tb(s2, "Jerry Jacob - Who He Is",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s2, Inches(1.35), color=GOLD)

stats = [
    ("3029", "Top of Postcode"),
    ("$190K", "Above Median (Truganina)"),
    ("7 days", "vs 50-day market"),
    ("24-26", "YPA Offices"),
    ("300+", "YPA Team Members"),
]
cw = Inches(2.3); ch = Inches(2.2); gap = Inches(0.18)
for i, (num, label) in enumerate(stats):
    l = Inches(0.55) + i * (cw + gap)
    rect(s2, l, Inches(1.8), cw, ch, DARK_BLUE)
    tb(s2, num, l + Inches(0.1), Inches(1.95), cw - Inches(0.2), Inches(1.1),
       size=26, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s2, label, l + Inches(0.1), Inches(3.05), cw - Inches(0.2), Inches(0.7),
       size=12, color=LGREY, align=PP_ALIGN.CENTER)

rect(s2, Inches(0.55), Inches(4.2), Inches(12.23), Inches(0.6), DARK_BLUE)
tb(s2, "Tarneit / Truganina / Williams Landing  -  YPA Founded 2007  -  VIC + QLD network  -  Co-Director with Jalpa Patel  -  Featured in Elite Agent",
   Inches(0.65), Inches(4.28), Inches(12), Inches(0.45), size=12, color=LGREY, align=PP_ALIGN.CENTER)

tb(s2, "WHAT SETS THE OFFICE APART", Inches(0.55), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=GOLD)
bullets(s2, ["Top-performing team in Western Melbourne's largest growth corridor (3029)",
             "Jerry-Jalpa partnership protects vendor outcomes over volume",
             "Track record of beating market by weeks and median by six figures"],
        Inches(0.55), Inches(5.35), Inches(5.8), Inches(1.5), size=13, color=LGREY)

tb(s2, "WHAT THE OFFICE COVERS", Inches(6.8), Inches(5.0), Inches(5.8), Inches(0.35),
   size=11, bold=True, color=GOLD)
bullets(s2, ["Tarneit, Truganina, Williams Landing - 3029 Western Melbourne corridor",
             "One of Australia's fastest-growing residential markets",
             "Multi-service: residential sales + property management"],
        Inches(6.8), Inches(5.35), Inches(5.8), Inches(1.5), size=13, color=LGREY)


# SLIDE 3 - WHAT'S WORKING
s3 = prs.slides.add_slide(BLANK)
bg(s3)
rect(s3, 0, 0, Inches(0.12), H, GREEN)
tb(s3, "WHAT'S WORKING", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s3, "Strengths to Build On",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s3, Inches(1.35), color=GREEN)

strengths = [
    ("Top of 3029 - Western Melbourne's Largest Growth Corridor",
     "Tarneit, Truganina, and Williams Landing sit at the centre of one of Australia's fastest-growing residential markets. Being top-performing in the 3029 postcode means leading the corridor with the largest expanding vendor pool in the state. The demographic tailwind alone is the kind of structural advantage agents in mature markets don't get."),
    ("Concrete, Quotable Records That Most Agents Can't Match",
     "Truganina sale at $190K above median, after previous agents failed. One Tarneit campaign closed in seven days against a market averaging 50. These aren't soft credentials. They're the exact stories vendors need to read before they decide who to call. The proof exists - it's just not on the office page."),
    ("Elite Agent Media Coverage Gives Independent Validation",
     "Elite Agent featuring Jerry's Truganina sale provides the kind of third-party validation that an in-house testimonial can't match. Industry trade publications carry weight with vendors who've done their research. Right now the article lives off-site, when it could be the centerpiece case study on the office page."),
    ("Jerry-Jalpa Partnership Model Is a Real Differentiator",
     "Two directors working every campaign together is rare in residential. Most agencies put one agent on a listing. The dual-director attention is what produces the records and is what builds long-term client relationships. The website doesn't lead with the partnership - it should."),
    ("YPA Franchise Scale Provides Shared Infrastructure",
     "24-26 offices, 300+ team members, brand recognition across Victoria and growth in Queensland. The franchise model means Jerry's office benefits from network referrals, shared marketing infrastructure, and brand equity that boutique agencies build over decades. The platform exists - the local digital layer needs to catch up."),
]
for i, (title, desc) in enumerate(strengths):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s3, Inches(0.55), t, Inches(0.06), Inches(0.7), GREEN)
    tb(s3, title, Inches(0.75), t, Inches(11.8), Inches(0.38), size=14, bold=True, color=WHITE)
    tb(s3, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# SLIDE 4 - WHERE LEADS ARE LEAKING
s4 = prs.slides.add_slide(BLANK)
bg(s4)
rect(s4, 0, 0, Inches(0.12), H, RED)
tb(s4, "WHERE LEADS ARE LEAKING", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=RED)
tb(s4, "The Vendor Search at 9pm Doesn't Become a YPA Lead.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=28, bold=True, color=WHITE)
aline(s4, Inches(1.35), color=RED)

issues = [
    ("APPRAISAL FORM IS A REQUEST, NOT AN INSTANT VALUATION", "HIGH",
     "YPA's appraisal form asks the vendor for their details and to wait. The next agency in their search results gives them an instant range. Every comparable agent now offers an instant valuation tool with comparable-based pricing. YPA has the sold data and the local expertise. The site is asking the vendor to do the wrong thing first."),
    ("OFFICE PAGE DOESN'T SHOW WHAT THE OFFICE ACTUALLY DOES", "HIGH",
     "Tarneit-Truganina office page reads like an agent roster, not a 3029 suburb authority. No recent sold data embedded, no neighborhood reports, no agent-suburb matching tool, no story about the Truganina sale that beat the market by $190K. The page that should be doing the heavy lifting for the corridor's leading team is doing none of it."),
    ("ELITE AGENT FEATURE LIVES OFF-SITE", "HIGH",
     "Jerry's Truganina sale at $190K above median was featured by Elite Agent. That's the kind of third-party validation vendors trust. Right now it sits on a publisher's site rather than as a permanent case study on the YPA office page where it would close trust gaps before the first call. The asset exists - the placement is wrong."),
    ("NO AI CHAT FOR AFTER-HOURS QUESTIONS", "MEDIUM",
     "Vendor questions about appraisal availability, buyer questions about inspection times, tenant questions about maintenance and rent - all arrive when offices are closed. No after-hours response means most of those visitors are gone by morning. AI chat trained on YPA's process closes that gap without staffing the night shift."),
    ("NO PAST-CLIENT RE-ENGAGEMENT AUTOMATION", "MEDIUM",
     "Australians sell every 6-9 years on average. The buyer Jerry settled in Tarneit in 2018 is starting to think about it now. Without anniversary sequences, current home value updates, and suburb activity reports, that past-client list is a dormant CRM table. Programmed properly, it's the recurring vendor pipeline that powers next year's listings."),
]
for i, (title, pri, desc) in enumerate(issues):
    t = Inches(1.7) + i * Inches(1.05)
    pc = RED if pri == "HIGH" else GOLD
    rect(s4, Inches(0.55), t, Inches(0.06), Inches(0.7), pc)
    tb(s4, title, Inches(0.75), t, Inches(9.2), Inches(0.38), size=12, bold=True, color=WHITE)
    tag(s4, pri, Inches(10.1), t + Inches(0.05), fill=pc, size=9)
    tb(s4, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# SLIDE 5 - WEBSITE DESIGN AUDIT
s5 = prs.slides.add_slide(BLANK)
bg(s5)
rect(s5, 0, 0, Inches(0.12), H, PURPLE)
tb(s5, "WEBSITE DESIGN AUDIT", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=PURPLE)
tb(s5, "Top of 3029. The Office Page Doesn't Show It.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.8), size=32, bold=True, color=WHITE)
aline(s5, Inches(1.45), color=PURPLE, w=Inches(1.5))

rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(3.6), DARK_BLUE)
rect(s5, Inches(0.55), Inches(1.9), Inches(3.8), Inches(0.42), RED)
tb(s5, "YPA TARNEIT-TRUGANINA TODAY", Inches(0.65), Inches(1.92), Inches(3.6), Inches(0.38),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Appraisal form is a request, not instant",
    "Office page = agent roster, not authority",
    "$190K Elite Agent story not on office page",
    "No AI chat / after-hours self-service",
    "No tenant portal for PM clients",
    "No recent sold data by suburb",
    "No past-client re-engagement automation",
    "Awards / case studies not on primary pages",
], Inches(0.7), Inches(2.42), Inches(3.6), Inches(2.9), size=11, color=LGREY, dot=RED)

tb(s5, "->", Inches(4.45), Inches(3.6), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(3.6), DARK_BLUE)
rect(s5, Inches(5.0), Inches(1.9), Inches(3.6), Inches(0.42), GOLD)
tb(s5, "WHAT TOP AU AGENCY SITES DO", Inches(5.1), Inches(1.92), Inches(3.4), Inches(0.38),
   size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Instant valuation tool as primary CTA",
    "Office page = postcode authority",
    "Case study wall on every primary page",
    "AI chat 24/7 for vendor / buyer / tenant",
    "Tenant portal: applications, maintenance, rent",
    "Recent sold data embedded by suburb",
    "Past-client anniversary nurture sequences",
    "Vendor weekly suburb-specific market reports",
], Inches(5.1), Inches(2.42), Inches(3.4), Inches(2.9), size=11, color=LGREY, dot=GOLD)

tb(s5, "->", Inches(8.68), Inches(3.6), Inches(0.5), Inches(0.5),
   size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(3.6), DARK_BLUE)
rect(s5, Inches(9.25), Inches(1.9), Inches(3.6), Inches(0.42), GREEN)
tb(s5, "WHAT GENOS AI BUILDS", Inches(9.35), Inches(1.92), Inches(3.4), Inches(0.38),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s5, [
    "Instant valuation tool with suburb routing",
    "Office page rebuilt as 3029 authority",
    "$190K Truganina story as permanent case",
    "AI chat: appraisals, inspections, PM 24/7",
    "Tenant portal: applications + maintenance",
    "Recent sold data + suburb reports embedded",
    "Past-client anniversary sequences (1/3/5/7yr)",
    "Vendor weekly Tarneit-Truganina reports",
], Inches(9.35), Inches(2.42), Inches(3.4), Inches(2.9), size=11, color=LGREY, dot=GREEN)

rect(s5, Inches(0.55), Inches(5.72), Inches(12.23), Inches(1.28), DARK_BLUE)
tb(s5, "WHY IT MATTERS:", Inches(0.7), Inches(5.78), Inches(2.2), Inches(0.38),
   size=11, bold=True, color=PURPLE)
stats2 = [
    ("3-5x", "vendor lead capture\nvs. plain appraisal request"),
    ("60-70%", "after-hours questions\nresolvable by AI chat"),
    ("6-9 yrs", "AU resale window\npast-client = future vendor"),
    ("$190K", "above median story\ndoing zero work today"),
]
sw = Inches(2.7); sl = Inches(2.9)
for i, (num, label) in enumerate(stats2):
    xl = sl + i * sw
    tb(s5, num, xl, Inches(5.75), Inches(2.5), Inches(0.55),
       size=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s5, label, xl, Inches(6.22), Inches(2.5), Inches(0.7),
       size=10, color=LGREY, align=PP_ALIGN.CENTER)


# SLIDE 6 - AI AUTOMATION OPPORTUNITIES
s6 = prs.slides.add_slide(BLANK)
bg(s6)
rect(s6, 0, 0, Inches(0.12), H, ACCENT)
tb(s6, "AI AUTOMATION OPPORTUNITIES", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s6, "Capturing the 3029 Vendor Searching at 9pm.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s6, Inches(1.35), color=ACCENT)

opps = [
    ("Instant Valuation\n+ Suburb Routing", GOLD,
     "Vendor enters Tarneit, Truganina, or Williams Landing address, gets instant comparable-based range from YPA's sold data, optionally requests refined appraisal. Lead routes automatically to Jerry, Jalpa, or matched agent. Replaces the current request-and-wait appraisal form with a same-second number that captures the lead."),
    ("AI Chat\n24/7 Self-Service", ACCENT,
     "Trained on YPA's appraisal process, current Tarneit-Truganina listings, suburb market data, and PM workflows. Handles vendor, buyer, and tenant questions after hours. Books appraisals and inspections. Triages PM maintenance with photo upload. Routes complex deals to Jerry or Jalpa with conversation context."),
    ("Vendor Nurture\n+ Market Reports", GREEN,
     "Weekly Tarneit-Truganina market reports automated to vendor prospects in the 6-12 month pre-listing window. Sold prices, days-on-market, supply data for the 3029 corridor. Stays in front of vendors until they're ready, then captures the appraisal request when they are."),
    ("Past-Client\nRe-Engagement", RED,
     "Buyer settled in 2018, 2020, 2022 receives anniversary sequences with current home value, suburb activity, and a soft appraisal CTA. Australians sell every 6-9 years on average. Past-client list activated as a recurring vendor pipeline rather than a dormant CRM table."),
    ("Office Page\nas 3029 Authority", PURPLE,
     "Tarneit-Truganina page rebuilt around the suburbs Jerry and Jalpa actually own. Recent sold data, suburb reports, $190K-above-median Truganina case study, agent matching, embedded testimonials. The office page becomes the suburb authority and ranks for postcode searches."),
]
cw2 = Inches(2.35); gap2 = Inches(0.14)
for i, (title, color, desc) in enumerate(opps):
    l2 = Inches(0.55) + i * (cw2 + gap2)
    rect(s6, l2, Inches(1.8), cw2, Inches(0.06), color)
    rect(s6, l2, Inches(1.86), cw2, Inches(4.6), DARK_BLUE)
    tb(s6, title, l2 + Inches(0.12), Inches(1.98), cw2 - Inches(0.24), Inches(0.65),
       size=13, bold=True, color=color)
    tb(s6, desc, l2 + Inches(0.12), Inches(2.65), cw2 - Inches(0.24), Inches(3.6),
       size=11, color=LGREY)

rect(s6, Inches(0.55), Inches(6.7), Inches(12.23), Inches(0.45), DARK_BLUE)
tb(s6, "Top of 3029 + a $190K story + an Elite Agent feature - the assets are there. The site needs to use them.",
   Inches(0.7), Inches(6.75), Inches(12), Inches(0.35), size=11, italic=True, color=GOLD)


# SLIDE 7 - WHAT THIS IS COSTING YOU
s7 = prs.slides.add_slide(BLANK)
bg(s7)
rect(s7, 0, 0, Inches(0.12), H, RED)
tb(s7, "WHAT THIS IS COSTING YOU", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=RED)
tb(s7, "Domain and realestate.com.au Capture What You Don't.",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=30, bold=True, color=WHITE)
aline(s7, Inches(1.35), color=RED)

costs = [
    ("No instant valuation = vendors hand their data to portals, not to YPA",
     "Domain and realestate.com.au capture every Tarneit-Truganina vendor punching in their address for an instant value range. They monetise that data back to whoever pays for leads. YPA has the same comparable sold data and the top-performing 3029 team. The site just doesn't ask for the lead first. Every vendor who values online without filling out a YPA form is a future listing that may not be YPA's."),
    ("Appraisal request that asks vendors to wait = lower conversion than instant",
     "The vendor researching at 9pm wants a number now. Filling out a form and waiting for an agent to call back tomorrow is the same friction every other agency offers - and the agencies offering instant ranges win that lead before YPA's appraisal call ever happens. The conversion gap between instant and request is consistently 3-5x in residential."),
    ("$190K story + Elite Agent feature off-site = trust signals doing zero work",
     "The Truganina sale at $190K above median is exactly the kind of proof that closes vendor decisions. Right now it lives on Elite Agent's site rather than the office page where vendors actually decide. As a permanent case study on the Tarneit-Truganina office page, it shortens every vendor sales cycle and removes the credibility question before the first call."),
    ("Past clients not re-engaged = the 2018 buyer becomes someone else's 2026 vendor",
     "Australians sell every 6-9 years on average. The 2018 settled list is starting to think about listing now. Without programmed anniversary sequences with current home value, suburb activity, and a soft appraisal CTA, those past clients hear from whoever's in their inbox - and that's not YPA. Every settled sale is a future opportunity that needs to be programmed, not hoped for."),
    ("No AI chat or PM portal = staff time on questions automation handles",
     "Maintenance requests, lease balance and renewal questions, inspection availability, application status. PM teams field these all day. Most are answered the same way every time. AI chat trained on YPA's process plus a tenant portal handles 60-70% of routine volume and frees the team for actual exception management."),
]
for i, (title, desc) in enumerate(costs):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s7, Inches(0.55), t, Inches(0.06), Inches(0.7), RED)
    tb(s7, title, Inches(0.75), t, Inches(11.8), Inches(0.38), size=14, bold=True, color=WHITE)
    tb(s7, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# SLIDE 8 - WHAT WE'D BUILD
s8 = prs.slides.add_slide(BLANK)
bg(s8)
rect(s8, 0, 0, Inches(0.12), H, GREEN)
tb(s8, "WHAT WE'D BUILD", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=GREEN)
tb(s8, "Genos AI for YPA Tarneit-Truganina",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s8, Inches(1.35), color=GREEN)

builds = [
    ("Instant Property Valuation + Suburb-Routed Lead Capture",
     "Embedded on the YPA homepage and the Tarneit-Truganina office page. Vendor enters Tarneit, Truganina, or Williams Landing address, gets instant comparable-based range from YPA's sold data, optionally books a refined appraisal. Lead routes automatically to Jerry, Jalpa, or matched specialist with calculator inputs already in CRM."),
    ("AI Chat - 24/7 Vendor / Buyer / Tenant Self-Service",
     "Trained on YPA's appraisal process, current Tarneit-Truganina listings, 3029 market data, PM workflows, and tenant procedures. Handles after-hours questions, books appraisals and inspections, triages maintenance with photo upload. Routes complex inquiries to Jerry or Jalpa with full conversation context."),
    ("Office Page Rebuild as 3029 Suburb Authority",
     "Tarneit-Truganina page rebuilt around the suburbs Jerry and Jalpa actually own: Tarneit, Truganina, Williams Landing. Recent sold data, suburb market reports, $190K-above-median Truganina sale as permanent case study with the Elite Agent feature embedded, agent matching for buyers, partnership story foregrounded."),
    ("Past-Client Re-Engagement at 1, 3, 5, 7 Years",
     "Buyer settled in 2024 enters automated 1-year anniversary sequence with current home value and 3029 activity, then 3-year, 5-year, 7-year sequences with appraisal CTA. Past-client list activated as a recurring vendor pipeline. Every settled sale becomes a future opportunity programmed."),
    ("Tenant Portal + PM Automation + Vendor Nurture",
     "Online rental applications, maintenance triage with photo upload, lease renewal sequences 90 days out, behavioural rent reminders. Plus weekly Tarneit-Truganina market reports for the 6-12 month pre-listing window. Plus awards and case study wall on homepage and office page."),
]
for i, (title, desc) in enumerate(builds):
    t = Inches(1.7) + i * Inches(1.05)
    rect(s8, Inches(0.55), t, Inches(0.06), Inches(0.7), GREEN)
    tb(s8, title, Inches(0.75), t, Inches(11.8), Inches(0.38), size=14, bold=True, color=WHITE)
    tb(s8, desc, Inches(0.75), t + Inches(0.35), Inches(11.8), Inches(0.55), size=12, color=LGREY)


# SLIDE 9 - 90-DAY OUTCOMES
s9 = prs.slides.add_slide(BLANK)
bg(s9)
rect(s9, 0, 0, Inches(0.12), H, ACCENT)
tb(s9, "90-DAY OUTCOMES", Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
   size=11, bold=True, color=ACCENT)
tb(s9, "What Changes in 90 Days",
   Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=32, bold=True, color=WHITE)
aline(s9, Inches(1.35))

phases = [
    ("Month 1", GOLD, [
        "Instant valuation tool live on YPA homepage and Tarneit-Truganina office page",
        "AI chat deployed for vendor, buyer, and tenant after-hours questions",
        "Office page rebuild begins with embedded sold data and case study integration",
        "$190K Truganina story surfaced on office page with Elite Agent feature embedded",
    ]),
    ("Month 2", ACCENT, [
        "Vendor nurture sequences running with weekly 3029 market reports",
        "Tenant portal live - online applications, maintenance triage, lease renewals",
        "Past-client re-engagement sequences active for 1, 3, 5, 7-year anniversaries",
        "First measurable inbound appraisal requests from instant valuation tool",
    ]),
    ("Month 3", GREEN, [
        "Tarneit-Truganina office page acting as 3029 suburb authority and ranking for postcode",
        "Buyer alerts segmented by suburb and property type running",
        "PM automation deflecting routine tenant questions, freeing team for exceptions",
        "Vendor pipeline pre-listing measurable, past-client appraisal requests landing",
    ]),
]
pw = Inches(3.9); gap3 = Inches(0.18)
for i, (month, color, pts) in enumerate(phases):
    l3 = Inches(0.55) + i * (pw + gap3)
    rect(s9, l3, Inches(1.8), pw, Inches(0.42), color)
    tb(s9, month.upper(), l3 + Inches(0.1), Inches(1.84), pw - Inches(0.2), Inches(0.36),
       size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    rect(s9, l3, Inches(2.22), pw, Inches(4.7), DARK_BLUE)
    bullets(s9, pts, l3 + Inches(0.15), Inches(2.35), pw - Inches(0.3), Inches(4.4),
            size=12, color=LGREY, dot=color)


# SLIDE 10 - NEXT STEPS
s10 = prs.slides.add_slide(BLANK)
bg(s10)
rect(s10, 0, 0, Inches(0.12), H, GOLD)
rect(s10, Inches(7.5), 0, Inches(5.83), H, DARK_BLUE)
tb(s10, "NEXT STEPS", Inches(0.55), Inches(0.3), Inches(6.5), Inches(0.5),
   size=11, bold=True, color=GOLD)
tb(s10, "Let's Talk for\n15 Minutes.",
   Inches(0.55), Inches(0.8), Inches(6.8), Inches(2.0), size=48, bold=True, color=WHITE)
aline(s10, Inches(2.8), color=GOLD, w=Inches(2.5))
tb(s10, "No pitch deck required. Reply or book a call directly.\nWe'll walk through what fits YPA Tarneit-Truganina.",
   Inches(0.55), Inches(3.05), Inches(6.5), Inches(0.8), size=14, color=LGREY)

steps = [
    "Reply to this email - we'll set up a time",
    "Book directly: www.genosai.tech/call",
    "Walk through the audit findings in 15 minutes",
    "You decide what, if anything, makes sense to move on",
]
bullets(s10, steps, Inches(0.55), Inches(3.95), Inches(6.5), Inches(2.0), size=14, color=WHITE, dot=GOLD)

rect(s10, 0, Inches(6.5), Inches(7.5), Inches(1.0), DARK_BLUE)
tb(s10, "Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699  |  www.genosai.tech",
   Inches(0.55), Inches(6.6), Inches(6.8), Inches(0.5), size=12, color=MGREY)

tb(s10, "WHAT WE BUILD FOR\nYPA TARNEIT", Inches(7.85), Inches(0.6), Inches(5.0), Inches(0.9),
   size=18, bold=True, color=GOLD)
aline(s10, Inches(1.45), color=GOLD)

summary_items = [
    ("Valuation", "Instant tool with Tarneit-Truganina suburb routing"),
    ("AI Chat", "24/7 vendor / buyer / tenant self-service"),
    ("Office Page", "3029 postcode authority with $190K case study"),
    ("Past Clients", "1/3/5/7 yr anniversary sequences = vendor pipeline"),
    ("PM Portal", "Applications, maintenance, lease renewals automated"),
    ("Vendor Nurture", "Weekly Tarneit-Truganina market reports - 6-12 mo window"),
]
for i, (label, val) in enumerate(summary_items):
    t2 = Inches(1.6) + i * Inches(0.88)
    rect(s10, Inches(7.85), t2, Inches(1.5), Inches(0.72), NAVY)
    tb(s10, label, Inches(7.9), t2 + Inches(0.05), Inches(1.4), Inches(0.32),
       size=10, bold=True, color=GOLD)
    tb(s10, val, Inches(9.5), t2 + Inches(0.05), Inches(3.6), Inches(0.65),
       size=11, color=LGREY)


prs.save(ROOT / "output" / "decks" / "YPATarneit_Audit_GenosAI.pptx")
print("Saved: YPATarneit_Audit_GenosAI.pptx")
