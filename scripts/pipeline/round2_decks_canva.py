"""
Round 2 deck builder using Canva-exported template.
Personalises 4 slide positions per lead:
  - Slide 1:  date -> "Prepared for {company} | {name}, {title}"
  - Slide 2:  score textbox -> actual score
  - Slide 4:  3 issue textboxes -> actual audit issues
  - Slide 11: signature textboxes -> Rohan Malik (LinkedIn link) + genosai.tech (link)
All other slides (3,5-10) are generic Canva design - preserved as-is.
Output: decks/GenosAI_Audit_NN_*.pptx
"""
import json, re, shutil
from copy import copy
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from lxml import etree

ROOT       = Path(r"c:/Users/91638/Desktop/GenosApollp clients")
TEMPLATE   = ROOT / "output" / "decks" / "GenosAI_Template_Canva.pptx"
DECKS_OUT  = ROOT / "output" / "decks"
LEADS      = ROOT / "data" / "pipeline" / "_round2_leads.json"
AUDITS     = ROOT / "data" / "pipeline" / "_round2_audits.json"

LINKEDIN   = "https://www.linkedin.com/in/rohanxmalik/"
WEBSITE    = "https://www.genosai.tech"

def slugify(s):
    s = re.sub(r"[^A-Za-z0-9]+", "_", s).strip("_")
    return s[:40] or "lead"

def short_company(c):
    return c.split(",")[0].split("|")[0].strip()[:60]

def capture_font(text_frame):
    """Return dict of first run's font attributes."""
    try:
        run = text_frame.paragraphs[0].runs[0]
        f = run.font
        color = None
        try:
            if f.color and f.color.type:
                color = f.color.rgb
        except Exception:
            pass
        return {"name": f.name, "size": f.size, "bold": f.bold,
                "italic": f.italic, "color": color}
    except (IndexError, AttributeError):
        return {}

def apply_font(run, font_dict):
    f = run.font
    if font_dict.get("name"):   f.name   = font_dict["name"]
    if font_dict.get("size"):   f.size   = font_dict["size"]
    if font_dict.get("bold") is not None:   f.bold   = font_dict["bold"]
    if font_dict.get("italic") is not None: f.italic = font_dict["italic"]
    if font_dict.get("color"):
        try: f.color.rgb = font_dict["color"]
        except Exception: pass

def clear_tf(text_frame):
    """Remove all runs and extra paragraphs, return first paragraph."""
    para0 = text_frame.paragraphs[0]
    for extra in list(text_frame.paragraphs[1:]):
        extra._p.getparent().remove(extra._p)
    for run in list(para0.runs):
        run._r.getparent().remove(run._r)
    return para0

def set_simple_text(text_frame, text, font_dict=None):
    """Replace text_frame content with text (newline → new paragraph)."""
    font_dict = font_dict or {}
    lines = text.split("\n")
    para0 = clear_tf(text_frame)
    run = para0.add_run()
    run.text = lines[0]
    apply_font(run, font_dict)
    for line in lines[1:]:
        new_p = text_frame.add_paragraph()
        new_p.alignment = para0.alignment
        r = new_p.add_run()
        r.text = line
        apply_font(r, font_dict)

def add_hyperlink_to_run(run, slide, url):
    """Wire a click hyperlink onto an existing run's XML."""
    rId = slide.part.relate_to(url, RT.HYPERLINK, is_external=True)
    rPr = run._r.get_or_add_rPr()
    # remove any existing hlinkClick
    for old in rPr.findall(qn("a:hlinkClick")):
        rPr.remove(old)
    hl = etree.SubElement(rPr, qn("a:hlinkClick"))
    hl.set(qn("r:id"), rId)

def set_linked_text(text_frame, slide, text, url, font_dict=None):
    """Replace text_frame with a single run linked to url."""
    font_dict = font_dict or {}
    para0 = clear_tf(text_frame)
    run = para0.add_run()
    run.text = text
    apply_font(run, font_dict)
    add_hyperlink_to_run(run, slide, url)

def shape_by_id(slide, sid):
    for sh in slide.shapes:
        if sh.shape_id == sid:
            return sh
    return None

def shape_by_id_deep(shapes, sid):
    """Recursively search through shapes and groups."""
    for sh in shapes:
        if sh.shape_id == sid:
            return sh
        if hasattr(sh, 'shapes'):
            found = shape_by_id_deep(sh.shapes, sid)
            if found:
                return found
    return None

def build_deck(lead, audit, out_path):
    shutil.copyfile(TEMPLATE, out_path)
    prs = Presentation(out_path)
    slides = list(prs.slides)

    company = short_company(lead["company"])
    full    = lead["full_name"]
    title   = lead["title"]
    score   = f"{audit['score']} / 10" if audit else "- / 10"
    issues  = audit["issues"] if audit else [
        "Site fetch failed - manual review needed",
        "Likely missing chat / CRM / AI layer based on category norms",
        "Recommend a quick walkthrough before the call",
    ]

    # ---- Slide 1: company + "Prepared for" ----
    s1   = slides[0]
    tb11 = shape_by_id(s1, 5)   # "May 5, 2026" — id changed to 5 after Canva edit
    if tb11 and tb11.has_text_frame:
        fnt = capture_font(tb11.text_frame)
        set_simple_text(
            tb11.text_frame,
            f"Prepared for {company}\n{full}, {title}",
            fnt
        )

    # ---- Slide 2: score + website URL ----
    s2   = slides[1]
    tb6  = shape_by_id(s2, 6)   # "4.0 / 10"
    if tb6 and tb6.has_text_frame:
        fnt = capture_font(tb6.text_frame)
        set_simple_text(tb6.text_frame, score, fnt)

    # Website URL inside Group 7 → id=8
    website_display = lead.get("website", "").replace("https://", "").replace("http://", "").rstrip("/")
    tb_website = shape_by_id_deep(s2.shapes, 8)
    if tb_website and tb_website.has_text_frame:
        fnt2 = capture_font(tb_website.text_frame)
        set_simple_text(tb_website.text_frame, website_display or "genosai.tech", fnt2)

    # ---- Slide 4: issues (ordered top-to-bottom by shape id position) ----
    s4 = slides[3]
    # shape ids sorted by top position: 8 (top), 6 (mid), 5 (bot)
    issue_shape_ids = [8, 6, 5]
    for idx, sid in enumerate(issue_shape_ids):
        sh = shape_by_id(s4, sid)
        if sh and sh.has_text_frame:
            fnt = capture_font(sh.text_frame)
            text = issues[idx] if idx < len(issues) else ""
            set_simple_text(sh.text_frame, text, fnt)

    # ---- Slide 11: signature ----
    s11  = slides[10]
    tb2  = shape_by_id(s11, 2)   # "Thank you!"
    tb3  = shape_by_id(s11, 3)   # contact info

    if tb2 and tb2.has_text_frame:
        fnt = capture_font(tb2.text_frame)
        # "Rohan Malik" as clickable LinkedIn link
        para0 = clear_tf(tb2.text_frame)
        run = para0.add_run()
        run.text = "Rohan Malik"
        apply_font(run, fnt)
        add_hyperlink_to_run(run, s11, LINKEDIN)
        # second line: title
        new_p = tb2.text_frame.add_paragraph()
        new_p.alignment = para0.alignment
        r2 = new_p.add_run()
        r2.text = "Co-Founder & CEO, GenosAI"
        apply_font(r2, {**fnt, "size": int(fnt.get("size", 400000) * 0.7) if fnt.get("size") else None})

    if tb3 and tb3.has_text_frame:
        fnt = capture_font(tb3.text_frame)
        # line 1: www.genosai.tech (linked)
        para0 = clear_tf(tb3.text_frame)
        run1 = para0.add_run()
        run1.text = "www.genosai.tech"
        apply_font(run1, fnt)
        add_hyperlink_to_run(run1, s11, WEBSITE)
        # separator
        sep = para0.add_run()
        sep.text = "  ·  hello@genosai.tech"
        apply_font(sep, fnt)

    prs.save(out_path)

def main():
    leads  = json.loads(LEADS.read_text())
    audits = json.loads(AUDITS.read_text())
    DECKS_OUT.mkdir(exist_ok=True)
    paths = {}
    ok = err = 0
    for i, lead in enumerate(leads, 1):
        a_entry = audits.get(lead["email"], {})
        audit   = a_entry.get("audit")
        slug    = slugify(lead["company"]).replace("__", "_")
        fname   = f"GenosAI_Audit_{lead['num']:02d}_{slug}.pptx"
        out     = DECKS_OUT / fname
        try:
            build_deck(lead, audit, out)
            paths[lead["email"]] = fname
            print(f"[{i}/{len(leads)}] {fname}", flush=True)
            ok += 1
        except Exception as e:
            print(f"[{i}/{len(leads)}] FAIL {fname}: {e}", flush=True)
            err += 1

    (ROOT / "data" / "pipeline" / "_round2_deck_paths.json").write_text(json.dumps(paths, indent=2))
    print(f"\nDone. {ok} ok, {err} errors. Paths -> scripts/_round2_deck_paths.json")

if __name__ == "__main__":
    main()
