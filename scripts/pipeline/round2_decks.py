"""
Round 2 deck builder: clones decks/GenosAI_Audit_01_S2A_Modular.pptx for each
of the 46 new leads with per-lead content. Output: decks/GenosAI_Audit_NN_*.pptx.
"""
import json, copy, re, shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt

ROOT = Path(__file__).parent.parent.parent
DECK_TEMPLATE = ROOT / "output" / "decks" / "GenosAI_Audit_01_S2A_Modular.pptx"
DECKS_OUT = ROOT / "output" / "decks"

LEADS = ROOT / "data" / "pipeline" / "_round2_leads.json"
AUDITS = ROOT / "data" / "pipeline" / "_round2_audits.json"
CONTENT = ROOT / "data" / "pipeline" / "_round2_content.json"

def slugify(s):
    s = re.sub(r"[^A-Za-z0-9]+", "_", s).strip("_")
    return s[:40] or "lead"

def replace_text_preserve_first_run(text_frame, new_text):
    """Replace text_frame's text with new_text, preserving the first paragraph's
    first run's font (size, color, bold, name)."""
    # capture first run formatting
    first_para = text_frame.paragraphs[0]
    first_run = first_para.runs[0] if first_para.runs else None
    src_font = None
    if first_run is not None:
        src_font = {
            "name": first_run.font.name,
            "size": first_run.font.size,
            "bold": first_run.font.bold,
            "italic": first_run.font.italic,
            "color_rgb": None,
        }
        try:
            if first_run.font.color and first_run.font.color.rgb is not None:
                src_font["color_rgb"] = first_run.font.color.rgb
        except Exception:
            pass

    # remove all paragraphs except the first; clear the first
    p = first_para
    # remove sibling paragraphs
    for extra in list(text_frame.paragraphs[1:]):
        extra._p.getparent().remove(extra._p)
    # clear runs in first paragraph
    for run in list(p.runs):
        run._r.getparent().remove(run._r)

    # split new_text by \n and add as separate paragraphs
    lines = new_text.split("\n")
    # first line goes into the (now-empty) first paragraph
    run = p.add_run()
    run.text = lines[0]
    if src_font:
        if src_font["name"]: run.font.name = src_font["name"]
        if src_font["size"]: run.font.size = src_font["size"]
        if src_font["bold"] is not None: run.font.bold = src_font["bold"]
        if src_font["italic"] is not None: run.font.italic = src_font["italic"]
        if src_font["color_rgb"] is not None:
            try:
                run.font.color.rgb = src_font["color_rgb"]
            except Exception:
                pass
    # remaining lines as new paragraphs
    for ln in lines[1:]:
        new_p = text_frame.add_paragraph()
        # copy paragraph alignment
        new_p.alignment = first_para.alignment
        nr = new_p.add_run()
        nr.text = ln
        if src_font:
            if src_font["name"]: nr.font.name = src_font["name"]
            if src_font["size"]: nr.font.size = src_font["size"]
            if src_font["bold"] is not None: nr.font.bold = src_font["bold"]
            if src_font["italic"] is not None: nr.font.italic = src_font["italic"]
            if src_font["color_rgb"] is not None:
                try:
                    nr.font.color.rgb = src_font["color_rgb"]
                except Exception:
                    pass

def find_shape_with_text(slide, needle):
    for sh in slide.shapes:
        if sh.has_text_frame and needle in sh.text_frame.text:
            return sh
    return None

def shape_by_name(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None

def build_deck(lead, content, audit, out_path):
    # copy template file then mutate
    shutil.copyfile(DECK_TEMPLATE, out_path)
    p = Presentation(out_path)
    company = (lead["company"].split(",")[0].split("|")[0]).strip()
    full = lead["full_name"]
    title = lead["title"]
    site = lead["website"]
    if audit:
        platform = audit["platform"]
        score = f"{audit['score']} / 10"
        priority = audit["priority"]
        top_finding = audit["issues"][0] if audit.get("issues") else "Conversion-side gaps are the biggest lever."
    else:
        platform = "Custom"
        score = "- / 10"
        priority = "MED"
        top_finding = "Site fetch failed during audit - manual review needed."

    slides = list(p.slides)

    # ---- Slide 1: cover ----
    s1 = slides[0]
    sh = find_shape_with_text(s1, "S2A Modular")
    if sh: replace_text_preserve_first_run(sh.text_frame, company)
    sh = find_shape_with_text(s1, "Prepared for Brian Kuzdas")
    if sh: replace_text_preserve_first_run(sh.text_frame, f"Prepared for {full}")
    sh = find_shape_with_text(s1, "CEO and Co-Founder")
    if sh: replace_text_preserve_first_run(sh.text_frame, title)
    # tagline + "By GenosAI" line stay generic

    # ---- Slide 2: at a glance ----
    s2 = slides[1]
    # body prose (the long shape with "decent shape" / "core pieces are missing")
    body_shape = None
    for sh in s2.shapes:
        if sh.has_text_frame and ("decent shape" in sh.text_frame.text or "Top finding" in sh.text_frame.text):
            body_shape = sh
            break
    if body_shape:
        # build prose from priority + top_finding
        if priority == "LOW":
            opening = "The site is in decent shape. The remaining wins are mostly AI and conversion-side, not rebuild-the-house work."
        elif priority == "MED":
            opening = "Several pieces are working, but core lead-capture infrastructure is missing or thin."
        else:
            opening = "Several core pieces are missing. Lead capture, follow-up, and AI are the biggest unaddressed levers."
        prose = f"{opening}\n\nTop finding: {top_finding}"
        replace_text_preserve_first_run(body_shape.text_frame, prose)

    # SITE / TECH STACK / SCORE / PRIORITY boxes — find the box right after each label
    label_to_value = [
        ("SITE", site),
        ("TECH STACK", platform),
        ("AUDIT SCORE", score),
        ("PRIORITY", priority),
    ]
    shapes = list(s2.shapes)
    for lbl, val in label_to_value:
        # find label index, then the next text shape with non-label content
        for idx, sh in enumerate(shapes):
            if sh.has_text_frame and sh.text_frame.text.strip() == lbl:
                # next shape with text frame
                for j in range(idx + 1, len(shapes)):
                    nxt = shapes[j]
                    if nxt.has_text_frame:
                        replace_text_preserve_first_run(nxt.text_frame, val)
                        break
                break

    # ---- Slides 3-8: replace the long body shape on each ----
    # We use content["slides"] indexes 2..7 which are: working, leaks, ai, cost, build, 90day
    # Each slide has exactly one big body text frame; we identify it as the largest text frame.
    def biggest_body_shape(slide):
        # the body is the wide tall TextBox with multi-line content
        candidates = []
        for sh in slide.shapes:
            if sh.has_text_frame and sh.text_frame.text.count("\n") >= 1 and len(sh.text_frame.text) > 80:
                candidates.append(sh)
        if not candidates:
            return None
        return max(candidates, key=lambda s: s.width * s.height)

    # build per-slide body content from the audit/content
    # slide indexes in content["slides"] (0-based): 2='working', 3='leaks', 4='ai opps', 5='cost', 6='build', 7='90day'
    slide_contents = content["slides"]
    # Use the body content directly, stripping the leading heading line that
    # duplicates the slide's heading shape.
    def strip_heading(s):
        # strip the first line (heading) + blank line that follows
        lines = s.split("\n", 2)
        if len(lines) >= 3 and lines[1] == "":
            return lines[2]
        return s

    for idx_slide, content_idx in [(2, 2), (3, 3), (4, 4), (5, 5), (6, 6), (7, 7)]:
        body = strip_heading(slide_contents[content_idx])
        body_shape = biggest_body_shape(slides[idx_slide])
        if body_shape:
            replace_text_preserve_first_run(body_shape.text_frame, body)

    # ---- Slide 9: next-steps prose only (signature stays as-is in template) ----
    s9 = slides[8]
    body_shape = None
    for sh in s9.shapes:
        if sh.has_text_frame and "15-minute call" in sh.text_frame.text:
            body_shape = sh
            break
    if body_shape:
        body = (
            f"A 15-minute call. Bring your top two or three questions. I'll bring the pricing model, "
            f"a relevant case study, and a rough timeline tailored to {company}.\n\n"
            f"If now isn't the right window, the deck is yours either way."
        )
        replace_text_preserve_first_run(body_shape.text_frame, body)

    p.save(out_path)

def main():
    leads = json.loads(LEADS.read_text())
    audits = json.loads(AUDITS.read_text())
    contents = json.loads(CONTENT.read_text())
    DECKS_OUT.mkdir(exist_ok=True)
    paths = {}
    for i, lead in enumerate(leads, 1):
        c = contents[lead["email"]]
        a = audits.get(lead["email"], {}).get("audit")
        slug = slugify(lead["company"]).replace("__", "_")
        fname = f"GenosAI_Audit_{lead['num']:02d}_{slug}.pptx"
        out_path = DECKS_OUT / fname
        try:
            build_deck(lead, c, a, out_path)
            paths[lead["email"]] = fname
            print(f"[{i}/{len(leads)}] {fname}", flush=True)
        except Exception as e:
            print(f"[{i}/{len(leads)}] FAIL {fname}: {e}", flush=True)
            raise

    # write back the path map for the workbook step
    (ROOT / "data" / "pipeline" / "_round2_deck_paths.json").write_text(json.dumps(paths, indent=2))
    print(f"\nWrote {len(paths)} decks. Path map -> scripts/_round2_deck_paths.json")

if __name__ == "__main__":
    main()
