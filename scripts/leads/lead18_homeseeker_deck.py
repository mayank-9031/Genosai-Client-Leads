"""
Build HomeSeeker_Audit_GenosAI.pptx from scratch.
Design: docs/BRAND_ASSETS.md — #0A0A0F canvas, Instrument Serif headlines, Satoshi body.
12 slides covering: cover, audit, vision, systems overview, 5 modules, benefits,
implementation, and closing.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).parent.parent.parent
OUT  = ROOT / "output" / "decks" / "HomeSeeker_Audit_GenosAI.pptx"

# ── Brand colors (BRAND_ASSETS.md §3) ───────────────────────────────────────
BG     = RGBColor(0x0A, 0x0A, 0x0F)   # canvas
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)   # headlines
BODY   = RGBColor(0x9D, 0x9D, 0x9F)   # white/60
SUB    = RGBColor(0x86, 0x86, 0x8A)   # white/50
MUTED  = RGBColor(0x6B, 0x6B, 0x6E)   # white/40
DIM    = RGBColor(0x44, 0x44, 0x47)   # white/25
VIOLET = RGBColor(0xA7, 0x8B, 0xFA)   # eyebrow #A78BFA
CARD   = RGBColor(0x11, 0x11, 0x18)   # card surface #111118
BORDER = RGBColor(0x22, 0x22, 0x2A)   # card border ≈ white/6

# ── Fonts (BRAND_ASSETS.md §4) ───────────────────────────────────────────────
SERIF = "Instrument Serif"   # headlines 400
SANS  = "Satoshi"            # body/UI; fallback: Calibri

# ── Slide geometry ───────────────────────────────────────────────────────────
SW = Inches(13.33)   # 16:9 width
SH = Inches(7.5)     # 16:9 height
P  = Inches(0.75)    # outer padding (≈96px at 1920px ref)
TOTAL = 12


# ── Core helpers ─────────────────────────────────────────────────────────────

def new_slide(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def tb(slide, l, t, w, h, text, font=SANS, size=14, color=BODY,
       bold=False, italic=False, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text        = text
    r.font.name   = font
    r.font.size   = Pt(size)
    r.font.color.rgb = color
    r.font.bold   = bold
    r.font.italic = italic
    return box


def tb_lines(slide, l, t, w, h, lines, font=SANS, size=14, color=BODY,
             bold=False, align=PP_ALIGN.LEFT, line_spacing=None):
    """Text box with multiple paragraphs (list of strings; '' = blank line)."""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name  = font
        r.font.size  = Pt(size)
        r.font.color.rgb = color
        r.font.bold  = bold
    return box


def card(slide, l, t, w, h, rounded=True):
    """Dark card background with border."""
    shape_type = 5 if rounded else 1
    s = slide.shapes.add_shape(shape_type, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = CARD
    s.line.color.rgb = BORDER
    s.line.width = Pt(0.75)
    return s


def hline(slide, l, t, w, color=BORDER):
    """Thin horizontal rule."""
    s = slide.shapes.add_shape(1, l, t, w, Pt(0.75))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def eyebrow(slide, text, l, t, w=Inches(6)):
    return tb(slide, l, t, w, Inches(0.3),
              text.upper(), SANS, 10, VIOLET, bold=True)


def headline(slide, text, l, t, w=Inches(10), size=42):
    return tb(slide, l, t, w, Inches(1.6), text, SERIF, size, WHITE)


def slide_header(slide, section, page):
    """Consistent header + footer on every content slide."""
    tb(slide, P, Inches(0.22), Inches(6), Inches(0.28),
       "GENOSAI  ×  HOME SEEKERS", SANS, 9, MUTED, bold=True)
    tb(slide, Inches(7.5), Inches(0.22), Inches(5.08), Inches(0.28),
       section, SANS, 9, MUTED, bold=True, align=PP_ALIGN.RIGHT)
    hline(slide, P, Inches(0.58), SW - 2 * P)
    tb(slide, P, SH - Inches(0.42), Inches(8), Inches(0.28),
       "AI Automation & Growth Proposal  ·  Home Seekers  ·  2026",
       SANS, 8, DIM)
    tb(slide, Inches(9.5), SH - Inches(0.42), Inches(3.08), Inches(0.28),
       f"{page:02d} / {TOTAL}", SANS, 8, DIM, align=PP_ALIGN.RIGHT)


def bullet(slide, dot_char, title, desc, l, t,
           col_w=Inches(5.0), title_size=13, desc_size=12):
    """Dot + bold title + grey description. Returns next top position."""
    tb(slide, l, t, Inches(0.25), Inches(0.26),
       dot_char, SANS, 13, VIOLET, bold=True)
    tb(slide, l + Inches(0.3), t, col_w - Inches(0.3), Inches(0.26),
       title, SANS, title_size, WHITE, bold=True)
    tb(slide, l + Inches(0.3), t + Inches(0.27), col_w - Inches(0.3), Inches(0.38),
       desc, SANS, desc_size, BODY)
    return t + Inches(0.72)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH

s1 = new_slide(prs)

# Thin top accent bar (violet)
accent = s1.shapes.add_shape(1, Inches(0), Inches(0), SW, Pt(3))
accent.fill.solid()
accent.fill.fore_color.rgb = VIOLET
accent.line.fill.background()

# Left column — company identity
tb(s1, P, Inches(1.0), Inches(5.5), Inches(0.35),
   "PREPARED FOR", SANS, 10, MUTED, bold=True)
hline(s1, P, Inches(1.42), Inches(5.5), BORDER)

tb(s1, P, Inches(1.6), Inches(7.5), Inches(1.4),
   "HOME SEEKERS", SERIF, 72, WHITE)

tb(s1, P, Inches(3.15), Inches(5.5), Inches(0.32),
   "Fixed-Fee Real Estate Platform  ·  South Africa", SANS, 13, BODY)

tb(s1, P, Inches(3.65), Inches(5.5), Inches(0.28),
   "Pretoria  ·  Gauteng", SANS, 11, MUTED)

hline(s1, P, Inches(4.1), Inches(5.5), BORDER)

# Proposal type
tb(s1, P, Inches(4.3), Inches(6), Inches(0.28),
   "AI AUTOMATION  &  GROWTH PROPOSAL", SANS, 10, VIOLET, bold=True)
tb_lines(s1, P, Inches(4.7), Inches(6.5), Inches(1.1),
         ["An intelligent automation layer for Home Seekers —",
          "built to capture buyer inquiries 24/7, qualify leads",
          "automatically, and scale agent productivity without",
          "adding headcount."],
         SANS, 13, SUB)

# Right column — "PREPARED BY" card
card(s1, Inches(9.0), Inches(1.2), Inches(3.58), Inches(2.5), rounded=True)
tb(s1, Inches(9.25), Inches(1.5), Inches(3.0), Inches(0.28),
   "PREPARED BY", SANS, 9, MUTED, bold=True)
tb(s1, Inches(9.25), Inches(1.88), Inches(3.0), Inches(0.32),
   "GenosAI Tech", SERIF, 20, WHITE)
tb(s1, Inches(9.25), Inches(2.3), Inches(3.0), Inches(0.26),
   "Rohan Malik  ·  Founder & CEO", SANS, 11, BODY)
tb(s1, Inches(9.25), Inches(2.65), Inches(3.0), Inches(0.26),
   "hello@genosai.tech  ·  www.genosai.tech", SANS, 10, MUTED)

# GenosAI wordmark (bottom of cover, per brand rules)
tb(s1, P, SH - Inches(0.65), Inches(4), Inches(0.45),
   "GenosAI", SERIF, 22, WHITE)
tb(s1, Inches(0.75) + Inches(1.65), SH - Inches(0.62), Inches(3), Inches(0.3),
   "v1.0  ·  2026  ·  Confidential", SANS, 9, DIM)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AUDIT: WHERE THEY STAND TODAY
# ══════════════════════════════════════════════════════════════════════════════
s2 = new_slide(prs)
slide_header(s2, "01  ·  AUDIT", 2)

eyebrow(s2, "WHERE THE BRAND STANDS TODAY", P, Inches(0.85))
headline(s2, "Current Digital Presence", P, Inches(1.18), size=38)
tb(s2, P, Inches(2.2), Inches(8.5), Inches(0.42),
   "Home Seekers' platform proposition is strong — the model is differentiated, production is real. "
   "The gap is entirely on the lead infrastructure: inquiries arriving after hours go unanswered until "
   "someone manually follows up.",
   SANS, 13, SUB)

hline(s2, P, Inches(2.82), SW - 2 * P, BORDER)

# Two columns
CL = P                         # left column start
CR = Inches(7.3)               # right column start
CW = Inches(5.5)               # column width

# LEFT — strengths
card(s2, CL, Inches(3.0), CW, Inches(3.65), rounded=True)
tb(s2, CL + Inches(0.3), Inches(3.15), CW - Inches(0.4), Inches(0.28),
   "WHAT'S ALREADY STRONG", SANS, 10, VIOLET, bold=True)
t = Inches(3.55)
for title, desc in [
    ("Platform model",
     "Agents keep 100% commission — differentiated from every franchise in SA."),
    ("Active production",
     "Rode Potgieter: 8 deals, R11.44M in a single month."),
    ("416+ active listings",
     "Strong traction for a platform launched in 2025."),
    ("Family-driven leadership",
     "Thomas & Rode Potgieter — 23+ combined years in SA real estate."),
]:
    t = bullet(s2, "·", title, desc, CL + Inches(0.3), t, CW - Inches(0.45), 12, 11)

# RIGHT — gaps
card(s2, CR, Inches(3.0), CW, Inches(3.65), rounded=True)
tb(s2, CR + Inches(0.3), Inches(3.15), CW - Inches(0.4), Inches(0.28),
   "ROOM TO ELEVATE", SANS, 10, RGBColor(0xF8, 0x71, 0x71), bold=True)
t = Inches(3.55)
for title, desc in [
    ("No WhatsApp AI",
     "Buyer inquiries sit unanswered while agents are at viewings."),
    ("No lead routing",
     "No automated routing from Property24 to the right agent."),
    ("No viewing scheduler",
     "Booking a site visit requires manual coordination."),
    ("No FICA automation",
     "Agents chase compliance documents instead of closing deals."),
]:
    t = bullet(s2, "·", title, desc, CR + Inches(0.3), t, CW - Inches(0.45), 12, 11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — VISION: THE OPPORTUNITY
# ══════════════════════════════════════════════════════════════════════════════
s3 = new_slide(prs)
slide_header(s3, "02  ·  VISION", 3)

eyebrow(s3, "THE OPPORTUNITY", P, Inches(0.85))
headline(s3, "An AI Automation Layer\nfor Home Seekers", P, Inches(1.18), size=38)

# Quote block
card(s3, P, Inches(2.85), Inches(7.5), Inches(1.25), rounded=True)
tb(s3, P + Inches(0.4), Inches(2.98), Inches(7.0), Inches(0.4),
   "“", SERIF, 32, VIOLET)
tb(s3, P + Inches(0.75), Inches(3.02), Inches(6.5), Inches(0.88),
   "We’re not building another chatbot. We’re building the automated infrastructure "
   "that makes every agent on your platform more productive than a franchise agent next door.",
   SERIF, 15, WHITE, italic=True)

# 4 numbered pillars
col_w = Inches(2.8)
cols  = [P, P + col_w + Inches(0.15), P + 2*(col_w + Inches(0.15)), P + 3*(col_w + Inches(0.15))]
pillars = [
    ("01", "24/7 lead\nresponse",
     "Every buyer inquiry answered instantly, on WhatsApp."),
    ("02", "Intelligent\nqualification",
     "Budget, area, type, timeline captured before agent contact."),
    ("03", "Platform\ncompliance",
     "Automated FICA collection removes admin from agent workflows."),
    ("04", "Agent\nrecruitment",
     "Automated intake and nurture at agents.homeseeker.co.za."),
]
for (num, title, desc), x in zip(pillars, cols):
    card(s3, x, Inches(4.35), col_w, Inches(2.65), rounded=True)
    tb(s3, x + Inches(0.25), Inches(4.55), Inches(0.6), Inches(0.42),
       num, SERIF, 22, VIOLET)
    tb(s3, x + Inches(0.25), Inches(5.05), col_w - Inches(0.4), Inches(0.65),
       title, SERIF, 16, WHITE)
    tb(s3, x + Inches(0.25), Inches(5.78), col_w - Inches(0.4), Inches(0.88),
       desc, SANS, 11, BODY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SYSTEMS OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
s4 = new_slide(prs)
slide_header(s4, "03  ·  AUTOMATION SYSTEMS", 4)

eyebrow(s4, "WHAT THE AUTOMATION LAYER BRINGS", P, Inches(0.85))
headline(s4, "Five AI Systems for Home Seekers", P, Inches(1.18), size=36)

systems = [
    ("A", "WhatsApp AI Buyer Agent",   "24/7 buyer inquiry response."),
    ("B", "Lead Routing Engine",        "Property24 to right agent, automatically."),
    ("C", "FICA Collection Flow",       "Compliance docs — no manual chasing."),
    ("D", "Agent Recruitment AI",       "Screening and onboarding automation."),
    ("E", "Buyer Nurture System",       "Re-engage cold leads, book more viewings."),
]
sw_each = Inches(2.25)
gap     = Inches(0.17)
sx      = P

for letter, name, tagline in systems:
    card(s4, sx, Inches(2.6), sw_each, Inches(4.1), rounded=True)
    # Letter badge
    badge = s4.shapes.add_shape(1, sx + Inches(0.25), Inches(2.85), Inches(0.55), Inches(0.55))
    badge.fill.solid()
    badge.fill.fore_color.rgb = VIOLET
    badge.line.fill.background()
    tb(s4, sx + Inches(0.25), Inches(2.87), Inches(0.55), Inches(0.48),
       letter, SANS, 16, BG, bold=True, align=PP_ALIGN.CENTER)
    # Name
    tb(s4, sx + Inches(0.2), Inches(3.55), sw_each - Inches(0.4), Inches(0.75),
       name, SERIF, 15, WHITE)
    # Tagline
    tb(s4, sx + Inches(0.2), Inches(4.38), sw_each - Inches(0.4), Inches(0.55),
       tagline, SANS, 12, BODY)
    # Feature label
    tb(s4, sx + Inches(0.2), Inches(5.05), sw_each - Inches(0.4), Inches(0.28),
       "FEATURE", SANS, 8, MUTED, bold=True)
    sx += sw_each + gap


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — MODULE A: WHATSAPP AI BUYER AGENT
# ══════════════════════════════════════════════════════════════════════════════
s5 = new_slide(prs)
slide_header(s5, "03A  ·  WHATSAPP AI", 5)

eyebrow(s5, "MODULE A", P, Inches(0.85))
headline(s5, "WhatsApp AI Buyer Agent", P, Inches(1.18), size=36)
tb(s5, P, Inches(2.2), Inches(6.5), Inches(0.42),
   "Responds in seconds. Qualifies in minutes. Books viewings without the agent being online.",
   SANS, 14, SUB)

hline(s5, P, Inches(2.75), SW - 2 * P, BORDER)

# Left: feature bullets
t = Inches(3.0)
for title, desc in [
    ("24/7 instant response",
     "No inquiry goes unanswered — ever, including Sundays at 9pm."),
    ("Buyer pre-qualification",
     "Budget, area, property type, timeline captured before agent contact."),
    ("Viewing booking",
     "AI books the slot directly — no manual back-and-forth."),
    ("Lead routing",
     "Qualified leads routed to the right agent by area and type."),
    ("Portal integration",
     "Works with Property24, Private Property, and homeseeker.co.za."),
    ("WhatsApp-native",
     "The channel SA buyers already use — no app download."),
]:
    t = bullet(s5, "·", title, desc, P, t, Inches(5.5), 13, 12)

# Right: stat cards + chat mockup
card(s5, Inches(7.5), Inches(3.0), Inches(5.0), Inches(3.7), rounded=True)
tb(s5, Inches(7.8), Inches(3.2), Inches(4.5), Inches(0.28),
   "BUYER  ·  QUALIFICATION FLOW", SANS, 9, MUTED, bold=True)
tb(s5, Inches(7.8), Inches(3.6), Inches(4.5), Inches(0.42),
   "From inquiry", SERIF, 22, WHITE)
tb(s5, Inches(7.8), Inches(4.08), Inches(4.5), Inches(0.42),
   "to booked viewing.", SERIF, 22, VIOLET)

# Three stats
for i, (label, value) in enumerate([
    ("LISTINGS", "416+"), ("COMMISSION", "100%"), ("FOUNDED", "2025")
]):
    sx = Inches(7.8) + i * Inches(1.58)
    tb(s5, sx, Inches(5.0), Inches(1.4), Inches(0.28), label, SANS, 8, MUTED, bold=True)
    tb(s5, sx, Inches(5.32), Inches(1.4), Inches(0.42), value, SERIF, 20, WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — MODULE B & C: LEAD ROUTING + FICA
# ══════════════════════════════════════════════════════════════════════════════
s6 = new_slide(prs)
slide_header(s6, "03B / 03C  ·  ROUTING & FICA", 6)

# Left half — Lead Routing
eyebrow(s6, "MODULE B", P, Inches(0.85))
headline(s6, "Lead Routing Engine", P, Inches(1.18), size=28)
tb(s6, P, Inches(2.05), Inches(5.5), Inches(0.38),
   "Every inquiry goes to the right agent — automatically.", SANS, 13, SUB)

t = Inches(2.6)
for title, desc in [
    ("Residential",      "For-sale inquiries routed to area specialists."),
    ("Commercial",       "Office and retail leases to commercial agents."),
    ("Rental / To-Let",  "To-let inquiries to rental-specialist agents."),
    ("Farm & Land",      "Niche types routed to the right specialist."),
    ("Missed Lead Recovery", "Unanswered inquiries re-routed after 30 min."),
]:
    t = bullet(s6, "·", title, desc, P, t, Inches(5.5), 12, 11)

hline(s6, P + Inches(0.3), Inches(2.6) + Inches(5 * 0.72) + Inches(0.15),
      Inches(4.5), BORDER)
tb(s6, P + Inches(0.3), t + Inches(0.05), Inches(4.5), Inches(0.28),
   "EACH ROUTED LEAD INCLUDES", SANS, 9, MUTED, bold=True)
tb(s6, P + Inches(0.3), t + Inches(0.38), Inches(4.5), Inches(0.28),
   "Lead source  ·  Property interest  ·  Budget  ·  Timeline  ·  WhatsApp thread",
   SANS, 11, BODY)

# Vertical divider
hline(s6, Inches(6.67) - Pt(0.5), Inches(0.85), Pt(1), BORDER)

# Right half — FICA
CR = Inches(6.85)
eyebrow(s6, "MODULE C", CR, Inches(0.85))
headline(s6, "FICA & Compliance Collection", CR, Inches(1.18), size=28)
tb(s6, CR, Inches(2.05), Inches(5.7), Inches(0.38),
   "The biggest time-drain for agents — automated end-to-end.", SANS, 13, SUB)

t = Inches(2.6)
for title, desc in [
    ("Saves 3-5 hours per deal",    "No more chasing the same document four times."),
    ("Zero manual follow-up",        "Reminders at 24h, 48h, 72h — agents never touch it."),
    ("Compliance confidence",        "Buyers understand exactly what's needed and why."),
    ("Agent satisfaction",           "Less admin, more deals — agents stay on the platform."),
    ("Full audit trail",             "Every request logged with timestamps."),
]:
    t = bullet(s6, "·", title, desc, CR, t, Inches(5.7), 12, 11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — MODULE D & E: RECRUITMENT + NURTURE
# ══════════════════════════════════════════════════════════════════════════════
s7 = new_slide(prs)
slide_header(s7, "03D / 03E  ·  RECRUITMENT & NURTURE", 7)

# Left — Recruitment
eyebrow(s7, "MODULE D", P, Inches(0.85))
headline(s7, "Agent Recruitment Automation", P, Inches(1.18), size=28)
tb(s7, P, Inches(2.05), Inches(5.5), Inches(0.38),
   "agents.homeseeker.co.za — AI intake, screening, and scheduling.", SANS, 13, SUB)

t = Inches(2.6)
for title, desc in [
    ("Automated application intake",
     "Every applicant enters an AI flow — no manual reading."),
    ("Qualifying questions",
     "AI asks: volume, areas, current agency, why switching."),
    ("Intro call scheduling",
     "Qualified applicants book a call directly in WhatsApp."),
    ("Nurture for undecideds",
     "3-touch sequence for candidates who don't book immediately."),
    ("Handover to Thomas",
     "Qualified leads land as clean summaries on your dashboard."),
]:
    t = bullet(s7, "·", title, desc, P, t, Inches(5.5), 12, 11)

# Right — Buyer Nurture
CR = Inches(6.85)
eyebrow(s7, "MODULE E", CR, Inches(0.85))
headline(s7, "Buyer Nurture System", CR, Inches(1.18), size=28)
tb(s7, CR, Inches(2.05), Inches(5.7), Inches(0.38),
   "Most inquiries don't convert first contact. Automated re-engagement changes that.",
   SANS, 13, SUB)

t = Inches(2.6)
for title, desc in [
    ("Re-engagement after 48h",
     "Non-converting buyers get a follow-up with a new matching listing."),
    ("Price drop alerts",
     "When an inquired listing drops in price, buyer is notified."),
    ("New listing matching",
     "New listings matching buyer criteria trigger a WhatsApp alert."),
    ("Viewing reschedule",
     "No-show buyers are automatically offered a new slot."),
    ("Post-viewing follow-up",
     "AI asks for feedback and offer intent after each viewing."),
    ("Post-sale review request",
     "Automated Google and Property24 review request at transfer date."),
]:
    t = bullet(s7, "·", title, desc, CR, t, Inches(5.7), 12, 11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — LEAD FLOW: INQUIRY → BOOKED VIEWING
# ══════════════════════════════════════════════════════════════════════════════
s8 = new_slide(prs)
slide_header(s8, "04  ·  LEAD FLOW", 8)

eyebrow(s8, "INQUIRY  →  BOOKED VIEWING", P, Inches(0.85))
headline(s8, "Automated Buyer Journey", P, Inches(1.18), size=36)
tb(s8, P, Inches(2.2), Inches(10), Inches(0.38),
   "Same flow for every inquiry — but each buyer feels like they are the only one.",
   SANS, 13, SUB)

hline(s8, P, Inches(2.72), SW - 2 * P, BORDER)

steps = [
    ("01", "Inquiry\nReceived",     "Property24 /\nhomeseeker.co.za"),
    ("02", "WhatsApp AI\nResponds", "Under 60 seconds"),
    ("03", "Qualification",         "Budget, area, type,\ntimeline captured"),
    ("04", "Agent\nMatched",        "Right agent notified\ninstantly"),
    ("05", "Viewing\nBooked",       "Confirmed in agent\ncalendar"),
    ("06", "FICA\nStarted",         "Document collection\nflow initiated"),
    ("07", "Nurture\nActive",       "Auto re-engagement\nif buyer goes cold"),
    ("08", "Deal\nClosed",          "Agent logs outcome;\nmetrics updated"),
]

step_w = Inches(1.4)
step_gap = Inches(0.12)
sx = P

for num, title, detail in steps:
    card(s8, sx, Inches(3.1), step_w, Inches(3.5), rounded=True)
    tb(s8, sx + Inches(0.12), Inches(3.28), step_w - Inches(0.2), Inches(0.38),
       num, SERIF, 20, VIOLET)
    tb(s8, sx + Inches(0.12), Inches(3.75), step_w - Inches(0.2), Inches(0.65),
       title, SERIF, 13, WHITE)
    tb(s8, sx + Inches(0.12), Inches(4.5), step_w - Inches(0.2), Inches(0.75),
       detail, SANS, 10, BODY)
    sx += step_w + step_gap
    if num not in ("08",):
        tb(s8, sx - step_gap + Inches(0.02), Inches(4.2), step_gap + Inches(0.05), Inches(0.28),
           "›", SANS, 14, MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — BENEFITS
# ══════════════════════════════════════════════════════════════════════════════
s9 = new_slide(prs)
slide_header(s9, "05  ·  BENEFITS", 9)

eyebrow(s9, "WHY THIS WORKS FOR HOME SEEKERS", P, Inches(0.85))
headline(s9, "What Changes After Launch", P, Inches(1.18), size=36)
hline(s9, P, Inches(2.72), SW - 2 * P, BORDER)

benefits = [
    ("Zero missed inquiries",      "Every buyer gets an instant response — including Sundays at 9pm."),
    ("Faster deal cycles",          "Pre-qualified buyers book viewings the same day they inquire."),
    ("Agent satisfaction",          "Less admin, more deal flow — agents stay on the platform longer."),
    ("FICA on autopilot",           "Document collection runs without agent involvement — every deal."),
    ("Recruitment on autopilot",    "Agent applications screened and nurtured 24/7."),
    ("Platform differentiation",    "AI infrastructure that franchise competitors can't replicate overnight."),
    ("Buyer re-engagement",         "Cold leads reactivated automatically with matching new listings."),
    ("Scales with your roster",     "More agents and listings don't add manual overhead."),
]

col1 = benefits[:4]
col2 = benefits[4:]
CL, CR = P, Inches(7.1)
t1, t2 = Inches(2.95), Inches(2.95)

for title, desc in col1:
    card(s9, CL, t1, Inches(5.75), Inches(0.82), rounded=True)
    tb(s9, CL + Inches(0.25), t1 + Inches(0.1), Inches(0.35), Inches(0.3),
       "✓", SANS, 14, VIOLET, bold=True)
    tb(s9, CL + Inches(0.65), t1 + Inches(0.08), Inches(4.8), Inches(0.3),
       title, SANS, 13, WHITE, bold=True)
    tb(s9, CL + Inches(0.65), t1 + Inches(0.4), Inches(4.8), Inches(0.35),
       desc, SANS, 11, BODY)
    t1 += Inches(0.95)

for title, desc in col2:
    card(s9, CR, t2, Inches(5.75), Inches(0.82), rounded=True)
    tb(s9, CR + Inches(0.25), t2 + Inches(0.1), Inches(0.35), Inches(0.3),
       "✓", SANS, 14, VIOLET, bold=True)
    tb(s9, CR + Inches(0.65), t2 + Inches(0.08), Inches(4.8), Inches(0.3),
       title, SANS, 13, WHITE, bold=True)
    tb(s9, CR + Inches(0.65), t2 + Inches(0.4), Inches(4.8), Inches(0.35),
       desc, SANS, 11, BODY)
    t2 += Inches(0.95)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — FUTURE ROADMAP
# ══════════════════════════════════════════════════════════════════════════════
s10 = new_slide(prs)
slide_header(s10, "06  ·  ROADMAP", 10)

eyebrow(s10, "BEYOND THE LAUNCH", P, Inches(0.85))
headline(s10, "Future Scalability", P, Inches(1.18), size=36)
tb(s10, P, Inches(2.2), Inches(10), Inches(0.38),
   "The same automation infrastructure scales as Home Seekers expands into new areas and service lines.",
   SANS, 13, SUB)
hline(s10, P, Inches(2.72), SW - 2 * P, BORDER)

future = [
    ("Seller Lead AI",          "Automate home valuation requests and seller onboarding."),
    ("Rental Management AI",    "Tenant screening and lease renewal reminders — automated."),
    ("Agent Performance Dashboard", "Real-time lead response times and conversion rates per agent."),
    ("Multi-city Expansion",    "Same AI infrastructure across Johannesburg, Cape Town, Durban."),
    ("Market Report AI",        "Automated monthly market updates to past buyers and sellers."),
    ("Multilingual AI",         "English, Afrikaans, Zulu — meet clients in their preferred language."),
    ("Investor Lead Nurture",   "Identify and nurture investment property buyers with a dedicated flow."),
    ("Developer Partnerships",  "Off-plan lead capture and qualification for new development projects."),
]

col1, col2 = future[:4], future[4:]
CL, CR = P, Inches(7.1)
t1, t2 = Inches(2.95), Inches(2.95)

for title, desc in col1:
    t1 = bullet(s10, "·", title, desc, CL, t1, Inches(5.75))
for title, desc in col2:
    t2 = bullet(s10, "·", title, desc, CR, t2, Inches(5.75))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — IMPLEMENTATION
# ══════════════════════════════════════════════════════════════════════════════
s11 = new_slide(prs)
slide_header(s11, "07  ·  IMPLEMENTATION", 11)

eyebrow(s11, "BUILT IN THREE CLEAN PHASES", P, Inches(0.85))
headline(s11, "Implementation Phases", P, Inches(1.18), size=36)
tb(s11, P, Inches(2.2), Inches(10), Inches(0.38),
   "Each phase ships visible value — no big-bang risk, no months of silence.",
   SANS, 13, SUB)
hline(s11, P, Inches(2.72), SW - 2 * P, BORDER)

phases = [
    ("PHASE 01  ·  Week 1",
     "WhatsApp AI Buyer Agent",
     "WhatsApp AI trained on HomeSeeker listings, areas, and pricing. "
     "Buyer inquiry intake, viewing booking, and agent notification live.",
     ["AI training on listings data", "Buyer inquiry intake flow",
      "Viewing booking integration", "Property24 / Private Property routing",
      "Agent notification system", "Go-live and testing"]),
    ("PHASE 02  ·  Weeks 2-3",
     "FICA & Lead Routing",
     "FICA document collection workflow and lead routing by area and property type activated.",
     ["FICA document flow design", "Compliance requirements mapping",
      "Document request automation", "Agent area mapping",
      "Platform dashboard", "Reminder sequence tuning"]),
    ("PHASE 03  ·  Weeks 4-5",
     "Nurture & Recruitment AI",
     "Buyer nurture sequences and agent recruitment intake go live at agents.homeseeker.co.za.",
     ["Buyer nurture drip sequences", "New listing match alerts",
      "Post-viewing follow-up flow", "Agent recruitment intake AI",
      "Review request sequence", "Training & handover"]),
]

ph_w = Inches(3.8)
ph_gap = Inches(0.17)
px = P

for phase_label, phase_name, phase_desc, deliverables in phases:
    card(s11, px, Inches(3.0), ph_w, Inches(3.75), rounded=True)
    tb(s11, px + Inches(0.25), Inches(3.15), ph_w - Inches(0.4), Inches(0.28),
       phase_label, SANS, 9, VIOLET, bold=True)
    tb(s11, px + Inches(0.25), Inches(3.5), ph_w - Inches(0.4), Inches(0.42),
       phase_name, SERIF, 15, WHITE)
    tb(s11, px + Inches(0.25), Inches(4.0), ph_w - Inches(0.4), Inches(0.52),
       phase_desc, SANS, 11, BODY)
    t = Inches(4.62)
    for d in deliverables[:4]:
        tb(s11, px + Inches(0.25), t, ph_w - Inches(0.4), Inches(0.26),
           f"›  {d}", SANS, 10, MUTED)
        t += Inches(0.3)
    px += ph_w + ph_gap


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — CLOSING
# ══════════════════════════════════════════════════════════════════════════════
s12 = new_slide(prs)

# Subtle top accent
acc = s12.shapes.add_shape(1, Inches(0), Inches(0), SW, Pt(3))
acc.fill.solid()
acc.fill.fore_color.rgb = VIOLET
acc.line.fill.background()

eyebrow(s12, "WHO IS BEHIND THIS", P, Inches(0.85))
headline(s12, "GenosAI Tech", P, Inches(1.18), size=40)
tb(s12, P, Inches(2.22), Inches(6.5), Inches(0.66),
   "GenosAI builds AI automation systems for real estate platforms, agencies, and property developers. "
   "We specialise in WhatsApp AI, lead qualification flows, compliance automations, and recruitment systems "
   "— production-grade from week one.",
   SANS, 13, SUB)

hline(s12, P, Inches(3.02), Inches(6), BORDER)

# What we build
tb(s12, P, Inches(3.2), Inches(4), Inches(0.28),
   "WHAT WE BUILD", SANS, 9, MUTED, bold=True)
t = Inches(3.58)
for item, desc in [
    ("WhatsApp AI Agents",        "Trained on your listings, areas, and pricing."),
    ("Lead Qualification Flows",  "Pre-qualify budget, area, type, timeline."),
    ("FICA & Compliance AI",      "Document collection without manual chasing."),
    ("Recruitment Intake AI",     "Screen and nurture agent applicants."),
    ("Buyer Nurture Systems",     "Re-engage cold leads with matching listings."),
]:
    t = bullet(s12, "◆", item, desc, P, t, Inches(5.5), 12, 11)

# RIGHT — next steps + contact
card(s12, Inches(7.5), Inches(1.1), Inches(5.05), Inches(5.75), rounded=True)

# Closing quote
tb(s12, Inches(7.75), Inches(1.35), Inches(4.6), Inches(0.35),
   "CLOSING NOTE", SANS, 9, MUTED, bold=True)
tb(s12, Inches(7.6), Inches(1.72), Inches(0.4), Inches(0.55),
   "“", SERIF, 28, VIOLET)
tb(s12, Inches(7.95), Inches(1.78), Inches(4.3), Inches(1.0),
   "Built to make every Home Seekers agent the fastest to respond, in every area they operate.",
   SERIF, 14, WHITE, italic=True)

hline(s12, Inches(7.75), Inches(2.88), Inches(4.45), BORDER)

tb(s12, Inches(7.75), Inches(3.08), Inches(4.5), Inches(0.28),
   "NEXT STEPS", SANS, 9, VIOLET, bold=True)
t = Inches(3.45)
for step, desc in [
    ("01", "15-minute call to confirm scope and integration points."),
    ("02", "Phase 1 kick-off — WhatsApp AI buyer agent live within 7 days."),
    ("03", "First inquiry routed automatically — results visible from day one."),
]:
    tb(s12, Inches(7.75), t, Inches(0.4), Inches(0.3), step, SERIF, 14, VIOLET)
    tb(s12, Inches(8.2), t + Inches(0.02), Inches(4.15), Inches(0.5),
       desc, SANS, 12, BODY)
    t += Inches(0.62)

hline(s12, Inches(7.75), Inches(5.45), Inches(4.45), BORDER)

tb(s12, Inches(7.75), Inches(5.65), Inches(4.5), Inches(0.28),
   "Rohan Malik  ·  CEO, Genos AI", SANS, 12, WHITE, bold=True)
tb(s12, Inches(7.75), Inches(6.0), Inches(4.5), Inches(0.28),
   "hello@genosai.tech  ·  +91 638-714-2699", SANS, 11, BODY)
tb(s12, Inches(7.75), Inches(6.32), Inches(4.5), Inches(0.28),
   "www.genosai.tech", SANS, 11, MUTED)

# Footer
tb(s12, P, SH - Inches(0.42), Inches(8), Inches(0.28),
   "Confidential  ·  Home Seekers  ·  2026", SANS, 8, DIM)
tb(s12, Inches(9.5), SH - Inches(0.42), Inches(3.08), Inches(0.28),
   f"{TOTAL:02d} / {TOTAL}", SANS, 8, DIM, align=PP_ALIGN.RIGHT)

# GenosAI wordmark — bottom
tb(s12, P, SH - Inches(0.65), Inches(3), Inches(0.4),
   "GenosAI", SERIF, 20, WHITE)


# ── Save ─────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved: {OUT}  ({TOTAL} slides)")
