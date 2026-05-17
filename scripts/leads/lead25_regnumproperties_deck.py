"""
Build RegnumProperties_Pitch_GenosAI.pptx from scratch.
Track: Pitch-only (user specified automation; score 6.5).
Design: docs/BRAND_ASSETS.md -- #0A0A0F canvas, Instrument Serif headlines, Satoshi body.
12 slides: cover, developer overview, vision, systems, 5 module slides,
buyer journey flow, implementation, closing.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).parent.parent.parent
OUT  = ROOT / "output" / "decks" / "RegnumProperties_Pitch_GenosAI.pptx"

# -- Brand colors --------------------------------------------------------------
BG     = RGBColor(0x0A, 0x0A, 0x0F)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BODY   = RGBColor(0x9D, 0x9D, 0x9F)
SUB    = RGBColor(0x86, 0x86, 0x8A)
MUTED  = RGBColor(0x6B, 0x6B, 0x6E)
DIM    = RGBColor(0x44, 0x44, 0x47)
VIOLET = RGBColor(0xA7, 0x8B, 0xFA)
CARD   = RGBColor(0x11, 0x11, 0x18)
BORDER = RGBColor(0x22, 0x22, 0x2A)

SERIF = "Instrument Serif"
SANS  = "Satoshi"

SW    = Inches(13.33)
SH    = Inches(7.5)
P     = Inches(0.75)
TOTAL = 12


# -- Helpers -------------------------------------------------------------------

def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def tb(slide, l, t, w, h, text, font=SANS, size=14, color=BODY,
       bold=False, italic=False, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text = text
    r.font.name   = font
    r.font.size   = Pt(size)
    r.font.color.rgb = color
    r.font.bold   = bold
    r.font.italic = italic
    return box


def tb_lines(slide, l, t, w, h, lines, font=SANS, size=14, color=BODY,
             bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name  = font
        r.font.size  = Pt(size)
        r.font.color.rgb = color
        r.font.bold  = bold
    return box


def card(slide, l, t, w, h, rounded=True):
    s = slide.shapes.add_shape(5 if rounded else 1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = CARD
    s.line.color.rgb = BORDER
    s.line.width = Pt(0.75)
    return s


def hline(slide, l, t, w, color=BORDER):
    s = slide.shapes.add_shape(1, l, t, w, Pt(0.75))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def eyebrow(slide, text, l, t, w=Inches(6)):
    return tb(slide, l, t, w, Inches(0.3),
              text.upper(), SANS, 10, VIOLET, bold=True)


def headline(slide, text, l, t, w=Inches(10), size=42):
    return tb(slide, l, t, w, Inches(1.6), text, SERIF, size, WHITE)


def slide_header(slide, section, page):
    tb(slide, P, Inches(0.22), Inches(6), Inches(0.28),
       "GENOSAI  x  REGNUM PROPERTIES", SANS, 9, MUTED, bold=True)
    tb(slide, Inches(7.5), Inches(0.22), Inches(5.08), Inches(0.28),
       section, SANS, 9, MUTED, bold=True, align=PP_ALIGN.RIGHT)
    hline(slide, P, Inches(0.58), SW - 2 * P)
    tb(slide, P, SH - Inches(0.42), Inches(8), Inches(0.28),
       "AI Automation Proposal  .  Regnum Properties  .  2026", SANS, 8, DIM)
    tb(slide, Inches(9.5), SH - Inches(0.42), Inches(3.08), Inches(0.28),
       f"{page:02d} / {TOTAL}", SANS, 8, DIM, align=PP_ALIGN.RIGHT)


def bullet(slide, dot_char, title, desc, l, t,
           col_w=Inches(5.0), title_size=13, desc_size=12):
    tb(slide, l, t, Inches(0.25), Inches(0.26),
       dot_char, SANS, 13, VIOLET, bold=True)
    tb(slide, l + Inches(0.3), t, col_w - Inches(0.3), Inches(0.26),
       title, SANS, title_size, WHITE, bold=True)
    tb(slide, l + Inches(0.3), t + Inches(0.27), col_w - Inches(0.3), Inches(0.38),
       desc, SANS, desc_size, BODY)
    return t + Inches(0.72)


# =============================================================================
# SLIDE 1 -- COVER
# =============================================================================
prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH

s1 = new_slide(prs)

accent = s1.shapes.add_shape(1, Inches(0), Inches(0), SW, Pt(3))
accent.fill.solid()
accent.fill.fore_color.rgb = VIOLET
accent.line.fill.background()

tb(s1, P, Inches(1.0), Inches(5.5), Inches(0.35),
   "PREPARED FOR", SANS, 10, MUTED, bold=True)
hline(s1, P, Inches(1.42), Inches(5.5), BORDER)

tb(s1, P, Inches(1.6), Inches(8.5), Inches(1.4),
   "Regnum Properties", SERIF, 64, WHITE)

tb(s1, P, Inches(3.1), Inches(7), Inches(0.32),
   "Luxury Residential Development  .  East Legon, Accra  .  Diaspora & International Buyers",
   SANS, 13, BODY)

tb(s1, P, Inches(3.58), Inches(6), Inches(0.28),
   "Accra, Ghana  .  The Alexandria  .  The Arlington", SANS, 11, MUTED)

hline(s1, P, Inches(4.04), Inches(6), BORDER)

tb(s1, P, Inches(4.24), Inches(7.5), Inches(0.28),
   "AI AUTOMATION PROPOSAL", SANS, 10, VIOLET, bold=True)
tb_lines(s1, P, Inches(4.64), Inches(6.5), Inches(1.4),
         ["An AI automation layer for Regnum Properties --",
          "built to qualify diaspora buyers, nurture international investors,",
          "book viewings 24/7, and close the gap between",
          "overseas interest and signed deals."],
         SANS, 13, SUB)

card(s1, Inches(9.0), Inches(1.2), Inches(3.58), Inches(2.5), rounded=True)
tb(s1, Inches(9.25), Inches(1.5), Inches(3.0), Inches(0.28),
   "PREPARED BY", SANS, 9, MUTED, bold=True)
tb(s1, Inches(9.25), Inches(1.88), Inches(3.0), Inches(0.32),
   "GenosAI Tech", SERIF, 20, WHITE)
tb(s1, Inches(9.25), Inches(2.3), Inches(3.0), Inches(0.26),
   "Rohan Malik  .  Founder & CEO", SANS, 11, BODY)
tb(s1, Inches(9.25), Inches(2.65), Inches(3.0), Inches(0.26),
   "hello@genosai.tech  .  www.genosai.tech", SANS, 10, MUTED)

tb(s1, P, SH - Inches(0.65), Inches(4), Inches(0.45),
   "GenosAI", SERIF, 22, WHITE)
tb(s1, P + Inches(1.65), SH - Inches(0.62), Inches(3), Inches(0.3),
   "v1.0  .  2026  .  Confidential", SANS, 9, DIM)


# =============================================================================
# SLIDE 2 -- THE DEVELOPER AT A GLANCE
# =============================================================================
s2 = new_slide(prs)
slide_header(s2, "01  .  THE DEVELOPER", 2)

eyebrow(s2, "REGNUM PROPERTIES AT A GLANCE", P, Inches(0.85))
headline(s2, "Luxury Homes in Accra's\nPremium Corridor. Built for a Global Buyer.", P, Inches(1.18), size=34)
tb(s2, P, Inches(2.32), Inches(9), Inches(0.38),
   "Regnum Properties develops gated luxury communities in East Legon -- "
   "Accra's most prestigious residential corridor -- and markets them directly "
   "to Ghanaian diaspora investors and international buyers. "
   "Smart home integrations, 24/7 security, and fully managed communities. "
   "The product is positioned globally. The sales infrastructure is not.",
   SANS, 13, SUB)

hline(s2, P, Inches(2.95), SW - 2 * P, BORDER)

stats = [
    ("LOCATION",     "East Legon",    "Trasacco Phase 1 corridor, Accra"),
    ("PROJECTS",     "2 Active",      "The Alexandria & The Arlington"),
    ("BUYER FOCUS",  "Diaspora",      "International & overseas Ghanaian investors"),
    ("TECHNOLOGY",   "Smart Homes",   "Integrated controls in all units"),
]
sw_each = Inches(2.8)
gap     = Inches(0.19)
sx      = P

for label, value, sub in stats:
    card(s2, sx, Inches(3.15), sw_each, Inches(2.0), rounded=True)
    tb(s2, sx + Inches(0.25), Inches(3.33), sw_each - Inches(0.4), Inches(0.28),
       label, SANS, 9, MUTED, bold=True)
    tb(s2, sx + Inches(0.25), Inches(3.7), sw_each - Inches(0.4), Inches(0.55),
       value, SERIF, 22, WHITE)
    tb(s2, sx + Inches(0.25), Inches(4.38), sw_each - Inches(0.4), Inches(0.45),
       sub, SANS, 11, BODY)
    sx += sw_each + gap

tb(s2, P, Inches(5.42), Inches(11), Inches(0.28),
   "PROJECT PORTFOLIO", SANS, 9, MUTED, bold=True)
hline(s2, P, Inches(5.74), SW - 2 * P, BORDER)
tb(s2, P, Inches(5.84), SW - 2 * P, Inches(0.55),
   "The Alexandria -- 4.5-bedroom smart townhouses, Trasacco Phase 1 area, gated community, pool, backup power  .  "
   "The Arlington -- 3.5-bedroom gated townhouses, 5-unit community, smart home integrations, 24/7 security",
   SANS, 12, BODY)


# =============================================================================
# SLIDE 3 -- VISION
# =============================================================================
s3 = new_slide(prs)
slide_header(s3, "02  .  VISION", 3)

eyebrow(s3, "THE SALES GAP", P, Inches(0.85))
headline(s3, "The Product Is Global.\nThe Sales Process Is Not.", P, Inches(1.18), size=38)

card(s3, P, Inches(2.85), Inches(7.5), Inches(1.22), rounded=True)
tb(s3, P + Inches(0.4), Inches(2.98), Inches(7.0), Inches(0.4),
   '"', SERIF, 32, VIOLET)
tb(s3, P + Inches(0.75), Inches(3.02), Inches(6.5), Inches(0.85),
   "A diaspora buyer in London sees The Arlington at 11 PM, fills in a contact form, "
   "and wakes up to silence. The interest was real. The follow-up was not. "
   "AI automation closes that gap -- at any hour, from any time zone.",
   SERIF, 15, WHITE, italic=True)

col_w = Inches(2.8)
cols  = [P, P + col_w + Inches(0.15), P + 2*(col_w + Inches(0.15)), P + 3*(col_w + Inches(0.15))]
pillars = [
    ("01", "24/7 Lead\nCapture",
     "Every enquiry qualified the moment it arrives -- regardless of time zone."),
    ("02", "Investor\nNurture",
     "International buyers followed up automatically for 12 weeks after first contact."),
    ("03", "Social\nAudience",
     "Consistent content builds the diaspora buyer audience without founder time."),
    ("04", "CRM\nPipeline",
     "Every buyer conversation tracked in one place as the portfolio scales."),
]
for (num, title, desc), x in zip(pillars, cols):
    card(s3, x, Inches(4.35), col_w, Inches(2.65), rounded=True)
    tb(s3, x + Inches(0.25), Inches(4.55), Inches(0.6), Inches(0.42),
       num, SERIF, 22, VIOLET)
    tb(s3, x + Inches(0.25), Inches(5.05), col_w - Inches(0.4), Inches(0.65),
       title, SERIF, 16, WHITE)
    tb(s3, x + Inches(0.25), Inches(5.78), col_w - Inches(0.4), Inches(0.88),
       desc, SANS, 11, BODY)


# =============================================================================
# SLIDE 4 -- SYSTEMS OVERVIEW
# =============================================================================
s4 = new_slide(prs)
slide_header(s4, "03  .  AUTOMATION SYSTEMS", 4)

eyebrow(s4, "FIVE AI SYSTEMS FOR REGNUM PROPERTIES", P, Inches(0.85))
headline(s4, "The Automation Blueprint", P, Inches(1.18), size=36)

systems = [
    ("A", "WhatsApp\nLead AI",        "24/7 buyer qualification via WhatsApp."),
    ("B", "Investor\nNurture AI",     "12-week automated follow-up sequences."),
    ("C", "Social Media\nContent AI", "LinkedIn & Instagram on autopilot."),
    ("D", "Property\nChatbot",        "24/7 inquiry and booking on the website."),
    ("E", "CRM\nPipeline AI",         "Automated deal tracking and follow-ups."),
]
sw_each = Inches(2.25)
gap     = Inches(0.17)
sx      = P

for letter, name, tagline in systems:
    card(s4, sx, Inches(2.6), sw_each, Inches(4.1), rounded=True)
    badge = s4.shapes.add_shape(1, sx + Inches(0.25), Inches(2.85), Inches(0.55), Inches(0.55))
    badge.fill.solid()
    badge.fill.fore_color.rgb = VIOLET
    badge.line.fill.background()
    tb(s4, sx + Inches(0.25), Inches(2.87), Inches(0.55), Inches(0.48),
       letter, SANS, 16, BG, bold=True, align=PP_ALIGN.CENTER)
    tb(s4, sx + Inches(0.2), Inches(3.55), sw_each - Inches(0.4), Inches(0.75),
       name, SERIF, 15, WHITE)
    tb(s4, sx + Inches(0.2), Inches(4.38), sw_each - Inches(0.4), Inches(0.55),
       tagline, SANS, 12, BODY)
    sx += sw_each + gap


# =============================================================================
# SLIDE 5 -- MODULE A: WHATSAPP LEAD AI
# =============================================================================
s5 = new_slide(prs)
slide_header(s5, "03A  .  WHATSAPP LEAD AI", 5)

eyebrow(s5, "MODULE A", P, Inches(0.85))
headline(s5, "WhatsApp Lead Qualification AI", P, Inches(1.18), size=34)
tb(s5, P, Inches(2.2), Inches(6.5), Inches(0.42),
   "A WhatsApp AI number displayed on the Regnum website and social profiles "
   "qualifies every buyer the moment they reach out -- budget, property, timeline -- "
   "at any hour, from any city in the world.",
   SANS, 14, SUB)

hline(s5, P, Inches(2.75), SW - 2 * P, BORDER)

t = Inches(3.0)
for title, desc in [
    ("Instant 24/7 response",
     "When a buyer in the UK or Canada sends a WhatsApp message at midnight Accra time, "
     "the AI responds immediately -- capturing the moment instead of losing it to silence."),
    ("Structured qualification intake",
     "AI captures budget range, preferred project (Alexandria or Arlington), "
     "number of bedrooms, purchase timeline, and whether the buyer is cash or finance. "
     "Full context before any agent is involved."),
    ("Viewing booking automation",
     "Qualified buyers with a 0-3 month timeline receive an automated viewing invite "
     "and booking link. The agent arrives prepared -- not to an unqualified cold enquiry."),
    ("International buyer routing",
     "Diaspora buyers in London, Houston, or Toronto routed to a virtual viewing "
     "sequence: floor plan delivery, video walkthrough link, and a scheduled video call "
     "with Laud or the sales team."),
    ("Warm lead handover",
     "When a buyer needs human involvement, the AI transfers the full conversation "
     "summary to the sales team via WhatsApp or email -- no context lost, no repeated questions."),
]:
    t = bullet(s5, ".", title, desc, P, t, Inches(5.5), 13, 12)

card(s5, Inches(7.5), Inches(3.0), Inches(5.0), Inches(3.7), rounded=True)
tb(s5, Inches(7.8), Inches(3.2), Inches(4.5), Inches(0.28),
   "THE TIME ZONE PROBLEM", SANS, 9, MUTED, bold=True)
tb(s5, Inches(7.8), Inches(3.6), Inches(4.5), Inches(0.55),
   "Diaspora buyers enquire at 11 PM", SERIF, 17, WHITE)
tb(s5, Inches(7.8), Inches(4.18), Inches(4.5), Inches(0.55),
   "in their time zone -- not yours.", SERIF, 16, VIOLET)
tb(s5, Inches(7.8), Inches(4.88), Inches(4.5), Inches(0.72),
   "Without a 24/7 qualification system, "
   "every after-hours enquiry either waits until Monday "
   "or is lost to a competitor who responded first.", SANS, 11, BODY)


# =============================================================================
# SLIDE 6 -- MODULE B: INVESTOR NURTURE SEQUENCES
# =============================================================================
s6 = new_slide(prs)
slide_header(s6, "03B  .  INVESTOR NURTURE AI", 6)

eyebrow(s6, "MODULE B", P, Inches(0.85))
headline(s6, "Investor Nurture Sequences", P, Inches(1.18), size=36)
tb(s6, P, Inches(2.2), Inches(6.5), Inches(0.42),
   "Most luxury property deals in Ghana close on the third or fourth contact. "
   "Most developers follow up once, then stop. "
   "Automated nurture sequences run for 12 weeks without any manual effort.",
   SANS, 14, SUB)

hline(s6, P, Inches(2.75), SW - 2 * P, BORDER)

t = Inches(3.0)
for title, desc in [
    ("Construction progress updates",
     "Buyers who have enquired but not committed receive automated updates "
     "as each project milestone is reached -- foundations, framing, finishes. "
     "The project feels real and moving to an investor 5,000 km away."),
    ("Payment plan education sequence",
     "A 4-step automated sequence explaining payment structure, deposit requirements, "
     "completion timeline, and title deed process -- removing every friction point "
     "in the international buyer's decision without a single sales call."),
    ("Accra market context delivery",
     "Monthly AI-generated market update sent to all active leads: "
     "East Legon price trends, rental yield data, diaspora investment statistics. "
     "Positions Regnum as the expert, not just the seller."),
    ("Reactivation sequences",
     "Leads who went cold at 60 or 90 days are automatically reactivated "
     "with a new angle: a price update, a unit availability reminder, "
     "or a construction milestone that changes the proposition."),
    ("Referral trigger",
     "Buyers who complete a purchase receive an automated referral prompt at "
     "60 days post-handover -- the highest satisfaction moment in the relationship."),
]:
    t = bullet(s6, ".", title, desc, P, t, Inches(5.5), 13, 12)

card(s6, Inches(7.5), Inches(3.0), Inches(5.0), Inches(3.7), rounded=True)
tb(s6, Inches(7.8), Inches(3.2), Inches(4.5), Inches(0.28),
   "THE FOLLOW-UP REALITY", SANS, 9, MUTED, bold=True)
tb(s6, Inches(7.8), Inches(3.62), Inches(4.5), Inches(0.6),
   "Deals close on the 3rd or 4th contact.", SERIF, 17, WHITE)
tb(s6, Inches(7.8), Inches(4.28), Inches(4.5), Inches(0.55),
   "Most developers stop at the first.", SERIF, 16, VIOLET)
tb(s6, Inches(7.8), Inches(4.95), Inches(4.5), Inches(0.68),
   "A 12-week automated nurture sequence "
   "ensures Regnum stays top of mind for every active lead -- "
   "without Laud or Maxwell sending a single manual message.", SANS, 11, BODY)


# =============================================================================
# SLIDE 7 -- MODULE C: SOCIAL MEDIA CONTENT AI
# =============================================================================
s7 = new_slide(prs)
slide_header(s7, "03C  .  SOCIAL MEDIA AI", 7)

eyebrow(s7, "MODULE C", P, Inches(0.85))
headline(s7, "Social Media & Content AI", P, Inches(1.18), size=34)
tb(s7, P, Inches(2.2), Inches(6.5), Inches(0.42),
   "Regnum is actively posting on LinkedIn and Instagram -- the intent is right. "
   "An AI content layer generates consistent, audience-specific posts "
   "on a weekly schedule without consuming founder time.",
   SANS, 14, SUB)

hline(s7, P, Inches(2.75), SW - 2 * P, BORDER)

t = Inches(3.0)
for title, desc in [
    ("Weekly LinkedIn content calendar",
     "AI generates 3 posts per week: an Accra market insight, a project construction update, "
     "and a diaspora buyer guide. Scheduled and published automatically -- "
     "growing the audience from its current base without manual effort."),
    ("Instagram property showcase posts",
     "AI-written captions for construction photos and finished unit imagery, "
     "optimised for diaspora buyer search terms and hashtags. "
     "Consistent visual storytelling builds the brand before a buyer ever enquires."),
    ("Buyer guide content series",
     "'How to buy property in Ghana from the UK' / 'Ghana property investment for diaspora' -- "
     "AI-generated long-form guides published as LinkedIn articles and Instagram carousels, "
     "positioning Laud as the expert voice in diaspora real estate investment."),
    ("Project milestone announcements",
     "Every construction milestone converted into a social post automatically: "
     "foundation complete, roofing complete, smart home installation -- "
     "keeping the audience engaged through the build cycle."),
    ("Engagement and comment responses",
     "AI monitors LinkedIn and Instagram comments and drafts responses for review, "
     "ensuring no inbound interest from a follower is missed."),
]:
    t = bullet(s7, ".", title, desc, P, t, Inches(5.5), 13, 12)

card(s7, Inches(7.5), Inches(3.0), Inches(5.0), Inches(3.7), rounded=True)
tb(s7, Inches(7.8), Inches(3.2), Inches(4.5), Inches(0.28),
   "THE AUDIENCE OPPORTUNITY", SANS, 9, MUTED, bold=True)
tb(s7, Inches(7.8), Inches(3.62), Inches(4.5), Inches(0.6),
   "The Ghanaian diaspora is 3M+ strong", SERIF, 17, WHITE)
tb(s7, Inches(7.8), Inches(4.28), Inches(4.5), Inches(0.55),
   "and actively investing in Accra.", SERIF, 16, VIOLET)
tb(s7, Inches(7.8), Inches(4.95), Inches(4.5), Inches(0.68),
   "Consistent social content positions Regnum as the developer "
   "diaspora buyers find first -- "
   "before they find a competitor.", SANS, 11, BODY)


# =============================================================================
# SLIDE 8 -- MODULE D & E: CHATBOT + CRM
# =============================================================================
s8 = new_slide(prs)
slide_header(s8, "03D / 03E  .  CHATBOT & CRM", 8)

eyebrow(s8, "MODULE D", P, Inches(0.85))
headline(s8, "Property Inquiry Chatbot", P, Inches(1.18), size=28)
tb(s8, P, Inches(2.05), Inches(5.5), Inches(0.38),
   "A 24/7 AI chat widget on theregnumproperties.com -- floor plans, pricing, "
   "availability, and viewing bookings handled without human involvement.",
   SANS, 13, SUB)

t = Inches(2.6)
for title, desc in [
    ("Property finder conversation",
     "A visitor describes what they want -- bedrooms, location, budget -- "
     "and the AI recommends the right Regnum project with specific unit details."),
    ("Pricing and payment plan delivery",
     "Exact pricing, payment schedule, and deposit requirements delivered "
     "on request -- removing the need for a buyer to chase a quote."),
    ("Virtual tour and floor plan delivery",
     "Floor plans and video walkthrough links sent automatically within the chat, "
     "enabling diaspora buyers to explore units fully before a viewing."),
    ("Viewing booking",
     "In-person and virtual viewings booked directly in the chat -- "
     "confirmed in the buyer's calendar and Laud's, no back-and-forth."),
]:
    t = bullet(s8, ".", title, desc, P, t, Inches(5.5), 12, 11)

hline(s8, Inches(6.67) - Pt(0.5), Inches(0.85), Pt(1), BORDER)

CR = Inches(6.85)
eyebrow(s8, "MODULE E", CR, Inches(0.85))
headline(s8, "CRM Pipeline Automation", CR, Inches(1.18), size=28)
tb(s8, CR, Inches(2.05), Inches(5.7), Inches(0.38),
   "Every buyer conversation in one pipeline -- from first WhatsApp to signed contract. "
   "Automated reminders ensure no deal stalls due to a missed follow-up.",
   SANS, 13, SUB)

t = Inches(2.6)
for title, desc in [
    ("Automated deal stage tracking",
     "Each buyer moves through stages automatically: Enquired, Qualified, Viewing Booked, "
     "Offer Received, Under Contract. No manual updates required."),
    ("Follow-up reminder sequences",
     "When a buyer has not responded in 5 days, the CRM triggers an automated "
     "follow-up message -- preventing deals from stalling silently."),
    ("Communication history",
     "Every WhatsApp message, email, and call note logged in one buyer record. "
     "Laud and Maxwell always have full context when engaging a buyer."),
    ("Pipeline reporting",
     "Weekly automated summary: new leads this week, viewings scheduled, "
     "offers pending, projected completions. One view, always current."),
]:
    t = bullet(s8, ".", title, desc, CR, t, Inches(5.7), 12, 11)


# =============================================================================
# SLIDE 9 -- BUYER JOURNEY FLOW
# =============================================================================
s9 = new_slide(prs)
slide_header(s9, "04  .  BUYER JOURNEY", 9)

eyebrow(s9, "ENQUIRY  ->  QUALIFIED  ->  NURTURED  ->  VIEWED  ->  SIGNED", P, Inches(0.85))
headline(s9, "Automated Diaspora Buyer Journey", P, Inches(1.18), size=32)
tb(s9, P, Inches(2.2), Inches(10), Inches(0.38),
   "From the first WhatsApp message to a signed contract -- every stage automated, "
   "every buyer followed up, no deal lost to a missed response.",
   SANS, 13, SUB)

hline(s9, P, Inches(2.72), SW - 2 * P, BORDER)

steps = [
    ("01", "Buyer\nEnquires",        "WhatsApp or\nwebsite chatbot"),
    ("02", "AI\nQualifies",          "Budget, timeline,\nproperty captured"),
    ("03", "Floor Plan\nDelivered",  "Auto-sent with\nvirtual tour link"),
    ("04", "Viewing\nBooked",        "In-person or\nvideo call"),
    ("05", "Nurture\nSequence",      "12-week updates\nif not yet signed"),
    ("06", "Market\nUpdate",         "Monthly Accra\ndata delivered"),
    ("07", "Offer\nReceived",        "CRM stages\nupdated auto"),
    ("08", "Contract\nSigned",       "Referral prompt\nat day 60"),
]

step_w   = Inches(1.4)
step_gap = Inches(0.12)
sx = P

for num, title, detail in steps:
    card(s9, sx, Inches(3.1), step_w, Inches(3.5), rounded=True)
    tb(s9, sx + Inches(0.12), Inches(3.28), step_w - Inches(0.2), Inches(0.38),
       num, SERIF, 20, VIOLET)
    tb(s9, sx + Inches(0.12), Inches(3.75), step_w - Inches(0.2), Inches(0.65),
       title, SERIF, 13, WHITE)
    tb(s9, sx + Inches(0.12), Inches(4.5), step_w - Inches(0.2), Inches(0.75),
       detail, SANS, 10, BODY)
    sx += step_w + step_gap
    if num != "08":
        tb(s9, sx - step_gap + Inches(0.02), Inches(4.2),
           step_gap + Inches(0.05), Inches(0.28),
           ">", SANS, 14, MUTED, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 10 -- BENEFITS
# =============================================================================
s10 = new_slide(prs)
slide_header(s10, "05  .  BENEFITS", 10)

eyebrow(s10, "WHAT CHANGES AFTER LAUNCH", P, Inches(0.85))
headline(s10, "What the Automation Layer Delivers", P, Inches(1.18), size=34)
hline(s10, P, Inches(2.72), SW - 2 * P, BORDER)

benefits = [
    ("No enquiry goes cold",          "Every buyer is followed up automatically -- whether they enquired at 9 AM or 11 PM."),
    ("Diaspora buyers convert",       "International leads receive floor plans, payment plans, and video tours before a single call."),
    ("12 weeks of nurture per lead",  "Every enquiry followed up for 3 months without Laud sending a single manual message."),
    ("24/7 website qualification",    "The website works as a sales tool at night and on weekends -- not just a brochure."),
    ("Social presence built passively", "LinkedIn and Instagram grow consistently without consuming founder time."),
    ("Pipeline fully visible",        "Maxwell and Laud see every buyer, every stage, every follow-up due -- in one view."),
    ("Viewing bookings automated",    "No scheduling back-and-forth -- buyers book directly, calendar updated automatically."),
    ("Ready to scale",                "The same automation infrastructure supports The Alexandria, The Arlington, and every project after."),
]

col1, col2 = benefits[:4], benefits[4:]
CL, CR = P, Inches(7.1)
t1, t2 = Inches(2.95), Inches(2.95)

for title, desc in col1:
    card(s10, CL, t1, Inches(5.75), Inches(0.82), rounded=True)
    tb(s10, CL + Inches(0.25), t1 + Inches(0.1), Inches(0.35), Inches(0.3),
       "v", SANS, 14, VIOLET, bold=True)
    tb(s10, CL + Inches(0.65), t1 + Inches(0.08), Inches(4.8), Inches(0.3),
       title, SANS, 13, WHITE, bold=True)
    tb(s10, CL + Inches(0.65), t1 + Inches(0.4), Inches(4.8), Inches(0.35),
       desc, SANS, 11, BODY)
    t1 += Inches(0.95)

for title, desc in col2:
    card(s10, CR, t2, Inches(5.75), Inches(0.82), rounded=True)
    tb(s10, CR + Inches(0.25), t2 + Inches(0.1), Inches(0.35), Inches(0.3),
       "v", SANS, 14, VIOLET, bold=True)
    tb(s10, CR + Inches(0.65), t2 + Inches(0.08), Inches(4.8), Inches(0.3),
       title, SANS, 13, WHITE, bold=True)
    tb(s10, CR + Inches(0.65), t2 + Inches(0.4), Inches(4.8), Inches(0.35),
       desc, SANS, 11, BODY)
    t2 += Inches(0.95)


# =============================================================================
# SLIDE 11 -- IMPLEMENTATION
# =============================================================================
s11 = new_slide(prs)
slide_header(s11, "06  .  IMPLEMENTATION", 11)

eyebrow(s11, "BUILT IN THREE CLEAN PHASES", P, Inches(0.85))
headline(s11, "Implementation Phases", P, Inches(1.18), size=36)
tb(s11, P, Inches(2.2), Inches(10), Inches(0.38),
   "WhatsApp lead AI first -- the highest daily revenue impact. "
   "Nurture and social follow. CRM last, with everything feeding into it.",
   SANS, 13, SUB)
hline(s11, P, Inches(2.72), SW - 2 * P, BORDER)

phases = [
    ("PHASE 01  .  Week 1",
     "WhatsApp Lead AI + Property Chatbot",
     "WhatsApp qualification AI live, property chatbot deployed on the website, "
     "viewing booking integrated, international buyer routing configured.",
     ["WhatsApp AI number setup", "Buyer qualification flow",
      "Website chatbot deployment", "Virtual tour and floor plan delivery",
      "Viewing booking integration", "24/7 response coverage live"]),
    ("PHASE 02  .  Weeks 2-3",
     "Investor Nurture + Social Content AI",
     "12-week nurture sequences for all active leads, construction update automation, "
     "Accra market data delivery, LinkedIn and Instagram content calendar.",
     ["12-week nurture sequence build", "Construction milestone triggers",
      "Market update automation", "LinkedIn content calendar",
      "Instagram caption generator", "Reactivation sequences"]),
    ("PHASE 03  .  Weeks 4-5",
     "CRM Pipeline + Reporting",
     "Full CRM setup with automated deal stage tracking, follow-up reminders, "
     "communication history, and weekly pipeline reporting for Laud and Maxwell.",
     ["CRM pipeline configuration", "Automated stage tracking",
      "Follow-up reminder sequences", "Communication history logging",
      "Weekly pipeline report", "Referral trigger at post-sale"]),
]

ph_w   = Inches(3.8)
ph_gap = Inches(0.17)
px = P

for phase_label, phase_name, phase_desc, deliverables in phases:
    card(s11, px, Inches(3.0), ph_w, Inches(3.75), rounded=True)
    tb(s11, px + Inches(0.25), Inches(3.15), ph_w - Inches(0.4), Inches(0.28),
       phase_label, SANS, 9, VIOLET, bold=True)
    tb(s11, px + Inches(0.25), Inches(3.5), ph_w - Inches(0.4), Inches(0.42),
       phase_name, SERIF, 15, WHITE)
    tb(s11, px + Inches(0.25), Inches(4.0), ph_w - Inches(0.4), Inches(0.52),
       phase_desc, SANS, 11, BODY)
    t = Inches(4.62)
    for d in deliverables[:4]:
        tb(s11, px + Inches(0.25), t, ph_w - Inches(0.4), Inches(0.26),
           f">  {d}", SANS, 10, MUTED)
        t += Inches(0.3)
    px += ph_w + ph_gap


# =============================================================================
# SLIDE 12 -- CLOSING
# =============================================================================
s12 = new_slide(prs)

acc = s12.shapes.add_shape(1, Inches(0), Inches(0), SW, Pt(3))
acc.fill.solid()
acc.fill.fore_color.rgb = VIOLET
acc.line.fill.background()

eyebrow(s12, "WHO IS BEHIND THIS", P, Inches(0.85))
headline(s12, "GenosAI Tech", P, Inches(1.18), size=40)
tb(s12, P, Inches(2.22), Inches(6.5), Inches(0.66),
   "GenosAI builds AI automation for real estate developers targeting diaspora and "
   "international buyers. We build WhatsApp qualification AI, investor nurture systems, "
   "social content automation, property chatbots, and CRM pipelines -- "
   "production-grade from week one, built to close the gap between overseas interest and signed deals.",
   SANS, 13, SUB)

hline(s12, P, Inches(3.02), Inches(6), BORDER)

tb(s12, P, Inches(3.2), Inches(4), Inches(0.28),
   "WHAT WE BUILD", SANS, 9, MUTED, bold=True)
t = Inches(3.58)
for item, desc in [
    ("WhatsApp Lead AI",         "24/7 buyer qualification -- any time zone, any hour."),
    ("Investor Nurture Systems", "12-week automated sequences that close on the 4th contact."),
    ("Social Content AI",        "LinkedIn and Instagram on autopilot -- audience built passively."),
    ("Property Chatbot",         "24/7 inquiry, floor plans, and viewing bookings on the website."),
    ("CRM Pipeline Automation",  "Every buyer tracked from first message to signed contract."),
]:
    t = bullet(s12, "o", item, desc, P, t, Inches(5.5), 12, 11)

card(s12, Inches(7.5), Inches(1.1), Inches(5.05), Inches(5.75), rounded=True)

tb(s12, Inches(7.75), Inches(1.35), Inches(4.6), Inches(0.35),
   "CLOSING NOTE", SANS, 9, MUTED, bold=True)
tb(s12, Inches(7.6), Inches(1.72), Inches(0.4), Inches(0.55),
   '"', SERIF, 28, VIOLET)
tb(s12, Inches(7.95), Inches(1.78), Inches(4.3), Inches(1.05),
   "Regnum Properties is built on a globally ambitious thesis -- "
   "luxury Accra real estate for the diaspora. "
   "The automation layer ensures that thesis converts into signed deals, "
   "not just inbound interest.",
   SERIF, 14, WHITE, italic=True)

hline(s12, Inches(7.75), Inches(2.92), Inches(4.45), BORDER)

tb(s12, Inches(7.75), Inches(3.12), Inches(4.5), Inches(0.28),
   "NEXT STEPS", SANS, 9, VIOLET, bold=True)
t = Inches(3.5)
for step, desc in [
    ("01", "20-minute call to confirm scope and integration points."),
    ("02", "Phase 1 kick-off -- WhatsApp AI live within 7 days."),
    ("03", "First overseas buyer qualified automatically -- visible from day one."),
]:
    tb(s12, Inches(7.75), t, Inches(0.4), Inches(0.3), step, SERIF, 14, VIOLET)
    tb(s12, Inches(8.2), t + Inches(0.02), Inches(4.15), Inches(0.5),
       desc, SANS, 12, BODY)
    t += Inches(0.62)

hline(s12, Inches(7.75), Inches(5.45), Inches(4.45), BORDER)

tb(s12, Inches(7.75), Inches(5.65), Inches(4.5), Inches(0.28),
   "Rohan Malik  .  CEO, Genos AI", SANS, 12, WHITE, bold=True)
tb(s12, Inches(7.75), Inches(6.0), Inches(4.5), Inches(0.28),
   "hello@genosai.tech  .  +91 638-714-2699", SANS, 11, BODY)
tb(s12, Inches(7.75), Inches(6.32), Inches(4.5), Inches(0.28),
   "www.genosai.tech", SANS, 11, MUTED)

tb(s12, P, SH - Inches(0.42), Inches(8), Inches(0.28),
   "Confidential  .  Regnum Properties  .  2026", SANS, 8, DIM)
tb(s12, Inches(9.5), SH - Inches(0.42), Inches(3.08), Inches(0.28),
   f"{TOTAL:02d} / {TOTAL}", SANS, 8, DIM, align=PP_ALIGN.RIGHT)

tb(s12, P, SH - Inches(0.65), Inches(3), Inches(0.4),
   "GenosAI", SERIF, 20, WHITE)


# -- Save ----------------------------------------------------------------------
prs.save(OUT)
print(f"Saved: {OUT}  ({TOTAL} slides)")
