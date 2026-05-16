"""
Build SwisszacoEstates_Audit_GenosAI.pptx from scratch.
Design: docs/BRAND_ASSETS.md — #0A0A0F canvas, Instrument Serif headlines, Satoshi body.
12 slides: cover, audit, vision, systems overview, 5 module slides,
lead flow, implementation, closing.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).parent.parent.parent
OUT  = ROOT / "output" / "decks" / "SwisszacoEstates_Audit_GenosAI.pptx"

# ── Brand colors (BRAND_ASSETS.md §3) ───────────────────────────────────────
BG     = RGBColor(0x0A, 0x0A, 0x0F)   # #0A0A0F canvas
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)   # headlines
BODY   = RGBColor(0x9D, 0x9D, 0x9F)   # white/60
SUB    = RGBColor(0x86, 0x86, 0x8A)   # white/50
MUTED  = RGBColor(0x6B, 0x6B, 0x6E)   # white/40
DIM    = RGBColor(0x44, 0x44, 0x47)   # white/25
VIOLET = RGBColor(0xA7, 0x8B, 0xFA)   # eyebrow #A78BFA
CARD   = RGBColor(0x11, 0x11, 0x18)   # card surface #111118
BORDER = RGBColor(0x22, 0x22, 0x2A)   # card border
RED    = RGBColor(0xF8, 0x71, 0x71)   # gap/issue highlights

# ── Fonts (BRAND_ASSETS.md §4) ───────────────────────────────────────────────
SERIF = "Instrument Serif"
SANS  = "Satoshi"

# ── Slide geometry ───────────────────────────────────────────────────────────
SW    = Inches(13.33)
SH    = Inches(7.5)
P     = Inches(0.75)
TOTAL = 12


# ── Core helpers ─────────────────────────────────────────────────────────────

def new_slide(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def tb(slide, l, t, w, h, text, font=SANS, size=14, color=BODY,
       bold=False, italic=False, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text        = text
    r.font.name   = font
    r.font.size   = Pt(size)
    r.font.color.rgb = color
    r.font.bold   = bold
    r.font.italic = italic
    return box


def tb_lines(slide, l, t, w, h, lines, font=SANS, size=14, color=BODY,
             bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name  = font
        r.font.size  = Pt(size)
        r.font.color.rgb = color
        r.font.bold  = bold
    return box


def card(slide, l, t, w, h, rounded=True):
    shape_type = 5 if rounded else 1
    s = slide.shapes.add_shape(shape_type, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = CARD
    s.line.color.rgb = BORDER
    s.line.width = Pt(0.75)
    return s


def hline(slide, l, t, w, color=BORDER):
    s = slide.shapes.add_shape(1, l, t, w, Pt(0.75))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def eyebrow(slide, text, l, t, w=Inches(8)):
    return tb(slide, l, t, w, Inches(0.3),
              text.upper(), SANS, 10, VIOLET, bold=True)


def headline(slide, text, l, t, w=Inches(10), size=42):
    return tb(slide, l, t, w, Inches(1.6), text, SERIF, size, WHITE)


def slide_header(slide, section, page):
    tb(slide, P, Inches(0.22), Inches(6), Inches(0.28),
       "GENOSAI  x  SWISSZACO ESTATES", SANS, 9, MUTED, bold=True)
    tb(slide, Inches(7.5), Inches(0.22), Inches(5.08), Inches(0.28),
       section, SANS, 9, MUTED, bold=True, align=PP_ALIGN.RIGHT)
    hline(slide, P, Inches(0.58), SW - 2 * P)
    tb(slide, P, SH - Inches(0.42), Inches(8), Inches(0.28),
       "AI Automation & Growth Proposal  ·  Swisszaco Estates  ·  2026",
       SANS, 8, DIM)
    tb(slide, Inches(9.5), SH - Inches(0.42), Inches(3.08), Inches(0.28),
       f"{page:02d} / {TOTAL}", SANS, 8, DIM, align=PP_ALIGN.RIGHT)


def bullet(slide, dot_char, title, desc, l, t,
           col_w=Inches(5.0), title_size=13, desc_size=12):
    tb(slide, l, t, Inches(0.25), Inches(0.26),
       dot_char, SANS, 13, VIOLET, bold=True)
    tb(slide, l + Inches(0.3), t, col_w - Inches(0.3), Inches(0.26),
       title, SANS, title_size, WHITE, bold=True)
    tb(slide, l + Inches(0.3), t + Inches(0.27), col_w - Inches(0.3), Inches(0.38),
       desc, SANS, desc_size, BODY)
    return t + Inches(0.72)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH

s1 = new_slide(prs)

# Top violet accent bar
accent = s1.shapes.add_shape(1, Inches(0), Inches(0), SW, Pt(3))
accent.fill.solid()
accent.fill.fore_color.rgb = VIOLET
accent.line.fill.background()

# Left — company identity
tb(s1, P, Inches(1.0), Inches(6), Inches(0.32),
   "PREPARED FOR", SANS, 10, MUTED, bold=True)
hline(s1, P, Inches(1.42), Inches(6), BORDER)

tb(s1, P, Inches(1.6), Inches(8), Inches(1.5),
   "SWISSZACO\nESTATES", SERIF, 58, WHITE)

tb(s1, P, Inches(3.3), Inches(6), Inches(0.32),
   "Real Estate Development  ·  Nigeria", SANS, 13, BODY)
tb(s1, P, Inches(3.72), Inches(6), Inches(0.28),
   "Construction  ·  Architecture  ·  Urban Development  ·  Estate Development",
   SANS, 11, MUTED)

hline(s1, P, Inches(4.15), Inches(6), BORDER)

tb(s1, P, Inches(4.38), Inches(7), Inches(0.28),
   "AI AUTOMATION  &  GROWTH PROPOSAL", SANS, 10, VIOLET, bold=True)
tb_lines(s1, P, Inches(4.78), Inches(6.5), Inches(1.1),
         ["An intelligent automation layer for Swisszaco Estates —",
          "built to respond to property inquiries 24/7, qualify buyers",
          "automatically, and scale agent productivity without",
          "adding headcount."],
         SANS, 13, SUB)

# Right — prepared by card
card(s1, Inches(9.2), Inches(1.2), Inches(3.38), Inches(2.5), rounded=True)
tb(s1, Inches(9.45), Inches(1.5), Inches(2.9), Inches(0.28),
   "PREPARED BY", SANS, 9, MUTED, bold=True)
tb(s1, Inches(9.45), Inches(1.88), Inches(2.9), Inches(0.38),
   "GenosAI Tech", SERIF, 20, WHITE)
tb(s1, Inches(9.45), Inches(2.32), Inches(2.9), Inches(0.28),
   "Rohan Malik  ·  Founder & CEO", SANS, 11, BODY)
tb(s1, Inches(9.45), Inches(2.65), Inches(2.9), Inches(0.26),
   "hello@genosai.tech  ·  www.genosai.tech", SANS, 10, MUTED)

# GenosAI wordmark bottom
tb(s1, P, SH - Inches(0.65), Inches(4), Inches(0.45),
   "GenosAI", SERIF, 22, WHITE)
tb(s1, P + Inches(1.7), SH - Inches(0.62), Inches(3), Inches(0.3),
   "v1.0  ·  2026  ·  Confidential", SANS, 9, DIM)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AUDIT: CURRENT DIGITAL PRESENCE
# ══════════════════════════════════════════════════════════════════════════════
s2 = new_slide(prs)
slide_header(s2, "01  ·  AUDIT", 2)

eyebrow(s2, "WHERE THE BRAND STANDS TODAY", P, Inches(0.85))
headline(s2, "Current Digital Presence", P, Inches(1.18), size=38)
tb(s2, P, Inches(2.2), Inches(10), Inches(0.42),
   "Swisszaco Estates has 13 years in the Nigerian market and a multi-national holding structure. "
   "The operational credibility is real. The gap is the digital front door: property buyers in Nigeria "
   "expect a WhatsApp reply within the hour — a contact form delivers it in 24-48.",
   SANS, 13, SUB)

hline(s2, P, Inches(2.85), SW - 2 * P, BORDER)

CL = P
CR = Inches(7.3)
CW = Inches(5.5)

# Left — strengths
card(s2, CL, Inches(3.05), CW, Inches(3.6), rounded=True)
tb(s2, CL + Inches(0.3), Inches(3.2), CW - Inches(0.4), Inches(0.28),
   "WHAT'S ALREADY STRONG", SANS, 10, VIOLET, bold=True)
t = Inches(3.6)
for title, desc in [
    ("13 years in market",
     "CAC-registered March 2012 — proven track record in Nigerian real estate."),
    ("Multi-national structure",
     "Associate of Ogorstar Conglomerates Inc — credibility for large-scale deals."),
    ("Full-stack service offering",
     "Construction, architecture, renovation, urban development, surveying — one firm."),
    ("Active agent network",
     "Agent recruitment program live — scalable distribution model already in motion."),
]:
    t = bullet(s2, "·", title, desc, CL + Inches(0.3), t, CW - Inches(0.45), 12, 11)

# Right — gaps
card(s2, CR, Inches(3.05), CW, Inches(3.6), rounded=True)
tb(s2, CR + Inches(0.3), Inches(3.2), CW - Inches(0.4), Inches(0.28),
   "WHERE LEADS ARE LOST", SANS, 10, RED, bold=True)
t = Inches(3.6)
for title, desc in [
    ("No WhatsApp integration",
     "Nigerian buyers expect instant WhatsApp responses — site offers a contact form."),
    ("No live chat or chatbot",
     "After-hours and weekend inquiries go unanswered."),
    ("No viewing / site visit booking",
     "No way to book a property visit — requires manual coordination."),
    ("No email marketing or nurture",
     "Leads who visit the site and don't convert are never followed up automatically."),
]:
    t = bullet(s2, "·", title, desc, CR + Inches(0.3), t, CW - Inches(0.45), 12, 11)

# Audit score
card(s2, Inches(0.75), Inches(6.75), Inches(11.83), Inches(0.45), rounded=True)
tb(s2, Inches(1.1), Inches(6.82), Inches(4), Inches(0.28),
   "WEBSITE AUDIT SCORE: 3 / 10  ·  PRIORITY: HIGH", SANS, 11, WHITE, bold=True)
tb(s2, Inches(6.5), Inches(6.82), Inches(5.5), Inches(0.28),
   "Mobile: Yes  ·  Analytics: Yes  ·  Remarketing: Yes  ·  Booking: No  ·  Chat: No  ·  Email Mkt: No  ·  CDN: No",
   SANS, 10, MUTED, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — VISION: THE OPPORTUNITY
# ══════════════════════════════════════════════════════════════════════════════
s3 = new_slide(prs)
slide_header(s3, "02  ·  VISION", 3)

eyebrow(s3, "THE OPPORTUNITY", P, Inches(0.85))
headline(s3, "An AI Automation Layer\nfor Swisszaco Estates", P, Inches(1.18), size=38)

# Quote block
card(s3, P, Inches(2.9), Inches(7.8), Inches(1.2), rounded=True)
tb(s3, P + Inches(0.4), Inches(3.02), Inches(0.5), Inches(0.45),
   "“", SERIF, 32, VIOLET)
tb(s3, P + Inches(0.78), Inches(3.06), Inches(7.0), Inches(0.88),
   "We don't build chatbots. We build the automated infrastructure that ensures "
   "Swisszaco is always the first developer to respond — and the best qualified to close.",
   SERIF, 15, WHITE, italic=True)

# 4 pillars
col_w = Inches(2.85)
cols = [P, P + col_w + Inches(0.14), P + 2*(col_w + Inches(0.14)), P + 3*(col_w + Inches(0.14))]
pillars = [
    ("01", "24/7 inquiry\nresponse",
     "Every WhatsApp message answered instantly — day or night."),
    ("02", "Automatic\nqualification",
     "Budget, property type, location, timeline captured before a human intervenes."),
    ("03", "Agent network\nautomation",
     "Agent recruitment and onboarding running without manual handling."),
    ("04", "Buyer\nre-engagement",
     "Cold leads automatically followed up with matching new listings."),
]
for (num, title, desc), x in zip(pillars, cols):
    card(s3, x, Inches(4.35), col_w, Inches(2.7), rounded=True)
    tb(s3, x + Inches(0.25), Inches(4.55), Inches(0.6), Inches(0.42),
       num, SERIF, 22, VIOLET)
    tb(s3, x + Inches(0.25), Inches(5.05), col_w - Inches(0.4), Inches(0.65),
       title, SERIF, 16, WHITE)
    tb(s3, x + Inches(0.25), Inches(5.78), col_w - Inches(0.4), Inches(0.88),
       desc, SANS, 11, BODY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SYSTEMS OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
s4 = new_slide(prs)
slide_header(s4, "03  ·  AUTOMATION SYSTEMS", 4)

eyebrow(s4, "WHAT THE AUTOMATION LAYER BRINGS", P, Inches(0.85))
headline(s4, "Five AI Systems for Swisszaco Estates", P, Inches(1.18), size=36)

systems = [
    ("A", "WhatsApp AI\nBuyer Agent",    "24/7 property inquiry response."),
    ("B", "Viewing Booking\nFlow",        "Schedule site visits automatically."),
    ("C", "Agent Network\nAutomation",    "Recruit and onboard agents hands-free."),
    ("D", "Lead Nurture\nSystem",         "Re-engage cold prospects with new listings."),
    ("E", "Review &\nTestimonial AI",     "Automate post-sale review collection."),
]
sw_each = Inches(2.25)
gap     = Inches(0.17)
sx      = P

for letter, name, tagline in systems:
    card(s4, sx, Inches(2.6), sw_each, Inches(4.1), rounded=True)
    badge = s4.shapes.add_shape(1, sx + Inches(0.25), Inches(2.85), Inches(0.55), Inches(0.55))
    badge.fill.solid()
    badge.fill.fore_color.rgb = VIOLET
    badge.line.fill.background()
    tb(s4, sx + Inches(0.25), Inches(2.87), Inches(0.55), Inches(0.48),
       letter, SANS, 16, BG, bold=True, align=PP_ALIGN.CENTER)
    tb(s4, sx + Inches(0.2), Inches(3.55), sw_each - Inches(0.4), Inches(0.9),
       name, SERIF, 15, WHITE)
    tb(s4, sx + Inches(0.2), Inches(4.5), sw_each - Inches(0.4), Inches(0.55),
       tagline, SANS, 12, BODY)
    sx += sw_each + gap


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — MODULE A: WHATSAPP AI BUYER AGENT
# ══════════════════════════════════════════════════════════════════════════════
s5 = new_slide(prs)
slide_header(s5, "03A  ·  WHATSAPP AI AGENT", 5)

eyebrow(s5, "MODULE A", P, Inches(0.85))
headline(s5, "WhatsApp AI Buyer Agent", P, Inches(1.18), size=36)
tb(s5, P, Inches(2.2), Inches(6.8), Inches(0.42),
   "Responds in under 60 seconds. Qualifies the buyer before an agent is involved. "
   "Books the site visit automatically. Runs 24/7 including weekends.",
   SANS, 14, SUB)

hline(s5, P, Inches(2.78), SW - 2 * P, BORDER)

t = Inches(3.05)
for title, desc in [
    ("Instant WhatsApp response",
     "No inquiry goes unanswered — construction, estate development, or renovation enquiry."),
    ("Buyer pre-qualification",
     "Property type, location, budget, and timeline captured before agent contact."),
    ("Site visit booking",
     "AI books the viewing slot directly — no manual back-and-forth required."),
    ("Lead routing by service",
     "Construction enquiry goes to sales, renovation to the renovation team."),
    ("Diaspora & international leads",
     "Nigerian diaspora (UK, US, Canada) handled with timezone-aware responses."),
    ("Multilingual ready",
     "English and Pidgin — meet clients where they communicate."),
]:
    t = bullet(s5, "·", title, desc, P, t, Inches(5.6), 13, 12)

# Right stat card
card(s5, Inches(7.6), Inches(3.05), Inches(4.9), Inches(3.65), rounded=True)
tb(s5, Inches(7.9), Inches(3.25), Inches(4.3), Inches(0.28),
   "PROPERTY INQUIRY  ·  RESPONSE FLOW", SANS, 9, MUTED, bold=True)
tb(s5, Inches(7.9), Inches(3.65), Inches(4.3), Inches(0.45),
   "From WhatsApp message", SERIF, 20, WHITE)
tb(s5, Inches(7.9), Inches(4.15), Inches(4.3), Inches(0.45),
   "to confirmed site visit.", SERIF, 20, VIOLET)

for i, (label, value) in enumerate([
    ("IN MARKET", "13 YRS"), ("SERVICES", "6+"), ("AGENT NETWORK", "ACTIVE")
]):
    sx = Inches(7.9) + i * Inches(1.55)
    tb(s5, sx, Inches(5.15), Inches(1.4), Inches(0.28), label, SANS, 8, MUTED, bold=True)
    tb(s5, sx, Inches(5.45), Inches(1.4), Inches(0.42), value, SERIF, 18, WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — MODULE B & C: VIEWING BOOKING + AGENT NETWORK
# ══════════════════════════════════════════════════════════════════════════════
s6 = new_slide(prs)
slide_header(s6, "03B / 03C  ·  BOOKING & AGENT NETWORK", 6)

# Left — Viewing Booking
eyebrow(s6, "MODULE B", P, Inches(0.85))
headline(s6, "Site Visit Booking Flow", P, Inches(1.18), size=28)
tb(s6, P, Inches(2.08), Inches(5.6), Inches(0.38),
   "Buyers book property viewings without anyone picking up the phone.", SANS, 13, SUB)

t = Inches(2.65)
for title, desc in [
    ("WhatsApp-native scheduling",
     "Buyer selects date and time — confirmation sent automatically."),
    ("Multi-service routing",
     "Construction site visits routed separately from renovation viewings."),
    ("Reminder sequence",
     "Automated reminders at 24h and 2h before — reduces no-shows."),
    ("Reschedule automation",
     "No-shows offered a new slot automatically — agent not involved."),
    ("Post-visit follow-up",
     "AI sends a follow-up message asking for intent within 24h."),
]:
    t = bullet(s6, "·", title, desc, P, t, Inches(5.6), 12, 11)

# Divider
hline(s6, Inches(6.67), Inches(0.75), Pt(1), BORDER)

# Right — Agent Network Automation
CR = Inches(6.9)
eyebrow(s6, "MODULE C", CR, Inches(0.85))
headline(s6, "Agent Network Automation", CR, Inches(1.18), size=28)
tb(s6, CR, Inches(2.08), Inches(5.7), Inches(0.38),
   "Agent recruitment and onboarding running automatically — no manual reading of every enquiry.",
   SANS, 13, SUB)

t = Inches(2.65)
for title, desc in [
    ("Automated agent intake",
     "Every applicant enters an AI flow — no manual screening needed."),
    ("Qualifying questions",
     "AI asks: experience, property types, areas covered, current agency."),
    ("WhatsApp onboarding",
     "Accepted agents get onboarding materials delivered via WhatsApp sequence."),
    ("Drip for undecided applicants",
     "3-touch nurture sequence for candidates who express interest but don't commit."),
    ("Recruitment funnel dashboard",
     "Applicants tracked from first contact to active agent status."),
]:
    t = bullet(s6, "·", title, desc, CR, t, Inches(5.7), 12, 11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — MODULE D & E: NURTURE + REVIEWS
# ══════════════════════════════════════════════════════════════════════════════
s7 = new_slide(prs)
slide_header(s7, "03D / 03E  ·  NURTURE & REVIEWS", 7)

# Left — Lead Nurture
eyebrow(s7, "MODULE D", P, Inches(0.85))
headline(s7, "Lead Nurture System", P, Inches(1.18), size=28)
tb(s7, P, Inches(2.08), Inches(5.6), Inches(0.38),
   "Most enquiries don't convert on first contact. Automated re-engagement recovers those leads.",
   SANS, 13, SUB)

t = Inches(2.65)
for title, desc in [
    ("48h re-engagement",
     "Non-converting buyers receive a WhatsApp follow-up with a new matching property."),
    ("New listing alerts",
     "Buyers who showed interest are notified when a new matching listing drops."),
    ("Price change notifications",
     "If a property they enquired about changes, they're told automatically."),
    ("Long-term investor nurture",
     "Estate development enquiries nurtured over weeks with project updates."),
    ("Diaspora investor sequences",
     "UK/US/Canada leads nurtured with investment case content via WhatsApp."),
]:
    t = bullet(s7, "·", title, desc, P, t, Inches(5.6), 12, 11)

# Right — Review Collection
CR = Inches(6.9)
eyebrow(s7, "MODULE E", CR, Inches(0.85))
headline(s7, "Review & Testimonial AI", CR, Inches(1.18), size=28)
tb(s7, CR, Inches(2.08), Inches(5.7), Inches(0.38),
   "13 years of completed projects. Almost none of them captured as public reviews — until now.",
   SANS, 13, SUB)

t = Inches(2.65)
for title, desc in [
    ("Post-handover review request",
     "Automated WhatsApp sent at project completion asking for a Google review."),
    ("Testimonial capture",
     "Buyers who respond positively are prompted to record a short video testimonial."),
    ("Negative feedback routing",
     "Dissatisfied clients flagged internally — not published — before escalation."),
    ("Agent-specific review collection",
     "Reviews collected per agent to build individual credibility profiles."),
    ("Review management dashboard",
     "All reviews tracked across platforms — Google, Facebook, Jora."),
]:
    t = bullet(s7, "·", title, desc, CR, t, Inches(5.7), 12, 11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — LEAD FLOW: INQUIRY → SITE VISIT → DEAL
# ══════════════════════════════════════════════════════════════════════════════
s8 = new_slide(prs)
slide_header(s8, "04  ·  LEAD FLOW", 8)

eyebrow(s8, "INQUIRY  ·  SITE VISIT  ·  DEAL", P, Inches(0.85))
headline(s8, "Automated Buyer Journey", P, Inches(1.18), size=36)
tb(s8, P, Inches(2.2), Inches(10.5), Inches(0.38),
   "The same flow handles every enquiry — construction, estate lots, or renovation — while each buyer feels like the only one.",
   SANS, 13, SUB)

hline(s8, P, Inches(2.75), SW - 2 * P, BORDER)

steps = [
    ("01", "Inquiry\nArrives",        "WhatsApp, website\nform, or referral"),
    ("02", "AI Responds\n< 60 sec",   "WhatsApp agent\nengages instantly"),
    ("03", "Qualification",           "Property type,\nbudget, timeline"),
    ("04", "Service\nRouting",        "Sales, renovation,\nor development team"),
    ("05", "Visit\nBooked",           "Calendar slot confirmed\nautomatically"),
    ("06", "Reminders\nSent",         "24h and 2h before\nthe site visit"),
    ("07", "Nurture\nActivated",      "If cold, auto\nre-engagement starts"),
    ("08", "Review\nRequested",       "Post-handover\nGoogle review ask"),
]

step_w   = Inches(1.4)
step_gap = Inches(0.12)
sx       = P

for num, title, detail in steps:
    card(s8, sx, Inches(3.1), step_w, Inches(3.5), rounded=True)
    tb(s8, sx + Inches(0.12), Inches(3.28), step_w - Inches(0.2), Inches(0.38),
       num, SERIF, 20, VIOLET)
    tb(s8, sx + Inches(0.12), Inches(3.75), step_w - Inches(0.2), Inches(0.65),
       title, SERIF, 13, WHITE)
    tb(s8, sx + Inches(0.12), Inches(4.5), step_w - Inches(0.2), Inches(0.75),
       detail, SANS, 10, BODY)
    sx += step_w + step_gap
    if num != "08":
        tb(s8, sx - step_gap + Inches(0.02), Inches(4.2), step_gap + Inches(0.06), Inches(0.28),
           ">", SANS, 14, MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — BENEFITS
# ══════════════════════════════════════════════════════════════════════════════
s9 = new_slide(prs)
slide_header(s9, "05  ·  BENEFITS", 9)

eyebrow(s9, "WHY THIS WORKS FOR SWISSZACO", P, Inches(0.85))
headline(s9, "What Changes After Launch", P, Inches(1.18), size=36)
hline(s9, P, Inches(2.72), SW - 2 * P, BORDER)

benefits = [
    ("Zero missed inquiries",
     "Every WhatsApp message answered within 60 seconds — day, night, or weekend."),
    ("Faster deal cycles",
     "Pre-qualified buyers arrive at site visits ready to decide — not just browsing."),
    ("Agent network growth",
     "Recruitment and onboarding running without your team reading every application."),
    ("Diaspora reach",
     "Nigerian diaspora investors in the UK, US, and Canada nurtured automatically."),
    ("Review & testimonial flywheel",
     "13 years of completed projects finally captured as public social proof."),
    ("Construction enquiry routing",
     "Each service line (sales, architecture, renovation) gets the right follow-up."),
    ("Long-term investor nurture",
     "Estate development prospects nurtured over weeks with project progress updates."),
    ("Scales with the business",
     "More listings, agents, and services add no manual overhead to your team."),
]

col1, col2 = benefits[:4], benefits[4:]
CL, CR = P, Inches(7.1)
t1 = t2 = Inches(2.95)

for title, desc in col1:
    card(s9, CL, t1, Inches(5.75), Inches(0.82), rounded=True)
    tb(s9, CL + Inches(0.25), t1 + Inches(0.1), Inches(0.35), Inches(0.3),
       "v", SANS, 13, VIOLET, bold=True)
    tb(s9, CL + Inches(0.65), t1 + Inches(0.08), Inches(4.8), Inches(0.3),
       title, SANS, 12, WHITE, bold=True)
    tb(s9, CL + Inches(0.65), t1 + Inches(0.4), Inches(4.8), Inches(0.35),
       desc, SANS, 11, BODY)
    t1 += Inches(0.95)

for title, desc in col2:
    card(s9, CR, t2, Inches(5.75), Inches(0.82), rounded=True)
    tb(s9, CR + Inches(0.25), t2 + Inches(0.1), Inches(0.35), Inches(0.3),
       "v", SANS, 13, VIOLET, bold=True)
    tb(s9, CR + Inches(0.65), t2 + Inches(0.08), Inches(4.8), Inches(0.3),
       title, SANS, 12, WHITE, bold=True)
    tb(s9, CR + Inches(0.65), t2 + Inches(0.4), Inches(4.8), Inches(0.35),
       desc, SANS, 11, BODY)
    t2 += Inches(0.95)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — FUTURE ROADMAP
# ══════════════════════════════════════════════════════════════════════════════
s10 = new_slide(prs)
slide_header(s10, "06  ·  ROADMAP", 10)

eyebrow(s10, "BEYOND THE LAUNCH", P, Inches(0.85))
headline(s10, "Future Scalability", P, Inches(1.18), size=36)
tb(s10, P, Inches(2.2), Inches(10.5), Inches(0.38),
   "The same AI infrastructure scales as Swisszaco expands into new regions, product lines, and service categories.",
   SANS, 13, SUB)
hline(s10, P, Inches(2.72), SW - 2 * P, BORDER)

future = [
    ("Mortgage & Finance Referral AI",
     "Connect buyers to finance partners automatically after qualification."),
    ("Rental & Property Management AI",
     "Tenant screening, lease renewal reminders, and maintenance request routing."),
    ("Project Progress Updates",
     "Automated WhatsApp construction updates to buyers — no manual calls."),
    ("Multi-state Expansion",
     "Same automation deployed across Lagos, Abuja, Port Harcourt, and beyond."),
    ("Off-plan Investor Flow",
     "Capture, qualify, and nurture investors for upcoming estate development launches."),
    ("Market Report Automation",
     "Automated quarterly market updates to past buyers and investors via WhatsApp."),
    ("Voice AI for Inbound Calls",
     "AI voice agent handles inbound calls — no missed call means no missed lead."),
    ("Social Media Lead Capture",
     "Instagram and Facebook inquiry automation — leads from DMs routed to WhatsApp flow."),
]

col1, col2 = future[:4], future[4:]
CL, CR = P, Inches(7.1)
t1 = t2 = Inches(2.95)

for title, desc in col1:
    t1 = bullet(s10, "·", title, desc, CL, t1, Inches(5.75))
for title, desc in col2:
    t2 = bullet(s10, "·", title, desc, CR, t2, Inches(5.75))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — IMPLEMENTATION
# ══════════════════════════════════════════════════════════════════════════════
s11 = new_slide(prs)
slide_header(s11, "07  ·  IMPLEMENTATION", 11)

eyebrow(s11, "BUILT IN THREE CLEAN PHASES", P, Inches(0.85))
headline(s11, "Implementation Phases", P, Inches(1.18), size=36)
tb(s11, P, Inches(2.2), Inches(10.5), Inches(0.38),
   "Each phase ships working automation — no months of silence before anything goes live.",
   SANS, 13, SUB)
hline(s11, P, Inches(2.72), SW - 2 * P, BORDER)

phases = [
    ("PHASE 01  ·  Week 1",
     "WhatsApp AI Buyer Agent",
     "WhatsApp AI trained on Swisszaco services, pricing, and locations. "
     "Property inquiry intake, qualification, and site visit booking live.",
     ["AI training on services & locations", "WhatsApp inquiry intake flow",
      "Site visit booking integration", "Lead routing by service type",
      "Diaspora/international lead routing", "Go-live and testing"]),
    ("PHASE 02  ·  Weeks 2-3",
     "Viewing Flow & Agent Network",
     "Site visit booking automation and agent recruitment intake deployed. "
     "No-show reschedule and post-visit follow-up active.",
     ["Viewing reminder sequence", "No-show reschedule automation",
      "Post-visit follow-up flow", "Agent recruitment AI intake",
      "WhatsApp onboarding sequence", "Agent nurture drip tuning"]),
    ("PHASE 03  ·  Weeks 4-5",
     "Nurture, Reviews & Roadmap",
     "Lead nurture sequences, review collection, and diaspora investor flow go live.",
     ["Buyer re-engagement sequences", "New listing match alerts",
      "Post-handover review requests", "Diaspora investor nurture flow",
      "Review management dashboard", "Training & team handover"]),
]

ph_w  = Inches(3.8)
ph_gap = Inches(0.17)
px    = P

for phase_label, phase_name, phase_desc, deliverables in phases:
    card(s11, px, Inches(3.0), ph_w, Inches(3.75), rounded=True)
    tb(s11, px + Inches(0.25), Inches(3.15), ph_w - Inches(0.4), Inches(0.28),
       phase_label, SANS, 9, VIOLET, bold=True)
    tb(s11, px + Inches(0.25), Inches(3.5), ph_w - Inches(0.4), Inches(0.42),
       phase_name, SERIF, 15, WHITE)
    tb(s11, px + Inches(0.25), Inches(4.0), ph_w - Inches(0.4), Inches(0.52),
       phase_desc, SANS, 11, BODY)
    t = Inches(4.62)
    for d in deliverables[:4]:
        tb(s11, px + Inches(0.25), t, ph_w - Inches(0.4), Inches(0.26),
           f">  {d}", SANS, 10, MUTED)
        t += Inches(0.3)
    px += ph_w + ph_gap


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — CLOSING
# ══════════════════════════════════════════════════════════════════════════════
s12 = new_slide(prs)

acc = s12.shapes.add_shape(1, Inches(0), Inches(0), SW, Pt(3))
acc.fill.solid()
acc.fill.fore_color.rgb = VIOLET
acc.line.fill.background()

eyebrow(s12, "WHO IS BEHIND THIS", P, Inches(0.85))
headline(s12, "GenosAI Tech", P, Inches(1.18), size=40)
tb(s12, P, Inches(2.22), Inches(6.5), Inches(0.66),
   "GenosAI builds AI automation systems for real estate developers, construction firms, and property agencies. "
   "We specialise in WhatsApp AI, lead qualification flows, agent network automation, and buyer nurture systems "
   "— production-grade, typically live within two weeks.",
   SANS, 13, SUB)

hline(s12, P, Inches(3.05), Inches(6), BORDER)

tb(s12, P, Inches(3.25), Inches(4), Inches(0.28),
   "WHAT WE BUILD", SANS, 9, MUTED, bold=True)
t = Inches(3.62)
for item, desc in [
    ("WhatsApp AI Agents",           "Trained on your properties, services, and pricing."),
    ("Lead Qualification Flows",     "Budget, type, location, timeline captured automatically."),
    ("Agent Network Automation",     "Recruitment intake, onboarding, and nurture — hands-free."),
    ("Buyer Nurture Systems",        "Re-engage cold leads with matching new listings."),
    ("Review & Reputation AI",       "Post-handover testimonial and review collection at scale."),
]:
    t = bullet(s12, "◆", item, desc, P, t, Inches(5.5), 12, 11)

# Right — contact & next steps card
card(s12, Inches(7.6), Inches(1.1), Inches(4.95), Inches(5.75), rounded=True)

tb(s12, Inches(7.85), Inches(1.35), Inches(4.5), Inches(0.35),
   "CLOSING NOTE", SANS, 9, MUTED, bold=True)
tb(s12, Inches(7.72), Inches(1.72), Inches(0.4), Inches(0.55),
   "“", SERIF, 28, VIOLET)
tb(s12, Inches(8.05), Inches(1.78), Inches(4.2), Inches(1.0),
   "Built to ensure Swisszaco Estates is always the first developer a Nigerian buyer hears back from.",
   SERIF, 14, WHITE, italic=True)

hline(s12, Inches(7.85), Inches(2.9), Inches(4.35), BORDER)

tb(s12, Inches(7.85), Inches(3.1), Inches(4.5), Inches(0.28),
   "NEXT STEPS", SANS, 9, VIOLET, bold=True)
t = Inches(3.48)
for step, desc in [
    ("01", "15-minute call to confirm scope and integration points."),
    ("02", "Phase 1 kick-off — WhatsApp AI buyer agent live within 7 days."),
    ("03", "First inquiry routed automatically — results visible from day one."),
]:
    tb(s12, Inches(7.85), t, Inches(0.4), Inches(0.3), step, SERIF, 14, VIOLET)
    tb(s12, Inches(8.3), t + Inches(0.02), Inches(4.05), Inches(0.5),
       desc, SANS, 12, BODY)
    t += Inches(0.62)

hline(s12, Inches(7.85), Inches(5.45), Inches(4.35), BORDER)

tb(s12, Inches(7.85), Inches(5.65), Inches(4.5), Inches(0.28),
   "Rohan Malik  ·  CEO, Genos AI", SANS, 12, WHITE, bold=True)
tb(s12, Inches(7.85), Inches(6.0), Inches(4.5), Inches(0.28),
   "hello@genosai.tech  ·  +91 638-714-2699", SANS, 11, BODY)
tb(s12, Inches(7.85), Inches(6.32), Inches(4.5), Inches(0.28),
   "www.genosai.tech", SANS, 11, MUTED)

tb(s12, P, SH - Inches(0.42), Inches(8), Inches(0.28),
   "Confidential  ·  Swisszaco Estates  ·  2026", SANS, 8, DIM)
tb(s12, Inches(9.5), SH - Inches(0.42), Inches(3.08), Inches(0.28),
   f"{TOTAL:02d} / {TOTAL}", SANS, 8, DIM, align=PP_ALIGN.RIGHT)

tb(s12, P, SH - Inches(0.65), Inches(3), Inches(0.4),
   "GenosAI", SERIF, 20, WHITE)


# ── Save ─────────────────────────────────────────────────────────────────────
OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"Saved: {OUT}  ({TOTAL} slides)")
