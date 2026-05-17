"""
Build EducationOutcomesFund_Audit_GenosAI.pptx from scratch.
Track: Audit (website score 4/10 -- critical operational automation gaps).
Design: docs/BRAND_ASSETS.md -- #0A0A0F canvas, Instrument Serif headlines, Satoshi body.
12 slides: cover, website audit, org overview, vision, systems, 4 module slides,
programme operations flow, implementation, closing.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).parent.parent.parent
OUT  = ROOT / "output" / "decks" / "EducationOutcomesFund_Audit_GenosAI.pptx"

# -- Brand colors --------------------------------------------------------------
BG     = RGBColor(0x0A, 0x0A, 0x0F)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BODY   = RGBColor(0x9D, 0x9D, 0x9F)
SUB    = RGBColor(0x86, 0x86, 0x8A)
MUTED  = RGBColor(0x6B, 0x6B, 0x6E)
DIM    = RGBColor(0x44, 0x44, 0x47)
VIOLET = RGBColor(0xA7, 0x8B, 0xFA)
CARD   = RGBColor(0x11, 0x11, 0x18)
BORDER = RGBColor(0x22, 0x22, 0x2A)
RED    = RGBColor(0xF8, 0x71, 0x71)

SERIF = "Instrument Serif"
SANS  = "Satoshi"

SW    = Inches(13.33)
SH    = Inches(7.5)
P     = Inches(0.75)
TOTAL = 12


# -- Helpers -------------------------------------------------------------------

def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def tb(slide, l, t, w, h, text, font=SANS, size=14, color=BODY,
       bold=False, italic=False, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text = text
    r.font.name   = font
    r.font.size   = Pt(size)
    r.font.color.rgb = color
    r.font.bold   = bold
    r.font.italic = italic
    return box


def tb_lines(slide, l, t, w, h, lines, font=SANS, size=14, color=BODY,
             bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name  = font
        r.font.size  = Pt(size)
        r.font.color.rgb = color
        r.font.bold  = bold
    return box


def card(slide, l, t, w, h, rounded=True):
    s = slide.shapes.add_shape(5 if rounded else 1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = CARD
    s.line.color.rgb = BORDER
    s.line.width = Pt(0.75)
    return s


def hline(slide, l, t, w, color=BORDER):
    s = slide.shapes.add_shape(1, l, t, w, Pt(0.75))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def eyebrow(slide, text, l, t, w=Inches(6)):
    return tb(slide, l, t, w, Inches(0.3),
              text.upper(), SANS, 10, VIOLET, bold=True)


def headline(slide, text, l, t, w=Inches(10), size=42):
    return tb(slide, l, t, w, Inches(1.6), text, SERIF, size, WHITE)


def slide_header(slide, section, page):
    tb(slide, P, Inches(0.22), Inches(6), Inches(0.28),
       "GENOSAI  x  EDUCATION OUTCOMES FUND", SANS, 9, MUTED, bold=True)
    tb(slide, Inches(7.5), Inches(0.22), Inches(5.08), Inches(0.28),
       section, SANS, 9, MUTED, bold=True, align=PP_ALIGN.RIGHT)
    hline(slide, P, Inches(0.58), SW - 2 * P)
    tb(slide, P, SH - Inches(0.42), Inches(8), Inches(0.28),
       "AI Audit & Automation Proposal  .  Education Outcomes Fund  .  2026",
       SANS, 8, DIM)
    tb(slide, Inches(9.5), SH - Inches(0.42), Inches(3.08), Inches(0.28),
       f"{page:02d} / {TOTAL}", SANS, 8, DIM, align=PP_ALIGN.RIGHT)


def bullet(slide, dot_char, title, desc, l, t,
           col_w=Inches(5.0), title_size=13, desc_size=12):
    tb(slide, l, t, Inches(0.25), Inches(0.26),
       dot_char, SANS, 13, VIOLET, bold=True)
    tb(slide, l + Inches(0.3), t, col_w - Inches(0.3), Inches(0.26),
       title, SANS, title_size, WHITE, bold=True)
    tb(slide, l + Inches(0.3), t + Inches(0.27), col_w - Inches(0.3), Inches(0.38),
       desc, SANS, desc_size, BODY)
    return t + Inches(0.72)


# =============================================================================
# SLIDE 1 -- COVER
# =============================================================================
prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH

s1 = new_slide(prs)

accent = s1.shapes.add_shape(1, Inches(0), Inches(0), SW, Pt(3))
accent.fill.solid()
accent.fill.fore_color.rgb = VIOLET
accent.line.fill.background()

tb(s1, P, Inches(1.0), Inches(6), Inches(0.35),
   "PREPARED FOR", SANS, 10, MUTED, bold=True)
hline(s1, P, Inches(1.42), Inches(6.5), BORDER)

tb(s1, P, Inches(1.6), Inches(9.5), Inches(1.6),
   "Education Outcomes Fund", SERIF, 54, WHITE)

tb(s1, P, Inches(3.1), Inches(8), Inches(0.32),
   "Outcomes-Based Education Financing  .  Sierra Leone  .  Nigeria  .  Rwanda  .  South Africa  .  Tunisia",
   SANS, 12, BODY)

tb(s1, P, Inches(3.55), Inches(6), Inches(0.28),
   "London, UK  .  Hosted by UNICEF  .  Founded 2018", SANS, 11, MUTED)

hline(s1, P, Inches(4.0), Inches(6.5), BORDER)

tb(s1, P, Inches(4.2), Inches(7.5), Inches(0.28),
   "AI AUDIT & AUTOMATION PROPOSAL", SANS, 10, VIOLET, bold=True)
tb_lines(s1, P, Inches(4.6), Inches(6.5), Inches(1.4),
         ["An AI automation blueprint for EOF's operational layer --",
          "from partner intake and donor reporting",
          "to knowledge management and country team coordination,",
          "built to scale impact without scaling headcount."],
         SANS, 13, SUB)

card(s1, Inches(9.3), Inches(1.2), Inches(3.28), Inches(2.5), rounded=True)
tb(s1, Inches(9.55), Inches(1.5), Inches(2.8), Inches(0.28),
   "PREPARED BY", SANS, 9, MUTED, bold=True)
tb(s1, Inches(9.55), Inches(1.88), Inches(2.8), Inches(0.32),
   "GenosAI Tech", SERIF, 20, WHITE)
tb(s1, Inches(9.55), Inches(2.3), Inches(2.8), Inches(0.26),
   "Rohan Malik  .  Founder & CEO", SANS, 11, BODY)
tb(s1, Inches(9.55), Inches(2.65), Inches(2.8), Inches(0.26),
   "hello@genosai.tech  .  www.genosai.tech", SANS, 10, MUTED)

tb(s1, P, SH - Inches(0.65), Inches(4), Inches(0.45),
   "GenosAI", SERIF, 22, WHITE)
tb(s1, P + Inches(1.65), SH - Inches(0.62), Inches(3), Inches(0.3),
   "v1.0  .  2026  .  Confidential", SANS, 9, DIM)


# =============================================================================
# SLIDE 2 -- WEBSITE AUDIT
# =============================================================================
s2 = new_slide(prs)
slide_header(s2, "01  .  WEBSITE AUDIT", 2)

eyebrow(s2, "WEBSITE AUDIT  --  educationoutcomesfund.org", P, Inches(0.85))
headline(s2, "Website Score: 4 / 10  --  Critical Gaps", P, Inches(1.18), size=34)
tb(s2, P, Inches(2.2), Inches(9), Inches(0.38),
   "EOF's website presents the mission clearly -- but the operational layer is absent. "
   "Inbound partner proposals, donor enquiries, and implementer applications "
   "all route through a single generic form with no triage, no automation, and no follow-up.",
   SANS, 13, SUB)

badge = s2.shapes.add_shape(5, Inches(11.6), Inches(1.2), Inches(1.1), Inches(1.1))
badge.fill.solid()
badge.fill.fore_color.rgb = RED
badge.line.fill.background()
tb(s2, Inches(11.6), Inches(1.36), Inches(1.1), Inches(0.45),
   "4/10", SERIF, 22, BG, align=PP_ALIGN.CENTER)
tb(s2, Inches(11.6), Inches(1.82), Inches(1.1), Inches(0.28),
   "SCORE", SANS, 8, BG, bold=True, align=PP_ALIGN.CENTER)

hline(s2, P, Inches(2.72), SW - 2 * P, BORDER)

gaps = [
    ("No partner intake automation",
     "Implementer and investor proposals arrive via a generic Google Form. "
     "No triage, no scoring, no routing -- every submission handled manually."),
    ("No meeting booking system",
     "No Calendly or scheduling link for inbound donor, government, or partner calls. "
     "All contact routed through a single form, creating unnecessary back-and-forth."),
    ("No email automation or audience segmentation",
     "Newsletter exists but no visible segmentation. Government partners, donors, "
     "implementers, and researchers each need different content -- none is differentiated."),
    ("No gated content or lead capture",
     "EOF publishes high-value practice notes and technical briefs -- none behind "
     "an email capture form. Every download is anonymous; no nurture pathway exists."),
    ("No searchable knowledge base",
     "Institutional knowledge lives in PDFs and blog posts. "
     "No AI-searchable layer means the same questions are answered repeatedly."),
    ("Built on Wix -- limited automation capability",
     "The site's CMS limits native integration with CRM, automation, or AI tools. "
     "A headless or API-first approach is needed for the operational layer to function."),
]

CL, CR = P, Inches(7.1)
t1, t2 = Inches(2.95), Inches(2.95)

for i, (title, desc) in enumerate(gaps):
    if i < 3:
        card(s2, CL, t1, Inches(5.75), Inches(0.9), rounded=True)
        tb(s2, CL + Inches(0.2), t1 + Inches(0.08), Inches(0.35), Inches(0.3),
           "!", SANS, 13, RED, bold=True)
        tb(s2, CL + Inches(0.6), t1 + Inches(0.07), Inches(4.9), Inches(0.3),
           title, SANS, 12, WHITE, bold=True)
        tb(s2, CL + Inches(0.6), t1 + Inches(0.4), Inches(4.9), Inches(0.42),
           desc, SANS, 10, BODY)
        t1 += Inches(1.02)
    else:
        card(s2, CR, t2, Inches(5.75), Inches(0.9), rounded=True)
        tb(s2, CR + Inches(0.2), t2 + Inches(0.08), Inches(0.35), Inches(0.3),
           "!", SANS, 13, RED, bold=True)
        tb(s2, CR + Inches(0.6), t2 + Inches(0.07), Inches(4.9), Inches(0.3),
           title, SANS, 12, WHITE, bold=True)
        tb(s2, CR + Inches(0.6), t2 + Inches(0.4), Inches(4.9), Inches(0.42),
           desc, SANS, 10, BODY)
        t2 += Inches(1.02)


# =============================================================================
# SLIDE 3 -- EOF AT A GLANCE
# =============================================================================
s3 = new_slide(prs)
slide_header(s3, "02  .  ORGANISATION OVERVIEW", 3)

eyebrow(s3, "THE FUND AT A GLANCE", P, Inches(0.85))
headline(s3, "5x Programme Growth in 4 Years.\n45 People. One Operational Layer Missing.", P, Inches(1.18), size=32)
tb(s3, P, Inches(2.32), Inches(9), Inches(0.38),
   "EOF designs and manages Outcomes Funds -- pooled financing vehicles where payment is tied to "
   "independently verified results. From one active programme in 2022 to five in 2026, "
   "with Namibia in development. The mission scales. The operations need to keep up.",
   SANS, 13, SUB)

hline(s3, P, Inches(2.88), SW - 2 * P, BORDER)

stats = [
    ("STAFF",          "~45",          "London HQ + 6 country teams"),
    ("MOBILISED",      "$130M+",       "across active Outcomes Funds"),
    ("BENEFICIARIES",  "500K+",        "children and youth reached"),
    ("PROGRAMMES",     "5 active",     "Sierra Leone, Nigeria, Rwanda, SA, Tunisia"),
]
sw_each = Inches(2.8)
gap     = Inches(0.19)
sx      = P

for label, value, sub in stats:
    card(s3, sx, Inches(3.1), sw_each, Inches(2.0), rounded=True)
    tb(s3, sx + Inches(0.25), Inches(3.28), sw_each - Inches(0.4), Inches(0.28),
       label, SANS, 9, MUTED, bold=True)
    tb(s3, sx + Inches(0.25), Inches(3.65), sw_each - Inches(0.4), Inches(0.62),
       value, SERIF, 24, WHITE)
    tb(s3, sx + Inches(0.25), Inches(4.38), sw_each - Inches(0.4), Inches(0.45),
       sub, SANS, 11, BODY)
    sx += sw_each + gap

tb(s3, P, Inches(5.4), Inches(11), Inches(0.28),
   "ACTIVE PROGRAMME FOOTPRINT", SANS, 9, MUTED, bold=True)
hline(s3, P, Inches(5.72), SW - 2 * P, BORDER)
tb(s3, P, Inches(5.82), SW - 2 * P, Inches(0.42),
   "Sierra Leone (SLEIC, $18M, 100K+ children)  .  Nigeria LEAF ($25M, 200K children, launched March 2026)  .  "
   "Rwanda (ECD)  .  South Africa  .  Tunisia (youth employment)  .  Namibia (in development)",
   SANS, 12, BODY)


# =============================================================================
# SLIDE 4 -- VISION
# =============================================================================
s4 = new_slide(prs)
slide_header(s4, "03  .  VISION", 4)

eyebrow(s4, "THE OPERATIONAL OPPORTUNITY", P, Inches(0.85))
headline(s4, "Scale Impact Without\nScaling Headcount", P, Inches(1.18), size=38)

card(s4, P, Inches(2.85), Inches(7.5), Inches(1.25), rounded=True)
tb(s4, P + Inches(0.4), Inches(2.98), Inches(7.0), Inches(0.4),
   '"', SERIF, 32, VIOLET)
tb(s4, P + Inches(0.75), Inches(3.02), Inches(6.5), Inches(0.85),
   "Real change happens when we stop paying for promises and start paying for results. "
   "The same principle applies to operations: every manual process EOF runs "
   "is a promise to scale impact that costs more than it should.",
   SERIF, 15, WHITE, italic=True)

col_w = Inches(2.8)
cols  = [P, P + col_w + Inches(0.15), P + 2*(col_w + Inches(0.15)), P + 3*(col_w + Inches(0.15))]
pillars = [
    ("01", "Intake\nAutomation",
     "Partner and implementer proposals qualified and routed without manual review."),
    ("02", "Donor\nReporting",
     "Verified outcomes data auto-compiled into funder-specific report drafts."),
    ("03", "Comms\nAutomation",
     "Stakeholder updates segmented by role and delivered at scale."),
    ("04", "Knowledge\nManagement",
     "Institutional knowledge searchable by any team member in any country."),
]
for (num, title, desc), x in zip(pillars, cols):
    card(s4, x, Inches(4.35), col_w, Inches(2.65), rounded=True)
    tb(s4, x + Inches(0.25), Inches(4.55), Inches(0.6), Inches(0.42),
       num, SERIF, 22, VIOLET)
    tb(s4, x + Inches(0.25), Inches(5.05), col_w - Inches(0.4), Inches(0.65),
       title, SERIF, 16, WHITE)
    tb(s4, x + Inches(0.25), Inches(5.78), col_w - Inches(0.4), Inches(0.88),
       desc, SANS, 11, BODY)


# =============================================================================
# SLIDE 5 -- SYSTEMS OVERVIEW
# =============================================================================
s5 = new_slide(prs)
slide_header(s5, "04  .  AUTOMATION SYSTEMS", 5)

eyebrow(s5, "FIVE AI SYSTEMS FOR EOF", P, Inches(0.85))
headline(s5, "The Automation Blueprint", P, Inches(1.18), size=36)

systems = [
    ("A", "Partner\nIntake AI",        "Smart triage for all inbound proposals."),
    ("B", "Donor\nReporting AI",       "Auto-generate funder-specific reports."),
    ("C", "Stakeholder\nComms AI",     "Segmented updates by audience."),
    ("D", "Knowledge\nManagement AI",  "RAG over institutional knowledge base."),
    ("E", "Country Team\nCoord. AI",   "Async updates + scheduling automation."),
]
sw_each = Inches(2.25)
gap     = Inches(0.17)
sx      = P

for letter, name, tagline in systems:
    card(s5, sx, Inches(2.6), sw_each, Inches(4.1), rounded=True)
    badge = s5.shapes.add_shape(1, sx + Inches(0.25), Inches(2.85), Inches(0.55), Inches(0.55))
    badge.fill.solid()
    badge.fill.fore_color.rgb = VIOLET
    badge.line.fill.background()
    tb(s5, sx + Inches(0.25), Inches(2.87), Inches(0.55), Inches(0.48),
       letter, SANS, 16, BG, bold=True, align=PP_ALIGN.CENTER)
    tb(s5, sx + Inches(0.2), Inches(3.55), sw_each - Inches(0.4), Inches(0.75),
       name, SERIF, 15, WHITE)
    tb(s5, sx + Inches(0.2), Inches(4.38), sw_each - Inches(0.4), Inches(0.55),
       tagline, SANS, 12, BODY)
    sx += sw_each + gap


# =============================================================================
# SLIDE 6 -- MODULE A: PARTNER INTAKE AI
# =============================================================================
s6 = new_slide(prs)
slide_header(s6, "04A  .  PARTNER INTAKE AI", 6)

eyebrow(s6, "MODULE A", P, Inches(0.85))
headline(s6, "Partner & Implementer Intake AI", P, Inches(1.18), size=34)
tb(s6, P, Inches(2.2), Inches(6.5), Inches(0.42),
   "Every inbound proposal -- NGO, government partner, investor, or implementer -- "
   "enters a structured AI intake system that qualifies, scores, and routes "
   "before a programme manager is involved.",
   SANS, 14, SUB)

hline(s6, P, Inches(2.75), SW - 2 * P, BORDER)

t = Inches(3.0)
for title, desc in [
    ("Smart intake routing",
     "Inbound enquiries from the website categorised by type on submission: "
     "implementer proposal, government partnership, donor enquiry, researcher request, media. "
     "Each routed to the appropriate team member automatically."),
    ("Implementer scoring",
     "For NGO and implementer proposals, an AI intake form captures organisation type, "
     "country context, target population, capacity, and track record. "
     "A scoring model ranks fit against EOF's active programme criteria."),
    ("Automated acknowledgement sequences",
     "Every submission receives an immediate, context-aware acknowledgement. "
     "Qualified proposals receive a follow-up timeline; unqualified ones a respectful redirect "
     "to the implementer database -- without programme team involvement."),
    ("Meeting scheduling for qualified leads",
     "Proposals that pass the scoring threshold trigger an automated meeting invite "
     "to Eugenio Donadio or the relevant Programme Lead. "
     "No back-and-forth scheduling required."),
    ("CRM pipeline integration",
     "All intake data flows into a CRM pipeline -- partner stage, geography, programme fit. "
     "Max McCabe and Milena Castellnou see the full inbound picture in one view."),
]:
    t = bullet(s6, ".", title, desc, P, t, Inches(5.5), 13, 12)

card(s6, Inches(7.5), Inches(3.0), Inches(5.0), Inches(3.7), rounded=True)
tb(s6, Inches(7.8), Inches(3.2), Inches(4.5), Inches(0.28),
   "THE CURRENT STATE", SANS, 9, MUTED, bold=True)
tb(s6, Inches(7.8), Inches(3.6), Inches(4.5), Inches(0.55),
   "A Google Form with no triage.", SERIF, 18, WHITE)
tb(s6, Inches(7.8), Inches(4.18), Inches(4.5), Inches(0.55),
   "Every submission handled manually.", SERIF, 16, RED)
tb(s6, Inches(7.8), Inches(4.88), Inches(4.5), Inches(0.72),
   "With five active programmes and a sixth in development, "
   "unqualified proposals absorb programme team time "
   "that belongs on Sierra Leone, Nigeria, and Namibia.", SANS, 11, BODY)


# =============================================================================
# SLIDE 7 -- MODULE B: DONOR REPORTING AUTOMATION
# =============================================================================
s7 = new_slide(prs)
slide_header(s7, "04B  .  DONOR REPORTING AI", 7)

eyebrow(s7, "MODULE B", P, Inches(0.85))
headline(s7, "Donor Reporting Automation", P, Inches(1.18), size=34)
tb(s7, P, Inches(2.2), Inches(6.5), Inches(0.42),
   "EOF has five active programmes, each with multiple funders. "
   "CIFF, UBS Optimus Foundation, and the Japan Ministry of Foreign Affairs "
   "each have different reporting requirements. "
   "AI compiles the data -- the team validates and signs off.",
   SANS, 14, SUB)

hline(s7, P, Inches(2.75), SW - 2 * P, BORDER)

t = Inches(3.0)
for title, desc in [
    ("Verified outcomes data aggregation",
     "Results from independent enumerators and government inspectors aggregated "
     "automatically by programme, cohort, and measurement period. "
     "No manual compilation from raw verification data."),
    ("Funder-specific report drafts",
     "AI generates a report draft formatted to each donor's template: "
     "CIFF narrative format, UBS Optimus financial summary, Japan MoFA results table. "
     "Programme Analysts review and finalise -- not build from scratch."),
    ("Anomaly flagging",
     "When verified outcomes data deviates significantly from projections, "
     "the AI flags it before the report is drafted -- giving the team time to "
     "contextualise the finding rather than explain it reactively."),
    ("Progress milestone alerts",
     "When a programme hits a verified milestone (e.g. 50K of 100K SLEIC children reached), "
     "the relevant funder receives an automated milestone update -- "
     "maintaining donor confidence between formal reporting periods."),
    ("Cross-programme results dashboard",
     "A live dashboard for Amel and Max: verified results across all five programmes, "
     "updated after each enumeration cycle. One view, always current."),
]:
    t = bullet(s7, ".", title, desc, P, t, Inches(5.5), 13, 12)

card(s7, Inches(7.5), Inches(3.0), Inches(5.0), Inches(3.7), rounded=True)
tb(s7, Inches(7.8), Inches(3.2), Inches(4.5), Inches(0.28),
   "THE REPORTING CHALLENGE", SANS, 9, MUTED, bold=True)
tb(s7, Inches(7.8), Inches(3.6), Inches(4.5), Inches(0.55),
   "5 programmes x multiple funders", SERIF, 18, WHITE)
tb(s7, Inches(7.8), Inches(4.18), Inches(4.5), Inches(0.55),
   "= compounding reporting load.", SERIF, 16, RED)
tb(s7, Inches(7.8), Inches(4.88), Inches(4.5), Inches(0.72),
   "Each new programme launch multiplies the donor reporting obligation. "
   "AI automation keeps that load flat "
   "as EOF scales toward its 2030 goal of 15 active partnerships.", SANS, 11, BODY)


# =============================================================================
# SLIDE 8 -- MODULE C: STAKEHOLDER COMMUNICATIONS AI
# =============================================================================
s8 = new_slide(prs)
slide_header(s8, "04C  .  STAKEHOLDER COMMS AI", 8)

eyebrow(s8, "MODULE C", P, Inches(0.85))
headline(s8, "Stakeholder Communications AI", P, Inches(1.18), size=32)
tb(s8, P, Inches(2.2), Inches(6.5), Inches(0.42),
   "19,581 LinkedIn followers and a newsletter database with no visible segmentation. "
   "Government partners, donors, implementers, and researchers need fundamentally "
   "different content. One broadcast to all is wasted on most.",
   SANS, 14, SUB)

hline(s8, P, Inches(2.75), SW - 2 * P, BORDER)

t = Inches(3.0)
for title, desc in [
    ("Audience segmentation",
     "Contacts segmented by role: government partner, institutional donor, "
     "implementer/NGO, researcher/academic, journalist/media. "
     "Each receives content relevant to their context."),
    ("Automated update sequences by geography",
     "When a new LEAF Nigeria result is published, Lagos State officials and "
     "Nigeria-based implementers receive a targeted update. "
     "UK donors receive a funder-perspective version. One event, four audiences."),
    ("Gated content with lead capture",
     "Practice notes, evaluation reports, and technical briefs placed behind "
     "an email capture form. Every download becomes a contactable stakeholder "
     "in the appropriate segment."),
    ("LinkedIn content scheduling",
     "EOF's high-quality programme content published on a consistent schedule "
     "via automated LinkedIn posting -- reducing the manual effort behind "
     "maintaining 19,581 followers with regular updates."),
    ("Conference and event follow-up sequences",
     "After GPE summits, UNESCO conferences, or bilateral meetings, "
     "AI sends tailored follow-up sequences to contacts met at the event "
     "within 24 hours -- maintaining momentum from in-person conversations."),
]:
    t = bullet(s8, ".", title, desc, P, t, Inches(5.5), 13, 12)

card(s8, Inches(7.5), Inches(3.0), Inches(5.0), Inches(3.7), rounded=True)
tb(s8, Inches(7.8), Inches(3.2), Inches(4.5), Inches(0.28),
   "THE COMMS OPPORTUNITY", SANS, 9, MUTED, bold=True)
tb(s8, Inches(7.8), Inches(3.6), Inches(4.5), Inches(0.55),
   "EOF's evidence base is world-class.", SERIF, 17, WHITE)
tb(s8, Inches(7.8), Inches(4.18), Inches(4.5), Inches(0.55),
   "Its distribution is not.", SERIF, 17, RED)
tb(s8, Inches(7.8), Inches(4.88), Inches(4.5), Inches(0.72),
   "SLEIC Year 3 results and the LEAF Nigeria launch represent globally "
   "significant education data. Automated, segmented distribution "
   "ensures that evidence reaches the people who can act on it.", SANS, 11, BODY)


# =============================================================================
# SLIDE 9 -- MODULE D & E: KNOWLEDGE MANAGEMENT + COORDINATION
# =============================================================================
s9 = new_slide(prs)
slide_header(s9, "04D / 04E  .  KNOWLEDGE & COORDINATION", 9)

eyebrow(s9, "MODULE D", P, Inches(0.85))
headline(s9, "Knowledge Management AI", P, Inches(1.18), size=28)
tb(s9, P, Inches(2.05), Inches(5.5), Inches(0.38),
   "EOF's institutional knowledge lives in PDFs and blog posts. "
   "A RAG AI makes it searchable for every team member in every country.",
   SANS, 13, SUB)

t = Inches(2.6)
for title, desc in [
    ("AI-searchable knowledge base",
     "All published practice notes, technical briefs, evaluation reports, and "
     "case studies indexed into a searchable RAG system. "
     "Any team member queries in plain language -- 'what worked in Rwanda ECD' "
     "-- and gets a sourced answer in seconds."),
    ("New joiner onboarding",
     "Programme Analysts joining for Sierra Leone or Namibia oriented via an AI "
     "that answers context-specific questions from the knowledge base, "
     "rather than relying on senior staff to explain the same background."),
    ("Meeting synthesis and capture",
     "Programme team meetings and country calls recorded and auto-summarised. "
     "Key decisions, action items, and learnings captured to the knowledge base "
     "without manual note-taking -- closing the 'Growing in Spirals' loop."),
]:
    t = bullet(s9, ".", title, desc, P, t, Inches(5.5), 12, 11)

hline(s9, Inches(6.67) - Pt(0.5), Inches(0.85), Pt(1), BORDER)

CR = Inches(6.85)
eyebrow(s9, "MODULE E", CR, Inches(0.85))
headline(s9, "Country Team Coordination AI", CR, Inches(1.18), size=28)
tb(s9, CR, Inches(2.05), Inches(5.7), Inches(0.38),
   "Six country teams (Sierra Leone, Nigeria, Rwanda, South Africa, Tunisia, Namibia) "
   "reporting to London. Async coordination at this scale needs automation.",
   SANS, 13, SUB)

t = Inches(2.6)
for title, desc in [
    ("Weekly async status synthesis",
     "Country leads submit structured weekly updates via a simple form or WhatsApp. "
     "AI synthesises into a single executive summary for Milena Castellnou and Max McCabe -- "
     "no Monday coordination call required to understand the portfolio status."),
    ("Scheduling automation",
     "Inbound requests for calls with programme leads managed by an AI scheduling layer. "
     "Government officials in Freetown and donors in London book directly "
     "without going through the London operations team."),
    ("Risk and delay flagging",
     "When a country programme reports a verification delay or government approval lag, "
     "the AI flags it to the relevant Programme Lead and suggests a mitigation note "
     "for the affected donor's reporting timeline."),
    ("Cross-programme learning alerts",
     "When Rwanda reports a finding that is relevant to the Namibia ECD programme "
     "in development, the AI surfaces the connection automatically "
     "for the Namibia team -- institutional knowledge transfer without meetings."),
]:
    t = bullet(s9, ".", title, desc, CR, t, Inches(5.7), 12, 11)


# =============================================================================
# SLIDE 10 -- PROGRAMME OPERATIONS FLOW
# =============================================================================
s10 = new_slide(prs)
slide_header(s10, "05  .  OPERATIONS FLOW", 10)

eyebrow(s10, "INTAKE  ->  DESIGN  ->  LAUNCH  ->  VERIFY  ->  REPORT  ->  SCALE", P, Inches(0.85))
headline(s10, "Automated Programme Operations Lifecycle", P, Inches(1.18), size=30)
tb(s10, P, Inches(2.2), Inches(10), Inches(0.38),
   "From the first partner enquiry to a published outcomes report -- every stage automated.",
   SANS, 13, SUB)

hline(s10, P, Inches(2.72), SW - 2 * P, BORDER)

steps = [
    ("01", "Inbound\nProposal",       "AI intake form\nscores and routes"),
    ("02", "Partner\nQualification",  "Scoring model\nflags fit/no-fit"),
    ("03", "Meeting\nScheduled",      "Auto-invite to\nProgramme Lead"),
    ("04", "Programme\nDesigned",     "Knowledge base\nconsulted by AI"),
    ("05", "Launch\nComms",           "Segmented alerts\nto stakeholders"),
    ("06", "Verification\nCycle",     "Data aggregated\nautomatically"),
    ("07", "Anomaly\nFlagged",        "AI alerts team\nbefore report draft"),
    ("08", "Donor\nReport Sent",      "AI draft reviewed\nand dispatched"),
]

step_w   = Inches(1.4)
step_gap = Inches(0.12)
sx = P

for num, title, detail in steps:
    card(s10, sx, Inches(3.1), step_w, Inches(3.5), rounded=True)
    tb(s10, sx + Inches(0.12), Inches(3.28), step_w - Inches(0.2), Inches(0.38),
       num, SERIF, 20, VIOLET)
    tb(s10, sx + Inches(0.12), Inches(3.75), step_w - Inches(0.2), Inches(0.65),
       title, SERIF, 13, WHITE)
    tb(s10, sx + Inches(0.12), Inches(4.5), step_w - Inches(0.2), Inches(0.75),
       detail, SANS, 10, BODY)
    sx += step_w + step_gap
    if num != "08":
        tb(s10, sx - step_gap + Inches(0.02), Inches(4.2),
           step_gap + Inches(0.05), Inches(0.28),
           ">", SANS, 14, MUTED, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 11 -- IMPLEMENTATION
# =============================================================================
s11 = new_slide(prs)
slide_header(s11, "06  .  IMPLEMENTATION", 11)

eyebrow(s11, "THREE PHASES  --  HIGHEST IMPACT FIRST", P, Inches(0.85))
headline(s11, "Implementation Roadmap", P, Inches(1.18), size=36)
tb(s11, P, Inches(2.2), Inches(10), Inches(0.38),
   "Partner intake addressed first -- the highest daily time cost. "
   "Reporting and comms follow. Knowledge management last, running in parallel.",
   SANS, 13, SUB)
hline(s11, P, Inches(2.72), SW - 2 * P, BORDER)

phases = [
    ("PHASE 01  .  Weeks 1-2",
     "Partner Intake AI + Scheduling",
     "Smart intake form replacing the generic Google Form, implementer scoring model, "
     "automated acknowledgement sequences, and meeting booking for qualified leads.",
     ["Intake form and routing logic", "Implementer scoring model",
      "Automated acknowledgement sequences", "Calendly/scheduling integration",
      "CRM pipeline setup", "Go-live for all inbound channels"]),
    ("PHASE 02  .  Weeks 3-4",
     "Donor Reporting + Comms AI",
     "Verified outcomes data aggregation pipeline, funder-specific report draft generation, "
     "and stakeholder comms segmentation with gated content for lead capture.",
     ["Outcomes data aggregation", "Funder-specific report templates",
      "Anomaly flagging logic", "Audience segmentation model",
      "Gated content email capture", "LinkedIn scheduling setup"]),
    ("PHASE 03  .  Week 5",
     "Knowledge Management + Country Coordination",
     "RAG system over published knowledge base, meeting synthesis automation, "
     "weekly async status synthesis, and cross-programme learning alerts.",
     ["Knowledge base indexing", "RAG query interface",
      "Meeting synthesis pipeline", "Weekly status synthesis flow",
      "Country risk and delay flagging", "Cross-programme learning alerts"]),
]

ph_w   = Inches(3.8)
ph_gap = Inches(0.17)
px = P

for phase_label, phase_name, phase_desc, deliverables in phases:
    card(s11, px, Inches(3.0), ph_w, Inches(3.75), rounded=True)
    tb(s11, px + Inches(0.25), Inches(3.15), ph_w - Inches(0.4), Inches(0.28),
       phase_label, SANS, 9, VIOLET, bold=True)
    tb(s11, px + Inches(0.25), Inches(3.5), ph_w - Inches(0.4), Inches(0.42),
       phase_name, SERIF, 15, WHITE)
    tb(s11, px + Inches(0.25), Inches(4.0), ph_w - Inches(0.4), Inches(0.52),
       phase_desc, SANS, 11, BODY)
    t = Inches(4.62)
    for d in deliverables[:4]:
        tb(s11, px + Inches(0.25), t, ph_w - Inches(0.4), Inches(0.26),
           f">  {d}", SANS, 10, MUTED)
        t += Inches(0.3)
    px += ph_w + ph_gap


# =============================================================================
# SLIDE 12 -- CLOSING
# =============================================================================
s12 = new_slide(prs)

acc = s12.shapes.add_shape(1, Inches(0), Inches(0), SW, Pt(3))
acc.fill.solid()
acc.fill.fore_color.rgb = VIOLET
acc.line.fill.background()

eyebrow(s12, "WHO IS BEHIND THIS", P, Inches(0.85))
headline(s12, "GenosAI Tech", P, Inches(1.18), size=40)
tb(s12, P, Inches(2.22), Inches(6.5), Inches(0.66),
   "GenosAI builds AI automation systems for mission-driven organisations, impact funds, "
   "and multi-country operations. We build partner intake flows, reporting automation, "
   "stakeholder communications AI, knowledge management systems, and coordination tools -- "
   "production-grade from week one, built to scale impact without scaling headcount.",
   SANS, 13, SUB)

hline(s12, P, Inches(3.02), Inches(6), BORDER)

tb(s12, P, Inches(3.2), Inches(4), Inches(0.28),
   "WHAT WE BUILD", SANS, 9, MUTED, bold=True)
t = Inches(3.58)
for item, desc in [
    ("Partner Intake AI",         "Smart triage and routing for all inbound proposals."),
    ("Donor Reporting Automation","Verified outcomes data compiled into funder-specific drafts."),
    ("Stakeholder Comms AI",      "Segmented updates by audience, geography, and role."),
    ("Knowledge Management AI",   "RAG over institutional knowledge -- searchable by any team member."),
    ("Country Coordination AI",   "Async updates and scheduling for distributed teams."),
]:
    t = bullet(s12, "o", item, desc, P, t, Inches(5.5), 12, 11)

card(s12, Inches(7.5), Inches(1.1), Inches(5.05), Inches(5.75), rounded=True)

tb(s12, Inches(7.75), Inches(1.35), Inches(4.6), Inches(0.35),
   "CLOSING NOTE", SANS, 9, MUTED, bold=True)
tb(s12, Inches(7.6), Inches(1.72), Inches(0.4), Inches(0.55),
   '"', SERIF, 28, VIOLET)
tb(s12, Inches(7.95), Inches(1.78), Inches(4.3), Inches(1.05),
   "EOF's 2030 goal is 15 active partnerships reaching 2 million children. "
   "That ambition requires an operational layer that scales with the programmes -- "
   "not one that adds a new hire for every new country.",
   SERIF, 14, WHITE, italic=True)

hline(s12, Inches(7.75), Inches(2.92), Inches(4.45), BORDER)

tb(s12, Inches(7.75), Inches(3.12), Inches(4.5), Inches(0.28),
   "NEXT STEPS", SANS, 9, VIOLET, bold=True)
t = Inches(3.5)
for step, desc in [
    ("01", "20-minute call to walk through the audit and confirm scope."),
    ("02", "Phase 1 kick-off -- partner intake AI live within 7 days."),
    ("03", "First qualified implementer routed automatically -- visible from week one."),
]:
    tb(s12, Inches(7.75), t, Inches(0.4), Inches(0.3), step, SERIF, 14, VIOLET)
    tb(s12, Inches(8.2), t + Inches(0.02), Inches(4.15), Inches(0.5),
       desc, SANS, 12, BODY)
    t += Inches(0.62)

hline(s12, Inches(7.75), Inches(5.45), Inches(4.45), BORDER)

tb(s12, Inches(7.75), Inches(5.65), Inches(4.5), Inches(0.28),
   "Rohan Malik  .  CEO, Genos AI", SANS, 12, WHITE, bold=True)
tb(s12, Inches(7.75), Inches(6.0), Inches(4.5), Inches(0.28),
   "hello@genosai.tech  .  +91 638-714-2699", SANS, 11, BODY)
tb(s12, Inches(7.75), Inches(6.32), Inches(4.5), Inches(0.28),
   "www.genosai.tech", SANS, 11, MUTED)

tb(s12, P, SH - Inches(0.42), Inches(8), Inches(0.28),
   "Confidential  .  Education Outcomes Fund  .  2026", SANS, 8, DIM)
tb(s12, Inches(9.5), SH - Inches(0.42), Inches(3.08), Inches(0.28),
   f"{TOTAL:02d} / {TOTAL}", SANS, 8, DIM, align=PP_ALIGN.RIGHT)

tb(s12, P, SH - Inches(0.65), Inches(3), Inches(0.4),
   "GenosAI", SERIF, 20, WHITE)


# -- Save ----------------------------------------------------------------------
prs.save(OUT)
print(f"Saved: {OUT}  ({TOTAL} slides)")
