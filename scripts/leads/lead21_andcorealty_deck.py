"""
Build AndCoRealty_Pitch_GenosAI.pptx from scratch.
Track: Pitch-only (automation focus, no website gap analysis).
Design: docs/BRAND_ASSETS.md — #0A0A0F canvas, Instrument Serif headlines, Satoshi body.
12 slides: cover, group overview, vision, systems, 4 module slides, agent journey,
benefits, implementation, closing.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).parent.parent.parent
OUT  = ROOT / "output" / "decks" / "AndCoRealty_Pitch_GenosAI.pptx"

# ── Brand colors ─────────────────────────────────────────────────────────────
BG     = RGBColor(0x0A, 0x0A, 0x0F)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BODY   = RGBColor(0x9D, 0x9D, 0x9F)
SUB    = RGBColor(0x86, 0x86, 0x8A)
MUTED  = RGBColor(0x6B, 0x6B, 0x6E)
DIM    = RGBColor(0x44, 0x44, 0x47)
VIOLET = RGBColor(0xA7, 0x8B, 0xFA)
CARD   = RGBColor(0x11, 0x11, 0x18)
BORDER = RGBColor(0x22, 0x22, 0x2A)

SERIF = "Instrument Serif"
SANS  = "Satoshi"

SW    = Inches(13.33)
SH    = Inches(7.5)
P     = Inches(0.75)
TOTAL = 12


# ── Helpers ───────────────────────────────────────────────────────────────────

def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def tb(slide, l, t, w, h, text, font=SANS, size=14, color=BODY,
       bold=False, italic=False, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text = text
    r.font.name   = font
    r.font.size   = Pt(size)
    r.font.color.rgb = color
    r.font.bold   = bold
    r.font.italic = italic
    return box


def tb_lines(slide, l, t, w, h, lines, font=SANS, size=14, color=BODY,
             bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name  = font
        r.font.size  = Pt(size)
        r.font.color.rgb = color
        r.font.bold  = bold
    return box


def card(slide, l, t, w, h, rounded=True):
    s = slide.shapes.add_shape(5 if rounded else 1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = CARD
    s.line.color.rgb = BORDER
    s.line.width = Pt(0.75)
    return s


def hline(slide, l, t, w, color=BORDER):
    s = slide.shapes.add_shape(1, l, t, w, Pt(0.75))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def eyebrow(slide, text, l, t, w=Inches(6)):
    return tb(slide, l, t, w, Inches(0.3),
              text.upper(), SANS, 10, VIOLET, bold=True)


def headline(slide, text, l, t, w=Inches(10), size=42):
    return tb(slide, l, t, w, Inches(1.6), text, SERIF, size, WHITE)


def slide_header(slide, section, page):
    tb(slide, P, Inches(0.22), Inches(6), Inches(0.28),
       "GENOSAI  ×  ANDCO REALTY GROUP", SANS, 9, MUTED, bold=True)
    tb(slide, Inches(7.5), Inches(0.22), Inches(5.08), Inches(0.28),
       section, SANS, 9, MUTED, bold=True, align=PP_ALIGN.RIGHT)
    hline(slide, P, Inches(0.58), SW - 2 * P)
    tb(slide, P, SH - Inches(0.42), Inches(8), Inches(0.28),
       "AI Automation Proposal  ·  AndCo Realty Group  ·  2026",
       SANS, 8, DIM)
    tb(slide, Inches(9.5), SH - Inches(0.42), Inches(3.08), Inches(0.28),
       f"{page:02d} / {TOTAL}", SANS, 8, DIM, align=PP_ALIGN.RIGHT)


def bullet(slide, dot_char, title, desc, l, t,
           col_w=Inches(5.0), title_size=13, desc_size=12):
    tb(slide, l, t, Inches(0.25), Inches(0.26),
       dot_char, SANS, 13, VIOLET, bold=True)
    tb(slide, l + Inches(0.3), t, col_w - Inches(0.3), Inches(0.26),
       title, SANS, title_size, WHITE, bold=True)
    tb(slide, l + Inches(0.3), t + Inches(0.27), col_w - Inches(0.3), Inches(0.38),
       desc, SANS, desc_size, BODY)
    return t + Inches(0.72)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH

s1 = new_slide(prs)

accent = s1.shapes.add_shape(1, Inches(0), Inches(0), SW, Pt(3))
accent.fill.solid()
accent.fill.fore_color.rgb = VIOLET
accent.line.fill.background()

tb(s1, P, Inches(1.0), Inches(5.5), Inches(0.35),
   "PREPARED FOR", SANS, 10, MUTED, bold=True)
hline(s1, P, Inches(1.42), Inches(5.5), BORDER)

tb(s1, P, Inches(1.6), Inches(8.5), Inches(1.4),
   "AndCo Realty Group", SERIF, 64, WHITE)

tb(s1, P, Inches(3.1), Inches(6.5), Inches(0.32),
   "Fixed-Fee Real Estate Group  ·  17 Independent Offices  ·  New Zealand",
   SANS, 13, BODY)

tb(s1, P, Inches(3.6), Inches(5.5), Inches(0.28),
   "Dunedin  ·  Queenstown  ·  Wellington  ·  Bay of Plenty  ·  Taranaki  ·  Whanganui",
   SANS, 11, MUTED)

hline(s1, P, Inches(4.06), Inches(5.5), BORDER)

tb(s1, P, Inches(4.26), Inches(7), Inches(0.28),
   "AI AUTOMATION PROPOSAL", SANS, 10, VIOLET, bold=True)
tb_lines(s1, P, Inches(4.66), Inches(6.5), Inches(1.2),
         ["An AI automation layer for AndCo Realty Group —",
          "built to onboard new office operators consistently,",
          "route buyer inquiries to the right agent,",
          "and scale support across 17 offices without adding headcount."],
         SANS, 13, SUB)

card(s1, Inches(9.0), Inches(1.2), Inches(3.58), Inches(2.5), rounded=True)
tb(s1, Inches(9.25), Inches(1.5), Inches(3.0), Inches(0.28),
   "PREPARED BY", SANS, 9, MUTED, bold=True)
tb(s1, Inches(9.25), Inches(1.88), Inches(3.0), Inches(0.32),
   "GenosAI Tech", SERIF, 20, WHITE)
tb(s1, Inches(9.25), Inches(2.3), Inches(3.0), Inches(0.26),
   "Rohan Malik  ·  Founder & CEO", SANS, 11, BODY)
tb(s1, Inches(9.25), Inches(2.65), Inches(3.0), Inches(0.26),
   "hello@genosai.tech  ·  www.genosai.tech", SANS, 10, MUTED)

tb(s1, P, SH - Inches(0.65), Inches(4), Inches(0.45),
   "GenosAI", SERIF, 22, WHITE)
tb(s1, P + Inches(1.65), SH - Inches(0.62), Inches(3), Inches(0.3),
   "v1.0  ·  2026  ·  Confidential", SANS, 9, DIM)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — ANDCO AT A GLANCE
# ══════════════════════════════════════════════════════════════════════════════
s2 = new_slide(prs)
slide_header(s2, "01  ·  THE GROUP", 2)

eyebrow(s2, "THE MODEL AT A GLANCE", P, Inches(0.85))
headline(s2, "17 Offices. One Group. Agent-Owned.", P, Inches(1.18), size=38)
tb(s2, P, Inches(2.2), Inches(8.5), Inches(0.42),
   "AndCo Realty Group operates a fixed-fee, agent-owned model across New Zealand — "
   "agents run their own branded businesses under a shared compliance, training, and "
   "technology framework. Scaling that framework across 17 offices is the operational challenge.",
   SANS, 13, SUB)

hline(s2, P, Inches(2.82), SW - 2 * P, BORDER)

stats = [
    ("OFFICES",         "17",       "independent across NZ"),
    ("MODEL",           "Fixed-fee","agents keep full commission"),
    ("FOCUS",           "$500K+",   "residential property"),
    ("COMPLIANCE",      "REAA 2008","licensed real estate agency"),
]
sw_each = Inches(2.8)
gap     = Inches(0.19)
sx      = P

for label, value, sub in stats:
    card(s2, sx, Inches(3.1), sw_each, Inches(2.0), rounded=True)
    tb(s2, sx + Inches(0.25), Inches(3.28), sw_each - Inches(0.4), Inches(0.28),
       label, SANS, 9, MUTED, bold=True)
    tb(s2, sx + Inches(0.25), Inches(3.65), sw_each - Inches(0.4), Inches(0.62),
       value, SERIF, 28, WHITE)
    tb(s2, sx + Inches(0.25), Inches(4.38), sw_each - Inches(0.4), Inches(0.45),
       sub, SANS, 11, BODY)
    sx += sw_each + gap

tb(s2, P, Inches(5.4), Inches(11), Inches(0.28),
   "OFFICE LOCATIONS", SANS, 9, MUTED, bold=True)
hline(s2, P, Inches(5.72), SW - 2 * P, BORDER)
tb(s2, P, Inches(5.82), SW - 2 * P, Inches(0.42),
   "Dunedin  ·  Queenstown  ·  Wellington  ·  Bay of Plenty  ·  Taranaki  ·  Whanganui  ·  "
   "and 11 further locations nationwide",
   SANS, 12, BODY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — VISION
# ══════════════════════════════════════════════════════════════════════════════
s3 = new_slide(prs)
slide_header(s3, "02  ·  VISION", 3)

eyebrow(s3, "THE OPPORTUNITY", P, Inches(0.85))
headline(s3, "An AI Automation Layer\nfor AndCo Realty Group", P, Inches(1.18), size=38)

card(s3, P, Inches(2.85), Inches(7.5), Inches(1.25), rounded=True)
tb(s3, P + Inches(0.4), Inches(2.98), Inches(7.0), Inches(0.4),
   "“", SERIF, 32, VIOLET)
tb(s3, P + Inches(0.75), Inches(3.02), Inches(6.5), Inches(0.88),
   "We're not replacing how your offices operate. We're building the infrastructure that lets "
   "every new office operator get the same quality of support, leads, and compliance guidance "
   "— without your team running it manually for each one.",
   SERIF, 15, WHITE, italic=True)

col_w = Inches(2.8)
cols  = [P, P + col_w + Inches(0.15), P + 2*(col_w + Inches(0.15)), P + 3*(col_w + Inches(0.15))]
pillars = [
    ("01", "Consistent\nonboarding",
     "Every new office operator gets the same structured setup flow."),
    ("02", "Buyer lead\nrouting",
     "Inquiries qualified and sent to the right agent, automatically."),
    ("03", "Vendor\nconversion",
     "Warm appraisal leads nurtured until they list — without manual follow-up."),
    ("04", "Compliance\non autopilot",
     "REAA and CPD requirements tracked for all 17 offices without admin overhead."),
]
for (num, title, desc), x in zip(pillars, cols):
    card(s3, x, Inches(4.35), col_w, Inches(2.65), rounded=True)
    tb(s3, x + Inches(0.25), Inches(4.55), Inches(0.6), Inches(0.42),
       num, SERIF, 22, VIOLET)
    tb(s3, x + Inches(0.25), Inches(5.05), col_w - Inches(0.4), Inches(0.65),
       title, SERIF, 16, WHITE)
    tb(s3, x + Inches(0.25), Inches(5.78), col_w - Inches(0.4), Inches(0.88),
       desc, SANS, 11, BODY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SYSTEMS OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
s4 = new_slide(prs)
slide_header(s4, "03  ·  AUTOMATION SYSTEMS", 4)

eyebrow(s4, "WHAT THE AUTOMATION LAYER BRINGS", P, Inches(0.85))
headline(s4, "Five AI Systems for AndCo Realty", P, Inches(1.18), size=36)

systems = [
    ("A", "Agent Onboarding\nAI",       "Consistent setup for every new office."),
    ("B", "Buyer Inquiry\nRouting AI",  "Right agent, right area, every time."),
    ("C", "Vendor Nurture\nSystem",     "Warm appraisals converted to listings."),
    ("D", "Compliance &\nCPD AI",       "REAA tracking across all 17 offices."),
    ("E", "Review &\nReferral AI",      "Social proof built at settlement."),
]
sw_each = Inches(2.25)
gap     = Inches(0.17)
sx      = P

for letter, name, tagline in systems:
    card(s4, sx, Inches(2.6), sw_each, Inches(4.1), rounded=True)
    badge = s4.shapes.add_shape(1, sx + Inches(0.25), Inches(2.85), Inches(0.55), Inches(0.55))
    badge.fill.solid()
    badge.fill.fore_color.rgb = VIOLET
    badge.line.fill.background()
    tb(s4, sx + Inches(0.25), Inches(2.87), Inches(0.55), Inches(0.48),
       letter, SANS, 16, BG, bold=True, align=PP_ALIGN.CENTER)
    tb(s4, sx + Inches(0.2), Inches(3.55), sw_each - Inches(0.4), Inches(0.75),
       name, SERIF, 15, WHITE)
    tb(s4, sx + Inches(0.2), Inches(4.38), sw_each - Inches(0.4), Inches(0.55),
       tagline, SANS, 12, BODY)
    sx += sw_each + gap


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — MODULE A: AGENT ONBOARDING AI
# ══════════════════════════════════════════════════════════════════════════════
s5 = new_slide(prs)
slide_header(s5, "03A  ·  AGENT ONBOARDING", 5)

eyebrow(s5, "MODULE A", P, Inches(0.85))
headline(s5, "Agent Onboarding AI", P, Inches(1.18), size=36)
tb(s5, P, Inches(2.2), Inches(6.5), Inches(0.42),
   "Every new AndCo office operator goes through the same structured onboarding — "
   "compliance, marketing, technology — without your team running it manually each time.",
   SANS, 14, SUB)

hline(s5, P, Inches(2.75), SW - 2 * P, BORDER)

t = Inches(3.0)
for title, desc in [
    ("WhatsApp onboarding flow",
     "New operators receive a structured sequence: compliance orientation, branding setup, "
     "platform training, and milestone check-ins — delivered via WhatsApp over 2 weeks."),
    ("Compliance orientation",
     "REAA obligations, trust account requirements, and disclosure rules delivered "
     "as a structured sequence with acknowledgment tracking."),
    ("Marketing setup guidance",
     "Step-by-step setup for Trade Me, listing profiles, social accounts, "
     "and Google Business — automated delivery with completion tracking."),
    ("Progress milestones",
     "AI tracks where each new operator is in onboarding and flags anyone who stalls."),
    ("Handover to Craig / Cameron",
     "Completed onboarding summary delivered to your team — no manual check-ins needed."),
]:
    t = bullet(s5, "·", title, desc, P, t, Inches(5.5), 13, 12)

card(s5, Inches(7.5), Inches(3.0), Inches(5.0), Inches(3.7), rounded=True)
tb(s5, Inches(7.8), Inches(3.2), Inches(4.5), Inches(0.28),
   "SCALE  ·  WITHOUT SCALING TEAM", SANS, 9, MUTED, bold=True)
tb(s5, Inches(7.8), Inches(3.6), Inches(4.5), Inches(0.42),
   "17 offices.", SERIF, 28, WHITE)
tb(s5, Inches(7.8), Inches(4.15), Inches(4.5), Inches(0.42),
   "One onboarding system.", SERIF, 20, VIOLET)
tb(s5, Inches(7.8), Inches(4.72), Inches(4.5), Inches(0.55),
   "New operators get the same quality of setup on day one — whether you're adding "
   "your 18th office or your 50th.",
   SANS, 11, BODY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — MODULE B & C: BUYER ROUTING + VENDOR NURTURE
# ══════════════════════════════════════════════════════════════════════════════
s6 = new_slide(prs)
slide_header(s6, "03B / 03C  ·  BUYER ROUTING & VENDOR NURTURE", 6)

eyebrow(s6, "MODULE B", P, Inches(0.85))
headline(s6, "Buyer Inquiry Routing AI", P, Inches(1.18), size=28)
tb(s6, P, Inches(2.05), Inches(5.5), Inches(0.38),
   "Every buyer inquiry from the AndCo website reaches the right regional agent — automatically.",
   SANS, 13, SUB)

t = Inches(2.6)
for title, desc in [
    ("Area-based routing",
     "Buyer specifies region; AI routes to the right office operator."),
    ("Buyer pre-qualification",
     "Budget, property type, and timeline captured before agent contact."),
    ("WhatsApp response",
     "AI responds on WhatsApp instantly — no buyer waits for an email reply."),
    ("Viewing booking",
     "AI books the viewing directly in the agent's calendar once lead is qualified."),
    ("Unrouted lead fallback",
     "If no agent is available, lead held and notified when an agent responds."),
]:
    t = bullet(s6, "·", title, desc, P, t, Inches(5.5), 12, 11)

hline(s6, Inches(6.67) - Pt(0.5), Inches(0.85), Pt(1), BORDER)

CR = Inches(6.85)
eyebrow(s6, "MODULE C", CR, Inches(0.85))
headline(s6, "Vendor Nurture System", CR, Inches(1.18), size=28)
tb(s6, CR, Inches(2.05), Inches(5.7), Inches(0.38),
   "Sellers who requested an appraisal but didn't list — automatically re-engaged over time.",
   SANS, 13, SUB)

t = Inches(2.6)
for title, desc in [
    ("Post-appraisal sequence",
     "Sellers receive a structured follow-up sequence at 1 week, 1 month, and 3 months."),
    ("Market update delivery",
     "Relevant local sales data sent to keep the agent top of mind while the seller decides."),
    ("Price movement alerts",
     "If comparable properties sell nearby, AI sends a timely market movement note."),
    ("Seasonal re-engagement",
     "Dormant appraisal leads re-contacted at spring and pre-summer listing peaks."),
    ("Agent attribution",
     "Every touchpoint tracked back to the originating agent — no lead confusion."),
]:
    t = bullet(s6, "·", title, desc, CR, t, Inches(5.7), 12, 11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — MODULE D & E: COMPLIANCE + REVIEWS
# ══════════════════════════════════════════════════════════════════════════════
s7 = new_slide(prs)
slide_header(s7, "03D / 03E  ·  COMPLIANCE & REVIEWS", 7)

eyebrow(s7, "MODULE D", P, Inches(0.85))
headline(s7, "Compliance & CPD AI", P, Inches(1.18), size=28)
tb(s7, P, Inches(2.05), Inches(5.5), Inches(0.38),
   "REAA licensing and CPD requirements tracked per agent across all 17 offices — automated.",
   SANS, 13, SUB)

t = Inches(2.6)
for title, desc in [
    ("License renewal tracking",
     "Each agent's REAA license expiry tracked; reminders at 90, 30, and 7 days."),
    ("CPD completion monitoring",
     "Annual CPD hours tracked per agent; AI flags anyone falling behind on schedule."),
    ("Trust account compliance",
     "Recurring trust account obligation reminders delivered to each office operator."),
    ("Disclosure obligation prompts",
     "Automated reminders for key disclosure requirements at relevant deal stages."),
    ("Compliance dashboard",
     "Jason and Craig receive a monthly compliance status summary across all offices."),
]:
    t = bullet(s7, "·", title, desc, P, t, Inches(5.5), 12, 11)

CR = Inches(6.85)
eyebrow(s7, "MODULE E", CR, Inches(0.85))
headline(s7, "Review & Referral AI", CR, Inches(1.18), size=28)
tb(s7, CR, Inches(2.05), Inches(5.7), Inches(0.38),
   "Post-settlement review requests and referral asks — automated across all 17 offices.",
   SANS, 13, SUB)

t = Inches(2.6)
for title, desc in [
    ("Settlement-triggered request",
     "At settlement date, buyer and seller each receive an automated review request."),
    ("Google and Trade Me reviews",
     "Request links to both platforms — one message, two social proof channels."),
    ("Referral ask sequence",
     "7 days after settlement, AI asks for a referral contact — while satisfaction is high."),
    ("Testimonial capture",
     "Positive reviewers offered the option to provide a written testimonial for AndCo website."),
    ("Review aggregation",
     "Monthly review summary for each office and group-level performance tracking."),
]:
    t = bullet(s7, "·", title, desc, CR, t, Inches(5.7), 12, 11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — AGENT / OPERATOR JOURNEY FLOW
# ══════════════════════════════════════════════════════════════════════════════
s8 = new_slide(prs)
slide_header(s8, "04  ·  OPERATOR JOURNEY", 8)

eyebrow(s8, "NEW OFFICE  →  PRODUCTIVE OPERATOR  →  GROWING OFFICE", P, Inches(0.85))
headline(s8, "Automated Operator Journey", P, Inches(1.18), size=36)
tb(s8, P, Inches(2.2), Inches(10), Inches(0.38),
   "From the moment a new office operator joins AndCo to a fully productive, compliant, "
   "and growing office — every support step automated.",
   SANS, 13, SUB)

hline(s8, P, Inches(2.72), SW - 2 * P, BORDER)

steps = [
    ("01", "Operator\nJoins",          "Signs up under\nAndCo franchise"),
    ("02", "Onboarding\nFlow Starts",  "WhatsApp sequence\nbegins day one"),
    ("03", "Compliance\nOrientation",  "REAA obligations\ndelivered and tracked"),
    ("04", "Marketing\nSetup",         "Trade Me, Google,\nsocial channels"),
    ("05", "First Buyer\nInquiry",     "AI routes lead\nto their office"),
    ("06", "Vendor\nAppraisal",        "AI nurtures\nunconverted sellers"),
    ("07", "CPD\nTracking",            "Annual compliance\nmonitored continuously"),
    ("08", "Reviews &\nReferrals",     "Post-settlement\nautomation active"),
]

step_w   = Inches(1.4)
step_gap = Inches(0.12)
sx = P

for num, title, detail in steps:
    card(s8, sx, Inches(3.1), step_w, Inches(3.5), rounded=True)
    tb(s8, sx + Inches(0.12), Inches(3.28), step_w - Inches(0.2), Inches(0.38),
       num, SERIF, 20, VIOLET)
    tb(s8, sx + Inches(0.12), Inches(3.75), step_w - Inches(0.2), Inches(0.65),
       title, SERIF, 13, WHITE)
    tb(s8, sx + Inches(0.12), Inches(4.5), step_w - Inches(0.2), Inches(0.75),
       detail, SANS, 10, BODY)
    sx += step_w + step_gap
    if num != "08":
        tb(s8, sx - step_gap + Inches(0.02), Inches(4.2),
           step_gap + Inches(0.05), Inches(0.28),
           "›", SANS, 14, MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — BENEFITS
# ══════════════════════════════════════════════════════════════════════════════
s9 = new_slide(prs)
slide_header(s9, "05  ·  BENEFITS", 9)

eyebrow(s9, "WHY THIS WORKS FOR ANDCO REALTY", P, Inches(0.85))
headline(s9, "What Changes After Launch", P, Inches(1.18), size=36)
hline(s9, P, Inches(2.72), SW - 2 * P, BORDER)

benefits = [
    ("Consistent onboarding at any scale",  "Every new office operator gets the same quality setup — whether you add 2 or 20 offices."),
    ("No buyer inquiry falls through",       "Every inquiry from the website reaches the right agent within minutes."),
    ("Vendor pipeline never goes cold",      "Appraisal leads nurtured automatically over weeks and months."),
    ("Compliance across 17 offices",         "REAA and CPD tracked per agent — no manual checking by Jason or Craig."),
    ("Reviews without chasing",              "Post-settlement review requests sent automatically for every transaction."),
    ("Referral pipeline activated",          "Every settlement is an opportunity for a referred lead — automated."),
    ("Agent retention",                      "Operators who feel supported stay — consistent AI-delivered guidance reduces churn."),
    ("Scales with growth",                   "Adding offices adds no linear admin overhead — the same system handles 17 or 170."),
]

col1, col2 = benefits[:4], benefits[4:]
CL, CR = P, Inches(7.1)
t1, t2 = Inches(2.95), Inches(2.95)

for title, desc in col1:
    card(s9, CL, t1, Inches(5.75), Inches(0.82), rounded=True)
    tb(s9, CL + Inches(0.25), t1 + Inches(0.1), Inches(0.35), Inches(0.3),
       "✓", SANS, 14, VIOLET, bold=True)
    tb(s9, CL + Inches(0.65), t1 + Inches(0.08), Inches(4.8), Inches(0.3),
       title, SANS, 13, WHITE, bold=True)
    tb(s9, CL + Inches(0.65), t1 + Inches(0.4), Inches(4.8), Inches(0.35),
       desc, SANS, 11, BODY)
    t1 += Inches(0.95)

for title, desc in col2:
    card(s9, CR, t2, Inches(5.75), Inches(0.82), rounded=True)
    tb(s9, CR + Inches(0.25), t2 + Inches(0.1), Inches(0.35), Inches(0.3),
       "✓", SANS, 14, VIOLET, bold=True)
    tb(s9, CR + Inches(0.65), t2 + Inches(0.08), Inches(4.8), Inches(0.3),
       title, SANS, 13, WHITE, bold=True)
    tb(s9, CR + Inches(0.65), t2 + Inches(0.4), Inches(4.8), Inches(0.35),
       desc, SANS, 11, BODY)
    t2 += Inches(0.95)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — FUTURE ROADMAP
# ══════════════════════════════════════════════════════════════════════════════
s10 = new_slide(prs)
slide_header(s10, "06  ·  ROADMAP", 10)

eyebrow(s10, "BEYOND THE LAUNCH", P, Inches(0.85))
headline(s10, "Future Scalability", P, Inches(1.18), size=36)
tb(s10, P, Inches(2.2), Inches(10), Inches(0.38),
   "The same automation infrastructure scales as AndCo grows to 30, 50, and 100 offices.",
   SANS, 13, SUB)
hline(s10, P, Inches(2.72), SW - 2 * P, BORDER)

future = [
    ("New office recruitment AI",       "AI nurtures prospective office operators from first enquiry to signing."),
    ("WhatsApp AI for each office",      "Localised WhatsApp AI per office — each operator's own AI agent."),
    ("Seller valuation lead magnet",     "Automated AVM-style appraisal request capture and nurture."),
    ("Property management AI",          "If rental management is added — lease renewals, maintenance flows."),
    ("Agent performance analytics",     "Monthly performance summaries per office delivered to Jason automatically."),
    ("Group-wide market reports",        "Automated monthly NZ market update sent to all vendors and past clients."),
    ("International buyer AI",          "Separate inquiry flow for offshore/expat buyers — timezone-aware follow-up."),
    ("Group podcast / content AI",       "Social media scheduling and distribution for group-level content."),
]

col1, col2 = future[:4], future[4:]
CL, CR = P, Inches(7.1)
t1, t2 = Inches(2.95), Inches(2.95)

for title, desc in col1:
    t1 = bullet(s10, "·", title, desc, CL, t1, Inches(5.75))
for title, desc in col2:
    t2 = bullet(s10, "·", title, desc, CR, t2, Inches(5.75))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — IMPLEMENTATION
# ══════════════════════════════════════════════════════════════════════════════
s11 = new_slide(prs)
slide_header(s11, "07  ·  IMPLEMENTATION", 11)

eyebrow(s11, "BUILT IN THREE CLEAN PHASES", P, Inches(0.85))
headline(s11, "Implementation Phases", P, Inches(1.18), size=36)
tb(s11, P, Inches(2.2), Inches(10), Inches(0.38),
   "Each phase ships visible value — live results before the next phase begins.",
   SANS, 13, SUB)
hline(s11, P, Inches(2.72), SW - 2 * P, BORDER)

phases = [
    ("PHASE 01  ·  Week 1",
     "Agent Onboarding AI",
     "WhatsApp onboarding flow for new office operators — compliance orientation, "
     "marketing setup guidance, and milestone tracking — live and running.",
     ["Onboarding flow design with Kelly/Craig", "Compliance sequence content",
      "Marketing setup delivery sequence", "Milestone tracking setup",
      "First operator goes through live flow", "Testing and refinement"]),
    ("PHASE 02  ·  Weeks 2-3",
     "Buyer Routing + Vendor Nurture",
     "Buyer inquiry routing AI from the AndCo website, plus vendor nurture "
     "sequences for all existing warm appraisal leads.",
     ["Website inquiry integration", "Area-based routing configuration",
      "WhatsApp AI qualification flow", "Vendor appraisal sequence design",
      "Market update delivery setup", "Go-live across all 17 offices"]),
    ("PHASE 03  ·  Weeks 4-5",
     "Compliance + Review AI",
     "REAA and CPD tracking per agent across all offices, and automated "
     "post-settlement review and referral sequences.",
     ["Agent license data import", "CPD tracking configuration",
      "Compliance reminder sequences", "Post-settlement review flow",
      "Referral ask sequence", "Compliance dashboard for Jason/Craig"]),
]

ph_w   = Inches(3.8)
ph_gap = Inches(0.17)
px = P

for phase_label, phase_name, phase_desc, deliverables in phases:
    card(s11, px, Inches(3.0), ph_w, Inches(3.75), rounded=True)
    tb(s11, px + Inches(0.25), Inches(3.15), ph_w - Inches(0.4), Inches(0.28),
       phase_label, SANS, 9, VIOLET, bold=True)
    tb(s11, px + Inches(0.25), Inches(3.5), ph_w - Inches(0.4), Inches(0.42),
       phase_name, SERIF, 15, WHITE)
    tb(s11, px + Inches(0.25), Inches(4.0), ph_w - Inches(0.4), Inches(0.52),
       phase_desc, SANS, 11, BODY)
    t = Inches(4.62)
    for d in deliverables[:4]:
        tb(s11, px + Inches(0.25), t, ph_w - Inches(0.4), Inches(0.26),
           f"›  {d}", SANS, 10, MUTED)
        t += Inches(0.3)
    px += ph_w + ph_gap


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — CLOSING
# ══════════════════════════════════════════════════════════════════════════════
s12 = new_slide(prs)

acc = s12.shapes.add_shape(1, Inches(0), Inches(0), SW, Pt(3))
acc.fill.solid()
acc.fill.fore_color.rgb = VIOLET
acc.line.fill.background()

eyebrow(s12, "WHO IS BEHIND THIS", P, Inches(0.85))
headline(s12, "GenosAI Tech", P, Inches(1.18), size=40)
tb(s12, P, Inches(2.22), Inches(6.5), Inches(0.66),
   "GenosAI builds AI automation systems for real estate groups, franchises, and independent agencies. "
   "We specialise in agent onboarding flows, buyer routing AI, vendor nurture sequences, compliance "
   "tracking, and review automation — production-grade from week one.",
   SANS, 13, SUB)

hline(s12, P, Inches(3.02), Inches(6), BORDER)

tb(s12, P, Inches(3.2), Inches(4), Inches(0.28),
   "WHAT WE BUILD", SANS, 9, MUTED, bold=True)
t = Inches(3.58)
for item, desc in [
    ("Agent Onboarding AI",       "Structured WhatsApp setup flows for every new operator."),
    ("Buyer Routing AI",          "Inquiries qualified and sent to the right office, instantly."),
    ("Vendor Nurture System",     "Warm appraisal leads converted to listings over time."),
    ("Compliance & CPD AI",       "REAA tracking for all agents across all offices."),
    ("Review & Referral AI",      "Post-settlement social proof built automatically."),
]:
    t = bullet(s12, "◆", item, desc, P, t, Inches(5.5), 12, 11)

card(s12, Inches(7.5), Inches(1.1), Inches(5.05), Inches(5.75), rounded=True)

tb(s12, Inches(7.75), Inches(1.35), Inches(4.6), Inches(0.35),
   "CLOSING NOTE", SANS, 9, MUTED, bold=True)
tb(s12, Inches(7.6), Inches(1.72), Inches(0.4), Inches(0.55),
   "“", SERIF, 28, VIOLET)
tb(s12, Inches(7.95), Inches(1.78), Inches(4.3), Inches(1.0),
   "Built to make every AndCo office operator as supported as if Jason, Kelly, and Craig "
   "were personally managing their onboarding, leads, and compliance.",
   SERIF, 14, WHITE, italic=True)

hline(s12, Inches(7.75), Inches(2.88), Inches(4.45), BORDER)

tb(s12, Inches(7.75), Inches(3.08), Inches(4.5), Inches(0.28),
   "NEXT STEPS", SANS, 9, VIOLET, bold=True)
t = Inches(3.45)
for step, desc in [
    ("01", "20-minute call to confirm scope and integration points."),
    ("02", "Phase 1 kick-off — onboarding AI live within 7 days."),
    ("03", "First operator goes through the flow — results visible from day one."),
]:
    tb(s12, Inches(7.75), t, Inches(0.4), Inches(0.3), step, SERIF, 14, VIOLET)
    tb(s12, Inches(8.2), t + Inches(0.02), Inches(4.15), Inches(0.5),
       desc, SANS, 12, BODY)
    t += Inches(0.62)

hline(s12, Inches(7.75), Inches(5.45), Inches(4.45), BORDER)

tb(s12, Inches(7.75), Inches(5.65), Inches(4.5), Inches(0.28),
   "Rohan Malik  ·  CEO, Genos AI", SANS, 12, WHITE, bold=True)
tb(s12, Inches(7.75), Inches(6.0), Inches(4.5), Inches(0.28),
   "hello@genosai.tech  ·  +91 638-714-2699", SANS, 11, BODY)
tb(s12, Inches(7.75), Inches(6.32), Inches(4.5), Inches(0.28),
   "www.genosai.tech", SANS, 11, MUTED)

tb(s12, P, SH - Inches(0.42), Inches(8), Inches(0.28),
   "Confidential  ·  AndCo Realty Group  ·  2026", SANS, 8, DIM)
tb(s12, Inches(9.5), SH - Inches(0.42), Inches(3.08), Inches(0.28),
   f"{TOTAL:02d} / {TOTAL}", SANS, 8, DIM, align=PP_ALIGN.RIGHT)

tb(s12, P, SH - Inches(0.65), Inches(3), Inches(0.4),
   "GenosAI", SERIF, 20, WHITE)


# ── Save ─────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved: {OUT}  ({TOTAL} slides)")
