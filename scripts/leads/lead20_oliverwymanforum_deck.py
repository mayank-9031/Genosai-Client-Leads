"""
Build OliverWymanForum_Pitch_GenosAI.pptx from scratch.
Track: Pitch-only (website score 7.5/10 — no website gap analysis).
Design: docs/BRAND_ASSETS.md — #0A0A0F canvas, Instrument Serif headlines, Satoshi body.
12 slides: cover, forum overview, vision, systems, 4 module slides, executive journey,
benefits, implementation, closing.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).parent.parent.parent
OUT  = ROOT / "output" / "decks" / "OliverWymanForum_Pitch_GenosAI.pptx"

# ── Brand colors ─────────────────────────────────────────────────────────────
BG     = RGBColor(0x0A, 0x0A, 0x0F)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BODY   = RGBColor(0x9D, 0x9D, 0x9F)
SUB    = RGBColor(0x86, 0x86, 0x8A)
MUTED  = RGBColor(0x6B, 0x6B, 0x6E)
DIM    = RGBColor(0x44, 0x44, 0x47)
VIOLET = RGBColor(0xA7, 0x8B, 0xFA)
CARD   = RGBColor(0x11, 0x11, 0x18)
BORDER = RGBColor(0x22, 0x22, 0x2A)

SERIF = "Instrument Serif"
SANS  = "Satoshi"

SW    = Inches(13.33)
SH    = Inches(7.5)
P     = Inches(0.75)
TOTAL = 12


# ── Helpers ───────────────────────────────────────────────────────────────────

def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def tb(slide, l, t, w, h, text, font=SANS, size=14, color=BODY,
       bold=False, italic=False, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text = text
    r.font.name   = font
    r.font.size   = Pt(size)
    r.font.color.rgb = color
    r.font.bold   = bold
    r.font.italic = italic
    return box


def tb_lines(slide, l, t, w, h, lines, font=SANS, size=14, color=BODY,
             bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name  = font
        r.font.size  = Pt(size)
        r.font.color.rgb = color
        r.font.bold  = bold
    return box


def card(slide, l, t, w, h, rounded=True):
    s = slide.shapes.add_shape(5 if rounded else 1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = CARD
    s.line.color.rgb = BORDER
    s.line.width = Pt(0.75)
    return s


def hline(slide, l, t, w, color=BORDER):
    s = slide.shapes.add_shape(1, l, t, w, Pt(0.75))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def eyebrow(slide, text, l, t, w=Inches(6)):
    return tb(slide, l, t, w, Inches(0.3),
              text.upper(), SANS, 10, VIOLET, bold=True)


def headline(slide, text, l, t, w=Inches(10), size=42):
    return tb(slide, l, t, w, Inches(1.6), text, SERIF, size, WHITE)


def slide_header(slide, section, page):
    tb(slide, P, Inches(0.22), Inches(6.5), Inches(0.28),
       "GENOSAI  ×  OLIVER WYMAN FORUM", SANS, 9, MUTED, bold=True)
    tb(slide, Inches(7.5), Inches(0.22), Inches(5.08), Inches(0.28),
       section, SANS, 9, MUTED, bold=True, align=PP_ALIGN.RIGHT)
    hline(slide, P, Inches(0.58), SW - 2 * P)
    tb(slide, P, SH - Inches(0.42), Inches(8), Inches(0.28),
       "AI Automation Proposal  ·  Oliver Wyman Forum  ·  2026",
       SANS, 8, DIM)
    tb(slide, Inches(9.5), SH - Inches(0.42), Inches(3.08), Inches(0.28),
       f"{page:02d} / {TOTAL}", SANS, 8, DIM, align=PP_ALIGN.RIGHT)


def bullet(slide, dot_char, title, desc, l, t,
           col_w=Inches(5.0), title_size=13, desc_size=12):
    tb(slide, l, t, Inches(0.25), Inches(0.26),
       dot_char, SANS, 13, VIOLET, bold=True)
    tb(slide, l + Inches(0.3), t, col_w - Inches(0.3), Inches(0.26),
       title, SANS, title_size, WHITE, bold=True)
    tb(slide, l + Inches(0.3), t + Inches(0.27), col_w - Inches(0.3), Inches(0.38),
       desc, SANS, desc_size, BODY)
    return t + Inches(0.72)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH

s1 = new_slide(prs)

accent = s1.shapes.add_shape(1, Inches(0), Inches(0), SW, Pt(3))
accent.fill.solid()
accent.fill.fore_color.rgb = VIOLET
accent.line.fill.background()

tb(s1, P, Inches(1.0), Inches(5.5), Inches(0.35),
   "PREPARED FOR", SANS, 10, MUTED, bold=True)
hline(s1, P, Inches(1.42), Inches(5.5), BORDER)

tb(s1, P, Inches(1.6), Inches(9.5), Inches(1.4),
   "Oliver Wyman Forum", SERIF, 64, WHITE)

tb(s1, P, Inches(3.1), Inches(6.5), Inches(0.32),
   "Executive Leadership Communities  ·  Think Tank  ·  Global Research",
   SANS, 13, BODY)

tb(s1, P, Inches(3.6), Inches(5.5), Inches(0.28),
   "New York  ·  Global", SANS, 11, MUTED)

hline(s1, P, Inches(4.06), Inches(5.5), BORDER)

tb(s1, P, Inches(4.26), Inches(7), Inches(0.28),
   "AI AUTOMATION PROPOSAL", SANS, 10, VIOLET, bold=True)
tb_lines(s1, P, Inches(4.66), Inches(6.5), Inches(1.2),
         ["An AI automation layer for the Forum —",
          "built to manage executive community relationships,",
          "distribute research, and coordinate convenings",
          "at the scale of global market capitalization."],
         SANS, 13, SUB)

card(s1, Inches(9.0), Inches(1.2), Inches(3.58), Inches(2.5), rounded=True)
tb(s1, Inches(9.25), Inches(1.5), Inches(3.0), Inches(0.28),
   "PREPARED BY", SANS, 9, MUTED, bold=True)
tb(s1, Inches(9.25), Inches(1.88), Inches(3.0), Inches(0.32),
   "GenosAI Tech", SERIF, 20, WHITE)
tb(s1, Inches(9.25), Inches(2.3), Inches(3.0), Inches(0.26),
   "Rohan Malik  ·  Founder & CEO", SANS, 11, BODY)
tb(s1, Inches(9.25), Inches(2.65), Inches(3.0), Inches(0.26),
   "hello@genosai.tech  ·  www.genosai.tech", SANS, 10, MUTED)

tb(s1, P, SH - Inches(0.65), Inches(4), Inches(0.45),
   "GenosAI", SERIF, 22, WHITE)
tb(s1, P + Inches(1.65), SH - Inches(0.62), Inches(3), Inches(0.3),
   "v1.0  ·  2026  ·  Confidential", SANS, 9, DIM)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE FORUM AT A GLANCE
# (Pitch track — no website audit; replace with org overview)
# ══════════════════════════════════════════════════════════════════════════════
s2 = new_slide(prs)
slide_header(s2, "01  ·  THE FORUM", 2)

eyebrow(s2, "THE ORGANIZATION AT A GLANCE", P, Inches(0.85))
headline(s2, "Scale, Mandate, and Programs", P, Inches(1.18), size=38)
tb(s2, P, Inches(2.2), Inches(8.5), Inches(0.42),
   "Oliver Wyman Forum convenes the most senior leadership communities in business. "
   "The programs span CEO and CFO agendas representing double-digit percentages of global market "
   "capitalization — and the operational complexity behind managing those relationships is significant.",
   SANS, 13, SUB)

hline(s2, P, Inches(2.82), SW - 2 * P, BORDER)

# Stat cards row
stats = [
    ("CEO AGENDA 2026",    "10%",    "of global market cap"),
    ("CFO AGENDA 2026",    "12%",    "of global market cap"),
    ("CONSUMER SURVEY",    "300K+",  "respondents annually"),
    ("UNCHARTED",          "Weekly", "data-driven insights series"),
]
sw_each = Inches(2.8)
gap     = Inches(0.19)
sx      = P

for label, value, sub in stats:
    card(s2, sx, Inches(3.1), sw_each, Inches(2.0), rounded=True)
    tb(s2, sx + Inches(0.25), Inches(3.28), sw_each - Inches(0.4), Inches(0.28),
       label, SANS, 9, MUTED, bold=True)
    tb(s2, sx + Inches(0.25), Inches(3.65), sw_each - Inches(0.4), Inches(0.62),
       value, SERIF, 30, WHITE)
    tb(s2, sx + Inches(0.25), Inches(4.38), sw_each - Inches(0.4), Inches(0.45),
       sub, SANS, 11, BODY)
    sx += sw_each + gap

# Community programs
tb(s2, P, Inches(5.4), Inches(11), Inches(0.28),
   "KEY PROGRAMS", SANS, 9, MUTED, bold=True)
hline(s2, P, Inches(5.72), SW - 2 * P, BORDER)
programs_text = (
   "Leadership Reimagined  ·  CEO Agenda  ·  CFO Agenda  ·  Performance Reimagined  ·  "
   "Equal Societies  ·  Business in Society  ·  Markets & Mobility  ·  Global Economic Change"
)
tb(s2, P, Inches(5.82), SW - 2 * P, Inches(0.42),
   programs_text, SANS, 12, BODY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — VISION
# ══════════════════════════════════════════════════════════════════════════════
s3 = new_slide(prs)
slide_header(s3, "02  ·  VISION", 3)

eyebrow(s3, "THE OPPORTUNITY", P, Inches(0.85))
headline(s3, "An AI Automation Layer\nfor the Forum's Operations", P, Inches(1.18), size=38)

card(s3, P, Inches(2.85), Inches(7.5), Inches(1.25), rounded=True)
tb(s3, P + Inches(0.4), Inches(2.98), Inches(7.0), Inches(0.4),
   "“", SERIF, 32, VIOLET)
tb(s3, P + Inches(0.75), Inches(3.02), Inches(6.5), Inches(0.88),
   "The Forum publishes on AI’s transformative impact on business. The automation layer "
   "for managing the relationships and research distribution behind that work is the natural next step.",
   SERIF, 15, WHITE, italic=True)

col_w = Inches(2.8)
cols  = [P, P + col_w + Inches(0.15), P + 2*(col_w + Inches(0.15)), P + 3*(col_w + Inches(0.15))]
pillars = [
    ("01", "Community\nmanagement",
     "Automated touchpoints for every CEO and CFO community member."),
    ("02", "Research\ndistribution",
     "Reports and insights delivered to the right audience, automatically."),
    ("03", "Event\ncoordination",
     "Convenings managed end-to-end without manual admin overhead."),
    ("04", "Relationship\ncontinuity",
     "No C-suite relationship goes cold because someone missed a follow-up."),
]
for (num, title, desc), x in zip(pillars, cols):
    card(s3, x, Inches(4.35), col_w, Inches(2.65), rounded=True)
    tb(s3, x + Inches(0.25), Inches(4.55), Inches(0.6), Inches(0.42),
       num, SERIF, 22, VIOLET)
    tb(s3, x + Inches(0.25), Inches(5.05), col_w - Inches(0.4), Inches(0.65),
       title, SERIF, 16, WHITE)
    tb(s3, x + Inches(0.25), Inches(5.78), col_w - Inches(0.4), Inches(0.88),
       desc, SANS, 11, BODY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SYSTEMS OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
s4 = new_slide(prs)
slide_header(s4, "03  ·  AUTOMATION SYSTEMS", 4)

eyebrow(s4, "WHAT THE AUTOMATION LAYER BRINGS", P, Inches(0.85))
headline(s4, "Five AI Systems for the Forum", P, Inches(1.18), size=36)

systems = [
    ("A", "Community\nManagement AI",   "CEO/CFO relationship automation."),
    ("B", "Research\nDistribution AI",  "Segmented delivery at scale."),
    ("C", "Event &\nConvening AI",      "End-to-end convening ops."),
    ("D", "Member\nQualification AI",   "Leadership Reimagined intake."),
    ("E", "CRM &\nNurture AI",          "No C-suite relationship goes cold."),
]
sw_each = Inches(2.25)
gap     = Inches(0.17)
sx      = P

for letter, name, tagline in systems:
    card(s4, sx, Inches(2.6), sw_each, Inches(4.1), rounded=True)
    badge = s4.shapes.add_shape(1, sx + Inches(0.25), Inches(2.85), Inches(0.55), Inches(0.55))
    badge.fill.solid()
    badge.fill.fore_color.rgb = VIOLET
    badge.line.fill.background()
    tb(s4, sx + Inches(0.25), Inches(2.87), Inches(0.55), Inches(0.48),
       letter, SANS, 16, BG, bold=True, align=PP_ALIGN.CENTER)
    tb(s4, sx + Inches(0.2), Inches(3.55), sw_each - Inches(0.4), Inches(0.75),
       name, SERIF, 15, WHITE)
    tb(s4, sx + Inches(0.2), Inches(4.38), sw_each - Inches(0.4), Inches(0.55),
       tagline, SANS, 12, BODY)
    sx += sw_each + gap


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — MODULE A: EXECUTIVE COMMUNITY MANAGEMENT AI
# ══════════════════════════════════════════════════════════════════════════════
s5 = new_slide(prs)
slide_header(s5, "03A  ·  COMMUNITY MANAGEMENT", 5)

eyebrow(s5, "MODULE A", P, Inches(0.85))
headline(s5, "Executive Community Management AI", P, Inches(1.18), size=34)
tb(s5, P, Inches(2.2), Inches(6.5), Inches(0.42),
   "Every CEO and CFO community member receives timely, personalized touchpoints — "
   "without anyone tracking a spreadsheet to make it happen.",
   SANS, 14, SUB)

hline(s5, P, Inches(2.75), SW - 2 * P, BORDER)

t = Inches(3.0)
for title, desc in [
    ("Onboarding sequences",
     "New community members receive a structured welcome flow with program context and next steps."),
    ("Event invitation automation",
     "Invitations triggered by calendar — personalized to each member's focus area and region."),
    ("Pre-event briefing delivery",
     "Relevant research, agenda, and attendee context delivered automatically before each convening."),
    ("Post-event follow-up",
     "Personalized follow-up within 48 hours of each convening — no manual drafting."),
    ("Renewal and re-engagement",
     "Members approaching inactivity receive a tailored re-engagement sequence before lapse."),
    ("Engagement scoring",
     "AI tracks open rates, event attendance, and content downloads to identify the most active members."),
]:
    t = bullet(s5, "·", title, desc, P, t, Inches(5.5), 13, 12)

card(s5, Inches(7.5), Inches(3.0), Inches(5.0), Inches(3.7), rounded=True)
tb(s5, Inches(7.8), Inches(3.2), Inches(4.5), Inches(0.28),
   "COMMUNITY  ·  AT A GLANCE", SANS, 9, MUTED, bold=True)
tb(s5, Inches(7.8), Inches(3.6), Inches(4.5), Inches(0.42),
   "CEO + CFO community", SERIF, 20, WHITE)
tb(s5, Inches(7.8), Inches(4.08), Inches(4.5), Inches(0.42),
   "managed at scale.", SERIF, 20, VIOLET)

for i, (label, value) in enumerate([
    ("CEO AGENDA", "10% mktcap"), ("CFO AGENDA", "12% mktcap"), ("PROGRAMS", "10+")
]):
    sx = Inches(7.8) + i * Inches(1.58)
    tb(s5, sx, Inches(5.0), Inches(1.5), Inches(0.28), label, SANS, 8, MUTED, bold=True)
    tb(s5, sx, Inches(5.32), Inches(1.5), Inches(0.42), value, SERIF, 16, WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — MODULE B & C: RESEARCH DISTRIBUTION + EVENT AI
# ══════════════════════════════════════════════════════════════════════════════
s6 = new_slide(prs)
slide_header(s6, "03B / 03C  ·  RESEARCH & EVENTS", 6)

eyebrow(s6, "MODULE B", P, Inches(0.85))
headline(s6, "Research Distribution AI", P, Inches(1.18), size=28)
tb(s6, P, Inches(2.05), Inches(5.5), Inches(0.38),
   "Every report and insight reaches the right audience — automatically, every time.",
   SANS, 13, SUB)

t = Inches(2.6)
for title, desc in [
    ("Audience segmentation",
     "Subscribers segmented by industry, region, and stated interests — automated."),
    ("Uncharted delivery",
     "Weekly insights distributed to segmented lists; open-rate triggers follow-up content."),
    ("Report download follow-up",
     "Executives who download a report receive a relevant follow-on piece 7 days later."),
    ("Inactive subscriber re-engagement",
     "Subscribers who haven't opened in 60 days receive a curated re-engagement sequence."),
    ("Cross-program content bridging",
     "AI matches research topics to the right community program for amplification."),
]:
    t = bullet(s6, "·", title, desc, P, t, Inches(5.5), 12, 11)

hline(s6, Inches(6.67) - Pt(0.5), Inches(0.85), Pt(1), BORDER)

CR = Inches(6.85)
eyebrow(s6, "MODULE C", CR, Inches(0.85))
headline(s6, "Event & Convening AI", CR, Inches(1.18), size=28)
tb(s6, CR, Inches(2.05), Inches(5.7), Inches(0.38),
   "Every executive convening coordinated end-to-end — no manual admin.",
   SANS, 13, SUB)

t = Inches(2.6)
for title, desc in [
    ("Registration flow",
     "RSVP management, dietary requirements, and confirmation — fully automated."),
    ("Speaker briefing delivery",
     "Speaker bios, run-of-show, and prep materials sent on a set schedule."),
    ("Attendee reminder sequences",
     "Tiered reminders at 2 weeks, 3 days, and day-of — including practical logistics."),
    ("Post-event summary distribution",
     "Key insights and next steps from each convening delivered to attendees within 24h."),
    ("Feedback collection",
     "Automated post-event survey with structured follow-up for qualitative responses."),
]:
    t = bullet(s6, "·", title, desc, CR, t, Inches(5.7), 12, 11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — MODULE D & E: MEMBER QUALIFICATION + CRM
# ══════════════════════════════════════════════════════════════════════════════
s7 = new_slide(prs)
slide_header(s7, "03D / 03E  ·  QUALIFICATION & CRM", 7)

eyebrow(s7, "MODULE D", P, Inches(0.85))
headline(s7, "Member Qualification AI", P, Inches(1.18), size=28)
tb(s7, P, Inches(2.05), Inches(5.5), Inches(0.38),
   "Every Leadership Reimagined application handled with the same quality — automatically.",
   SANS, 13, SUB)

t = Inches(2.6)
for title, desc in [
    ("Application intake flow",
     "Executives who express interest enter a structured AI qualification sequence."),
    ("Profile qualification",
     "AI assesses organization size, role seniority, industry, and program fit."),
    ("Fit routing",
     "Qualified applicants routed to CEO community, CFO community, or research-only track."),
    ("Nurture for borderline applicants",
     "Not-yet-ready applicants receive a content sequence to develop program fit."),
    ("Application status communication",
     "Applicants receive timely, personalized updates — no one left waiting in silence."),
]:
    t = bullet(s7, "·", title, desc, P, t, Inches(5.5), 12, 11)

CR = Inches(6.85)
eyebrow(s7, "MODULE E", CR, Inches(0.85))
headline(s7, "CRM & Relationship AI", CR, Inches(1.18), size=28)
tb(s7, CR, Inches(2.05), Inches(5.7), Inches(0.38),
   "No C-suite relationship goes cold because someone missed a follow-up.",
   SANS, 13, SUB)

t = Inches(2.6)
for title, desc in [
    ("Relationship health scoring",
     "AI flags members who haven't engaged in 30+ days for proactive outreach."),
    ("Milestone-triggered outreach",
     "Leadership changes, company news, or research publication trigger relevant touchpoints."),
    ("Cross-program referral automation",
     "Members in one program recommended to adjacent programs based on engagement history."),
    ("Advisory and speaker pipeline",
     "High-engagement members identified and nurtured toward advisory or speaker roles."),
    ("Annual review sequences",
     "Automated annual relationship review prompt — keeps high-value connections active."),
]:
    t = bullet(s7, "·", title, desc, CR, t, Inches(5.7), 12, 11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — EXECUTIVE JOURNEY FLOW
# ══════════════════════════════════════════════════════════════════════════════
s8 = new_slide(prs)
slide_header(s8, "04  ·  EXECUTIVE JOURNEY", 8)

eyebrow(s8, "INTEREST  →  MEMBER  →  ENGAGED  →  ADVOCATE", P, Inches(0.85))
headline(s8, "Automated Executive Journey", P, Inches(1.18), size=36)
tb(s8, P, Inches(2.2), Inches(10), Inches(0.38),
   "From first contact to a fully engaged, renewing community member — every step automated.",
   SANS, 13, SUB)

hline(s8, P, Inches(2.72), SW - 2 * P, BORDER)

steps = [
    ("01", "Interest\nSignal",         "Newsletter / report\ndownload / referral"),
    ("02", "Qualification\nFlow",      "AI assesses fit\nfor program track"),
    ("03", "Community\nOnboarding",    "Welcome sequence\nand program brief"),
    ("04", "First\nConvening",         "Invitation, prep\nmaterials, logistics"),
    ("05", "Research\nDelivery",       "Segmented reports\nmatched to interests"),
    ("06", "Engagement\nTracking",     "Opens, attendance,\ncontent downloads"),
    ("07", "Re-engagement\nor Upsell", "Quiet members\nor program upgrade"),
    ("08", "Renewal &\nAdvocacy",      "Annual renewal\nand referral ask"),
]

step_w   = Inches(1.4)
step_gap = Inches(0.12)
sx = P

for num, title, detail in steps:
    card(s8, sx, Inches(3.1), step_w, Inches(3.5), rounded=True)
    tb(s8, sx + Inches(0.12), Inches(3.28), step_w - Inches(0.2), Inches(0.38),
       num, SERIF, 20, VIOLET)
    tb(s8, sx + Inches(0.12), Inches(3.75), step_w - Inches(0.2), Inches(0.65),
       title, SERIF, 13, WHITE)
    tb(s8, sx + Inches(0.12), Inches(4.5), step_w - Inches(0.2), Inches(0.75),
       detail, SANS, 10, BODY)
    sx += step_w + step_gap
    if num != "08":
        tb(s8, sx - step_gap + Inches(0.02), Inches(4.2),
           step_gap + Inches(0.05), Inches(0.28),
           "›", SANS, 14, MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — BENEFITS
# ══════════════════════════════════════════════════════════════════════════════
s9 = new_slide(prs)
slide_header(s9, "05  ·  BENEFITS", 9)

eyebrow(s9, "WHY THIS WORKS FOR THE FORUM", P, Inches(0.85))
headline(s9, "What Changes After Launch", P, Inches(1.18), size=36)
hline(s9, P, Inches(2.72), SW - 2 * P, BORDER)

benefits = [
    ("No missed follow-ups",          "Every post-event and post-download touchpoint triggered automatically."),
    ("Consistent member experience",  "Every community member receives the same quality of engagement — regardless of team bandwidth."),
    ("Research reaches the right desk","Reports delivered to executives whose interests and roles match the topic."),
    ("Events run without admin drag",  "Registration, reminders, and post-event distribution handled end-to-end."),
    ("Community health visibility",    "Engagement scores surface which relationships need attention before they lapse."),
    ("Scales with program growth",     "New programs and new members don’t add linear admin overhead."),
    ("Leadership pipeline automation", "High-engagement members identified and moved toward advisory and speaker roles."),
    ("Institutional memory",           "AI logs every touchpoint, so no relationship context is lost when team members change."),
]

col1, col2 = benefits[:4], benefits[4:]
CL, CR = P, Inches(7.1)
t1, t2 = Inches(2.95), Inches(2.95)

for title, desc in col1:
    card(s9, CL, t1, Inches(5.75), Inches(0.82), rounded=True)
    tb(s9, CL + Inches(0.25), t1 + Inches(0.1), Inches(0.35), Inches(0.3),
       "✓", SANS, 14, VIOLET, bold=True)
    tb(s9, CL + Inches(0.65), t1 + Inches(0.08), Inches(4.8), Inches(0.3),
       title, SANS, 13, WHITE, bold=True)
    tb(s9, CL + Inches(0.65), t1 + Inches(0.4), Inches(4.8), Inches(0.35),
       desc, SANS, 11, BODY)
    t1 += Inches(0.95)

for title, desc in col2:
    card(s9, CR, t2, Inches(5.75), Inches(0.82), rounded=True)
    tb(s9, CR + Inches(0.25), t2 + Inches(0.1), Inches(0.35), Inches(0.3),
       "✓", SANS, 14, VIOLET, bold=True)
    tb(s9, CR + Inches(0.65), t2 + Inches(0.08), Inches(4.8), Inches(0.3),
       title, SANS, 13, WHITE, bold=True)
    tb(s9, CR + Inches(0.65), t2 + Inches(0.4), Inches(4.8), Inches(0.35),
       desc, SANS, 11, BODY)
    t2 += Inches(0.95)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — FUTURE ROADMAP
# ══════════════════════════════════════════════════════════════════════════════
s10 = new_slide(prs)
slide_header(s10, "06  ·  ROADMAP", 10)

eyebrow(s10, "BEYOND THE LAUNCH", P, Inches(0.85))
headline(s10, "Future Scalability", P, Inches(1.18), size=36)
tb(s10, P, Inches(2.2), Inches(10), Inches(0.38),
   "The same automation infrastructure scales as the Forum expands programs and geographies.",
   SANS, 13, SUB)
hline(s10, P, Inches(2.72), SW - 2 * P, BORDER)

future = [
    ("Voice AI for executive outreach",   "AI-assisted outbound calls for high-priority relationship touchpoints."),
    ("Survey data processing AI",         "300K+ Consumer Sentiment responses processed and personalized at scale."),
    ("AI-generated research briefings",   "Custom executive briefings auto-generated from survey and report data."),
    ("Multilingual distribution",         "Research delivered in executive’s preferred language — automated translation and delivery."),
    ("Regional program AI",               "Separate community management flows for EMEA, APAC, and Americas programs."),
    ("Sponsor and partner automation",    "Sponsor relationship management and deliverables tracking — automated."),
    ("Social media distribution AI",      "LinkedIn and X research distribution scheduled and sequenced automatically."),
    ("AI-powered policy advisory",        "Structured briefing generation for government and policy advisory programs."),
]

col1, col2 = future[:4], future[4:]
CL, CR = P, Inches(7.1)
t1, t2 = Inches(2.95), Inches(2.95)

for title, desc in col1:
    t1 = bullet(s10, "·", title, desc, CL, t1, Inches(5.75))
for title, desc in col2:
    t2 = bullet(s10, "·", title, desc, CR, t2, Inches(5.75))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — IMPLEMENTATION
# ══════════════════════════════════════════════════════════════════════════════
s11 = new_slide(prs)
slide_header(s11, "07  ·  IMPLEMENTATION", 11)

eyebrow(s11, "BUILT IN THREE CLEAN PHASES", P, Inches(0.85))
headline(s11, "Implementation Phases", P, Inches(1.18), size=36)
tb(s11, P, Inches(2.2), Inches(10), Inches(0.38),
   "Each phase ships visible value before the next begins — no months of silence.",
   SANS, 13, SUB)
hline(s11, P, Inches(2.72), SW - 2 * P, BORDER)

phases = [
    ("PHASE 01  ·  Week 1",
     "Community Management AI",
     "Executive community management sequences for CEO and CFO Agenda members: "
     "onboarding, event invitations, post-event follow-ups, and renewal triggers live.",
     ["AI trained on Forum programs and member data", "Onboarding sequence deployment",
      "Event invitation and follow-up automation", "Member engagement scoring",
      "Post-event summary delivery", "Go-live and testing"]),
    ("PHASE 02  ·  Weeks 2-3",
     "Research Distribution & Events",
     "Automated research distribution with audience segmentation, and full event "
     "coordination flows for upcoming convenings.",
     ["Subscriber segmentation by interest/region", "Uncharted delivery automation",
      "Report download follow-up sequences", "Convening registration flow",
      "Speaker briefing distribution", "Attendee reminder sequences"]),
    ("PHASE 03  ·  Weeks 4-5",
     "Qualification & CRM AI",
     "Member qualification flow for Leadership Reimagined, and CRM automation "
     "for ongoing relationship health and cross-program referrals.",
     ["Leadership Reimagined intake flow", "Profile qualification AI",
      "Relationship health scoring", "Cross-program referral automation",
      "Annual review sequences", "Training and handover"]),
]

ph_w   = Inches(3.8)
ph_gap = Inches(0.17)
px = P

for phase_label, phase_name, phase_desc, deliverables in phases:
    card(s11, px, Inches(3.0), ph_w, Inches(3.75), rounded=True)
    tb(s11, px + Inches(0.25), Inches(3.15), ph_w - Inches(0.4), Inches(0.28),
       phase_label, SANS, 9, VIOLET, bold=True)
    tb(s11, px + Inches(0.25), Inches(3.5), ph_w - Inches(0.4), Inches(0.42),
       phase_name, SERIF, 15, WHITE)
    tb(s11, px + Inches(0.25), Inches(4.0), ph_w - Inches(0.4), Inches(0.52),
       phase_desc, SANS, 11, BODY)
    t = Inches(4.62)
    for d in deliverables[:4]:
        tb(s11, px + Inches(0.25), t, ph_w - Inches(0.4), Inches(0.26),
           f"›  {d}", SANS, 10, MUTED)
        t += Inches(0.3)
    px += ph_w + ph_gap


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — CLOSING
# ══════════════════════════════════════════════════════════════════════════════
s12 = new_slide(prs)

acc = s12.shapes.add_shape(1, Inches(0), Inches(0), SW, Pt(3))
acc.fill.solid()
acc.fill.fore_color.rgb = VIOLET
acc.line.fill.background()

eyebrow(s12, "WHO IS BEHIND THIS", P, Inches(0.85))
headline(s12, "GenosAI Tech", P, Inches(1.18), size=40)
tb(s12, P, Inches(2.22), Inches(6.5), Inches(0.66),
   "GenosAI builds AI automation systems for research-led organizations, executive communities, "
   "and professional services firms. We specialise in community management AI, research distribution "
   "automation, event coordination, and CRM flows — production-grade from week one.",
   SANS, 13, SUB)

hline(s12, P, Inches(3.02), Inches(6), BORDER)

tb(s12, P, Inches(3.2), Inches(4), Inches(0.28),
   "WHAT WE BUILD", SANS, 9, MUTED, bold=True)
t = Inches(3.58)
for item, desc in [
    ("Community Management AI",    "Onboarding, events, follow-ups, renewals — automated."),
    ("Research Distribution AI",   "Segmented delivery to the right executive, every time."),
    ("Event & Convening AI",       "End-to-end event coordination without admin overhead."),
    ("Member Qualification Flows", "Leadership community intake, routed and nurtured."),
    ("CRM & Relationship AI",      "No senior relationship goes cold between convenings."),
]:
    t = bullet(s12, "◆", item, desc, P, t, Inches(5.5), 12, 11)

card(s12, Inches(7.5), Inches(1.1), Inches(5.05), Inches(5.75), rounded=True)

tb(s12, Inches(7.75), Inches(1.35), Inches(4.6), Inches(0.35),
   "CLOSING NOTE", SANS, 9, MUTED, bold=True)
tb(s12, Inches(7.6), Inches(1.72), Inches(0.4), Inches(0.55),
   "“", SERIF, 28, VIOLET)
tb(s12, Inches(7.95), Inches(1.78), Inches(4.3), Inches(1.0),
   "Built to ensure every executive who engages with the Forum receives a response, "
   "a follow-up, and a reason to stay engaged — without anyone managing it manually.",
   SERIF, 14, WHITE, italic=True)

hline(s12, Inches(7.75), Inches(2.88), Inches(4.45), BORDER)

tb(s12, Inches(7.75), Inches(3.08), Inches(4.5), Inches(0.28),
   "NEXT STEPS", SANS, 9, VIOLET, bold=True)
t = Inches(3.45)
for step, desc in [
    ("01", "15-minute call to confirm scope and integration points."),
    ("02", "Phase 1 kick-off — community management sequences live within 7 days."),
    ("03", "First executive touchpoint automated — results visible from day one."),
]:
    tb(s12, Inches(7.75), t, Inches(0.4), Inches(0.3), step, SERIF, 14, VIOLET)
    tb(s12, Inches(8.2), t + Inches(0.02), Inches(4.15), Inches(0.5),
       desc, SANS, 12, BODY)
    t += Inches(0.62)

hline(s12, Inches(7.75), Inches(5.45), Inches(4.45), BORDER)

tb(s12, Inches(7.75), Inches(5.65), Inches(4.5), Inches(0.28),
   "Rohan Malik  ·  CEO, Genos AI", SANS, 12, WHITE, bold=True)
tb(s12, Inches(7.75), Inches(6.0), Inches(4.5), Inches(0.28),
   "hello@genosai.tech  ·  +91 638-714-2699", SANS, 11, BODY)
tb(s12, Inches(7.75), Inches(6.32), Inches(4.5), Inches(0.28),
   "www.genosai.tech", SANS, 11, MUTED)

tb(s12, P, SH - Inches(0.42), Inches(8), Inches(0.28),
   "Confidential  ·  Oliver Wyman Forum  ·  2026", SANS, 8, DIM)
tb(s12, Inches(9.5), SH - Inches(0.42), Inches(3.08), Inches(0.28),
   f"{TOTAL:02d} / {TOTAL}", SANS, 8, DIM, align=PP_ALIGN.RIGHT)

tb(s12, P, SH - Inches(0.65), Inches(3), Inches(0.4),
   "GenosAI", SERIF, 20, WHITE)


# ── Save ─────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved: {OUT}  ({TOTAL} slides)")
