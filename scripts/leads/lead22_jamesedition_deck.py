"""
Build JamesEdition_Pitch_GenosAI.pptx from scratch.
Track: Pitch-only (website score 8/10 — no website gap analysis).
Design: docs/BRAND_ASSETS.md — #0A0A0F canvas, Instrument Serif headlines, Satoshi body.
12 slides: cover, marketplace overview, vision, systems, 4 module slides,
seller lifecycle flow, benefits, implementation, closing.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).parent.parent.parent
OUT  = ROOT / "output" / "decks" / "JamesEdition_Pitch_GenosAI.pptx"

# ── Brand colors ─────────────────────────────────────────────────────────────
BG     = RGBColor(0x0A, 0x0A, 0x0F)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BODY   = RGBColor(0x9D, 0x9D, 0x9F)
SUB    = RGBColor(0x86, 0x86, 0x8A)
MUTED  = RGBColor(0x6B, 0x6B, 0x6E)
DIM    = RGBColor(0x44, 0x44, 0x47)
VIOLET = RGBColor(0xA7, 0x8B, 0xFA)
CARD   = RGBColor(0x11, 0x11, 0x18)
BORDER = RGBColor(0x22, 0x22, 0x2A)

SERIF = "Instrument Serif"
SANS  = "Satoshi"

SW    = Inches(13.33)
SH    = Inches(7.5)
P     = Inches(0.75)
TOTAL = 12


# ── Helpers ───────────────────────────────────────────────────────────────────

def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def tb(slide, l, t, w, h, text, font=SANS, size=14, color=BODY,
       bold=False, italic=False, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text = text
    r.font.name   = font
    r.font.size   = Pt(size)
    r.font.color.rgb = color
    r.font.bold   = bold
    r.font.italic = italic
    return box


def tb_lines(slide, l, t, w, h, lines, font=SANS, size=14, color=BODY,
             bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name  = font
        r.font.size  = Pt(size)
        r.font.color.rgb = color
        r.font.bold  = bold
    return box


def card(slide, l, t, w, h, rounded=True):
    s = slide.shapes.add_shape(5 if rounded else 1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = CARD
    s.line.color.rgb = BORDER
    s.line.width = Pt(0.75)
    return s


def hline(slide, l, t, w, color=BORDER):
    s = slide.shapes.add_shape(1, l, t, w, Pt(0.75))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def eyebrow(slide, text, l, t, w=Inches(6)):
    return tb(slide, l, t, w, Inches(0.3),
              text.upper(), SANS, 10, VIOLET, bold=True)


def headline(slide, text, l, t, w=Inches(10), size=42):
    return tb(slide, l, t, w, Inches(1.6), text, SERIF, size, WHITE)


def slide_header(slide, section, page):
    tb(slide, P, Inches(0.22), Inches(6), Inches(0.28),
       "GENOSAI  ×  JAMESEDITION", SANS, 9, MUTED, bold=True)
    tb(slide, Inches(7.5), Inches(0.22), Inches(5.08), Inches(0.28),
       section, SANS, 9, MUTED, bold=True, align=PP_ALIGN.RIGHT)
    hline(slide, P, Inches(0.58), SW - 2 * P)
    tb(slide, P, SH - Inches(0.42), Inches(8), Inches(0.28),
       "AI Automation Proposal  ·  JamesEdition  ·  2026", SANS, 8, DIM)
    tb(slide, Inches(9.5), SH - Inches(0.42), Inches(3.08), Inches(0.28),
       f"{page:02d} / {TOTAL}", SANS, 8, DIM, align=PP_ALIGN.RIGHT)


def bullet(slide, dot_char, title, desc, l, t,
           col_w=Inches(5.0), title_size=13, desc_size=12):
    tb(slide, l, t, Inches(0.25), Inches(0.26),
       dot_char, SANS, 13, VIOLET, bold=True)
    tb(slide, l + Inches(0.3), t, col_w - Inches(0.3), Inches(0.26),
       title, SANS, title_size, WHITE, bold=True)
    tb(slide, l + Inches(0.3), t + Inches(0.27), col_w - Inches(0.3), Inches(0.38),
       desc, SANS, desc_size, BODY)
    return t + Inches(0.72)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH

s1 = new_slide(prs)

accent = s1.shapes.add_shape(1, Inches(0), Inches(0), SW, Pt(3))
accent.fill.solid()
accent.fill.fore_color.rgb = VIOLET
accent.line.fill.background()

tb(s1, P, Inches(1.0), Inches(5.5), Inches(0.35),
   "PREPARED FOR", SANS, 10, MUTED, bold=True)
hline(s1, P, Inches(1.42), Inches(5.5), BORDER)

tb(s1, P, Inches(1.6), Inches(8.5), Inches(1.4),
   "JamesEdition", SERIF, 72, WHITE)

tb(s1, P, Inches(3.1), Inches(7), Inches(0.32),
   "Global Luxury Marketplace  ·  Real Estate  ·  Yachts  ·  Jets  ·  Automobiles",
   SANS, 13, BODY)

tb(s1, P, Inches(3.6), Inches(5.5), Inches(0.28),
   "Amsterdam, Netherlands  ·  120+ Countries", SANS, 11, MUTED)

hline(s1, P, Inches(4.06), Inches(5.5), BORDER)

tb(s1, P, Inches(4.26), Inches(7), Inches(0.28),
   "AI AUTOMATION PROPOSAL", SANS, 10, VIOLET, bold=True)
tb_lines(s1, P, Inches(4.66), Inches(6.5), Inches(1.2),
         ["An AI automation layer for JamesEdition —",
          "built to onboard sellers at scale, prevent churn,",
          "trigger upsells, and pre-qualify buyers",
          "across 40,000+ sellers and 120 countries."],
         SANS, 13, SUB)

card(s1, Inches(9.0), Inches(1.2), Inches(3.58), Inches(2.5), rounded=True)
tb(s1, Inches(9.25), Inches(1.5), Inches(3.0), Inches(0.28),
   "PREPARED BY", SANS, 9, MUTED, bold=True)
tb(s1, Inches(9.25), Inches(1.88), Inches(3.0), Inches(0.32),
   "GenosAI Tech", SERIF, 20, WHITE)
tb(s1, Inches(9.25), Inches(2.3), Inches(3.0), Inches(0.26),
   "Rohan Malik  ·  Founder & CEO", SANS, 11, BODY)
tb(s1, Inches(9.25), Inches(2.65), Inches(3.0), Inches(0.26),
   "hello@genosai.tech  ·  www.genosai.tech", SANS, 10, MUTED)

tb(s1, P, SH - Inches(0.65), Inches(4), Inches(0.45),
   "GenosAI", SERIF, 22, WHITE)
tb(s1, P + Inches(1.65), SH - Inches(0.62), Inches(3), Inches(0.3),
   "v1.0  ·  2026  ·  Confidential", SANS, 9, DIM)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — JAMESEDITION AT A GLANCE
# ══════════════════════════════════════════════════════════════════════════════
s2 = new_slide(prs)
slide_header(s2, "01  ·  THE MARKETPLACE", 2)

eyebrow(s2, "THE PLATFORM AT A GLANCE", P, Inches(0.85))
headline(s2, "Global Scale. Lean Team. High Stakes.", P, Inches(1.18), size=38)
tb(s2, P, Inches(2.2), Inches(8.5), Inches(0.42),
   "JamesEdition operates one of the world's largest luxury marketplaces — "
   "770,000+ listings across real estate, yachts, jets, and automobiles — "
   "with a 35-person team. The subscription revenue model means seller success "
   "and retention are directly tied to every automation decision.",
   SANS, 13, SUB)

hline(s2, P, Inches(2.82), SW - 2 * P, BORDER)

stats = [
    ("LISTINGS",  "770K+",  "active on platform"),
    ("SELLERS",   "40K+",   "trusted across 120 countries"),
    ("TEAM",      "35",     "employees worldwide"),
    ("SUBSCRIPTIONS", "$299–$1,199", "per seller per month"),
]
sw_each = Inches(2.8)
gap     = Inches(0.19)
sx      = P

for label, value, sub in stats:
    card(s2, sx, Inches(3.1), sw_each, Inches(2.0), rounded=True)
    tb(s2, sx + Inches(0.25), Inches(3.28), sw_each - Inches(0.4), Inches(0.28),
       label, SANS, 9, MUTED, bold=True)
    tb(s2, sx + Inches(0.25), Inches(3.65), sw_each - Inches(0.4), Inches(0.62),
       value, SERIF, 26, WHITE)
    tb(s2, sx + Inches(0.25), Inches(4.38), sw_each - Inches(0.4), Inches(0.45),
       sub, SANS, 11, BODY)
    sx += sw_each + gap

tb(s2, P, Inches(5.4), Inches(11), Inches(0.28),
   "ASSET CATEGORIES", SANS, 9, MUTED, bold=True)
hline(s2, P, Inches(5.72), SW - 2 * P, BORDER)
tb(s2, P, Inches(5.82), SW - 2 * P, Inches(0.42),
   "Real Estate  ·  Automobiles  ·  Yachts  ·  Private Jets  ·  Watches  ·  "
   "Motorcycles  ·  Helicopters  ·  Jewellery  ·  Collectibles",
   SANS, 12, BODY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — VISION
# ══════════════════════════════════════════════════════════════════════════════
s3 = new_slide(prs)
slide_header(s3, "02  ·  VISION", 3)

eyebrow(s3, "THE OPPORTUNITY", P, Inches(0.85))
headline(s3, "The Seller Success Layer\nThat Makes 35 People Feel Like 350", P, Inches(1.18), size=36)

card(s3, P, Inches(2.85), Inches(7.5), Inches(1.25), rounded=True)
tb(s3, P + Inches(0.4), Inches(2.98), Inches(7.0), Inches(0.4),
   "“", SERIF, 32, VIOLET)
tb(s3, P + Inches(0.75), Inches(3.02), Inches(6.5), Inches(0.88),
   "The marketplace handles discovery. AI handles the seller relationship — "
   "from day one onboarding through to renewal, upsell, and re-engagement. "
   "No headcount added. No seller left to figure it out alone.",
   SERIF, 15, WHITE, italic=True)

col_w = Inches(2.8)
cols  = [P, P + col_w + Inches(0.15), P + 2*(col_w + Inches(0.15)), P + 3*(col_w + Inches(0.15))]
pillars = [
    ("01", "Seller\nonboarding",
     "Every new seller gets a structured setup — listing guidance, tips, milestones."),
    ("02", "Churn\nprevention",
     "Sellers who go quiet are flagged and re-engaged before they cancel."),
    ("03", "Revenue\nupsell",
     "Basic plan sellers showing strong lead volume prompted to upgrade automatically."),
    ("04", "Buyer\npre-qualification",
     "Luxury inquiries pre-qualified before reaching the seller — no time wasted."),
]
for (num, title, desc), x in zip(pillars, cols):
    card(s3, x, Inches(4.35), col_w, Inches(2.65), rounded=True)
    tb(s3, x + Inches(0.25), Inches(4.55), Inches(0.6), Inches(0.42),
       num, SERIF, 22, VIOLET)
    tb(s3, x + Inches(0.25), Inches(5.05), col_w - Inches(0.4), Inches(0.65),
       title, SERIF, 16, WHITE)
    tb(s3, x + Inches(0.25), Inches(5.78), col_w - Inches(0.4), Inches(0.88),
       desc, SANS, 11, BODY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SYSTEMS OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
s4 = new_slide(prs)
slide_header(s4, "03  ·  AUTOMATION SYSTEMS", 4)

eyebrow(s4, "WHAT THE AUTOMATION LAYER BRINGS", P, Inches(0.85))
headline(s4, "Five AI Systems for JamesEdition", P, Inches(1.18), size=36)

systems = [
    ("A", "Seller\nOnboarding AI",    "Structured setup for every new seller."),
    ("B", "Churn\nPrevention AI",     "Re-engage before cancellation."),
    ("C", "Upsell\nTrigger AI",       "Basic to premium — automated."),
    ("D", "Buyer\nQualification AI",  "Pre-qualify $500K+ inquiries."),
    ("E", "Support\nAutomation",      "Tier-1 queries handled by AI."),
]
sw_each = Inches(2.25)
gap     = Inches(0.17)
sx      = P

for letter, name, tagline in systems:
    card(s4, sx, Inches(2.6), sw_each, Inches(4.1), rounded=True)
    badge = s4.shapes.add_shape(1, sx + Inches(0.25), Inches(2.85), Inches(0.55), Inches(0.55))
    badge.fill.solid()
    badge.fill.fore_color.rgb = VIOLET
    badge.line.fill.background()
    tb(s4, sx + Inches(0.25), Inches(2.87), Inches(0.55), Inches(0.48),
       letter, SANS, 16, BG, bold=True, align=PP_ALIGN.CENTER)
    tb(s4, sx + Inches(0.2), Inches(3.55), sw_each - Inches(0.4), Inches(0.75),
       name, SERIF, 15, WHITE)
    tb(s4, sx + Inches(0.2), Inches(4.38), sw_each - Inches(0.4), Inches(0.55),
       tagline, SANS, 12, BODY)
    sx += sw_each + gap


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — MODULE A: SELLER ONBOARDING AI
# ══════════════════════════════════════════════════════════════════════════════
s5 = new_slide(prs)
slide_header(s5, "03A  ·  SELLER ONBOARDING", 5)

eyebrow(s5, "MODULE A", P, Inches(0.85))
headline(s5, "Seller Onboarding AI", P, Inches(1.18), size=36)
tb(s5, P, Inches(2.2), Inches(6.5), Inches(0.42),
   "Every new JamesEdition seller gets a structured, high-quality setup experience — "
   "whether you add 10 sellers this week or 1,000. No human required.",
   SANS, 14, SUB)

hline(s5, P, Inches(2.75), SW - 2 * P, BORDER)

t = Inches(3.0)
for title, desc in [
    ("Day-one welcome sequence",
     "New sellers receive an immediate structured welcome with platform overview, "
     "listing best practices, and a 7-day activation checklist."),
    ("Listing optimisation guidance",
     "AI reviews the seller's first listing and sends specific improvement prompts — "
     "photos, pricing, description completeness — before it goes live."),
    ("Performance benchmark delivery",
     "At day 7 and day 30, seller receives a personalised benchmark: "
     "how their listings compare to similar active sellers in their category and region."),
    ("Milestone check-ins",
     "First inquiry, first saved listing, first contact — each milestone triggers "
     "a congratulatory nudge and a next-step prompt to build momentum."),
    ("Handover trigger",
     "Sellers who complete onboarding without receiving an inquiry in 30 days "
     "are flagged to the Customer Success team with full context."),
]:
    t = bullet(s5, "·", title, desc, P, t, Inches(5.5), 13, 12)

card(s5, Inches(7.5), Inches(3.0), Inches(5.0), Inches(3.7), rounded=True)
tb(s5, Inches(7.8), Inches(3.2), Inches(4.5), Inches(0.28),
   "THE ONBOARDING PROBLEM", SANS, 9, MUTED, bold=True)
tb(s5, Inches(7.8), Inches(3.6), Inches(4.5), Inches(0.55),
   "Sellers who leave in month one", SERIF, 18, WHITE)
tb(s5, Inches(7.8), Inches(4.22), Inches(4.5), Inches(0.55),
   "leave because nobody helped them get started.", SERIF, 16, VIOLET)
tb(s5, Inches(7.8), Inches(4.9), Inches(4.5), Inches(0.65),
   "At 40,000+ sellers and 35 employees, structured onboarding "
   "only happens if it's automated.", SANS, 11, BODY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — MODULE B & C: CHURN PREVENTION + UPSELL
# ══════════════════════════════════════════════════════════════════════════════
s6 = new_slide(prs)
slide_header(s6, "03B / 03C  ·  CHURN & UPSELL", 6)

eyebrow(s6, "MODULE B", P, Inches(0.85))
headline(s6, "Churn Prevention AI", P, Inches(1.18), size=28)
tb(s6, P, Inches(2.05), Inches(5.5), Inches(0.38),
   "Sellers who are about to cancel — identified and re-engaged before they do.",
   SANS, 13, SUB)

t = Inches(2.6)
for title, desc in [
    ("Engagement scoring",
     "AI tracks login frequency, listing updates, and inquiry response rates per seller."),
    ("Early churn signal detection",
     "Sellers inactive for 14+ days or with zero leads in 30 days flagged automatically."),
    ("Re-engagement sequence",
     "Flagged sellers receive a targeted sequence: listing optimisation prompt, "
     "market data relevant to their category, and a direct offer to speak with support."),
    ("Cancellation intent intercept",
     "If a seller initiates cancellation, AI triggers an immediate value-delivery sequence "
     "before the cancellation completes."),
    ("Renewal reminder sequences",
     "Annual and monthly subscribers receive tiered renewal reminders at 30, 7, and 1 day."),
]:
    t = bullet(s6, "·", title, desc, P, t, Inches(5.5), 12, 11)

hline(s6, Inches(6.67) - Pt(0.5), Inches(0.85), Pt(1), BORDER)

CR = Inches(6.85)
eyebrow(s6, "MODULE C", CR, Inches(0.85))
headline(s6, "Upsell Trigger AI", CR, Inches(1.18), size=28)
tb(s6, CR, Inches(2.05), Inches(5.7), Inches(0.38),
   "Basic plan sellers ready to upgrade — identified and prompted automatically.",
   SANS, 13, SUB)

t = Inches(2.6)
for title, desc in [
    ("Lead volume threshold trigger",
     "Sellers who hit 5+ inquiries on a basic plan receive an upgrade prompt "
     "with specific ROI context from their own data."),
    ("Featured placement prompt",
     "Sellers in high-competition categories shown data on how premium placement "
     "performs against their current visibility."),
    ("Seasonal upgrade windows",
     "Category-specific upgrade prompts timed to high-demand periods — "
     "yacht season, spring real estate market, auto show season."),
    ("Upgrade abandonment recovery",
     "Sellers who start but don't complete an upgrade receive a follow-up "
     "within 24 hours with a specific reason to complete."),
    ("Annual plan conversion",
     "Monthly subscribers showing consistent engagement prompted to switch "
     "to annual billing with a calculated saving."),
]:
    t = bullet(s6, "·", title, desc, CR, t, Inches(5.7), 12, 11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — MODULE D & E: BUYER QUALIFICATION + SUPPORT
# ══════════════════════════════════════════════════════════════════════════════
s7 = new_slide(prs)
slide_header(s7, "03D / 03E  ·  BUYER QUAL & SUPPORT", 7)

eyebrow(s7, "MODULE D", P, Inches(0.85))
headline(s7, "Buyer Pre-Qualification AI", P, Inches(1.18), size=28)
tb(s7, P, Inches(2.05), Inches(5.5), Inches(0.38),
   "Only qualified buyers reach the seller — protecting seller time on $500K+ assets.",
   SANS, 13, SUB)

t = Inches(2.6)
for title, desc in [
    ("Budget confirmation",
     "AI confirms the buyer's available budget before connecting to the seller."),
    ("Finance vs cash qualification",
     "Buyer clarifies purchase method — mortgage pre-approval, cash, or investment fund."),
    ("Timeline qualification",
     "Buyers with a 3-month purchase timeline prioritised over speculative browsers."),
    ("Seriousness scoring",
     "Profile completeness, inquiry history, and response speed used to score each buyer."),
    ("Seller notification with context",
     "Seller receives a pre-qualified lead summary — budget, timeline, purchase method — "
     "not just a raw inquiry email."),
]:
    t = bullet(s7, "·", title, desc, P, t, Inches(5.5), 12, 11)

CR = Inches(6.85)
eyebrow(s7, "MODULE E", CR, Inches(0.85))
headline(s7, "Tier-1 Support Automation", CR, Inches(1.18), size=28)
tb(s7, CR, Inches(2.05), Inches(5.7), Inches(0.38),
   "21,000+ trusted sellers generate support volume that 35 people can't absorb manually.",
   SANS, 13, SUB)

t = Inches(2.6)
for title, desc in [
    ("Listing issue resolution",
     "AI diagnoses and resolves common listing problems (photo rejection, "
     "missing fields, pricing errors) without human escalation."),
    ("Account and billing queries",
     "Subscription status, invoice requests, payment failures — handled by AI "
     "with Pipedrive CRM integration for context."),
    ("Platform how-to support",
     "Dashboard questions, search optimisation, feature usage — AI answers "
     "from the help centre knowledge base."),
    ("Escalation routing",
     "Queries that require human handling are routed to the right team member "
     "with full context — no cold handover."),
    ("Multilingual support",
     "AI responds in the seller's language — critical for a 120-country marketplace."),
]:
    t = bullet(s7, "·", title, desc, CR, t, Inches(5.7), 12, 11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — SELLER LIFECYCLE FLOW
# ══════════════════════════════════════════════════════════════════════════════
s8 = new_slide(prs)
slide_header(s8, "04  ·  SELLER LIFECYCLE", 8)

eyebrow(s8, "SIGN-UP  →  ONBOARDED  →  ENGAGED  →  RENEWED  →  UPGRADED", P, Inches(0.85))
headline(s8, "Automated Seller Lifecycle", P, Inches(1.18), size=34)
tb(s8, P, Inches(2.2), Inches(10), Inches(0.38),
   "From the moment a seller signs up to a paying, active, renewing subscriber — every stage automated.",
   SANS, 13, SUB)

hline(s8, P, Inches(2.72), SW - 2 * P, BORDER)

steps = [
    ("01", "Seller\nSigns Up",        "Basic or premium\nplan selected"),
    ("02", "Onboarding\nSequence",    "Day-1 welcome,\nlisting guidance"),
    ("03", "First Listing\nLive",     "AI reviews and\noptimisation sent"),
    ("04", "First Inquiry\nReceived", "Milestone trigger +\nnext-step prompt"),
    ("05", "Engagement\nTracking",    "Login frequency,\nlead response rate"),
    ("06", "Churn Signal\nDetected",  "Re-engagement\nsequence triggers"),
    ("07", "Upsell\nWindow",          "Lead volume hits\nthreshold — upgrade"),
    ("08", "Renewal\nor Upgrade",     "Renewal reminder\nor plan upgrade"),
]

step_w   = Inches(1.4)
step_gap = Inches(0.12)
sx = P

for num, title, detail in steps:
    card(s8, sx, Inches(3.1), step_w, Inches(3.5), rounded=True)
    tb(s8, sx + Inches(0.12), Inches(3.28), step_w - Inches(0.2), Inches(0.38),
       num, SERIF, 20, VIOLET)
    tb(s8, sx + Inches(0.12), Inches(3.75), step_w - Inches(0.2), Inches(0.65),
       title, SERIF, 13, WHITE)
    tb(s8, sx + Inches(0.12), Inches(4.5), step_w - Inches(0.2), Inches(0.75),
       detail, SANS, 10, BODY)
    sx += step_w + step_gap
    if num != "08":
        tb(s8, sx - step_gap + Inches(0.02), Inches(4.2),
           step_gap + Inches(0.05), Inches(0.28),
           "›", SANS, 14, MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — BENEFITS
# ══════════════════════════════════════════════════════════════════════════════
s9 = new_slide(prs)
slide_header(s9, "05  ·  BENEFITS", 9)

eyebrow(s9, "WHY THIS WORKS FOR JAMESEDITION", P, Inches(0.85))
headline(s9, "What Changes After Launch", P, Inches(1.18), size=36)
hline(s9, P, Inches(2.72), SW - 2 * P, BORDER)

benefits = [
    ("Lower month-one churn",         "Sellers who go through structured onboarding cancel less — they understand the platform's value."),
    ("Subscription revenue protected", "Churn prevention AI catches at-risk sellers before they cancel — not after."),
    ("Upsell revenue unlocked",        "High-performing basic plan sellers converted to premium automatically — no sales rep needed."),
    ("Seller success at 35-person scale", "40,000+ sellers receive consistent, high-quality support without headcount increase."),
    ("Qualified buyers only",          "Sellers on $500K+ listings spend time on serious buyers, not speculative inquiries."),
    ("Support load reduced",           "Tier-1 queries handled by AI — team focuses on complex, high-value seller issues."),
    ("Multilingual by default",        "120-country seller base supported in their own language without hiring translators."),
    ("Pipedrive CRM enriched",         "Every seller interaction logged back to Pipedrive — full pipeline visibility for the team."),
]

col1, col2 = benefits[:4], benefits[4:]
CL, CR = P, Inches(7.1)
t1, t2 = Inches(2.95), Inches(2.95)

for title, desc in col1:
    card(s9, CL, t1, Inches(5.75), Inches(0.82), rounded=True)
    tb(s9, CL + Inches(0.25), t1 + Inches(0.1), Inches(0.35), Inches(0.3),
       "✓", SANS, 14, VIOLET, bold=True)
    tb(s9, CL + Inches(0.65), t1 + Inches(0.08), Inches(4.8), Inches(0.3),
       title, SANS, 13, WHITE, bold=True)
    tb(s9, CL + Inches(0.65), t1 + Inches(0.4), Inches(4.8), Inches(0.35),
       desc, SANS, 11, BODY)
    t1 += Inches(0.95)

for title, desc in col2:
    card(s9, CR, t2, Inches(5.75), Inches(0.82), rounded=True)
    tb(s9, CR + Inches(0.25), t2 + Inches(0.1), Inches(0.35), Inches(0.3),
       "✓", SANS, 14, VIOLET, bold=True)
    tb(s9, CR + Inches(0.65), t2 + Inches(0.08), Inches(4.8), Inches(0.3),
       title, SANS, 13, WHITE, bold=True)
    tb(s9, CR + Inches(0.65), t2 + Inches(0.4), Inches(4.8), Inches(0.35),
       desc, SANS, 11, BODY)
    t2 += Inches(0.95)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — FUTURE ROADMAP
# ══════════════════════════════════════════════════════════════════════════════
s10 = new_slide(prs)
slide_header(s10, "06  ·  ROADMAP", 10)

eyebrow(s10, "BEYOND THE LAUNCH", P, Inches(0.85))
headline(s10, "Future Scalability", P, Inches(1.18), size=36)
tb(s10, P, Inches(2.2), Inches(10), Inches(0.38),
   "The same automation infrastructure scales across every category and region as JamesEdition grows.",
   SANS, 13, SUB)
hline(s10, P, Inches(2.72), SW - 2 * P, BORDER)

future = [
    ("Voice AI for seller outreach",     "Outbound calls to at-risk sellers before cancellation — AI-led, human-feeling."),
    ("Category-specific onboarding",     "Separate onboarding sequences for real estate vs yacht vs jet sellers."),
    ("Editorial content personalisation","AI-matched Journal articles and market insights delivered to sellers by category."),
    ("AI-generated listing copy",        "AI drafts listing descriptions from seller-provided specs — reduces friction to list."),
    ("Buyer matching AI",                "Buyers saved-search criteria matched to new listings and notified automatically."),
    ("Broker partnership automation",    "New broker partnerships onboarded via structured sequence — no manual setup."),
    ("Investor alert sequences",         "High-value buyers (HNWIs) matched to new off-market listings in their criteria."),
    ("Marketplace analytics AI",         "Natural language queries on platform performance delivered to Eric's team weekly."),
]

col1, col2 = future[:4], future[4:]
CL, CR = P, Inches(7.1)
t1, t2 = Inches(2.95), Inches(2.95)

for title, desc in col1:
    t1 = bullet(s10, "·", title, desc, CL, t1, Inches(5.75))
for title, desc in col2:
    t2 = bullet(s10, "·", title, desc, CR, t2, Inches(5.75))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — IMPLEMENTATION
# ══════════════════════════════════════════════════════════════════════════════
s11 = new_slide(prs)
slide_header(s11, "07  ·  IMPLEMENTATION", 11)

eyebrow(s11, "BUILT IN THREE CLEAN PHASES", P, Inches(0.85))
headline(s11, "Implementation Phases", P, Inches(1.18), size=36)
tb(s11, P, Inches(2.2), Inches(10), Inches(0.38),
   "Each phase protects revenue immediately — no months of silence before results.",
   SANS, 13, SUB)
hline(s11, P, Inches(2.72), SW - 2 * P, BORDER)

phases = [
    ("PHASE 01  ·  Week 1",
     "Seller Onboarding AI",
     "Structured onboarding sequences for all new sellers — day-one welcome, "
     "listing optimisation guidance, performance benchmarks, and milestone triggers.",
     ["Onboarding sequence design", "Listing optimisation prompt flow",
      "Day-7 and day-30 benchmark delivery", "Milestone trigger configuration",
      "Pipedrive CRM integration", "Go-live for all new sign-ups"]),
    ("PHASE 02  ·  Weeks 2-3",
     "Churn Prevention + Upsell AI",
     "Churn risk scoring for all active sellers, automated re-engagement sequences, "
     "and upsell triggers for high-performing basic plan subscribers.",
     ["Engagement scoring model", "Churn signal threshold configuration",
      "Re-engagement sequence deployment", "Upsell trigger logic",
      "Cancellation intercept flow", "Renewal reminder sequences"]),
    ("PHASE 03  ·  Weeks 4-5",
     "Buyer Qualification + Support AI",
     "Buyer pre-qualification for luxury inquiries and tier-1 support automation "
     "integrated with the JamesEdition help centre and Pipedrive.",
     ["Buyer qualification flow", "Budget and timeline capture",
      "Seller notification enrichment", "Support AI knowledge base setup",
      "Escalation routing logic", "Multilingual response configuration"]),
]

ph_w   = Inches(3.8)
ph_gap = Inches(0.17)
px = P

for phase_label, phase_name, phase_desc, deliverables in phases:
    card(s11, px, Inches(3.0), ph_w, Inches(3.75), rounded=True)
    tb(s11, px + Inches(0.25), Inches(3.15), ph_w - Inches(0.4), Inches(0.28),
       phase_label, SANS, 9, VIOLET, bold=True)
    tb(s11, px + Inches(0.25), Inches(3.5), ph_w - Inches(0.4), Inches(0.42),
       phase_name, SERIF, 15, WHITE)
    tb(s11, px + Inches(0.25), Inches(4.0), ph_w - Inches(0.4), Inches(0.52),
       phase_desc, SANS, 11, BODY)
    t = Inches(4.62)
    for d in deliverables[:4]:
        tb(s11, px + Inches(0.25), t, ph_w - Inches(0.4), Inches(0.26),
           f"›  {d}", SANS, 10, MUTED)
        t += Inches(0.3)
    px += ph_w + ph_gap


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — CLOSING
# ══════════════════════════════════════════════════════════════════════════════
s12 = new_slide(prs)

acc = s12.shapes.add_shape(1, Inches(0), Inches(0), SW, Pt(3))
acc.fill.solid()
acc.fill.fore_color.rgb = VIOLET
acc.line.fill.background()

eyebrow(s12, "WHO IS BEHIND THIS", P, Inches(0.85))
headline(s12, "GenosAI Tech", P, Inches(1.18), size=40)
tb(s12, P, Inches(2.22), Inches(6.5), Inches(0.66),
   "GenosAI builds AI automation systems for marketplace businesses, subscription platforms, "
   "and high-volume B2B operations. We specialise in seller onboarding flows, churn prevention AI, "
   "upsell trigger systems, and support automation — production-grade from week one.",
   SANS, 13, SUB)

hline(s12, P, Inches(3.02), Inches(6), BORDER)

tb(s12, P, Inches(3.2), Inches(4), Inches(0.28),
   "WHAT WE BUILD", SANS, 9, MUTED, bold=True)
t = Inches(3.58)
for item, desc in [
    ("Seller Onboarding AI",      "Structured setup for every new seller — at any volume."),
    ("Churn Prevention AI",       "At-risk sellers identified and re-engaged automatically."),
    ("Upsell Trigger Systems",    "Revenue upgrade prompts triggered by seller behaviour."),
    ("Buyer Qualification AI",    "Luxury inquiries pre-qualified before reaching sellers."),
    ("Support Automation",        "Tier-1 queries handled by AI with CRM integration."),
]:
    t = bullet(s12, "◆", item, desc, P, t, Inches(5.5), 12, 11)

card(s12, Inches(7.5), Inches(1.1), Inches(5.05), Inches(5.75), rounded=True)

tb(s12, Inches(7.75), Inches(1.35), Inches(4.6), Inches(0.35),
   "CLOSING NOTE", SANS, 9, MUTED, bold=True)
tb(s12, Inches(7.6), Inches(1.72), Inches(0.4), Inches(0.55),
   "“", SERIF, 28, VIOLET)
tb(s12, Inches(7.95), Inches(1.78), Inches(4.3), Inches(1.0),
   "Built to make 35 people feel like 350 — every seller onboarded, every "
   "churn signal caught, every upsell triggered, without adding headcount.",
   SERIF, 14, WHITE, italic=True)

hline(s12, Inches(7.75), Inches(2.88), Inches(4.45), BORDER)

tb(s12, Inches(7.75), Inches(3.08), Inches(4.5), Inches(0.28),
   "NEXT STEPS", SANS, 9, VIOLET, bold=True)
t = Inches(3.45)
for step, desc in [
    ("01", "15-minute call to confirm scope and integration points."),
    ("02", "Phase 1 kick-off — seller onboarding AI live within 7 days."),
    ("03", "First new seller goes through the automated flow — visible from day one."),
]:
    tb(s12, Inches(7.75), t, Inches(0.4), Inches(0.3), step, SERIF, 14, VIOLET)
    tb(s12, Inches(8.2), t + Inches(0.02), Inches(4.15), Inches(0.5),
       desc, SANS, 12, BODY)
    t += Inches(0.62)

hline(s12, Inches(7.75), Inches(5.45), Inches(4.45), BORDER)

tb(s12, Inches(7.75), Inches(5.65), Inches(4.5), Inches(0.28),
   "Rohan Malik  ·  CEO, Genos AI", SANS, 12, WHITE, bold=True)
tb(s12, Inches(7.75), Inches(6.0), Inches(4.5), Inches(0.28),
   "hello@genosai.tech  ·  +91 638-714-2699", SANS, 11, BODY)
tb(s12, Inches(7.75), Inches(6.32), Inches(4.5), Inches(0.28),
   "www.genosai.tech", SANS, 11, MUTED)

tb(s12, P, SH - Inches(0.42), Inches(8), Inches(0.28),
   "Confidential  ·  JamesEdition  ·  2026", SANS, 8, DIM)
tb(s12, Inches(9.5), SH - Inches(0.42), Inches(3.08), Inches(0.28),
   f"{TOTAL:02d} / {TOTAL}", SANS, 8, DIM, align=PP_ALIGN.RIGHT)

tb(s12, P, SH - Inches(0.65), Inches(3), Inches(0.4),
   "GenosAI", SERIF, 20, WHITE)


# ── Save ─────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved: {OUT}  ({TOTAL} slides)")
