"""
Build 24-7Security_Audit_GenosAI.pptx from scratch.
Design: docs/BRAND_ASSETS.md — #0A0A0F canvas, Instrument Serif headlines, Satoshi body.
12 slides: cover, audit, vision, systems overview, 5 modules, prospect flow,
implementation, closing.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).parent.parent.parent
OUT  = ROOT / "output" / "decks" / "24-7Security_Audit_GenosAI.pptx"

# ── Brand colors (BRAND_ASSETS.md §3) ───────────────────────────────────────
BG     = RGBColor(0x0A, 0x0A, 0x0F)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BODY   = RGBColor(0x9D, 0x9D, 0x9F)   # white/60
SUB    = RGBColor(0x86, 0x86, 0x8A)   # white/50
MUTED  = RGBColor(0x6B, 0x6B, 0x6E)   # white/40
DIM    = RGBColor(0x44, 0x44, 0x47)   # white/25
VIOLET = RGBColor(0xA7, 0x8B, 0xFA)   # eyebrow #A78BFA
RED    = RGBColor(0xF8, 0x71, 0x71)   # gap indicators
CARD   = RGBColor(0x11, 0x11, 0x18)
BORDER = RGBColor(0x22, 0x22, 0x2A)

SERIF = "Instrument Serif"
SANS  = "Satoshi"

SW    = Inches(13.33)
SH    = Inches(7.5)
P     = Inches(0.75)
TOTAL = 12


# ── Core helpers ─────────────────────────────────────────────────────────────

def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def tb(slide, l, t, w, h, text, font=SANS, size=14, color=BODY,
       bold=False, italic=False, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text        = text
    r.font.name   = font
    r.font.size   = Pt(size)
    r.font.color.rgb = color
    r.font.bold   = bold
    r.font.italic = italic
    return box


def tb_lines(slide, l, t, w, h, lines, font=SANS, size=14, color=BODY,
             bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name  = font
        r.font.size  = Pt(size)
        r.font.color.rgb = color
        r.font.bold  = bold
    return box


def card(slide, l, t, w, h, rounded=True):
    s = slide.shapes.add_shape(5 if rounded else 1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = CARD
    s.line.color.rgb = BORDER
    s.line.width = Pt(0.75)
    return s


def hline(slide, l, t, w, color=BORDER):
    s = slide.shapes.add_shape(1, l, t, w, Pt(0.75))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def eyebrow(slide, text, l, t, w=Inches(6)):
    return tb(slide, l, t, w, Inches(0.3),
              text.upper(), SANS, 10, VIOLET, bold=True)


def headline(slide, text, l, t, w=Inches(10), size=42):
    return tb(slide, l, t, w, Inches(1.6), text, SERIF, size, WHITE)


def slide_header(slide, section, page):
    tb(slide, P, Inches(0.22), Inches(6), Inches(0.28),
       "GENOSAI  ×  24/7 SECURITY SERVICES", SANS, 9, MUTED, bold=True)
    tb(slide, Inches(7.5), Inches(0.22), Inches(5.08), Inches(0.28),
       section, SANS, 9, MUTED, bold=True, align=PP_ALIGN.RIGHT)
    hline(slide, P, Inches(0.58), SW - 2 * P)
    tb(slide, P, SH - Inches(0.42), Inches(8), Inches(0.28),
       "AI Automation & Growth Proposal  ·  24/7 Security Services  ·  2026",
       SANS, 8, DIM)
    tb(slide, Inches(9.5), SH - Inches(0.42), Inches(3.08), Inches(0.28),
       f"{page:02d} / {TOTAL}", SANS, 8, DIM, align=PP_ALIGN.RIGHT)


def bullet(slide, dot_char, title, desc, l, t,
           col_w=Inches(5.0), title_size=13, desc_size=12):
    tb(slide, l, t, Inches(0.25), Inches(0.26),
       dot_char, SANS, 13, VIOLET, bold=True)
    tb(slide, l + Inches(0.3), t, col_w - Inches(0.3), Inches(0.26),
       title, SANS, title_size, WHITE, bold=True)
    tb(slide, l + Inches(0.3), t + Inches(0.27), col_w - Inches(0.3), Inches(0.38),
       desc, SANS, desc_size, BODY)
    return t + Inches(0.72)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH

s1 = new_slide(prs)

accent = s1.shapes.add_shape(1, Inches(0), Inches(0), SW, Pt(3))
accent.fill.solid()
accent.fill.fore_color.rgb = VIOLET
accent.line.fill.background()

tb(s1, P, Inches(1.0), Inches(5.5), Inches(0.35),
   "PREPARED FOR", SANS, 10, MUTED, bold=True)
hline(s1, P, Inches(1.42), Inches(5.5), BORDER)

tb(s1, P, Inches(1.6), Inches(8.5), Inches(1.4),
   "24/7 Security Services", SERIF, 62, WHITE)

tb(s1, P, Inches(3.12), Inches(5.5), Inches(0.32),
   "Armed Reaction  ·  Physical Guarding  ·  Drone Force  ·  South Africa",
   SANS, 13, BODY)

tb(s1, P, Inches(3.62), Inches(5.5), Inches(0.28),
   "Johannesburg  ·  Tshwane  ·  Cape Town", SANS, 11, MUTED)

hline(s1, P, Inches(4.08), Inches(5.5), BORDER)

tb(s1, P, Inches(4.28), Inches(6), Inches(0.28),
   "AI AUTOMATION  &  GROWTH PROPOSAL", SANS, 10, VIOLET, bold=True)
tb_lines(s1, P, Inches(4.68), Inches(6.5), Inches(1.2),
         ["An AI automation layer for 24/7 Security Services —",
          "built to capture prospect inquiries 24/7, automate guard",
          "operations, and deliver client reporting without manual effort."],
         SANS, 13, SUB)

card(s1, Inches(9.0), Inches(1.2), Inches(3.58), Inches(2.5), rounded=True)
tb(s1, Inches(9.25), Inches(1.5), Inches(3.0), Inches(0.28),
   "PREPARED BY", SANS, 9, MUTED, bold=True)
tb(s1, Inches(9.25), Inches(1.88), Inches(3.0), Inches(0.32),
   "GenosAI Tech", SERIF, 20, WHITE)
tb(s1, Inches(9.25), Inches(2.3), Inches(3.0), Inches(0.26),
   "Rohan Malik  ·  Founder & CEO", SANS, 11, BODY)
tb(s1, Inches(9.25), Inches(2.65), Inches(3.0), Inches(0.26),
   "hello@genosai.tech  ·  www.genosai.tech", SANS, 10, MUTED)

tb(s1, P, SH - Inches(0.65), Inches(4), Inches(0.45),
   "GenosAI", SERIF, 22, WHITE)
tb(s1, P + Inches(1.65), SH - Inches(0.62), Inches(3), Inches(0.3),
   "v1.0  ·  2026  ·  Confidential", SANS, 9, DIM)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AUDIT: WHERE THEY STAND TODAY
# ══════════════════════════════════════════════════════════════════════════════
s2 = new_slide(prs)
slide_header(s2, "01  ·  AUDIT", 2)

eyebrow(s2, "WHERE THE BUSINESS STANDS TODAY", P, Inches(0.85))
headline(s2, "Current Digital Presence", P, Inches(1.18), size=38)
tb(s2, P, Inches(2.2), Inches(8.5), Inches(0.42),
   "24/7 Security Services has strong operational credentials — SASA Gold, ISO 9001, Level 1 B-BBEE, "
   "30+ years in the market, and a Drone Force unit no competitor can easily replicate. "
   "The gap is on the digital front: prospects land on the website and find a contact form.",
   SANS, 13, SUB)

hline(s2, P, Inches(2.82), SW - 2 * P, BORDER)

CL = P
CR = Inches(7.3)
CW = Inches(5.5)

card(s2, CL, Inches(3.0), CW, Inches(3.65), rounded=True)
tb(s2, CL + Inches(0.3), Inches(3.15), CW - Inches(0.4), Inches(0.28),
   "WHAT'S ALREADY STRONG", SANS, 10, VIOLET, bold=True)
t = Inches(3.55)
for title, desc in [
    ("30+ years in market",
     "Deep operational track record — SASA Gold and ISO 9001 certified."),
    ("Multi-service platform",
     "Armed reaction, guarding, surveillance, drone force under one brand."),
    ("24/7 Drone Force",
     "Autonomous perimeter drones for residential estates — a clear differentiator."),
    ("Level 1 B-BBEE",
     "Preferred supplier status in public sector and corporate procurement."),
]:
    t = bullet(s2, "·", title, desc, CL + Inches(0.3), t, CW - Inches(0.45), 12, 11)

card(s2, CR, Inches(3.0), CW, Inches(3.65), rounded=True)
tb(s2, CR + Inches(0.3), Inches(3.15), CW - Inches(0.4), Inches(0.28),
   "ROOM TO ELEVATE", SANS, 10, RED, bold=True)
t = Inches(3.55)
for title, desc in [
    ("No WhatsApp channel",
     "Prospects can't reach you on the channel they use to compare providers."),
    ("Contact form only",
     "No live chat, no booking system — inquiries wait until someone reads email."),
    ("No online quote flow",
     "Prospects can't start a quote without speaking to a person first."),
    ("No client portal or reporting",
     "Monthly reports and incident summaries sent manually — or not at all."),
]:
    t = bullet(s2, "·", title, desc, CR + Inches(0.3), t, CW - Inches(0.45), 12, 11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — VISION
# ══════════════════════════════════════════════════════════════════════════════
s3 = new_slide(prs)
slide_header(s3, "02  ·  VISION", 3)

eyebrow(s3, "THE OPPORTUNITY", P, Inches(0.85))
headline(s3, "An AI Operations Layer\nfor 24/7 Security Services", P, Inches(1.18), size=38)

card(s3, P, Inches(2.85), Inches(7.5), Inches(1.25), rounded=True)
tb(s3, P + Inches(0.4), Inches(2.98), Inches(7.0), Inches(0.4),
   "“", SERIF, 32, VIOLET)
tb(s3, P + Inches(0.75), Inches(3.02), Inches(6.5), Inches(0.88),
   "We're not adding a chatbot. We're building the automation infrastructure that makes your "
   "operations — from prospect inquiry to monthly client report — run without manual intervention.",
   SERIF, 15, WHITE, italic=True)

col_w = Inches(2.8)
cols  = [P, P + col_w + Inches(0.15), P + 2*(col_w + Inches(0.15)), P + 3*(col_w + Inches(0.15))]
pillars = [
    ("01", "24/7 prospect\nresponse",
     "Every security quote inquiry answered on WhatsApp instantly."),
    ("02", "Automated\noperations",
     "Guard scheduling, attendance, and PSIRA renewals — off your plate."),
    ("03", "Client\ncommunication",
     "Monthly reports and incident updates delivered automatically."),
    ("04", "Lead\nconversion",
     "Inquiries that don't convert immediately are nurtured — not lost."),
]
for (num, title, desc), x in zip(pillars, cols):
    card(s3, x, Inches(4.35), col_w, Inches(2.65), rounded=True)
    tb(s3, x + Inches(0.25), Inches(4.55), Inches(0.6), Inches(0.42),
       num, SERIF, 22, VIOLET)
    tb(s3, x + Inches(0.25), Inches(5.05), col_w - Inches(0.4), Inches(0.65),
       title, SERIF, 16, WHITE)
    tb(s3, x + Inches(0.25), Inches(5.78), col_w - Inches(0.4), Inches(0.88),
       desc, SANS, 11, BODY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SYSTEMS OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
s4 = new_slide(prs)
slide_header(s4, "03  ·  AUTOMATION SYSTEMS", 4)

eyebrow(s4, "WHAT THE AUTOMATION LAYER BRINGS", P, Inches(0.85))
headline(s4, "Five AI Systems for 24/7 Security", P, Inches(1.18), size=36)

systems = [
    ("A", "WhatsApp Inquiry\nAgent",       "24/7 prospect response."),
    ("B", "Quote & Onboarding\nAutomation","From inquiry to contract."),
    ("C", "Guard Operations\nAI",          "Shift, attendance, PSIRA."),
    ("D", "Client Reporting\nAI",          "Monthly reports — automated."),
    ("E", "Lead Nurture\nSystem",          "Re-engage non-converting leads."),
]
sw_each = Inches(2.25)
gap     = Inches(0.17)
sx      = P

for letter, name, tagline in systems:
    card(s4, sx, Inches(2.6), sw_each, Inches(4.1), rounded=True)
    badge = s4.shapes.add_shape(1, sx + Inches(0.25), Inches(2.85), Inches(0.55), Inches(0.55))
    badge.fill.solid()
    badge.fill.fore_color.rgb = VIOLET
    badge.line.fill.background()
    tb(s4, sx + Inches(0.25), Inches(2.87), Inches(0.55), Inches(0.48),
       letter, SANS, 16, BG, bold=True, align=PP_ALIGN.CENTER)
    tb(s4, sx + Inches(0.2), Inches(3.55), sw_each - Inches(0.4), Inches(0.75),
       name, SERIF, 15, WHITE)
    tb(s4, sx + Inches(0.2), Inches(4.38), sw_each - Inches(0.4), Inches(0.55),
       tagline, SANS, 12, BODY)
    sx += sw_each + gap


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — MODULE A: WHATSAPP INQUIRY AGENT
# ══════════════════════════════════════════════════════════════════════════════
s5 = new_slide(prs)
slide_header(s5, "03A  ·  WHATSAPP INQUIRY AGENT", 5)

eyebrow(s5, "MODULE A", P, Inches(0.85))
headline(s5, "WhatsApp AI Inquiry Agent", P, Inches(1.18), size=36)
tb(s5, P, Inches(2.2), Inches(6.5), Inches(0.42),
   "Every prospect inquiry answered instantly — residential estate manager, commercial property owner, "
   "or individual homeowner. 24/7, no human required.",
   SANS, 14, SUB)

hline(s5, P, Inches(2.75), SW - 2 * P, BORDER)

t = Inches(3.0)
for title, desc in [
    ("24/7 instant response",
     "No inquiry goes unanswered — including Sunday night after a break-in scare."),
    ("Prospect qualification",
     "Property type, size, location, current security setup, risk level — captured automatically."),
    ("Service routing",
     "Residential to residential team, commercial to commercial. Drone Force to the right specialist."),
    ("Site assessment booking",
     "AI books the physical security assessment in the right team member's calendar."),
    ("Lead source tracking",
     "Every inquiry logged with source, time, and qualification outcome."),
    ("Competitor comparison handling",
     "AI trained to address common objections and position 24/7's certifications."),
]:
    t = bullet(s5, "·", title, desc, P, t, Inches(5.5), 13, 12)

card(s5, Inches(7.5), Inches(3.0), Inches(5.0), Inches(3.7), rounded=True)
tb(s5, Inches(7.8), Inches(3.2), Inches(4.5), Inches(0.28),
   "PROSPECT  ·  QUALIFICATION FLOW", SANS, 9, MUTED, bold=True)
tb(s5, Inches(7.8), Inches(3.6), Inches(4.5), Inches(0.42),
   "From inquiry", SERIF, 22, WHITE)
tb(s5, Inches(7.8), Inches(4.08), Inches(4.5), Inches(0.42),
   "to site assessment.", SERIF, 22, VIOLET)

for i, (label, value) in enumerate([
    ("EXPERIENCE", "30+ yrs"), ("CERT", "SASA Gold"), ("BBEE", "Level 1")
]):
    sx = Inches(7.8) + i * Inches(1.58)
    tb(s5, sx, Inches(5.0), Inches(1.4), Inches(0.28), label, SANS, 8, MUTED, bold=True)
    tb(s5, sx, Inches(5.32), Inches(1.4), Inches(0.42), value, SERIF, 18, WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — MODULE B & C: QUOTE AUTOMATION + GUARD OPS
# ══════════════════════════════════════════════════════════════════════════════
s6 = new_slide(prs)
slide_header(s6, "03B / 03C  ·  QUOTE & GUARD OPS", 6)

eyebrow(s6, "MODULE B", P, Inches(0.85))
headline(s6, "Quote & Client Onboarding", P, Inches(1.18), size=28)
tb(s6, P, Inches(2.05), Inches(5.5), Inches(0.38),
   "Prospect to signed contract — without manual data entry at any step.", SANS, 13, SUB)

t = Inches(2.6)
for title, desc in [
    ("Assessment input capture",
     "Site assessment notes fed into an AI that drafts the quote structure."),
    ("Quote template generation",
     "Ops team receives a pre-filled proposal — just reviews and sends."),
    ("Contract trigger",
     "On acceptance, onboarding flow activates automatically."),
    ("Deployment checklist",
     "Guard assignments, access codes, equipment — all tracked per client."),
    ("Handover documentation",
     "Automated client handover pack sent on first day of service."),
]:
    t = bullet(s6, "·", title, desc, P, t, Inches(5.5), 12, 11)

hline(s6, Inches(6.67) - Pt(0.5), Inches(0.85), Pt(1), BORDER)

CR = Inches(6.85)
eyebrow(s6, "MODULE C", CR, Inches(0.85))
headline(s6, "Guard Operations AI", CR, Inches(1.18), size=28)
tb(s6, CR, Inches(2.05), Inches(5.7), Inches(0.38),
   "Shift management, attendance, and PSIRA compliance — automated.", SANS, 13, SUB)

t = Inches(2.6)
for title, desc in [
    ("Daily shift confirmation",
     "Guards confirm shift via WhatsApp — exceptions flagged immediately."),
    ("Missed check-in alerts",
     "Site supervisor notified within minutes of a missed checkpoint."),
    ("PSIRA renewal tracking",
     "Each guard's certification expiry tracked; reminders at 90, 30, 7 days."),
    ("Incident logging",
     "Guards log incidents via WhatsApp — structured report generated automatically."),
    ("Overtime and leave tracking",
     "Hours logged per guard, leave requests captured in the same WhatsApp flow."),
]:
    t = bullet(s6, "·", title, desc, CR, t, Inches(5.7), 12, 11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — MODULE D & E: CLIENT REPORTING + LEAD NURTURE
# ══════════════════════════════════════════════════════════════════════════════
s7 = new_slide(prs)
slide_header(s7, "03D / 03E  ·  REPORTING & NURTURE", 7)

eyebrow(s7, "MODULE D", P, Inches(0.85))
headline(s7, "Client Reporting AI", P, Inches(1.18), size=28)
tb(s7, P, Inches(2.05), Inches(5.5), Inches(0.38),
   "Monthly security reports delivered to each client account — without your team writing them.",
   SANS, 13, SUB)

t = Inches(2.6)
for title, desc in [
    ("Monthly activity report",
     "Patrols, incidents, response times — auto-compiled and sent on the 1st."),
    ("Incident summary delivery",
     "Every logged incident summarised and delivered to the client contact."),
    ("Renewal reminder sequences",
     "Contract renewal and SLA review reminders at 90, 60, 30 days before expiry."),
    ("Client satisfaction check-in",
     "Quarterly satisfaction survey sent via WhatsApp — responses logged."),
    ("Reduces churn",
     "Clients who see value in writing cancel less — visible reporting builds retention."),
]:
    t = bullet(s7, "·", title, desc, P, t, Inches(5.5), 12, 11)

CR = Inches(6.85)
eyebrow(s7, "MODULE E", CR, Inches(0.85))
headline(s7, "Lead Nurture System", CR, Inches(1.18), size=28)
tb(s7, CR, Inches(2.05), Inches(5.7), Inches(0.38),
   "Prospects who inquired but didn't convert — reactivated automatically.",
   SANS, 13, SUB)

t = Inches(2.6)
for title, desc in [
    ("3-touch follow-up sequence",
     "Automated messages at 3, 7, and 14 days after inquiry with no response."),
    ("Seasonal re-engagement",
     "Residential prospects re-contacted before high-crime periods (Dec, April)."),
    ("Service update alerts",
     "Drone Force expansion or new service offerings pushed to past inquiries."),
    ("Reference case delivery",
     "AI sends a relevant client story — estate manager to estate manager."),
    ("Review and referral request",
     "Post-contract clients asked for a Google review and a referral contact."),
]:
    t = bullet(s7, "·", title, desc, CR, t, Inches(5.7), 12, 11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — MODULE E: LEAD FLOW (PROSPECT JOURNEY)
# ══════════════════════════════════════════════════════════════════════════════
s8 = new_slide(prs)
slide_header(s8, "04  ·  PROSPECT FLOW", 8)

eyebrow(s8, "INQUIRY  →  CONTRACT  →  ACTIVE CLIENT", P, Inches(0.85))
headline(s8, "Automated Prospect Journey", P, Inches(1.18), size=36)
tb(s8, P, Inches(2.2), Inches(10), Inches(0.38),
   "From the moment a prospect contacts 24/7 to their first active patrol — every step automated.",
   SANS, 13, SUB)

hline(s8, P, Inches(2.72), SW - 2 * P, BORDER)

steps = [
    ("01", "Inquiry\nReceived",      "Website form /\nWhatsApp / referral"),
    ("02", "WhatsApp AI\nResponds",  "Under 60 seconds"),
    ("03", "Qualification",          "Property, size,\nrisk, service need"),
    ("04", "Site Assessment\nBooked","Right team member\ncalendar confirmed"),
    ("05", "Quote\nGenerated",       "Auto-drafted from\nassessment inputs"),
    ("06", "Contract\nSigned",       "Digital acceptance\ntriggers onboarding"),
    ("07", "Deployment\nActive",     "Guard assigned,\nclient notified"),
    ("08", "Reporting\nLive",        "Monthly reports\nstart automatically"),
]

step_w   = Inches(1.4)
step_gap = Inches(0.12)
sx = P

for num, title, detail in steps:
    card(s8, sx, Inches(3.1), step_w, Inches(3.5), rounded=True)
    tb(s8, sx + Inches(0.12), Inches(3.28), step_w - Inches(0.2), Inches(0.38),
       num, SERIF, 20, VIOLET)
    tb(s8, sx + Inches(0.12), Inches(3.75), step_w - Inches(0.2), Inches(0.65),
       title, SERIF, 13, WHITE)
    tb(s8, sx + Inches(0.12), Inches(4.5), step_w - Inches(0.2), Inches(0.75),
       detail, SANS, 10, BODY)
    sx += step_w + step_gap
    if num != "08":
        tb(s8, sx - step_gap + Inches(0.02), Inches(4.2), step_gap + Inches(0.05), Inches(0.28),
           "›", SANS, 14, MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — BENEFITS
# ══════════════════════════════════════════════════════════════════════════════
s9 = new_slide(prs)
slide_header(s9, "05  ·  BENEFITS", 9)

eyebrow(s9, "WHY THIS WORKS FOR 24/7 SECURITY", P, Inches(0.85))
headline(s9, "What Changes After Launch", P, Inches(1.18), size=36)
hline(s9, P, Inches(2.72), SW - 2 * P, BORDER)

benefits = [
    ("Zero missed inquiries",         "Every security quote request gets an instant WhatsApp response — 24/7."),
    ("Faster deal cycles",             "Prospects go from first contact to booked site assessment in minutes."),
    ("Guard compliance on autopilot",  "PSIRA renewals tracked per guard — no admin overhead."),
    ("Client retention",               "Automatic monthly reports make your service value tangible."),
    ("No manual shift chasing",        "Attendance and check-ins handled via WhatsApp — exceptions surfaced instantly."),
    ("Scales without headcount",       "More clients and guards don't mean more admin staff."),
    ("Prospect re-engagement",         "Non-converting inquiries automatically followed up — not forgotten."),
    ("Differentiated positioning",     "AI-powered operations alongside Drone Force — a clear distance from competitors."),
]

col1, col2 = benefits[:4], benefits[4:]
CL, CR = P, Inches(7.1)
t1, t2 = Inches(2.95), Inches(2.95)

for title, desc in col1:
    card(s9, CL, t1, Inches(5.75), Inches(0.82), rounded=True)
    tb(s9, CL + Inches(0.25), t1 + Inches(0.1), Inches(0.35), Inches(0.3),
       "✓", SANS, 14, VIOLET, bold=True)
    tb(s9, CL + Inches(0.65), t1 + Inches(0.08), Inches(4.8), Inches(0.3),
       title, SANS, 13, WHITE, bold=True)
    tb(s9, CL + Inches(0.65), t1 + Inches(0.4), Inches(4.8), Inches(0.35),
       desc, SANS, 11, BODY)
    t1 += Inches(0.95)

for title, desc in col2:
    card(s9, CR, t2, Inches(5.75), Inches(0.82), rounded=True)
    tb(s9, CR + Inches(0.25), t2 + Inches(0.1), Inches(0.35), Inches(0.3),
       "✓", SANS, 14, VIOLET, bold=True)
    tb(s9, CR + Inches(0.65), t2 + Inches(0.08), Inches(4.8), Inches(0.3),
       title, SANS, 13, WHITE, bold=True)
    tb(s9, CR + Inches(0.65), t2 + Inches(0.4), Inches(4.8), Inches(0.35),
       desc, SANS, 11, BODY)
    t2 += Inches(0.95)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — FUTURE ROADMAP
# ══════════════════════════════════════════════════════════════════════════════
s10 = new_slide(prs)
slide_header(s10, "06  ·  ROADMAP", 10)

eyebrow(s10, "BEYOND THE LAUNCH", P, Inches(0.85))
headline(s10, "Future Scalability", P, Inches(1.18), size=36)
tb(s10, P, Inches(2.2), Inches(10), Inches(0.38),
   "The same automation infrastructure scales as 24/7 expands into new geographies and service lines.",
   SANS, 13, SUB)
hline(s10, P, Inches(2.72), SW - 2 * P, BORDER)

future = [
    ("Voice AI for inbound calls",    "AI answers calls, qualifies prospects, and routes to the right team."),
    ("Drone Force inquiry AI",        "Separate WhatsApp flow for residential estate drone inquiries."),
    ("Guard performance tracking",    "Response times and incident rates per guard — auto-reported to ops."),
    ("Multi-province expansion",      "Same AI infrastructure replicated across all three offices."),
    ("Subcontractor management AI",   "Automate briefings, shifts, and compliance for sub-contracted guards."),
    ("Emergency protocol AI",         "Automated escalation flows for armed response activation."),
    ("Client upsell sequences",       "Existing clients offered Drone Force upgrade based on area profile."),
    ("Tender and RFP automation",     "AI compiles tender documents from a structured proposal template."),
]

col1, col2 = future[:4], future[4:]
CL, CR = P, Inches(7.1)
t1, t2 = Inches(2.95), Inches(2.95)

for title, desc in col1:
    t1 = bullet(s10, "·", title, desc, CL, t1, Inches(5.75))
for title, desc in col2:
    t2 = bullet(s10, "·", title, desc, CR, t2, Inches(5.75))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — IMPLEMENTATION
# ══════════════════════════════════════════════════════════════════════════════
s11 = new_slide(prs)
slide_header(s11, "07  ·  IMPLEMENTATION", 11)

eyebrow(s11, "BUILT IN THREE CLEAN PHASES", P, Inches(0.85))
headline(s11, "Implementation Phases", P, Inches(1.18), size=36)
tb(s11, P, Inches(2.2), Inches(10), Inches(0.38),
   "Each phase ships visible value before the next one starts — no big-bang risk.",
   SANS, 13, SUB)
hline(s11, P, Inches(2.72), SW - 2 * P, BORDER)

phases = [
    ("PHASE 01  ·  Week 1",
     "WhatsApp Inquiry Agent",
     "WhatsApp AI trained on 24/7 services, certifications, and service areas. "
     "Prospect inquiry intake, qualification, and site assessment booking live.",
     ["AI trained on 24/7 service portfolio", "Prospect qualification flow",
      "Service routing (residential/commercial)", "Site assessment booking integration",
      "Lead logging and handover", "Go-live and testing"]),
    ("PHASE 02  ·  Weeks 2-3",
     "Guard Operations AI",
     "Shift confirmation, attendance tracking, PSIRA renewal reminders, and "
     "incident logging via WhatsApp — activated for all active guards.",
     ["WhatsApp shift confirmation flow", "Missed check-in alerting",
      "PSIRA expiry tracking per guard", "Incident logging via WhatsApp",
      "Ops supervisor dashboard", "Client reporting pipeline"]),
    ("PHASE 03  ·  Weeks 4-5",
     "Reporting & Lead Nurture",
     "Monthly client activity reports, renewal reminders, and lead nurture "
     "sequences for non-converting prospects — all automated.",
     ["Monthly report auto-compilation", "Client delivery flow",
      "Contract renewal reminders", "Lead nurture sequences",
      "Review and referral requests", "Training & handover"]),
]

ph_w  = Inches(3.8)
ph_gap = Inches(0.17)
px = P

for phase_label, phase_name, phase_desc, deliverables in phases:
    card(s11, px, Inches(3.0), ph_w, Inches(3.75), rounded=True)
    tb(s11, px + Inches(0.25), Inches(3.15), ph_w - Inches(0.4), Inches(0.28),
       phase_label, SANS, 9, VIOLET, bold=True)
    tb(s11, px + Inches(0.25), Inches(3.5), ph_w - Inches(0.4), Inches(0.42),
       phase_name, SERIF, 15, WHITE)
    tb(s11, px + Inches(0.25), Inches(4.0), ph_w - Inches(0.4), Inches(0.52),
       phase_desc, SANS, 11, BODY)
    t = Inches(4.62)
    for d in deliverables[:4]:
        tb(s11, px + Inches(0.25), t, ph_w - Inches(0.4), Inches(0.26),
           f"›  {d}", SANS, 10, MUTED)
        t += Inches(0.3)
    px += ph_w + ph_gap


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — CLOSING
# ══════════════════════════════════════════════════════════════════════════════
s12 = new_slide(prs)

acc = s12.shapes.add_shape(1, Inches(0), Inches(0), SW, Pt(3))
acc.fill.solid()
acc.fill.fore_color.rgb = VIOLET
acc.line.fill.background()

eyebrow(s12, "WHO IS BEHIND THIS", P, Inches(0.85))
headline(s12, "GenosAI Tech", P, Inches(1.18), size=40)
tb(s12, P, Inches(2.22), Inches(6.5), Inches(0.66),
   "GenosAI builds AI automation systems for security companies, managed service businesses, and "
   "operations-heavy industries. We specialise in WhatsApp AI, guard management automation, client "
   "reporting systems, and lead qualification flows — production-grade from week one.",
   SANS, 13, SUB)

hline(s12, P, Inches(3.02), Inches(6), BORDER)

tb(s12, P, Inches(3.2), Inches(4), Inches(0.28),
   "WHAT WE BUILD", SANS, 9, MUTED, bold=True)
t = Inches(3.58)
for item, desc in [
    ("WhatsApp AI Agents",         "Trained on your services, areas, and pricing."),
    ("Guard Operations Automation","Shift, attendance, PSIRA — off your plate."),
    ("Client Reporting AI",        "Monthly reports delivered without manual effort."),
    ("Lead Qualification Flows",   "Prospects qualified and routed automatically."),
    ("Lead Nurture Systems",       "Non-converting inquiries re-engaged over time."),
]:
    t = bullet(s12, "◆", item, desc, P, t, Inches(5.5), 12, 11)

card(s12, Inches(7.5), Inches(1.1), Inches(5.05), Inches(5.75), rounded=True)

tb(s12, Inches(7.75), Inches(1.35), Inches(4.6), Inches(0.35),
   "CLOSING NOTE", SANS, 9, MUTED, bold=True)
tb(s12, Inches(7.6), Inches(1.72), Inches(0.4), Inches(0.55),
   "“", SERIF, 28, VIOLET)
tb(s12, Inches(7.95), Inches(1.78), Inches(4.3), Inches(1.0),
   "Built to ensure 24/7 Security responds to every prospect first — and serves every client "
   "without a single manual report.",
   SERIF, 14, WHITE, italic=True)

hline(s12, Inches(7.75), Inches(2.88), Inches(4.45), BORDER)

tb(s12, Inches(7.75), Inches(3.08), Inches(4.5), Inches(0.28),
   "NEXT STEPS", SANS, 9, VIOLET, bold=True)
t = Inches(3.45)
for step, desc in [
    ("01", "20-minute call to confirm scope and integration points."),
    ("02", "Phase 1 kick-off — WhatsApp inquiry agent live within 7 days."),
    ("03", "First prospect qualified automatically — results visible from day one."),
]:
    tb(s12, Inches(7.75), t, Inches(0.4), Inches(0.3), step, SERIF, 14, VIOLET)
    tb(s12, Inches(8.2), t + Inches(0.02), Inches(4.15), Inches(0.5),
       desc, SANS, 12, BODY)
    t += Inches(0.62)

hline(s12, Inches(7.75), Inches(5.45), Inches(4.45), BORDER)

tb(s12, Inches(7.75), Inches(5.65), Inches(4.5), Inches(0.28),
   "Rohan Malik  ·  CEO, Genos AI", SANS, 12, WHITE, bold=True)
tb(s12, Inches(7.75), Inches(6.0), Inches(4.5), Inches(0.28),
   "hello@genosai.tech  ·  +91 638-714-2699", SANS, 11, BODY)
tb(s12, Inches(7.75), Inches(6.32), Inches(4.5), Inches(0.28),
   "www.genosai.tech", SANS, 11, MUTED)

tb(s12, P, SH - Inches(0.42), Inches(8), Inches(0.28),
   "Confidential  ·  24/7 Security Services  ·  2026", SANS, 8, DIM)
tb(s12, Inches(9.5), SH - Inches(0.42), Inches(3.08), Inches(0.28),
   f"{TOTAL:02d} / {TOTAL}", SANS, 8, DIM, align=PP_ALIGN.RIGHT)

tb(s12, P, SH - Inches(0.65), Inches(3), Inches(0.4),
   "GenosAI", SERIF, 20, WHITE)


# ── Save ─────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved: {OUT}  ({TOTAL} slides)")
