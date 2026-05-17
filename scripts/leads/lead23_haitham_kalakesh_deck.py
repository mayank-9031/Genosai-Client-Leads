"""
Build HaithamKalakesh_Audit_GenosAI.pptx from scratch.
Track: Dual Audit (Reve Real Estate 3/10 + Riversong Technology 6/10).
Design: docs/BRAND_ASSETS.md -- #0A0A0F canvas, Instrument Serif headlines, Satoshi body.
12 slides: cover, reve audit, riversong audit, vision, systems overview,
4 module slides (2 per company), cross-business module, implementation, closing.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).parent.parent.parent
OUT  = ROOT / "output" / "decks" / "HaithamKalakesh_Audit_GenosAI.pptx"

# -- Brand colors --------------------------------------------------------------
BG     = RGBColor(0x0A, 0x0A, 0x0F)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BODY   = RGBColor(0x9D, 0x9D, 0x9F)
SUB    = RGBColor(0x86, 0x86, 0x8A)
MUTED  = RGBColor(0x6B, 0x6B, 0x6E)
DIM    = RGBColor(0x44, 0x44, 0x47)
VIOLET = RGBColor(0xA7, 0x8B, 0xFA)
CARD   = RGBColor(0x11, 0x11, 0x18)
BORDER = RGBColor(0x22, 0x22, 0x2A)
RED    = RGBColor(0xF8, 0x71, 0x71)
AMBER  = RGBColor(0xFB, 0xBF, 0x24)

SERIF = "Instrument Serif"
SANS  = "Satoshi"

SW    = Inches(13.33)
SH    = Inches(7.5)
P     = Inches(0.75)
TOTAL = 12


# -- Helpers -------------------------------------------------------------------

def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def tb(slide, l, t, w, h, text, font=SANS, size=14, color=BODY,
       bold=False, italic=False, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text = text
    r.font.name   = font
    r.font.size   = Pt(size)
    r.font.color.rgb = color
    r.font.bold   = bold
    r.font.italic = italic
    return box


def tb_lines(slide, l, t, w, h, lines, font=SANS, size=14, color=BODY,
             bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name  = font
        r.font.size  = Pt(size)
        r.font.color.rgb = color
        r.font.bold  = bold
    return box


def card(slide, l, t, w, h, rounded=True):
    s = slide.shapes.add_shape(5 if rounded else 1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = CARD
    s.line.color.rgb = BORDER
    s.line.width = Pt(0.75)
    return s


def hline(slide, l, t, w, color=BORDER):
    s = slide.shapes.add_shape(1, l, t, w, Pt(0.75))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def eyebrow(slide, text, l, t, w=Inches(6)):
    return tb(slide, l, t, w, Inches(0.3),
              text.upper(), SANS, 10, VIOLET, bold=True)


def headline(slide, text, l, t, w=Inches(10), size=42):
    return tb(slide, l, t, w, Inches(1.6), text, SERIF, size, WHITE)


def slide_header(slide, section, page):
    tb(slide, P, Inches(0.22), Inches(6), Inches(0.28),
       "GENOSAI  x  HAITHAM KALAKESH", SANS, 9, MUTED, bold=True)
    tb(slide, Inches(7.5), Inches(0.22), Inches(5.08), Inches(0.28),
       section, SANS, 9, MUTED, bold=True, align=PP_ALIGN.RIGHT)
    hline(slide, P, Inches(0.58), SW - 2 * P)
    tb(slide, P, SH - Inches(0.42), Inches(8), Inches(0.28),
       "AI Audit & Automation Proposal  .  Reve Real Estate & Riversong Technology  .  2026",
       SANS, 8, DIM)
    tb(slide, Inches(9.5), SH - Inches(0.42), Inches(3.08), Inches(0.28),
       f"{page:02d} / {TOTAL}", SANS, 8, DIM, align=PP_ALIGN.RIGHT)


def bullet(slide, dot_char, title, desc, l, t,
           col_w=Inches(5.0), title_size=13, desc_size=12):
    tb(slide, l, t, Inches(0.25), Inches(0.26),
       dot_char, SANS, 13, VIOLET, bold=True)
    tb(slide, l + Inches(0.3), t, col_w - Inches(0.3), Inches(0.26),
       title, SANS, title_size, WHITE, bold=True)
    tb(slide, l + Inches(0.3), t + Inches(0.27), col_w - Inches(0.3), Inches(0.38),
       desc, SANS, desc_size, BODY)
    return t + Inches(0.72)


def score_badge(slide, score, l, t, color=RED):
    badge = slide.shapes.add_shape(5, l, t, Inches(1.1), Inches(1.1))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tb(slide, l, t + Inches(0.16), Inches(1.1), Inches(0.45),
       f"{score}/10", SERIF, 22, BG, bold=False, align=PP_ALIGN.CENTER)
    tb(slide, l, t + Inches(0.65), Inches(1.1), Inches(0.28),
       "SCORE", SANS, 8, BG, bold=True, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 1 -- COVER
# =============================================================================
prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH

s1 = new_slide(prs)

accent = s1.shapes.add_shape(1, Inches(0), Inches(0), SW, Pt(3))
accent.fill.solid()
accent.fill.fore_color.rgb = VIOLET
accent.line.fill.background()

tb(s1, P, Inches(1.0), Inches(6), Inches(0.35),
   "PREPARED FOR", SANS, 10, MUTED, bold=True)
hline(s1, P, Inches(1.42), Inches(6), BORDER)

tb(s1, P, Inches(1.6), Inches(9), Inches(1.4),
   "Haitham Kalakesh", SERIF, 60, WHITE)

tb(s1, P, Inches(2.95), Inches(7.5), Inches(0.32),
   "Founder & CEO  .  Reve Real Estate  .  Riversong Technology",
   SANS, 13, BODY)

tb(s1, P, Inches(3.42), Inches(6), Inches(0.28),
   "Dubai, UAE  .  London  .  Beirut  .  Amman  .  Cairo", SANS, 11, MUTED)

hline(s1, P, Inches(3.88), Inches(6), BORDER)

tb(s1, P, Inches(4.08), Inches(7), Inches(0.28),
   "AI AUDIT & AUTOMATION PROPOSAL", SANS, 10, VIOLET, bold=True)
tb_lines(s1, P, Inches(4.48), Inches(6.5), Inches(1.4),
         ["A full website audit and AI automation blueprint for",
          "two businesses -- Reve Real Estate and Riversong Technology --",
          "covering inquiry capture, lead nurture, support automation,",
          "and cross-business CRM unification."],
         SANS, 13, SUB)

card(s1, Inches(9.2), Inches(1.2), Inches(3.38), Inches(2.5), rounded=True)
tb(s1, Inches(9.45), Inches(1.5), Inches(2.8), Inches(0.28),
   "PREPARED BY", SANS, 9, MUTED, bold=True)
tb(s1, Inches(9.45), Inches(1.88), Inches(2.8), Inches(0.32),
   "GenosAI Tech", SERIF, 20, WHITE)
tb(s1, Inches(9.45), Inches(2.3), Inches(2.8), Inches(0.26),
   "Rohan Malik  .  Founder & CEO", SANS, 11, BODY)
tb(s1, Inches(9.45), Inches(2.65), Inches(2.8), Inches(0.26),
   "hello@genosai.tech  .  www.genosai.tech", SANS, 10, MUTED)

tb(s1, P, SH - Inches(0.65), Inches(4), Inches(0.45),
   "GenosAI", SERIF, 22, WHITE)
tb(s1, P + Inches(1.65), SH - Inches(0.62), Inches(3), Inches(0.3),
   "v1.0  .  2026  .  Confidential", SANS, 9, DIM)


# =============================================================================
# SLIDE 2 -- REVE REAL ESTATE WEBSITE AUDIT
# =============================================================================
s2 = new_slide(prs)
slide_header(s2, "01  .  REVE REAL ESTATE AUDIT", 2)

eyebrow(s2, "WEBSITE AUDIT  --  REVE REAL ESTATE", P, Inches(0.85))
headline(s2, "reverealestates.com  --  Score: 3 / 10", P, Inches(1.18), size=34)
tb(s2, P, Inches(2.2), Inches(9), Inches(0.38),
   "Dubai luxury real estate portal with offices in London, Beirut, and Amman. "
   "Developer partners include Emaar, Damac, Meraas, and Ellington. "
   "The brand brings serious traffic -- but the website loses almost all of it.",
   SANS, 13, SUB)

score_badge(s2, 3, Inches(11.6), Inches(1.2), RED)

hline(s2, P, Inches(2.72), SW - 2 * P, BORDER)

gaps = [
    ("No inquiry capture system",
     "No web chat, no WhatsApp button, no lead form above the fold. "
     "A buyer lands on the homepage and has no immediate way to engage."),
    ("No appointment booking",
     "No online scheduling for property viewings or consultation calls. "
     "Prospective clients must find a phone number or send a blind email."),
    ("No email capture or lead magnet",
     "No newsletter, no market report download, no Golden Visa guide. "
     "Every visitor who leaves is lost with no follow-up pathway."),
    ("No social proof or testimonials",
     "No client reviews, no transaction counts, no sold property showcases. "
     "International investors making $500K+ decisions need trust signals."),
    ("No Golden Visa conversion pathway",
     "Golden Visa consultation is listed as a service but has no dedicated "
     "funnel, no qualifying form, and no automated education sequence."),
    ("Contact information not visible",
     "Phone numbers and office addresses require navigation to a separate page. "
     "High-intent visitors in London or Beirut have no instant contact point."),
]

CL, CR = P, Inches(7.1)
t1, t2 = Inches(2.95), Inches(2.95)

for i, (title, desc) in enumerate(gaps):
    if i < 3:
        card(s2, CL, t1, Inches(5.75), Inches(0.9), rounded=True)
        tb(s2, CL + Inches(0.2), t1 + Inches(0.08), Inches(0.35), Inches(0.3),
           "!", SANS, 13, RED, bold=True)
        tb(s2, CL + Inches(0.6), t1 + Inches(0.07), Inches(4.9), Inches(0.3),
           title, SANS, 12, WHITE, bold=True)
        tb(s2, CL + Inches(0.6), t1 + Inches(0.4), Inches(4.9), Inches(0.42),
           desc, SANS, 10, BODY)
        t1 += Inches(1.02)
    else:
        card(s2, CR, t2, Inches(5.75), Inches(0.9), rounded=True)
        tb(s2, CR + Inches(0.2), t2 + Inches(0.08), Inches(0.35), Inches(0.3),
           "!", SANS, 13, RED, bold=True)
        tb(s2, CR + Inches(0.6), t2 + Inches(0.07), Inches(4.9), Inches(0.3),
           title, SANS, 12, WHITE, bold=True)
        tb(s2, CR + Inches(0.6), t2 + Inches(0.4), Inches(4.9), Inches(0.42),
           desc, SANS, 10, BODY)
        t2 += Inches(1.02)


# =============================================================================
# SLIDE 3 -- RIVERSONG TECHNOLOGY WEBSITE AUDIT
# =============================================================================
s3 = new_slide(prs)
slide_header(s3, "02  .  RIVERSONG TECHNOLOGY AUDIT", 3)

eyebrow(s3, "WEBSITE AUDIT  --  RIVERSONG TECHNOLOGY", P, Inches(0.85))
headline(s3, "riversongtechnology.com  --  Score: 6 / 10", P, Inches(1.18), size=34)
tb(s3, P, Inches(2.2), Inches(9), Inches(0.38),
   "Consumer electronics brand (smart watches, earbuds, phones, TVs, ACs) "
   "with distribution across Dubai, London, and Cairo. "
   "A functional website -- but the B2B and support layers are missing entirely.",
   SANS, 13, SUB)

score_badge(s3, 6, Inches(11.6), Inches(1.2), AMBER)

hline(s3, P, Inches(2.72), SW - 2 * P, BORDER)

gaps3 = [
    ("No distributor or retailer inquiry flow",
     "No B2B partnership form, no trade inquiry pathway, no qualification for "
     "new market entry. Distribution enquiries arrive cold with no routing."),
    ("No live chat or WhatsApp support",
     "Consumer electronics buyers expect instant support access. "
     "No chat widget, no WhatsApp link, no FAQ bot on product pages."),
    ("No customer support automation",
     "Warranty queries, return requests, and product support handled manually. "
     "With a range spanning 5+ categories, the inbound volume is significant."),
    ("No product review or social proof layer",
     "No customer reviews on product pages, no star ratings, no user-generated "
     "content. Buyers comparing against established brands need trust signals."),
    ("No post-purchase automation",
     "No automated warranty registration, no setup guide delivery, no review "
     "request. The customer relationship ends at purchase."),
    ("Limited B2B market entry pathway",
     "Riversong is expanding into new markets (Levant, North Africa). "
     "No partner onboarding flow or territory qualification process exists."),
]

CL, CR = P, Inches(7.1)
t1, t2 = Inches(2.95), Inches(2.95)

for i, (title, desc) in enumerate(gaps3):
    if i < 3:
        card(s3, CL, t1, Inches(5.75), Inches(0.9), rounded=True)
        tb(s3, CL + Inches(0.2), t1 + Inches(0.08), Inches(0.35), Inches(0.3),
           "!", SANS, 13, AMBER, bold=True)
        tb(s3, CL + Inches(0.6), t1 + Inches(0.07), Inches(4.9), Inches(0.3),
           title, SANS, 12, WHITE, bold=True)
        tb(s3, CL + Inches(0.6), t1 + Inches(0.4), Inches(4.9), Inches(0.42),
           desc, SANS, 10, BODY)
        t1 += Inches(1.02)
    else:
        card(s3, CR, t2, Inches(5.75), Inches(0.9), rounded=True)
        tb(s3, CR + Inches(0.2), t2 + Inches(0.08), Inches(0.35), Inches(0.3),
           "!", SANS, 13, AMBER, bold=True)
        tb(s3, CR + Inches(0.6), t2 + Inches(0.07), Inches(4.9), Inches(0.3),
           title, SANS, 12, WHITE, bold=True)
        tb(s3, CR + Inches(0.6), t2 + Inches(0.4), Inches(4.9), Inches(0.42),
           desc, SANS, 10, BODY)
        t2 += Inches(1.02)


# =============================================================================
# SLIDE 4 -- VISION
# =============================================================================
s4 = new_slide(prs)
slide_header(s4, "03  .  VISION", 4)

eyebrow(s4, "THE COMBINED OPPORTUNITY", P, Inches(0.85))
headline(s4, "One Automation Layer.\nTwo Businesses. Every Inquiry Captured.", P, Inches(1.18), size=34)

card(s4, P, Inches(2.85), Inches(7.5), Inches(1.2), rounded=True)
tb(s4, P + Inches(0.4), Inches(2.98), Inches(7.0), Inches(0.4),
   '"', SERIF, 32, VIOLET)
tb(s4, P + Inches(0.75), Inches(3.02), Inches(6.5), Inches(0.85),
   "Both businesses share the same problem: serious buyers and partners arrive, "
   "find no structured intake, and leave. AI automation closes that gap -- "
   "in real estate, in consumer electronics, and across both simultaneously.",
   SERIF, 15, WHITE, italic=True)

col_w = Inches(2.8)
cols  = [P, P + col_w + Inches(0.15), P + 2*(col_w + Inches(0.15)), P + 3*(col_w + Inches(0.15))]
pillars = [
    ("01", "Reve:\nInquiry Capture",
     "Every property buyer and investor qualified before reaching an agent."),
    ("02", "Reve:\nInvestor Nurture",
     "Leads who did not book a viewing followed up automatically for 6 weeks."),
    ("03", "Riversong:\nB2B Routing",
     "Distributor and retailer enquiries qualified and routed without manual work."),
    ("04", "Riversong:\nSupport AI",
     "Product support and warranty queries resolved by AI across all categories."),
]
for (num, title, desc), x in zip(pillars, cols):
    card(s4, x, Inches(4.35), col_w, Inches(2.65), rounded=True)
    tb(s4, x + Inches(0.25), Inches(4.55), Inches(0.6), Inches(0.42),
       num, SERIF, 22, VIOLET)
    tb(s4, x + Inches(0.25), Inches(5.05), col_w - Inches(0.4), Inches(0.65),
       title, SERIF, 16, WHITE)
    tb(s4, x + Inches(0.25), Inches(5.78), col_w - Inches(0.4), Inches(0.88),
       desc, SANS, 11, BODY)


# =============================================================================
# SLIDE 5 -- SYSTEMS OVERVIEW
# =============================================================================
s5 = new_slide(prs)
slide_header(s5, "04  .  AUTOMATION SYSTEMS", 5)

eyebrow(s5, "FIVE AI SYSTEMS ACROSS BOTH BUSINESSES", P, Inches(0.85))
headline(s5, "The Automation Blueprint", P, Inches(1.18), size=36)

systems = [
    ("A", "Reve:\nInquiry AI",         "WhatsApp & web chat lead capture."),
    ("B", "Reve:\nNurture & Visa AI",  "Investor follow-up and Golden Visa funnel."),
    ("C", "Riversong:\nB2B Inquiry AI","Distributor & retailer qualification."),
    ("D", "Riversong:\nSupport AI",    "Product support and post-purchase."),
    ("E", "Cross-Business:\nCRM AI",   "Unified pipeline and executive reporting."),
]
sw_each = Inches(2.25)
gap     = Inches(0.17)
sx      = P

for letter, name, tagline in systems:
    card(s5, sx, Inches(2.6), sw_each, Inches(4.1), rounded=True)
    badge = s5.shapes.add_shape(1, sx + Inches(0.25), Inches(2.85), Inches(0.55), Inches(0.55))
    badge.fill.solid()
    badge.fill.fore_color.rgb = VIOLET
    badge.line.fill.background()
    tb(s5, sx + Inches(0.25), Inches(2.87), Inches(0.55), Inches(0.48),
       letter, SANS, 16, BG, bold=True, align=PP_ALIGN.CENTER)
    tb(s5, sx + Inches(0.2), Inches(3.55), sw_each - Inches(0.4), Inches(0.75),
       name, SERIF, 15, WHITE)
    tb(s5, sx + Inches(0.2), Inches(4.38), sw_each - Inches(0.4), Inches(0.55),
       tagline, SANS, 12, BODY)
    sx += sw_each + gap


# =============================================================================
# SLIDE 6 -- MODULE A: REVE PROPERTY INQUIRY AI
# =============================================================================
s6 = new_slide(prs)
slide_header(s6, "04A  .  REVE INQUIRY AI", 6)

eyebrow(s6, "MODULE A  --  REVE REAL ESTATE", P, Inches(0.85))
headline(s6, "Property Inquiry AI", P, Inches(1.18), size=36)
tb(s6, P, Inches(2.2), Inches(6.5), Inches(0.42),
   "Every buyer or investor who lands on Reve Real Estate's website or WhatsApp "
   "enters a structured intake -- before a single agent is involved.",
   SANS, 14, SUB)

hline(s6, P, Inches(2.75), SW - 2 * P, BORDER)

t = Inches(3.0)
for title, desc in [
    ("WhatsApp inquiry agent",
     "A WhatsApp AI number displayed on the Reve website captures buyer intent: "
     "budget range, preferred developer (Emaar, Damac, Meraas), location, timeline."),
    ("Web chat lead intake",
     "An AI chat widget on every property listing page captures visitor details "
     "and property interest before the visitor exits."),
    ("Buyer qualification scoring",
     "Leads scored by budget, readiness, and property type. "
     "Qualified leads routed to an agent with full context; warm leads enter nurture."),
    ("Developer launch alerts",
     "Buyers who expressed interest in Emaar or Damac launches "
     "receive automated alerts when new projects in their criteria are released."),
    ("Multi-language intake",
     "Arabic, English, and French intake flows for the London, Beirut, and Dubai markets -- "
     "each routed to the appropriate regional team."),
]:
    t = bullet(s6, ".", title, desc, P, t, Inches(5.5), 13, 12)

card(s6, Inches(7.5), Inches(3.0), Inches(5.0), Inches(3.7), rounded=True)
tb(s6, Inches(7.8), Inches(3.2), Inches(4.5), Inches(0.28),
   "THE REVE WEBSITE GAP", SANS, 9, MUTED, bold=True)
tb(s6, Inches(7.8), Inches(3.6), Inches(4.5), Inches(0.55),
   "A buyer landing on reverealestates.com", SERIF, 16, WHITE)
tb(s6, Inches(7.8), Inches(4.18), Inches(4.5), Inches(0.55),
   "has no immediate way to inquire.", SERIF, 16, RED)
tb(s6, Inches(7.8), Inches(4.85), Inches(4.5), Inches(0.75),
   "No chat, no WhatsApp, no form above the fold. "
   "Every high-intent visitor who leaves is lost. "
   "The inquiry AI fixes this from week one.", SANS, 11, BODY)


# =============================================================================
# SLIDE 7 -- MODULE B: REVE INVESTOR NURTURE + GOLDEN VISA
# =============================================================================
s7 = new_slide(prs)
slide_header(s7, "04B  .  REVE NURTURE & GOLDEN VISA", 7)

eyebrow(s7, "MODULE B  --  REVE REAL ESTATE", P, Inches(0.85))
headline(s7, "Investor Nurture & Golden Visa Funnel", P, Inches(1.18), size=32)
tb(s7, P, Inches(2.15), Inches(6.5), Inches(0.42),
   "Most real estate deals close on the fourth or fifth follow-up. "
   "Most Golden Visa enquiries require education before conversion. "
   "Both need an automated system -- not a spreadsheet.",
   SANS, 13, SUB)

hline(s7, P, Inches(2.7), SW - 2 * P, BORDER)

CL, CR = P, Inches(6.85)

eyebrow(s7, "INVESTOR NURTURE SEQUENCES", CL, Inches(2.9))
t = Inches(3.2)
for title, desc in [
    ("Viewing no-show recovery",
     "Leads who booked a viewing but did not attend receive an automated "
     "re-engagement within 2 hours -- reschedule prompt with a different property."),
    ("Cold lead reactivation",
     "Enquiries older than 30 days with no response receive a new sequence: "
     "market update, price movement data, and a relevant project recommendation."),
    ("Quarterly investor update",
     "All leads in the pipeline receive a quarterly Dubai market data summary "
     "maintaining Reve's presence without agent effort."),
    ("Referral trigger",
     "Clients who completed a purchase receive an automated referral prompt "
     "at 60 days post-handover -- the highest NPS moment in the transaction."),
]:
    t = bullet(s7, ".", title, desc, CL, t, Inches(5.7), 12, 11)

eyebrow(s7, "GOLDEN VISA CONSULTATION FUNNEL", CR, Inches(2.9))
t = Inches(3.2)
for title, desc in [
    ("Qualification intake",
     "AI chat or form captures investment level, property type interest, "
     "nationality, and timeline for UAE residency by investment."),
    ("Automated education sequence",
     "Qualified leads receive a 4-step email sequence: "
     "eligibility, property threshold ($545K+), process timeline, and FAQs."),
    ("Call booking trigger",
     "Leads who complete the education sequence receive an automated "
     "invitation to book a 30-minute Golden Visa consultation call."),
    ("CRM handover",
     "Qualified leads transferred to the Reve team with full intake data, "
     "education sequence completion status, and investment readiness score."),
]:
    t = bullet(s7, ".", title, desc, CR, t, Inches(5.7), 12, 11)


# =============================================================================
# SLIDE 8 -- MODULE C: RIVERSONG B2B INQUIRY AI
# =============================================================================
s8 = new_slide(prs)
slide_header(s8, "04C  .  RIVERSONG B2B INQUIRY AI", 8)

eyebrow(s8, "MODULE C  --  RIVERSONG TECHNOLOGY", P, Inches(0.85))
headline(s8, "Distributor & B2B Inquiry AI", P, Inches(1.18), size=34)
tb(s8, P, Inches(2.2), Inches(6.5), Inches(0.42),
   "Riversong is expanding into the Levant, North Africa, and beyond. "
   "Every new distribution or retail partnership enquiry currently arrives cold. "
   "The B2B inquiry AI qualifies and routes them before any human time is spent.",
   SANS, 13, SUB)

hline(s8, P, Inches(2.75), SW - 2 * P, BORDER)

t = Inches(3.0)
for title, desc in [
    ("Territory and category qualification",
     "AI chat captures enquiring company type (distributor, retailer, e-commerce), "
     "target territory, relevant product categories, and estimated order volume."),
    ("Partnership tier routing",
     "Qualified leads routed to the appropriate Riversong regional contact. "
     "Levant enquiries go to Mutasem Shehadeh's pipeline; others routed by region."),
    ("Onboarding automation",
     "New approved partners receive a structured onboarding sequence: "
     "product catalogue, pricing tiers, marketing assets, and compliance documents."),
    ("Product training sequences",
     "New distributors and retailers receive automated product knowledge "
     "sequences for each category they carry -- reducing sales team load."),
    ("Partner performance check-ins",
     "Quarterly automated check-ins to active partners: sell-through data, "
     "new product announcements, co-marketing opportunities."),
]:
    t = bullet(s8, ".", title, desc, P, t, Inches(5.5), 13, 12)

card(s8, Inches(7.5), Inches(3.0), Inches(5.0), Inches(3.7), rounded=True)
tb(s8, Inches(7.8), Inches(3.2), Inches(4.5), Inches(0.28),
   "THE B2B GAP AT RIVERSONG", SANS, 9, MUTED, bold=True)
tb(s8, Inches(7.8), Inches(3.6), Inches(4.5), Inches(0.55),
   "Expanding into new markets requires", SERIF, 16, WHITE)
tb(s8, Inches(7.8), Inches(4.18), Inches(4.5), Inches(0.55),
   "a structured partner intake.", SERIF, 16, AMBER)
tb(s8, Inches(7.8), Inches(4.85), Inches(4.5), Inches(0.75),
   "Without automation, every new territory enquiry depends on "
   "a sales rep picking it up manually. "
   "The B2B AI ensures nothing is missed and every partner is qualified.", SANS, 11, BODY)


# =============================================================================
# SLIDE 9 -- MODULE D: RIVERSONG CUSTOMER SUPPORT AUTOMATION
# =============================================================================
s9 = new_slide(prs)
slide_header(s9, "04D  .  RIVERSONG SUPPORT AI", 9)

eyebrow(s9, "MODULE D  --  RIVERSONG TECHNOLOGY", P, Inches(0.85))
headline(s9, "Customer Support & Post-Purchase Automation", P, Inches(1.18), size=30)
tb(s9, P, Inches(2.15), Inches(6.5), Inches(0.42),
   "Smart watches, earbuds, phones, TVs, ACs -- five product categories each generate "
   "support queries that cannot be handled manually at scale.",
   SANS, 13, SUB)

hline(s9, P, Inches(2.7), SW - 2 * P, BORDER)

CL, CR = P, Inches(6.85)

eyebrow(s9, "CUSTOMER SUPPORT AUTOMATION", CL, Inches(2.9))
t = Inches(3.2)
for title, desc in [
    ("AI support chat on product pages",
     "A chat widget on each product page handles setup queries, "
     "compatibility questions, and feature explanations before escalation."),
    ("Warranty claim routing",
     "Customers submit warranty claims via AI chat. "
     "Serial number, purchase date, and fault description captured automatically."),
    ("Return and exchange flow",
     "AI guides customers through the return policy and "
     "initiates the exchange process without human intervention."),
    ("Escalation with context",
     "Complex queries escalated to the support team with full "
     "conversation history, product details, and purchase record attached."),
]:
    t = bullet(s9, ".", title, desc, CL, t, Inches(5.7), 12, 11)

eyebrow(s9, "POST-PURCHASE AUTOMATION", CR, Inches(2.9))
t = Inches(3.2)
for title, desc in [
    ("Setup guide delivery",
     "Immediately after purchase, customers receive an automated setup guide "
     "specific to the product they bought -- reducing support queries by category."),
    ("Warranty registration",
     "Automated warranty registration prompt at 48 hours post-purchase. "
     "Captures customer data and extends the relationship."),
    ("Review request sequence",
     "At day 14 post-purchase, automated review request sent to the "
     "platform where the purchase was made (Amazon, noon, website)."),
    ("Reorder and accessory prompts",
     "Earbuds buyers prompted for replacement tips at 6 months. "
     "Smart watch buyers prompted for straps and chargers at 3 months."),
]:
    t = bullet(s9, ".", title, desc, CR, t, Inches(5.7), 12, 11)


# =============================================================================
# SLIDE 10 -- MODULE E: CROSS-BUSINESS CRM & REPORTING
# =============================================================================
s10 = new_slide(prs)
slide_header(s10, "04E  .  CROSS-BUSINESS CRM AI", 10)

eyebrow(s10, "MODULE E  --  BOTH BUSINESSES", P, Inches(0.85))
headline(s10, "Unified CRM & Executive Reporting AI", P, Inches(1.18), size=32)
tb(s10, P, Inches(2.2), Inches(10), Inches(0.42),
   "Running two businesses across five markets from Dubai means two separate lead pipelines, "
   "two support queues, and zero unified visibility. "
   "The cross-business AI layer consolidates both into one operational view.",
   SANS, 13, SUB)

hline(s10, P, Inches(2.72), SW - 2 * P, BORDER)

features = [
    ("Unified CRM pipeline",
     "Reve Real Estate and Riversong leads, deals, and support tickets consolidated "
     "in a single CRM view -- one login for both businesses."),
    ("Weekly executive summary",
     "Every Monday, an automated report delivers: new leads by source, "
     "pipeline movement, support volume, and follow-up actions due this week."),
    ("Lead source attribution",
     "Every inquiry traced to its source -- WhatsApp, website, referral, or developer event. "
     "Haitham sees which channel drives qualified leads for each business."),
    ("Cross-business lead flagging",
     "A real estate buyer with investment capital may also be a Riversong B2B partner. "
     "The AI flags overlap and routes the opportunity appropriately."),
    ("Team performance reporting",
     "Agent and sales rep follow-up compliance tracked automatically -- "
     "no manual check-ins required to know who is active in the pipeline."),
    ("Monthly investor report delivery",
     "For Reve Real Estate, an automated monthly market update delivered "
     "to the full buyer and investor database -- built from CRM and market data."),
    ("Revenue forecasting",
     "Pipeline weighted by deal stage, deal size, and historical close rate -- "
     "a live revenue forecast for both businesses without manual spreadsheets."),
    ("Alerts and anomaly detection",
     "AI flags when lead volume drops sharply, when a deal stalls past its "
     "expected close date, or when a VIP contact has not been followed up."),
]

col1, col2 = features[:4], features[4:]
CL, CR = P, Inches(7.1)
t1, t2 = Inches(2.95), Inches(2.95)

for title, desc in col1:
    t1 = bullet(s10, ".", title, desc, CL, t1, Inches(5.75))
for title, desc in col2:
    t2 = bullet(s10, ".", title, desc, CR, t2, Inches(5.75))


# =============================================================================
# SLIDE 11 -- IMPLEMENTATION
# =============================================================================
s11 = new_slide(prs)
slide_header(s11, "05  .  IMPLEMENTATION", 11)

eyebrow(s11, "THREE PHASES ACROSS BOTH BUSINESSES", P, Inches(0.85))
headline(s11, "Implementation Roadmap", P, Inches(1.18), size=36)
tb(s11, P, Inches(2.2), Inches(10), Inches(0.38),
   "Reve Real Estate addressed first -- highest urgency (score 3/10). "
   "Riversong and cross-business systems follow in parallel.",
   SANS, 13, SUB)
hline(s11, P, Inches(2.72), SW - 2 * P, BORDER)

phases = [
    ("PHASE 01  .  Weeks 1-2",
     "Reve: Inquiry & Lead Capture",
     "WhatsApp inquiry AI, web chat lead intake, buyer qualification scoring, "
     "and developer launch alert system live on reverealestates.com.",
     ["WhatsApp AI number setup", "Web chat widget deployment",
      "Buyer qualification scoring model", "Developer alert sequence",
      "Multi-language flow (EN / AR / FR)", "CRM integration for Reve"]),
    ("PHASE 02  .  Weeks 3-4",
     "Reve: Nurture + Riversong: Support AI",
     "Investor nurture sequences and Golden Visa funnel for Reve; "
     "customer support AI, warranty routing, and post-purchase automation for Riversong.",
     ["Investor nurture sequence", "Golden Visa funnel build",
      "Riversong support chat widget", "Warranty and return flow",
      "Post-purchase sequences", "Review request automation"]),
    ("PHASE 03  .  Week 5",
     "Riversong B2B + Cross-Business CRM",
     "B2B distributor inquiry AI, partner onboarding sequences, "
     "and unified CRM with executive reporting across both businesses.",
     ["B2B inquiry qualification flow", "Partner onboarding sequence",
      "Unified CRM pipeline setup", "Weekly executive report",
      "Lead source attribution", "Revenue forecast automation"]),
]

ph_w   = Inches(3.8)
ph_gap = Inches(0.17)
px = P

for phase_label, phase_name, phase_desc, deliverables in phases:
    card(s11, px, Inches(3.0), ph_w, Inches(3.75), rounded=True)
    tb(s11, px + Inches(0.25), Inches(3.15), ph_w - Inches(0.4), Inches(0.28),
       phase_label, SANS, 9, VIOLET, bold=True)
    tb(s11, px + Inches(0.25), Inches(3.5), ph_w - Inches(0.4), Inches(0.42),
       phase_name, SERIF, 15, WHITE)
    tb(s11, px + Inches(0.25), Inches(4.0), ph_w - Inches(0.4), Inches(0.52),
       phase_desc, SANS, 11, BODY)
    t = Inches(4.62)
    for d in deliverables[:4]:
        tb(s11, px + Inches(0.25), t, ph_w - Inches(0.4), Inches(0.26),
           f">  {d}", SANS, 10, MUTED)
        t += Inches(0.3)
    px += ph_w + ph_gap


# =============================================================================
# SLIDE 12 -- CLOSING
# =============================================================================
s12 = new_slide(prs)

acc = s12.shapes.add_shape(1, Inches(0), Inches(0), SW, Pt(3))
acc.fill.solid()
acc.fill.fore_color.rgb = VIOLET
acc.line.fill.background()

eyebrow(s12, "WHO IS BEHIND THIS", P, Inches(0.85))
headline(s12, "GenosAI Tech", P, Inches(1.18), size=40)
tb(s12, P, Inches(2.22), Inches(6.5), Inches(0.66),
   "GenosAI builds AI automation systems for real estate groups, consumer electronics brands, "
   "and multi-market businesses. We build inquiry capture flows, investor nurture systems, "
   "B2B qualification AI, support automation, and executive reporting -- "
   "production-grade from week one.",
   SANS, 13, SUB)

hline(s12, P, Inches(3.02), Inches(6), BORDER)

tb(s12, P, Inches(3.2), Inches(4), Inches(0.28),
   "WHAT WE BUILD", SANS, 9, MUTED, bold=True)
t = Inches(3.58)
for item, desc in [
    ("Property Inquiry AI",       "WhatsApp and web chat intake for every buyer and investor."),
    ("Investor Nurture Systems",  "Automated follow-up sequences that close deals on the 4th touch."),
    ("B2B Qualification AI",      "Distributor and retailer enquiries routed without manual work."),
    ("Support Automation",        "Product support, warranty, and post-purchase handled by AI."),
    ("Unified CRM & Reporting",   "One pipeline view across both businesses for Haitham."),
]:
    t = bullet(s12, "o", item, desc, P, t, Inches(5.5), 12, 11)

card(s12, Inches(7.5), Inches(1.1), Inches(5.05), Inches(5.75), rounded=True)

tb(s12, Inches(7.75), Inches(1.35), Inches(4.6), Inches(0.35),
   "CLOSING NOTE", SANS, 9, MUTED, bold=True)
tb(s12, Inches(7.6), Inches(1.72), Inches(0.4), Inches(0.55),
   '"', SERIF, 28, VIOLET)
tb(s12, Inches(7.95), Inches(1.78), Inches(4.3), Inches(1.0),
   "Two businesses, five markets, one automation layer -- "
   "every inquiry captured, every investor followed up, "
   "every partner qualified, without adding headcount.",
   SERIF, 14, WHITE, italic=True)

hline(s12, Inches(7.75), Inches(2.88), Inches(4.45), BORDER)

tb(s12, Inches(7.75), Inches(3.08), Inches(4.5), Inches(0.28),
   "NEXT STEPS", SANS, 9, VIOLET, bold=True)
t = Inches(3.45)
for step, desc in [
    ("01", "20-minute call to walk through both audits and confirm scope."),
    ("02", "Phase 1 kick-off -- Reve inquiry AI live within 7 days."),
    ("03", "First captured inquiry routed to Reve agent -- visible from day one."),
]:
    tb(s12, Inches(7.75), t, Inches(0.4), Inches(0.3), step, SERIF, 14, VIOLET)
    tb(s12, Inches(8.2), t + Inches(0.02), Inches(4.15), Inches(0.5),
       desc, SANS, 12, BODY)
    t += Inches(0.62)

hline(s12, Inches(7.75), Inches(5.45), Inches(4.45), BORDER)

tb(s12, Inches(7.75), Inches(5.65), Inches(4.5), Inches(0.28),
   "Rohan Malik  .  CEO, Genos AI", SANS, 12, WHITE, bold=True)
tb(s12, Inches(7.75), Inches(6.0), Inches(4.5), Inches(0.28),
   "hello@genosai.tech  .  +91 638-714-2699", SANS, 11, BODY)
tb(s12, Inches(7.75), Inches(6.32), Inches(4.5), Inches(0.28),
   "www.genosai.tech", SANS, 11, MUTED)

tb(s12, P, SH - Inches(0.42), Inches(8), Inches(0.28),
   "Confidential  .  Reve Real Estate & Riversong Technology  .  2026", SANS, 8, DIM)
tb(s12, Inches(9.5), SH - Inches(0.42), Inches(3.08), Inches(0.28),
   f"{TOTAL:02d} / {TOTAL}", SANS, 8, DIM, align=PP_ALIGN.RIGHT)

tb(s12, P, SH - Inches(0.65), Inches(3), Inches(0.4),
   "GenosAI", SERIF, 20, WHITE)


# -- Save ----------------------------------------------------------------------
prs.save(OUT)
print(f"Saved: {OUT}  ({TOTAL} slides)")
