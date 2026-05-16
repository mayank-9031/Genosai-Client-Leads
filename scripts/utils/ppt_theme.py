"""
GenosAI PPT Design System
Reference design: clients/AquaTechnic/AquaTechnic_Audit_GenosAI [Repaired].pptx
Brand assets:    docs/BRAND_ASSETS.md

Import in every client PPT generator:

    import sys
    from pathlib import Path
    sys.path.insert(0, str(Path(__file__).parent.parent / "utils"))
    from ppt_theme import *
"""

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Color Palette ─────────────────────────────────────────────────────────────
# Canvas / surfaces
NAVY      = RGBColor(0x0A, 0x0F, 0x2E)   # primary slide background
DARK_BLUE = RGBColor(0x10, 0x1B, 0x4A)   # card / panel surface

# Per-slide accent colors — pick ONE per slide; use for left strip + eyebrow + dots
ACCENT    = RGBColor(0x4F, 0x8E, 0xFF)   # electric blue  (default)
GOLD      = RGBColor(0xF5, 0xC5, 0x18)   # amber / gold
RED       = RGBColor(0xFF, 0x4D, 0x4D)   # alert / critical gap
GREEN     = RGBColor(0x2E, 0xCC, 0x71)   # positive / completed
PURPLE    = RGBColor(0xA0, 0x55, 0xFF)   # design / website
ORANGE    = RGBColor(0xFF, 0x8C, 0x00)   # urgency / warning
TEAL      = RGBColor(0x00, 0xB4, 0xD8)   # info / secondary
CYAN      = RGBColor(0x00, 0xD4, 0xAA)   # primary CTA / WhatsApp

# Text
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)   # headlines, primary text
LGREY     = RGBColor(0xC8, 0xD4, 0xEE)  # body text on dark
MGREY     = RGBColor(0x8A, 0x9B, 0xBB)  # muted / secondary

# ── Slide Dimensions (16:9 widescreen) ───────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)


# ── Primitive Helpers ─────────────────────────────────────────────────────────

def bg(slide, color=NAVY):
    """Fill entire slide with a solid color."""
    s = slide.shapes.add_shape(1, 0, 0, W, H)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def rect(slide, l, t, w, h, fill):
    """Filled rectangle with no border."""
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def tb(slide, text, l, t, w, h, size=16, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """Single-run text box."""
    b = slide.shapes.add_textbox(l, t, w, h)
    tf = b.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return b


def aline(slide, t, color=ACCENT, w=Inches(1.2)):
    """Thin horizontal accent bar anchored to the left margin."""
    s = slide.shapes.add_shape(1, Inches(0.55), t, w, Inches(0.04))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()


def bullets(slide, items, l, t, w, h, size=13, color=LGREY, dot=ACCENT):
    """Bullet list: colored dot (●) + body text."""
    b = slide.shapes.add_textbox(l, t, w, h)
    tf = b.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r0 = p.add_run(); r0.text = "  "; r0.font.size = Pt(size - 2)
        r1 = p.add_run(); r1.text = "● "; r1.font.size = Pt(size - 2)
        r1.font.color.rgb = dot; r1.font.bold = True
        r2 = p.add_run(); r2.text = item
        r2.font.size = Pt(size); r2.font.color.rgb = color
        if i < len(items) - 1:
            p.space_after = Pt(6)


def tag(slide, text, l, t, fill=ACCENT, tc=WHITE, size=10):
    """Rounded-rectangle label badge."""
    w = Inches(len(text) * 0.105 + 0.4); h = Inches(0.3)
    s = slide.shapes.add_shape(9, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    tf = s.text_frame; p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = tc


# ── Composite Layout Helpers ──────────────────────────────────────────────────

def slide_header(slide, eyebrow, headline, accent_color=ACCENT,
                 aline_w=Inches(1.2), headline_size=28):
    """
    Standard slide header (every content slide):
      • 0.12" left accent strip in accent_color
      • 11pt UPPERCASE eyebrow label
      • Headline in white
      • Thin accent underline
    """
    rect(slide, 0, 0, Inches(0.12), H, accent_color)
    tb(slide, eyebrow.upper(), Inches(0.55), Inches(0.3), Inches(12), Inches(0.5),
       size=11, bold=True, color=accent_color)
    tb(slide, headline, Inches(0.55), Inches(0.7), Inches(12), Inches(0.7),
       size=headline_size, bold=True, color=WHITE)
    aline(slide, Inches(1.35), color=accent_color, w=aline_w)


def slide_footer(slide,
                 left_text="Rohan Malik  |  CEO, Genos AI  |  hello@genosai.tech  |  +91 638-714-2699",
                 right_text="May 2026  |  CONFIDENTIAL"):
    """Dark footer bar at bottom with contact info."""
    rect(slide, 0, Inches(6.5), W, Inches(1.0), DARK_BLUE)
    tb(slide, left_text, Inches(0.55), Inches(6.6), Inches(9), Inches(0.5),
       size=12, color=MGREY)
    tb(slide, right_text, Inches(9.5), Inches(6.6), Inches(3.3), Inches(0.5),
       size=12, color=MGREY, align=PP_ALIGN.RIGHT)


def score_badge(slide, score, label="AUDIT SCORE",
                priority_text="HIGH PRIORITY", priority_color=RED):
    """Audit score block for the cover slide (top-right quadrant)."""
    rect(slide, Inches(10.5), Inches(3.3), Inches(2.3), Inches(2.5), DARK_BLUE)
    tb(slide, label, Inches(10.55), Inches(3.4), Inches(2.2), Inches(0.5),
       size=11, bold=True, color=MGREY, align=PP_ALIGN.CENTER)
    tb(slide, str(score), Inches(10.55), Inches(3.75), Inches(2.2), Inches(1.1),
       size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(slide, "/ 10", Inches(10.55), Inches(4.7), Inches(2.2), Inches(0.4),
       size=18, color=LGREY, align=PP_ALIGN.CENTER)
    tag(slide, priority_text, Inches(10.75), Inches(5.1), fill=priority_color, size=10)


def stat_cards(slide, stats, top=Inches(1.8), card_w=Inches(2.3), card_h=Inches(2.2),
               gap=Inches(0.18), start=Inches(0.55), num_color=CYAN, label_color=LGREY):
    """
    Row of up to 5 stat cards (large number + label).
    stats: list of (number_str, label_str)
    """
    for i, (num, label) in enumerate(stats):
        l = start + i * (card_w + gap)
        rect(slide, l, top, card_w, card_h, DARK_BLUE)
        tb(slide, num, l + Inches(0.1), top + Inches(0.15),
           card_w - Inches(0.2), Inches(1.1),
           size=22, bold=True, color=num_color, align=PP_ALIGN.CENTER)
        tb(slide, label, l + Inches(0.1), top + Inches(1.25),
           card_w - Inches(0.2), Inches(0.7),
           size=11, color=label_color, align=PP_ALIGN.CENTER)


def top_bar_card(slide, l, t, w, h, accent_color, title, desc,
                 bar_h=Inches(0.06), title_size=13, desc_size=11):
    """
    Card with a thin colored top bar.
    Layout: [thin accent bar][DARK_BLUE body: bold title + grey description]
    """
    rect(slide, l, t, w, bar_h, accent_color)
    rect(slide, l, t + bar_h, w, h - bar_h, DARK_BLUE)
    tb(slide, title, l + Inches(0.12), t + Inches(0.18),
       w - Inches(0.24), Inches(0.65),
       size=title_size, bold=True, color=accent_color)
    tb(slide, desc, l + Inches(0.12), t + Inches(0.85),
       w - Inches(0.24), h - Inches(0.95),
       size=desc_size, color=LGREY)


def phase_card(slide, l, t, w, h, accent_color, phase_label, timing, items):
    """
    Implementation phase card.
    items: list of (title_str, desc_str)
    """
    rect(slide, l, t, w, Inches(0.52), accent_color)
    tb(slide, phase_label, l + Inches(0.12), t + Inches(0.04),
       w * 0.5 - Inches(0.12), Inches(0.42), size=13, bold=True, color=NAVY)
    tb(slide, timing, l + w * 0.45, t + Inches(0.08),
       w * 0.55 - Inches(0.15), Inches(0.38),
       size=11, color=NAVY, align=PP_ALIGN.RIGHT)
    rect(slide, l, t + Inches(0.52), w, h - Inches(0.52), DARK_BLUE)
    y = t + Inches(0.62)
    for title, desc in items:
        rect(slide, l + Inches(0.12), y + Inches(0.06), Inches(0.05), Inches(0.52), accent_color)
        tb(slide, title, l + Inches(0.25), y, w - Inches(0.38), Inches(0.32),
           size=12, bold=True, color=accent_color)
        tb(slide, desc, l + Inches(0.25), y + Inches(0.3), w - Inches(0.38), Inches(0.45),
           size=11, color=LGREY)
        y += Inches(0.95)


def before_after(slide, comparisons,
                 today_label="TODAY", after_label="AFTER 90 DAYS"):
    """
    Two-column before/after comparison table.
    comparisons: list of (topic_str, before_str, after_str)
    """
    lw = Inches(5.4)
    rect(slide, Inches(0.55), Inches(1.65), lw, Inches(0.35), RED)
    tb(slide, today_label, Inches(0.65), Inches(1.68), lw - Inches(0.2), Inches(0.28),
       size=11, bold=True, color=WHITE)
    rect(slide, Inches(6.45), Inches(1.65), lw + Inches(0.43), Inches(0.35), GREEN)
    tb(slide, after_label, Inches(6.55), Inches(1.68), lw, Inches(0.28),
       size=11, bold=True, color=WHITE)
    for i, (topic, before, after) in enumerate(comparisons):
        t = Inches(2.1) + i * Inches(1.02)
        rect(slide, Inches(0.55), t, Inches(12.23), Inches(0.22), DARK_BLUE)
        tb(slide, topic.upper(), Inches(0.65), t + Inches(0.02), Inches(12), Inches(0.2),
           size=9, bold=True, color=MGREY)
        tb(slide, before, Inches(0.65), t + Inches(0.24), lw - Inches(0.2), Inches(0.66),
           size=11, color=LGREY)
        tb(slide, after, Inches(6.55), t + Inches(0.24), lw, Inches(0.66),
           size=11, color=WHITE)


def genos_logo_block(slide, l=Inches(9.2), t_logo=Inches(0.4), t_url=Inches(1.05)):
    """Top-right GenosAI wordmark + URL (used on cover slides)."""
    rect(slide, Inches(8.5), 0, Inches(4.83), Inches(2.2), DARK_BLUE)
    tb(slide, "GENOS AI", l, t_logo, Inches(3.5), Inches(0.8),
       size=28, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
    tb(slide, "www.genosai.tech", l, t_url, Inches(3.5), Inches(0.5),
       size=12, color=MGREY, align=PP_ALIGN.RIGHT)
